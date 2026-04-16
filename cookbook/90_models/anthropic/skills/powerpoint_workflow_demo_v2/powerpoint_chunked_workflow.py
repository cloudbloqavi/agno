"""
Agno Workflow: Chunked PowerPoint Generation Pipeline.

A chunked workflow that overcomes Claude API limitations for large presentations
by splitting generation into manageable chunks, then merging the results.

Problem solved: Single Claude API calls fail for 10+ slide presentations;
               Claude PPTX skill is also prone to throttling and timeouts.
Solution: Generate slides in configurable chunks (default: 1 slide per call),
          then merge all chunks into one final presentation.
          A 3-tier fallback ensures production reliability when the primary
          Claude PPTX skill is unavailable or too slow.

Architecture:
  This file is a thin orchestration layer built on top of powerpoint_template_workflow.py.
  It imports all helpers, agents, Pydantic models, and step functions from that file via
  a wildcard import, then adds the chunked orchestration logic on top.

  powerpoint_template_workflow.py  — Core pipeline: content gen, images, template assembly,
                                    visual review, all helper functions (~6500 lines)
  powerpoint_chunked_workflow.py   — Chunked orchestration layer (~3200 lines, this file)

Chunk generation uses a 3-tier fallback hierarchy per chunk:
  Tier 1      Claude PPTX Skill         - Primary Haiku 4.5; on throttle/overload falls to
                                          Sonnet (PPTX skill), then Gemini/OpenAI code-gen
  Tier 2      LLM Code Gen              - Haiku → Gemini Pro → Gemini Flash
                                          → GPT-5.4 → o3-mini (code-gen chain)
  Tier 3      python-pptx Direct        - Last resort; text-only slides; 100% reliable

Brand/Style-Aware Query Parsing:
  Before the optimizer step, the workflow analyzes the user prompt for branding
  or styling intent using a dedicated brand_style_analyzer agent (default Claude Haiku 4.5
  but swappable to OpenAI or Gemini via --llm-provider). This agent:
    - Detects brand directives (e.g. "using Nike branding", "in the style of Apple")
    - Decides autonomously whether to search for brand guidelines (colors, tone, fonts)
    - Returns structured BrandStyleIntent (brand_name, color_palette, tone, typography)
  When a template file is provided, the workflow extracts styling from the template's
  theme XML (colors, fonts, company name heuristics) and overrides any query-level
  branding, with a descriptive [BRAND OVERRIDE] log explaining the decision.
  Brand context is propagated to Tier 1 and Tier 2 chunk prompts; Tier 3 is unaffected
  (no LLM call). Downstream steps (image pipeline, visual review, merge) are unchanged.

Relationship between the two files:
  - powerpoint_template_workflow.py acts strictly as a core foundational library 
    and cannot be run independently.
  - powerpoint_chunked_workflow.py wraps the same template/image/review logic via wildcard
    import so that large presentations (8-15+ slides) are split into chunks and merged.
  - Do NOT modify powerpoint_template_workflow.py to add chunking logic; keep them separate.

Workflow steps:
  Step 0  Brand/Style Parse  - (within Step 1) Detects brand intent, optionally searches
                                for brand guidelines, handles template override
  Step 1  Optimize & Plan    - LLM analyzes prompt, decides slide count, creates storyboard;
                                brand context is injected into the optimizer prompt and search
  Step 2  Generate Chunks    - Call Claude pptx skill (Tier 1) for each chunk;
                               auto-escalates to Tier 2 (Haiku code gen) on timeout/
                               failure, then Tier 3 (text-only) if Tier 2 fails.
                               Brand context is included in Tier 1 and Tier 2 prompts.
  Step 3  Process Chunks     - Apply template + image pipeline per chunk (if template provided)
  Step 4  Visual Review      - Optional per-chunk visual QA (if --visual-review + template)
  Step 5  Merge Chunks       - Merge all processed chunks into the final PPTX
  Step 6  Usage Summary      - (Verbosity specific) Displays total token consumption and 
                                estimated cost across all agents/providers


Key Models:
  BrandStyleIntent      - Structured brand/style data (name, colors, tone, fonts, style)
  SlideStoryboard       - Per-slide storyboard entry (title, type, key points, visual)
  StoryboardPlan        - Complete storyboard plan with global context and per-slide entries

Key Agents (all except Content Generator are swappable via --llm-provider):
  brand_style_analyzer  - Default: Claude Haiku 4.5 + web_search (max 2); detects and enriches brand intent.
                          OpenAI: gpt-5-mini + web_search_preview
                          Gemini: gemini-3-flash-preview + search=True
                          (Fallback: brand_style_analyzer_fallback uses a complementary provider model)
  query_optimizer       - Default: Claude Opus + web_search (max 5); creates researched storyboard.
                          OpenAI: gpt-5.2 + web_search_preview
                          Gemini: gemini-3.1-pro-preview + search=True
                          (Fallback: query_optimizer_fallback uses a complementary provider model)
  fallback_code_agent   - Default: Claude Sonnet + PythonTools; Tier 2 code generation primary.
                          OpenAI: gpt-5.2 + PythonTools
                          Gemini: gemini-3.1-pro-preview + PythonTools
  fallback_code_agent_lite - Default: Claude Haiku + PythonTools; Tier 2 fallback.
                          OpenAI: gpt-5-mini + PythonTools
                          Gemini: gemini-3-flash-preview + PythonTools
  image_planner         - Default: Gemini gemini-3-flash-preview; per-slide image decisions.
                          OpenAI: gpt-5-mini
                          (Fallback: image_planner_fallback uses a complementary provider model)
  slide_quality_reviewer- Default: Gemini 2.5 Flash; visual defect detection + correction.
                          OpenAI: gpt-5-mini
                          (Fallback: slide_quality_reviewer_fallback uses a complementary provider model)
  (LOCKED)              - chunk_agent + content_agent always use Claude (native PPTX skill dep.)

Prerequisites:
- uv pip install agno anthropic openai google-genai python-pptx pillow lxml python-dotenv \
    openinference-instrumentation-agno opentelemetry-exporter-otlp-proto-http opentelemetry-sdk
- export ANTHROPIC_API_KEY="your_anthropic_key"    # ALWAYS required (Content Generator)
- export OPENAI_API_KEY="your_openai_key"           # Required for brand analysis / --llm-provider openai
- export GOOGLE_API_KEY="your_google_key"           # Required for vision QA / --llm-provider gemini (also for image gen)
- export LANGFUSE_PUBLIC_KEY="..."                  # Optional: Observability
- export LANGFUSE_SECRET_KEY="..."
- export OTEL_EXPORTER_OTLP_ENDPOINT="..."
- A .pptx template file (optional)
- LibreOffice (required for --visual-review step: `sudo apt-get install -y libreoffice`)
- poppler-utils (required for per-slide PNG rendering: `sudo apt-get install -y poppler-utils`)

Template Quality Safeguards (active when --template is provided):
  - Per-slide rendering: PPTX→PDF→PNG pipeline via pdftoppm renders all slides for visual review
  - Background detection: 6-layer detection (shape→slide→layout→master→theme→large shapes)
  - Minimum font size: 10pt body / 14pt title floor prevents unreadable 4pt text from fit_text()
  - Layout sanitization: 3-pass boundary clamping, min size enforcement, and shape overlap reflow
  - Template-aware LLM prompts: Tier 2 code gen prompt includes template bg color/text color constraints
  - Template visual references: Renders template slides as PNGs and injects into chunk prompts for layout context

Usage:
    # 1. Basic (auto-decides slide count):
    python powerpoint_chunked_workflow.py \
        -p "AI in Healthcare" --chunk-size 1 --start-tier 2 --verbose

    # 2. Branded Executive Pitch (Live Research):
    python powerpoint_chunked_workflow.py \
        -p "Create a 7-slide presentation about AI trends using Nike branding" \
        --chunk-size 1 --start-tier 2 --no-images --verbose \
        --visual-review --visual-passes 3 --llm-provider openai

    # 3. Branded Corporate Template (Large Deck):
    python powerpoint_chunked_workflow.py \
        -t templates/my_template.pptx \
        -p "15-slide enterprise AI strategy for the board" \
        --chunk-size 1 --start-tier 2 --no-images --verbose \
        --visual-review --visual-passes 3 --llm-provider openai

    # 4. Use Gemini for all swappable agents + template:
    python powerpoint_chunked_workflow.py \
        -t templates/my_template.pptx \
        -p "10-slide AI strategy deck" --llm-provider gemini

    # 5. Skip directly to Tier 2 LLM code-gen (instant charts, high reliable):
    python powerpoint_chunked_workflow.py \
        -p "Quarterly review deck" --start-tier 2

CLI Flags:
    --template, -t       Path to .pptx template (optional). Without it, skips
                         template assembly and visual review; just merges raw chunks.
                         When provided, template styling overrides any query-level
                         branding with a [BRAND OVERRIDE] log.
    --output, -o         Output filename (default: presentation_chunked.pptx).
    --prompt, -p         User prompt describing the presentation.
                         Supports brand directives like "using X branding" or
                         "in the style of X".
    --llm-provider       LLM provider for all swappable agents (default: claude).
                         Choices: claude | openai | gemini
                         The Claude Content Generator agent is ALWAYS used regardless
                         of this flag, as it depends on Claude's native PPTX skill.
                         Requires the corresponding API key:
                           claude  → ANTHROPIC_API_KEY (always needed)
                           openai  → OPENAI_API_KEY
                           gemini  → GOOGLE_API_KEY
                         Examples:
                           --llm-provider openai   # Swaps auxiliary agents to GPT-5.2/gpt-5-mini
                           --llm-provider gemini   # Swaps auxiliary agents to gemini-3-pro/flash
    --no-images          Skip AI image generation.
    --no-stream          Disable streaming mode for Claude agent.
    --min-images         Minimum slides that must have images (default: 1).
    --visual-review      Enable visual QA with swappable vision agent per chunk.
    --footer-text        Footer text for all slides.
    --date-text          Date text for footer date placeholder.
    --show-slide-numbers Preserve slide number placeholder on all slides.
    --verbose, -v        Enable verbose/debug logging. Also triggers the final 
                         Token Usage & Cost Summary report at the end of the run.

    --chunk-size         Number of slides per LLM API chunk call (default: 1).
    --max-retries        Max retries per chunk on failure (default: 2).
    NOTE: When all retries fail or a timeout (300s) occurs, the system
          automatically switches to Tier 2 (LLM code gen) fallback,
          then Tier 3 (text-only) if Tier 2 also fails.
    --visual-passes      Maximum visual inspection passes per chunk (default: 3).
    --start-tier         Starting tier for chunk generation (default: 1):
                         1 = Claude PPTX skill (best quality, native charts/tables)
                         2 = LLM code generation (80-92% quality, faster, python-pptx native charts)
                         3 = Text-only (structural only, instant, no API calls)
                         Fallback chain continues from selected tier (e.g., Tier 2 → Tier 3).
    --inter-chunk-delay-min
                         Minimum delay in ms between chunks (default: provider-specific).
    --inter-chunk-delay-max
                         Maximum delay in ms between chunks (default: provider-specific).
    --no-web-search      Disable web search for query optimization.

Logging conventions:
    Always printed:
        [STEP_NAME] Starting / result messages
        [TIMING] step_XXX completed in X.Xs
        [ERROR] ...
        [WARNING] ...
        [VISUAL REVIEW MISSING FIX] ...  (always, per spec)
        [BRAND] Brand detection and extraction results
        [BRAND OVERRIDE] Template overriding query-level branding (with reason)
    Verbose-only (requires --verbose / -v):
        [VERBOSE] detailed debug information
        [TOKEN SUMMARY] Final table of token counts and USD cost per model

Usage Examples:
    # Basic (default 1 slide per chunk, tier 2 code-gen):
    python powerpoint_chunked_workflow.py -p "AI Trends" --start-tier 2

    # High-fidelity with template, visual review, and verbose logs:
    python powerpoint_chunked_workflow.py -t templates/corp.pptx \
        -p "Annual Report" --visual-review --verbose

    # OpenAI-powered auxiliary agents:
    python powerpoint_chunked_workflow.py -p "Market Analysis" --llm-provider openai
"""

import argparse
import asyncio
import concurrent.futures
import copy
import json
import os
import uuid

try:
    # Ensure environment variables are loaded and override existing exported variables
    from dotenv import load_dotenv  # type: ignore

    load_dotenv(override=True)
except ImportError:
    pass

import random
import shutil
import sys
import time
import traceback
import zipfile
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from agno.run.agent import RunOutput  # type: ignore

# === WILDCARD IMPORT: Reuse all helpers, agents, models, and step functions ===
# This gives us access to all ~6500 lines of helper logic without duplication.
# Specifically imports: SlideImageDecision, ImagePlan, ShapeIssue, SlideQualityReport,
# PresentationQualityReport, all dataclasses, image_planner, slide_quality_reviewer,
# step_plan_images, step_generate_images, step_assemble_template, step_visual_quality_review,
# and all _helper functions.
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from agno.agent import Agent  # type: ignore
from lib_patches.tools_python import PythonTools  # type: ignore
from agno.workflow.step import Step  # type: ignore
from agno.workflow.types import StepInput, StepOutput  # type: ignore
from agno.workflow.workflow import Workflow  # type: ignore
from anthropic import Anthropic  # type: ignore
from file_download_helper import download_skill_files  # type: ignore
from lib_patches.anthropic.claude import Claude  # type: ignore
from powerpoint_template_workflow import *  # type: ignore # noqa: F401, F403, E402
from powerpoint_template_workflow import (  # type: ignore
    clean_presentation_visual_noise_and_contrast,
    enforce_final_contrast,
    inject_template_footer_band,
    sanitize_llm_shapes,
    sanitize_presentation,
    step_assemble_template,
    step_generate_images,
    step_visual_quality_review,
)
from pptx import Presentation  # type: ignore
from pptx.dml.color import RGBColor  # type: ignore
from pptx.util import Inches, Pt  # type: ignore
from pydantic import BaseModel, Field  # type: ignore

# Default inter-chunk delays in milliseconds based on Tier 2 / Pay-as-you-go rate limits
# (End of imports)
DEFAULT_INTER_CHUNK_DELAYS_MS = {
    "claude": {"min": 2000, "max": 5000},  # 1K RPM, 450K TPM
    "openai": {"min": 1000, "max": 2000},  # 5K RPM, 2M TPM
    "gemini": {"min": 1000, "max": 2000},  # 1-2K RPM, multi-million TPM
}

# Global verbose flag (overridden in main() or by session_state)
VERBOSE = False

# === TELEMETRY SETUP (Langfuse via OpenInference) ===


def setup_langfuse_telemetry():
    """Setup Langfuse telemetry via OpenInference and OpenTelemetry following best practices."""
    public_key = os.getenv("LANGFUSE_PUBLIC_KEY")
    secret_key = os.getenv("LANGFUSE_SECRET_KEY")
    endpoint = os.getenv(
        "OTEL_EXPORTER_OTLP_ENDPOINT", "https://cloud.langfuse.com/api/public/otel"
    )

    if not public_key or not secret_key:
        if VERBOSE:
            print("[TELEMETRY] Langfuse keys missing. Observability disabled.")
        return None

    if not os.getenv("OTEL_EXPORTER_OTLP_HEADERS"):
        import base64

        auth = base64.b64encode(f"{public_key}:{secret_key}".encode()).decode()
        os.environ["OTEL_EXPORTER_OTLP_HEADERS"] = f"Authorization=Basic {auth}"

    if not os.getenv("OTEL_EXPORTER_OTLP_ENDPOINT"):
        os.environ["OTEL_EXPORTER_OTLP_ENDPOINT"] = endpoint

    try:
        if VERBOSE:
            import sys
        from openinference.instrumentation.agno import AgnoInstrumentor
        from opentelemetry.exporter.otlp.proto.http.trace_exporter import (
            OTLPSpanExporter,
        )
        from opentelemetry.sdk.resources import Resource
        from opentelemetry.sdk.trace import TracerProvider
        from opentelemetry.sdk.trace.export import BatchSpanProcessor

        # Best Practice: Define Service Resource
        resource = Resource(
            attributes={
                "service.name": "agno-pptx-workflow",
                "environment": os.getenv("APP_ENV", "development"),
            }
        )

        tracer_provider = TracerProvider(resource=resource)
        # Best Practice: BatchSpanProcessor for efficiency
        tracer_provider.add_span_processor(BatchSpanProcessor(OTLPSpanExporter()))

        # Instrument Agno
        AgnoInstrumentor().instrument(tracer_provider=tracer_provider)

        if VERBOSE:
            print(
                f"[TELEMETRY] Langfuse initialized (Service: agno-pptx-workflow, Endpoint: {os.environ['OTEL_EXPORTER_OTLP_ENDPOINT']})"
            )

        return tracer_provider
    except ImportError as e:
        if VERBOSE:
            print(f"[TELEMETRY] Telemetry packages missing ({e}). Tracing disabled.")
        return None


# === NEW PYDANTIC MODELS FOR CHUNKED WORKFLOW ===


# =============================================================================
# RATE LIMIT TRACKER
# Monitors all sequential Anthropic Claude API calls in this workflow run,
# estimates cumulative token usage within the rolling 60-second window, and
# auto-inserts cooldown sleeps before calls that would exceed the per-minute limit.
#
# Tier 1 rate limits (as of 2026-03, Anthropic):
#   claude-sonnet-4-6  : 50 RPM | 30,000 input tokens/min | 8,000 output tokens/min
#   claude-opus-4-6    : 50 RPM | 30,000 input tokens/min | 8,000 output tokens/min
#   claude-haiku-4-5   : 50 RPM | 50,000 input tokens/min | 10,000 output tokens/min
#   (gpt-4o-mini is OpenAI — entirely separate rate-limit pool, never tracked here)
# =============================================================================

# Per-model input-token-per-minute limits (Tier 1 defaults)
_ANTHROPIC_TOKEN_LIMITS: Dict[str, int] = {
    "claude-opus-4-6": 30_000,
    "claude-sonnet-4-6": 30_000,
    "claude-haiku-4-5": 50_000,
    "claude-haiku-3-5": 50_000,
    # fallback for unknown Claude models
    "default": 30_000,
}


class _RateLimitTracker:
    """Track sequential Anthropic Claude API calls and their estimated token usage.

    Maintains a rolling list of (timestamp, model, tokens) entries.  Before
    each registered call, ``check_and_wait`` removes entries older than 60 s
    and checks if adding the new call would exceed the model's per-minute limit.
    If so, it sleeps for the remaining time in the current window.

    Token estimation: len(prompt_chars) // 4  (~4 chars per token — conservative).
    The tracker is a module-level singleton initialised by ``_get_rate_tracker()``.
    """

    def __init__(self) -> None:
        self._calls: List[
            Dict
        ] = []  # list of {"ts": float, "model": str, "tokens": int}

    def _prune(self) -> None:
        """Remove entries older than 60 seconds from the rolling window."""
        cutoff = time.time() - 60.0
        self._calls = [c for c in self._calls if c["ts"] >= cutoff]

    def _window_tokens(self, model: str) -> int:
        """Sum tokens used by the given model in the current 60-second window."""
        return sum(c["tokens"] for c in self._calls if c["model"] == model)

    def _limit(self, model: str) -> int:
        """Return the input-token-per-minute limit for a given model."""
        for key, val in _ANTHROPIC_TOKEN_LIMITS.items():
            if key in model:
                return val
        return _ANTHROPIC_TOKEN_LIMITS["default"]

    def check_and_wait(self, model: str, prompt: str, caller: str = "") -> None:
        """Check token budget and sleep if needed, then log the call registration.

        Args:
            model:   Claude model ID string (e.g. 'claude-haiku-4-5').
            prompt:  The full prompt string (used to estimate token count).
            caller:  Human-readable name of the calling function (for logging).
        """
        self._prune()
        estimated = max(1, len(prompt) // 4)
        limit = self._limit(model)
        window_used = self._window_tokens(model)

        caller_tag = (" [%s]" % caller) if caller else ""
        print(
            "[RATE TRACKER]%s %s — ~%d estimated input tokens | "
            "window so far: ~%d / %d tokens/min"
            % (caller_tag, model, estimated, window_used, limit)
        )

        if window_used + estimated > limit:
            # Find oldest call within the window to determine when the window resets
            oldest_ts = min(
                (c["ts"] for c in self._calls if c["model"] == model),
                default=time.time(),
            )
            sleep_secs = max(1.0, 61.0 - (time.time() - oldest_ts))
            print(
                "[RATE TRACKER] Estimated token budget would be exceeded (%d + %d > %d). "
                "Sleeping %.0fs to reset the 60s window..."
                % (window_used, estimated, limit, sleep_secs)
            )
            _countdown_sleep(sleep_secs, label="[RATE TRACKER] Cooldown")
            self._prune()  # prune again after sleep

        # Register the call
        self._calls.append({"ts": time.time(), "model": model, "tokens": estimated})

    def record_done(self, model: str, prompt: str) -> None:
        """Alias — call after the API call completes to update the tracker log.

        In practice check_and_wait already registers the entry; this is a no-op
        convenience hook for future use (e.g. recording actual token usage from
        response headers).
        """
        pass  # Entry already recorded in check_and_wait


# Module-level singleton — shared across all step functions
_RATE_TRACKER: Optional["_RateLimitTracker"] = None


def _get_rate_tracker() -> "_RateLimitTracker":
    """Return (or create) the module-level rate-limit tracker singleton."""
    global _RATE_TRACKER
    if _RATE_TRACKER is None:
        _RATE_TRACKER = _RateLimitTracker()
    return _RATE_TRACKER


def _reset_rate_tracker() -> None:
    """Reset the rate tracker — call once at workflow start."""
    global _RATE_TRACKER
    _RATE_TRACKER = _RateLimitTracker()


# =============================================================================
# INTER-CHUNK DELAY HELPERS
# =============================================================================


def _countdown_sleep(seconds: float, label: str = "[GENERATE]") -> None:
    """Sleep for `seconds` with periodic log messages every 15 seconds.

    Args:
        seconds: Total number of seconds to sleep.
        label:   Prefix for log messages (default '[GENERATE]').
    """
    remaining = seconds
    while remaining > 0:
        tick = min(15.0, remaining)
        if remaining > 15:
            print(
                "%s Waiting... %.0fs remaining (%.0fs total)"
                % (label, remaining, seconds)
            )
        else:
            print("%s Final %.0fs..." % (label, remaining))
        time.sleep(tick)
        remaining -= tick


def _is_throttling_error(error_str: str) -> bool:
    """Detect any Anthropic throttling, capacity, or transient server error.

    Covers all known error variants that indicate the model is temporarily
    unavailable and warrants a fallback to another model:

      - HTTP 429  rate_limit_error   (account RPM/TPM quota exceeded)
      - HTTP 529  overloaded_error   (Anthropic server capacity exhausted)
      - Status 200 with 'overloaded' in body (mid-stream capacity failure)
      - HTTP 500  api_error          (internal server error / outage)
      - HTTP 402  billing_error      (credit exhausted — can't proceed)

    Args:
        error_str: The stringified exception from an API call.

    Returns:
        True if the error matches any known throttling/capacity pattern.
    """
    s = error_str.lower()
    return any(
        kw in s
        for kw in (
            "overloaded",
            "rate_limit",
            "rate limit",
            "429",
            "529",
            "billing_error",
            "402",
            "api_error",
        )
    )


def _inter_chunk_sleep(
    chunk_idx: int,
    total_chunks: int,
    min_delay: float = 2000.0,
    max_delay: float = 5000.0,
    rate_limit_hit: bool = False,
) -> None:
    """Sleep a random delay between chunks to respect provider rate limits.

    Uses a random value in [min_delay, max_delay] milliseconds. When a 429 rate-limit
    error was encountered in the current chunk, forces `max_delay` immediately.

    Logs a header line, a live countdown through `_countdown_sleep`, and a
    completion line so the user knows the process is not stuck.

    Args:
        chunk_idx:       0-based index of the chunk just completed.
        total_chunks:    Total number of chunks in the run.
        min_delay:       Lower bound of random delay range in ms (default 2000).
        max_delay:       Upper bound of random delay range in ms (default 5000).
        rate_limit_hit:  If True, skip the random choice and use max_delay.
    """
    if rate_limit_hit:
        delay = max_delay
        reason = "(max delay — rate limit hit)"
    else:
        delay = random.uniform(min_delay, max_delay)
        reason = "(rate limit safety jitter: %.0f–%.0fms range)" % (
            min_delay,
            max_delay,
        )

    delay_sec = delay / 1000.0

    print(
        "[GENERATE] --- Inter-chunk delay before Chunk %d/%d: %.1fs %s ---"
        % (chunk_idx + 2, total_chunks, delay_sec, reason)
    )
    _countdown_sleep(delay_sec, label="[GENERATE]")
    print(
        "[GENERATE] Inter-chunk delay complete. Resuming Chunk %d/%d."
        % (chunk_idx + 2, total_chunks)
    )


from agents._shared import (  # type: ignore
    LayoutConstraints,
    SlideStoryboard,
    StoryboardPlan,
)

# === TEMPLATE VISUAL PROFILE DATACLASSES ===
# Programmatic analysis of template layout characteristics, computed before
# storyboard generation so the optimizer can produce template-aware
# visual_suggestion and layout_constraints values.


@dataclass
class SlideLayoutProfile:
    """Visual characteristics of a single template slide layout.

    Attributes:
        slide_index: Position in analyzed template.
        slide_type_hint: Heuristic classification (title, section, content, blank).
        placeholder_count: Native placeholder count.
        content_zone_left_pct: Safe content region start (0-100).
        content_zone_top_pct: Safe content region start (0-100).
        content_zone_width_pct: Safe content region width (0-100).
        content_zone_height_pct: Safe content region height (0-100).
        decorative_shape_count: Count of non-placeholder visual elements.
        has_background_image: Boolean for picture/image background.
        has_gradient_background: Boolean for gradient background.
        text_box_count: Count of free-floating text boxes.
        total_shape_count: Total shape count on slide.
        usable_width_pct: Heuristic width after decorative interference.
        usable_height_pct: Heuristic height after decorative interference.
    """

    slide_index: int = 0
    slide_type_hint: str = "content"
    placeholder_count: int = 0
    content_zone_left_pct: float = 5.0
    content_zone_top_pct: float = 12.0
    content_zone_width_pct: float = 90.0
    content_zone_height_pct: float = 76.0
    decorative_shape_count: int = 0
    has_background_image: bool = False
    has_gradient_background: bool = False
    text_box_count: int = 0
    total_shape_count: int = 0
    usable_width_pct: float = 90.0
    usable_height_pct: float = 76.0
    accent_shapes: list = field(default_factory=list)


@dataclass
class TemplateVisualProfile:
    """Aggregated visual characteristics of an entire PPTX template.

    Attributes:
        aspect_ratio: Ratio (16:9, 4:3, etc).
        slide_count: Total slides analyzed.
        avg_placeholder_count: Mean placeholders per slide.
        avg_decorative_shapes: Mean decorative shapes per slide.
        avg_content_zone_width_pct: Mean horizontal content coverage.
        avg_content_zone_height_pct: Mean vertical content coverage.
        layout_density: Heuristic density (sparse, balanced, dense).
        dominant_layout_style: Heuristic style (sidebar, split, overlapping, full).
        max_comfortable_bullets: Suggester bullet limit.
        recommended_text_weight: Suggested text length (light, balanced, dense).
        has_charts_in_template: Presence of charts in any analyzed slide.
        has_tables_in_template: Presence of tables in any analyzed slide.
        has_images_in_template: Presence of images in any analyzed slide.
        has_smartart_shapes: Presence of SmartArt in any analyzed slide.
        slide_width_emu: EMU width.
        slide_height_emu: EMU height.
    """

    aspect_ratio: str = "16:9"
    slide_count: int = 0
    avg_placeholder_count: float = 2.0
    avg_decorative_shapes: float = 0.0
    avg_content_zone_width_pct: float = 90.0
    avg_content_zone_height_pct: float = 76.0
    layout_density: str = "balanced"
    dominant_layout_style: str = "full"
    max_comfortable_bullets: int = 5
    recommended_text_weight: str = "balanced"
    has_charts_in_template: bool = False
    has_tables_in_template: bool = False
    has_images_in_template: bool = False
    has_smartart_shapes: bool = False
    slide_width_emu: int = 9144000
    slide_height_emu: int = 5143500
    has_accent_lines: bool = False
    accent_pattern: dict = field(default_factory=dict)


# === BRAND/STYLE INTENT MODEL ===


class BrandStyleIntent(BaseModel):
    """Parsed branding/styling intent extracted from a user query or template file.

    Produced by the brand_style_analyzer agent (query-based) or by
    extract_style_from_template() (template-based).  Propagated through the
    workflow via session_state["brand_style_intent"] and injected into
    optimizer, Tier 1, and Tier 2 prompts.
    """

    has_branding: bool = Field(
        False,
        description=(
            "True when the user query contains an identifiable branding or styling "
            "intent (e.g. 'Nike branding', 'in the style of Apple').  False when "
            "the query is purely topical with no brand/style directive."
        ),
    )
    brand_name: str = Field(
        "",
        description="Brand name extracted from query (e.g. 'Nike', 'Tesla').",
    )
    style_keywords: List[str] = Field(
        default_factory=list,
        description=(
            "Style descriptors inferred or researched for the brand "
            "(e.g. ['bold', 'sporty', 'minimalist'])."
        ),
    )
    color_palette: List[str] = Field(
        default_factory=list,
        description=(
            "Specific color names or hex codes associated with the brand "
            "(e.g. ['#FF6600', 'black', 'white']).  Prefer hex when known."
        ),
    )
    tone_override: str = Field(
        "",
        description=(
            "Tone suggested by the brand identity "
            "(e.g. 'empowering', 'innovative', 'luxurious')."
        ),
    )
    typography_hints: List[str] = Field(
        default_factory=list,
        description=(
            "Font families or typographic style hints associated with the brand "
            "(e.g. ['Futura', 'Helvetica Neue', 'sans-serif bold'])."
        ),
    )
    content_query: str = Field(
        "",
        description=(
            "The user's original query with branding clauses removed, preserving "
            "only the content/topic portion.  Empty when has_branding is False."
        ),
    )
    source: str = Field(
        "query",
        description="'query' or 'template' — where the intent was derived.",
    )
    source_detail: str = Field(
        "",
        description=(
            "Human-readable detail about the source "
            "(e.g. template filename, or 'user query')."
        ),
    )
    brand_voice: str = Field(
        "",
        description="The brand's voice and personality (e.g. 'direct', 'playful', 'authoritative').",
    )
    target_audience: str = Field(
        "General",
        description=(
            "Primary audience (e.g. 'Potential clients', 'Internal team', "
            "'Industry peers'). Extracted per Rule 1 of RULES.md."
        ),
    )
    theme_definition: Optional[Dict[str, Any]] = Field(
        default=None,
        description=(
            "The autonomously selected or generated Theme metadata containing "
            "color palette (hex codes) and typography for Tier 2 extraction."
        ),
    )


# === MODULE-LEVEL AGENTS ===
# Swappable agents (brand_style_analyzer, query_optimizer, fallback_code_agent,
# image_planner, slide_quality_reviewer) are loaded from the agents/ package via
# get_agents(provider) lazily inside each step function.
# Agents are NOT stored in session_state to avoid deepcopy failures at workflow startup
# (Agent objects contain PythonTools which hold Python module references that cannot
# be pickled).  Instead session_state["llm_provider"] carries the provider name string,
# and each step calls get_agents() on first use.
# See agents/__init__.py for the factory.
#
# The Content Generator (chunk_agent) remains defined locally because it has a
# hard dependency on Claude's PPTX skill and cannot be swapped.


# === AGENT TRACEABILITY ===


def _log_agent_banner(
    agent_name: str,
    model_id: str,
    provider: str,
    step_name: str,
) -> None:
    """Print an always-on traceability banner before every agent invocation.

    Shown regardless of --verbose so the user always knows which agent/LLM
    is executing.  Format:
        ┌─────────────────────────────────────────────────
        │ 🤖 AGENT: Brand Style Analyzer
        │ 📡 MODEL: gpt-4o-mini [OpenAI]
        │ 📋 STEP:  step_optimize_and_plan / Brand Parse
        └─────────────────────────────────────────────────
    """
    line = "─" * 50
    print(
        "\n┌%s\n"
        "│ 🤖 AGENT: %s\n"
        "│ 📡 MODEL: %s [%s]\n"
        "│ 📋 STEP:  %s\n"
        "└%s" % (line, agent_name, model_id, provider, step_name, line)
    )


# === TEMPLATE VISUAL REFERENCE (Concern 1+6) ===
# Renders template slides to PNG and builds visual reference sections for
# chunk prompts, giving the LLM actual visual context of the template layout.


def _render_template_slides_to_png(
    template_path: str, output_dir: str
) -> Dict[int, str]:
    """Render all template slides to PNG using the LibreOffice→PDF→PNG pipeline.

    Uses the same pipeline as the visual review step. If LibreOffice or
    pdftoppm is unavailable, returns an empty dict (graceful degradation).

    PNGs are rendered at 72 DPI to produce ~80k chars of base64 per slide,
    keeping the LLM prompt well within context window limits for all models
    (including OpenAI fallback agents).

    Args:
        template_path: Path to the .pptx template file.
        output_dir:    Directory to write PNG files into (a 'template_pngs'
                       subdirectory is created automatically).

    Returns:
        Dict mapping 0-based slide index → PNG file path.
        Empty dict if rendering is unavailable or fails.
    """
    import glob
    import shutil as _shutil
    import subprocess

    png_dir = os.path.join(output_dir, "template_pngs")
    os.makedirs(png_dir, exist_ok=True)

    lo_cmd = _shutil.which("libreoffice") or _shutil.which("soffice")
    if not lo_cmd:
        if VERBOSE:
            print(
                "[VERBOSE] [TEMPLATE REF] LibreOffice not found — skipping template visual references. "
                "Install with: apt-get install libreoffice"
            )
        return {}

    pdftoppm_cmd = _shutil.which("pdftoppm")
    if not pdftoppm_cmd:
        if VERBOSE:
            print(
                "[VERBOSE] [TEMPLATE REF] pdftoppm not found — skipping template visual references. "
                "Install with: apt-get install poppler-utils"
            )
        return {}

    try:
        # Step 1: PPTX → PDF
        pdf_result = subprocess.run(
            [
                lo_cmd,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                png_dir,
                template_path,
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )
        if pdf_result.returncode != 0:
            if VERBOSE:
                err_str = str(pdf_result.stderr or "")[:200]  # type: ignore
                print(f"[VERBOSE] [TEMPLATE REF] PDF conversion failed: {err_str}")
            return {}

        base = os.path.splitext(os.path.basename(template_path))[0]
        pdf_path = os.path.join(png_dir, base + ".pdf")
        if not os.path.isfile(pdf_path):
            if VERBOSE:
                print("[VERBOSE] [TEMPLATE REF] PDF file not created.")
            return {}

        # Step 2: PDF → PNGs via pdftoppm (72 DPI for optimal context budget)
        if VERBOSE:
            print(
                "[VERBOSE] [PIPELINE] PPTX -> PDF -> PNG: Rendering per-slide placeholders at 72 DPI..."
            )
        png_prefix = os.path.join(png_dir, "tmpl")
        ppm_result = subprocess.run(
            [pdftoppm_cmd, "-png", "-r", "72", pdf_path, png_prefix],
            capture_output=True,
            text=True,
            timeout=120,
        )

        # Cleanup PDF
        try:
            os.remove(pdf_path)
        except OSError:
            pass

        if ppm_result.returncode != 0:
            if VERBOSE:
                err_str = str(ppm_result.stderr or "")[:200]  # type: ignore
                print(f"[VERBOSE] [TEMPLATE REF] pdftoppm failed: {err_str}")
            return {}

        pngs = sorted(glob.glob(os.path.join(png_dir, "tmpl-*.png")))
        result = {}
        for idx, png_path in enumerate(pngs):
            result[idx] = png_path

        if result:
            if VERBOSE:
                print(
                    "[VERBOSE] [TEMPLATE REF] Rendered %d template slide(s) as visual references."
                    % len(result)
                )
        return result

    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] [TEMPLATE REF] Template rendering failed: %s" % e)
        return {}


def _match_storyboard_to_template_slide(
    slide_type: str, template_pngs: Dict[int, str]
) -> Optional[str]:
    """Map a storyboard slide_type to the best-matching template slide PNG.

    Matching strategy:
        - 'title' / 'hero'    → template slide 0 (title slide)
        - 'closing'           → last template slide
        - 'section' / 'agenda' → template slide 2 (section header), or 0
        - 'data' / 'metrics'  → template slide 3 if exists, else 1
        - 'content' / default → template slide 1 (content layout), or 0

    Args:
        slide_type:     The semantic slide type from the storyboard.
        template_pngs:  Dict of {slide_index: png_path} from rendering.

    Returns:
        Path to the best-matching template PNG, or None if no match.
    """
    if not template_pngs:
        return None

    n = len(template_pngs)
    slide_type_lower = (slide_type or "content").lower().strip()

    def _get(idx: int) -> Optional[str]:
        return template_pngs.get(idx)

    if slide_type_lower in ("title", "hero"):
        return _get(0) or _get(min(template_pngs.keys()))
    elif slide_type_lower in ("closing",):
        return _get(n - 1) or _get(0)
    elif slide_type_lower in ("section", "agenda"):
        return _get(2) or _get(0)
    elif slide_type_lower in ("data", "metrics"):
        return _get(3) or _get(1) or _get(0)
    else:  # content, comparative, sequential, etc.
        return _get(1) or _get(0)


def _build_visual_reference_section(
    chunk_slides: List,
    template_pngs: Dict[int, str],
    brand_style_intent: Optional["BrandStyleIntent"] = None,
) -> str:
    """Build a markdown prompt section with base64-encoded template slide references.

    NOTE: This is only called if session_state["template_visuals"] is True.

    For each slide in the chunk, finds the best-matching template slide and
    encodes it as a base64 data URI for inclusion in the LLM prompt.

    With the default --chunk-size 1, this sends exactly one template image
    per chunk prompt (~80k chars at 72 DPI), matching the Manus/Claude Addon
    per-slide pattern for optimal template reproduction.

    When brand_style_intent is provided, a compact textual metadata block
    (theme colors, fonts) is prepended to give the LLM precise reproduction
    parameters alongside the visual reference.

    Args:
        chunk_slides:       List of SlideSpec objects from the storyboard.
        template_pngs:      Dict from _render_template_slides_to_png().
        brand_style_intent: Optional BrandStyleIntent with extracted theme
                            colors and fonts for textual metadata injection.

    Returns:
        Markdown section string (may be empty if no template or rendering unavailable).
    """
    import base64

    if not template_pngs:
        return ""

    sections = []
    used_refs = set()  # Deduplicate same template slide for multiple chunk slides

    for slide_spec in chunk_slides:
        slide_type = getattr(slide_spec, "slide_type", "content")
        slide_num = getattr(slide_spec, "slide_number", "?")
        ref_path = _match_storyboard_to_template_slide(slide_type, template_pngs)

        if ref_path and ref_path not in used_refs and os.path.isfile(ref_path):
            used_refs.add(ref_path)
            try:
                if VERBOSE:
                    print(
                        "[VERBOSE] [IMAGE] Encoding base64 reference for slide: %s"
                        % ref_path
                    )
                with open(ref_path, "rb") as f:
                    img_data = base64.b64encode(f.read()).decode("ascii")
                sections.append(
                    "### Template Reference for Slide %s (type: %s)\n"
                    "![Template slide](data:image/png;base64,%s)\n"
                    "Use this template as a VISUAL STYLE reference ONLY. "
                    "Replicate the color scheme, font styles, decorative shapes, "
                    "background design, and layout structure. "
                    "Do NOT copy or paraphrase ANY text content from the template image. "
                    "Your slide text must come ONLY from the storyboard content above.\n"
                    % (slide_num, slide_type, img_data)
                )
            except Exception:
                continue

    if not sections:
        return ""

    # Build optional theme metadata block for precise style reproduction
    theme_metadata = ""
    if brand_style_intent:
        meta_parts: list[str] = []
        if getattr(brand_style_intent, "color_palette", None):
            meta_parts.append(
                "- **Theme Colors:** %s"
                % ", ".join(list(brand_style_intent.color_palette))  # type: ignore
            )
        if getattr(brand_style_intent, "typography_hints", None):
            meta_parts.append(
                "- **Theme Fonts:** %s"
                % ", ".join(list(brand_style_intent.typography_hints))  # type: ignore
            )
        if getattr(brand_style_intent, "brand_name", None):
            meta_parts.append(
                "- **Company/Brand:** %s" % str(brand_style_intent.brand_name)  # type: ignore
            )
        if meta_parts:
            theme_metadata = (
                "\n### Template Theme Metadata\n"
                "Use these EXACT values when reproducing the template style:\n"
                + "\n".join(meta_parts)
                + "\n"
            )

    return (
        "\n## VISUAL REFERENCE — Template Slides\n"
        "The following images show what the template slides look like. "
        "Use them as visual reference for layout, positioning, and styling.\n"
        + theme_metadata
        + "\n"
        + "\n".join(sections)
    )


# === BRAND/STYLE HELPER FUNCTIONS ===


def _resolve_theme_with_factory(
    brand_intent: "BrandStyleIntent", agent_model
) -> "BrandStyleIntent":
    """
    Uses the Agno Skills Architecture to either select an existing theme or dynamically generate
    a custom theme if a template wasn't provided but branding is required.
    """
    try:
        import os

        from agno.agent import Agent  # type: ignore
        from agno.skills import LocalSkills, Skills  # type: ignore
        from pydantic import BaseModel, Field

        class ThemeDefinition(BaseModel):
            name: str = Field(..., description="Name of the theme.")
            source: str = Field(
                ...,
                description="Must be exactly 'predefined' if citing a theme from the folder, or 'custom' if you generated it from scratch.",
            )
            description: str = Field(..., description="Description of the theme usage.")
            palette: dict[str, str] = Field(
                ...,
                description="Dictionary mapping color names (e.g., 'accent1', 'dk1', 'lt1', 'lt2') to hex codes.",
            )
            typography: dict[str, str] = Field(
                ..., description="Dictionary with 'major' and 'minor' font names."
            )

        # Point LocalSkills to the directory containing 'theme-factory/SKILL.md'
        skills_dir = os.path.dirname(os.path.abspath(__file__))

        print(
            f"[BRAND] Activating Theme Factory to resolve presentation styling (fallback agent: {getattr(agent_model, 'id', 'unknown')})..."
        )

        theme_agent = Agent(
            name="Theme Selector Agent",
            model=agent_model,
            skills=Skills(loaders=[LocalSkills(skills_dir)]),
            instructions=[
                "You are an expert design director. Select or generate the best Theme for a presentation.",
                "Use the 'theme-factory' skill from your available skills.",
                "1. Read the instructions for the 'theme-factory' skill.",
                "2. Check available themes using get_skill_script with 'theme-factory' and 'scripts/list_themes.py'.",
                "3. Use get_skill_reference to read details of a specific theme if needed.",
                f"Extracted brand intent: {brand_intent.model_dump()}",
                "RULES:",
                "1. If the user explicitly asks for a specific preset theme OR if a preset matches the basic mood, output that strictly and set source='predefined'.",
                "2. If the brand is highly specific (e.g., real-world distinct brand colors) AND no preset is a perfect match, you MUST generate a new custom theme and set source='custom'.",
                "Your final output MUST be a valid JSON matching the ThemeDefinition schema.",
                "Always output valid JSON.",
            ],
            output_schema=ThemeDefinition,
            markdown=False,
        )

        response = theme_agent.run(
            "Resolve the best theme for the given brand intent.", stream=False
        )

        if (
            response
            and response.content
            and isinstance(response.content, ThemeDefinition)
        ):
            brand_intent.theme_definition = response.content.model_dump()

            source_type = brand_intent.theme_definition.get("source", "unknown").upper()
            theme_name = brand_intent.theme_definition.get("name", "Unknown")
            print(
                f"[BRAND] Theme Factory successfully resolved a [{source_type}] theme: {theme_name}"
            )

            # --- CRITICAL: Propagate theme palette back into brand_intent fields ---
            # This ensures _build_no_template_design_system and
            # _format_brand_context_for_prompt use the resolved theme hex codes
            # instead of the initial Brand Analyzer's generic color names.
            theme_palette = brand_intent.theme_definition.get("palette", {})
            if theme_palette:
                brand_intent.color_palette = list(theme_palette.values())
                print(
                    "[BRAND] Propagated theme palette → brand_intent.color_palette: %s"
                    % brand_intent.color_palette
                )
            theme_typo = brand_intent.theme_definition.get("typography", {})
            if theme_typo:
                brand_intent.typography_hints = list(theme_typo.values())
                print(
                    "[BRAND] Propagated theme typography → brand_intent.typography_hints: %s"
                    % brand_intent.typography_hints
                )

            try:
                import json

                # Always show a high-level summary of the decision
                print(
                    "[BRAND] Theme Palette Colors: %s"
                    % str(list(theme_palette.values()))
                )
                print("[BRAND] Theme Typography: %s" % str(list(theme_typo.values())))

                # If VERBOSE is enabled globally, dump the full structure
                if "VERBOSE" in globals() and globals()["VERBOSE"]:
                    print(
                        "[VERBOSE] [BRAND] Detailed Theme Metadata injected into layout prompt:"
                    )
                    print(
                        "[VERBOSE]\n%s"
                        % json.dumps(brand_intent.theme_definition, indent=2)
                    )
            except Exception:
                pass
        else:
            print("[BRAND] Theme Factory did not return a valid ThemeDefinition.")

    except Exception as e:
        print(f"[WARNING] Theme Factory resolution failed: {e}")

    return brand_intent


def parse_brand_style_intent(
    user_prompt: str, brand_agent: "Agent" = None, brand_agent_fallback: "Agent" = None
) -> "BrandStyleIntent":
    """Extract branding/styling intent from the user query via a two-stage approach.

    Stage 1 — Keyword pre-check (zero cost, zero tokens):
        Scans the prompt for explicit brand terms ('brand', 'logo', 'color scheme',
        'style guide') AND implicit style signals ('corporate look', 'professional
        feel', 'using X aesthetic', 'dark theme', 'minimalist', etc.).
        This step logs whether explicit intent was found.

    Stage 2 — gpt-4o-mini classification and extraction (always runs):
        Even if Stage 1 finds no explicit signals, the LLM is always invoked
        to catch implicit styling intent that keywords might have missed.
        Uses the brand_style_analyzer agent (configured as gpt-4o-mini) which
        runs on OpenAI's separate rate-limit pool and does NOT consume Anthropic
        input tokens. This preserves the claude-haiku-4-5 budget for chunk generation.

    Args:
        user_prompt: The raw user prompt string.
        brand_agent: The brand_style_analyzer Agent instance (gpt-4o-mini via OpenAI).

    Returns:
        BrandStyleIntent with extracted brand data, or BrandStyleIntent() if no
        branding intent was detected or the agent call failed).
    """
    import re as _re

    print("[BRAND] Analyzing query for branding/styling intent...")

    # === STAGE 1: KEYWORD PRE-CHECK (zero cost) ===
    # Covers explicit brand terms AND implicit style/look-and-feel signals.
    _BRAND_PATTERNS = [
        # Explicit branding directives
        r"\b(brand(?:ing)?)\b",
        r"\b(style\s+guide\b|brand\s+guide\b|brand\s+identity\b)",
        r"\b(logo|logotype|wordmark)\b",
        r"\b(color\s+sch(?:eme|emes?)|colour\s+sch(?:eme|emes?))\b",
        r"\b(corporate\s+(?:identity|colors?|colours?|palette|look))\b",
        r"\b(color\s+palette|colour\s+palette|color\s+theme|colour\s+theme)\b",
        r"\b(visual\s+identity|brand\s+voice|brand\s+tone)\b",
        # Implicit style / look-and-feel signals
        r"\b(corporate\s+(?:feel|design|style|aesthetic))\b",
        r"\b(professional\s+(?:look|feel|design|theme|style))\b",
        r"\b(minimalist|minimalistic|flat\s+design|clean\s+aesthetic)\b",
        r"\b(dark\s+(?:mode|theme)|light\s+(?:mode|theme))\b",
        r"\b(modern\s+(?:design|look|feel|style|aesthetic))\b",
        r"\b(using\s+\w+(?:\s+\w+)?\s+(?:branding|style|colors?|colours?|identity|aesthetic|look))\b",
        r"\b(in\s+(?:the\s+)?style\s+of\b)",
        r"\b(following\s+\w+(?:\s+\w+)?\s+(?:brand\s+)?guidelines)\b",
        r"\b(match(?:ing)?\s+(?:the|our|their|its)\s+(?:brand|style|identity|theme))\b",
        r"\b(typography|typeface|font\s+(?:family|choice|style))\b",
        r"\b(accent\s+color|primary\s+color|secondary\s+color|hex\s+code)\b",
    ]

    prompt_lower = user_prompt.lower()
    has_brand_signal = any(
        _re.search(pattern, prompt_lower) for pattern in _BRAND_PATTERNS
    )

    if not has_brand_signal:
        print(
            "[BRAND] No explicit branding keywords detected, but analyzing prompt "
            "with LLM (gpt-4o-mini) to check for implicit styling intent..."
        )
    else:
        print(
            "[BRAND] Brand/style signal detected in query — "
            "calling gpt-4o-mini (OpenAI, off Anthropic quota)..."
        )

    # === STAGE 2: LLM extraction via gpt-4o-mini (OpenAI rate-limit pool) ===
    if brand_agent is None:
        print("[WARNING] Brand agent not provided; skipping LLM brand analysis.")
        return BrandStyleIntent()

    try:
        _log_agent_banner(
            agent_name=getattr(brand_agent, "name", "Brand Style Analyzer"),
            model_id=getattr(getattr(brand_agent, "model", None), "id", "unknown"),
            provider=getattr(getattr(brand_agent, "model", None), "provider", "OpenAI"),
            step_name="step_optimize_and_plan / Brand Parse",
        )
        response = brand_agent.run(user_prompt, stream=False)

        if response and response.content:
            content = response.content
            if isinstance(content, BrandStyleIntent):
                intent = content
            elif isinstance(content, dict):
                intent = BrandStyleIntent(**content)
            elif isinstance(content, BaseModel):
                intent = BrandStyleIntent(**content.model_dump())
            else:
                # Try JSON parse from string
                text = str(content).strip()
                fence = _re.search(r"```(?:json)?\s*([\s\S]+?)```", text)
                if fence:
                    text = fence.group(1).strip()
                obj = _re.search(r"\{[\s\S]+\}", text)
                if obj:
                    text = obj.group(0)
                intent = BrandStyleIntent.model_validate_json(text)

            if intent.has_branding:
                print(
                    "[BRAND] Detected brand intent: '%s' | style: %s | colors: %s"
                    % (
                        intent.brand_name,
                        list(intent.style_keywords)[:3],  # type: ignore
                        list(intent.color_palette)[:4],  # type: ignore
                    )
                )
                if intent.tone_override:
                    print("[BRAND] Tone override: '%s'" % intent.tone_override)
            else:
                print("[BRAND] No branding intent confirmed by primary agent.")

            return intent

    except Exception as e:
        print("[WARNING] Primary brand style analysis failed: %s" % str(e))
        if brand_agent_fallback is not None:
            print("[BRAND] Attempting fallback brand style analysis...")
            try:
                _log_agent_banner(
                    agent_name=getattr(
                        brand_agent_fallback, "name", "Brand Style Analyzer (Fallback)"
                    ),
                    model_id=getattr(
                        getattr(brand_agent_fallback, "model", None), "id", "unknown"
                    ),
                    provider=getattr(
                        getattr(brand_agent_fallback, "model", None),
                        "provider",
                        "Fallback",
                    ),
                    step_name="step_optimize_and_plan / Brand Parse (Fallback)",
                )
                response = brand_agent_fallback.run(user_prompt, stream=False)

                if response and response.content:
                    content = response.content
                    if isinstance(content, BrandStyleIntent):
                        intent = content
                    elif isinstance(content, dict):
                        intent = BrandStyleIntent(**content)
                    elif isinstance(content, BaseModel):
                        intent = BrandStyleIntent(**content.model_dump())
                    else:
                        text = str(content).strip()
                        fence = _re.search(r"```(?:json)?\s*([\s\S]+?)```", text)
                        if fence:
                            text = fence.group(1).strip()
                        obj = _re.search(r"\{[\s\S]+\}", text)
                        if obj:
                            text = obj.group(0)
                        intent = BrandStyleIntent.model_validate_json(text)
                    return intent
            except Exception as fallback_e:
                print(
                    "[WARNING] Fallback brand style analysis failed: %s"
                    % str(fallback_e)
                )

        if VERBOSE:  # noqa: F405
            traceback.print_exc()

    return BrandStyleIntent()


def extract_style_from_template(template_path: str) -> "BrandStyleIntent":
    """Extract branding/styling information from a .pptx template file.

    Reads the template's theme XML to extract:
    - Color palette (theme colors as hex values)
    - Font scheme (major/minor theme fonts)
    - Company name heuristics (from title slide placeholder text)

    This is a structural read — it captures color palettes and font families
    but does not perform OCR or pixel-level analysis.

    Args:
        template_path: Absolute or relative path to a .pptx template file.

    Returns:
        BrandStyleIntent with source='template' and extracted styling data.
        Returns an empty BrandStyleIntent on any read error.
    """
    print("[BRAND] Extracting style from template: %s" % template_path)

    intent = BrandStyleIntent(  # type: ignore
        source="template",
        source_detail=str(os.path.basename(template_path)),  # type: ignore
        has_branding=True,  # type: ignore
    )

    try:
        prs = Presentation(template_path)

        # --- Extract theme colors ---
        colors = []
        try:
            theme_el = prs.slide_masters[0].element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}theme"
            )
            if theme_el is None:
                # Try via the slide master's part
                theme_part = prs.slide_masters[0].part.slide_master.element
                ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                # Search for clrScheme in the theme
                for clr_scheme in theme_part.iter(ns_a + "clrScheme"):
                    for child in clr_scheme:
                        tag = (
                            child.tag.split("}")[-1] if "}" in child.tag else child.tag
                        )
                        if tag in (
                            "dk1",
                            "dk2",
                            "lt1",
                            "lt2",
                            "accent1",
                            "accent2",
                            "accent3",
                            "accent4",
                            "accent5",
                            "accent6",
                        ):
                            for color_el in child:
                                val = color_el.get("val", "")
                                last_clr = color_el.get("lastClr", "")
                                hex_val = val if len(val) == 6 else last_clr
                                if hex_val and len(hex_val) == 6:
                                    colors.append("#%s" % hex_val.upper())
        except Exception as e:
            if VERBOSE:  # noqa: F405
                print("[VERBOSE] Theme color extraction error: %s" % e)

        if colors:
            intent.color_palette = list(dict.fromkeys(colors))[
                :8
            ]  # dedupe, max 8  # type: ignore
            print("[BRAND] Template colors: %s" % intent.color_palette)

        # --- Extract theme fonts ---
        fonts = []
        try:
            theme_part = prs.slide_masters[0].part.slide_master.element
            ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            for font_scheme in theme_part.iter(ns_a + "fontScheme"):
                for font_tag in ("majorFont", "minorFont"):
                    font_el = font_scheme.find(ns_a + font_tag)
                    if font_el is not None:
                        latin = font_el.find(ns_a + "latin")
                        if latin is not None:
                            typeface = latin.get("typeface", "")
                            if typeface and typeface not in fonts:
                                fonts.append(typeface)
        except Exception as e:
            if VERBOSE:  # noqa: F405
                print("[VERBOSE] Theme font extraction error: %s" % e)

        if fonts:
            intent.typography_hints = fonts[:3]  # type: ignore
            print("[BRAND] Template fonts: %s" % intent.typography_hints)

        # --- Extract company name heuristic ---
        try:
            if prs.slides:
                first_slide = prs.slides[0]
                for shape in first_slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        # Heuristic: short text (1-4 words) in the first slide
                        # that is not a common placeholder is likely a company name
                        if 0 < len(text) < 60 and len(text.split()) <= 4:
                            lower = text.lower()
                            skip = {
                                "click to add title",
                                "click to add subtitle",
                                "click to add text",
                                "title",
                                "subtitle",
                            }
                            if lower not in skip:
                                intent.brand_name = text
                                print(
                                    "[BRAND] Template company name heuristic: '%s'"
                                    % text
                                )
                                break
        except Exception as e:
            if VERBOSE:  # noqa: F405
                print("[VERBOSE] Company name extraction error: %s" % e)

    except Exception as e:
        print("[WARNING] Template style extraction failed: %s" % str(e))
        return BrandStyleIntent(
            source="template", source_detail=str(os.path.basename(template_path))
        )  # type: ignore

    return intent


# === TEMPLATE VISUAL PROFILE ANALYSIS ===
# Programmatic analysis of template layout geometry, shape density, and
# visual patterns.  Runs BEFORE storyboard generation (zero LLM calls)
# so the optimizer can produce template-aware visual_suggestion and
# layout_constraints values.


def _analyze_template_visual_profile(
    template_path: str,
) -> "TemplateVisualProfile":
    """Analyze a template's visual layout characteristics programmatically.

    Opens the PPTX with python-pptx and inspects each slide for:
    - Placeholder geometry → content zone bounding boxes
    - Decorative shape count (non-placeholder, non-text shapes)
    - Background fills (solid, gradient, image)
    - Chart, table, image, and SmartArt presence
    - Slide type heuristics (title vs content vs section vs blank)

    Aggregates into a TemplateVisualProfile with layout density,
    dominant layout style, max comfortable bullets, and recommended
    text weight.

    This is complementary to (not a replacement for):
    - extract_style_from_template()   → colors/fonts only
    - _render_template_slides_to_png() → pixel-level, chunk-time only
    - _extract_template_styles()       → deep XML, assembly-time only

    Args:
        template_path: Absolute or relative path to a .pptx template file.

    Returns:
        TemplateVisualProfile with computed layout characteristics.
        Returns a default profile on any read error (graceful degradation).
    """
    if VERBOSE:
        print(
            "[VERBOSE] [VISUAL PROFILE] Starting template analysis: %s" % template_path
        )

    profile = TemplateVisualProfile()

    try:
        prs = Presentation(template_path)
    except Exception as e:
        print("[WARNING] [VISUAL PROFILE] Failed to open template: %s" % str(e))
        return profile

    # --- Slide dimensions and aspect ratio ---
    profile.slide_width_emu = prs.slide_width
    profile.slide_height_emu = prs.slide_height
    w_inches = prs.slide_width / 914400.0
    h_inches = prs.slide_height / 914400.0

    # Classify aspect ratio
    if abs(w_inches / h_inches - 16.0 / 9.0) < 0.1:
        profile.aspect_ratio = "16:9"
    elif abs(w_inches / h_inches - 4.0 / 3.0) < 0.1:
        profile.aspect_ratio = "4:3"
    elif abs(w_inches / h_inches - 16.0 / 10.0) < 0.1:
        profile.aspect_ratio = "16:10"
    else:
        profile.aspect_ratio = "custom"

    if VERBOSE:
        print(
            "[VERBOSE] [VISUAL PROFILE] Slide dimensions: %.1f x %.1f inches (%s)"
            % (w_inches, h_inches, profile.aspect_ratio)
        )

    slides = list(prs.slides)
    profile.slide_count = len(slides)
    if not slides:
        if VERBOSE:
            print(
                "[VERBOSE] [VISUAL PROFILE] Template has no slides — returning defaults."
            )
        return profile

    # Footer/utility placeholder indices to exclude from content zone
    FOOTER_PH_IDXS = {10, 11, 12, 13}  # date, footer, slide number, etc.

    slide_profiles = []
    total_shapes_all = 0
    content_w_pcts = []
    content_h_pcts = []

    for s_idx, slide in enumerate(slides):
        sp = SlideLayoutProfile(slide_index=s_idx)
        sp.total_shape_count = len(slide.shapes)
        total_shapes_all += sp.total_shape_count

        # --- Placeholder inspection ---
        content_ph_left = None
        content_ph_top = None
        content_ph_right = None
        content_ph_bottom = None
        has_title_ph = False

        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx
                if ph_idx in FOOTER_PH_IDXS:
                    continue
                sp.placeholder_count += 1

                if ph_idx == 0:
                    has_title_ph = True

                # Build content zone bounding box (union of all content placeholders)
                s_left = shape.left
                s_top = shape.top
                s_right = s_left + shape.width
                s_bottom = s_top + shape.height

                if content_ph_left is None or s_left < content_ph_left:
                    content_ph_left = s_left
                if content_ph_top is None or s_top < content_ph_top:
                    content_ph_top = s_top
                if content_ph_right is None or s_right > content_ph_right:
                    content_ph_right = s_right
                if content_ph_bottom is None or s_bottom > content_ph_bottom:
                    content_ph_bottom = s_bottom
            else:
                # Non-placeholder shape classification
                if (
                    getattr(shape, "has_text_frame", False)
                    and shape.text_frame.text.strip()
                ):
                    sp.text_box_count += 1
                elif shape.has_chart:
                    profile.has_charts_in_template = True
                elif shape.has_table:
                    profile.has_tables_in_template = True
                elif shape.shape_type is not None:
                    shape_type_val = int(shape.shape_type) if shape.shape_type else 0
                    # MSO_SHAPE_TYPE.PICTURE = 13, LINKED_PICTURE = 11
                    if shape_type_val in (13, 11):
                        profile.has_images_in_template = True
                    # MSO_SHAPE_TYPE.SMART_ART = 24, SMART_ART_GRAPHIC = 25
                    elif shape_type_val in (24, 25):
                        profile.has_smartart_shapes = True
                    else:
                        sp.decorative_shape_count += 1
                else:
                    sp.decorative_shape_count += 1

        # --- Content zone percentages ---
        if (
            content_ph_left is not None
            and content_ph_top is not None
            and content_ph_right is not None
            and content_ph_bottom is not None
        ):
            sp.content_zone_left_pct = round(  # type: ignore
                float(content_ph_left) / float(prs.slide_width) * 100.0,
                1,  # type: ignore
            )
            sp.content_zone_top_pct = round(  # type: ignore
                float(content_ph_top) / float(prs.slide_height) * 100.0,
                1,  # type: ignore
            )
            zone_w = content_ph_right - content_ph_left
            zone_h = content_ph_bottom - content_ph_top
            sp.content_zone_width_pct = round(  # type: ignore
                float(zone_w) / float(prs.slide_width) * 100.0,
                1,  # type: ignore
            )
            sp.content_zone_height_pct = round(  # type: ignore
                float(zone_h) / float(prs.slide_height) * 100.0,
                1,  # type: ignore
            )
            # Effective usable area (reduced by decorative shapes heuristically)
            deco_reduction = min(float(sp.decorative_shape_count) * 5.0, 25.0)
            sp.usable_width_pct = max(40.0, sp.content_zone_width_pct - deco_reduction)
            sp.usable_height_pct = max(40.0, sp.content_zone_height_pct)

        content_w_pcts.append(sp.content_zone_width_pct)
        content_h_pcts.append(sp.content_zone_height_pct)

        # --- Slide type heuristic ---
        if s_idx == 0 and has_title_ph and sp.placeholder_count <= 2:
            sp.slide_type_hint = "title"
        elif sp.placeholder_count == 0:
            sp.slide_type_hint = "blank"
        elif has_title_ph and sp.placeholder_count == 1:
            sp.slide_type_hint = "section"
        else:
            sp.slide_type_hint = "content"

        # --- Background analysis ---
        try:
            bg = slide.background
            fill = bg.fill
            if fill.type is not None:
                fill_type = str(fill.type)
                if "PICTURE" in fill_type.upper():
                    sp.has_background_image = True
                elif "GRADIENT" in fill_type.upper():
                    sp.has_gradient_background = True
        except Exception:
            pass  # Background inspection is best-effort

        # --- Accent line / bar detection ---
        # Accent shapes are thin decorative bars/lines. Heuristic: height < 3% or
        # width < 3% of slide dimension, non-placeholder, non-picture, non-chart.
        _height_thresh_emu = int(prs.slide_height * 0.03)
        _width_thresh_emu = int(prs.slide_width * 0.03)
        for shape in slide.shapes:
            if shape.is_placeholder:
                continue
            try:
                s_w = shape.width
                s_h = shape.height
                if s_w is None or s_h is None:
                    continue
                is_horizontal_bar = s_h < _height_thresh_emu and s_w > _width_thresh_emu
                is_vertical_bar = s_w < _width_thresh_emu and s_h > _height_thresh_emu
                if not (is_horizontal_bar or is_vertical_bar):
                    continue
                # Skip shapes with text (likely labels, not accents)
                if (
                    getattr(shape, "has_text_frame", False)
                    and shape.text_frame.text.strip()
                ):
                    continue
                # Skip pictures and charts
                if getattr(shape, "has_chart", False) or getattr(
                    shape, "has_table", False
                ):
                    continue
                shape_type_val = int(shape.shape_type) if shape.shape_type else 0
                if shape_type_val in (13, 11):  # picture types
                    continue

                # Determine orientation and position
                orientation = "horizontal" if is_horizontal_bar else "vertical"
                top_pct = round(shape.top / prs.slide_height * 100.0, 1)
                left_pct = round(shape.left / prs.slide_width * 100.0, 1)
                width_pct = round(s_w / prs.slide_width * 100.0, 1)
                height_pct = round(s_h / prs.slide_height * 100.0, 1)

                # Extract color if available
                accent_color = None
                try:
                    fill = shape.fill
                    if fill.type is not None and hasattr(fill, "fore_color"):
                        color_obj = fill.fore_color
                        if hasattr(color_obj, "rgb") and color_obj.rgb:
                            accent_color = str(color_obj.rgb)
                except Exception:
                    pass

                # Classify region
                if top_pct < 15:
                    region = "top"
                elif top_pct > 85:
                    region = "bottom"
                elif left_pct < 10:
                    region = "left"
                elif left_pct > 85:
                    region = "right"
                else:
                    region = "middle"

                sp.accent_shapes.append(
                    {
                        "orientation": orientation,
                        "region": region,
                        "top_pct": top_pct,
                        "left_pct": left_pct,
                        "width_pct": width_pct,
                        "height_pct": height_pct,
                        "color": accent_color,
                    }
                )
            except Exception:
                pass  # Accent detection is best-effort

        slide_profiles.append(sp)

        if VERBOSE:
            print(
                "[VERBOSE] [VISUAL PROFILE] Slide %d: %s | %d placeholders | "
                "%d decorative | %d text boxes | zone: %.0f%%-%.0f%% x %.0f%%-%.0f%%"
                % (
                    s_idx,
                    sp.slide_type_hint,
                    sp.placeholder_count,
                    sp.decorative_shape_count,
                    sp.text_box_count,
                    sp.content_zone_left_pct,
                    sp.content_zone_left_pct + sp.content_zone_width_pct,
                    sp.content_zone_top_pct,
                    sp.content_zone_top_pct + sp.content_zone_height_pct,
                )
            )

    profile.slide_count = len(slide_profiles)

    # --- Aggregate metrics ---
    n = float(len(slide_profiles))
    profile.avg_placeholder_count = round(  # type: ignore
        float(sum(sp.placeholder_count for sp in slide_profiles)) / n,
        1,  # type: ignore
    )
    profile.avg_decorative_shapes = round(  # type: ignore
        float(sum(sp.decorative_shape_count for sp in slide_profiles)) / n,
        1,  # type: ignore
    )
    profile.avg_content_zone_width_pct = round(float(sum(content_w_pcts)) / n, 1)  # type: ignore
    profile.avg_content_zone_height_pct = round(float(sum(content_h_pcts)) / n, 1)  # type: ignore

    # --- Layout density classification ---
    avg_total_shapes = total_shapes_all / n
    if avg_total_shapes < 3:
        profile.layout_density = "sparse"
    elif avg_total_shapes <= 6:
        profile.layout_density = "balanced"
    else:
        profile.layout_density = "dense"

    if VERBOSE:
        print(
            "[VERBOSE] [VISUAL PROFILE] Avg shapes/slide: %.1f -> density: %s"
            % (avg_total_shapes, profile.layout_density)
        )

    # --- Dominant layout style ---
    # If content zone is notably narrower than full width, it's a sidebar/split layout
    if profile.avg_content_zone_width_pct < 55:
        profile.dominant_layout_style = "sidebar"
    elif profile.avg_content_zone_width_pct < 75:
        profile.dominant_layout_style = "split"
    elif profile.avg_decorative_shapes > 3:
        profile.dominant_layout_style = "overlapping"
    else:
        profile.dominant_layout_style = "full"

    if VERBOSE:
        print(
            "[VERBOSE] [VISUAL PROFILE] Content zone avg: %.0f%% width x %.0f%% height -> style: %s"
            % (
                profile.avg_content_zone_width_pct,
                profile.avg_content_zone_height_pct,
                profile.dominant_layout_style,
            )
        )

    # --- Max comfortable bullets ---
    # Heuristic: each bullet needs ~15% of content zone height at comfortable font sizes
    profile.max_comfortable_bullets = max(
        2, min(6, int(profile.avg_content_zone_height_pct / 15))
    )

    # --- Recommended text weight ---
    if profile.avg_decorative_shapes > 2 or profile.has_images_in_template:
        profile.recommended_text_weight = "light"
    elif profile.layout_density == "dense":
        profile.recommended_text_weight = "dense"
    else:
        profile.recommended_text_weight = "balanced"

    # --- Accent pattern aggregation ---
    # Gather accent shapes across all slides to detect consistent patterns
    all_accents = []
    slides_with_accents = 0
    for sp in slide_profiles:
        accents = getattr(sp, "accent_shapes", [])
        if accents:
            slides_with_accents += 1  # type: ignore
            all_accents.extend(accents)

    if slides_with_accents >= 2 or (
        slides_with_accents == 1 and len(slide_profiles) <= 3
    ):
        profile.has_accent_lines = True
        # Find dominant orientation and region
        orientations = [a["orientation"] for a in all_accents]
        regions = [a["region"] for a in all_accents]
        colors = [a["color"] for a in all_accents if a.get("color")]

        from collections import Counter

        dominant_orientation = (
            Counter(orientations).most_common(1)[0][0] if orientations else "horizontal"
        )
        dominant_region = Counter(regions).most_common(1)[0][0] if regions else "top"
        dominant_color = Counter(colors).most_common(1)[0][0] if colors else None

        # Average geometry for the dominant pattern
        matching = [
            a
            for a in all_accents
            if a["orientation"] == dominant_orientation
            and a["region"] == dominant_region
        ]
        avg_top = (
            round(float(sum(a["top_pct"] for a in matching)) / float(len(matching)), 1)
            if matching
            else 0
        )  # type: ignore
        avg_left = (
            round(float(sum(a["left_pct"] for a in matching)) / float(len(matching)), 1)
            if matching
            else 0
        )  # type: ignore
        avg_width = (
            round(
                float(sum(a["width_pct"] for a in matching)) / float(len(matching)), 1
            )
            if matching
            else 0
        )  # type: ignore
        avg_height = (
            round(
                float(sum(a["height_pct"] for a in matching)) / float(len(matching)), 1
            )
            if matching
            else 0
        )  # type: ignore

        profile.accent_pattern = {
            "orientation": dominant_orientation,
            "region": dominant_region,
            "color": dominant_color,
            "avg_top_pct": avg_top,
            "avg_left_pct": avg_left,
            "avg_width_pct": avg_width,
            "avg_height_pct": avg_height,
            "slide_coverage_pct": round(
                float(slides_with_accents) / float(len(slide_profiles)) * 100.0, 0
            ),  # type: ignore
            "total_accent_count": len(all_accents),
        }

        if VERBOSE:
            print(
                "[VERBOSE] [VISUAL PROFILE] Accent pattern: %s %s bar (%d found across %d/%d slides, color=%s)"
                % (
                    dominant_orientation,
                    dominant_region,
                    len(all_accents),
                    slides_with_accents,
                    len(slide_profiles),
                    dominant_color or "auto",
                )
            )

    # --- Always-on summary log ---
    print(
        "[VISUAL PROFILE] Template: %s | %d slides | %s | density=%s | style=%s | "
        "max_bullets=%d | text_weight=%s"
        % (
            os.path.basename(template_path),
            profile.slide_count,
            profile.aspect_ratio,
            profile.layout_density,
            profile.dominant_layout_style,
            profile.max_comfortable_bullets,
            profile.recommended_text_weight,
        )
    )

    if VERBOSE:
        template_contents = []
        if profile.has_charts_in_template:
            template_contents.append("charts")
        if profile.has_tables_in_template:
            template_contents.append("tables")
        if profile.has_images_in_template:
            template_contents.append("images")
        if profile.has_smartart_shapes:
            template_contents.append("SmartArt")
        any_bg_image = any(sp.has_background_image for sp in slide_profiles)
        any_bg_gradient = any(sp.has_gradient_background for sp in slide_profiles)
        if any_bg_image:
            template_contents.append("background images")
        if any_bg_gradient:
            template_contents.append("gradient backgrounds")
        print(
            "[VERBOSE] [VISUAL PROFILE] Template contains: %s"
            % (", ".join(template_contents) if template_contents else "none detected")
        )

    return profile


def _build_brand_override_log(
    query_intent: "BrandStyleIntent",
    template_intent: "BrandStyleIntent",
) -> str:
    """Build a structured log message when template styling overrides query branding.

    Called only when a template file is provided AND the user query contained
    an explicit branding directive.  The log captures the specific reason for
    ignoring the query-level styling intent.

    Args:
        query_intent: BrandStyleIntent extracted from the user's natural language query.
        template_intent: BrandStyleIntent extracted from the template file.

    Returns:
        Multi-line log string suitable for printing to stdout.
    """
    template_name = template_intent.source_detail or "provided template"
    lines = [
        "[BRAND OVERRIDE] User specified '%s branding' in query, but a template file "
        "was provided (%s)." % (query_intent.brand_name, template_name),
        "[BRAND OVERRIDE] Styling will be derived from the template file. "
        "Query-level branding intent has been disregarded.",
        "[BRAND OVERRIDE] Reason: Explicit template file takes precedence over "
        "natural language branding directives per workflow specification.",
    ]
    if template_intent.color_palette:
        lines.append(
            "[BRAND OVERRIDE] Template colors: %s"
            % ", ".join(list(template_intent.color_palette)[:6])  # type: ignore
        )
    if template_intent.typography_hints:
        lines.append(
            "[BRAND OVERRIDE] Template fonts: %s"
            % ", ".join(list(template_intent.typography_hints))
        )
    return "\n".join(lines)


def _format_brand_context_for_prompt(brand_intent: "BrandStyleIntent") -> str:
    """Format a BrandStyleIntent as a markdown section for injection into LLM prompts.

    Produces a concise, structured block that can be appended to the optimizer prompt,
    Tier 1 chunk prompt, or Tier 2 code-gen prompt to guide brand-aware generation.

    Args:
        brand_intent: BrandStyleIntent to format.

    Returns:
        Markdown string with brand context, or empty string if no branding is present.
    """
    if not brand_intent or not brand_intent.has_branding:
        return ""

    sections = ["## Brand/Style Guidance\n"]

    if brand_intent.brand_name:
        sections.append("**Brand:** %s" % brand_intent.brand_name)
    if brand_intent.style_keywords:
        sections.append(
            "**Style:** %s" % ", ".join(brand_intent.style_keywords[:5])  # type: ignore
        )
    if brand_intent.color_palette:
        sections.append(
            "**Color Palette:** %s" % ", ".join(brand_intent.color_palette[:6])  # type: ignore
        )
    if brand_intent.tone_override:
        sections.append("**Tone:** %s" % brand_intent.tone_override)
    if brand_intent.typography_hints:
        sections.append(
            "**Typography:** %s" % ", ".join(brand_intent.typography_hints[:3])  # type: ignore
        )

    # Include resolved Theme Factory definition if available
    theme_def = getattr(brand_intent, "theme_definition", None)
    if theme_def and isinstance(theme_def, dict):
        sections.append("")
        sections.append("### Resolved Theme Definition (from Theme Factory)")
        sections.append("**Theme Name:** %s" % theme_def.get("name", "Unknown"))
        sections.append("**Theme Source:** %s" % theme_def.get("source", "unknown"))
        palette = theme_def.get("palette", {})
        if palette:
            palette_str = ", ".join("%s: %s" % (k, v) for k, v in palette.items())
            sections.append("**Theme Palette (hex):** %s" % palette_str)
        typo = theme_def.get("typography", {})
        if typo:
            typo_str = ", ".join("%s: %s" % (k, v) for k, v in typo.items())
            sections.append("**Theme Typography:** %s" % typo_str)
        sections.append(
            "\nYou MUST use the Theme Palette hex codes above for ALL slide backgrounds, "
            "text colors, accent fills, and chart series colors. These hex codes take "
            "precedence over any generic color names mentioned elsewhere.\n"
        )

    sections.append(
        "\nUse these brand guidelines to inform visual direction, tone, terminology, "
        "and content framing throughout the presentation. Reflect the brand's identity "
        "in slide language, suggested color references, and overall aesthetic.\n"
    )

    return "\n".join(sections)


def _format_visual_profile_for_prompt(
    profile: "TemplateVisualProfile",
) -> str:
    """Format a TemplateVisualProfile as a markdown section for the optimizer prompt.

    Produces a concise block describing the template's layout characteristics
    so the storyboard agent can generate template-aware visual_suggestion
    and layout_constraints values.

    Args:
        profile: TemplateVisualProfile from _analyze_template_visual_profile().

    Returns:
        Markdown string with visual profile context, or empty string if no profile.
    """
    if not profile or profile.slide_count == 0:
        return ""

    sections = ["## Template Visual Profile\n"]

    sections.append("- **Aspect Ratio:** %s" % profile.aspect_ratio)
    sections.append(
        "- **Layout Density:** %s (avg %.1f shapes/slide)"
        % (
            profile.layout_density,
            profile.avg_placeholder_count + profile.avg_decorative_shapes,
        )
    )
    sections.append(
        "- **Content Zone:** ~%.0f%% width, ~%.0f%% height"
        % (profile.avg_content_zone_width_pct, profile.avg_content_zone_height_pct)
    )
    if profile.dominant_layout_style != "full":
        sections.append(
            "  _(template uses %s layout — decorative elements reduce usable area)_"
            % profile.dominant_layout_style
        )
    sections.append(
        "- **Max Comfortable Bullets:** %d per slide" % profile.max_comfortable_bullets
    )
    sections.append(
        "- **Recommended Text Weight:** %s" % profile.recommended_text_weight
    )

    # Template element inventory
    elements = []
    if profile.has_charts_in_template:
        elements.append("charts")
    if profile.has_tables_in_template:
        elements.append("tables")
    if profile.has_images_in_template:
        elements.append("images")
    if profile.has_smartart_shapes:
        elements.append("SmartArt")
    deco_note = ""
    if profile.avg_decorative_shapes > 1:
        deco_note = (
            ", decorative shapes (avg %.0f/slide)" % profile.avg_decorative_shapes
        )
    sections.append(
        "- **Template Contains:** %s%s"
        % (", ".join(elements) if elements else "text placeholders only", deco_note)
    )

    # Template-aware constraints for the optimizer
    sections.append(
        "\nTEMPLATE-AWARE CONSTRAINTS (apply these when generating the storyboard):\n"
        "- Limit key_points to %d or fewer per slide (based on actual content zone height).\n"
        "- Set layout_constraints.content_zone_top_pct and content_zone_bottom_pct to "
        "%.0f and %.0f respectively.\n"
        "- Set layout_constraints.text_weight to '%s' unless the slide is purely text.\n"
        "- Prefer visual_suggestions that complement existing template decorations.\n"
        % (
            profile.max_comfortable_bullets,
            max(
                8,
                100
                - profile.avg_content_zone_height_pct
                - (100 - profile.avg_content_zone_width_pct) / 2,
            ),  # type: ignore
            min(92, profile.avg_content_zone_height_pct + 12),  # type: ignore
            profile.recommended_text_weight,
        )
    )

    if profile.has_charts_in_template:
        sections.append(
            "- For chart slides, prefer chart types already present in the template.\n"
        )

    # --- Accent pattern metadata injection ---
    # Gives the LLM explicit, structured instructions about accent lines so
    # Tier 2 code-gen can reproduce them without relying solely on visual PNGs.
    if profile.has_accent_lines and profile.accent_pattern:
        ap = profile.accent_pattern
        color_note = " in color #%s" % ap["color"] if ap.get("color") else ""
        sections.append(
            "\nACCENT LINE PATTERN (MUST REPLICATE ON EVERY CONTENT SLIDE):\n"
            "The template uses a consistent **%s accent bar** in the **%s** region%s.\n"
            "- Position: top=%.1f%%, left=%.1f%% of slide dimensions\n"
            "- Size: width=%.1f%%, height=%.1f%% of slide dimensions\n"
            "- Coverage: present on %.0f%% of template slides (%d total accent shapes detected)\n"
            "- **Action**: For every content slide in the storyboard, include this accent bar "
            "at the same relative position and size using `slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left=..., top=..., width=..., height=...)`. "
            "You MUST use `RECTANGLE` with `.line.fill.background()` (no border) to create a solid bar. DO NOT use `LINE` to avoid slanting diagonal lines. "
            "Use the exact color if provided, otherwise use the template's primary accent color.\n"
            % (
                ap.get("orientation", "horizontal"),
                ap.get("region", "top"),
                color_note,
                ap.get("avg_top_pct", 0),
                ap.get("avg_left_pct", 0),
                ap.get("avg_width_pct", 0),
                ap.get("avg_height_pct", 0),
                ap.get("slide_coverage_pct", 0),
                ap.get("total_accent_count", 0),
            )
        )

    return "\n".join(sections)


# === HELPER: STORYBOARD MARKDOWN FORMATTING ===


def _format_slide_markdown(slide: SlideStoryboard) -> str:
    """Format a SlideStoryboard as a markdown string for the pptx agent.

    Excludes transition_note (planning meta-info, not useful for content generation)
    to keep context size lean. Includes type, key points, visual suggestion,
    and semantic intelligence fields.

    Args:
        slide: SlideStoryboard instance to format.

    Returns:
        Markdown string with slide number, title, type, key points, visual suggestion, and semantic hints.
    """
    lines = [
        "## Slide %d" % slide.slide_number,
        "**Title:** %s" % slide.slide_title,
        "**Type:** %s" % slide.slide_type,
        "**Semantic Type:** %s" % getattr(slide, "semantic_type", "default"),
    ]
    if getattr(slide, "key_metrics", []):
        lines.append("**Key Metrics:** %s" % ", ".join(slide.key_metrics))

    lines.append("**Key Points:**")
    for point in slide.key_points:
        lines.append("- %s" % point)
    lines.append("**Visual Suggestion:** %s" % slide.visual_suggestion)
    lc = getattr(slide, "layout_constraints", None)
    if lc:
        lines.append(
            "**Layout Constraints:** max %d content blocks | min %dpt font | "
            "content zone %d%%-%d%% | text weight: %s"
            % (
                lc.max_content_blocks,
                lc.min_font_pt,
                lc.content_zone_top_pct,
                lc.content_zone_bottom_pct,
                lc.text_weight,
            )
        )
    return "\n".join(lines) + "\n"


def _format_global_context_markdown(plan: StoryboardPlan) -> str:
    """Format the global context as a markdown string for the pptx agent.

    Kept concise: title, audience, tone, brand voice, visual style, content balance,
    and the 2-3 sentence global context.
    This file is included in every chunk prompt so brevity matters for context size.

    Args:
        plan: StoryboardPlan instance containing global presentation metadata.

    Returns:
        Markdown string with presentation title, audience, tone, brand voice,
        visual style, content balance, and context.
    """
    return (
        "# Presentation: %s\n\n"
        "Audience: %s | Tone: %s | Brand Voice: %s | Visual Style: %s | Content Balance: %s\n\n"
        "## Context\n%s\n"
    ) % (
        plan.presentation_title,
        plan.target_audience,
        plan.tone,
        plan.brand_voice,
        plan.visual_style,
        plan.content_balance,
        plan.global_context,
    )


def _build_no_template_design_system(
    visual_style: str,
    brand_intent: Optional[BrandStyleIntent] = None,
) -> str:
    """Build a VISUAL DESIGN SYSTEM prompt section for no-template chunk generation.

    Maps the visual_style inferred by the optimizer to concrete python-pptx
    design instructions: background color, accent line, title bar, footer,
    color palette, and typography.  Returns an empty string when
    visual_style == 'template_driven' (template already dictates styling).

    Per SCRATCHPAD Note: "always follow a minimal, standard design outline
    with Accent lines, Title/Header, Footer and a Logo image (optional)."

    Args:
        visual_style: One of 'bold_modern', 'clean_minimal',
                      'creative_experimental', 'corporate_professional',
                      or 'template_driven'.
        brand_intent: Optional brand context — when present with specific
                      color_palette, those colors override the default palette.

    Returns:
        Multi-line prompt string with python-pptx design instructions,
        or empty string for template_driven.
    """
    # When visual_style == "template_driven" but there is NO actual template file,
    # we still need to build a design system. The theme_definition from the Theme
    # Factory provides the concrete hex codes. If theme_definition is absent AND
    # visual_style is template_driven, fall back to clean_minimal.
    if visual_style == "template_driven":
        if brand_intent and getattr(brand_intent, "theme_definition", None):
            # Theme Factory resolved a theme — build design system from it
            td = brand_intent.theme_definition
            palette = td.get("palette", {})
            typo = td.get("typography", {})
            # Map theme palette keys to design tokens
            dk1 = palette.get("dk1", "1A1A2E").lstrip("#")
            accent1 = palette.get("accent1", "00D4AA").lstrip("#")
            lt1 = palette.get("lt1", "CCCCCC").lstrip("#")
            lt2 = palette.get("lt2", "FFFFFF").lstrip("#")
            major_font = typo.get("major", "Segoe UI")
            minor_font = typo.get("minor", "Calibri")

            # Determine if dark or light background
            try:
                lum = (
                    int(dk1[:2], 16) * 0.299
                    + int(dk1[2:4], 16) * 0.587
                    + int(dk1[4:6], 16) * 0.114
                ) / 255
            except Exception:
                lum = 0.1
            text_hex = lt2 if lum < 0.4 else "333333"
            text_label = "light/white" if lum < 0.4 else "dark"

            brand_name_text = (
                brand_intent.brand_name.upper() if brand_intent.brand_name else "LOGO"
            )

            if VERBOSE:
                print(
                    "[VERBOSE] [DESIGN SYSTEM] Building from Theme Factory definition: "
                    "bg=#%s, accent=#%s, text=#%s, fonts=%s/%s"
                    % (dk1, accent1, text_hex, major_font, minor_font)
                )

            return (
                "\nVISUAL DESIGN SYSTEM (MANDATORY — follow these styling rules for every slide):\n"
                "Style: Theme Factory — %s\n\n"
                "BACKGROUND:\n"
                "  Set the slide background for EVERY slide to #%s. Do NOT use white or any other color.\n"
                "  Code: from pptx.util import Inches, Pt, Emu\n"
                "        from pptx.dml.color import RGBColor\n"
                "        slide.background.fill.solid()\n"
                "        slide.background.fill.fore_color.rgb = RGBColor(0x%s, 0x%s, 0x%s)\n\n"
                "TEXT COLORS (ADAPTIVE):\n"
                "  Body text: #%s (%s) — ensures readability on the dark background\n"
                "  Title text: #%s at 32pt using font '%s'\n"
                "  Body font: '%s' at 16pt minimum\n\n"
                "ACCENT COLOR:\n"
                "  Primary accent: #%s — use for accent bars, shape fills, chart series, highlights\n"
                "  Secondary/soft accent: #%s — use for subtle highlights, secondary elements\n\n"
                "LOGO (ADD TO EVERY SLIDE):\n"
                "  Add a logo placeholder text at the top right corner of each slide:\n"
                "    logo_box = slide.shapes.add_textbox(Inches(8.0), Inches(0.3), Inches(1.5), Inches(0.4))\n"
                "    logo_tf = logo_box.text_frame\n"
                "    logo_p = logo_tf.paragraphs[0]\n"
                "    logo_p.text = '%s'\n"
                "    logo_p.font.size = Pt(14)\n"
                "    logo_p.font.bold = True\n"
                "    logo_p.font.color.rgb = RGBColor(0x%s, 0x%s, 0x%s)\n\n"
                "COLOR RULES:\n"
                "  - NEVER use white (#FFFFFF) as background. The theme background is #%s.\n"
                "  - ALL text must be readable against the #%s background.\n"
                "  - Use #%s for accent bars, dividers, chart fills, and highlighted shapes.\n"
                "  - Use #%s for soft highlights and secondary fills.\n"
                "  - For charts: use theme accent colors for data series, NOT default Office colors.\n"
                "  - For tables: header row fill=#%s with text=#%s; body rows alternate #%s and background.\n"
            ) % (
                td.get("name", "Custom Theme"),
                dk1,
                dk1[:2],
                dk1[2:4],
                dk1[4:6],
                text_hex,
                text_label,
                lt2,
                major_font,
                minor_font,
                accent1,
                lt1,
                brand_name_text,
                accent1[:2],
                accent1[2:4],
                accent1[4:6],
                dk1,
                dk1,
                accent1,
                lt1,
                accent1,
                lt2,
                lt1,
            )
        else:
            # No theme definition and no template — fall back to clean_minimal
            visual_style = "clean_minimal"

    # --- Style-specific design tokens ---
    STYLE_MAP = {
        "bold_modern": {
            "bg_hex": "1A1A2E",
            "bg_label": "dark navy",
            "text_color": "FFFFFF",
            "text_label": "white",
            "accent_hex": "00D4AA",
            "accent_label": "teal/cyan",
            "secondary_hex": "E94560",
            "title_size": 32,
            "body_size": 16,
            "font_family": "Segoe UI",
            "accent_bar_height_pct": 1.5,
            "accent_bar_top_pct": 10.0,
            "description": "Bold & Modern: Dark background with vibrant accents, high contrast",
        },
        "clean_minimal": {
            "bg_hex": "FFFFFF",
            "bg_label": "white",
            "text_color": "333333",
            "text_label": "dark charcoal",
            "accent_hex": "4A90D9",
            "accent_label": "muted blue",
            "secondary_hex": "7F8C8D",
            "title_size": 28,
            "body_size": 14,
            "font_family": "Calibri",
            "accent_bar_height_pct": 0.8,
            "accent_bar_top_pct": 10.0,
            "description": "Clean & Minimal: Light background, whitespace, muted palette",
        },
        "creative_experimental": {
            "bg_hex": "0F0E17",
            "bg_label": "deep dark",
            "text_color": "FFFFFE",
            "text_label": "off-white",
            "accent_hex": "FF8906",
            "accent_label": "warm orange",
            "secondary_hex": "E53170",
            "title_size": 34,
            "body_size": 16,
            "font_family": "Segoe UI",
            "accent_bar_height_pct": 2.0,
            "accent_bar_top_pct": 8.0,
            "description": "Creative & Experimental: Deep dark bg, vibrant warm accents, dynamic feel",
        },
        "corporate_professional": {
            "bg_hex": "F5F5F5",
            "bg_label": "light gray",
            "text_color": "2C3E50",
            "text_label": "navy/dark blue",
            "accent_hex": "2980B9",
            "accent_label": "professional blue",
            "secondary_hex": "27AE60",
            "title_size": 28,
            "body_size": 14,
            "font_family": "Calibri",
            "accent_bar_height_pct": 1.0,
            "accent_bar_top_pct": 10.0,
            "description": "Corporate Professional: Structured, conservative, data-focused",
        },
    }

    tokens = STYLE_MAP.get(visual_style, STYLE_MAP["clean_minimal"])

    # Override palette from brand intent if specific colors were detected
    brand_section = ""
    if brand_intent and brand_intent.has_branding:
        tokens = dict(tokens)  # copy

        # --- Theme Factory full override (highest priority) ---
        # When a theme_definition is present, it contains the authoritative
        # palette (dk1=background, accent1=accent, lt1=soft, lt2=text).
        # Override ALL design tokens from the theme definition.
        theme_def = getattr(brand_intent, "theme_definition", None)
        if theme_def and isinstance(theme_def, dict):
            td_palette = theme_def.get("palette", {})
            td_typo = theme_def.get("typography", {})

            # Background: dk1
            dk1 = td_palette.get("dk1", "").lstrip("#")
            if len(dk1) == 6:
                tokens["bg_hex"] = dk1
                tokens["bg_label"] = "theme background (%s)" % theme_def.get(
                    "name", "custom"
                )

            # Text color: lt1 (or lt2 as fallback)
            lt1 = td_palette.get("lt1", "").lstrip("#")
            lt2 = td_palette.get("lt2", "").lstrip("#")
            text_hex = lt1 if len(lt1) == 6 else (lt2 if len(lt2) == 6 else "")
            if len(text_hex) == 6:
                tokens["text_color"] = text_hex
                tokens["text_label"] = "theme text"

            # Accent: accent1
            accent1 = td_palette.get("accent1", "").lstrip("#")
            if len(accent1) == 6:
                tokens["accent_hex"] = accent1
                tokens["accent_label"] = "theme accent"

            # Secondary: lt2 (or lt1 as fallback for soft accent)
            soft = lt2 if len(lt2) == 6 else (lt1 if len(lt1) == 6 else "")
            if len(soft) == 6:
                tokens["secondary_hex"] = soft

            # Typography from theme
            if td_typo.get("major"):
                tokens["font_family"] = td_typo["major"]

            if VERBOSE:
                print(
                    "[VERBOSE] [DESIGN SYSTEM] Theme Factory FULL override: "
                    "bg=#%s, text=#%s, accent=#%s, secondary=#%s, font=%s (theme: %s)"
                    % (
                        tokens["bg_hex"],
                        tokens["text_color"],
                        tokens["accent_hex"],
                        tokens["secondary_hex"],
                        tokens["font_family"],
                        theme_def.get("name", "unknown"),
                    )
                )
        elif brand_intent.color_palette:
            # Fallback: no theme_definition, use raw brand colors for accent/secondary only
            palette = brand_intent.color_palette
            if len(palette) >= 1:
                raw = palette[0].lstrip("#")
                if len(raw) == 6:
                    tokens["accent_hex"] = raw
            if len(palette) >= 2:
                raw2 = palette[1].lstrip("#")
                if len(raw2) == 6:
                    tokens["secondary_hex"] = raw2
            if VERBOSE:
                print(
                    "[VERBOSE] [DESIGN SYSTEM] Brand palette override (no theme): "
                    "accent=#%s, secondary=#%s (from brand: %s)"
                    % (
                        tokens["accent_hex"],
                        tokens["secondary_hex"],
                        brand_intent.brand_name,
                    )
                )

        # Apply brand typography if available (only if not already set by theme)
        if brand_intent.typography_hints and not (
            theme_def and isinstance(theme_def, dict)
        ):
            tokens["font_family"] = brand_intent.typography_hints[0]
            if VERBOSE:
                print(
                    "[VERBOSE] [DESIGN SYSTEM] Brand font override: '%s' (from brand: %s)"
                    % (tokens["font_family"], brand_intent.brand_name)
                )

        # Build robust branding contextual string for the LLM
        brand_name = brand_intent.brand_name
        tone = brand_intent.tone_override or "optimistic"
        style_kw = ", ".join(brand_intent.style_keywords)

        brand_section = (
            f"BRANDING & TONE (CRITICAL):\n"
            f"  Brand Name: {brand_name}\n"
            f"  Visual Tone: {tone.capitalize()} (Keywords: {style_kw})\n"
            f"  Adaptive Layout: Blend all visual elements, shapes, typography, and alignments flawlessly with the '{visual_style}' theme, while strictly maintaining '{brand_name}' aesthetics.\n"
            f"  Adaptive Color Contrast: Ensure the brand's primary colors pop and maintain perfect contrast. Adjust shape and text colors if they clash with the background, but do NOT change the background color itself (maintain the baseline theme background).\n\n"
        )

    bg_hex = str(tokens["bg_hex"])
    bg_r, bg_g, bg_b = bg_hex[:2], bg_hex[2:4], bg_hex[4:6]
    txt_hex = str(tokens["text_color"])
    txt_r, txt_g, txt_b = txt_hex[:2], txt_hex[2:4], txt_hex[4:6]
    acc_hex = str(tokens["accent_hex"])
    acc_r, acc_g, acc_b = acc_hex[:2], acc_hex[2:4], acc_hex[4:6]
    sec_hex = str(tokens["secondary_hex"])

    brand_name_text = (
        brand_intent.brand_name.upper()
        if brand_intent and brand_intent.brand_name
        else "LOGO"
    )

    return (
        "\nVISUAL DESIGN SYSTEM (MANDATORY — follow these styling rules for every slide):\n"
        "Style: %s\n\n"
        "%s"
        "BACKGROUND:\n"
        "  Set the slide background for EVERY slide to #%s (%s). Do NOT change this color.\n"
        "  Code: from pptx.util import Inches, Pt, Emu\n"
        "        from pptx.dml.color import RGBColor\n"
        "        slide.background.fill.solid()\n"
        "        slide.background.fill.fore_color.rgb = RGBColor(0x%s, 0x%s, 0x%s)\n\n"
        "TEXT COLORS (ADAPTIVE):\n"
        "  Suggested body text: #%s (%s)\n"
        "  Suggested title text: #%s (%s) at %dpt\n"
        "  Body font size: %dpt minimum\n"
        "  Font family: %s\n\n"
        "LOGO (ADD TO EVERY SLIDE):\n"
        "  Add a logo placeholder text at the top right corner of each slide:\n"
        "    logo_box = slide.shapes.add_textbox(Inches(8.5), Inches(0.2), Inches(1.2), Inches(0.4))\n"
        "    p_logo = logo_box.text_frame.paragraphs[0]\n"
        "    p_logo.text = '%s'\n"
        "    p_logo.font.bold = True\n"
        "    p_logo.font.size = Pt(12)\n"
        "    p_logo.font.color.rgb = RGBColor(0x%s, 0x%s, 0x%s)\n\n"
        "ACCENT LINE (MUST ADD TO EVERY CONTENT SLIDE):\n"
        "  Add a thin horizontal accent bar near the top of each content slide (not title slide):\n"
        "    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE\n"
        "    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,\n"
        "        left=Inches(0.5), top=Inches(%.2f),\n"
        "        width=Inches(9.0), height=Inches(%.2f))\n"
        "    bar.fill.solid()\n"
        "    bar.fill.fore_color.rgb = RGBColor(0x%s, 0x%s, 0x%s)\n"
        "    bar.line.fill.background()  # no border\n"
        "  This accent bar is a REQUIRED design element per the design system.\n\n"
        "FOOTER (ADD TO EVERY SLIDE):\n"
        "  Add a small footer text box at the bottom of each slide:\n"
        "    from pptx.util import Inches, Pt\n"
        "    footer = slide.shapes.add_textbox(\n"
        "        Inches(0.5), Inches(6.9), Inches(9.0), Inches(0.35))\n"
        "    tf = footer.text_frame\n"
        "    tf.word_wrap = True\n"
        "    p = tf.paragraphs[0]\n"
        "    p.text = 'Slide N of M'  # Replace N with slide number, M with total\n"
        "    p.font.size = Pt(9)\n"
        "    p.font.color.rgb = RGBColor(0x%s, 0x%s, 0x%s) # (Adjust RGB if needed)\n"
        "    p.alignment = PP_ALIGN.RIGHT  # from pptx.enum.text import PP_ALIGN\n\n"
        "TITLE BAR (title slides only):\n"
        "  For the first/title slide, add a full-width accent banner:\n"
        "    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,\n"
        "        Inches(0), Inches(2.5), Inches(10), Inches(2.5))\n"
        "    banner.fill.solid()\n"
        "    banner.fill.fore_color.rgb = RGBColor(0x%s, 0x%s, 0x%s)\n"
        "    banner.line.fill.background()\n"
        "  Place the title text ON TOP of this banner with contrasting color.\n\n"
        "COLOR PALETTE:\n"
        "  Primary accent: #%s    Secondary accent: #%s\n"
        "  Use these for chart series, table headers, card fills, and metric highlights.\n"
        "  Do NOT introduce arbitrary colors outside this palette.\n"
    ) % (
        tokens["description"],
        brand_section,
        bg_hex,
        tokens["bg_label"],
        bg_r,
        bg_g,
        bg_b,
        txt_hex,
        tokens["text_label"],
        txt_hex,
        tokens["text_label"],
        tokens["title_size"],
        tokens["body_size"],
        tokens["font_family"],
        brand_name_text,
        txt_r,
        txt_g,
        txt_b,
        float(tokens["accent_bar_top_pct"])
        / 100.0
        * 7.5,  # Convert pct to inches (7.5" slide)
        float(tokens["accent_bar_height_pct"]) / 100.0 * 7.5,
        acc_r,
        acc_g,
        acc_b,
        txt_r,
        txt_g,
        txt_b,  # footer text color = body text color
        acc_r,
        acc_g,
        acc_b,  # title banner = accent color
        acc_hex,
        sec_hex,
    )


# === HELPER: SAVE PROMPT TO FILE ===


def _save_prompt_to_file(
    prompt: str, step_name: str, output_dir: str, extra: str = ""
) -> str:
    """Save a prompt string to a timestamped .txt file inside output_dir.

    Files are written to output_dir directly (which is output_chunked/chunked_workflow_work/).
    Filenames follow the pattern: prompt_<step_name>[_<extra>]_<timestamp_ms>.txt

    Args:
        prompt: The full prompt text to save.
        step_name: Short identifier for the workflow step (e.g. 'optimize_and_plan', 'chunk').
        output_dir: Directory in which to write the file.
        extra: Optional extra qualifier appended between step_name and timestamp.

    Returns:
        Absolute path of the saved file, or empty string if saving failed.
    """
    os.makedirs(output_dir, exist_ok=True)
    timestamp_ms = int(time.time() * 1000)
    if extra:
        filename = "prompt_%s_%s_%d.txt" % (step_name, extra, timestamp_ms)
    else:
        filename = "prompt_%s_%d.txt" % (step_name, timestamp_ms)
    filepath = os.path.join(output_dir, filename)
    try:
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(prompt)
    except Exception as e:
        print("[WARNING] Failed to save prompt file %s: %s" % (filepath, e))
        return ""
    return filepath


# === HELPER: EXTRACT SLIDES DATA FROM A CHUNK PPTX ===


def _extract_chunk_slides_data(chunk_file: str) -> List[dict]:
    """Extract basic slide metadata from a PPTX chunk file.

    Returns a list of dicts compatible with the session_state['slides_data'] format
    used by step_plan_images and step_generate_images.
    """
    slides_data = []
    try:
        prs = Presentation(chunk_file)
        for idx, slide in enumerate(prs.slides):
            slide_info: dict = {
                "index": idx,
                "title": "",
                "body": "",
                "has_table": False,
                "has_chart": False,
                "has_image": False,
            }
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for ph_attr in ["placeholder_format"]:
                        ph = getattr(shape, ph_attr, None)
                        if ph and hasattr(ph, "idx") and ph.idx == 0:
                            slide_info["title"] = shape.text_frame.text.strip()
                if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                    slide_info["has_table"] = True
                if shape.shape_type == 3:  # CHART
                    slide_info["has_chart"] = True
            slides_data.append(slide_info)
    except Exception as e:
        print("[WARNING] Could not extract slides data from %s: %s" % (chunk_file, e))
    return slides_data


# === WORKFLOW STEP 1: OPTIMIZE AND PLAN ===


def step_optimize_and_plan(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 1: Parse brand intent, enhance the user prompt, and generate a per-slide storyboard.

    Substeps:
      0. Brand/Style Parsing — calls the brand_style_analyzer agent
         to detect branding directives in the user prompt (e.g. "using Nike branding").
         The agent uses tools (built-in search or DuckDuckGo) to look up brand colors,
         tone, and typography if needed.
         If a template file is provided, extracts styling from the template's theme XML
         and overrides any query-level branding with a [BRAND OVERRIDE] log.
         Stores the effective BrandStyleIntent in session_state["brand_style_intent"].
      1. Query Optimization — calls the query_optimizer agent with the
         user prompt enriched with brand context and brand-aware search guidance.
         This agent uses web search tools (e.g., DuckDuckGo) to gather current facts
         before producing a grounded StoryboardPlan.
         Produces a StoryboardPlan with optimal slide count, global context, per-slide
         storyboard, tone, and brand voice.

    Uses the query_optimizer agent to produce a StoryboardPlan with:
    - Optimal slide count (respects user-specified count, otherwise picks 8-15)
    - Global context applicable to all slides
    - Per-slide storyboard with title, type, key points, visual suggestions
    - Presentation tone and brand voice

    Saves storyboard to individual markdown files in {output_dir}/storyboard/.

    Args:
        step_input: Workflow step input (not used directly; context comes from session_state).
        session_state: Shared workflow state containing user_prompt, output_dir, chunk_size,
                       max_retries, template_path, and template_visuals flag.

    Returns:
        StepOutput with success=True and a summary string when a valid storyboard is produced,
        or success=False with an error message if the optimizer fails or returns invalid JSON.

    Side effects:
        - session_state["brand_style_intent"] is set to the effective BrandStyleIntent.
        - session_state["storyboard"] is set to the validated StoryboardPlan.
        - Storyboard markdown files are written to {output_dir}/storyboard/.
    """
    step_start = time.time()

    user_prompt = session_state.get("user_prompt", "")
    output_dir = session_state.get("output_dir", ".")
    chunk_size = session_state.get("chunk_size", 1)
    max_retries = session_state.get("max_retries", 2)
    template_path = session_state.get("template_path", "")

    print("=" * 60)
    print("Step 1: Optimizing query and generating storyboard...")
    print("=" * 60)
    print("User prompt: %s" % user_prompt[:200])

    # === BRAND/STYLE PARSING ===
    # Parse branding intent from user query via the brand_style_analyzer agent.
    # If a template file is provided, extract template styling and override query branding.
    brand_parse_start = time.time()
    # Lazily load provider-specific agents inside the step (not in session_state)
    # so that Agent objects never enter the deepcopy path at workflow startup.
    from agents import get_agents as _get_agents  # type: ignore

    _provider = session_state.get("llm_provider", "claude")
    agents = _get_agents(_provider)
    query_brand_intent = parse_brand_style_intent(
        user_prompt,
        brand_agent=agents.get("brand_style_analyzer"),
        brand_agent_fallback=agents.get("brand_style_analyzer_fallback"),
    )
    brand_intent = query_brand_intent  # default: use query-derived intent

    if template_path and os.path.isfile(template_path):
        template_intent = extract_style_from_template(template_path)
        if query_brand_intent.has_branding and query_brand_intent.brand_name:
            # Template overrides query-level branding — log the decision
            override_log = _build_brand_override_log(
                query_brand_intent, template_intent
            )
            print(override_log)
        brand_intent = template_intent
    else:
        # No template provided. If there is branding intent or default theme requested, use Theme Factory.
        # Haiku 4.5 is extremely fast and effective for theme resolution
        model = (
            agents.get("fallback_code_agent_lite").model
            if agents.get("fallback_code_agent_lite")
            else None
        )
        if model:
            brand_intent = _resolve_theme_with_factory(brand_intent, model)

    session_state["brand_style_intent"] = brand_intent
    brand_parse_elapsed = time.time() - brand_parse_start
    print("[TIMING] Brand/style parsing completed in %.1fs" % brand_parse_elapsed)

    # Render template slides to PNG for visual reference in chunk prompts
    if template_path and os.path.isfile(template_path):
        # Always invoke rendering so 'template_pngs' folder is created for the session
        print("[STEP 1] Rendering template slides to PNG...")
        template_pngs = _render_template_slides_to_png(template_path, output_dir)

        # Only store in session_state for prompt injection if flag is set
        if session_state.get("template_visuals"):
            session_state["template_slide_pngs"] = template_pngs
        else:
            session_state["template_slide_pngs"] = {}
    else:
        session_state["template_slide_pngs"] = {}

    # Build brand context section for injection into the optimizer prompt
    brand_context_section = _format_brand_context_for_prompt(brand_intent)

    # Analyze template visual profile for storyboard-aware planning
    visual_profile = None
    visual_profile_section = ""
    if template_path and os.path.isfile(template_path):
        print("[STEP 1] Analyzing template visual profile...")
        visual_profile = _analyze_template_visual_profile(template_path)
        session_state["template_visual_profile"] = visual_profile
        visual_profile_section = _format_visual_profile_for_prompt(visual_profile)
        if VERBOSE and visual_profile_section:  # noqa: F405
            print(
                "[VERBOSE] [VISUAL PROFILE] Profile prompt section (%d chars) will be "
                "injected into query optimizer prompt" % len(visual_profile_section)
            )
    else:
        session_state["template_visual_profile"] = None

    # Build brand-enriched search guidance
    brand_search_guidance = ""
    if brand_intent.has_branding and brand_intent.brand_name:
        brand_search_guidance = (
            "BRAND-ENRICHED SEARCH GUIDANCE:\n"
            "The presentation has a brand context: '%s'. When constructing search queries,\n"
            "include at least one query that combines the core topic with the brand name\n"
            "(e.g., '%s [core topic]' or '[core topic] %s innovation').\n"
            "This helps ground the presentation in brand-relevant context.\n\n"
        ) % (brand_intent.brand_name, brand_intent.brand_name, brand_intent.brand_name)

    storyboard_dir = os.path.join(output_dir, "storyboard")
    os.makedirs(storyboard_dir, exist_ok=True)

    optimizer_prompt = (
        "Analyze the following user request for a PowerPoint presentation and create an optimized storyboard.\n\n"
        "USER REQUEST:\n%s\n\n"
        "%s"
        "%s"
        "**CRITICAL INSTRUCTIONS FOR STORYBOARD GENERATION:**\n"
        "1. **TEXT FOOTPRINT MINIMIZATION**: Prioritize visual communication. Keep all text "
        "   (titles, key points, global context) as concise as possible. Avoid verbose "
        "   descriptions. The goal is a visually-driven presentation, not a document.\n"
        "2. **HIGH-IMPACT VISUAL GUIDANCE**: Propose bold visual elements (e.g., hero images, "
        "   process diagrams, interconnected nodes). Dictate exactly what visual structure "
        "   should accompany the text, ensuring visual variety across the presentation. "
        "   Crucially, when proposing a slide with charts or large visuals, keep the text "
        "   extremely minimal to prevent overlapping issues in the final layout.\n"
        "3. **AESTHETIC CONSTRAINTS**: Prevent 'wall-of-text' layouts. Limit bullet points to 3-4 "
        "   per slide, with concise phrases. Use 'key_metrics' arrays for data highlights rather "
        "   than burying metrics inside dense paragraphs. For the final conclusion/summary "
        "   slide, enforce a visually striking wrap-up (e.g., three pillars, a bold quote block, "
        "   or a simple next-steps diagram) instead of a standard bulleted list.\n"
        "4. **STORY ARC**: Ensure the narrative flows logically and matches the specified tone.\n"
        "5. **AUDIENCE INFERENCE (MANDATORY)**:\n"
        "   From the given prompt input, understand and determine 'Who is the primary audience for this presentation?'\n"
        "   Set `target_audience` to a **specific** description. Examples:\n"
        "   - 'Potential clients' (Pitching expertise to prospects)\n"
        "   - 'Internal team / founders' (Strategic planning for building the agency)\n"
        "   - 'Industry peers / conference' (Thought leadership on trends & agency models)\n"
        "   - 'Investors / board members' (Metrics-driven, concise, high-level)\n"
        "   Tailor depth, phrasing, and visual complexity to the inferred audience.\n"
        "6. **LAYOUT CONSTRAINTS**: For each slide, include a `layout_constraints` object:\n"
        "   - `max_content_blocks`: Maximum text/shape groups (default 4; use 2-3 for chart/visual slides, 4-5 for text slides).\n"
        "   - `min_font_pt`: Minimum font size (default 14; use 18+ for title slides, 14 for content).\n"
        "   - `content_zone_top_pct` / `content_zone_bottom_pct`: Safe content area (default 12%%-88%%).\n"
        "   - `text_weight`: 'light' for visual-heavy slides, 'balanced' for mixed, 'dense' only for text-heavy slides.\n"
        "7. **VISUAL ANTI-PATTERNS (AVOID AT ALL COSTS)**:\n"
        "   - Never plan more content than can fit without overlapping at 14pt+ font size.\n"
        "   - Never suggest layouts that would push content outside slide boundaries.\n"
        "   - For chart/infographic slides, plan at most 1-2 text blocks alongside the visual.\n"
        "   - Title + footer consume ~24%% of vertical space — plan content accordingly.\n"
        "8. **VISUAL STYLE INFERENCE (when NO template is provided)**:\n"
        "   From the given prompt, if not explicitly mentioned, try to understand what tone and visual style\n"
        "   is applicable based on the prompt context and audience.\n"
        "   **COMBINE audience (from step 5) + tone/style to determine the best cohesive pptx THEME:**\n"
        "   - Primary and secondary colors (accent palette)\n"
        "   - Suitable font style, font family, and font color\n"
        "   - Layouts and alignments appropriate to the audience\n"
        "   - Titles, subtitles, and footer styling\n"
        "   Set `visual_style` in the output JSON. Options:\n"
        "   - 'bold_modern': Dark background, vibrant accents, high contrast (best for: creative pitches, tech demos, startup decks)\n"
        "   - 'clean_minimal': Light background, lots of whitespace, muted palette (best for: executive briefs, strategy docs, investor updates)\n"
        "   - 'creative_experimental': Gradients, asymmetric layouts, dynamic feel (best for: design showcases, trend reports, thought leadership)\n"
        "   - 'corporate_professional': Structured, conservative, data-focused (best for: board meetings, compliance, financial reports)\n"
        "   If any Brand name, company name, or some explicit style or theme name is mentioned, you MUST follow it closely and blend it into the chosen style.\n"
        "   If a template IS provided, set `visual_style` to 'template_driven'.\n"
        "9. **CONTENT BALANCE (multi-topic prompts)**:\n"
        "   When the prompt covers multiple topics or domains, explicitly decide and set \n"
        "   `content_balance` to describe the distribution strategy. Examples:\n"
        "   - 'equal' — Even split across all topics\n"
        "   - '<main_topic>-heavy' — Primary topic gets ~60-70%% of slides\n"
        "   - 'context-then-detail' — 30%% context-setting, 70%% deep-dive\n"
        "   For single-topic prompts, set it to 'focused'.\n"
        "10. **TEMPLATE VISUAL REUSE (when template provided):**\n"
        "    When a template file is provided, analyse its slides for reusable visual\n"
        "    structures (charts, infographics, icon grids, SmartArt, image layouts).\n"
        "    For EACH generated slide, if a template slide's visual structure is\n"
        "    contextually applicable to the new content, set `reuse_template_slide_idx`\n"
        "    to that template slide's 0-based index. This tells the assembler to\n"
        "    clone that template slide and preserve its visual elements while only\n"
        "    replacing textual content (titles, labels, body text).\n"
        "    Rules:\n"
        "    - Only reference template slides whose visual structure MATCHES the\n"
        "      semantic intent of the generated slide (e.g., a data-comparison\n"
        "      template slide for a generated 'data' slide with comparisons).\n"
        "    - Do NOT reuse a template slide just because it exists — only when\n"
        "      its chart/infographic/layout genuinely fits the new content.\n"
        "    - Set to null if no template slide is a good visual match.\n"
        "    - The same template slide index CAN be referenced by multiple\n"
        "      generated slides if its visual structure is broadly applicable.\n"
        "STEP 1 — RESEARCH FIRST:\n"
        "Extract and define ONE clear Search Topic from the user request before calling the search tool.\n"
        "Generate 1-3 highly specific search queries aimed at finding verified data points.\n"
        "Before planning slides, use the search tool with those queries to find 2-4 relevant facts, statistics, "
        "and recency-based data about the topic.\n"
        "STRICTLY cite the data source and year in the slide content (e.g., 'Source: IEA 2024').\n"
        "If the search tool returns weak or conflicting data, do NOT invent facts. Use conservative language "
        "and reduce numeric specificity rather than fabricating values.\n\n"
        "%s"
        "STEP 1B — RUN WEB SEARCH USING THAT SEARCH TOPIC:\n"
        "Before planning slides, use web_search with those queries to find 2-4 relevant facts, statistics, "
        "or examples for the Search Topic. Prioritize recent, credible sources and specific numbers "
        "(e.g., market size, CAGR, adoption rates, trend changes). Use these findings to ground the storyboard "
        "in real, specific information. For each fact you use, internally track source name + publication year "
        "and prefer the newest credible source when sources conflict.\n"
        "If web_search returns weak or conflicting data, do NOT invent facts. Use conservative language "
        "and reduce numeric specificity rather than fabricating values.\n\n"
        "STEP 2 — BUILD THE STORYBOARD:\n"
        "1. If the user specifies a slide count (e.g. '12 slides', '10-slide deck'), honor it exactly.\n"
        "2. If not specified, infer the optimal count from topic complexity and audience:\n"
        "   - **Short (5-7 slides)**: High-level overview, single focused topic, quick pitch\n"
        "   - **Medium (10-12 slides)**: Deeper dive, multiple sub-topics, standard business deck\n"
        "   - **Long (15-20 slides)**: Comprehensive coverage, multiple domains, full playbook\n"
        "   Use the audience to calibrate: executives prefer shorter; technical teams accept longer.\n"
        "   Do NOT add unnecessary slides; quality over quantity.\n"
        "3. Define a clear tone and brand voice appropriate to the topic and audience.\n"
        "   Consider: Is the topic formal/informal? Is the audience expecting data or inspiration?\n"
        "4. Write global_context as 2-3 focused sentences covering: the core topic, target audience, "
        "and central theme or key message. Include a specific fact or statistic from your research.\n"
        "5. For each slide, provide:\n"
        "   - A concise, descriptive title (5-8 words max)\n"
        "   - slide_type: one of title, agenda, content, data, closing\n"
        "   - 3-4 key_points: each a single sentence (10-20 words), specific and actionable. "
        "     Include real data or examples where relevant. Avoid vague filler bullets.\n"
        "   - visual_suggestion: one concrete line — specify chart type + data, image concept, or 'none'. "
        "     Example: 'bar chart: AI adoption rate by industry 2023' not just 'chart'.\n"
        "   - transition_note: one brief sentence connecting this slide to the next.\n"
        "6. Ensure continuity: the storyboard should feel like a coherent narrative arc.\n"
        "7. Use professional language. Do not add emojis or overly casual language.\n\n"
        "STEP 3 — OUTPUT FORMAT:\n"
        "Respond with ONLY a valid JSON object matching this exact schema (no markdown fences, "
        "no extra commentary before or after the JSON).\n"
        "The JSON must be syntactically valid and parseable with strict JSON parsers.\n"
        "Ensure slides length == total_slides and slide_number values are contiguous from 1..total_slides:\n"
        "{\n"
        '  "total_slides": <integer>,\n'
        '  "presentation_title": "<string>",\n'
        '  "search_topic": "<string>",\n'
        '  "target_audience": "<string>",\n'
        '  "tone": "<string>",\n'
        '  "brand_voice": "<string>",\n'
        '  "visual_style": "<bold_modern|clean_minimal|creative_experimental|corporate_professional|template_driven>",\n'
        '  "content_balance": "<string>",\n'
        '  "global_context": "<string>",\n'
        '  "slides": [\n'
        "    {\n"
        '      "slide_number": <integer>,\n'
        '      "slide_title": "<string>",\n'
        '      "slide_type": "<title|agenda|content|data|closing>",\n'
        '      "semantic_type": "<sequential|comparative|metrics|hero|default>",\n'
        '      "key_points": ["<string>", ...],\n'
        '      "key_metrics": ["<string>", ...],\n'
        '      "visual_suggestion": "<string>",\n'
        '      "layout_constraints": {\n'
        '        "max_content_blocks": <integer 2-6>,\n'
        '        "min_font_pt": <integer 12-24>,\n'
        '        "content_zone_top_pct": <integer 8-15>,\n'
        '        "content_zone_bottom_pct": <integer 82-92>,\n'
        '        "text_weight": "<light|balanced|dense>"\n'
        "      },\n"
        '      "reuse_template_slide_idx": <integer or null>,\n'
        '      "transition_note": "<string>"\n'
        "    },\n"
        "    ...\n"
        "  ]\n"
        "}\n"
    ) % (
        user_prompt,
        brand_context_section,
        visual_profile_section,
        brand_search_guidance,
    )

    # Strongly enforce valid JSON output and completion for all models (especially Gemini)
    optimizer_prompt += (
        "\n\nCRITICAL MUST FOLLOW RULE: You MUST return a single, complete, properly formatted JSON object. "
        "Do not truncate your response. Do not output anything before or after the JSON. "
        "Make sure every list and object is properly enclosed with closing brackets!"
    )

    prompt_file = _save_prompt_to_file(
        optimizer_prompt, "optimize_and_plan", output_dir
    )
    if prompt_file:
        print("[PROMPT] Optimizer prompt saved to: %s" % prompt_file)

    try:
        response = None
        from agents import get_agents as _get_agents  # type: ignore

        _query_optimizer = _get_agents(session_state.get("llm_provider", "claude")).get(
            "query_optimizer"
        )
        provider = getattr(
            getattr(_query_optimizer, "model", None),
            "provider",
            session_state.get("llm_provider", "claude"),
        )
        if provider == "Anthropic":
            # Fallback to claude-haiku-4-5 if attribute is missing
            actual_model_id = getattr(
                getattr(_query_optimizer, "model", None), "id", "claude-haiku-4-5"
            )
        else:
            actual_model_id = getattr(
                getattr(_query_optimizer, "model", None), "id", "claude-haiku-4-5"
            )
        _log_agent_banner(
            agent_name=getattr(_query_optimizer, "name", "Presentation Strategist"),
            model_id=actual_model_id,
            provider=getattr(
                getattr(_query_optimizer, "model", None),
                "provider",
                session_state.get("llm_provider", "claude"),
            ),  # type: ignore
            step_name="step_optimize_and_plan / Storyboard Generation",
        )
        # Check rate limits before calling
        _get_rate_tracker().check_and_wait(
            model=actual_model_id,
            prompt=optimizer_prompt,
            caller="step_optimize_and_plan/query_optimizer",
        )
        response = _query_optimizer.run(optimizer_prompt, stream=False)
    except Exception as e:
        error_msg = str(e)
        status_code = getattr(getattr(e, "response", None), "status_code", None)

        # Check for API capacity, credit limits, or rate limits
        is_capacity_error = (
            status_code in (400, 429, 529, 500, 503)
            or "credit" in error_msg.lower()
            or "rate limit" in error_msg.lower()
            or "overloaded" in error_msg.lower()
            or "429" in error_msg
            or "529" in error_msg
            or "status 400" in error_msg.lower()
        )

        if is_capacity_error:
            provider = session_state.get("llm_provider", "claude")
            print(
                "\n[FALLBACK AGENT ENGAGED] Primary provider (%s) hit capacity/credit error (%s) on query_optimizer."
                % (provider, status_code or "Unknown")
            )
            print(
                "[FALLBACK TRIGGERED] Engaging fallback agent for storyboard generation..."
            )
            try:
                from agents import get_agents as _get_agents  # type: ignore

                _fallback_query_optimizer = _get_agents(provider).get(
                    "query_optimizer_fallback"
                )

                if not _fallback_query_optimizer:
                    raise ValueError("No fallback defined for query_optimizer")

                _log_agent_banner(
                    agent_name=getattr(
                        _fallback_query_optimizer,
                        "name",
                        "Presentation Strategist (Fallback)",
                    ),
                    model_id=getattr(
                        getattr(_fallback_query_optimizer, "model", None),
                        "id",
                        "unknown",
                    ),
                    provider=getattr(
                        getattr(_fallback_query_optimizer, "model", None),
                        "provider",
                        "Fallback",
                    ),
                    step_name="step_optimize_and_plan / Storyboard Generation (Fallback)",
                )

                response = _fallback_query_optimizer.run(optimizer_prompt, stream=False)

            except Exception as fallback_e:
                print(
                    "[ERROR] Fallback query optimizer failed during exception fallback: %s"
                    % str(fallback_e)
                )
                if VERBOSE:  # noqa: F405
                    traceback.print_exc()
                return StepOutput(
                    content="Query optimization failed on both primary and fallback: %s"
                    % str(fallback_e),
                    success=False,
                )
        else:
            print("[ERROR] Query optimizer failed: %s" % str(e))
            if VERBOSE:  # noqa: F405
                traceback.print_exc()
            return StepOutput(
                content="Query optimization failed: %s" % str(e), success=False
            )

    # Agno catches API 400/429 limits internally and yields no response, bypassing exceptions.
    # Therefore, we must also apply the fallback logic if response is None after execution.
    if response is None:
        provider = session_state.get("llm_provider", "claude")
        print(
            "\n[FALLBACK AGENT ENGAGED] Primary provider (%s) produced no output (likely hit capacity/credit error)."
            % provider
        )
        print(
            "[FALLBACK TRIGGERED] Engaging fallback agent for storyboard generation..."
        )
        try:
            from agents import get_agents as _get_agents  # type: ignore

            _fallback_query_optimizer = _get_agents(provider).get(
                "query_optimizer_fallback"
            )

            if not _fallback_query_optimizer:
                raise ValueError("No fallback defined for query_optimizer")

                _log_agent_banner(
                    agent_name=getattr(
                        _fallback_query_optimizer,
                        "name",
                        "Presentation Strategist (Fallback)",
                    ),
                    model_id=getattr(
                        getattr(_fallback_query_optimizer, "model", None),
                        "id",
                        "unknown",
                    ),
                    provider=getattr(
                        getattr(_fallback_query_optimizer, "model", None),
                        "provider",
                        "Fallback",
                    ),
                    step_name="step_optimize_and_plan / Storyboard Generation (Fallback)",
                )

            response = _fallback_query_optimizer.run(optimizer_prompt, stream=False)

        except Exception as fallback_e:
            print("[ERROR] Fallback query optimizer failed: %s" % str(fallback_e))
            if VERBOSE:  # noqa: F405
                traceback.print_exc()
            return StepOutput(
                content="Query optimization failed on both primary and fallback: %s"
                % str(fallback_e),
                success=False,
            )

    # Parse the StoryboardPlan from response.
    # Without output_schema the agent returns plain text; extract JSON from it.
    plan: Optional[StoryboardPlan] = None
    if response and response.content:
        content = response.content
        if isinstance(content, StoryboardPlan):
            plan = content
        elif hasattr(content, "model_dump"):
            try:
                plan = StoryboardPlan(**content.model_dump())
            except Exception as e:
                print("[ERROR] Failed to parse StoryboardPlan from model: %s" % e)
        elif isinstance(content, dict):
            try:
                plan = StoryboardPlan(**content)
            except Exception as e:
                print("[ERROR] Failed to parse StoryboardPlan from dict: %s" % e)
        elif isinstance(content, str):
            # Strip markdown code fences if present (```json ... ``` or ``` ... ```)
            import re as _re

            json_text = content.strip()
            fence_match = _re.search(r"```(?:json)?\s*([\s\S]+?)```", json_text)
            if fence_match:
                json_text = fence_match.group(1).strip()
            # Locate the outermost JSON object in the text
            obj_match = _re.search(r"\{[\s\S]+\}", json_text)
            if obj_match:
                json_text = obj_match.group(0)
            try:
                plan = StoryboardPlan.model_validate_json(json_text)
            except Exception as e:
                print(
                    f"[WARN] Failed to parse StoryboardPlan from JSON string. Initiating fallback... ({e})"
                )
                if VERBOSE:  # noqa: F405
                    print(
                        "[VERBOSE] Raw optimizer response (first 2000 chars):\n%s"
                        % content[:2000]  # type: ignore
                    )

    if not plan:
        provider = session_state.get("llm_provider", "claude")
        print(
            "\n[FALLBACK TRIGGERED] Primary provider (%s) produced invalid or truncated JSON."
            % provider
        )
        print(
            "[FALLBACK TRIGGERED] Engaging fallback agent for storyboard generation..."
        )
        try:
            from agents import get_agents as _get_agents  # type: ignore

            _fallback_query_optimizer = _get_agents(provider).get(
                "query_optimizer_fallback"
            )

            if not _fallback_query_optimizer:
                raise ValueError("No fallback defined for query_optimizer")

            _log_agent_banner(
                agent_name=getattr(
                    _fallback_query_optimizer,
                    "name",
                    "Presentation Strategist (Fallback)",
                ),
                model_id=getattr(
                    getattr(_fallback_query_optimizer, "model", None), "id", "unknown"
                ),
                provider=getattr(
                    getattr(_fallback_query_optimizer, "model", None),
                    "provider",
                    "Fallback",
                ),
                step_name="step_optimize_and_plan / Storyboard Generation (JSON Fallback)",
            )

            fb_resp = _fallback_query_optimizer.run(optimizer_prompt, stream=False)
            if fb_resp and fb_resp.content:
                fb_content = fb_resp.content
                if isinstance(fb_content, StoryboardPlan):
                    plan = fb_content
                elif hasattr(fb_content, "model_dump"):
                    plan = StoryboardPlan(**fb_content.model_dump())
                elif isinstance(fb_content, dict):
                    plan = StoryboardPlan(**fb_content)
                elif isinstance(fb_content, str):
                    import re as _re

                    json_text = fb_content.strip()
                    fence_match = _re.search(r"```(?:json)?\s*([\s\S]+?)```", json_text)
                    if fence_match:
                        json_text = fence_match.group(1).strip()
                    obj_match = _re.search(r"\{[\s\S]+\}", json_text)
                    if obj_match:
                        json_text = obj_match.group(0)
                    plan = StoryboardPlan.model_validate_json(json_text)
        except Exception as fb_e:
            print("[ERROR] Fallback query optimizer failed on JSON fallback: %s" % fb_e)

    if not plan:
        print("[ERROR] No valid storyboard plan produced.")
        return StepOutput(content="No storyboard plan produced.", success=False)

    print(
        "Storyboard plan: '%s' (%d slides, tone: %s)"
        % (plan.presentation_title, plan.total_slides, plan.tone)
    )

    if VERBOSE:  # noqa: F405
        print("[VERBOSE] Full storyboard JSON:\n%s" % plan.model_dump_json(indent=2))

    # === LAYOUT CONSTRAINTS ENRICHMENT FROM VISUAL PROFILE ===
    # If a template visual profile exists, override default layout_constraints
    # with profile-derived values for slides whose constraints are still generic.
    if visual_profile and visual_profile.slide_count > 0:
        enriched_count = 0
        profile_top_pct = int(
            max(
                8.0,
                float(
                    100.0
                    - visual_profile.avg_content_zone_height_pct
                    - (100.0 - visual_profile.avg_content_zone_width_pct) / 2.0
                ),
            )
        )
        profile_bottom_pct = int(
            min(92.0, float(visual_profile.avg_content_zone_height_pct + 12.0))
        )
        profile_max_blocks = min(6, visual_profile.max_comfortable_bullets)
        profile_text_w = visual_profile.recommended_text_weight

        for slide in plan.slides:
            lc = slide.layout_constraints
            if lc is None:
                slide.layout_constraints = LayoutConstraints(
                    content_zone_top_pct=profile_top_pct,
                    content_zone_bottom_pct=profile_bottom_pct,
                    max_content_blocks=profile_max_blocks,
                    text_weight=profile_text_w,
                )
                enriched_count += 1
            else:
                # Only override fields that still have their generic defaults
                changed = False
                if lc.content_zone_top_pct == 12 and profile_top_pct != 12:
                    lc.content_zone_top_pct = profile_top_pct
                    changed = True
                if lc.content_zone_bottom_pct == 88 and profile_bottom_pct != 88:
                    lc.content_zone_bottom_pct = profile_bottom_pct
                    changed = True
                if lc.max_content_blocks == 4 and profile_max_blocks != 4:
                    lc.max_content_blocks = profile_max_blocks
                    changed = True
                if lc.text_weight == "balanced" and profile_text_w != "balanced":
                    lc.text_weight = profile_text_w
                    changed = True
                if changed:
                    enriched_count += 1

        if enriched_count > 0:
            print(
                "[VISUAL PROFILE] Enriched layout_constraints for %d/%d slides "
                "(top=%d%%, bottom=%d%%, max_blocks=%d, text_weight=%s)"
                % (
                    enriched_count,
                    len(plan.slides),
                    profile_top_pct,
                    profile_bottom_pct,
                    profile_max_blocks,
                    profile_text_w,
                )
            )
        if VERBOSE:  # noqa: F405
            print(
                "[VERBOSE] [VISUAL PROFILE] Layout enrichment details: "
                "profile_top=%d, profile_bottom=%d, profile_max_blocks=%d, "
                "profile_text_weight=%s, slides_enriched=%d/%d"
                % (
                    profile_top_pct,
                    profile_bottom_pct,
                    profile_max_blocks,
                    profile_text_w,
                    enriched_count,
                    len(plan.slides),
                )
            )

    # Save global context markdown
    global_context_path = os.path.join(storyboard_dir, "global_context.md")
    with open(global_context_path, "w", encoding="utf-8") as f:
        f.write(_format_global_context_markdown(plan))
    print("Saved global context: %s" % global_context_path)

    # Save per-slide storyboard markdown files
    for slide in plan.slides:
        slide_path = os.path.join(storyboard_dir, "slide_%03d.md" % slide.slide_number)
        slide_md = _format_slide_markdown(slide)
        if VERBOSE:  # noqa: F405
            print("[VERBOSE] Slide %d storyboard:\n%s" % (slide.slide_number, slide_md))
        with open(slide_path, "w", encoding="utf-8") as f:
            f.write(slide_md)
    print("Saved %d slide storyboard files to: %s" % (len(plan.slides), storyboard_dir))

    # Store in session_state
    session_state["storyboard"] = plan
    session_state["total_slides"] = plan.total_slides
    session_state["storyboard_dir"] = storyboard_dir
    session_state["chunk_size"] = chunk_size
    session_state["max_retries"] = max_retries

    step_elapsed = time.time() - step_start
    print("[TIMING] step_optimize_and_plan completed in %.1fs" % step_elapsed)

    summary = (
        "Storyboard created: '%s' | %d slides | tone: %s | brand voice: %s | chunk size: %d | Duration: %.1fs"
    ) % (
        plan.presentation_title,
        plan.total_slides,
        plan.tone,
        plan.brand_voice,
        chunk_size,
        step_elapsed,
    )
    return StepOutput(content=summary, success=True)


CHUNK_TIMEOUT_SECONDS = 300  # 5 minutes


def _run_chunk_agent(chunk_agent, chunk_prompt):
    """Run the chunk agent with streaming and collect the final RunOutput.

    Intended to be submitted to a ThreadPoolExecutor so the caller can enforce a
    wall-clock timeout via Future.result(timeout=...).

    Args:
        chunk_agent: Configured Agent instance to run.
        chunk_prompt: Full prompt string to send to the agent.

    Returns:
        Tuple of (RunOutput or None, int) where the second element is the total
        number of streaming events received. RunOutput is None if no RunOutput
        event was emitted during the stream.
    """
    response = None
    event_count = 0
    for event in chunk_agent.run(chunk_prompt, stream=True, yield_run_output=True):
        event_count += 1
        if isinstance(event, RunOutput):
            response = event
    return response, event_count


# === HELPER: GENERATE A SINGLE CHUNK VIA CLAUDE PPTX SKILL ===


def generate_chunk_pptx(
    chunk_slides: List[SlideStoryboard],
    session_state: Dict,
    chunk_idx: int,
) -> Optional[str]:
    """Call the Claude PPTX skill for a chunk of slides with retry logic (Tier 1).

    Creates a fresh agent per call (not reused across chunks) and applies
    exponential backoff on retries.

    Applies a 300-second (CHUNK_TIMEOUT_SECONDS) wall-clock timeout per attempt
    via ThreadPoolExecutor. On timeout, activates the session-level fallback flag
    and returns None immediately (no further retries for this chunk).

    *Progressive Fallback on Throttling*: On the first throttling/capacity error
    (429 rate limit, 529 overloaded, status-200 overloaded in body, 500 api_error,
    402 billing_error), remaining Opus retries are skipped and the function
    immediately tries a progressive fallback chain (single attempt each):
      1. Sonnet PPTX Skill — same capability, different GPU pool
      2. Gemini code-gen   — separate provider (Google)
      3. OpenAI code-gen   — separate provider (GPT-5.4)
    If all fail, None is returned and the caller drops to Tier 2.

    For non-throttling errors, the function retries up to max_retries attempts,
    then tries OpenAI code-gen before returning None.

    Throttling detection uses the shared ``_is_throttling_error()`` helper.

    Brand/style context from session_state["brand_style_intent"] is injected into
    the chunk prompt as a '## Brand/Style Guidance' section when branding is present.

    Args:
        chunk_slides: List of SlideStoryboard objects for this chunk.
        session_state: Shared workflow session state (must contain brand_style_intent
                       key set by step_optimize_and_plan).
        chunk_idx: 0-based chunk index (used for file naming and logging).

    Returns:
        Path to the generated chunk PPTX file, or None if all attempts failed.
        When None is returned, session_state["use_fallback_generator"] is set to
        True, causing all subsequent chunks to bypass Tier 1 and use Tier 2/3.
    """
    storyboard: StoryboardPlan = session_state["storyboard"]
    storyboard_dir = session_state["storyboard_dir"]
    output_dir = session_state["output_dir"]
    max_retries = session_state.get("max_retries", 2)

    # Load global context
    global_context_path = os.path.join(storyboard_dir, "global_context.md")
    global_context = ""
    if os.path.exists(global_context_path):
        with open(global_context_path, encoding="utf-8") as f:
            global_context = f.read()

    # Load per-slide markdown for this chunk
    slide_details = []
    for s in chunk_slides:
        md_path = os.path.join(storyboard_dir, "slide_%03d.md" % s.slide_number)
        if os.path.exists(md_path):
            with open(md_path, encoding="utf-8") as f:
                slide_details.append(str(f.read()))  # type: ignore
        else:
            # Fallback: format inline
            slide_details.append(str(_format_slide_markdown(s)))  # type: ignore

    first_slide = chunk_slides[0].slide_number
    last_slide = chunk_slides[-1].slide_number

    chunk_prompt = (
        "## Global Presentation Context\n"
        "%s\n\n"
        "## Task: Generate slides %d through %d of %d\n\n"
        "You are generating a CHUNK of a larger presentation. "
        "This chunk contains %d slides.\n"
        "Maintain the presentation's tone (%s) and brand voice (%s).\n"
        "Design for the PRIMARY AUDIENCE: %s (Rule 1, RULES.md).\n"
        'These are slides %d-%d of the full %d-slide deck titled "%s".\n\n'
        "## Per-Slide Content for This Chunk:\n\n"
        "%s\n\n"
        "%s"
        "%s"
        "Please generate EXACTLY %d slides for this chunk with the content described above.\n"
        "Do not add extra slides. Do not include slide numbers outside the range %d-%d.\n"
        "Use clean formatting without custom fonts or colors. "
        "Include native data-vis (tables/charts/infographics/diagrams) only where explicitly suggested.\n"
        "For any chart: use native PPTX chart objects only (bar, column, line, or pie) — "
        "do NOT use matplotlib or embed chart images.\n"
        "For any table: use native PPTX table objects only — do NOT embed a table as an image "
        "or use matplotlib/PIL to render one.\n"
        "For infographics or diagrams: use native PowerPoint shapes (rectangles, arrows, text boxes) "
        "or a native table to approximate the visual — do NOT insert images for infographics or diagrams.\n"
        "SEMANTIC LAYOUT RULES: If the slide Semantic Type is 'sequential', 'comparative', 'metrics', or 'hero', "
        "DO NOT use standard bullet points. Instead, create native shapes (chevrons, cards, metric grids, banners) "
        "that visually represent that semantic classification.\n"
        "If a requested visual cannot be represented exactly, preserve the slide structure and add a concise "
        "native textbox note; do NOT fail the chunk and do NOT use image-based substitutes.\n\n"
        "STRICT VISUAL RULES (MANDATORY — ZERO TOLERANCE):\n"
        "1. ZERO TEXT OVERLAP: No text shape may overlap any other text shape. "
        "Place each text block in its own clear space with at least 10px gap between shapes. "
        "If space is insufficient, REDUCE the amount of content — do NOT overlap.\n"
        "2. NO TINY TEXT: Every text element must use at least 14pt font size. "
        "If content does not fit at 14pt, shorten the text — do NOT shrink the font below 14pt. "
        "Never create text boxes smaller than 1 inch wide or 0.4 inches tall.\n"
        "3. STAY WITHIN SLIDE BOUNDS: Every shape (text, chart, table, infographic) "
        "must be fully contained within slide dimensions (0,0 to slide_width, slide_height). "
        "No shape edge may exceed these boundaries.\n"
        "4. MATCH TEMPLATE FONTS: Use ONLY the template's font family for all text. "
        "Do NOT introduce Arial, Calibri, Helvetica, or any other font unless it IS the template font. "
        "Use prs.slide_masters[0] font theme if available.\n"
        "5. CONTENT ZONE: Place all content shapes between 12%%-88%% of slide height. "
        "Reserve the top for the title area and the bottom for footer/page number.\n"
        "6. SHAPE DENSITY: Create AT MOST 4 content shapes per slide (excluding title and footer). "
        "Fewer shapes means cleaner layout and zero overlap risk.\n\n"
        "TEMPLATE FIDELITY RULES (when template is provided):\n"
        "7. HEADER/FOOTER & LOGO RULES: The template's built-in slide layouts already contain footers, slide numbers, and logos. "
        "DO NOT generate new shapes for them, as they will overlap, clip, or look misaligned. Keep the bottom footer area free of content shapes.\n"
        "8. LOGO REPLACEMENT: If the Brand/Style Guidance provides a brand name, you MUST replace the template's default Logo image. "
        "To achieve this, create a solid `RECTANGLE` matching the slide background color, place it over the original logo position "
        "(usually top-left or top-right) to hide it, and then add a text box with the brand name on top.\n"
        "9. SHAPES AND VISUAL ELEMENTS: Study the template's decorative shapes (lines, rectangles, "
        "accent bars, chevrons). When creating content shapes (cards, callout boxes, metric panels), "
        "use similar border radius, fill colors, and positioning patterns as the template.\n"
        "10. CHARTS AND INFOGRAPHICS: When creating charts, use the template's theme colors "
        "for series fills. Match the template's chart style (flat vs. 3D, border vs. borderless, "
        "data label position). If the template has no charts, use the template's accent colors.\n"
        "11. COLOR THEME ADHERENCE: Use the template's color palette for ALL colored elements — "
        "shape fills, text accents, chart series, table headers. If a template color is unsuitable, "
        "pick the closest matching color from the palette — do NOT introduce arbitrary colors.\n"
        "12. SMARTART AND DIAGRAMS: If the template contains SmartArt or process diagrams AND "
        "it fits or is relevant to the content, try to mimic that visual language using native shapes. "
        "If the SmartArt style is not relevant to the content, use simpler shapes with the same color theme.\n\n"
        "Save the output as 'chunk_%03d.pptx'."
    ) % (
        global_context,
        first_slide,
        last_slide,
        storyboard.total_slides,
        len(chunk_slides),
        storyboard.tone,
        storyboard.brand_voice,
        session_state.get("brand_style_intent", BrandStyleIntent()).target_audience,
        first_slide,
        last_slide,
        storyboard.total_slides,
        storyboard.presentation_title,
        "\n\n---\n\n".join(slide_details),
        _format_brand_context_for_prompt(
            session_state.get("brand_style_intent", BrandStyleIntent())
        ),
        _build_no_template_design_system(
            storyboard.visual_style if storyboard else "clean_minimal",
            session_state.get("brand_style_intent"),
        )
        if not session_state.get("template_path")
        else "",
        len(chunk_slides),
        first_slide,
        last_slide,
        chunk_idx,
    )

    # Append template visual reference section (if template was rendered)
    template_pngs = session_state.get("template_slide_pngs", {})
    if template_pngs and session_state.get("template_visuals"):
        visual_ref = _build_visual_reference_section(
            chunk_slides,
            template_pngs,
            brand_style_intent=session_state.get("brand_style_intent"),
        )
        if visual_ref:
            chunk_prompt += "\n" + visual_ref
            if VERBOSE:  # noqa: F405
                print(
                    "[VERBOSE] Chunk %d: appended %d-char visual reference to prompt."
                    % (chunk_idx, len(visual_ref))
                )

    prompt_file = _save_prompt_to_file(
        chunk_prompt, "chunk", output_dir, "chunk_%03d" % chunk_idx
    )
    if prompt_file:
        print("[PROMPT] Chunk %d prompt saved to: %s" % (chunk_idx, prompt_file))

    chunk_output_path = os.path.join(output_dir, "chunk_%03d.pptx" % chunk_idx)

    # Create a fresh agent per chunk call — do NOT reuse across calls.
    # betas + max_tokens=128000: the PPTX skill generates full multi-slide decks whose
    # output can be large; 128k tokens and the context-1m beta are both safe here because
    # this agent is always invoked with stream=True (see _run_chunk_agent).
    chunk_agent = Agent(
        name="Chunk Generator %d" % chunk_idx,
        model=Claude(
            id="claude-haiku-4-5",
            betas=["context-1m-2025-08-07"],
            skills=[{"type": "anthropic", "skill_id": "pptx", "version": "latest"}],
            max_tokens=128000,
        ),
        instructions=[
            "You are a structured content generator for PowerPoint presentations.",
            "Generate EXACTLY the number of slides specified in the task.",
            "Use one clear title per slide with concise bullet points.",
            "Do NOT apply custom fonts, colors, or theme styling.",
            "Do NOT add animations, transitions, or speaker notes.",
            "Keep tables to max 6 rows x 5 columns.",
            "For charts: use only native PPTX chart objects (bar, column, line, or pie). Do NOT use matplotlib, PIL, or any image-based approach for charts.",
            "For any chart mentioned in a visual_suggestion: produce a native Office chart with synthesized data — never embed a chart as an image.",
            "For tables: use ONLY native PPTX table objects — never embed a table as an image or use matplotlib/PIL to render one.",
            "For infographics or diagrams: use native PowerPoint shapes (rectangles, arrows, text boxes) or a native table to approximate the visual — never insert an image for an infographic or diagram.",
            "SEMANTIC LAYOUT RULES: If the slide Semantic Type is 'sequential', 'comparative', 'metrics', or 'hero', "
            "DO NOT use standard bullet points. Instead, create native shapes (chevrons, cards, metric grids, banners) "
            "that visually represent that semantic classification.\n"
            "If a requested visual cannot be represented exactly, preserve the slide structure and add a concise "
            "native textbox note; do NOT fail the chunk and do NOT use image-based substitutes.",
        ],
        markdown=True,
    )

    client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

    for attempt in range(max_retries + 1):
        attempt_start = time.time()

        if attempt > 0:
            delay_secs = 60.0 + random.uniform(
                0, 30.0
            )  # 60-90s jitter to respect 30K tokens/min window
            print(
                f"[CHUNK {chunk_idx}] Retry {attempt}/{max_retries} — cooling down for {delay_secs:.0f}s (rate limit window reset)..."
            )
            _countdown_sleep(delay_secs, label=f"[CHUNK {chunk_idx} RETRY]")  # type: ignore

        print(
            f"[CHUNK {chunk_idx}] API call attempt {attempt + 1}/{max_retries + 1} (slides {first_slide}-{last_slide})..."
        )

        _log_agent_banner(
            agent_name=f"Chunk Generator {chunk_idx}",  # type: ignore
            model_id="claude-haiku-4-5",
            provider=str(
                getattr(getattr(chunk_agent, "model", None), "provider", "Anthropic")
            ),  # type: ignore
            step_name=f"step_generate_chunks / Tier 1 PPTX Skill (attempt {attempt + 1}/{max_retries + 1})",
        )

        # Register this call with the rate tracker before executing.
        # check_and_wait will auto-sleep if the 30K token/min window would be exceeded.
        _get_rate_tracker().check_and_wait(
            model="claude-haiku-4-5",
            prompt=chunk_prompt,
            caller="generate_chunk_pptx/Tier1",
        )

        try:
            response = None
            event_count = 0
            timed_out = False
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
                future = executor.submit(_run_chunk_agent, chunk_agent, chunk_prompt)  # type: ignore
                try:
                    response, event_count = future.result(timeout=CHUNK_TIMEOUT_SECONDS)
                except concurrent.futures.TimeoutError:
                    print(
                        f"[CHUNK {chunk_idx}] Attempt {attempt + 1}/{max_retries + 1} timed out after {CHUNK_TIMEOUT_SECONDS}s. Activating fallback generator."
                    )
                    session_state["use_fallback_generator"] = True
                    timed_out = True
                except Exception as e:
                    raise

            if timed_out:
                return None

            if response is None:
                print(
                    f"[CHUNK {chunk_idx}] No RunOutput received after {event_count} events."
                )
                attempt_elapsed = time.time() - attempt_start
                print(
                    f"[TIMING] Chunk {chunk_idx} attempt {attempt + 1}/{max_retries + 1}: {attempt_elapsed:.1f}s (no output)"
                )
                continue

            msg_list = getattr(response, "messages", [])
            has_messages = bool(msg_list)

            if VERBOSE:  # noqa: F405
                msg_count = len(msg_list) if has_messages else 0
                print(
                    f"[VERBOSE] Chunk {chunk_idx} attempt {attempt + 1}: received {event_count} events, {msg_count} messages"
                )
                if has_messages:
                    for m_idx, msg in enumerate(msg_list):
                        print(
                            f"[VERBOSE] Chunk {chunk_idx} message {m_idx}: type={type(msg).__name__} role={getattr(msg, 'role', 'N/A')} has_provider_data={bool(getattr(msg, 'provider_data', None))}"
                        )

        except Exception as e:
            err_str = str(e)
            attempt_elapsed = time.time() - attempt_start
            is_throttling = _is_throttling_error(err_str)

            if is_throttling:
                # Transient capacity/rate error — mark it on session_state but
                # do NOT permanently activate use_fallback_generator yet.  The
                # inter-chunk delay in step_generate_chunks will apply max_delay
                # before the next chunk.
                print(
                    "[CHUNK %d] Attempt %d/%d hit rate/capacity limit: %s"
                    % (chunk_idx, attempt + 1, max_retries + 1, err_str[:200])  # type: ignore
                )
                session_state["rate_limit_hit"] = True
            else:
                print(
                    "[CHUNK %d] Attempt %d/%d failed with error: %s"
                    % (chunk_idx, attempt + 1, max_retries + 1, err_str)
                )

            print(
                "[TIMING] Chunk %d attempt %d/%d: %.1fs (%s)"
                % (
                    chunk_idx,
                    attempt + 1,
                    max_retries + 1,
                    attempt_elapsed,
                    "rate/capacity limit" if is_throttling else "error",
                )
            )

            # =================================================================
            # PROGRESSIVE FALLBACK CHAIN (on throttling/overload)
            #
            # When any throttling/capacity error is detected we immediately
            # skip remaining Opus retries and try (single attempt each):
            #   1. Sonnet (PPTX skill) — same capability, different GPU pool
            #   2. Gemini code-gen     — separate provider (Google)
            #   3. OpenAI code-gen     — separate provider (GPT-5.4)
            # If all fail we return None and the caller drops to Tier 2.
            # =================================================================
            if is_throttling:
                print(
                    "[CHUNK %d] Throttling detected. Engaging progressive "  # type: ignore
                    "fallback chain (Sonnet → Gemini → OpenAI)..." % chunk_idx
                )

                # --- Fallback 1: Sonnet PPTX Skill (same capability, smaller model) ---
                print(
                    "[CHUNK %d SONNET-FALLBACK] Attempting Claude Sonnet "  # type: ignore
                    "PPTX skill fallback..." % chunk_idx
                )
                try:
                    chunk_agent_sonnet = Agent(
                        name="Chunk Generator %d (Haiku Fallback)" % chunk_idx,  # type: ignore
                        model=Claude(
                            id="claude-haiku-4-5",
                            betas=["context-1m-2025-08-07"],
                            skills=[
                                {
                                    "type": "anthropic",
                                    "skill_id": "pptx",
                                    "version": "latest",
                                }
                            ],
                            max_tokens=128000,
                        ),
                        instructions=chunk_agent.instructions,  # type: ignore
                        markdown=True,
                    )
                    _log_agent_banner(
                        agent_name="Chunk Generator %d (Haiku Fallback)"  # type: ignore
                        % chunk_idx,
                        model_id="claude-haiku-4-5",
                        provider="Anthropic",
                        step_name="step_generate_chunks / Tier 1 Haiku "  # type: ignore
                        "Fallback (chunk %d)" % chunk_idx,
                    )
                    _get_rate_tracker().check_and_wait(
                        model="claude-haiku-4-5",
                        prompt=chunk_prompt,
                        caller="generate_chunk_pptx/Tier1-Haiku-Fallback",
                    )
                    t1_sonnet_start = time.time()
                    with concurrent.futures.ThreadPoolExecutor(
                        max_workers=1
                    ) as sonnet_executor:
                        sonnet_future = sonnet_executor.submit(
                            _run_chunk_agent,
                            chunk_agent_sonnet,
                            chunk_prompt,  # type: ignore
                        )
                        try:
                            sonnet_response, sonnet_events = sonnet_future.result(
                                timeout=CHUNK_TIMEOUT_SECONDS
                            )
                        except concurrent.futures.TimeoutError:
                            print(
                                "[CHUNK %d SONNET-FALLBACK] Timed out after "
                                "%ds." % (chunk_idx, CHUNK_TIMEOUT_SECONDS)
                            )
                            sonnet_response = None

                    if sonnet_response and sonnet_response.messages:  # type: ignore
                        sonnet_file = None
                        for msg in sonnet_response.messages:  # type: ignore
                            if hasattr(msg, "provider_data") and msg.provider_data:
                                try:
                                    files = download_skill_files(
                                        msg.provider_data,
                                        client,
                                        output_dir=output_dir,
                                    )
                                    for f in files:
                                        if not f.endswith(".pptx"):
                                            continue
                                        try:
                                            Presentation(f)
                                            sonnet_file = f
                                            break
                                        except Exception:
                                            continue
                                except Exception:
                                    pass
                            if sonnet_file:
                                break

                        # Also try response-level model_provider_data
                        if (
                            not sonnet_file
                            and hasattr(sonnet_response, "model_provider_data")
                            and sonnet_response.model_provider_data  # type: ignore
                        ):
                            try:
                                files = download_skill_files(
                                    sonnet_response.model_provider_data,  # type: ignore
                                    client,
                                    output_dir=output_dir,
                                )
                                for f in files:
                                    if not f.endswith(".pptx"):
                                        continue
                                    try:
                                        Presentation(f)
                                        sonnet_file = f
                                        break
                                    except Exception:
                                        continue
                            except Exception:
                                pass

                        if sonnet_file and os.path.exists(sonnet_file):
                            # Normalize name + cleanup
                            if sonnet_file != chunk_output_path:
                                shutil.copy2(sonnet_file, chunk_output_path)
                            try:
                                prs = Presentation(chunk_output_path)
                                clean_presentation_visual_noise_and_contrast(prs)
                                sanitize_presentation(prs)  # noqa: F405
                                prs.save(chunk_output_path)
                            except Exception as clnp_e:
                                print(
                                    "[CHUNK %d] Sonnet cleanup failed: %s"
                                    % (chunk_idx, clnp_e)
                                )
                            print(
                                "[TIMING] Chunk %d Sonnet PPTX fallback: "
                                "%.1fs (success)"
                                % (chunk_idx, time.time() - t1_sonnet_start)
                            )
                            print(
                                "[CHUNK %d SONNET-FALLBACK] Successfully "
                                "generated: %s" % (chunk_idx, chunk_output_path)
                            )
                            return chunk_output_path

                    print(
                        "[CHUNK %d SONNET-FALLBACK] Sonnet produced no "
                        "valid file." % chunk_idx
                    )
                except Exception as e_sonnet:
                    print(
                        "[CHUNK %d SONNET-FALLBACK] Sonnet PPTX skill "
                        "failed: %s" % (chunk_idx, str(e_sonnet)[:200])  # type: ignore
                    )

                # --- Fallback 2: Gemini code-gen (separate provider) ---
                print(
                    "[CHUNK %d GEMINI-FALLBACK] Attempting Gemini code-gen "
                    "fallback..." % chunk_idx
                )
                try:
                    from agents import get_agents as _get_agents  # type: ignore

                    _all_agents = _get_agents(
                        session_state.get("llm_provider", "claude")
                    )
                    _gemini_agent = _all_agents.get("fallback_code_agent_fallback")
                    if _gemini_agent is not None:
                        # Build a code-gen prompt (reuse Tier 2 approach)
                        _gemini_prompt = (
                            "Generate a complete Python script using "
                            "python-pptx to create a PowerPoint file at "
                            "this EXACT path: %s\n\n"
                            "Use the following chunk prompt as content "
                            "guidance:\n%s\n\n"
                            "Execute the script using "
                            "save_to_file_and_run.\n"
                        ) % (chunk_output_path, chunk_prompt[:8000])  # type: ignore
                        _log_agent_banner(
                            agent_name=getattr(
                                _gemini_agent,
                                "name",
                                "PPTX Code Generator (Gemini Fallback)",
                            ),
                            model_id=getattr(
                                getattr(_gemini_agent, "model", None),
                                "id",
                                "gemini-3.1-pro-preview",
                            ),
                            provider="Google",
                            step_name="step_generate_chunks / Tier 1 "
                            "Gemini Fallback (chunk %d)" % chunk_idx,
                        )
                        t1_gemini_start = time.time()
                        for _ in _gemini_agent.run(_gemini_prompt, stream=True):
                            pass

                        if os.path.exists(chunk_output_path):
                            try:
                                prs = Presentation(chunk_output_path)
                                clean_presentation_visual_noise_and_contrast(prs)
                                sanitize_presentation(prs)  # noqa: F405
                                prs.save(chunk_output_path)
                            except Exception as clnp_e:
                                print(
                                    "[CHUNK %d] Gemini cleanup failed: %s"
                                    % (chunk_idx, clnp_e)
                                )
                            print(
                                "[TIMING] Chunk %d Gemini fallback: %.1fs "
                                "(success)" % (chunk_idx, time.time() - t1_gemini_start)
                            )
                            print(
                                "[CHUNK %d GEMINI-FALLBACK] Successfully "
                                "generated: %s" % (chunk_idx, chunk_output_path)
                            )
                            return chunk_output_path
                        else:
                            print(
                                "[CHUNK %d GEMINI-FALLBACK] No file "
                                "produced." % chunk_idx
                            )
                    else:
                        print(
                            "[CHUNK %d GEMINI-FALLBACK] No Gemini fallback "
                            "agent available." % chunk_idx
                        )
                except Exception as e_gemini:
                    print(
                        "[CHUNK %d GEMINI-FALLBACK] Gemini code-gen "
                        "failed: %s" % (chunk_idx, str(e_gemini)[:200])  # type: ignore
                    )

                # --- Fallback 3: OpenAI code-gen (existing Universal Fallback) ---
                try:
                    from agents.fallback_openai_agents import (  # type: ignore
                        get_openai_fallback_agents,
                    )

                    _fallback_agents = get_openai_fallback_agents()
                    _universal_generator = _fallback_agents[
                        "fallback_content_generator"
                    ]

                    print(
                        "[OPENAI FALLBACK TRIGGERED] Chunk %d: Attempting "
                        "GPT-5.4 Pro Tier 1 Fallback..." % chunk_idx
                    )
                    t1_univ_start = time.time()
                    _log_agent_banner(
                        agent_name=_universal_generator.name,
                        model_id=_universal_generator.model.id,
                        provider=getattr(
                            _universal_generator.model, "provider", "openai"
                        ),
                        step_name="step_generate_chunks / Tier 1 Universal "
                        "Fallback (chunk %d)" % chunk_idx,
                    )
                    for _ in _universal_generator.run(chunk_prompt, stream=True):
                        pass

                    if os.path.exists(chunk_output_path):
                        print(
                            "[TIMING] Chunk %d Universal Tier 1 fallback: "
                            "%.1fs" % (chunk_idx, time.time() - t1_univ_start)
                        )
                        try:
                            prs = Presentation(chunk_output_path)
                            clean_presentation_visual_noise_and_contrast(prs)
                            sanitize_presentation(prs)  # noqa: F405
                            prs.save(chunk_output_path)
                        except Exception as clnp_e:
                            print(
                                "[CHUNK %d] Fallback Cleanup failed: %s"
                                % (chunk_idx, clnp_e)
                            )
                        print(
                            "[CHUNK %d TIER1-FALLBACK] Successfully "
                            "generated via Universal Fallback: %s"
                            % (chunk_idx, chunk_output_path)
                        )
                        return chunk_output_path

                except Exception as e_univ:
                    print(
                        "[CHUNK %d TIER1-FALLBACK] Universal OpenAI "
                        "Fallback failed: %s" % (chunk_idx, str(e_univ))
                    )

                # All fallbacks exhausted for this throttling error.
                # Break the retry loop — no point retrying Opus.
                print(
                    "[CHUNK %d] All Tier 1 fallbacks exhausted. "
                    "Dropping to Tier 2." % chunk_idx
                )
                break

            # Non-throttling error on final retry → try OpenAI fallback
            if attempt == max_retries:
                print(
                    "[CHUNK %d] All %d attempts failed (non-throttling). "
                    "Engaging Universal OpenAI Fallback..."
                    % (chunk_idx, max_retries + 1)
                )
                try:
                    from agents.fallback_openai_agents import (  # type: ignore
                        get_openai_fallback_agents,
                    )

                    _fallback_agents = get_openai_fallback_agents()
                    _universal_generator = _fallback_agents[
                        "fallback_content_generator"
                    ]

                    print(
                        "[OPENAI FALLBACK TRIGGERED] Chunk %d: Attempting "
                        "GPT-5.4 Pro Tier 1 Fallback..." % chunk_idx
                    )
                    t1_univ_start = time.time()
                    _log_agent_banner(
                        agent_name=_universal_generator.name,
                        model_id=_universal_generator.model.id,
                        provider=getattr(
                            _universal_generator.model, "provider", "openai"
                        ),
                        step_name="step_generate_chunks / Tier 1 Universal "
                        "Fallback (chunk %d)" % chunk_idx,
                    )
                    for _ in _universal_generator.run(chunk_prompt, stream=True):
                        pass

                    if os.path.exists(chunk_output_path):
                        print(
                            "[TIMING] Chunk %d Universal Tier 1 fallback: "
                            "%.1fs" % (chunk_idx, time.time() - t1_univ_start)
                        )
                        try:
                            prs = Presentation(chunk_output_path)
                            clean_presentation_visual_noise_and_contrast(prs)
                            sanitize_presentation(prs)  # noqa: F405
                            prs.save(chunk_output_path)
                        except Exception as clnp_e:
                            print(
                                "[CHUNK %d] Fallback Cleanup failed: %s"
                                % (chunk_idx, clnp_e)
                            )
                        print(
                            "[CHUNK %d TIER1-FALLBACK] Successfully "
                            "generated via Universal Fallback: %s"
                            % (chunk_idx, chunk_output_path)
                        )
                        return chunk_output_path

                except Exception as e_univ:
                    print(
                        "[CHUNK %d TIER1-FALLBACK] Universal OpenAI "
                        "Fallback failed: %s" % (chunk_idx, str(e_univ))
                    )

            continue

        # Try to download the generated file from message provider_data
        generated_file = None

        if response.messages:  # type: ignore
            for msg in response.messages:  # type: ignore
                if hasattr(msg, "provider_data") and msg.provider_data:
                    if VERBOSE:  # noqa: F405
                        print(
                            "[VERBOSE] Chunk %d: attempting file download from message provider_data..."  # type: ignore
                            % chunk_idx
                        )
                    try:
                        files = download_skill_files(
                            msg.provider_data, client, output_dir=output_dir
                        )
                    except Exception as e:
                        print(
                            "[CHUNK %d] download_skill_files (message) failed: %s"
                            % (chunk_idx, e)
                        )
                        files = []

                    if VERBOSE:  # noqa: F405
                        print(
                            "[VERBOSE] Chunk %d: download returned files: %s"
                            % (chunk_idx, files)
                        )

                    if files:
                        for f in files:
                            if not f.endswith(".pptx"):
                                continue
                            try:
                                Presentation(f)
                                generated_file = f
                                break
                            except Exception:
                                continue
                    if generated_file:
                        break

        # Fallback: try response.model_provider_data
        if (
            not generated_file
            and hasattr(response, "model_provider_data")
            and response.model_provider_data  # type: ignore
        ):
            if VERBOSE:  # noqa: F405
                print(
                    "[VERBOSE] Chunk %d: trying fallback model_provider_data download..."  # type: ignore
                    % chunk_idx
                )
            try:
                files = download_skill_files(
                    response.model_provider_data,
                    client,
                    output_dir=output_dir,  # type: ignore
                )
                if VERBOSE:  # noqa: F405
                    print(
                        "[VERBOSE] Chunk %d: fallback download returned files: %s"
                        % (chunk_idx, files)
                    )
                for f in files:
                    if not f.endswith(".pptx"):
                        continue
                    try:
                        Presentation(f)
                        generated_file = f
                        break
                    except Exception:
                        continue
            except Exception as e:
                print(
                    "[CHUNK %d] download_skill_files (fallback) failed: %s"
                    % (chunk_idx, e)
                )

        attempt_elapsed = time.time() - attempt_start

        if generated_file and os.path.exists(generated_file):
            # Normalize to standard chunk name
            if generated_file != chunk_output_path:
                shutil.copy2(generated_file, chunk_output_path)
                generated_file = chunk_output_path

            # Clean up empty placeholders and hardcoded contrast issues
            try:
                prs = Presentation(generated_file)
                clean_presentation_visual_noise_and_contrast(prs)
                sanitize_presentation(prs)  # noqa: F405
                prs.save(generated_file)
            except Exception as e:
                print("[CHUNK %d] Cleanup failed: %s" % (chunk_idx, e))

            print(
                "[TIMING] Chunk %d attempt %d/%d: %.1fs (success)"
                % (chunk_idx, attempt + 1, max_retries + 1, attempt_elapsed)
            )
            print("[CHUNK %d] Successfully generated: %s" % (chunk_idx, generated_file))
            return generated_file
        else:
            print(
                "[TIMING] Chunk %d attempt %d/%d: %.1fs (no file returned)"
                % (chunk_idx, attempt + 1, max_retries + 1, attempt_elapsed)
            )
            print(
                "[CHUNK %d] Attempt %d/%d produced no file."
                % (chunk_idx, attempt + 1, max_retries + 1)
            )

    print(
        "[CHUNK %d] All %d attempts failed. Skipping chunk."
        % (chunk_idx, max_retries + 1)
    )
    session_state["use_fallback_generator"] = True
    return None


# Slide type -> python-pptx layout index mapping for the fallback generator.
# Index 0: Title Slide (large title + subtitle, used for opening/closing)
# Index 1: Title and Content (standard bullet slide, the most common layout)
# Index 2: Section Header (bold title only, ideal for agenda/divider slides)
# Index 3: Two Content (two side-by-side content areas, good for data/comparison)
# Fallback for unknown types: index 1 (Title and Content)
FALLBACK_SLIDE_LAYOUT_MAP = {
    "title": 0,
    "agenda": 2,
    "content": 1,
    "data": 3,
    "closing": 0,
}


def _detect_chart_type_from_suggestion(visual: str) -> Optional[object]:
    """Return an XL_CHART_TYPE enum value when the visual_suggestion describes a chart.

    Scans the visual_suggestion string for chart-type keywords.
    Returns None if no recognizable chart keyword is found (e.g. an image,
    icon, or diagram description), in which case the caller should fall back
    to a textbox annotation.

    Keyword mapping:
      pie                        -> XL_CHART_TYPE.PIE
      line                       -> XL_CHART_TYPE.LINE
      bar (horizontal context)   -> XL_CHART_TYPE.BAR_CLUSTERED
      column / chart (generic)   -> XL_CHART_TYPE.COLUMN_CLUSTERED
    """
    from pptx.enum.chart import XL_CHART_TYPE  # type: ignore

    v = visual.lower()
    if "pie" in v:
        return XL_CHART_TYPE.PIE
    if "line" in v:
        return XL_CHART_TYPE.LINE
    if "bar" in v:
        return XL_CHART_TYPE.BAR_CLUSTERED
    if "column" in v or "chart" in v:
        return XL_CHART_TYPE.COLUMN_CLUSTERED
    return None


def generate_chunk_pptx_fallback(
    chunk_slides: List[SlideStoryboard],
    session_state: Dict,
    chunk_idx: int,
) -> Optional[str]:
    """Tier 3 (last-resort) chunk generator using python-pptx directly.

    No Claude API call — generates slides programmatically from SlideStoryboard
    data with zero network I/O. Always produces a valid .pptx or returns None
    only on an extreme exception (e.g., disk full).

    This is the last tier in the fallback hierarchy:
      Tier 1: Primary provider PPTX skill (generate_chunk_pptx)
      Tier 2: Primary provider LLM code gen (generate_chunk_pptx_v2)
              -> Universal 4-step OpenAI fallback on failure
      Tier 3: This function — python-pptx direct, <100ms

    Output slides contain title text + bullet points. When a slide's
    visual_suggestion contains a chart-type keyword (bar, column, line, pie, chart),
    a native python-pptx chart is inserted using CategoryChartData + add_chart()
    (no matplotlib or image embedding). Non-chart visuals receive a small textbox
    annotation. The output is a structurally valid .pptx compatible with
    step_process_chunks() template assembly and _merge_pptx_zip_level() merging.

    Slide layout mapping:
    - slide_type == "title"   -> layout index 0 (TITLE slide)
    - all others              -> layout index 1 (TITLE_AND_CONTENT)

    Args:
        chunk_slides: List of SlideStoryboard objects for this chunk.
        session_state: Shared workflow session state.
        chunk_idx: 0-based chunk index (used for file naming and logging).

    Returns:
        Path to the generated chunk PPTX file, or None if generation failed.
    """
    output_dir = session_state.get("output_dir", ".")
    output_path = os.path.join(output_dir, "chunk_%03d.pptx" % chunk_idx)

    try:
        prs = Presentation()

        for slide in chunk_slides:
            layout_idx = FALLBACK_SLIDE_LAYOUT_MAP.get(slide.slide_type, 1)
            # Guard: some presentations may have fewer layouts than expected.
            # Fall back to index 1 (Title and Content) if the chosen index is out of range.
            if layout_idx >= len(prs.slide_layouts):
                layout_idx = min(1, len(prs.slide_layouts) - 1)
            slide_layout = prs.slide_layouts[layout_idx]

            pptx_slide = prs.slides.add_slide(slide_layout)

            # Set title
            if pptx_slide.shapes.title:
                pptx_slide.shapes.title.text = slide.slide_title

            if slide.slide_type == "title":
                # Set subtitle placeholder
                subtitle_text = (
                    slide.key_points[0] if slide.key_points else slide.slide_title
                )
                if subtitle_text:
                    for ph in pptx_slide.placeholders:
                        if ph.placeholder_format.idx == 1:
                            ph.text = subtitle_text
                            break
            else:
                # Set body content with key points
                body_ph = None
                for ph in pptx_slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        body_ph = ph
                        break
                if body_ph:
                    tf = body_ph.text_frame  # type: ignore
                    tf.word_wrap = True
                    tf.clear()
                    for i, point in enumerate(slide.key_points):
                        if i == 0:
                            tf.paragraphs[0].text = point  # type: ignore
                        else:
                            p = tf.add_paragraph()  # type: ignore
                            p.text = point
                        tf.paragraphs[i].level = 0  # type: ignore

            # Add visual element if applicable — prefer native chart over textbox label
            visual = slide.visual_suggestion
            if visual and visual.lower() != "none":
                chart_xl_type = _detect_chart_type_from_suggestion(visual)
                if chart_xl_type is not None:
                    # Insert a native python-pptx chart with synthesized sample data.
                    # CategoryChartData / add_chart() produce editable Office chart objects;
                    # no matplotlib or image embedding is used.
                    try:
                        from pptx.chart.data import CategoryChartData  # type: ignore

                        chart_data = CategoryChartData()
                        chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
                        chart_data.add_series("Series 1", (25, 40, 35, 55))
                        pptx_slide.shapes.add_chart(
                            chart_xl_type,
                            Inches(1),
                            Inches(2.5),
                            Inches(8),
                            Inches(3.5),
                            chart_data,
                        )
                    except Exception as chart_err:
                        print(
                            "[CHUNK FALLBACK] Native chart insertion failed: %s; "
                            "falling back to textbox label." % chart_err
                        )
                        txBox = pptx_slide.shapes.add_textbox(
                            Inches(7), Inches(5.5), Inches(2.5), Inches(0.5)
                        )
                        tf = txBox.text_frame
                        tf.text = "[Chart: %s]" % visual[:60]
                        if tf.paragraphs and tf.paragraphs[0].runs:
                            run = tf.paragraphs[0].runs[0]
                        elif tf.paragraphs:
                            run = tf.paragraphs[0].add_run()
                        else:
                            run = None
                        if run is not None:
                            run.font.size = Pt(8)
                            run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
                else:
                    v_lower = visual.lower()
                    if "table" in v_lower:
                        # Insert a native 2x2 placeholder table using python-pptx add_table().
                        # No matplotlib or image embedding — all native OOXML.
                        try:
                            table_shape = pptx_slide.shapes.add_table(
                                2, 2, Inches(1), Inches(2.5), Inches(8), Inches(3.0)
                            )
                            tbl = table_shape.table
                            tbl.cell(0, 0).text = "Item"
                            tbl.cell(0, 1).text = "Value"
                            tbl.cell(1, 0).text = "—"
                            tbl.cell(1, 1).text = "—"
                        except Exception as tbl_err:
                            print(
                                "[CHUNK FALLBACK] Native table insertion failed: %s; "
                                "falling back to textbox label." % tbl_err
                            )
                            txBox = pptx_slide.shapes.add_textbox(
                                Inches(7), Inches(5.5), Inches(2.5), Inches(0.5)
                            )
                            tf = txBox.text_frame
                            tf.text = "[Table: %s]" % visual[:60]
                            if tf.paragraphs and tf.paragraphs[0].runs:
                                run = tf.paragraphs[0].runs[0]
                            elif tf.paragraphs:
                                run = tf.paragraphs[0].add_run()
                            else:
                                run = None
                            if run is not None:
                                run.font.size = Pt(8)
                                run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
                    elif "infographic" in v_lower or "diagram" in v_lower:
                        # Insert labeled rectangle shapes to approximate an infographic/diagram.
                        # Uses native python-pptx shapes — no images inserted.
                        try:
                            from pptx.enum.shapes import (
                                MSO_AUTO_SHAPE_TYPE,  # type: ignore
                            )

                            labels = ["Step 1", "Step 2", "Step 3"]
                            box_w = Inches(2.2)
                            box_h = Inches(1.0)
                            top = Inches(2.8)
                            for li, label in enumerate(labels):
                                left = Inches(0.8 + li * 2.8)
                                shape = pptx_slide.shapes.add_shape(
                                    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                    left,
                                    top,
                                    box_w,
                                    box_h,
                                )
                                shape.text = label
                                shape.text_frame.paragraphs[0].font.size = Pt(11)
                        except Exception as inf_err:
                            print(
                                "[CHUNK FALLBACK] Infographic shape insertion failed: %s; "
                                "falling back to textbox label." % inf_err
                            )
                            txBox = pptx_slide.shapes.add_textbox(
                                Inches(7), Inches(5.5), Inches(2.5), Inches(0.5)
                            )
                            tf = txBox.text_frame
                            tf.text = "[Visual: %s]" % visual[:60]
                            if tf.paragraphs and tf.paragraphs[0].runs:
                                run = tf.paragraphs[0].runs[0]
                            elif tf.paragraphs:
                                run = tf.paragraphs[0].add_run()
                            else:
                                run = None
                            if run is not None:
                                run.font.size = Pt(8)
                                run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
                    else:
                        # Non-chart, non-table, non-diagram visual: add a label annotation
                        txBox = pptx_slide.shapes.add_textbox(
                            Inches(7), Inches(5.5), Inches(2.5), Inches(0.5)
                        )
                        tf = txBox.text_frame
                        tf.text = "[Visual: %s]" % visual[:60]
                        if tf.paragraphs and tf.paragraphs[0].runs:
                            run = tf.paragraphs[0].runs[0]
                        elif tf.paragraphs:
                            run = tf.paragraphs[0].add_run()
                        else:
                            run = None
                        if run is not None:
                            run.font.size = Pt(8)
                            run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

        # Clean up empty placeholders and hardcoded contrast issues
        clean_presentation_visual_noise_and_contrast(prs)
        sanitize_presentation(prs)  # noqa: F405
        prs.save(output_path)
        print(
            "[CHUNK %d FALLBACK] Generated %d slides via python-pptx fallback"
            % (chunk_idx, len(chunk_slides))
        )
        return output_path

    except Exception as e:
        print(
            "[CHUNK %d FALLBACK] Failed to generate fallback PPTX: %s" % (chunk_idx, e)
        )
        return None


# Instructions for the Tier 2 fallback code-generation agent.
# Kept as a module-level constant to avoid string duplication.
PPTX_CODE_GEN_INSTRUCTIONS = [
    "You are a Python code generator that creates PowerPoint presentations using python-pptx.",
    "When given slide specifications, write a COMPLETE, SELF-CONTAINED Python script that generates a .pptx file.",
    "The script must import all required libraries at the top.",
    "ALLOWED imports only: pptx, pptx.util, pptx.chart.data, pptx.enum.chart, pptx.dml.color, pptx.enum.text, io, os, os.path, collections, math.",
    "FORBIDDEN imports: matplotlib, matplotlib.pyplot, subprocess, socket, requests, urllib, httpx, shutil, glob, sys, importlib, __import__.",
    "For each slide, create one slide in the presentation using prs.slides.add_slide(prs.slide_layouts[N]).",
    "Slide layout indices: 0=Title Slide, 1=Title and Content, 2=Section Header, 3=Two Content.",
    "For CHARTS: ALWAYS use python-pptx native CategoryChartData or ChartData (creates editable Office charts). NEVER use matplotlib, PIL, or any image-based approach for charts.",
    "For python-pptx ChartData bar/column charts: from pptx.chart.data import ChartData; from pptx.enum.chart import XL_CHART_TYPE.",
    "ChartData example: chart_data = ChartData(); chart_data.categories = ['A','B','C']; chart_data.add_series('Series1', (10, 20, 30)); slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(1), Inches(1.5), Inches(8), Inches(4.5), chart_data).",
    "For line charts use XL_CHART_TYPE.LINE, for pie charts use XL_CHART_TYPE.PIE — always via ChartData, never via image embedding.",
    "CHART SIZING (CRITICAL): Charts MUST fill at least 60% of the slide content area. Default chart dimensions: Inches(1), Inches(1.5), Inches(8), Inches(4.5). NEVER create tiny charts — they must dominate the visual area.",
    "CHART DATA LABELS (CRITICAL): ALWAYS enable data labels on every chart: plot = chart.plots[0]; plot.has_data_labels = True; plot.data_labels.show_value = True. For PIE charts additionally: plot.data_labels.show_category_name = True; plot.data_labels.show_percentage = True; plot.data_labels.show_value = False.",
    "For TABLES: ALWAYS use slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(4.5)). Fill cells via table.cell(row, col).text = 'value'. NEVER embed a table as an image or use matplotlib/PIL to render one.",
    "For INFOGRAPHICS or DIAGRAMS: use native python-pptx shapes (add_shape with MSO_AUTO_SHAPE_TYPE.RECTANGLE, add_textbox, add_connector) or a native table to approximate the layout. NEVER insert an image for an infographic or diagram.",
    "Treat charts/tables/infographics/diagrams as native data-vis. Preserve data-vis intent even when exact styling cannot be replicated.",
    "Synthesize plausible, specific data values from the visual_suggestion and key_points descriptions. Do NOT use generic placeholder data.",
    "CHART AXIS SAFETY: When styling chart axes, always wrap axis.format.line.fill.background(), axis.format.line.color.rgb, and axis.major_gridlines.format.line.color.rgb calls in try/except blocks to handle cases where the underlying XML element does not yet exist. Example: try: chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0x80,0x80,0x80)\\nexcept Exception: pass",
    "CHART AXIS SAFETY: Similarly wrap axis.tick_labels.font.color.rgb, chart.legend.font.color.rgb, and series.data_labels.font.color.rgb in try/except blocks — these can raise 'NoneType object has no attribute attrib' if the XML element is absent.",
    "FILL COLOR API (CRITICAL): The correct python-pptx API for shape fill colors is shape.fill.fore_color (NOT shape.fill.foreground_color — that attribute does NOT exist and will crash). "
    "To set a solid fill: shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xFF,0x00,0x00). "
    "To read a fill color: rgb_val = shape.fill.fore_color.rgb. "
    "Always wrap fill color access in try/except since some shapes have no fill or use theme/pattern fills: try: shape.fill.fore_color.rgb = RGBColor(...) except Exception: pass.",
    "TITLE PLACEMENT (CRITICAL): DO NOT create a separate textbox for the slide title. The slide layout already has a title placeholder at index 0 — set title text on it: slide.shapes.title.text = 'Your Title'. Creating an extra title textbox causes duplicate/overlapping titles.",
    "Save the final presentation using prs.save('EXACT_OUTPUT_PATH') where EXACT_OUTPUT_PATH is the path given in the task.",
    "Execute the script using the save_to_file_and_run tool immediately after writing it.",
    "If the script has an error, fix it and re-run. Maximum 2 fix attempts.",
    "If a visual cannot be implemented exactly, keep the slide and add a concise native textbox note. Do not skip slides.",
    "Do not add speaker notes, animations, or transitions.",
    "Do not print to stdout or write any files other than the final prs.save() call.",
    "COLOR AND CONTRAST RULES:",
    "- Default slide background: white (#FFFFFF) or very light colors (luminance > 0.9).",
    "- Default body text: dark colors (#333333 or #000000) for readability.",
    "- If using a dark background (luminance < 0.3), use white (#FFFFFF) or very light text.",
    "- If using colored shape fills for headers/accents, ensure text color has sufficient contrast.",
    "- For charts: use medium-to-dark accent colors; data labels should contrast against their background.",
    "- For tables: header rows with dark fills should have white text; body rows with light fills should have dark text.",
    "- NEVER use dark text (#000000-#666666) on dark backgrounds (#000000-#555555).",
    "- NEVER use light text (#AAAAAA-#FFFFFF) on light backgrounds (#CCCCCC-#FFFFFF).",
    "- When in doubt, use white background with black text — readability is paramount.",
    "MINIMUM TEXT BOX SIZE (CRITICAL):",
    "- Any textbox containing body content MUST be at least Inches(3) wide and Inches(1.5) tall.",
    "- NEVER create tiny text boxes that cram dense content into a small area.",
    "- If space is limited, REDUCE the amount of text content rather than shrinking the text box.",
    "- Maximum 4 text-bearing shapes per slide (excluding title placeholder and chart/table labels).",
    "- Each text shape should contain at most 5 bullet points or 6 lines of text.",
    "SPATIAL POSITIONING RULES (CRITICAL):",
    "- Use a GRID SYSTEM: divide the slide into columns (2, 3, or 4) and place shapes on grid lines.",
    "- For multi-column layouts: calculate column_width = (Inches(11.5) - Inches(1)) / num_columns. "
    "Place column N at left = Inches(0.75) + N * column_width.",
    "- All shapes in the SAME ROW must share the SAME top position. All shapes in the SAME COLUMN "
    "must share the SAME left position and width.",
    "- NEVER place text shapes at arbitrary positions — always compute positions from a grid.",
    "- For vertical stacking: use consistent left margin (Inches(0.75)) and increment top by "
    "previous_shape.height + Inches(0.3) gap.",
    "DECORATIVE ELEMENT RULES:",
    "- Do NOT add decorative icons, emoji characters, or symbol shapes unless they directly "
    "illustrate the slide content (e.g., a chart icon next to a data section title).",
    "- Each shape on the slide must serve a clear informational purpose.",
    "- Do NOT add bell icons, star icons, trophy icons, checkmark decorations, or other generic decorative shapes.",
    "- If you want visual interest, use colored rectangles, accent bars, or the template's own shapes — not icon characters.",
    "DATA VISUALIZATION FALLBACK:",
    "- If a chart cannot be fully implemented, create a STRUCTURED table or KPI card layout instead.",
    "- NEVER dump raw data as a plain text block — always structure data into labeled rows/columns.",
    "- For metric slides: use large number + label pattern (e.g., Inches(3) wide rectangle "
    "with 36pt number and 14pt label below it).",
    "- Every data point on the slide must have a clear label and be visually distinct from other data points.",
]

# Tier 2 fallback agent is now loaded from the agents/ package via get_agents().
# Stored in session_state["agents"]["fallback_code_agent"].
# See agents/claude_agents.py (or openai/gemini variants) for the definition.


def generate_chunk_pptx_v2(
    chunk_slides: List[SlideStoryboard],
    session_state: Dict,
    chunk_idx: int,
) -> Optional[str]:
    """Tier 2 fallback chunk generator using LLM code generation + PythonTools execution.

    Prompts the fallback_code_agent (Claude Opus without PPTX skill, equipped with
    PythonTools) to write and execute a python-pptx script that creates the chunk
    slides with native charts (CategoryChartData/ChartData), tables, and rich visual content.
    matplotlib is FORBIDDEN — all charts must use python-pptx native chart objects.

    Quality level: 80-92% parity with Tier 1 (Claude PPTX skill). Charts and tables
    are generated via python-pptx native objects rather than the native PPTX skill,
    so visual fidelity may differ in edge cases.

    Brand/style context from session_state["brand_style_intent"] is appended to the
    GLOBAL CONTEXT section of the code generation prompt when branding is present.

    Hierarchy & Fallback Logic (6-stage chain):
      Stage 1: Primary agent (Sonnet code-gen)
      Stage 2: Lite agent (Haiku code-gen)
      Stage 2.5: Gemini Pro (gemini-3.1-pro-preview) code-gen
               → Gemini Flash (gemini-2.5-flash) code-gen
      Stage 3: OpenAI Pro (gpt-5.4) WITH visual base64 references
             → OpenAI Lite (o3-mini) WITH visual base64 references
             → OpenAI Pro WITHOUT visual references (payload stripped)
             → OpenAI Lite WITHOUT visual references

    This ensures context limits (TPM/RPM) don't break generation. Only after all
    stages fail does it escalate to Tier 3 text-only generation.

    Args:
        chunk_slides: List of SlideStoryboard objects for this chunk.
        session_state: Shared workflow session state (must contain brand_style_intent
                       key set by step_optimize_and_plan).
        chunk_idx: 0-based chunk index (used for file naming and logging).

    Returns:
        Path to the generated chunk PPTX file, or None if generation failed.
    """
    storyboard: StoryboardPlan = session_state.get("storyboard")
    output_dir = session_state.get("output_dir", ".")
    chunk_output_path = os.path.join(output_dir, "chunk_%03d.pptx" % chunk_idx)

    first_slide = chunk_slides[0].slide_number
    last_slide = chunk_slides[-1].slide_number

    print(
        "[CHUNK %d TIER2] Starting LLM code generation fallback (slides %d-%d)..."
        % (chunk_idx, first_slide, last_slide)
    )
    if VERBOSE:
        template_pngs = session_state.get("template_slide_pngs", {})
        print(
            "[VERBOSE] [TIER2] Visual references available: %d slide(s)"
            % len(template_pngs)
        )

    # Build the code generation prompt with full slide specifications
    slide_specs = []
    for slide in chunk_slides:
        spec = (
            "Slide %d (type=%s): title='%s'\n  Key points: %s\n  Visual suggestion: %s"
        ) % (
            slide.slide_number,
            slide.slide_type,
            slide.slide_title,
            "; ".join(slide.key_points),
            slide.visual_suggestion,
        )
        slide_specs.append(spec)

    global_ctx = ""
    if storyboard:
        global_ctx = (
            "Presentation: '%s' | Audience: %s | Tone: %s | Brand voice: %s\n"
            "Context: %s"
        ) % (
            storyboard.presentation_title,
            storyboard.target_audience,
            storyboard.tone,
            storyboard.brand_voice,
            storyboard.global_context,
        )

    # Append brand/style context if available
    brand_ctx = _format_brand_context_for_prompt(
        session_state.get("brand_style_intent", BrandStyleIntent())
    )
    if brand_ctx:
        global_ctx = global_ctx + "\n\n" + brand_ctx if global_ctx else brand_ctx

    # --- Fix 5: Build template constraints for LLM code gen ---
    # When a template exists, extract its styling. When no template exists,
    # inject the visual design system derived from the optimizer's visual_style.
    template_constraints = ""
    template_path_t2 = session_state.get("template_path", "")
    if template_path_t2 and os.path.isfile(template_path_t2):
        try:
            from powerpoint_template_workflow import (  # type: ignore
                _extract_template_styles,
                _get_shape_background_color,
                _hex_to_rgb,
                _relative_luminance,
            )

            _tmpl_prs = Presentation(template_path_t2)
            _tmpl_style = _extract_template_styles(_tmpl_prs)

            # Detect if template has a dark background
            bg_is_dark = False
            bg_hex = "FFFFFF"
            if _tmpl_prs.slides:
                first_slide = _tmpl_prs.slides[0]
                if first_slide.shapes:
                    bg_hex = _get_shape_background_color(
                        list(first_slide.shapes)[0], first_slide
                    )
                bg_rgb = _hex_to_rgb(bg_hex)
                bg_lum = _relative_luminance(*bg_rgb)
                bg_is_dark = bg_lum < 0.4

            if bg_is_dark:
                template_constraints = (
                    "\nTEMPLATE CONSTRAINTS (CRITICAL — your content will be placed on a template):\n"
                    "- The template has a DARK background (approx #%s). You MUST use WHITE or very "
                    "light text colors (e.g. RGBColor(0xFF, 0xFF, 0xFF)) for ALL text elements.\n"
                    "- Do NOT use dark blue, dark gray, or black text — it will be invisible.\n"
                    "- Keep each slide SIMPLE: maximum 4-5 content shapes per slide.\n"
                    "- Use large font sizes (minimum 14pt for body, 24pt for titles).\n"
                    "- Avoid overlapping shapes — stack content vertically with clear spacing.\n"
                    "- Do NOT add background shapes or colored rectangles — the template provides these.\n"
                ) % bg_hex
            else:
                template_constraints = (
                    "\nTEMPLATE CONSTRAINTS (your content will be placed on a template):\n"
                    "- Keep each slide SIMPLE: maximum 4-5 content shapes per slide.\n"
                    "- Use large font sizes (minimum 14pt for body, 24pt for titles).\n"
                    "- Avoid overlapping shapes — stack content vertically with clear spacing.\n"
                    "- Do NOT add background shapes — the template provides these.\n"
                )
            print(
                "  [TEMPLATE CTX] Template context injected into Tier 2 prompt "
                "(bg_dark=%s, bg_hex=#%s)." % (bg_is_dark, bg_hex)
            )
        except Exception as e:
            print("  [TEMPLATE CTX WARNING] Could not extract template context: %s" % e)

        # --- Accent line pattern for Tier 2 code-gen ---
        # Inject structured accent bar instructions so the code-gen LLM can
        # programmatically reproduce template accent lines via python-pptx.
        visual_profile = session_state.get("template_visual_profile")
        if visual_profile and getattr(visual_profile, "has_accent_lines", False):
            ap = getattr(visual_profile, "accent_pattern", {})
            # Fix 16: Only inject accent instructions if the template genuinely
            # has wide horizontal bars (avg_width_pct > 20%). Skip narrow
            # diagonal decorations that mislead the LLM.
            if ap and ap.get("avg_width_pct", 0) > 20:
                color_code = ap.get("color", "")
                color_instruction = (
                    "RGBColor(0x%s, 0x%s, 0x%s)"
                    % (color_code[:2], color_code[2:4], color_code[4:6])
                    if color_code and len(color_code) == 6
                    else "the template's primary accent color from the theme"
                )
                template_constraints += (
                    "- ACCENT LINE PATTERN (MUST ADD TO EVERY CONTENT SLIDE):\n"
                    "  The template has a consistent %s accent bar in the %s region.\n"
                    "  For each content slide, add a thin rectangle shape:\n"
                    "    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,\n"
                    "        left=Inches(%.2f), top=Inches(%.2f),\n"
                    "        width=Inches(%.2f), height=Inches(%.2f))\n"
                    "    shape.fill.solid()\n"
                    "    shape.fill.fore_color.rgb = %s\n"
                    "    shape.line.fill.background()\n"
                    "  This bar is present on %.0f%% of template slides.\n"
                ) % (
                    ap.get("orientation", "horizontal"),
                    ap.get("region", "top"),
                    ap.get("avg_left_pct", 0)
                    / 100.0
                    * 10.0,  # Convert pct to approx inches (10" slide)
                    ap.get("avg_top_pct", 0)
                    / 100.0
                    * 7.5,  # Convert pct to approx inches (7.5" slide)
                    ap.get("avg_width_pct", 0) / 100.0 * 10.0,
                    ap.get("avg_height_pct", 0) / 100.0 * 7.5,
                    color_instruction,
                    ap.get("slide_coverage_pct", 0),
                )
                if VERBOSE:
                    print(
                        "  [TEMPLATE CTX] Accent pattern injected into Tier 2 prompt "
                        "(orientation=%s, region=%s, color=%s)."
                        % (
                            ap.get("orientation"),
                            ap.get("region"),
                            color_code or "auto",
                        )
                    )
    else:
        # No template: inject design system from optimizer's visual_style
        template_constraints = _build_no_template_design_system(
            storyboard.visual_style if storyboard else "clean_minimal",
            session_state.get("brand_style_intent"),
        )
        if template_constraints and VERBOSE:  # noqa: F405
            print(
                "[VERBOSE] [TIER2] No-template design system injected "
                "(visual_style=%s, %d chars)"
                % (
                    storyboard.visual_style if storyboard else "clean_minimal",
                    len(template_constraints),
                )
            )

    code_gen_prompt = (
        "Generate a complete Python script using python-pptx to create "
        "a PowerPoint file at this EXACT path: %s\n\n"
        "GLOBAL CONTEXT:\n%s\n%s\n\n"
        "SLIDES TO GENERATE (%d slides):\n%s\n\n"
        "REQUIREMENTS:\n"
        "- Create exactly %d slides in the exact order listed above.\n"
        "- For 'title' type slides: use prs.slide_layouts[0] (Title Slide layout).\n"
        "- For 'agenda'/'section' type slides: use prs.slide_layouts[2] (Section Header).\n"
        "- For 'content' type slides: use prs.slide_layouts[1] (Title and Content).\n"
        "- For 'data' type slides: use prs.slide_layouts[3] if available, else [1].\n"
        "- For 'closing' type slides: use prs.slide_layouts[1] (Title and Content) to ensure body text placeholders exist.\n"
        "SEMANTIC LAYOUT RULES (CRITICAL):\n"
        "- If Semantic Type is 'sequential': DO NOT use standard bullets. Programmatically draw a horizontal chevron/arrow process flow using native python-pptx shapes.\n"
        "- If Semantic Type is 'comparative': DO NOT use standard bullets. Programmatically draw a 2-4 column card grid using filled rectangles and textboxes.\n"
        "- If Semantic Type is 'metrics': DO NOT use standard bullets. Draw a KPI Dashboard row with large, high-contrast numbers inside subtle native rectangles.\n"
        "- If Semantic Type is 'hero': DO NOT use bullets. Draw a massive, centered title and subtitle on a dark tinted rectangle spanning the full slide.\n"
        "- For any slide with a chart visual_suggestion: generate a REAL chart using "
        "  python-pptx native ChartData ONLY (e.g. CategoryChartData + slide.shapes.add_chart()).\n"
        "  Do NOT use matplotlib, PIL, or any image-based approach for charts.\n"
        "  Do NOT use matplotlib, PIL, or any image-based approach for charts.\n"
        "- Synthesize specific, plausible data values matching the visual_suggestion topic.\n"
        "- For any slide with a table visual_suggestion: generate a REAL native table using slide.shapes.add_table(). NEVER use matplotlib or embed a table as an image.\n"
        "- For any slide with an infographic or diagram visual_suggestion: use native python-pptx shapes (add_shape with MSO_AUTO_SHAPE_TYPE.RECTANGLE, add_textbox) or a native table to approximate it. NEVER insert an image for an infographic or diagram.\n"
        "- If exact visual styling is not feasible, preserve content and structure with a concise native textbox note. Never skip a requested slide.\n"
        "ALIGNMENT AND COMPLETENESS RULES (CRITICAL):\n"
        "- GRID ALIGNMENT: For multi-element slides, compute a grid. All shapes in a row share "
        "the SAME top. All shapes in a column share the SAME left and width. Use math to "
        "compute positions (e.g., col_left = Inches(0.75) + col_index * col_width).\n"
        "- VISUAL COMPLETENESS: If you create a visual element (diagram, flow, chart), include "
        "ALL parts. Do not leave orphaned arrows, incomplete flows, or disconnected labels.\n"
        "- NO RAW DATA DUMPS: Never create a textbox that contains a list of raw numbers/labels "
        "without structure. If data cannot be charted, format it as a native table or KPI cards "
        "(large number rectangles with labels below).\n"
        "- NO DECORATIVE ICONS: Do not add symbol/emoji shapes (bells, stars, trophies) for "
        "decoration. Every shape must carry meaningful content.\n"
        "- Save the file to: %s\n"
        "- Then immediately execute the script using save_to_file_and_run.\n"
        "- Return only tool execution needed to write and run the script; avoid extra narrative output.\n"
    ) % (
        chunk_output_path,
        global_ctx,
        template_constraints,
        len(chunk_slides),
        "\n\n".join(slide_specs),
        len(chunk_slides),
        chunk_output_path,
    )

    if VERBOSE:  # noqa: F405
        print(
            "[VERBOSE] Chunk %d Tier 2 code-gen prompt length: %d chars"
            % (chunk_idx, len(code_gen_prompt))
        )

    # Append template visual reference section for Tier 2 (if available)
    template_pngs = session_state.get("template_slide_pngs", {})
    if template_pngs and session_state.get("template_visuals"):
        visual_ref = _build_visual_reference_section(
            chunk_slides,
            template_pngs,
            brand_style_intent=session_state.get("brand_style_intent"),
        )
        if visual_ref:
            code_gen_prompt += "\n" + visual_ref
            if VERBOSE:  # noqa: F405
                print(
                    "[VERBOSE] Chunk %d Tier 2: appended %d-char visual reference."
                    % (chunk_idx, len(visual_ref))
                )

    # === Tier 2 Model Chain: Primary (Sonnet) → Lite (Haiku) → Tier 3 ===
    # stream=True is required: max_tokens=16384 causes the Anthropic SDK to enforce
    # streaming for calls that may take longer than 10 minutes.
    # Tier 2 only cares whether the file was created on disk, not about the response
    # content, so we iterate through the stream and discard events.
    t2_start = time.time()
    from agents import get_agents as _get_agents  # type: ignore

    _all_agents = _get_agents(session_state.get("llm_provider", "claude"))
    _fallback_agent = _all_agents.get("fallback_code_agent")
    _fallback_agent_lite = _all_agents.get("fallback_code_agent_lite")

    tier2_success = False

    # --- Stage 1: Primary agent (Sonnet) ---
    actual_model_id = getattr(
        getattr(_fallback_agent, "model", None), "id", "claude-haiku-4-5"
    )
    _log_agent_banner(
        agent_name=getattr(_fallback_agent, "name", "PPTX Code Generator"),
        model_id=actual_model_id,
        provider=getattr(
            getattr(_fallback_agent, "model", None),
            "provider",
            session_state.get("llm_provider", "claude"),
        ),  # type: ignore
        step_name="step_generate_chunks / Tier 2 Primary (chunk %d)" % chunk_idx,
    )
    _get_rate_tracker().check_and_wait(
        model=actual_model_id,
        prompt=code_gen_prompt,
        caller="generate_chunk_pptx_v2/Tier2-primary",
    )
    try:
        for _ in _fallback_agent.run(code_gen_prompt, stream=True):
            pass
        t2_elapsed = time.time() - t2_start
        print(
            "[TIMING] Chunk %d Tier 2 primary code generation: %.1fs"
            % (chunk_idx, t2_elapsed)
        )
        tier2_success = True
    except (Exception, SystemExit) as e:
        t2_elapsed = time.time() - t2_start
        print(
            "[CHUNK %d TIER2] Primary agent (Sonnet) failed after %.1fs: %s"
            % (chunk_idx, t2_elapsed, str(e))
        )

    # If primary produced no file, also treat as failure
    if tier2_success and not os.path.exists(chunk_output_path):
        print(
            "[CHUNK %d TIER2] Primary agent ran but produced no file. Trying lite agent..."
            % chunk_idx
        )
        tier2_success = False

    # --- Stage 2: Lite agent (Haiku/GPT-5-mini) — only if primary failed ---
    if not tier2_success and _fallback_agent_lite is not None:
        print(
            "[CHUNK %d TIER2] Retrying with lite agent (%s)..."
            % (chunk_idx, getattr(_fallback_agent_lite, "name", "Fallback Lite"))
        )
        t2_lite_start = time.time()
        _log_agent_banner(
            agent_name=getattr(
                _fallback_agent_lite, "name", "PPTX Code Generator (Lite)"
            ),
            model_id=getattr(
                getattr(_fallback_agent_lite, "model", None), "id", "unknown-lite"
            ),
            provider=getattr(
                getattr(_fallback_agent_lite, "model", None),
                "provider",
                session_state.get("llm_provider", "claude"),
            ),  # type: ignore
            step_name="step_generate_chunks / Tier 2 Lite Fallback (chunk %d)"
            % chunk_idx,
        )

        # Rate tracking only for Anthropic models
        if session_state.get("llm_provider", "claude") == "claude":
            _get_rate_tracker().check_and_wait(
                model="claude-haiku-4-5",
                prompt=code_gen_prompt,
                caller="generate_chunk_pptx_v2/Tier2-lite",
            )

        try:
            for _ in _fallback_agent_lite.run(code_gen_prompt, stream=True):
                pass
            t2_lite_elapsed = time.time() - t2_lite_start
            print(
                "[TIMING] Chunk %d Tier 2 lite code generation: %.1fs"
                % (chunk_idx, t2_lite_elapsed)
            )
            tier2_success = True
        except (Exception, SystemExit) as e2:
            t2_lite_elapsed = time.time() - t2_lite_start
            print(
                "[CHUNK %d TIER2] Lite agent also failed after %.1fs: %s"
                % (chunk_idx, t2_lite_elapsed, str(e2))
            )

    # If Lite also produced no file
    if tier2_success and not os.path.exists(chunk_output_path):
        print("[CHUNK %d TIER2] Lite agent ran but produced no file." % chunk_idx)
        tier2_success = False

    # --- Stage 2.5: Gemini Fallback (Pro → Flash) ---
    # Uses fallback_code_agent_fallback (gemini-3.1-pro-preview) and
    # fallback_code_agent_lite_fallback (gemini-2.5-flash) defined in
    # claude_agents.py. Single attempt each, code-gen path.
    if not tier2_success:
        _gemini_pro = _all_agents.get("fallback_code_agent_fallback")
        if _gemini_pro is not None:
            print(
                "[CHUNK %d TIER2] Attempting Gemini Pro code-gen fallback..."
                % chunk_idx
            )
            t2_gemini_start = time.time()
            _log_agent_banner(
                agent_name=getattr(
                    _gemini_pro,
                    "name",
                    "PPTX Code Generator (Gemini Fallback)",
                ),
                model_id=getattr(
                    getattr(_gemini_pro, "model", None),
                    "id",
                    "gemini-3.1-pro-preview",
                ),
                provider="Google",
                step_name="step_generate_chunks / Tier 2 Gemini Pro "
                "Fallback (chunk %d)" % chunk_idx,
            )
            try:
                for _ in _gemini_pro.run(code_gen_prompt, stream=True):
                    pass
                t2_gemini_elapsed = time.time() - t2_gemini_start
                print(
                    "[TIMING] Chunk %d Tier 2 Gemini Pro: %.1fs"
                    % (chunk_idx, t2_gemini_elapsed)
                )
                if os.path.exists(chunk_output_path):
                    tier2_success = True
                else:
                    print(
                        "[CHUNK %d TIER2] Gemini Pro ran but produced "
                        "no file." % chunk_idx
                    )
            except Exception as e_gemini:
                t2_gemini_elapsed = time.time() - t2_gemini_start
                print(
                    "[CHUNK %d TIER2] Gemini Pro failed after %.1fs: %s"
                    % (chunk_idx, t2_gemini_elapsed, str(e_gemini)[:200])  # type: ignore
                )

    if not tier2_success:
        _gemini_flash = _all_agents.get("fallback_code_agent_lite_fallback")
        if _gemini_flash is not None:
            print(
                "[CHUNK %d TIER2] Attempting Gemini Flash code-gen "
                "fallback..." % chunk_idx
            )
            t2_gflash_start = time.time()
            _log_agent_banner(
                agent_name=getattr(
                    _gemini_flash,
                    "name",
                    "PPTX Code Generator (Gemini Flash Fallback)",
                ),
                model_id=getattr(
                    getattr(_gemini_flash, "model", None),
                    "id",
                    "gemini-2.5-flash",
                ),
                provider="Google",
                step_name="step_generate_chunks / Tier 2 Gemini Flash "
                "Fallback (chunk %d)" % chunk_idx,
            )
            try:
                for _ in _gemini_flash.run(code_gen_prompt, stream=True):
                    pass
                t2_gflash_elapsed = time.time() - t2_gflash_start
                print(
                    "[TIMING] Chunk %d Tier 2 Gemini Flash: %.1fs"
                    % (chunk_idx, t2_gflash_elapsed)
                )
                if os.path.exists(chunk_output_path):
                    tier2_success = True
                else:
                    print(
                        "[CHUNK %d TIER2] Gemini Flash ran but produced "
                        "no file." % chunk_idx
                    )
            except Exception as e_gflash:
                t2_gflash_elapsed = time.time() - t2_gflash_start
                print(
                    "[CHUNK %d TIER2] Gemini Flash failed after %.1fs: "
                    "%s" % (chunk_idx, t2_gflash_elapsed, str(e_gflash)[:200])  # type: ignore
                )

    # --- Stage 3: Universal OpenAI Fallback (GPT-5.4 -> o3-mini) ---
    if not tier2_success:
        print(
            "[CHUNK %d TIER2] Primary and Lite agents failed. Engaging Universal OpenAI Fallback..."
            % chunk_idx
        )
        try:
            from agents.fallback_openai_agents import (
                get_openai_fallback_agents,  # type: ignore
            )

            _fallback_agents = get_openai_fallback_agents()
            _universal_pro = _fallback_agents["fallback_code_agent"]
            _universal_lite = _fallback_agents["fallback_code_agent_lite"]

            # 1. Try Universal Pro (GPT-5.4) WITH images
            print(
                "[OPENAI FALLBACK TRIGGERED] Chunk %d: Attempting GPT-5.4 Pro Fallback (with visual context)..."
                % chunk_idx
            )
            t2_univ_start = time.time()
            _log_agent_banner(
                agent_name=_universal_pro.name,
                model_id=_universal_pro.model.id,
                provider=getattr(_universal_pro.model, "provider", "openai"),
                step_name="step_generate_chunks / Universal Pro Fallback (chunk %d)"
                % chunk_idx,
            )
            try:
                for _ in _universal_pro.run(code_gen_prompt, stream=True):
                    pass
            except Exception as e_pro:
                print(
                    "[OPENAI FALLBACK TRIGGERED] Chunk %d: Pro failed with visual context: %s"
                    % (chunk_idx, str(e_pro))
                )

            if os.path.exists(chunk_output_path):
                print(
                    "[TIMING] Chunk %d Universal Pro fallback: %.1fs"
                    % (chunk_idx, time.time() - t2_univ_start)
                )
                tier2_success = True
            else:
                # 2. Try Universal Lite (o3-mini) WITH images
                print(
                    "[OPENAI FALLBACK TRIGGERED] Chunk %d: Pro produced no file. Attempting o3-mini Lite Fallback (with visual context)..."
                    % chunk_idx
                )
                t2_univ_lite_start = time.time()
                _log_agent_banner(
                    agent_name=_universal_lite.name,
                    model_id=_universal_lite.model.id,
                    provider=getattr(_universal_lite.model, "provider", "openai"),
                    step_name="step_generate_chunks / Universal Lite Fallback (chunk %d)"
                    % chunk_idx,
                )
                try:
                    for _ in _universal_lite.run(code_gen_prompt, stream=True):
                        pass
                except Exception as e_lite:
                    print(
                        "[OPENAI FALLBACK TRIGGERED] Chunk %d: Lite failed with visual context: %s"
                        % (chunk_idx, str(e_lite))
                    )

                if os.path.exists(chunk_output_path):
                    print(
                        "[TIMING] Chunk %d Universal Lite fallback: %.1fs"
                        % (chunk_idx, time.time() - t2_univ_lite_start)
                    )
                    tier2_success = True

            # If both failed with images, strip images and try again
            if not tier2_success and template_pngs and visual_ref:
                print(
                    "[OPENAI FALLBACK TRIGGERED] Chunk %d: Fallbacks failed with visual context. Retrying without images..."
                    % chunk_idx
                )
                openai_fallback_code_gen_prompt_stripped = code_gen_prompt.replace(
                    "\n" + visual_ref, ""
                )
                print(
                    "[OPENAI FALLBACK TRIGGERED] Chunk %d: Stripped %d chars of base-64 visual reference to limit payload."
                    % (chunk_idx, len(visual_ref))
                )

                # 3. Try Universal Pro (GPT-5.4) WITHOUT images
                print(
                    "[OPENAI FALLBACK TRIGGERED] Chunk %d: Attempting GPT-5.4 Pro Fallback (stripped context)..."
                    % chunk_idx
                )
                t2_univ_stripped_start = time.time()
                _log_agent_banner(
                    agent_name=_universal_pro.name,
                    model_id=_universal_pro.model.id,
                    provider=getattr(_universal_pro.model, "provider", "openai"),
                    step_name="step_generate_chunks / Universal Pro Fallback Stripped (chunk %d)"
                    % chunk_idx,
                )
                try:
                    for _ in _universal_pro.run(
                        openai_fallback_code_gen_prompt_stripped, stream=True
                    ):
                        pass
                except Exception as e_pro_stripped:
                    print(
                        "[OPENAI FALLBACK TRIGGERED] Chunk %d: Pro stripped failed: %s"
                        % (chunk_idx, str(e_pro_stripped))
                    )

                if os.path.exists(chunk_output_path):
                    print(
                        "[TIMING] Chunk %d Universal Pro fallback stripped: %.1fs"
                        % (chunk_idx, time.time() - t2_univ_stripped_start)
                    )
                    tier2_success = True
                else:
                    # 4. Try Universal Lite (o3-mini) WITHOUT images
                    print(
                        "[OPENAI FALLBACK TRIGGERED] Chunk %d: Pro stripped produced no file. Attempting o3-mini Lite Fallback (stripped context)..."
                        % chunk_idx
                    )
                    t2_univ_lite_stripped_start = time.time()
                    _log_agent_banner(
                        agent_name=_universal_lite.name,
                        model_id=_universal_lite.model.id,
                        provider=getattr(_universal_lite.model, "provider", "openai"),
                        step_name="step_generate_chunks / Universal Lite Fallback Stripped (chunk %d)"
                        % chunk_idx,
                    )
                    try:
                        for _ in _universal_lite.run(
                            openai_fallback_code_gen_prompt_stripped, stream=True
                        ):
                            pass
                    except Exception as e_lite_stripped:
                        print(
                            "[OPENAI FALLBACK TRIGGERED] Chunk %d: Lite stripped failed: %s"
                            % (chunk_idx, str(e_lite_stripped))
                        )

                    if os.path.exists(chunk_output_path):
                        print(
                            "[TIMING] Chunk %d Universal Lite fallback stripped: %.1fs"
                            % (chunk_idx, time.time() - t2_univ_lite_stripped_start)
                        )
                        tier2_success = True

        except Exception as e_univ:
            print(
                "[CHUNK %d TIER2] Universal OpenAI Fallback failed entirely: %s"
                % (chunk_idx, str(e_univ))
            )
        print("[CHUNK %d TIER2] Retrying with lite agent (Haiku)..." % chunk_idx)
        t2_lite_start = time.time()
        _log_agent_banner(
            agent_name=getattr(
                _fallback_agent_lite, "name", "PPTX Code Generator (Lite)"
            ),
            model_id=getattr(
                getattr(_fallback_agent_lite, "model", None), "id", "claude-haiku-4-5"
            ),
            provider=getattr(
                getattr(_fallback_agent_lite, "model", None),
                "provider",
                session_state.get("llm_provider", "claude"),
            ),  # type: ignore
            step_name="step_generate_chunks / Tier 2 Lite Fallback (chunk %d)"
            % chunk_idx,
        )
        _get_rate_tracker().check_and_wait(
            model="claude-haiku-4-5",
            prompt=code_gen_prompt,
            caller="generate_chunk_pptx_v2/Tier2-lite",
        )
        try:
            for _ in _fallback_agent_lite.run(code_gen_prompt, stream=True):
                pass
            t2_lite_elapsed = time.time() - t2_lite_start
            print(
                "[TIMING] Chunk %d Tier 2 lite code generation: %.1fs"
                % (chunk_idx, t2_lite_elapsed)
            )
            tier2_success = True
        except Exception as e2:
            t2_lite_elapsed = time.time() - t2_lite_start
            print(
                "[CHUNK %d TIER2] Lite agent (Haiku) also failed after %.1fs: %s"
                % (chunk_idx, t2_lite_elapsed, str(e2))
            )

    if not tier2_success:
        print(
            "[CHUNK %d TIER2] Both agents failed. Falling back to Tier 3 (text-only)."
            % chunk_idx
        )
        return None

    # Verify the file was actually written by the executed code
    if os.path.exists(chunk_output_path):
        try:
            prs = Presentation(chunk_output_path)
            # Clean up empty placeholders and hardcoded contrast issues
            clean_presentation_visual_noise_and_contrast(prs)
            # Fix 16: Remove LLM-generated LINE/freeform/diagonal shapes
            sanitize_llm_shapes(prs)
            sanitize_presentation(prs)  # noqa: F405
            prs.save(chunk_output_path)
            print(
                "[CHUNK %d TIER2] Successfully generated via LLM code execution: %s"
                % (chunk_idx, chunk_output_path)
            )
            return chunk_output_path
        except Exception as e:
            print(
                "[CHUNK %d TIER2] Generated file is invalid PPTX: %s — falling back to Tier 3."
                % (chunk_idx, str(e))
            )
            return None
    else:
        print(
            "[CHUNK %d TIER2] No file produced at %s — falling back to Tier 3."
            % (chunk_idx, chunk_output_path)
        )
        return None


# === WORKFLOW STEP 2: GENERATE CHUNKS ===


def step_generate_chunks(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 2: Orchestrate chunked PPTX generation across all slide groups.

    Splits the full storyboard into chunks of {chunk_size} slides, dispatches
    generation for each chunk using the 3-tier fallback hierarchy, and stores results.

    Each chunk is generated using a 3-tier fallback:
      - Tier 1 (Claude PPTX skill): attempted unless session flag is active
      - Tier 2 (LLM code generation): used when Tier 1 fails or is bypassed
      - Tier 3 (text-only python-pptx): used when Tier 2 also fails
    The session-level flag "use_fallback_generator" persists across chunks —
    once set, Tier 1 is skipped for all remaining chunks in the run.

    A 1-second inter-chunk delay is applied between calls to avoid rate limits.

    Args:
        step_input: Workflow step input (not used directly).
        session_state: Shared workflow state; must contain storyboard, chunk_size, output_dir,
                       max_retries, start_tier, and use_fallback_generator.

    Returns:
        StepOutput with success=True when at least one chunk was generated successfully,
        or success=False if no storyboard was found or all chunks failed.
    """
    step_start = time.time()

    storyboard: Optional[StoryboardPlan] = session_state.get("storyboard")
    if not storyboard:
        print("[ERROR] No storyboard found in session_state.")
        return StepOutput(content="No storyboard found.", success=False)

    chunk_size = session_state.get("chunk_size", 1)
    slides = storyboard.slides

    # Build chunk list
    chunks = [slides[i : i + chunk_size] for i in range(0, len(slides), chunk_size)]

    print("\n" + "=" * 60)
    print("Step 2: Generating presentation chunks...")
    print("=" * 60)
    print(
        "Total slides: %d | Chunk size: %d | Number of chunks: %d"
        % (len(slides), chunk_size, len(chunks))
    )

    if VERBOSE:  # noqa: F405
        for ci, chunk in enumerate(chunks):
            slide_nums = [s.slide_number for s in chunk]
            print("[VERBOSE] Chunk %d: slides %s" % (ci, slide_nums))

    chunk_files: List[Optional[str]] = []
    successful = 0
    total_chunks = len(chunks)
    start_tier = session_state.get("start_tier", 1)

    def _process_chunk(chunk_idx: int, chunk_slides: List[SlideStoryboard]):
        chunk_start = time.time()
        print(
            "[GENERATE] Chunk %d/%d: slides %d-%d"
            % (
                chunk_idx + 1,
                total_chunks,
                chunk_slides[0].slide_number,
                chunk_slides[-1].slide_number,
            )
        )

        # Determine effective starting tier for this chunk
        effective_tier = (
            max(start_tier, 2)
            if session_state.get("use_fallback_generator")
            else start_tier
        )

        chunk_file = None

        if effective_tier == 3:
            # Start with Tier 3: text-only python-pptx (no fallback needed)
            print(
                "[GENERATE] Chunk %d/%d: Starting at Tier 3 (text-only, instant)."
                % (chunk_idx + 1, total_chunks)
            )
            chunk_file = generate_chunk_pptx_fallback(
                chunk_slides, session_state, chunk_idx
            )

        elif effective_tier == 2:
            if session_state.get("use_fallback_generator") and start_tier == 1:
                if VERBOSE:  # noqa: F405
                    print(
                        f"    [FALLBACK AGENT ENGAGED] Switching Tier 1 chunk generation to Tier 2 (LLM code generation) due to prior provider capacity restrictions."
                    )
            # Start with Tier 2: LLM code generation, fallback to Tier 3
            print(
                "[GENERATE] Chunk %d/%d: Starting at Tier 2 (LLM code generation)."
                % (chunk_idx + 1, total_chunks)
            )
            chunk_file = generate_chunk_pptx_v2(chunk_slides, session_state, chunk_idx)
            if chunk_file is None:
                print(
                    "[GENERATE] Chunk %d/%d: Tier 2 failed. Falling back to Tier 3 (text-only)."
                    % (chunk_idx + 1, total_chunks)
                )
                chunk_file = generate_chunk_pptx_fallback(
                    chunk_slides, session_state, chunk_idx
                )

        else:  # effective_tier == 1
            # Start with Tier 1: Claude PPTX skill, fallback to Tier 2 → Tier 3
            print(
                "[GENERATE] Chunk %d/%d: Starting at Tier 1 (Claude PPTX skill)."
                % (chunk_idx + 1, total_chunks)
            )
            chunk_file = generate_chunk_pptx(chunk_slides, session_state, chunk_idx)

            if chunk_file is None and session_state.get("use_fallback_generator"):
                # Tier 1 failed for this chunk (and activated the session flag).
                # Try Tier 2 immediately for this chunk.
                print(
                    "[GENERATE] Chunk %d/%d: Tier 1 failed. Attempting Tier 2 (LLM code generation)..."
                    % (chunk_idx + 1, total_chunks)
                )
                chunk_file = generate_chunk_pptx_v2(
                    chunk_slides, session_state, chunk_idx
                )

                if chunk_file is None:
                    # Tier 2 also failed — fall through to Tier 3.
                    print(
                        "[GENERATE] Chunk %d/%d: Tier 2 also failed. "
                        "Using Tier 3 (text-only fallback)."
                        % (chunk_idx + 1, total_chunks)
                    )
                    chunk_file = generate_chunk_pptx_fallback(
                        chunk_slides, session_state, chunk_idx
                    )

        chunk_elapsed = time.time() - chunk_start
        if chunk_file:
            print(
                "[TIMING] Chunk %d/%d done in %.1fs -> %s"
                % (chunk_idx + 1, total_chunks, chunk_elapsed, chunk_file)
            )
        else:
            print(
                "[TIMING] Chunk %d/%d FAILED after %.1fs (skipping)"
                % (chunk_idx + 1, total_chunks, chunk_elapsed)
            )
            print(
                "[WARNING] Chunk %d failed (slides %d-%d). Continuing..."
                % (
                    chunk_idx,
                    chunk_slides[0].slide_number,
                    chunk_slides[-1].slide_number,
                )
            )

        # Handle rate limit flags per chunk to allow future jitter adjustments
        rate_limit_hit_this_chunk = bool(session_state.get("rate_limit_hit", False))
        if rate_limit_hit_this_chunk:
            session_state["rate_limit_hit"] = False

        return chunk_idx, chunk_file

    async def _generate_all_chunks_async():
        tasks = []
        for chunk_idx, chunk_slides in enumerate(chunks):
            # Stagger chunk initiation to spread out API requests
            if chunk_idx > 0:
                min_del = float(session_state.get("inter_chunk_delay_min", 2000.0))
                max_del = float(session_state.get("inter_chunk_delay_max", 5000.0))
                delay_sec = random.uniform(min_del, max_del) / 1000.0
                print(
                    "[GENERATE] --- Stagger delay before Chunk %d/%d: %.1fs ---"
                    % (chunk_idx + 1, total_chunks, delay_sec)
                )
                await asyncio.sleep(delay_sec)

            # Execute the heavy synchronous chunk generation in a background thread
            tasks.append(asyncio.to_thread(_process_chunk, chunk_idx, chunk_slides))  # type: ignore

        return await asyncio.gather(*tasks)

    # Run the async chunk generation process concurrently
    results = asyncio.run(_generate_all_chunks_async())

    # Reassemble results ensuring original order
    results.sort(key=lambda x: x[0])  # type: ignore
    chunk_files = [res[1] for res in results]  # type: ignore
    successful = sum(1 for f in chunk_files if f is not None)

    session_state["chunk_files"] = chunk_files
    session_state["chunk_slide_groups"] = chunks

    step_elapsed = time.time() - step_start
    failed = total_chunks - successful
    print(
        "\n[TIMING] step_generate_chunks completed in %.1fs (%d chunks: %d succeeded, %d failed)"
        % (step_elapsed, total_chunks, successful, failed)
    )

    summary = "%d of %d chunks generated successfully. Duration: %.1fs" % (
        successful,
        total_chunks,
        step_elapsed,
    )
    return StepOutput(content=summary, success=successful > 0)


# === WORKFLOW STEP 3: PROCESS CHUNKS (TEMPLATE + IMAGES) ===


def step_process_chunks(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 3: Apply template assembly and image pipeline to each chunk.

    For each successfully generated chunk, runs:
    1. Image planning (which slides need AI-generated images)
    2. Image generation (NanoBanana)
    3. Template assembly (if --template is provided)

    Each chunk is processed with a temporary session_state copy that adapts
    the existing step functions to work on individual chunk files.

    Args:
        step_input: Workflow step input (not used directly).
        session_state: Shared workflow state; must contain chunk_files, chunk_slide_groups,
                       template_path, output_dir, and no_images.

    Returns:
        StepOutput with success=True always (failures are logged but do not abort the step).
        Content string reports how many chunks were processed.
    """
    step_start = time.time()

    print("\n" + "=" * 60)
    print("Step 3: Processing chunks (images + template assembly)...")
    print("=" * 60)

    chunk_files: List[Optional[str]] = session_state.get("chunk_files", [])
    chunk_slide_groups: List[List[SlideStoryboard]] = session_state.get(
        "chunk_slide_groups", []
    )
    template_path = session_state.get("template_path", "")
    output_dir = session_state.get("output_dir", ".")
    no_images = session_state.get("no_images", False)

    processed_chunks: Dict[int, Optional[str]] = {}
    total_process_chunks = len(chunk_files)

    for chunk_idx, chunk_file in enumerate(chunk_files):
        chunk_proc_start = time.time()

        if chunk_file is None:
            print(
                "[PROCESS] Chunk %d (%d/%d): skipped (no file)."
                % (chunk_idx, chunk_idx + 1, total_process_chunks)
            )
            processed_chunks[chunk_idx] = None
            continue

        print(
            "\n[PROCESS] Chunk %d (%d/%d): processing %s"
            % (chunk_idx, chunk_idx + 1, total_process_chunks, chunk_file)
        )

        # Determine which slides are in this chunk
        chunk_slides = (
            chunk_slide_groups[chunk_idx] if chunk_idx < len(chunk_slide_groups) else []
        )
        slides_data = _extract_chunk_slides_data(chunk_file)

        # Enrich slides_data entries with storyboard visual_suggestion and has_data_vis.
        # _extract_chunk_slides_data only sees raw PPTX shapes (has_chart, has_table).
        # For Tier 2/3 chunks the storyboard visual_suggestion is the authoritative
        # signal for "this slide will carry an infographic/diagram/chart/table and must
        # NOT receive an external AI-generated image."
        # Without this enrichment the image_planner receives no keyword signal for
        # data-vis slides, so it may incorrectly plan images for them.
        _DATA_VIS_KW = ("chart", "table", "infographic", "diagram", "graph")
        for _sd_i, _sd in enumerate(slides_data):
            if _sd_i < len(chunk_slides):
                _vs = chunk_slides[_sd_i].visual_suggestion or ""  # type: ignore
                _sd["visual_suggestion"] = _vs
                # has_data_vis: true when the storyboard declares a data visual
                # OR the actual PPTX shape inspection found a chart or table.
                _sd["has_data_vis"] = (
                    _sd.get("has_chart", False)
                    or _sd.get("has_table", False)
                    or any(_kw in _vs.lower() for _kw in _DATA_VIS_KW)
                )
                _sd["global_slide_index"] = chunk_slides[_sd_i].slide_number - 1  # type: ignore

        total_chunk_slides = len(slides_data)

        assembled_path = os.path.join(
            output_dir, "chunk_%03d_assembled.pptx" % chunk_idx
        )

        # Build a temporary session_state for the existing step functions
        chunk_session = dict(session_state)
        chunk_session["generated_file"] = chunk_file
        chunk_session["total_slides"] = total_chunk_slides
        chunk_session["global_total_slides"] = session_state.get(
            "total_slides", total_chunk_slides
        )
        chunk_session["slides_data"] = slides_data
        chunk_session["output_path"] = assembled_path
        chunk_session["generated_images"] = {}
        # Use chunk-specific output subdirectory to avoid collisions
        chunk_out_subdir = os.path.join(output_dir, "chunk_%03d_work" % chunk_idx)
        os.makedirs(chunk_out_subdir, exist_ok=True)
        chunk_session["output_dir"] = chunk_out_subdir

        if VERBOSE:  # noqa: F405
            print(
                "[VERBOSE] Chunk %d session state keys: %s"
                % (chunk_idx, sorted(chunk_session.keys()))
            )

        # Adjust src_slide dimensions if not set
        if not chunk_session.get("src_slide_width"):
            try:
                prs = Presentation(chunk_file)
                chunk_session["src_slide_width"] = prs.slide_width
                chunk_session["src_slide_height"] = prs.slide_height
            except Exception:
                pass

        current_file = chunk_file

        # --- Image planning ---
        if not no_images:
            print("[PROCESS] Chunk %d: running image planning..." % chunk_idx)
            try:
                # Build slides JSON for image planner
                slides_json = json.dumps(slides_data, indent=2)
                user_prompt = session_state.get(
                    "user_prompt", "professional presentation"
                )
                combined_message = (
                    'Presentation topic: "%s"\n\nSlide metadata:\n%s\n\n'
                    "Analyze each slide and decide which ones need AI-generated images.\n"
                    "Treat chart/table/infographic/diagram slides as native data-vis and do NOT request image generation for those slides.\n"
                    "Consider the presentation topic when writing image prompts."
                ) % (user_prompt, slides_json)

                from agents import get_agents as _get_agents  # type: ignore

                _image_planner = _get_agents(
                    session_state.get("llm_provider", "claude")
                ).get("image_planner")
                _log_agent_banner(
                    agent_name=getattr(_image_planner, "name", "Image Planner"),
                    model_id=getattr(
                        getattr(_image_planner, "model", None),
                        "id",
                        "gemini-3-flash-preview",
                    ),
                    provider=getattr(
                        getattr(_image_planner, "model", None),
                        "provider",
                        session_state.get("llm_provider", "claude"),
                    ),  # type: ignore
                    step_name="step_process_chunks / Image Planning (chunk %d)"
                    % chunk_idx,
                )
                img_plan_response = _image_planner.run(combined_message, stream=False)

                if img_plan_response and img_plan_response.content:
                    content = img_plan_response.content

                    if VERBOSE:  # noqa: F405
                        if isinstance(content, BaseModel):
                            print(
                                "[VERBOSE] Chunk %d image plan:\n%s"
                                % (chunk_idx, content.model_dump_json(indent=2))
                            )
                        else:
                            print(
                                "[VERBOSE] Chunk %d image plan content: %s"
                                % (chunk_idx, str(content)[:500])  # type: ignore
                            )

                    if isinstance(content, BaseModel):
                        plan_json = content.model_dump_json()
                    elif isinstance(content, dict):
                        plan_json = json.dumps(content)
                    else:
                        plan_json = str(content)

                    # Create a mock StepInput for step_generate_images
                    mock_input = StepInput(
                        input=user_prompt,
                        previous_step_content=plan_json,
                    )
                    step_generate_images(mock_input, chunk_session)  # noqa: F405
                    print(
                        "[PROCESS] Chunk %d: images generated. Count: %d"
                        % (chunk_idx, len(chunk_session.get("generated_images", {})))
                    )
                else:
                    print(
                        "[PROCESS] Chunk %d: image planner returned no plan."
                        % chunk_idx
                    )

            except Exception as e:
                print("[PROCESS] Chunk %d: image pipeline failed: %s" % (chunk_idx, e))
                if session_state.get("verbose"):
                    traceback.print_exc()

        # --- Template assembly ---
        if template_path and os.path.isfile(template_path):
            print("[PROCESS] Chunk %d: running template assembly..." % chunk_idx)
            try:
                # Propagate generated images back to chunk_session
                mock_assemble_input = StepInput(
                    input=session_state.get("user_prompt", ""),
                    previous_step_content=json.dumps(slides_data),
                )
                step_assemble_template(mock_assemble_input, chunk_session)  # noqa: F405

                assembled_output = chunk_session.get("output_path", assembled_path)
                if assembled_output and os.path.isfile(assembled_output):
                    current_file = assembled_output
                    print(
                        "[PROCESS] Chunk %d: assembled -> %s"
                        % (chunk_idx, current_file)
                    )
                else:
                    print(
                        "[PROCESS] Chunk %d: template assembly produced no file; "
                        "keeping raw chunk." % chunk_idx
                    )
            except Exception as e:
                print(
                    "[PROCESS] Chunk %d: template assembly failed: %s" % (chunk_idx, e)
                )
                if session_state.get("verbose"):
                    traceback.print_exc()
        else:
            # No template: just copy raw chunk to assembled path name for consistency
            shutil.copy2(chunk_file, assembled_path)
            current_file = assembled_path
            print(
                "[PROCESS] Chunk %d: no template; raw chunk copied to %s"
                % (chunk_idx, assembled_path)
            )

        processed_chunks[chunk_idx] = current_file

        chunk_proc_elapsed = time.time() - chunk_proc_start
        print(
            "[TIMING] Chunk %d processing done in %.1fs"
            % (chunk_idx, chunk_proc_elapsed)
        )
        print("[PROCESS] Chunk %d: result -> %s" % (chunk_idx, current_file))

    session_state["processed_chunks"] = processed_chunks
    successful = sum(1 for v in processed_chunks.values() if v is not None)

    step_elapsed = time.time() - step_start
    print(
        "\n[TIMING] step_process_chunks completed in %.1fs (%d chunks processed)"
        % (step_elapsed, successful)
    )

    return StepOutput(
        content="%d of %d chunks processed. Duration: %.1fs"
        % (successful, len(chunk_files), step_elapsed),
        success=True,
    )


# === WORKFLOW STEP 4 (OPTIONAL): VISUAL REVIEW PER CHUNK ===


def step_visual_review_chunks(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 4 (Optional): Run visual inspection on each chunk's assembled PPTX.

    For each chunk:
    1. Render slides to PNG via LibreOffice.
    2. Call slide_quality_reviewer for each slide image.
    3. Apply programmatic corrections if needed.
    4. Repeat up to max_passes passes until no further changes are needed.

    This step is non-blocking: any failure silently returns success=True.
    If a programmatic fix is missing in Python, logs it to console regardless
    of --verbose setting.

    Args:
        step_input: Workflow step input (not used directly).
        session_state: Shared workflow state; must contain processed_chunks, output_dir,
                       template_path, and visual_passes.

    Returns:
        StepOutput with success=True always; errors per chunk are logged and do not abort.
        Content string reports how many chunks were reviewed.
    """
    step_start = time.time()

    print("\n" + "=" * 60)
    print("Step 4 (Optional): Visual review per chunk...")
    print("=" * 60)

    processed_chunks: Dict[int, Optional[str]] = session_state.get(
        "processed_chunks", {}
    )
    output_dir = session_state.get("output_dir", ".")
    template_path = session_state.get("template_path", "")
    max_passes = session_state.get("visual_passes", 3)
    reviewed_chunks: Dict[int, Optional[str]] = {}

    def _review_chunk(chunk_idx: int, assembled_path: str):
        chunk_review_start = time.time()

        if assembled_path is None or not os.path.isfile(assembled_path):
            print(
                "[VISUAL] Chunk %d: skipped (file not found: %s)."
                % (chunk_idx, assembled_path)
            )
            return chunk_idx, None

        print(
            "\n[VISUAL REVIEW] Chunk %d: starting review of %s"
            % (chunk_idx, assembled_path)
        )
        from agents import get_agents as _get_agents_vr  # type: ignore

        _vr_reviewer = _get_agents_vr(session_state.get("llm_provider", "claude")).get(
            "slide_quality_reviewer"
        )
        _log_agent_banner(
            agent_name=getattr(
                _vr_reviewer, "name", "Senior UI/UX Presentation Designer"
            ),
            model_id=getattr(
                getattr(_vr_reviewer, "model", None), "id", "gemini-2.5-flash"
            ),
            provider=getattr(
                getattr(_vr_reviewer, "model", None),
                "provider",
                session_state.get("llm_provider", "claude"),
            ),  # type: ignore
            step_name="step_visual_review_chunks / Visual QA (chunk %d)" % chunk_idx,
        )

        # Build a per-chunk session_state for the visual review step
        chunk_session = dict(session_state)
        chunk_session["output_path"] = assembled_path
        chunk_session["template_path"] = template_path

        current_path = assembled_path

        for pass_num in range(max_passes):
            pass_start = time.time()
            print(
                "[VISUAL REVIEW] Chunk %d: pass %d/%d starting..."
                % (chunk_idx, pass_num + 1, max_passes)
            )
            chunk_session["output_path"] = current_path

            try:
                mock_input = StepInput(
                    input=session_state.get("user_prompt", ""),
                    previous_step_content="",
                )
                result = step_visual_quality_review(mock_input, chunk_session)  # noqa: F405

                # Check if any actionable issues exist (mirrors _apply_visual_corrections logic).
                # PresentationQualityReport has 'total_critical_issues', NOT 'total_corrections_applied'.
                # We check slide_reports for critical/moderate issues with a real programmatic_fix.
                quality_report = chunk_session.get("quality_report", {})
                slide_reports_data = quality_report.get("slide_reports", [])

                if VERBOSE:  # noqa: F405
                    for r in slide_reports_data:
                        issues = r.get("issues", [])
                        print(
                            "[VERBOSE] Chunk %d pass %d slide %s: %d issues"
                            % (
                                chunk_idx,
                                pass_num + 1,
                                r.get("slide_index", "?"),
                                len(issues),
                            )
                        )
                        for issue in issues:
                            print(
                                "[VERBOSE]   severity=%s fix=%s desc=%s"
                                % (
                                    issue.get("severity", "?"),
                                    issue.get("programmatic_fix", "?"),
                                    str(issue.get("description", ""))[:80],  # type: ignore
                                )
                            )

                changes_applied = any(
                    any(
                        i.get("severity") in ("critical", "moderate")
                        and i.get("programmatic_fix") != "none"
                        for i in r.get("issues", [])
                    )
                    for r in slide_reports_data
                )

                pass_elapsed = time.time() - pass_start
                print(
                    "[TIMING] Chunk %d pass %d: %.1fs"
                    % (chunk_idx, pass_num + 1, pass_elapsed)
                )

                if not changes_applied:
                    print(
                        "[VISUAL REVIEW] Chunk %d: pass %d/%d — no changes needed. Done."
                        % (chunk_idx, pass_num + 1, max_passes)
                    )
                    break
                else:
                    print(
                        "[VISUAL REVIEW] Chunk %d: pass %d/%d — corrections applied. Re-checking..."
                        % (chunk_idx, pass_num + 1, max_passes)
                    )

            except Exception as e:
                pass_elapsed = time.time() - pass_start
                print(
                    "[TIMING] Chunk %d pass %d: %.1fs (error)"
                    % (chunk_idx, pass_num + 1, pass_elapsed)
                )
                print(
                    "[VISUAL] Chunk %d, pass %d: review failed: %s"
                    % (chunk_idx, pass_num + 1, e)
                )
                # Log missing programmatic fix to console regardless of verbose mode (per spec).
                # Any exception here means the visual review or correction logic is broken/missing.
                print(
                    "[VISUAL REVIEW MISSING FIX] Chunk %d, pass %d: exception during "
                    "visual correction step: %s" % (chunk_idx, pass_num + 1, str(e))
                )
                print(
                    "[SUGGESTION] Review step_visual_quality_review() and "
                    "_apply_visual_corrections() for the issue type that raised this error. "
                    "Add handling logic if a programmatic_fix type is missing."
                )
                break

        reviewed_path = chunk_session.get("output_path", current_path)
        chunk_review_elapsed = time.time() - chunk_review_start
        print(
            "[TIMING] Chunk %d total review: %.1fs" % (chunk_idx, chunk_review_elapsed)
        )
        print("[VISUAL REVIEW] Chunk %d: reviewed -> %s" % (chunk_idx, reviewed_path))
        return chunk_idx, reviewed_path

    async def _review_all_chunks_async():
        tasks = []
        for chunk_idx, assembled_path in sorted(processed_chunks.items()):
            tasks.append(asyncio.to_thread(_review_chunk, chunk_idx, assembled_path))  # type: ignore
        return await asyncio.gather(*tasks)

    results = asyncio.run(_review_all_chunks_async())

    for c_idx, r_path in results:  # type: ignore
        reviewed_chunks[c_idx] = r_path

    session_state["reviewed_chunks"] = reviewed_chunks
    reviewed_count = sum(1 for v in reviewed_chunks.values() if v is not None)

    step_elapsed = time.time() - step_start
    print(
        "\n[TIMING] step_visual_review_chunks completed in %.1fs (%d chunks reviewed)"
        % (step_elapsed, reviewed_count)
    )

    return StepOutput(
        content="%d of %d chunks visually reviewed. Duration: %.1fs"
        % (reviewed_count, len(processed_chunks), step_elapsed),
        success=True,
    )


# === HELPER: MERGE MULTIPLE PPTX FILES ===


def _merge_pptx_zip_level(pptx_paths: List[str], output_path: str) -> bool:
    """Merge multiple PPTX files by manipulating the ZIP structure directly.

    This is the most reliable approach — avoids OPC context issues that cause
    PowerPoint to report "found a problem with content" on the merged file.
    Binary parts (images, charts, workbooks) are copied at the raw bytes level,
    so there is no risk of the OPC package context dropping or corrupting data.

    Args:
        pptx_paths: List of valid, existing PPTX file paths to merge in order.
        output_path: Destination path for the merged presentation.

    Returns:
        True if merge succeeded, False otherwise.
    """
    import posixpath
    import re

    from lxml import etree  # type: ignore

    valid_paths = [str(p) for p in pptx_paths if p and os.path.exists(p)]
    if not valid_paths:
        print("[MERGE] No valid PPTX files to merge")
        return False

    if len(valid_paths) == 1:
        shutil.copy(valid_paths[0], output_path)
        print("[MERGE] Single file, copied directly: %s" % output_path)
        return True

    # Use first file as base — copy it to output
    shutil.copy(valid_paths[0], output_path)

    # Read the base presentation XML
    with zipfile.ZipFile(output_path, "r") as base_zip:
        base_prs_xml = base_zip.read("ppt/presentation.xml")
        base_prs_rels_xml = base_zip.read("ppt/_rels/presentation.xml.rels")
        base_content_types_xml = base_zip.read("[Content_Types].xml")

    base_prs_tree = etree.fromstring(base_prs_xml)
    base_prs_rels_tree = etree.fromstring(base_prs_rels_xml)
    base_ct_tree = etree.fromstring(base_content_types_xml)

    # Namespaces
    NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_RELS = "http://schemas.openxmlformats.org/package/2006/relationships"
    NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"

    def _get_slide_numbers(prs_rels_tree):
        """Extract slide XML numbers from ``presentation.xml.rels``.

        Args:
            prs_rels_tree: Parsed lxml tree for ``ppt/_rels/presentation.xml.rels``.

        Returns:
            List of integer slide numbers extracted from Target attributes of slide
            relationships in prs_rels_tree. Non-slide relationships are ignored.
        """
        slide_nums = []
        for rel in prs_rels_tree.findall("{%s}Relationship" % NS_RELS):
            target = rel.get("Target", "")
            m = re.match(r"slides/slide(\d+)\.xml", target)
            if m:
                slide_nums.append(int(m.group(1)))
        return slide_nums

    existing_slide_nums = _get_slide_numbers(base_prs_rels_tree)
    next_slide_num = max(existing_slide_nums) + 1 if existing_slide_nums else 1

    # Track next rel ID for presentation.xml.rels
    existing_rel_ids = [
        int(rel.get("Id", "rId0").replace("rId", ""))
        for rel in base_prs_rels_tree.findall("{%s}Relationship" % NS_RELS)
        if rel.get("Id", "").startswith("rId")
    ]
    next_rel_id = max(existing_rel_ids) + 1 if existing_rel_ids else 100

    # Open output as writable archive
    with zipfile.ZipFile(
        output_path, "a", compression=zipfile.ZIP_DEFLATED, allowZip64=True
    ) as out_zip:
        for src_path in valid_paths[1:]:  # type: ignore
            with zipfile.ZipFile(src_path, "r") as src_zip:
                src_names = set(src_zip.namelist())

                # Read source presentation rels to find slides
                src_prs_rels_xml = src_zip.read("ppt/_rels/presentation.xml.rels")
                src_prs_rels = etree.fromstring(src_prs_rels_xml)
                src_slide_nums = sorted(_get_slide_numbers(src_prs_rels))

                for src_slide_num in src_slide_nums:
                    new_slide_num = next_slide_num
                    next_slide_num += 1  # type: ignore

                    old_slide_name = "ppt/slides/slide%d.xml" % src_slide_num
                    new_slide_name = "ppt/slides/slide%d.xml" % new_slide_num
                    old_slide_rels_name = (
                        "ppt/slides/_rels/slide%d.xml.rels" % src_slide_num
                    )
                    new_slide_rels_name = (
                        "ppt/slides/_rels/slide%d.xml.rels" % new_slide_num
                    )

                    if old_slide_name not in src_names:
                        continue

                    # Copy slide XML
                    slide_xml_bytes = src_zip.read(old_slide_name)
                    out_zip.writestr(new_slide_name, slide_xml_bytes)

                    # Copy slide rels and rewrite media/chart refs with unique names
                    if old_slide_rels_name in src_names:
                        slide_rels_xml = src_zip.read(old_slide_rels_name)
                        slide_rels_tree = etree.fromstring(slide_rels_xml)

                        for rel in slide_rels_tree.findall(
                            "{%s}Relationship" % NS_RELS
                        ):
                            rel_type = rel.get("Type", "")
                            target = rel.get("Target", "")

                            # Skip slide layout — keep reference as-is
                            if "slideLayout" in rel_type:
                                continue

                            if not target.startswith(".."):
                                continue  # absolute or external refs

                            # Resolve the actual part path in the source zip
                            # e.g. "../media/image1.png" -> "ppt/media/image1.png"
                            # NOTE: lstrip("../") is character-based and strips individual
                            # '.' and '/' chars, giving wrong results for nested paths like
                            # "../charts/chart1.xml" (produces "ppt/slides/charts/chart1.xml"
                            # instead of "ppt/charts/chart1.xml"). Use posixpath.normpath
                            # to handle the ".." parent-dir segment correctly.
                            actual_old = posixpath.normpath("ppt/slides/" + target)

                            if actual_old not in src_names:
                                continue

                            # Generate unique name for target archive
                            basename = os.path.basename(actual_old)
                            stem, ext = os.path.splitext(basename)
                            new_part_name = actual_old
                            counter = 1
                            all_names = set(out_zip.namelist())
                            while new_part_name in all_names:
                                new_part_name = (
                                    os.path.dirname(actual_old)
                                    + "/"
                                    + stem
                                    + "_s%d_%d" % (new_slide_num, counter)
                                    + ext
                                )
                                counter += 1

                            # Copy the part
                            part_bytes = src_zip.read(actual_old)
                            out_zip.writestr(new_part_name, part_bytes)

                            # If chart, also copy its rels and embedded workbook
                            if "chart" in actual_old:
                                chart_basename = os.path.basename(actual_old)
                                chart_rels_old = (
                                    os.path.dirname(actual_old)
                                    + "/_rels/"
                                    + chart_basename
                                    + ".rels"
                                )
                                if chart_rels_old in src_names:
                                    cr_bytes = src_zip.read(chart_rels_old)
                                    cr_new = (
                                        os.path.dirname(new_part_name)
                                        + "/_rels/"
                                        + os.path.basename(new_part_name)
                                        + ".rels"
                                    )
                                    # Copy chart's embedded xlsx workbook and
                                    # build an updated chart rels XML that points
                                    # to the newly-renamed workbook file.
                                    # NOTE: wb_old uses posixpath.normpath to
                                    # correctly resolve "../embeddings/file.xlsx"
                                    # relative to ppt/charts/ (not ppt/slides/).
                                    # lstrip("../") would incorrectly strip leading
                                    # '.' and '/' characters individually.
                                    cr_tree = etree.fromstring(cr_bytes)
                                    for cr_rel in cr_tree.findall(
                                        "{%s}Relationship" % NS_RELS
                                    ):
                                        cr_target = cr_rel.get("Target", "")
                                        if cr_target.startswith(".."):
                                            wb_old = posixpath.normpath(
                                                os.path.dirname(actual_old)
                                                + "/"
                                                + cr_target
                                            )
                                            if wb_old in src_names:
                                                wb_bytes = src_zip.read(wb_old)
                                                wb_stem, wb_ext = os.path.splitext(
                                                    os.path.basename(wb_old)
                                                )
                                                wb_new = wb_old.replace(
                                                    os.path.basename(wb_old),
                                                    wb_stem
                                                    + "_s%d_wb" % new_slide_num  # type: ignore
                                                    + wb_ext,
                                                )
                                                if wb_new not in set(
                                                    out_zip.namelist()
                                                ):
                                                    out_zip.writestr(wb_new, wb_bytes)
                                                # Update the chart rels entry to point
                                                # to the renamed workbook.
                                                # The new target must be relative to the
                                                # chart's directory (ppt/charts/), so we
                                                # build a "../<subdir>/<name>" path.
                                                new_wb_rel_target = "../" + "/".join(
                                                    wb_new.split("/")[1:]  # type: ignore
                                                )
                                                cr_rel.set("Target", new_wb_rel_target)
                                    # Serialise the (possibly updated) rels tree.
                                    updated_cr_bytes = etree.tostring(
                                        cr_tree,
                                        xml_declaration=True,
                                        encoding="UTF-8",
                                        standalone=True,
                                    )
                                    out_zip.writestr(cr_new, updated_cr_bytes)

                            # Update the slide relationship target to point to the new
                            # unique part name.
                            # NOTE: split("/")[2:] was wrong — it drops the subdirectory,
                            # e.g. "ppt/charts/chart1_s5_1.xml".split("/")[2:] yields
                            # ["chart1_s5_1.xml"] (missing "charts"), giving the incorrect
                            # relative target "../chart1_s5_1.xml".
                            # split("/")[1:] yields ["charts", "chart1_s5_1.xml"],
                            # producing the correct "../charts/chart1_s5_1.xml".
                            new_rel_target = "../" + "/".join(
                                new_part_name.split("/")[1:]  # type: ignore
                            )
                            rel.set("Target", new_rel_target)

                        updated_rels_bytes = etree.tostring(
                            slide_rels_tree,
                            xml_declaration=True,
                            encoding="UTF-8",
                            standalone=True,
                        )
                        out_zip.writestr(new_slide_rels_name, updated_rels_bytes)

                    # Register slide in presentation.xml.rels
                    new_rel_id = "rId%d" % next_rel_id
                    next_rel_id += 1
                    new_prs_rel = etree.SubElement(
                        base_prs_rels_tree, "{%s}Relationship" % NS_RELS
                    )
                    new_prs_rel.set("Id", new_rel_id)
                    new_prs_rel.set(
                        "Type",
                        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
                    )
                    new_prs_rel.set("Target", "slides/slide%d.xml" % new_slide_num)

                    # Register slide in presentation.xml sldIdLst
                    sld_id_lst = base_prs_tree.find(".//{%s}sldIdLst" % NS_P)
                    if sld_id_lst is None:
                        sld_id_lst = etree.SubElement(
                            base_prs_tree, "{%s}sldIdLst" % NS_P
                        )

                    existing_ids = [
                        int(el.get("id", 256))
                        for el in sld_id_lst.findall("{%s}sldId" % NS_P)
                    ]
                    new_id = max(existing_ids) + 1 if existing_ids else 256
                    sld_id_el = etree.SubElement(sld_id_lst, "{%s}sldId" % NS_P)
                    sld_id_el.set("id", str(new_id))
                    sld_id_el.set("{%s}id" % NS_R, new_rel_id)

                    # Add content type entry for new slide
                    existing_ct_parts = {
                        el.get("PartName", "")
                        for el in base_ct_tree.findall("{%s}Override" % NS_CT)
                    }
                    new_part_uri = "/ppt/slides/slide%d.xml" % new_slide_num
                    if new_part_uri not in existing_ct_parts:
                        ct_el = etree.SubElement(base_ct_tree, "{%s}Override" % NS_CT)
                        ct_el.set("PartName", new_part_uri)
                        ct_el.set(
                            "ContentType",
                            "application/vnd.openxmlformats-officedocument.presentationml.slide+xml",
                        )

    # Write updated presentation.xml, rels, and content types back.
    # Python's zipfile does not support overwriting entries, so we copy to a temp file.
    import tempfile

    tmp_path = output_path + ".tmp"
    with zipfile.ZipFile(output_path, "r") as old_zip:
        with zipfile.ZipFile(
            tmp_path, "w", compression=zipfile.ZIP_DEFLATED, allowZip64=True
        ) as new_zip:
            for item in old_zip.namelist():
                if item == "ppt/presentation.xml":
                    new_zip.writestr(
                        item,
                        etree.tostring(
                            base_prs_tree,
                            xml_declaration=True,
                            encoding="UTF-8",
                            standalone=True,
                        ),
                    )
                elif item == "ppt/_rels/presentation.xml.rels":
                    new_zip.writestr(
                        item,
                        etree.tostring(
                            base_prs_rels_tree,
                            xml_declaration=True,
                            encoding="UTF-8",
                            standalone=True,
                        ),
                    )
                elif item == "[Content_Types].xml":
                    new_zip.writestr(
                        item,
                        etree.tostring(
                            base_ct_tree,
                            xml_declaration=True,
                            encoding="UTF-8",
                            standalone=True,
                        ),
                    )
                else:
                    new_zip.writestr(item, old_zip.read(item))

    os.replace(tmp_path, output_path)
    print("[MERGE] Saved merged presentation: %s" % output_path)
    return True


def _try_auto_repair_with_libreoffice(pptx_path: str) -> bool:
    """Attempt to auto-repair a PPTX by converting it through LibreOffice.

    Only runs if LibreOffice is available. Non-destructive on failure —
    the original file is left unchanged if the conversion fails.

    Args:
        pptx_path: Path to the PPTX file to repair in-place.

    Returns:
        True if LibreOffice repair succeeded, False otherwise.
    """
    import glob
    import subprocess
    import tempfile

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return False

    tmp_dir = tempfile.mkdtemp(prefix="pptx_repair_")
    try:
        result = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "pptx",
                "--outdir",
                tmp_dir,
                pptx_path,
            ],
            capture_output=True,
            timeout=180,
        )
        if result.returncode == 0:
            converted = glob.glob(os.path.join(tmp_dir, "*.pptx"))
            if converted:
                shutil.copy(converted[0], pptx_path)
                print("[MERGE] Auto-repair via LibreOffice succeeded: %s" % pptx_path)
                return True
    except Exception as e:
        print("[MERGE] Auto-repair via LibreOffice failed: %s" % e)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)
    return False


def merge_pptx_files(pptx_paths: List[str], output_path: str) -> bool:
    """Merge multiple PPTX files into a single presentation using ZIP-level manipulation.

    Uses _merge_pptx_zip_level() which copies all binary parts (images, charts,
    workbooks) at the raw bytes level, avoiding OPC package context issues that
    cause PowerPoint to report "found a problem with content" on the merged file.

    Args:
        pptx_paths: List of PPTX file paths to merge in order.
        output_path: Destination path for the merged presentation.

    Returns:
        True if merge succeeded, False otherwise.
    """
    merge_start = time.time()
    valid_paths = [p for p in pptx_paths if p and os.path.exists(p)]
    print("[MERGE] Merging %d PPTX files into %s" % (len(valid_paths), output_path))
    if VERBOSE:  # noqa: F405
        for i, p in enumerate(valid_paths):
            print("[VERBOSE][MERGE] Source %d: %s" % (i, p))
    result = _merge_pptx_zip_level(valid_paths, output_path)
    merge_elapsed = time.time() - merge_start
    print("[TIMING] merge_pptx_files completed in %.1fs" % merge_elapsed)
    return result


# === WORKFLOW STEP 5 (FINAL): MERGE ALL CHUNKS ===


def step_merge_chunks(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 5 (Final): Merge all processed/reviewed chunk PPTX files into the final output.

    Source selection priority (explicit, robust):
    1. Template + visual review + reviewed_chunks present -> use reviewed_chunks
    2. Template + processed_chunks present             -> use processed_chunks
    3. No template (raw mode)                          -> use raw chunk_files

    Chunks are merged in order (by chunk_idx).
    """
    step_start = time.time()

    print("\n" + "=" * 60)
    print("Step 5 (Final): Merging chunks into final presentation...")
    print("=" * 60)

    output_path = session_state.get("output_path", "presentation_chunked.pptx")
    has_template = bool(session_state.get("template_path"))
    visual_review = session_state.get("visual_review", False)
    chunk_files: List[Optional[str]] = session_state.get("chunk_files", [])
    processed_chunks: Dict[int, Optional[str]] = session_state.get(
        "processed_chunks", {}
    )
    reviewed_chunks: Dict[int, Optional[str]] = session_state.get("reviewed_chunks", {})

    # Determine which chunk paths to use (priority: reviewed > processed > raw)
    if has_template and visual_review and reviewed_chunks:
        source_label = "reviewed (template + visual review)"
        ordered_paths = [reviewed_chunks.get(i) for i in sorted(reviewed_chunks.keys())]
    elif has_template and processed_chunks:
        source_label = "processed (template-assembled)"
        ordered_paths = [
            processed_chunks.get(i) for i in sorted(processed_chunks.keys())
        ]
    else:
        # No template path: use raw chunk files directly
        source_label = "raw (no template)"
        ordered_paths = [f for f in chunk_files if f is not None]

    if not ordered_paths:
        print("[MERGE] No chunk files found to merge")
        return StepOutput(
            content="No files to merge",
            success=False,
        )

    print(
        "Merging from: %s (%d total, %d valid)"
        % (
            source_label,
            len(ordered_paths),
            sum(1 for p in ordered_paths if p and os.path.exists(p)),
        )
    )

    if VERBOSE:  # noqa: F405
        print("[VERBOSE] Ordered chunk files for merge:")
        for i, p in enumerate(ordered_paths):
            print("[VERBOSE]   %d. %s" % (i, p))

    success = merge_pptx_files(
        [str(p) for p in ordered_paths if p],
        output_path,
    )

    # Attempt optional auto-repair (only if LibreOffice is available)
    if success:
        _try_auto_repair_with_libreoffice(output_path)

    step_elapsed = time.time() - step_start
    final_file = os.path.basename(output_path)
    print(
        "[TIMING] step_merge_chunks completed in %.1fs (final: %s)"
        % (step_elapsed, final_file)
    )

    if success:
        summary = "Merged %d chunks (%s) -> %s. Duration: %.1fs" % (
            len([p for p in ordered_paths if p]),
            source_label,
            output_path,
            step_elapsed,
        )
        print("[MERGE] %s" % summary)
        return StepOutput(
            content=summary,
            success=True,
        )
    else:
        return StepOutput(
            content="Merge failed. No output file produced.",
            success=False,
        )


# === WORKFLOW BUILDER ===


def build_chunked_workflow(session_state: Dict) -> Workflow:
    """Build the chunked PPTX workflow with the appropriate set of steps.

    Steps included:
    - Step 1: Optimize & Plan   (always)
    - Step 2: Generate Chunks   (always)
    - Step 3: Process Chunks    (only when template_path is set)
    - Step 4: Visual Review     (when visual_review is set; works with or without template)
    - Step 5: Merge Chunks      (always)

    No-template pipeline: Step 1 -> Step 2 [-> Step 4] -> Step 5
    Template pipeline:    Step 1 -> Step 2 -> Step 3 [-> Step 4] -> Step 5

    Args:
        session_state: Shared workflow state; must contain template_path and visual_review
                       to determine which optional steps to include.

    Returns:
        Configured Workflow instance with the appropriate step sequence.
    """
    has_template = bool(session_state.get("template_path"))
    do_visual_review = bool(session_state.get("visual_review"))

    steps = [
        Step(name="Optimize and Plan", executor=step_optimize_and_plan),
        Step(name="Generate Chunks", executor=step_generate_chunks),
    ]

    # Template assembly + image pipeline only runs when a template is provided
    if has_template:
        steps.append(Step(name="Process Chunks", executor=step_process_chunks))

    # Visual review: runs with or without template (template-independent contrast checks available)
    if do_visual_review:
        steps.append(
            Step(name="Visual Review Chunks", executor=step_visual_review_chunks)
        )

    steps.append(Step(name="Merge Chunks", executor=step_merge_chunks))

    return Workflow(
        name="Chunked PPTX Workflow",
        steps=steps,
        session_state=session_state,
    )


# === MAIN ENTRY POINT ===


def main() -> None:
    """Parse CLI flags, build session state, and execute the chunked workflow.

    Validates required environment variables and the optional template path, creates
    output working directories under output_chunked/, runs the workflow end-to-end,
    and writes the final PPTX to the path specified by --output.

    All output (print statements, agno framework logging, and stderr) is redirected
    to OUTPUT.md in the script directory when run via __main__. The console stays
    silent. OUTPUT.md is cleared at the start of each run.
    """
    parser = argparse.ArgumentParser(
        description="Chunked PPTX generation workflow — overcomes Claude API limits for large presentations."
    )

    # Existing args (compatible with powerpoint_template_workflow.py)
    parser.add_argument(
        "--template",
        "-t",
        default=None,
        help="Path to .pptx template file (optional). Without it, skips template assembly.",
    )
    parser.add_argument(
        "--output",
        "-o",
        default="presentation_chunked.pptx",
        help="Output filename (default: presentation_chunked.pptx).",
    )
    parser.add_argument(
        "--prompt",
        "-p",
        default=None,
        help="User prompt describing the presentation topic.",
    )
    parser.add_argument(
        "--no-images",
        action="store_true",
        help="Skip AI image generation.",
    )
    parser.add_argument(
        "--no-stream",
        action="store_true",
        help="Disable streaming mode for Claude agent.",
    )
    parser.add_argument(
        "--min-images",
        type=int,
        default=1,
        help="Minimum slides that must have AI-generated images (default: 1).",
    )
    parser.add_argument(
        "--visual-review",
        action="store_true",
        help="Enable visual QA with Gemini vision per chunk (requires LibreOffice + template).",
    )
    parser.add_argument(
        "--template-visuals",
        "-tv",
        action="store_true",
        help="Inject base64 images of template slides into LLM prompts to improve visual accuracy (consumes more tokens).",
    )
    parser.add_argument(
        "--footer-text",
        default="",
        help="Footer text for all slides (idx=11 placeholder).",
    )
    parser.add_argument(
        "--date-text",
        default="",
        help="Date text for footer date placeholder (idx=10).",
    )
    parser.add_argument(
        "--show-slide-numbers",
        action="store_true",
        help="Preserve slide number placeholder (idx=12) on all slides.",
    )
    parser.add_argument(
        "--verbose",
        "-v",
        action="store_true",
        help="Enable verbose/debug logging.",
    )

    # LLM provider selection
    parser.add_argument(
        "--llm-provider",
        choices=["claude", "openai", "gemini"],
        default="claude",
        help=(
            "LLM provider for swappable agents (brand analyzer, query optimizer, "
            "fallback code gen, image planner, visual reviewer). "
            "The Content Generator always uses Claude (PPTX skill). Default: claude."
        ),
    )

    # New args for chunked workflow
    parser.add_argument(
        "--chunk-size",
        type=int,
        default=1,
        help="Number of slides per LLM API chunk call (default: 1). "
        "Using 1 ensures each chunk sends only the single best-matching "
        "template slide image, keeping prompts within all model context windows.",
    )
    parser.add_argument(
        "--max-retries",
        type=int,
        default=2,
        help="Max retries per chunk on failure (default: 2).",
    )
    parser.add_argument(
        "--visual-passes",
        type=int,
        default=3,
        help="Maximum visual inspection passes per chunk (default: 3).",
    )
    parser.add_argument(
        "--start-tier",
        type=int,
        choices=[1, 2, 3],
        default=1,
        help=(
            "Starting tier for chunk generation (default: 1). "
            "1=Claude PPTX skill (best quality), "
            "2=LLM code generation (80-92%% quality, faster, python-pptx native charts), "
            "3=text-only (structural, instant). "
            "Fallback continues from selected tier."
        ),
    )
    parser.add_argument(
        "--inter-chunk-delay-min",
        type=float,
        default=None,
        metavar="MS",
        help=(
            "Minimum inter-chunk delay in milliseconds (default: provider-specific). "
            "A random value in [min, max] is chosen between each chunk."
        ),
    )
    parser.add_argument(
        "--inter-chunk-delay-max",
        type=float,
        default=None,
        metavar="MS",
        help=(
            "Maximum inter-chunk delay in milliseconds (default: provider-specific). "
            "When a 429 rate-limit error is detected, max_delay is used directly."
        ),
    )

    args = parser.parse_args()

    # Update module-level VERBOSE (imported from powerpoint_template_workflow via *)
    global VERBOSE  # noqa: F405
    VERBOSE = args.verbose  # noqa: F405

    # Initialize observability via Langfuse (Agno / OpenInference)
    tracer_provider = setup_langfuse_telemetry()
    try:
        _run_main_workflow(args)
    finally:
        if tracer_provider:
            if VERBOSE:
                print("[TELEMETRY] Flushing and shutting down tracer...")
            tracer_provider.shutdown()


def _run_main_workflow(args):
    # Validate API keys
    # ANTHROPIC_API_KEY is always required (Content Generator is locked to Claude).
    if not os.getenv("ANTHROPIC_API_KEY"):
        print("Error: ANTHROPIC_API_KEY environment variable not set.")
        print("  (Required: Content Generator agent always uses Claude.)")
        sys.exit(1)
    # OPENAI_API_KEY is always required: brand_style_analyzer uses gpt-4o-mini (OpenAI)
    # regardless of --llm-provider, to avoid consuming Anthropic token budget.
    if not os.getenv("OPENAI_API_KEY"):
        print(
            "[WARNING] OPENAI_API_KEY not set — brand style analysis (gpt-4o-mini) will be skipped."
        )
        print(
            "[WARNING] Two-stage brand parsing will fall through to keyword-only detection."
        )
    if args.llm_provider == "openai" and not os.getenv("OPENAI_API_KEY"):
        print(
            "Error: OPENAI_API_KEY environment variable not set (required for --llm-provider openai)."
        )
        sys.exit(1)
    if args.llm_provider == "gemini" and not os.getenv("GOOGLE_API_KEY"):
        print(
            "Error: GOOGLE_API_KEY environment variable not set (required for --llm-provider gemini)."
        )
        sys.exit(1)

    # Validate template if provided
    if args.template is not None:
        if not os.path.isfile(args.template):
            print("Error: Template file not found: %s" % args.template)
            sys.exit(1)
        if not args.template.endswith(".pptx"):
            print("Error: Template file must be a .pptx file.")
            sys.exit(1)

    # Gemini API key validation guard for --visual-review
    # If the key is missing or blank, auto-disable visual review to avoid burning
    # time on guaranteed-to-fail Gemini API calls.
    google_key = os.getenv("GOOGLE_API_KEY", "")
    if getattr(args, "visual_review", False) and not google_key:
        print(
            "[WARNING] --visual-review requested but GOOGLE_API_KEY is not set. "
            "Visual review automatically disabled."
        )
        setattr(args, "visual_review", False)

    # Effective values: visual review and passes work with or without template
    effective_visual_review = bool(args.visual_review)
    effective_visual_passes = args.visual_passes

    # Reset the rate tracker singleton for this run
    _reset_rate_tracker()
    print(
        "[RATE TRACKER] Rate limit tracker initialised. "
        "Claude model limits: sonnet=30K, opus=30K, haiku=50K input tokens/min."
    )
    # Apply provider-specific millisecond defaults if not set
    provider_delays = DEFAULT_INTER_CHUNK_DELAYS_MS.get(
        args.llm_provider, DEFAULT_INTER_CHUNK_DELAYS_MS["claude"]
    )
    if getattr(args, "inter_chunk_delay_min", None) is None:
        setattr(args, "inter_chunk_delay_min", provider_delays["min"])
    if getattr(args, "inter_chunk_delay_max", None) is None:
        setattr(args, "inter_chunk_delay_max", provider_delays["max"])

    print(
        "[RATE TRACKER] Inter-chunk logic set to: random %.0f–%.0f ms "
        "(override with --inter-chunk-delay-min / --inter-chunk-delay-max)."
        % (
            getattr(args, "inter_chunk_delay_min", 2000),
            getattr(args, "inter_chunk_delay_max", 5000),
        )
    )

    # Setup output directory with unique session ID + timestamp per run
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_base = os.path.join(script_dir, "output_chunked")
    os.makedirs(output_base, exist_ok=True)

    # Generate unique session identifier: short uuid + timestamp
    session_id = uuid.uuid4().hex[:8]  # type: ignore
    session_timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    session_name = "session_%s_%s" % (session_id, session_timestamp)

    # Provider name is stored in session_state (a plain string, safely deep-copyable).
    # Agents are loaded lazily inside each step via get_agents() to avoid pickling errors.

    # Chunked workflow uses a session-specific working directory
    session_dir = os.path.join(output_base, "chunked_workflow_work", session_name)
    os.makedirs(session_dir, exist_ok=True)
    output_dir = session_dir

    # Resolve output path — place the final .pptx inside the session directory
    output_path = args.output
    if not os.path.isabs(output_path):
        output_path = os.path.join(session_dir, output_path)

    # Start dual-logging to session directory log.md
    session_log_path = os.path.join(session_dir, "log.md")
    if hasattr(sys.stdout, "add_file"):
        sys.stdout.add_file(session_log_path)  # type: ignore

    import logging

    _session_handler = logging.FileHandler(session_log_path, mode="a", encoding="utf-8")
    _session_handler.setLevel(logging.DEBUG)
    logging.root.addHandler(_session_handler)

    default_prompt = "Create a professional business presentation about AI transformation in enterprise companies"

    session_state = {
        # Core paths
        "template_path": args.template or "",
        "output_path": output_path,
        "output_dir": output_dir,
        # User inputs
        "user_prompt": args.prompt or default_prompt,
        "verbose": args.verbose,
        "stream": not args.no_stream,
        "no_images": args.no_images,
        "min_images": args.min_images,
        # visual_review and visual_passes (work with or without template)
        "visual_review": effective_visual_review,
        "footer_text": args.footer_text,
        "date_text": args.date_text,
        "show_slide_numbers": args.show_slide_numbers,
        # Chunked workflow settings
        "chunk_size": args.chunk_size,
        "max_retries": args.max_retries,
        "visual_passes": effective_visual_passes,
        "start_tier": args.start_tier,
        "template_visuals": args.template_visuals,
        # Inter-chunk delay range (seconds) — read by _inter_chunk_sleep()
        "inter_chunk_delay_min": args.inter_chunk_delay_min,
        "inter_chunk_delay_max": args.inter_chunk_delay_max,
        # Fields populated by steps
        "storyboard": None,
        "storyboard_dir": None,
        "total_slides": 0,
        "chunk_files": [],
        "chunk_slide_groups": [],
        "processed_chunks": {},
        "reviewed_chunks": {},
        "use_fallback_generator": False,
        # Transient rate-limit flag — set when a 429 is caught in generate_chunk_pptx.
        # Cleared after each inter-chunk delay. Does NOT permanently activate fallback.
        "rate_limit_hit": False,
        "brand_style_intent": None,
        # LLM provider name for swappable agents (agents are loaded lazily per step)
        "llm_provider": args.llm_provider,
        # Fields used by existing step helpers
        "generated_file": "",
        "slides_data": [],
        "generated_images": {},
        "src_slide_width": 0,
        "src_slide_height": 0,
        "assembly_knowledge": {},
        "quality_report": {},
    }

    workflow = build_chunked_workflow(session_state)

    print("=" * 60)
    print("Chunked PPTX Workflow")
    print("=" * 60)
    print("Provider:   %s" % args.llm_provider)
    print("Session:    %s" % session_name)
    print("Session dir: %s" % session_dir)
    print("Prompt:     %s" % (args.prompt or default_prompt)[:80])
    print("Output:     %s" % output_path)
    if args.template:
        print("Mode:       template-assisted generation")
        print("Template:   %s" % args.template)
        if effective_visual_review:
            print("Visual review: enabled (%d passes max)" % args.visual_passes)
        else:
            print("Visual review: disabled")
    else:
        print("Mode:       raw generation (no template)")
        if effective_visual_review:
            print(
                "Visual review: enabled, template-independent (%d passes max)"
                % args.visual_passes
            )
        else:
            print("Visual review: disabled")
    print("Chunk size: %d slides per API call" % args.chunk_size)
    print("Max retries per chunk: %d" % args.max_retries)
    print(
        "Start tier: %d (%s)"
        % (
            args.start_tier,
            {1: "Claude PPTX skill", 2: "LLM code generation", 3: "text-only"}[
                args.start_tier
            ],
        )
    )
    print("Images:     %s" % ("disabled" if args.no_images else "enabled"))
    if args.verbose:
        print("Verbose:    enabled")

    start_time = time.time()

    workflow.run()

    # Post-merge contrast enforcement (template-independent safety net)
    if os.path.isfile(output_path):
        enforce_final_contrast(output_path)  # noqa: F405
        # Fix 13A final safety net: run sanitize_presentation on the
        # merged output to catch any tiny text shapes that survived
        # template assembly or chunk merging. This is the last-resort
        # purge before the user sees the file.
        try:
            from pptx import Presentation as _Prs  # type: ignore # noqa: F811

            _final_prs = _Prs(output_path)
            sanitize_presentation(_final_prs)  # noqa: F405
            _final_prs.save(output_path)
            print("[POST-MERGE] Final sanitize_presentation pass completed.")
        except Exception as _e:
            print("[POST-MERGE] Final sanitize pass failed: %s" % _e)

    elapsed = time.time() - start_time
    print("\n" + "=" * 60)
    print("[TIMING] Total workflow: %.1fs" % elapsed)
    print("Output: %s" % output_path)
    print("=" * 60)

    # GLOBAL_TOKEN_TRACKER replaced by Langfuse/OpenInference telemetry.


if __name__ == "__main__":
    import glob
    import logging

    # Snapshot current python files in the directory for generic cleanup later
    _cwd_py_files_start = set(glob.glob("*.py"))

    # --- Output redirect: all output goes to OUTPUT.md instead of console ---
    # Captures print() (stdout), warnings/errors (stderr), and Python logging
    # (agno framework logs WARNING/ERROR/INFO via the logging module).
    _script_dir = os.path.dirname(os.path.abspath(__file__))
    _output_md_path = os.path.join(_script_dir, "OUTPUT.md")

    class _FileWriter:
        """Redirect a stream (stdout/stderr) to a file.

        Writes go exclusively to the file — nothing is printed to the console.
        The file is opened in append mode so both stdout and stderr can share it.
        This handles dual-logging to a primary script-level OUTPUT.md and an
        optional session-specific log.md.
        """

        def __init__(self, filepath: str):
            self._file = open(filepath, "a", encoding="utf-8")
            import typing

            self._file2: typing.Optional[typing.TextIO] = None

        def add_file(self, filepath: str) -> None:
            """Add a secondary file to stream all writes to."""
            self._file2 = open(filepath, "a", encoding="utf-8")

        def write(self, text: str) -> int:
            self._file.write(text)
            self._file.flush()
            if self._file2 is not None:
                getattr(self._file2, "write", lambda x: None)(text)  # type: ignore
                getattr(self._file2, "flush", lambda: None)()  # type: ignore
            return len(text)

        def flush(self) -> None:
            self._file.flush()
            if self._file2 is not None:
                getattr(self._file2, "flush", lambda: None)()  # type: ignore

        def fileno(self) -> int:
            return self._file.fileno()

        def close(self) -> None:
            self._file.close()
            if self._file2 is not None:
                getattr(self._file2, "close", lambda: None)()  # type: ignore

    # Clear the file at the start of each run
    with open(_output_md_path, "w", encoding="utf-8") as f:
        f.write("")

    _file_writer = _FileWriter(_output_md_path)

    # Redirect stdout and stderr
    sys.stdout = _file_writer  # type: ignore[assignment]
    sys.stderr = _file_writer  # type: ignore[assignment]

    # Redirect Python logging (captures agno WARNING/ERROR/INFO messages)
    _file_handler = logging.FileHandler(_output_md_path, mode="a", encoding="utf-8")
    _file_handler.setLevel(logging.DEBUG)
    logging.root.addHandler(_file_handler)
    # Remove any existing console (stream) handlers so logs don't try to write
    # to the original stderr which is now redirected
    for _h in list(logging.root.handlers):  # type: ignore
        if getattr(_h, "name", None) != _file_handler.name:
            try:
                logging.root.removeHandler(_h)  # type: ignore
            except Exception:
                pass

    try:
        main()
    finally:
        # Restore original streams on exit
        sys.stdout = sys.__stdout__
        sys.stderr = sys.__stderr__
        _file_writer.close()

        # Generic Cleanup: Identify and remove .py files created during this run
        import glob

        _is_verbose = any(arg in sys.argv for arg in ("--verbose", "-v"))
        _cwd_py_files_end = set(glob.glob("*.py"))
        _new_py_files = _cwd_py_files_end - _cwd_py_files_start

        if _new_py_files:
            if _is_verbose:
                print("\n[VERBOSE] Cleaning up dynamically created files:")
            for tmp_file in sorted(_new_py_files):
                try:
                    if _is_verbose:
                        print("  - Removing: %s" % tmp_file)
                    os.remove(tmp_file)
                except Exception as e:
                    if _is_verbose:
                        print("  [ERROR] Failed to remove %s: %s" % (tmp_file, e))
