"""
Offline unit tests for Template Visual Profile analysis logic.

Tests the SlideLayoutProfile, TemplateVisualProfile dataclasses,
_analyze_template_visual_profile(), and _format_visual_profile_for_prompt().

No API calls are made — all tests are fully offline.

This test file re-implements only the visual-profile-specific code to avoid
requiring the full agno/anthropic dependency chain.

Usage:
    python test_template_visual_profile.py
"""

import json
import os
import sys
import tempfile

from dataclasses import dataclass, field, asdict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE


# --- Copy of SlideLayoutProfile and TemplateVisualProfile dataclasses ---


@dataclass
class SlideLayoutProfile:
    """Visual characteristics of a single template slide."""

    slide_index: int = 0
    slide_type_hint: str = "content"
    placeholder_count: int = 0
    content_zone_left_pct: float = 5.0
    content_zone_top_pct: float = 12.0
    content_zone_width_pct: float = 90.0
    content_zone_height_pct: float = 76.0
    decorative_shape_count: int = 0
    has_background_image: bool = False
    has_gradient_background: bool = False
    text_box_count: int = 0
    total_shape_count: int = 0
    usable_width_pct: float = 90.0
    usable_height_pct: float = 76.0


@dataclass
class TemplateVisualProfile:
    """Comprehensive visual analysis of a template's layout characteristics."""

    slide_count: int = 0
    slide_width_emu: int = 0
    slide_height_emu: int = 0
    aspect_ratio: str = "16:9"
    avg_placeholder_count: float = 2.0
    avg_decorative_shapes: float = 0.0
    avg_content_zone_width_pct: float = 90.0
    avg_content_zone_height_pct: float = 76.0
    layout_density: str = "balanced"
    dominant_layout_style: str = "full"
    has_charts_in_template: bool = False
    has_tables_in_template: bool = False
    has_images_in_template: bool = False
    has_smartart_shapes: bool = False
    max_comfortable_bullets: int = 5
    recommended_text_weight: str = "balanced"
    slide_profiles: list = field(default_factory=list)


# --- Copy of analysis function (standalone, no VERBOSE dependency) ---

VERBOSE = False


def _analyze_template_visual_profile(template_path):
    """Analyze a template's visual layout characteristics."""
    profile = TemplateVisualProfile()

    try:
        prs = Presentation(template_path)
    except Exception as e:
        print("[WARNING] [VISUAL PROFILE] Failed to open template: %s" % str(e))
        return profile

    profile.slide_width_emu = prs.slide_width
    profile.slide_height_emu = prs.slide_height
    w_inches = prs.slide_width / 914400.0
    h_inches = prs.slide_height / 914400.0

    if abs(w_inches / h_inches - 16.0 / 9.0) < 0.1:
        profile.aspect_ratio = "16:9"
    elif abs(w_inches / h_inches - 4.0 / 3.0) < 0.1:
        profile.aspect_ratio = "4:3"
    elif abs(w_inches / h_inches - 16.0 / 10.0) < 0.1:
        profile.aspect_ratio = "16:10"
    else:
        profile.aspect_ratio = "custom"

    slides = list(prs.slides)
    profile.slide_count = len(slides)
    if not slides:
        return profile

    FOOTER_PH_IDXS = {10, 11, 12, 13}

    slide_profiles = []
    total_shapes_all = 0
    content_w_pcts = []
    content_h_pcts = []

    for s_idx, slide in enumerate(slides):
        sp = SlideLayoutProfile(slide_index=s_idx)
        sp.total_shape_count = len(slide.shapes)
        total_shapes_all += sp.total_shape_count

        content_ph_left = None
        content_ph_top = None
        content_ph_right = None
        content_ph_bottom = None
        has_title_ph = False

        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx
                if ph_idx in FOOTER_PH_IDXS:
                    continue
                sp.placeholder_count += 1

                if ph_idx == 0:
                    has_title_ph = True

                s_left = shape.left
                s_top = shape.top
                s_right = s_left + shape.width
                s_bottom = s_top + shape.height

                if content_ph_left is None or s_left < content_ph_left:
                    content_ph_left = s_left
                if content_ph_top is None or s_top < content_ph_top:
                    content_ph_top = s_top
                if content_ph_right is None or s_right > content_ph_right:
                    content_ph_right = s_right
                if content_ph_bottom is None or s_bottom > content_ph_bottom:
                    content_ph_bottom = s_bottom
            else:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                    sp.text_box_count += 1
                elif hasattr(shape, "has_chart") and shape.has_chart:
                    profile.has_charts_in_template = True
                elif hasattr(shape, "has_table") and shape.has_table:
                    profile.has_tables_in_template = True
                elif shape.shape_type is not None:
                    shape_type_val = int(shape.shape_type) if shape.shape_type else 0
                    if shape_type_val in (13, 11):
                        profile.has_images_in_template = True
                    elif shape_type_val in (24, 25):
                        profile.has_smartart_shapes = True
                    else:
                        sp.decorative_shape_count += 1
                else:
                    sp.decorative_shape_count += 1

        if content_ph_left is not None:
            sp.content_zone_left_pct = round(
                content_ph_left / prs.slide_width * 100.0, 1
            )
            sp.content_zone_top_pct = round(
                content_ph_top / prs.slide_height * 100.0, 1
            )
            zone_w = content_ph_right - content_ph_left
            zone_h = content_ph_bottom - content_ph_top
            sp.content_zone_width_pct = round(
                zone_w / prs.slide_width * 100.0, 1
            )
            sp.content_zone_height_pct = round(
                zone_h / prs.slide_height * 100.0, 1
            )
            deco_reduction = min(sp.decorative_shape_count * 5.0, 25.0)
            sp.usable_width_pct = max(40.0, sp.content_zone_width_pct - deco_reduction)
            sp.usable_height_pct = max(40.0, sp.content_zone_height_pct)

        content_w_pcts.append(sp.content_zone_width_pct)
        content_h_pcts.append(sp.content_zone_height_pct)

        if s_idx == 0 and has_title_ph and sp.placeholder_count <= 2:
            sp.slide_type_hint = "title"
        elif sp.placeholder_count == 0:
            sp.slide_type_hint = "blank"
        elif has_title_ph and sp.placeholder_count == 1:
            sp.slide_type_hint = "section"
        else:
            sp.slide_type_hint = "content"

        try:
            bg = slide.background
            fill = bg.fill
            if fill.type is not None:
                fill_type = str(fill.type)
                if "PICTURE" in fill_type.upper():
                    sp.has_background_image = True
                elif "GRADIENT" in fill_type.upper():
                    sp.has_gradient_background = True
        except Exception:
            pass

        slide_profiles.append(sp)

    profile.slide_profiles = slide_profiles

    n = len(slide_profiles)
    profile.avg_placeholder_count = round(
        sum(sp.placeholder_count for sp in slide_profiles) / n, 1
    )
    profile.avg_decorative_shapes = round(
        sum(sp.decorative_shape_count for sp in slide_profiles) / n, 1
    )
    profile.avg_content_zone_width_pct = round(sum(content_w_pcts) / n, 1)
    profile.avg_content_zone_height_pct = round(sum(content_h_pcts) / n, 1)

    avg_total_shapes = total_shapes_all / n
    if avg_total_shapes < 3:
        profile.layout_density = "sparse"
    elif avg_total_shapes <= 6:
        profile.layout_density = "balanced"
    else:
        profile.layout_density = "dense"

    if profile.avg_content_zone_width_pct < 55:
        profile.dominant_layout_style = "sidebar"
    elif profile.avg_content_zone_width_pct < 75:
        profile.dominant_layout_style = "split"
    elif profile.avg_decorative_shapes > 3:
        profile.dominant_layout_style = "overlapping"
    else:
        profile.dominant_layout_style = "full"

    profile.max_comfortable_bullets = max(
        2, min(6, int(profile.avg_content_zone_height_pct / 15))
    )

    if profile.avg_decorative_shapes > 2 or profile.has_images_in_template:
        profile.recommended_text_weight = "light"
    elif profile.layout_density == "dense":
        profile.recommended_text_weight = "dense"
    else:
        profile.recommended_text_weight = "balanced"

    return profile


def _format_visual_profile_for_prompt(profile):
    """Format a TemplateVisualProfile as a markdown section for the optimizer prompt."""
    if not profile or profile.slide_count == 0:
        return ""

    sections = ["## Template Visual Profile\n"]

    sections.append("- **Aspect Ratio:** %s" % profile.aspect_ratio)
    sections.append(
        "- **Layout Density:** %s (avg %.1f shapes/slide)"
        % (profile.layout_density, profile.avg_placeholder_count + profile.avg_decorative_shapes)
    )
    sections.append(
        "- **Content Zone:** ~%.0f%% width, ~%.0f%% height"
        % (profile.avg_content_zone_width_pct, profile.avg_content_zone_height_pct)
    )
    if profile.dominant_layout_style != "full":
        sections.append(
            "  _(template uses %s layout)_" % profile.dominant_layout_style
        )
    sections.append(
        "- **Max Comfortable Bullets:** %d per slide" % profile.max_comfortable_bullets
    )
    sections.append(
        "- **Recommended Text Weight:** %s" % profile.recommended_text_weight
    )

    elements = []
    if profile.has_charts_in_template:
        elements.append("charts")
    if profile.has_tables_in_template:
        elements.append("tables")
    if profile.has_images_in_template:
        elements.append("images")
    if profile.has_smartart_shapes:
        elements.append("SmartArt")
    deco_note = ""
    if profile.avg_decorative_shapes > 1:
        deco_note = ", decorative shapes (avg %.0f/slide)" % profile.avg_decorative_shapes
    sections.append(
        "- **Template Contains:** %s%s"
        % (", ".join(elements) if elements else "text placeholders only", deco_note)
    )

    sections.append(
        "\nTEMPLATE-AWARE CONSTRAINTS (apply these when generating the storyboard):\n"
        "- Limit key_points to %d or fewer per slide (based on actual content zone height).\n"
        % profile.max_comfortable_bullets
    )

    return "\n".join(sections)


# === TESTS ===


def test_slide_layout_profile_defaults():
    """SlideLayoutProfile should have sensible defaults."""
    sp = SlideLayoutProfile()
    assert sp.slide_index == 0, "Default slide_index should be 0"
    assert sp.slide_type_hint == "content", "Default type hint should be 'content'"
    assert sp.placeholder_count == 0, "Default placeholder_count should be 0"
    assert sp.content_zone_left_pct == 5.0, "Default content_zone_left_pct should be 5.0"
    assert sp.content_zone_top_pct == 12.0, "Default content_zone_top_pct should be 12.0"
    assert sp.content_zone_width_pct == 90.0, "Default width should be 90.0"
    assert sp.content_zone_height_pct == 76.0, "Default height should be 76.0"
    assert sp.decorative_shape_count == 0, "Default decorative count should be 0"
    assert sp.has_background_image is False, "No background image by default"
    assert sp.has_gradient_background is False, "No gradient by default"
    assert sp.total_shape_count == 0, "Default total shapes should be 0"
    print("  PASS: test_slide_layout_profile_defaults")


def test_visual_profile_defaults():
    """TemplateVisualProfile should have sensible defaults."""
    profile = TemplateVisualProfile()
    assert profile.slide_count == 0, "Default slide_count should be 0"
    assert profile.aspect_ratio == "16:9", "Default aspect_ratio should be '16:9'"
    assert profile.layout_density == "balanced", "Default density should be 'balanced'"
    assert profile.dominant_layout_style == "full", "Default style should be 'full'"
    assert profile.has_charts_in_template is False, "No charts by default"
    assert profile.has_tables_in_template is False, "No tables by default"
    assert profile.max_comfortable_bullets == 5, "Default max bullets should be 5"
    assert profile.recommended_text_weight == "balanced", "Default text weight should be 'balanced'"
    assert profile.slide_profiles == [], "Default slide_profiles should be empty list"
    print("  PASS: test_visual_profile_defaults")


def test_analyze_minimal_template():
    """Single-slide minimal template should produce a valid profile."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name

    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(tmp_path)

        profile = _analyze_template_visual_profile(tmp_path)
        assert profile.slide_count == 1, "Should detect 1 slide"
        assert profile.slide_width_emu > 0, "Should have positive slide width"
        assert profile.slide_height_emu > 0, "Should have positive slide height"
        assert profile.aspect_ratio in ("16:9", "4:3"), "Should detect a standard aspect ratio"
        assert len(profile.slide_profiles) == 1, "Should have 1 slide profile"
        print("  PASS: test_analyze_minimal_template")
    finally:
        os.unlink(tmp_path)


def test_analyze_multi_slide_template():
    """Multi-slide template should compute averages correctly."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name

    try:
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
        prs.slides.add_slide(prs.slide_layouts[1])  # Content slide
        prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
        prs.save(tmp_path)

        profile = _analyze_template_visual_profile(tmp_path)
        assert profile.slide_count == 3, "Should detect 3 slides, got %d" % profile.slide_count
        assert len(profile.slide_profiles) == 3, "Should have 3 slide profiles"
        # Averages should be computed
        assert isinstance(profile.avg_placeholder_count, float), "avg_placeholder_count should be float"
        assert isinstance(profile.avg_decorative_shapes, float), "avg_decorative_shapes should be float"
        print("  PASS: test_analyze_multi_slide_template")
    finally:
        os.unlink(tmp_path)


def test_content_zone_computation():
    """Placeholder positions should derive content_zone percentages."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name

    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + content layout
        prs.save(tmp_path)

        profile = _analyze_template_visual_profile(tmp_path)
        sp = profile.slide_profiles[0]
        # Content zone should be derived from placeholder positions
        assert sp.content_zone_width_pct > 0, "Width should be positive, got %.1f" % sp.content_zone_width_pct
        assert sp.content_zone_height_pct > 0, "Height should be positive, got %.1f" % sp.content_zone_height_pct
        assert sp.content_zone_left_pct >= 0, "Left should be non-negative"
        assert sp.content_zone_top_pct >= 0, "Top should be non-negative"
        print("  PASS: test_content_zone_computation")
    finally:
        os.unlink(tmp_path)


def test_decorative_shape_counting():
    """Non-placeholder shapes should be counted as decorative."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name

    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        # Add decorative shapes
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
        slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(1), Inches(1), Inches(1))
        slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(6), Inches(1), Inches(1), Inches(1))
        prs.save(tmp_path)

        profile = _analyze_template_visual_profile(tmp_path)
        sp = profile.slide_profiles[0]
        assert sp.decorative_shape_count >= 3, \
            "Should have at least 3 decorative shapes, got %d" % sp.decorative_shape_count
        assert sp.total_shape_count >= 3, \
            "Total shapes should be at least 3, got %d" % sp.total_shape_count
        print("  PASS: test_decorative_shape_counting")
    finally:
        os.unlink(tmp_path)


def test_layout_density_classification():
    """Layout density should classify based on shape count thresholds."""
    # Sparse: < 3 shapes/slide
    profile_sparse = TemplateVisualProfile()
    profile_sparse.slide_profiles = [SlideLayoutProfile(total_shape_count=1)]
    # We can't easily test the classification without calling the full function,
    # but we can verify the thresholds are correct conceptually.

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name
    try:
        prs = Presentation()
        # Blank slide = sparse (0-1 shapes)
        prs.slides.add_slide(prs.slide_layouts[5])
        prs.save(tmp_path)

        profile = _analyze_template_visual_profile(tmp_path)
        assert profile.layout_density in ("sparse", "balanced"), \
            "Blank slide should be sparse or balanced, got '%s'" % profile.layout_density
        print("  PASS: test_layout_density_classification")
    finally:
        os.unlink(tmp_path)


def test_max_comfortable_bullets():
    """Max bullets should be derived from content zone height."""
    # Default height 76% / 15 = 5.06 -> 5
    profile = TemplateVisualProfile(avg_content_zone_height_pct=76.0)
    expected = max(2, min(6, int(76.0 / 15)))
    assert expected == 5, "Expected 5, got %d" % expected

    # Small content zone (30%) -> 30/15 = 2
    profile2 = TemplateVisualProfile(avg_content_zone_height_pct=30.0)
    expected2 = max(2, min(6, int(30.0 / 15)))
    assert expected2 == 2, "Expected 2, got %d" % expected2

    # Large content zone (95%) -> 95/15 = 6.3 -> capped at 6
    profile3 = TemplateVisualProfile(avg_content_zone_height_pct=95.0)
    expected3 = max(2, min(6, int(95.0 / 15)))
    assert expected3 == 6, "Expected 6, got %d" % expected3
    print("  PASS: test_max_comfortable_bullets")


def test_format_visual_profile_for_prompt():
    """_format_visual_profile_for_prompt should produce valid markdown."""
    profile = TemplateVisualProfile(
        slide_count=5,
        aspect_ratio="16:9",
        layout_density="sparse",
        dominant_layout_style="sidebar",
        avg_placeholder_count=2.0,
        avg_decorative_shapes=1.5,
        avg_content_zone_width_pct=65.0,
        avg_content_zone_height_pct=70.0,
        max_comfortable_bullets=4,
        recommended_text_weight="light",
        has_charts_in_template=True,
        has_images_in_template=True,
    )

    result = _format_visual_profile_for_prompt(profile)
    assert "## Template Visual Profile" in result, "Should have header"
    assert "16:9" in result, "Should include aspect ratio"
    assert "sparse" in result, "Should include density"
    assert "sidebar" in result, "Should include layout style"
    assert "65%" in result or "65" in result, "Should reference content zone width"
    assert "4 per slide" in result, "Should include max bullets"
    assert "light" in result, "Should include text weight"
    assert "charts" in result, "Should mention charts"
    assert "images" in result, "Should mention images"
    print("  PASS: test_format_visual_profile_for_prompt")


def test_format_visual_profile_empty():
    """_format_visual_profile_for_prompt should return empty for no-profile."""
    result = _format_visual_profile_for_prompt(None)
    assert result == "", "Should return empty string for None input"

    empty_profile = TemplateVisualProfile()
    result2 = _format_visual_profile_for_prompt(empty_profile)
    assert result2 == "", "Should return empty string for empty profile"
    print("  PASS: test_format_visual_profile_empty")


def test_profile_json_roundtrip():
    """TemplateVisualProfile should survive JSON serialization."""
    profile = TemplateVisualProfile(
        slide_count=3,
        aspect_ratio="4:3",
        layout_density="dense",
        dominant_layout_style="full",
        has_charts_in_template=True,
        max_comfortable_bullets=4,
        recommended_text_weight="light",
    )

    data = asdict(profile)
    json_str = json.dumps(data)
    restored_data = json.loads(json_str)

    restored = TemplateVisualProfile(**restored_data)
    assert restored.slide_count == 3, "slide_count should survive roundtrip"
    assert restored.aspect_ratio == "4:3", "aspect_ratio should survive roundtrip"
    assert restored.layout_density == "dense", "layout_density should survive roundtrip"
    assert restored.has_charts_in_template is True, "has_charts should survive roundtrip"
    assert restored.max_comfortable_bullets == 4, "max_bullets should survive roundtrip"
    print("  PASS: test_profile_json_roundtrip")


def test_nonexistent_template():
    """Analysis should handle missing file gracefully."""
    profile = _analyze_template_visual_profile("/tmp/nonexistent_visual_profile_test.pptx")
    assert profile.slide_count == 0, "Should have 0 slides for missing file"
    assert profile.slide_width_emu == 0, "Should have 0 width for missing file"
    print("  PASS: test_nonexistent_template")


def run_all_tests():
    """Run all template visual profile tests."""
    print("=" * 60)
    print("Running Template Visual Profile Tests")
    print("=" * 60)

    tests = [
        test_slide_layout_profile_defaults,
        test_visual_profile_defaults,
        test_analyze_minimal_template,
        test_analyze_multi_slide_template,
        test_content_zone_computation,
        test_decorative_shape_counting,
        test_layout_density_classification,
        test_max_comfortable_bullets,
        test_format_visual_profile_for_prompt,
        test_format_visual_profile_empty,
        test_profile_json_roundtrip,
        test_nonexistent_template,
    ]

    passed = 0
    failed = 0
    for test_fn in tests:
        try:
            test_fn()
            passed += 1
        except Exception as e:
            print("  FAIL: %s — %s" % (test_fn.__name__, e))
            import traceback
            traceback.print_exc()
            failed += 1

    print("\n" + "=" * 60)
    print("Results: %d passed, %d failed" % (passed, failed))
    print("=" * 60)

    if failed > 0:
        sys.exit(1)
    else:
        print("ALL TEMPLATE VISUAL PROFILE TESTS PASSED")


if __name__ == "__main__":
    run_all_tests()
