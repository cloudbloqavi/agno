#!/usr/bin/env python3
"""Deep diagnostic of the existing PPTX — inspect actual slide backgrounds and all text runs."""

import os
import sys
from lxml import etree

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, SCRIPT_DIR)

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Emu

from powerpoint_template_workflow import (
    _get_shape_background_color,
    _hex_to_rgb,
    _relative_luminance,
    _contrast_ratio,
    _extract_color_from_solid_fill,
)


def deep_diagnose(pptx_path: str):
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    ns_p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    prs = PptxPresentation(pptx_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n{'='*70}")
        print(f"  SLIDE {slide_idx + 1}")
        print(f"{'='*70}")
        
        # Check slide background XML directly
        slide_elem = slide._element
        bg = slide_elem.find(ns_p + "bg")
        if bg is not None:
            print(f"  Slide background XML found:")
            print(f"    {etree.tostring(bg, pretty_print=True).decode()[:500]}")
        else:
            print(f"  No <p:bg> element in slide XML")
        
        # Check slide layout background
        layout = slide.slide_layout
        layout_bg = layout._element.find(ns_p + "bg")
        if layout_bg is not None:
            print(f"  Layout background XML found:")
            print(f"    {etree.tostring(layout_bg, pretty_print=True).decode()[:500]}")
        else:
            print(f"  No layout background")
        
        # Check slide master background  
        master = layout.slide_master
        master_bg = master._element.find(ns_p + "bg")
        if master_bg is not None:
            print(f"  Master background XML found:")
            print(f"    {etree.tostring(master_bg, pretty_print=True).decode()[:300]}")
        
        for shape in slide.shapes:
            bg_hex = _get_shape_background_color(shape, slide)
            bg_rgb = _hex_to_rgb(bg_hex)
            bg_lum = _relative_luminance(*bg_rgb)
            
            print(f"\n  Shape: '{getattr(shape, 'name', '?')}' type={type(shape).__name__}")
            print(f"    Position: {shape.left}, {shape.top}  Size: {shape.width}x{shape.height}")
            print(f"    Detected bg: #{bg_hex} (lum={bg_lum:.3f})")
            
            if getattr(shape, "has_chart", False):
                print(f"    HAS CHART — w={shape.width/Inches(1):.1f}\" h={shape.height/Inches(1):.1f}\"")
                try:
                    chart_elem = shape.chart._chartSpace
                    defRPrs = list(chart_elem.iter(ns_a + "defRPr"))
                    rPrs = list(chart_elem.iter(ns_a + "rPr"))
                    print(f"    Chart text elements: {len(defRPrs)} defRPr, {len(rPrs)} rPr")
                    for i, rPr in enumerate(defRPrs[:3]):
                        fill = rPr.find(ns_a + "solidFill")
                        if fill is not None:
                            color = _extract_color_from_solid_fill(fill)
                            print(f"      defRPr[{i}] color: #{color}")
                        else:
                            print(f"      defRPr[{i}] color: none (inherited)")
                except Exception as e:
                    print(f"    Chart analysis error: {e}")
            
            if getattr(shape, "has_text_frame", False):
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(para.runs):
                        if not run.text.strip():
                            continue
                        rPr = run._r.find(ns_a + "rPr")
                        text_hex = None
                        color_src = "no-rPr"
                        
                        if rPr is not None:
                            solidFill = rPr.find(ns_a + "solidFill")
                            if solidFill is not None:
                                text_hex = _extract_color_from_solid_fill(solidFill)
                                color_src = f"explicit:#{text_hex}"
                            else:
                                # Check for schemeClr directly on rPr
                                scheme = rPr.find(ns_a + "solidFill/" + ns_a + "schemeClr")
                                color_src = "rPr-no-fill"
                        
                        if text_hex:
                            ratio = _contrast_ratio(_hex_to_rgb(text_hex), bg_rgb)
                        else:
                            ratio = _contrast_ratio((0,0,0), bg_rgb)
                        
                        status = "✅" if ratio >= 3.0 else "❌"
                        print(f"    {status} Run[{para_idx}.{run_idx}]: color={color_src} ratio={ratio:.1f} text='{run.text[:50]}'")


if __name__ == "__main__":
    source = os.path.join(
        SCRIPT_DIR,
        "output_chunked/chunked_workflow_work/session_e3fa4b1d_20260305_051723/"
        "software_deck_template.pptx"
    )
    
    if not os.path.exists(source):
        print(f"ERROR: {source}")
        sys.exit(1)
    
    deep_diagnose(source)
