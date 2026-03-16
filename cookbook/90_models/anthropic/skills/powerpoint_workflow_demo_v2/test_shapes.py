import sys
from pptx import Presentation
prs = Presentation(sys.argv[1])
for i, slide in enumerate(prs.slides):
    print(f"Slide {i+1}")
    for j, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            print(f"  Shape {j}: {shape.name} x={shape.left} y={shape.top} w={shape.width} h={shape.height} type={shape.shape_type}")
            print(f"    Text: {shape.text[:50]}")
