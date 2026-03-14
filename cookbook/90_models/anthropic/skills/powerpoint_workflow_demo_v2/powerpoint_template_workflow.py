"""
Agno Workflow: PowerPoint Template Generation Pipeline.

A sequential workflow that generates presentations using Claude's pptx skill,
intelligently adds AI-generated images via NanoBanana, and
applies a custom .pptx template for professional styling. Auxiliary agents
like image planning and visual review are swappable (Claude, OpenAI, Gemini).

Workflow steps:
  Step 1  Content Generation  - Claude + pptx skill -> raw .pptx
  Step 2  Image Planning      - Swappable Agent decides which slides need images
  Step 3  Image Generation    - NanoBanana generates slide images
  Step 4  Template Assembly   - The most critical step. Before constructing any slide,
                                consolidates all context into a comprehensive knowledge
                                file from four mandatory inputs: (1) original user prompt,
                                (2) complete slide-by-slide content plan from Step 1,
                                (3) deep per-layout analysis of the template's full design
                                language (fonts, colors, placeholder positions, decorative
                                shapes, accent palette, motifs), and (4) all AI-generated
                                image assets with dimensions and target layouts. Only then
                                does the actual PPTX construction begin, governed entirely
                                by this knowledge file as the single source of truth.
  Step 5  Visual Quality Review (optional) - Swappable Vision Agent inspects rendered slides

Operating modes:
  With template (--template / -t):
    Runs all applicable steps (1-4, optionally 5). The .pptx template provides
    the visual design — fonts, colors, layouts, placeholders, decorative shapes.
    The final output is a professionally styled presentation that matches your
    corporate or personal brand.

  Without template (no --template flag):
    Runs only Steps 1, 2 (unless --no-images), and 3 (unless --no-images).
    Step 4 (Template Assembly) is skipped entirely. The raw Claude-generated
    .pptx file is copied directly to the output path. Useful for quickly
    previewing AI-generated content before applying a template, or when no
    template is available.

Prerequisites:
- uv pip install agno anthropic openai google-genai python-pptx pillow lxml python-dotenv
- export ANTHROPIC_API_KEY="your_api_key_here" (always required)
- export OPENAI_API_KEY="your_openai_api_key_here" (for --llm-provider openai)
- export GOOGLE_API_KEY="your_google_api_key_here" (for --llm-provider gemini)
- A .pptx template file (optional — omit to get raw Claude output)
- LibreOffice (required for --visual-review step: `sudo apt-get install -y libreoffice`)
- poppler-utils (required for per-slide PNG rendering: `sudo apt-get install -y poppler-utils`)

Template Quality Safeguards:
  When a template is provided, these safeguards prevent common styling issues:
  - Per-slide rendering: PPTX→PDF→PNG pipeline (via pdftoppm) renders every slide for visual review and prompt context
  - Background detection: 6-layer detection (shape→slide→layout→master→theme→large shapes)
  - Layout sanitization: 3-pass boundary clamping, min size enforcement, and shape overlap reflow
  - Template-aware LLM prompts: Tier 2 prompt includes template constraints (bg color, text color, max shapes)
  - Template visual references: Injects rendered template PNGs into chunk prompts for spatial layout context

Usage:
    # Basic usage with a template:
    /workspaces/agno/.venvs/demo/bin/python powerpoint_template_workflow.py \\
        --template my_template.pptx

    # Run without a template (raw Claude output only):
    /workspaces/agno/.venvs/demo/bin/python powerpoint_template_workflow.py

    # Run without a template, with a custom prompt and output file:
    /workspaces/agno/.venvs/demo/bin/python powerpoint_template_workflow.py \\
        -p "Create a 5-slide overview of renewable energy trends" -o energy.pptx

    # Full options: custom prompt, output path, verbose logging:
    /workspaces/agno/.venvs/demo/bin/python powerpoint_template_workflow.py \\
        -t my_template.pptx -o report.pptx -p "Create a 5-slide AI trends presentation" -v

    # Disable streaming (more reliable for shorter prompts, may timeout on complex ones):
    /workspaces/agno/.venvs/demo/bin/python powerpoint_template_workflow.py \\
        -t my_template.pptx --no-stream

    # Skip AI image generation entirely:
    /workspaces/agno/.venvs/demo/bin/python powerpoint_template_workflow.py \\
        -t my_template.pptx --no-images

    # Require at least 3 AI-generated images across slides:
    /workspaces/agno/.venvs/demo/bin/python powerpoint_template_workflow.py \\
        -t my_template.pptx --min-images 3

    # Let the image planner decide freely (no minimum enforced):
    /workspaces/agno/.venvs/demo/bin/python powerpoint_template_workflow.py \\
        -t my_template.pptx --min-images 0

    # Enable visual quality review with Gemini vision (requires LibreOffice):
    /workspaces/agno/.venvs/demo/bin/python powerpoint_template_workflow.py \\
        -t my_template.pptx --visual-review

    # Add footer text and slide numbers to all slides:
    /workspaces/agno/.venvs/demo/bin/python powerpoint_template_workflow.py \\
        -t my_template.pptx --footer-text "Confidential" --show-slide-numbers

CLI Flags:
    --template, -t       Path to the .pptx template file (optional).
                         With template: runs the full pipeline including Step 4
                         (Template Assembly), which applies the template's fonts,
                         colors, and layouts to produce a professionally branded .pptx.
                         Without template: skips Step 4; the raw Claude-generated
                         presentation is saved directly to the output path. Steps 2
                         and 3 (image generation) still run unless --no-images is set.
    --output, -o         Output filename (default: presentation_from_template.pptx).
    --prompt, -p         Custom prompt for the presentation content.
    --no-images          Skip AI image generation (Steps 2 and 3).
    --no-stream          Disable streaming mode for Claude agent (more reliable for
                         shorter prompts, but may timeout on complex presentations).
    --min-images         Minimum number of slides that must have AI-generated images
                         (default: 1). Use 0 to let the image planner decide freely.
    --visual-review      Enable Step 5 visual QA with Gemini 2.5 Flash vision.
                         Renders each slide to PNG via LibreOffice headless, inspects
                         for defects (overlap, text overflow, ghost text, low contrast,
                         visual blandness), and applies safe corrections for critical
                         issues. Non-blocking: skips if LibreOffice is unavailable.
    --footer-text        Footer text for all slides (template idx=11 placeholder).
    --date-text          Date/time text for footer date placeholder (idx=10).
    --show-slide-numbers Preserve slide number placeholder (idx=12) on all slides.
    --verbose, -v        Enable verbose/debug logging for troubleshooting.
"""

import argparse
import copy
import json
import os
import shutil
import sys
import time
import traceback
from dataclasses import dataclass, field
from enum import Enum
from io import BytesIO
from typing import Dict, List

from agno.agent import Agent
from agno.media import Image as AgnoImage
from lib_patches.anthropic.claude import Claude
from agno.models.google import Gemini
from agno.run.agent import RunOutput
from agno.tools.nano_banana import NanoBananaTools
from agno.workflow.step import Step
from agno.workflow.types import StepInput, StepOutput
from agno.workflow.workflow import Workflow
from anthropic import Anthropic
from file_download_helper import download_skill_files
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field, model_validator
from typing import Optional

# Module-level verbose flag (set from CLI args)
VERBOSE = False

# Minimum body font size in points — fit_text() and MSO_AUTO_SIZE can shrink text
# to as little as 4pt which is unreadable. This floor ensures legibility.
# When text overflows at this size, content is truncated with '…' instead.
_MIN_BODY_FONT_PT = 10
_MIN_TITLE_FONT_PT = 14


# ---------------------------------------------------------------------------
# Pydantic Models for Structured Data Flow
# ---------------------------------------------------------------------------


class SlideImageDecision(BaseModel):
    """Decision about whether a slide needs an AI-generated image."""

    slide_index: int = Field(description="Zero-based index of the slide")
    needs_image: bool = Field(
        description="Whether this slide would benefit from an image"
    )
    image_prompt: str = Field(
        default="",
        description="If needs_image is True, a detailed prompt for generating the image. "
        "Should describe a professional, clean illustration suitable for a business presentation. "
        "If needs_image is False, leave empty.",
    )
    reasoning: str = Field(
        default="",
        description="Brief explanation of why the slide does or does not need an image",
    )


class ImagePlan(BaseModel):
    """Plan for which slides need AI-generated images."""

    decisions: List[SlideImageDecision] = Field(
        description="List of image decisions, one per slide"
    )


# ---------------------------------------------------------------------------
# Pydantic Models for Visual Quality Review (Step 5)
# ---------------------------------------------------------------------------


class ShapeIssue(BaseModel):
    """A single design issue detected on a rendered slide image.

    Covers both structural defects (text overflow, overlapping elements) and
    design-quality improvements (typography hierarchy, color underutilization,
    spacing improvements, visual enrichment opportunities).
    """

    issue_type: str = Field(
        description=(
            "Category of issue. Structural: text_overflow, text_too_small, overlap, "
            "ghost_text, low_contrast, element_clipped, empty_placeholder, footer_inconsistent. "
            "Design quality: poor_spacing, alignment_off, typography_hierarchy, "
            "color_underutilized, visual_enrichment_needed, font_inconsistency."
        )
    )
    severity: str = Field(
        description=(
            "Impact level. One of: critical (broken / unreadable), "
            "moderate (clearly suboptimal), minor (subtle improvement opportunity)."
        )
    )
    description: str = Field(
        description="Specific, actionable description of the issue"
    )
    programmatic_fix: str = Field(
        description=(
            "Action to apply. "
            "Structural: reduce_font_size, increase_contrast, remove_element, clear_placeholder. "
            "Spacing/alignment: fix_spacing, fix_alignment, fix_body_paragraph_alignment. "
            "Typography: increase_title_font_size, enforce_typography_hierarchy. "
            "Color: apply_accent_color_title, apply_accent_color_body. "
            "Visual enrichment: apply_body_accent_border, "
            "enrich_header_bar, enrich_title_card, enrich_divider, enrich_accent_strip. "
            "Use 'none' only for issues requiring human judgment or AI-generated content."
        )
    )
    shape_description: str = Field(
        default="",
        description=(
            "Brief description of the affected element (e.g., 'title text box', "
            "'bullet list', 'slide background'). Leave empty if slide-wide."
        ),
    )


class SlideQualityReport(BaseModel):
    """Design and quality assessment for a single rendered slide image.

    Produced by a senior UI/UX designer persona who evaluates both structural
    correctness AND visual design quality (typography hierarchy, whitespace balance,
    color palette utilization, visual interest, and brand consistency).

    Robust schema: accepts 'design_score' as an alias for 'overall_quality' to handle
    vision model responses that may use either field name. Contains graceful defaults
    for all required fields so partial responses never crash the pipeline.
    """

    slide_index: int = Field(default=0, description="Zero-based index of the slide")
    overall_quality: str = Field(
        default="acceptable",
        description=(
            "Overall slide design quality from a professional designer perspective. "
            "One of: good (well-designed, minimal changes), "
            "acceptable (functional but improvable), poor (broken or significantly underdesigned)."
        ),
    )
    design_score: int = Field(
        default=5,
        description=(
            "Designer score 1-10. 9-10: excellent layout and visual hierarchy. "
            "7-8: good, minor improvements possible. 5-6: functional but generic. "
            "3-4: noticeably bland or has design problems. 1-2: broken or unusable."
        ),
    )
    is_visually_bland: bool = Field(
        default=False,
        description=(
            "True if the slide fails to use the template's visual vocabulary — "
            "e.g. text-only with no accent colors applied, no visual hierarchy, "
            "excessive unused whitespace, or content that looks like a bare draft."
        ),
    )
    blandness_reason: str = Field(
        default="",
        description=(
            "Specific description of why the slide is bland and what visual elements "
            "from the template palette could improve it. Leave empty if not bland."
        ),
    )
    issues: List[ShapeIssue] = Field(
        default_factory=list,
        description=(
            "List of detected issues ordered by severity. Include both structural "
            "defects and design quality improvements. Be specific and actionable."
        ),
    )

    @model_validator(mode="before")
    @classmethod
    def _reconcile_quality_fields(cls, data):
        """Accept design_score as alias for overall_quality and vice-versa.

        Vision models may return either field name. This validator normalises
        the data so both fields are always populated, preventing Pydantic
        validation errors that caused the visual review step to silently do nothing.
        """
        if not isinstance(data, dict):
            return data
        # If overall_quality is missing but design_score is present, synthesise it
        if not data.get("overall_quality") and "design_score" in data:
            score = data["design_score"]
            try:
                score = int(score)
            except (TypeError, ValueError):
                score = 5
            if score >= 8:
                data["overall_quality"] = "good"
            elif score >= 5:
                data["overall_quality"] = "acceptable"
            else:
                data["overall_quality"] = "poor"
        # If design_score is missing but overall_quality is present, synthesise it
        if not data.get("design_score") and data.get("overall_quality"):
            quality_map = {"good": 8, "acceptable": 6, "poor": 3}
            data["design_score"] = quality_map.get(
                str(data["overall_quality"]).lower(), 5
            )
        return data


class PresentationQualityReport(BaseModel):
    """Design and quality assessment for the entire presentation."""

    slide_reports: List[SlideQualityReport] = Field(
        description="Per-slide quality reports, one per rendered slide"
    )
    overall_pass: bool = Field(
        description="True if no slides are rated 'poor' AND average design_score >= 6"
    )
    total_critical_issues: int = Field(
        description="Total count of critical-severity issues across all slides"
    )
    average_design_score: float = Field(
        default=0.0,
        description="Mean design_score across all slides (1.0-10.0)",
    )
    recommendations: List[str] = Field(
        default_factory=list,
        description=(
            "Specific actionable recommendations for human follow-up — things that "
            "cannot be auto-corrected (e.g. 'Slide 3 would benefit from an AI-generated "
            "image of a rocket launch to match the SpaceTech theme')."
        ),
    )


# ---------------------------------------------------------------------------
# Dataclasses for Slide Content (reused from agent_with_powerpoint_template)
# ---------------------------------------------------------------------------


@dataclass
class TableData:
    """Extracted table data with position."""

    rows: list
    left: int
    top: int
    width: int
    height: int


@dataclass
class ImageData:
    """Extracted image data with position."""

    blob: bytes
    left: int
    top: int
    width: int
    height: int
    content_type: str = "image/png"


@dataclass
class ChartExtract:
    """Extracted chart data with position."""

    chart_type: int
    categories: list
    series: list
    left: int
    top: int
    width: int
    height: int


@dataclass
class SlideContent:
    """All extracted content from a single slide."""

    title: str = ""
    subtitle: str = ""
    body_paragraphs: list = field(default_factory=list)
    tables: list = field(default_factory=list)
    images: list = field(default_factory=list)
    charts: list = field(default_factory=list)
    shapes_xml: list = field(default_factory=list)
    text_shapes_xml: list = field(default_factory=list)
    text_box_paragraphs: list = field(default_factory=list)
    # Track picture placeholders detected in the slide layout
    has_image_placeholder: bool = False
    image_placeholder_indices: list = field(default_factory=list)
    # Tracks whether the slide contains any non-placeholder text boxes
    has_text_box: bool = False


@dataclass
class ContentArea:
    """Defines the safe content region on a template slide (all values in EMU)."""

    left: int
    top: int
    width: int
    height: int


class ContentMix(Enum):
    """Classification of what element types a slide contains."""

    TEXT_ONLY = "text_only"
    TEXT_AND_TABLE = "text_and_table"
    TEXT_AND_CHART = "text_and_chart"
    TEXT_AND_IMAGE = "text_and_image"
    TEXT_AND_GENERATED_IMAGE = "text_and_generated_image"
    VISUAL_ONLY = "visual_only"
    MIXED = "mixed"  # Multiple visual types


@dataclass
class RegionMap:
    """Defines separate regions for text and visual elements on a slide."""

    text_region: ContentArea
    visual_region: ContentArea
    layout_type: str = "full"  # "full", "split_horizontal", "split_vertical", "native"


@dataclass
class TemplateTheme:
    """Extracted theme colors and font scheme from a template."""

    accent_colors: list = field(default_factory=list)  # Up to 6 accent hex colors
    dk1: str = ""  # Dark 1 color (hex)
    dk2: str = ""  # Dark 2 color (hex)
    lt1: str = ""  # Light 1 color (hex)
    lt2: str = ""  # Light 2 color (hex)
    hyperlink: str = ""
    followed_hyperlink: str = ""
    major_font: str = "Calibri"  # Heading font
    minor_font: str = "Calibri"  # Body font


@dataclass
class TemplateTableStyle:
    """Styling extracted from a reference template table."""

    header_fill_rgb: str = ""  # hex like "4472C4"
    header_font_color_rgb: str = ""
    header_font_size_pt: int = 11
    header_font_bold: bool = True
    header_font_family: str = ""
    cell_font_size_pt: int = 10
    cell_font_color_rgb: str = ""
    cell_font_family: str = ""
    band_row_color_rgb: str = ""  # Alternating row fill
    border_color_rgb: str = ""
    border_width_emu: int = 12700  # 1pt default
    cell_fill_rgb: str = ""  # Default cell fill
    # Raw XML for reliable deep-copy transfer
    table_properties_xml: object = None  # lxml Element or None
    table_style_id: str = ""


@dataclass
class TemplateChartStyle:
    """Styling extracted from a reference template chart."""

    chart_type: int = -1
    series_fill_colors: list = field(default_factory=list)  # list of hex RGB
    series_line_colors: list = field(default_factory=list)
    axis_font_size_pt: int = 10
    axis_font_color_rgb: str = ""
    axis_font_family: str = ""
    axis_line_color_rgb: str = ""
    legend_font_size_pt: int = 10
    legend_font_color_rgb: str = ""
    legend_font_family: str = ""
    data_label_font_size_pt: int = 9
    data_label_font_color_rgb: str = ""
    plot_area_fill_rgb: str = ""
    chart_bg_fill_rgb: str = ""
    # Store the raw chartStyle XML element for full fidelity transfer
    chart_style_xml: object = None


@dataclass
class TemplateStyle:
    """Complete style information extracted from a template presentation."""

    theme: TemplateTheme = field(default_factory=TemplateTheme)
    table_styles: list = field(default_factory=list)
    chart_styles: list = field(default_factory=list)
    body_font_size_pt: int = 18
    body_font_color_rgb: str = ""
    body_font_family: str = ""
    title_font_size_pt: int = 28
    title_font_color_rgb: str = ""
    title_font_family: str = ""


# ---------------------------------------------------------------------------
# Template Style Extraction Functions
# ---------------------------------------------------------------------------


def _extract_theme_from_prs(prs) -> TemplateTheme:
    """Extract theme colors and font scheme from a presentation."""
    theme = TemplateTheme()
    try:
        # Access the theme through the slide master's relationships
        slide_master = prs.slide_masters[0]
        theme_part = None
        for rel in slide_master.part.rels.values():
            if "theme" in rel.reltype:
                theme_part = rel.target_part
                break
        if theme_part is None:
            return theme

        # theme_part.element may not be available in all python-pptx builds;
        # fall back to parsing from the raw blob.
        try:
            theme_xml = theme_part.element
        except AttributeError:
            from lxml import etree as _etree

            theme_xml = _etree.fromstring(theme_part.blob)
        ns = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        }

        # Extract color scheme
        clrScheme = theme_xml.find(".//a:clrScheme", ns)
        if clrScheme is not None:
            accent_colors = []
            for child in clrScheme:
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                # Get the color value from srgbClr or sysClr
                color_val = ""
                srgb = child.find("a:srgbClr", ns)
                if srgb is not None:
                    color_val = srgb.get("val", "")
                else:
                    sys_clr = child.find("a:sysClr", ns)
                    if sys_clr is not None:
                        color_val = sys_clr.get("lastClr", "") or sys_clr.get("val", "")

                if tag.startswith("accent"):
                    accent_colors.append(color_val)
                elif tag == "dk1":
                    theme.dk1 = color_val
                elif tag == "dk2":
                    theme.dk2 = color_val
                elif tag == "lt1":
                    theme.lt1 = color_val
                elif tag == "lt2":
                    theme.lt2 = color_val
                elif tag == "hlink":
                    theme.hyperlink = color_val
                elif tag == "folHlink":
                    theme.followed_hyperlink = color_val
            theme.accent_colors = accent_colors

        # Extract font scheme
        fontScheme = theme_xml.find(".//a:fontScheme", ns)
        if fontScheme is not None:
            major_font_elem = fontScheme.find(".//a:majorFont/a:latin", ns)
            if major_font_elem is not None:
                theme.major_font = major_font_elem.get("typeface", "Calibri")
            minor_font_elem = fontScheme.find(".//a:minorFont/a:latin", ns)
            if minor_font_elem is not None:
                theme.minor_font = minor_font_elem.get("typeface", "Calibri")
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Theme extraction error: %s" % str(e))
    return theme


def _resolve_scheme_color(scheme_val: str, theme: TemplateTheme) -> str:
    """Resolve a scheme color name to an RGB hex value using the theme."""
    scheme_map = {
        "dk1": theme.dk1,
        "dk2": theme.dk2,
        "lt1": theme.lt1,
        "lt2": theme.lt2,
        "tx1": theme.dk1,  # tx1 typically maps to dk1
        "tx2": theme.dk2,  # tx2 typically maps to dk2
        "bg1": theme.lt1,  # bg1 typically maps to lt1
        "bg2": theme.lt2,  # bg2 typically maps to lt2
        "accent1": theme.accent_colors[0] if len(theme.accent_colors) > 0 else "",
        "accent2": theme.accent_colors[1] if len(theme.accent_colors) > 1 else "",
        "accent3": theme.accent_colors[2] if len(theme.accent_colors) > 2 else "",
        "accent4": theme.accent_colors[3] if len(theme.accent_colors) > 3 else "",
        "accent5": theme.accent_colors[4] if len(theme.accent_colors) > 4 else "",
        "accent6": theme.accent_colors[5] if len(theme.accent_colors) > 5 else "",
        "hlink": theme.hyperlink,
        "folHlink": theme.followed_hyperlink,
    }
    return scheme_map.get(scheme_val, "")


def _extract_color_from_rPr(rPr_elem, theme: TemplateTheme) -> str:
    """Extract RGB hex color from an rPr or defRPr element, resolving scheme colors."""
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    solidFill = rPr_elem.find(ns_a + "solidFill")
    if solidFill is not None:
        srgb = solidFill.find(ns_a + "srgbClr")
        if srgb is not None:
            return srgb.get("val", "")
        schemeClr = solidFill.find(ns_a + "schemeClr")
        if schemeClr is not None:
            scheme_val = schemeClr.get("val", "")
            return _resolve_scheme_color(scheme_val, theme)
    return ""


def _hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color string to (R, G, B) tuple (0-255)."""
    hex_color = hex_color.strip().lstrip("#")
    if len(hex_color) != 6:
        return (0, 0, 0)
    try:
        return (
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )
    except ValueError:
        return (0, 0, 0)


def _relative_luminance(r: int, g: int, b: int) -> float:
    """Calculate relative luminance per WCAG 2.0 formula.

    Returns a value between 0.0 (black) and 1.0 (white).
    """

    def _linearize(c):
        c_srgb = c / 255.0
        if c_srgb <= 0.03928:
            return c_srgb / 12.92
        return ((c_srgb + 0.055) / 1.055) ** 2.4

    return 0.2126 * _linearize(r) + 0.7152 * _linearize(g) + 0.0722 * _linearize(b)


def _contrast_ratio(color1_rgb: tuple, color2_rgb: tuple) -> float:
    """Calculate WCAG contrast ratio between two RGB color tuples.

    Each argument should be an (R, G, B) tuple with values 0-255.
    Returns a ratio >= 1.0. WCAG AA requires >= 4.5 for normal text.
    """
    l1 = _relative_luminance(*color1_rgb)
    l2 = _relative_luminance(*color2_rgb)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _get_slide_background_color(
    slide, template_style: "TemplateStyle | None" = None
) -> str:
    """Determine the effective background color of a slide.

    Checks (in priority order):
    1. Slide's own background fill (solid fill)
    2. Slide layout's background fill
    3. Slide master's background fill
    4. Template theme lt1 (light 1, typically white)
    5. Default white "FFFFFF"

    Returns hex RGB string (e.g., "FFFFFF").
    """
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    ns_p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

    # Helper to extract color from a background element
    def _extract_bg_color(element):
        if element is None:
            return ""
        # Look for <p:bg><p:bgPr><a:solidFill>
        bgPr = element.find(ns_p + "bgPr")
        if bgPr is None:
            bgPr = element.find(ns_p + "bgRef")
        if bgPr is None:
            bgPr = element

        solidFill = bgPr.find(ns_a + "solidFill")
        if solidFill is not None:
            srgb = solidFill.find(ns_a + "srgbClr")
            if srgb is not None:
                return srgb.get("val", "")
            schemeClr = solidFill.find(ns_a + "schemeClr")
            if schemeClr is not None and template_style:
                return _resolve_scheme_color(
                    schemeClr.get("val", ""), template_style.theme
                )
        return ""

    try:
        # 1. Check slide background
        bg = slide._element.find(ns_p + "cSld/" + ns_p + "bg")
        color = _extract_bg_color(bg)
        if color:
            return color
    except Exception:
        pass

    try:
        # 2. Check slide layout background
        bg = slide.slide_layout._element.find(ns_p + "cSld/" + ns_p + "bg")
        color = _extract_bg_color(bg)
        if color:
            return color
    except Exception:
        pass

    try:
        # 3. Check slide master background
        bg = slide.slide_layout.slide_master._element.find(ns_p + "cSld/" + ns_p + "bg")
        color = _extract_bg_color(bg)
        if color:
            return color
    except Exception:
        pass

    # 4. Delegate to robust multi-layer detection (_get_shape_background_color)
    #    which handles blipFill (image), gradFill (gradient), bgRef (theme ref),
    #    large background-covering shapes, and dk1/dk2 theme heuristics.
    #    This is critical for templates with non-solidFill dark backgrounds.
    try:
        # Use any shape on the slide as a reference for per-shape detection
        for shape in slide.shapes:
            bg_color_robust = _get_shape_background_color(shape, slide)
            if bg_color_robust and bg_color_robust != "FFFFFF":
                print(
                    "  [BG DETECT] _get_slide_background_color: robust detection "
                    "found bg=%s via _get_shape_background_color fallback."
                    % bg_color_robust
                )
                return bg_color_robust
            break  # Only need to check once — background is slide-level
    except Exception:
        pass

    # 5. Use theme lt1 (usually white or near-white)
    if template_style and template_style.theme.lt1:
        return template_style.theme.lt1

    # 6. Default white
    return "FFFFFF"


def _ensure_text_contrast(slide, template_style: "TemplateStyle | None" = None) -> None:
    """Check all text elements on a slide and fix poor contrast.

    For each text run, checks if the text color has sufficient contrast
    against the slide's background. If contrast ratio < 3.0 (below WCAG AA),
    swaps the text to a high-contrast alternative.

    Uses theme dk1 (dark color) for text on light backgrounds, and
    theme lt1 (light color) for text on dark backgrounds.
    """
    if template_style is None:
        # Template-independent fallback: detect background and fix contrast
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                bg_hex = _get_shape_background_color(shape, slide)
                bg_rgb = _hex_to_rgb(bg_hex)
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        rPr = run._r.find(ns_a + "rPr")
                        if rPr is not None:
                            solidFill = rPr.find(ns_a + "solidFill")
                            if solidFill is not None:
                                text_hex = _extract_color_from_solid_fill(solidFill)
                                if text_hex:
                                    text_rgb = _hex_to_rgb(text_hex)
                                    ratio = _contrast_ratio(text_rgb, bg_rgb)
                                    if ratio < 3.0:
                                        _make_high_contrast_fill(rPr, bg_hex, existing_solidFill=solidFill)
                            else:
                                # Check default (black) against background
                                ratio = _contrast_ratio((0, 0, 0), bg_rgb)
                                if ratio < 3.0:
                                    _make_high_contrast_fill(rPr, bg_hex)
        return  # Done with template-independent path

    bg_color = _get_slide_background_color(slide, template_style)
    if not bg_color:
        return

    bg_r, bg_g, bg_b = _hex_to_rgb(bg_color)
    bg_luminance = _relative_luminance(bg_r, bg_g, bg_b)

    # Determine the contrasting color to use
    # If background is light (luminance > 0.5), use dark text
    # If background is dark (luminance <= 0.5), use light text
    theme = template_style.theme
    if bg_luminance > 0.5:
        # Light background: use dark text
        contrast_color = theme.dk1 or "000000"
    else:
        # Dark background: use light text
        contrast_color = theme.lt1 or "FFFFFF"

    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    CONTRAST_THRESHOLD = 3.0  # Minimum acceptable contrast ratio

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                # Get the current text color
                text_color = ""
                rPr = run._r.find(ns_a + "rPr")
                if rPr is not None:
                    solidFill = rPr.find(ns_a + "solidFill")
                    if solidFill is not None:
                        srgb = solidFill.find(ns_a + "srgbClr")
                        if srgb is not None:
                            text_color = srgb.get("val", "")
                        elif template_style:
                            schemeClr = solidFill.find(ns_a + "schemeClr")
                            if schemeClr is not None:
                                text_color = _resolve_scheme_color(
                                    schemeClr.get("val", ""), theme
                                )

                if not text_color:
                    # No explicit color — assume inherited default is dk1 (typically black).
                    # On dark backgrounds this creates invisible text. Fix it.
                    inherited_color = theme.dk1 or "000000"
                    ratio = _contrast_ratio(_hex_to_rgb(inherited_color), _hex_to_rgb(bg_color))
                    if ratio < CONTRAST_THRESHOLD:
                        # Inherited color has poor contrast — set explicit high-contrast color
                        if rPr is None:
                            rPr = etree.SubElement(run._r, ns_a + "rPr")
                        new_fill = etree.SubElement(rPr, ns_a + "solidFill")
                        srgb = etree.SubElement(new_fill, ns_a + "srgbClr")
                        srgb.set("val", contrast_color)
                    continue

                # Check contrast
                ratio = _contrast_ratio(_hex_to_rgb(text_color), _hex_to_rgb(bg_color))

                if ratio < CONTRAST_THRESHOLD:
                    if VERBOSE:
                        print(
                            "[VERBOSE] Low contrast detected: text=%s bg=%s ratio=%.1f, fixing to %s"
                            % (text_color, bg_color, ratio, contrast_color)
                        )

                    # Replace the text color with the contrasting color
                    if rPr is not None:
                        solidFill = rPr.find(ns_a + "solidFill")
                        if solidFill is not None:
                            rPr.remove(solidFill)
                        # Create new solidFill with contrast color
                        new_fill = etree.SubElement(rPr, ns_a + "solidFill")
                        srgb = etree.SubElement(new_fill, ns_a + "srgbClr")
                        srgb.set("val", contrast_color)


def _extract_table_styles_from_prs(prs) -> list:
    """Scan template slides for tables and extract their styling."""
    styles = []
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    for slide in prs.slides:
        resized_charts = []
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            style = TemplateTableStyle()

            # Extract tblPr XML (contains tblStyle reference and banding info)
            tbl_elem = table._tbl
            tblPr = tbl_elem.find("a:tblPr", ns)
            if tblPr is not None:
                style.table_properties_xml = copy.deepcopy(tblPr)
                style.table_style_id = tblPr.get("tblStyle", "") or ""

            # Extract header row styling (first row)
            if len(table.rows) > 0:
                first_row = table.rows[0]
                for cell in first_row.cells:
                    tf = cell.text_frame
                    if tf and tf.paragraphs:
                        for para in tf.paragraphs:
                            if para.runs:
                                run = para.runs[0]
                                rPr = run._r.find("a:rPr", ns)
                                if rPr is not None:
                                    # Font size
                                    sz = rPr.get("sz")
                                    if sz:
                                        style.header_font_size_pt = int(sz) // 100
                                    # Bold
                                    b = rPr.get("b")
                                    if b:
                                        style.header_font_bold = b in ("1", "true")
                                    # Font color
                                    solidFill = rPr.find("a:solidFill", ns)
                                    if solidFill is not None:
                                        srgb = solidFill.find("a:srgbClr", ns)
                                        if srgb is not None:
                                            style.header_font_color_rgb = srgb.get(
                                                "val", ""
                                            )
                                    # Font family
                                    latin = rPr.find("a:latin", ns)
                                    if latin is not None:
                                        style.header_font_family = latin.get(
                                            "typeface", ""
                                        )

                    # Cell fill
                    tc_elem = cell._tc
                    tcPr = tc_elem.find("a:tcPr", ns)
                    if tcPr is not None:
                        solidFill = tcPr.find("a:solidFill", ns)
                        if solidFill is not None:
                            srgb = solidFill.find("a:srgbClr", ns)
                            if srgb is not None:
                                style.header_fill_rgb = srgb.get("val", "")
                    break  # Only need first cell of header

            # Extract body cell styling (second row if exists)
            if len(table.rows) > 1:
                second_row = table.rows[1]
                for cell in second_row.cells:
                    tf = cell.text_frame
                    if tf and tf.paragraphs:
                        for para in tf.paragraphs:
                            if para.runs:
                                run = para.runs[0]
                                rPr = run._r.find("a:rPr", ns)
                                if rPr is not None:
                                    sz = rPr.get("sz")
                                    if sz:
                                        style.cell_font_size_pt = int(sz) // 100
                                    solidFill = rPr.find("a:solidFill", ns)
                                    if solidFill is not None:
                                        srgb = solidFill.find("a:srgbClr", ns)
                                        if srgb is not None:
                                            style.cell_font_color_rgb = srgb.get(
                                                "val", ""
                                            )
                                    latin = rPr.find("a:latin", ns)
                                    if latin is not None:
                                        style.cell_font_family = latin.get(
                                            "typeface", ""
                                        )

                    # Cell fill
                    tc_elem = cell._tc
                    tcPr = tc_elem.find("a:tcPr", ns)
                    if tcPr is not None:
                        solidFill = tcPr.find("a:solidFill", ns)
                        if solidFill is not None:
                            srgb = solidFill.find("a:srgbClr", ns)
                            if srgb is not None:
                                style.cell_fill_rgb = srgb.get("val", "")
                    break

            # Extract border styling from first cell
            if len(table.rows) > 0:
                for cell in table.rows[0].cells:
                    tc_elem = cell._tc
                    tcPr = tc_elem.find("a:tcPr", ns)
                    if tcPr is not None:
                        # Look for border elements (lnL, lnR, lnT, lnB)
                        for border_tag in ["a:lnB", "a:lnT", "a:lnL", "a:lnR"]:
                            ln = tcPr.find(border_tag, ns)
                            if ln is not None:
                                w = ln.get("w")
                                if w:
                                    style.border_width_emu = int(w)
                                solidFill = ln.find("a:solidFill", ns)
                                if solidFill is not None:
                                    srgb = solidFill.find("a:srgbClr", ns)
                                    if srgb is not None:
                                        style.border_color_rgb = srgb.get("val", "")
                                break  # One border is enough
                    break

            styles.append(style)
    return styles


def _extract_chart_styles_from_prs(prs) -> list:
    """Scan template slides for charts and extract their styling."""
    styles = []
    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
    }

    for slide in prs.slides:
        resized_charts = []
        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            chart = shape.chart
            style = TemplateChartStyle()

            try:
                style.chart_type = chart.chart_type
            except Exception:
                pass

            # Extract series colors from the chart XML
            chart_elem = chart._chartSpace

            # Find series elements and extract their fill colors
            series_elems = chart_elem.findall(".//c:ser", ns)
            for ser in series_elems:
                # Check for solidFill in the series
                spPr = ser.find("c:spPr", ns)
                if spPr is not None:
                    solidFill = spPr.find("a:solidFill", ns)
                    if solidFill is not None:
                        srgb = solidFill.find("a:srgbClr", ns)
                        if srgb is not None:
                            style.series_fill_colors.append(srgb.get("val", ""))
                        else:
                            schemeClr = solidFill.find("a:schemeClr", ns)
                            if schemeClr is not None:
                                style.series_fill_colors.append(
                                    "scheme:" + schemeClr.get("val", "")
                                )
                    # Line color
                    ln = spPr.find("a:ln", ns)
                    if ln is not None:
                        lnFill = ln.find("a:solidFill", ns)
                        if lnFill is not None:
                            srgb = lnFill.find("a:srgbClr", ns)
                            if srgb is not None:
                                style.series_line_colors.append(srgb.get("val", ""))

            # Extract axis styling
            for axis_tag in ["c:catAx", "c:valAx", "c:dateAx"]:
                ax = chart_elem.find(".//" + axis_tag, ns)
                if ax is not None:
                    # Axis text properties
                    txPr = ax.find("c:txPr", ns)
                    if txPr is not None:
                        defRPr = txPr.find(".//a:defRPr", ns)
                        if defRPr is not None:
                            sz = defRPr.get("sz")
                            if sz:
                                style.axis_font_size_pt = int(sz) // 100
                            solidFill = defRPr.find("a:solidFill", ns)
                            if solidFill is not None:
                                srgb = solidFill.find("a:srgbClr", ns)
                                if srgb is not None:
                                    style.axis_font_color_rgb = srgb.get("val", "")
                            latin = defRPr.find("a:latin", ns)
                            if latin is not None:
                                style.axis_font_family = latin.get("typeface", "")
                    # Axis line
                    spPr = ax.find("c:spPr", ns)
                    if spPr is not None:
                        ln = spPr.find("a:ln", ns)
                        if ln is not None:
                            sf = ln.find("a:solidFill", ns)
                            if sf is not None:
                                srgb = sf.find("a:srgbClr", ns)
                                if srgb is not None:
                                    style.axis_line_color_rgb = srgb.get("val", "")
                    break  # Use first axis found

            # Extract legend styling
            legend = chart_elem.find(".//c:legend", ns)
            if legend is not None:
                txPr = legend.find("c:txPr", ns)
                if txPr is not None:
                    defRPr = txPr.find(".//a:defRPr", ns)
                    if defRPr is not None:
                        sz = defRPr.get("sz")
                        if sz:
                            style.legend_font_size_pt = int(sz) // 100
                        solidFill = defRPr.find("a:solidFill", ns)
                        if solidFill is not None:
                            srgb = solidFill.find("a:srgbClr", ns)
                            if srgb is not None:
                                style.legend_font_color_rgb = srgb.get("val", "")
                        latin = defRPr.find("a:latin", ns)
                        if latin is not None:
                            style.legend_font_family = latin.get("typeface", "")

            # Extract data label styling
            dLbls = chart_elem.find(".//c:dLbls", ns)
            if dLbls is not None:
                txPr = dLbls.find("c:txPr", ns)
                if txPr is not None:
                    defRPr = txPr.find(".//a:defRPr", ns)
                    if defRPr is not None:
                        sz = defRPr.get("sz")
                        if sz:
                            style.data_label_font_size_pt = int(sz) // 100
                        solidFill = defRPr.find("a:solidFill", ns)
                        if solidFill is not None:
                            srgb = solidFill.find("a:srgbClr", ns)
                            if srgb is not None:
                                style.data_label_font_color_rgb = srgb.get("val", "")

            # Plot area fill
            plotArea = chart_elem.find(".//c:plotArea", ns)
            if plotArea is not None:
                spPr = plotArea.find("c:spPr", ns)
                if spPr is not None:
                    solidFill = spPr.find("a:solidFill", ns)
                    if solidFill is not None:
                        srgb = solidFill.find("a:srgbClr", ns)
                        if srgb is not None:
                            style.plot_area_fill_rgb = srgb.get("val", "")

            styles.append(style)
    return styles


def _extract_template_styles(template_prs) -> TemplateStyle:
    """Extract all styling information from a template presentation."""
    ts = TemplateStyle()
    try:
        ts.theme = _extract_theme_from_prs(template_prs)
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Theme extraction failed: %s" % str(e))
    try:
        ts.table_styles = _extract_table_styles_from_prs(template_prs)
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Table style extraction failed: %s" % str(e))
    try:
        ts.chart_styles = _extract_chart_styles_from_prs(template_prs)
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Chart style extraction failed: %s" % str(e))

    # Extract body/title font sizes from template placeholders
    try:
        for slide in template_prs.slides:
            for shape in slide.placeholders:
                ph_idx = shape.placeholder_format.idx
                if ph_idx == 0 and shape.has_text_frame:  # Title
                    for para in shape.text_frame.paragraphs:
                        if para.runs:
                            rPr = para.runs[0]._r.find(
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
                            )
                            if rPr is not None:
                                sz = rPr.get("sz")
                                if sz and not ts.title_font_size_pt:
                                    ts.title_font_size_pt = int(sz) // 100
                                color_rgb = _extract_color_from_rPr(rPr, ts.theme)
                                if color_rgb and not ts.title_font_color_rgb:
                                    ts.title_font_color_rgb = color_rgb
                                latin = rPr.find(
                                    "{http://schemas.openxmlformats.org/drawingml/2006/main}latin"
                                )
                                font_name = ""
                                if latin is not None:
                                    typeface = latin.get("typeface", "")
                                    if typeface and not typeface.startswith("+"):
                                        font_name = typeface
                                    elif typeface == "+mj-lt":
                                        font_name = ts.theme.major_font
                                    elif typeface == "+mn-lt":
                                        font_name = ts.theme.minor_font
                                if font_name and not ts.title_font_family:
                                    ts.title_font_family = font_name
                elif ph_idx == 1 and shape.has_text_frame:  # Body
                    for para in shape.text_frame.paragraphs:
                        if para.runs:
                            rPr = para.runs[0]._r.find(
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
                            )
                            if rPr is not None:
                                sz = rPr.get("sz")
                                if sz and not ts.body_font_size_pt:
                                    ts.body_font_size_pt = int(sz) // 100
                                color_rgb = _extract_color_from_rPr(rPr, ts.theme)
                                if color_rgb and not ts.body_font_color_rgb:
                                    ts.body_font_color_rgb = color_rgb
                                latin = rPr.find(
                                    "{http://schemas.openxmlformats.org/drawingml/2006/main}latin"
                                )
                                font_name = ""
                                if latin is not None:
                                    typeface = latin.get("typeface", "")
                                    if typeface and not typeface.startswith("+"):
                                        font_name = typeface
                                    elif typeface == "+mj-lt":
                                        font_name = ts.theme.major_font
                                    elif typeface == "+mn-lt":
                                        font_name = ts.theme.minor_font
                                if font_name and not ts.body_font_family:
                                    ts.body_font_family = font_name
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Placeholder style extraction failed: %s" % str(e))

    # If slide-level extraction didn't find explicit colors, scan slide layouts
    # Template placeholders often inherit their styling from the layout/master level
    if not ts.title_font_color_rgb or not ts.body_font_color_rgb:
        try:
            ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            for layout in template_prs.slide_layouts:
                for ph in layout.placeholders:
                    ph_idx = ph.placeholder_format.idx
                    if ph_idx in (0, 1) and ph.has_text_frame:
                        for para in ph.text_frame.paragraphs:
                            # Check default run properties on the paragraph level
                            pPr = para._p.find(ns_a + "pPr")
                            if pPr is not None:
                                defRPr = pPr.find(ns_a + "defRPr")
                                if defRPr is not None:
                                    color_rgb = _extract_color_from_rPr(
                                        defRPr, ts.theme
                                    )
                                    sz = defRPr.get("sz")
                                    if (
                                        ph_idx == 0
                                        and not ts.title_font_color_rgb
                                        and color_rgb
                                    ):
                                        ts.title_font_color_rgb = color_rgb
                                    if ph_idx == 0 and not ts.title_font_size_pt and sz:
                                        ts.title_font_size_pt = int(sz) // 100
                                    if (
                                        ph_idx == 1
                                        and not ts.body_font_color_rgb
                                        and color_rgb
                                    ):
                                        ts.body_font_color_rgb = color_rgb
                                    if ph_idx == 1 and not ts.body_font_size_pt and sz:
                                        ts.body_font_size_pt = int(sz) // 100
                                    # Also extract font family from defRPr
                                    latin = defRPr.find(ns_a + "latin")
                                    font_name = ""
                                    if latin is not None:
                                        typeface = latin.get("typeface", "")
                                        if typeface and not typeface.startswith("+"):
                                            font_name = typeface
                                        elif typeface == "+mj-lt":
                                            font_name = ts.theme.major_font
                                        elif typeface == "+mn-lt":
                                            font_name = ts.theme.minor_font
                                    if (
                                        ph_idx == 0
                                        and not ts.title_font_family
                                        and font_name
                                    ):
                                        ts.title_font_family = font_name
                                    if (
                                        ph_idx == 1
                                        and not ts.body_font_family
                                        and font_name
                                    ):
                                        ts.body_font_family = font_name
                            # Also check actual runs
                            if para.runs:
                                rPr = para.runs[0]._r.find(ns_a + "rPr")
                                if rPr is not None:
                                    color_rgb = _extract_color_from_rPr(rPr, ts.theme)
                                    sz = rPr.get("sz")
                                    if (
                                        ph_idx == 0
                                        and not ts.title_font_color_rgb
                                        and color_rgb
                                    ):
                                        ts.title_font_color_rgb = color_rgb
                                    if ph_idx == 0 and not ts.title_font_size_pt and sz:
                                        ts.title_font_size_pt = int(sz) // 100
                                    if (
                                        ph_idx == 1
                                        and not ts.body_font_color_rgb
                                        and color_rgb
                                    ):
                                        ts.body_font_color_rgb = color_rgb
                                    if ph_idx == 1 and not ts.body_font_size_pt and sz:
                                        ts.body_font_size_pt = int(sz) // 100
                                    # Also extract font family from run rPr
                                    latin = rPr.find(ns_a + "latin")
                                    font_name = ""
                                    if latin is not None:
                                        typeface = latin.get("typeface", "")
                                        if typeface and not typeface.startswith("+"):
                                            font_name = typeface
                                        elif typeface == "+mj-lt":
                                            font_name = ts.theme.major_font
                                        elif typeface == "+mn-lt":
                                            font_name = ts.theme.minor_font
                                    if (
                                        ph_idx == 0
                                        and not ts.title_font_family
                                        and font_name
                                    ):
                                        ts.title_font_family = font_name
                                    if (
                                        ph_idx == 1
                                        and not ts.body_font_family
                                        and font_name
                                    ):
                                        ts.body_font_family = font_name
        except Exception as e:
            if VERBOSE:
                print(
                    "[VERBOSE] Layout placeholder style extraction failed: %s" % str(e)
                )

    if VERBOSE:
        print("[VERBOSE] Extracted template styles:")
        print("[VERBOSE]   Theme accent colors: %s" % ts.theme.accent_colors)
        print(
            "[VERBOSE]   Theme fonts: major=%s minor=%s"
            % (ts.theme.major_font, ts.theme.minor_font)
        )
        print("[VERBOSE]   Reference tables found: %d" % len(ts.table_styles))
        print("[VERBOSE]   Reference charts found: %d" % len(ts.chart_styles))
        print("[VERBOSE]   Title font family: %s" % ts.title_font_family)
        print("[VERBOSE]   Body font family: %s" % ts.body_font_family)

    return ts


# ---------------------------------------------------------------------------
# Template Style Application Functions
# ---------------------------------------------------------------------------


def _apply_table_style(table_shape, template_style: TemplateStyle):
    """Apply template-derived styling to a newly created table."""
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    table = table_shape.table
    tbl = table._tbl

    # Get the best matching table style from the template
    ref_style = template_style.table_styles[0] if template_style.table_styles else None
    theme = template_style.theme

    # Apply tblPr from reference if available (includes tblStyle ID + banding flags)
    if ref_style and ref_style.table_properties_xml is not None:
        existing_tblPr = tbl.find(qn("a:tblPr"))
        if existing_tblPr is not None:
            tbl.remove(existing_tblPr)
        new_tblPr = copy.deepcopy(ref_style.table_properties_xml)
        tbl.insert(0, new_tblPr)

    # Determine font settings (reference table > theme > defaults)
    header_font_family = (
        ref_style.header_font_family
        if ref_style and ref_style.header_font_family
        else theme.minor_font or "Calibri"
    )
    cell_font_family = (
        ref_style.cell_font_family
        if ref_style and ref_style.cell_font_family
        else theme.minor_font or "Calibri"
    )
    header_font_size = Pt(
        ref_style.header_font_size_pt
        if ref_style and ref_style.header_font_size_pt
        else 11
    )
    cell_font_size = Pt(
        ref_style.cell_font_size_pt if ref_style and ref_style.cell_font_size_pt else 10
    )
    header_bold = ref_style.header_font_bold if ref_style else True

    # Color resolution helper
    def _resolve_color(rgb_hex, fallback=""):
        if rgb_hex:
            try:
                return RGBColor.from_string(rgb_hex)
            except Exception:
                pass
        if fallback:
            try:
                return RGBColor.from_string(fallback)
            except Exception:
                pass
        return None

    # Apply styling to each cell
    for r_idx, row in enumerate(table.rows):
        is_header = r_idx == 0
        for cell in row.cells:
            # Font formatting
            for para in cell.text_frame.paragraphs:
                para.font.size = header_font_size if is_header else cell_font_size
                para.font.bold = header_bold if is_header else False
                para.font.name = header_font_family if is_header else cell_font_family

                # Font color
                if is_header and ref_style and ref_style.header_font_color_rgb:
                    color = _resolve_color(ref_style.header_font_color_rgb)
                    if color:
                        para.font.color.rgb = color
                elif not is_header and ref_style and ref_style.cell_font_color_rgb:
                    color = _resolve_color(ref_style.cell_font_color_rgb)
                    if color:
                        para.font.color.rgb = color
                elif theme.dk1:
                    color = _resolve_color(theme.dk1)
                    if color:
                        para.font.color.rgb = color

            cell.text_frame.word_wrap = True

            # Cell fill (via XML for reliability)
            tc_elem = cell._tc
            tcPr = tc_elem.find(qn("a:tcPr"))
            if tcPr is None:
                tcPr = etree.SubElement(tc_elem, qn("a:tcPr"))

            if is_header and ref_style and ref_style.header_fill_rgb:
                # Remove existing fill
                for old_fill in tcPr.findall(qn("a:solidFill")):
                    tcPr.remove(old_fill)
                fill_elem = etree.SubElement(tcPr, qn("a:solidFill"))
                srgb = etree.SubElement(fill_elem, qn("a:srgbClr"))
                srgb.set("val", ref_style.header_fill_rgb)
            elif is_header and theme.accent_colors:
                # Use first accent color for header
                for old_fill in tcPr.findall(qn("a:solidFill")):
                    tcPr.remove(old_fill)
                fill_elem = etree.SubElement(tcPr, qn("a:solidFill"))
                srgb = etree.SubElement(fill_elem, qn("a:srgbClr"))
                srgb.set("val", theme.accent_colors[0])

            if not is_header and ref_style and ref_style.cell_fill_rgb:
                for old_fill in tcPr.findall(qn("a:solidFill")):
                    tcPr.remove(old_fill)
                fill_elem = etree.SubElement(tcPr, qn("a:solidFill"))
                srgb = etree.SubElement(fill_elem, qn("a:srgbClr"))
                srgb.set("val", ref_style.cell_fill_rgb)

            # Borders
            if ref_style and ref_style.border_color_rgb:
                border_color = ref_style.border_color_rgb
                border_w = str(ref_style.border_width_emu)
                for border_name in ["lnL", "lnR", "lnT", "lnB"]:
                    ln = tcPr.find(qn("a:" + border_name))
                    if ln is None:
                        ln = etree.SubElement(tcPr, qn("a:" + border_name))
                    ln.set("w", border_w)
                    for old_fill in ln.findall(qn("a:solidFill")):
                        ln.remove(old_fill)
                    bf = etree.SubElement(ln, qn("a:solidFill"))
                    srgb = etree.SubElement(bf, qn("a:srgbClr"))
                    srgb.set("val", border_color)


def _apply_chart_style(chart_shape, template_style: TemplateStyle):
    """Apply template-derived styling to a newly created chart."""
    from pptx.oxml.ns import qn

    chart = chart_shape.chart
    theme = template_style.theme

    # Find the best matching chart style (prefer same chart type)
    ref_style = None
    for cs in template_style.chart_styles:
        try:
            if cs.chart_type == chart.chart_type:
                ref_style = cs
                break
        except Exception:
            pass
    if ref_style is None and template_style.chart_styles:
        ref_style = template_style.chart_styles[0]

    chart_space = chart._chartSpace

    # Apply series colors
    colors_to_use = []
    if ref_style and ref_style.series_fill_colors:
        colors_to_use = ref_style.series_fill_colors
    elif theme.accent_colors:
        colors_to_use = theme.accent_colors

    if colors_to_use:
        series_elems = chart_space.findall(".//" + qn("c:ser"))
        for i, ser in enumerate(series_elems):
            color = colors_to_use[i % len(colors_to_use)]
            if color.startswith("scheme:"):
                continue  # Skip scheme colors for now

            spPr = ser.find(qn("c:spPr"))
            if spPr is None:
                spPr = etree.SubElement(ser, qn("c:spPr"))

            # Set fill color
            for old in spPr.findall(qn("a:solidFill")):
                spPr.remove(old)
            fill = etree.SubElement(spPr, qn("a:solidFill"))
            srgb = etree.SubElement(fill, qn("a:srgbClr"))
            srgb.set("val", color)

    # Apply axis styling
    font_family = (
        ref_style.axis_font_family
        if ref_style and ref_style.axis_font_family
        else theme.minor_font or "Calibri"
    )
    font_size = str(
        (
            ref_style.axis_font_size_pt
            if ref_style and ref_style.axis_font_size_pt
            else 10
        )
        * 100
    )
    font_color = ref_style.axis_font_color_rgb if ref_style else (theme.dk1 or "000000")

    for axis_tag in [qn("c:catAx"), qn("c:valAx"), qn("c:dateAx")]:
        for ax in chart_space.findall(".//" + axis_tag):
            # Create or update txPr (text properties)
            txPr = ax.find(qn("c:txPr"))
            if txPr is None:
                txPr = etree.SubElement(ax, qn("c:txPr"))
            # Clear existing
            for child in list(txPr):
                txPr.remove(child)
            # Build text properties
            etree.SubElement(txPr, qn("a:bodyPr"))
            etree.SubElement(txPr, qn("a:lstStyle"))
            p = etree.SubElement(txPr, qn("a:p"))
            pPr = etree.SubElement(p, qn("a:pPr"))
            defRPr = etree.SubElement(pPr, qn("a:defRPr"))
            defRPr.set("sz", font_size)
            if font_color:
                sf = etree.SubElement(defRPr, qn("a:solidFill"))
                srgb = etree.SubElement(sf, qn("a:srgbClr"))
                srgb.set("val", font_color)
            latin = etree.SubElement(defRPr, qn("a:latin"))
            latin.set("typeface", font_family)
            endParaRPr = etree.SubElement(p, qn("a:endParaRPr"))
            endParaRPr.set("lang", "en-US")

            # Axis line color
            if ref_style and ref_style.axis_line_color_rgb:
                spPr = ax.find(qn("c:spPr"))
                if spPr is None:
                    spPr = etree.SubElement(ax, qn("c:spPr"))
                ln = spPr.find(qn("a:ln"))
                if ln is None:
                    ln = etree.SubElement(spPr, qn("a:ln"))
                for old in ln.findall(qn("a:solidFill")):
                    ln.remove(old)
                sf = etree.SubElement(ln, qn("a:solidFill"))
                srgb = etree.SubElement(sf, qn("a:srgbClr"))
                srgb.set("val", ref_style.axis_line_color_rgb)

    # Apply legend styling
    legend = chart_space.find(".//" + qn("c:legend"))
    if legend is not None:
        leg_font_family = (
            ref_style.legend_font_family
            if ref_style and ref_style.legend_font_family
            else font_family
        )
        leg_font_size = str(
            (
                ref_style.legend_font_size_pt
                if ref_style and ref_style.legend_font_size_pt
                else 10
            )
            * 100
        )
        leg_font_color = ref_style.legend_font_color_rgb if ref_style else font_color

        txPr = legend.find(qn("c:txPr"))
        if txPr is None:
            txPr = etree.SubElement(legend, qn("c:txPr"))
        for child in list(txPr):
            txPr.remove(child)
        etree.SubElement(txPr, qn("a:bodyPr"))
        etree.SubElement(txPr, qn("a:lstStyle"))
        p = etree.SubElement(txPr, qn("a:p"))
        pPr = etree.SubElement(p, qn("a:pPr"))
        defRPr = etree.SubElement(pPr, qn("a:defRPr"))
        defRPr.set("sz", leg_font_size)
        if leg_font_color:
            sf = etree.SubElement(defRPr, qn("a:solidFill"))
            srgb = etree.SubElement(sf, qn("a:srgbClr"))
            srgb.set("val", leg_font_color)
        latin = etree.SubElement(defRPr, qn("a:latin"))
        latin.set("typeface", leg_font_family)
        endParaRPr = etree.SubElement(p, qn("a:endParaRPr"))
        endParaRPr.set("lang", "en-US")

    # Plot area fill
    if ref_style and ref_style.plot_area_fill_rgb:
        plotArea = chart_space.find(".//" + qn("c:plotArea"))
        if plotArea is not None:
            spPr = plotArea.find(qn("c:spPr"))
            if spPr is None:
                spPr = etree.SubElement(plotArea, qn("c:spPr"))
            for old in spPr.findall(qn("a:solidFill")):
                spPr.remove(old)
            fill = etree.SubElement(spPr, qn("a:solidFill"))
            srgb = etree.SubElement(fill, qn("a:srgbClr"))
            srgb.set("val", ref_style.plot_area_fill_rgb)


# ---------------------------------------------------------------------------
# Content Extraction Functions (from agent_with_powerpoint_template.py)
# ---------------------------------------------------------------------------


def _extract_slide_content(slide) -> SlideContent:
    """Extract all content from a slide including text, tables, images, charts, and shapes.

    Detects picture placeholders using multiple strategies for robustness:
      1. int(type) == 18  (the raw OOXML value for PICTURE placeholders)
      2. str(type) contains 'PICTURE (18)'
      3. PP_PLACEHOLDER.PICTURE enum comparison
      4. XML-level fallback: <p:ph type="pic"/>
    """
    content = SlideContent()

    # Detect picture placeholders in the slide layout
    for shape in slide.placeholders:
        ph_fmt = shape.placeholder_format
        if ph_fmt is not None:
            ph_type_val = ph_fmt.type
            # Strategy 1 & 2: enum / int / string comparison
            is_picture_ph = False
            try:
                if ph_type_val is not None and (
                    int(ph_type_val) == 18 or str(ph_type_val) == "PICTURE (18)"
                ):
                    is_picture_ph = True
            except (ValueError, TypeError):
                pass
            # Also try the enum directly
            if not is_picture_ph:
                try:
                    if ph_type_val == PP_PLACEHOLDER.PICTURE:
                        is_picture_ph = True
                except Exception as e:
                    if VERBOSE:
                        print("[VERBOSE] Exception suppressed: %s" % str(e))
            # Strategy 3: XML-level fallback
            if not is_picture_ph:
                nsmap = {
                    "p": "http://schemas.openxmlformats.org/presentationml/2006/main"
                }
                ph_elem = shape._element.find(".//p:ph", nsmap)
                if ph_elem is not None and ph_elem.get("type") == "pic":
                    is_picture_ph = True
            if is_picture_ph:
                content.has_image_placeholder = True
                content.image_placeholder_indices.append(ph_fmt.idx)

    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                rows_data.append(row_data)
            content.tables.append(
                TableData(
                    rows=rows_data,
                    left=shape.left,
                    top=shape.top,
                    width=shape.width,
                    height=shape.height,
                )
            )
            continue

        if shape.has_chart:
            try:
                chart = shape.chart
                chart_type_val = chart.chart_type
                categories = []
                series_data = []
                plot = chart.plots[0] if chart.plots else None
                if plot:
                    if plot.categories:
                        categories = list(plot.categories)
                    for series in plot.series:
                        name = series.name if hasattr(series, "name") else ""
                        values = (
                            list(series.values) if hasattr(series, "values") else []
                        )
                        series_data.append((name or "", values))
                if categories or series_data:
                    content.charts.append(
                        ChartExtract(
                            chart_type=chart_type_val,
                            categories=categories,
                            series=series_data,
                            left=shape.left,
                            top=shape.top,
                            width=shape.width,
                            height=shape.height,
                        )
                    )
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                ct = shape.image.content_type
                content.images.append(
                    ImageData(
                        blob=blob,
                        left=shape.left,
                        top=shape.top,
                        width=shape.width,
                        height=shape.height,
                        content_type=ct,
                    )
                )
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_xml = copy.deepcopy(shape._element)
            content.shapes_xml.append(group_xml)
            try:
                for grp_shape in shape.shapes:
                    if grp_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            blob = grp_shape.image.blob
                            ct = grp_shape.image.content_type
                            content.images.append(
                                ImageData(
                                    blob=blob,
                                    left=grp_shape.left,
                                    top=grp_shape.top,
                                    width=grp_shape.width,
                                    height=grp_shape.height,
                                    content_type=ct,
                                )
                            )
                        except Exception as e:
                            if VERBOSE:
                                print("[VERBOSE] Exception suppressed: %s" % str(e))
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))
            continue

        if shape.has_text_frame:
            text_frame = shape.text_frame
            # Guard: shape.is_placeholder can be True but placeholder_format may
            # still raise "shape is not a placeholder" for Tier-2 LLM-generated
            # slides with non-standard XML. Always access via try/except.
            ph_idx = None
            if shape.is_placeholder:
                try:
                    ph_fmt = shape.placeholder_format
                    ph_idx = ph_fmt.idx if ph_fmt is not None else None
                except Exception:
                    ph_idx = None  # treat as non-placeholder text box

            if ph_idx is not None:
                if ph_idx == 0:
                    content.title = text_frame.text.strip()
                elif ph_idx == 1:
                    paragraphs = text_frame.paragraphs
                    if len(paragraphs) == 1 and paragraphs[0].level == 0:
                        content.subtitle = paragraphs[0].text.strip()
                    else:
                        for para in paragraphs:
                            text = para.text.strip()
                            if text:
                                content.body_paragraphs.append((text, para.level))
                else:
                    for para in text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            content.body_paragraphs.append((text, para.level))
            else:
                # Non-placeholder text boxes: preserve styling via XML.
                text_val = text_frame.text.strip()
                if text_val:
                    content.has_text_box = True
                    for para in text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            content.text_box_paragraphs.append((t, para.level))
                shape_xml = copy.deepcopy(shape._element)
                content.text_shapes_xml.append(shape_xml)
        elif not shape.is_placeholder:
            shape_xml = copy.deepcopy(shape._element)
            content.shapes_xml.append(shape_xml)

    return content


# ---------------------------------------------------------------------------
# Layout Collision Resolution Helpers
# ---------------------------------------------------------------------------


def _is_picture_placeholder(shape) -> bool:
    """Check if a shape is a picture placeholder using multiple strategies."""
    ph_fmt = shape.placeholder_format
    if ph_fmt is None:
        return False
    ph_type_val = ph_fmt.type
    # Strategy 1: int/string comparison
    try:
        if ph_type_val is not None and (
            int(ph_type_val) == 18 or str(ph_type_val) == "PICTURE (18)"
        ):
            return True
    except (ValueError, TypeError):
        pass
    # Strategy 2: PP_PLACEHOLDER enum
    try:
        if ph_type_val == PP_PLACEHOLDER.PICTURE:
            return True
    except Exception:
        pass
    # Strategy 3: XML-level fallback
    try:
        nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        ph_elem = shape._element.find(".//p:ph", nsmap)
        if ph_elem is not None and ph_elem.get("type") in ("pic", "clipArt"):
            return True
    except Exception:
        pass
    return False


def _placeholder_area(ph) -> int:
    """Return placeholder area in EMU^2."""
    try:
        return int(ph.width) * int(ph.height)
    except Exception:
        return 0


def _is_text_placeholder_type(ph_type) -> bool:
    """Return True if placeholder type is intended for text content."""
    return ph_type in (
        PP_PLACEHOLDER.BODY,
        PP_PLACEHOLDER.OBJECT,
        PP_PLACEHOLDER.SUBTITLE,
    )


def _is_visual_placeholder_type(ph_type) -> bool:
    """Return True if placeholder type is intended for visual content."""
    return ph_type in (
        PP_PLACEHOLDER.PICTURE,
        PP_PLACEHOLDER.CHART,
        PP_PLACEHOLDER.TABLE,
    )


def _largest_placeholder(layout, types=None, min_area: int = 0):
    """Return the largest placeholder matching types and min_area, or None."""
    best = None
    best_area = 0
    for ph in layout.placeholders:
        try:
            ph_type = ph.placeholder_format.type
        except Exception:
            continue
        if types is not None and ph_type not in types:
            continue
        area = _placeholder_area(ph)
        if area < min_area:
            continue
        if area > best_area:
            best = ph
            best_area = area
    return best


def _layout_placeholder_summary(layout) -> dict:
    """Return counts of placeholder types for logging."""
    summary = {
        "title": 0,
        "subtitle": 0,
        "body": 0,
        "object": 0,
        "picture": 0,
        "chart": 0,
        "table": 0,
        "other": 0,
    }
    for ph in layout.placeholders:
        try:
            ph_type = ph.placeholder_format.type
        except Exception:
            summary["other"] += 1
            continue
        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            summary["title"] += 1
        elif ph_type == PP_PLACEHOLDER.SUBTITLE:
            summary["subtitle"] += 1
        elif ph_type == PP_PLACEHOLDER.BODY:
            summary["body"] += 1
        elif ph_type == PP_PLACEHOLDER.OBJECT:
            summary["object"] += 1
        elif ph_type == PP_PLACEHOLDER.PICTURE:
            summary["picture"] += 1
        elif ph_type == PP_PLACEHOLDER.CHART:
            summary["chart"] += 1
        elif ph_type == PP_PLACEHOLDER.TABLE:
            summary["table"] += 1
        else:
            summary["other"] += 1
    return summary


def _layout_richness_score(layout) -> int:
    """Heuristic for visually rich layouts (cards, shapes, picture slots)."""
    score = 0
    # Prefer layouts with multiple body placeholders (card-style layouts)
    body_count = 0
    for ph in layout.placeholders:
        try:
            if ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
                body_count += 1
        except Exception:
            continue
    if body_count >= 2:
        score += 10
    if body_count >= 3:
        score += 10

    # Prefer layouts with at least one picture placeholder (visual interest)
    has_pic = False
    for ph in layout.placeholders:
        if _is_picture_placeholder(ph):
            has_pic = True
            break
    if has_pic:
        score += 15

    # Prefer layouts with non-placeholder shapes (cards, colored blocks)
    non_placeholder_shapes = 0
    for shape in layout.shapes:
        try:
            if shape.is_placeholder:
                continue
            # Ignore lines
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                continue
            non_placeholder_shapes += 1
        except Exception:
            continue
    if non_placeholder_shapes >= 2:
        score += 15
    elif non_placeholder_shapes == 1:
        score += 5

    return score


def _rect_overlap(a: ContentArea, b: ContentArea) -> int:
    """Return overlap area between two rectangles (or 0)."""
    x1 = max(a.left, b.left)
    y1 = max(a.top, b.top)
    x2 = min(a.left + a.width, b.left + b.width)
    y2 = min(a.top + a.height, b.top + b.height)
    if x2 <= x1 or y2 <= y1:
        return 0
    return (x2 - x1) * (y2 - y1)


def _layout_non_placeholder_text_bounds(layout) -> list[ContentArea]:
    """Collect bounds of non-placeholder text boxes defined in the layout."""
    bounds = []
    for shape in layout.shapes:
        try:
            if shape.is_placeholder:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = getattr(shape, "text", "")
            if text and text.strip():
                bounds.append(
                    ContentArea(
                        left=shape.left,
                        top=shape.top,
                        width=shape.width,
                        height=shape.height,
                    )
                )
        except Exception:
            continue
    return bounds


def _best_visual_placeholder(
    layout,
    desired_types: set,
    min_area: int,
    reserved: list[ContentArea],
    allow_body_fallback: bool = True,
):
    """Pick a visual placeholder that minimizes overlap with reserved text areas."""
    candidates = []
    for ph in layout.placeholders:
        try:
            ph_type = ph.placeholder_format.type
        except Exception:
            continue
        if ph_type in desired_types:
            candidates.append(ph)
    if allow_body_fallback and not candidates:
        for ph in layout.placeholders:
            try:
                ph_type = ph.placeholder_format.type
            except Exception:
                continue
            if ph_type in (
                PP_PLACEHOLDER.BODY,
                PP_PLACEHOLDER.OBJECT,
                PP_PLACEHOLDER.SUBTITLE,
            ):
                candidates.append(ph)

    best = None
    best_score = None
    for ph in candidates:
        area = _placeholder_area(ph)
        if area < min_area and desired_types != {PP_PLACEHOLDER.PICTURE}:
            # allow small picture placeholders but penalize small chart/table
            pass
        ph_bounds = ContentArea(
            left=ph.left, top=ph.top, width=ph.width, height=ph.height
        )
        overlap = 0
        for r in reserved:
            overlap += _rect_overlap(ph_bounds, r)
        # Prefer area-sufficient placeholders first, then zero overlap, then
        # minimal overlap, then largest area.
        # Previously the tuple was (overlap==0, -overlap, area) which allowed
        # a tiny zero-overlap footer placeholder to outrank a large content
        # placeholder that had any layout overlap — causing images to be placed
        # in a small footer-left slot instead of the main visual region.
        score = (area >= min_area, overlap == 0, -overlap, area)
        if best_score is None or score > best_score:
            best_score = score
            best = ph
    return best


def _best_text_placeholder(layout, avoid: ContentArea | None = None):
    """Pick a text placeholder with minimal overlap to avoid region."""
    candidates = []
    for ph in layout.placeholders:
        try:
            ph_type = ph.placeholder_format.type
        except Exception:
            continue
        if ph_type in (
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.SUBTITLE,
        ):
            candidates.append(ph)
    if not candidates:
        return None
    best = None
    best_score = None
    for ph in candidates:
        area = _placeholder_area(ph)
        overlap = 0
        if avoid is not None:
            ph_bounds = ContentArea(
                left=ph.left, top=ph.top, width=ph.width, height=ph.height
            )
            overlap = _rect_overlap(ph_bounds, avoid)
        score = (overlap == 0, -overlap, area)
        if best_score is None or score > best_score:
            best_score = score
            best = ph
    return best


def _classify_content_mix(
    content: SlideContent, has_generated_image: bool = False
) -> ContentMix:
    """Classify the mix of element types in a slide."""
    has_text = bool(
        content.title
        or content.body_paragraphs
        or content.subtitle
        or content.has_text_box
    )
    has_table = len(content.tables) > 0
    has_chart = len(content.charts) > 0
    has_image = len(content.images) > 0

    visual_count = sum([has_table, has_chart, has_image, has_generated_image])

    if visual_count == 0:
        return ContentMix.TEXT_ONLY
    # If visuals exist but no text, still treat generated/user images as paired with text
    # so we reserve space for titles/subtitles in the template.
    if not has_text and has_generated_image:
        return ContentMix.TEXT_AND_GENERATED_IMAGE
    if not has_text and has_image:
        return ContentMix.TEXT_AND_IMAGE
    if not has_text and visual_count > 0:
        return ContentMix.VISUAL_ONLY
    if visual_count > 1:
        return ContentMix.MIXED
    if has_table:
        return ContentMix.TEXT_AND_TABLE
    if has_chart:
        return ContentMix.TEXT_AND_CHART
    if has_generated_image:
        return ContentMix.TEXT_AND_GENERATED_IMAGE
    if has_image:
        return ContentMix.TEXT_AND_IMAGE
    return ContentMix.TEXT_ONLY


def _compute_text_ratio(content: SlideContent, content_mix: ContentMix) -> float:
    """Compute how much of the content area should be allocated to text.

    Returns a ratio (0.0-1.0) representing the text portion.

    Factors in average line length to account for text wrapping: a slide with
    3 long bullets needs more height than one with 3 short bullets.
    """
    num_paragraphs = len(content.body_paragraphs)

    # Estimate wrap factor: average chars relative to ~45 chars-per-line threshold.
    # Long bullets wrap to multiple display lines and need more vertical space.
    CHARS_PER_LINE = 45
    wrap_factor = 1.0
    if content.body_paragraphs:
        total_chars = sum(len(text) for text, _level in content.body_paragraphs)
        avg_chars = total_chars / len(content.body_paragraphs)
        wrap_factor = max(1.0, avg_chars / CHARS_PER_LINE)

    # Effective line count accounts for text wrapping
    effective_lines = num_paragraphs * wrap_factor

    if content_mix in (
        ContentMix.TEXT_AND_TABLE,
        ContentMix.TEXT_AND_CHART,
        ContentMix.MIXED,
    ):
        # Top/bottom split: give text 25-50% based on effective line count
        if effective_lines <= 2:
            return 0.25
        elif effective_lines <= 4:
            return 0.35
        elif effective_lines <= 6:
            return 0.45
        else:
            return 0.50

    if content_mix in (ContentMix.TEXT_AND_IMAGE, ContentMix.TEXT_AND_GENERATED_IMAGE):
        # Left/right split: give text 40-60% based on effective line count
        if effective_lines <= 2:
            return 0.40
        elif effective_lines <= 4:
            return 0.48
        elif effective_lines <= 6:
            return 0.55
        else:
            return 0.60

    return 0.50  # Default even split


def _compute_max_font_size(
    region: ContentArea, num_paragraphs: int, is_title: bool = False
) -> int:
    """Compute the maximum font size that fits text in the given region.

    Returns font size in points, clamped between 10-28pt for titles and 10-18pt for body.
    """
    if num_paragraphs == 0:
        return 28 if is_title else 18

    # Rough estimate: each point of font size ~ 1 * EMU_PER_PT in EMU height per line
    # Use 1.8x factor to account for line spacing PLUS paragraph spacing above/below.
    # The original 1.5 was too low and caused text overflow on dense body slides.
    EMU_PER_PT = 12700  # 1 point = 12700 EMU
    LINE_SPACING_FACTOR = 1.8

    available_height = region.height
    lines_needed = max(num_paragraphs, 1)

    # font_size_pt <= available_height / (lines_needed * EMU_PER_PT * LINE_SPACING_FACTOR)
    max_size_pt = int(
        available_height / (lines_needed * EMU_PER_PT * LINE_SPACING_FACTOR)
    )

    if is_title:
        return max(10, min(28, max_size_pt))
    else:
        # Body text cap lowered from 18pt to 16pt as a defensive measure against
        # overflow on slides with many paragraphs or templates with larger line heights.
        return max(10, min(16, max_size_pt))


def _compute_region_map(
    layout,
    content_mix: ContentMix,
    slide_width: int,
    slide_height: int,
    content: SlideContent,
) -> RegionMap:
    """Compute separate text and visual regions based on content mix and layout."""
    slide_area = int(slide_width) * int(slide_height)
    min_area = int(slide_area * 0.10)
    content_area = _get_content_area(layout, slide_width, slide_height)
    layout_text_bounds = _layout_non_placeholder_text_bounds(layout)

    # For text-only or visual-only slides, use the full content area
    if content_mix in (ContentMix.TEXT_ONLY, ContentMix.VISUAL_ONLY):
        return RegionMap(
            text_region=content_area,
            visual_region=content_area,
            layout_type="full",
        )

    # Check if the layout natively has separate content placeholders
    # (e.g., "Two Content", "Content + Picture")
    content_placeholders = []
    picture_placeholders = []
    chart_placeholders = []
    table_placeholders = []
    for ph in layout.placeholders:
        if _is_picture_placeholder(ph):
            picture_placeholders.append(ph)
        elif ph.placeholder_format.type == PP_PLACEHOLDER.CHART:
            chart_placeholders.append(ph)
        elif ph.placeholder_format.type == PP_PLACEHOLDER.TABLE:
            table_placeholders.append(ph)
        elif ph.placeholder_format.idx > 0 and ph.has_text_frame:
            content_placeholders.append(ph)

    if content_mix in (ContentMix.TEXT_AND_IMAGE, ContentMix.TEXT_AND_GENERATED_IMAGE):
        # If layout has a picture placeholder + content placeholder, use native regions
        if picture_placeholders and content_placeholders:
            text_ph = content_placeholders[0]
            pic_ph = picture_placeholders[0]
            return RegionMap(
                text_region=ContentArea(
                    left=text_ph.left,
                    top=text_ph.top,
                    width=text_ph.width,
                    height=text_ph.height,
                ),
                visual_region=ContentArea(
                    left=pic_ph.left,
                    top=pic_ph.top,
                    width=pic_ph.width,
                    height=pic_ph.height,
                ),
                layout_type="native",
            )
        text_ph = _largest_placeholder(
            layout,
            {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE},
            min_area=0,
        )
        reserved = list(layout_text_bounds)
        if text_ph is not None:
            reserved.append(
                ContentArea(
                    left=text_ph.left,
                    top=text_ph.top,
                    width=text_ph.width,
                    height=text_ph.height,
                )
            )
        pic_ph = _best_visual_placeholder(
            layout,
            {PP_PLACEHOLDER.PICTURE},
            min_area=min_area,
            reserved=reserved,
            allow_body_fallback=True,
        )
        if pic_ph is not None:
            text_region = (
                ContentArea(
                    left=text_ph.left,
                    top=text_ph.top,
                    width=text_ph.width,
                    height=text_ph.height,
                )
                if text_ph is not None
                else _get_content_area(
                    layout,
                    slide_width,
                    slide_height,
                    preferred_types={
                        PP_PLACEHOLDER.BODY,
                        PP_PLACEHOLDER.OBJECT,
                        PP_PLACEHOLDER.SUBTITLE,
                    },
                )
            )
            return RegionMap(
                text_region=text_region,
                visual_region=ContentArea(
                    left=pic_ph.left,
                    top=pic_ph.top,
                    width=pic_ph.width,
                    height=pic_ph.height,
                ),
                layout_type="native",
            )

    # If content is chart/table/image and layout has a matching placeholder,
    # use that placeholder for the visual region and derive text elsewhere.
    if content_mix == ContentMix.TEXT_AND_CHART:
        if chart_placeholders and content_placeholders:
            chart_ph = chart_placeholders[0]
            text_ph = (
                _best_text_placeholder(
                    layout,
                    avoid=ContentArea(
                        left=chart_ph.left,
                        top=chart_ph.top,
                        width=chart_ph.width,
                        height=chart_ph.height,
                    ),
                )
                or content_placeholders[0]
            )
            return RegionMap(
                text_region=ContentArea(
                    left=text_ph.left,
                    top=text_ph.top,
                    width=text_ph.width,
                    height=text_ph.height,
                ),
                visual_region=ContentArea(
                    left=chart_ph.left,
                    top=chart_ph.top,
                    width=chart_ph.width,
                    height=chart_ph.height,
                ),
                layout_type="native",
            )
        text_ph = _largest_placeholder(
            layout,
            {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE},
            min_area=0,
        )
        reserved = list(layout_text_bounds)
        if text_ph is not None:
            reserved.append(
                ContentArea(
                    left=text_ph.left,
                    top=text_ph.top,
                    width=text_ph.width,
                    height=text_ph.height,
                )
            )
        chart_ph = _best_visual_placeholder(
            layout,
            {PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.TABLE},
            min_area=min_area,
            reserved=reserved,
            allow_body_fallback=True,
        )
        if chart_ph is not None:
            best_text = (
                _best_text_placeholder(
                    layout,
                    avoid=ContentArea(
                        left=chart_ph.left,
                        top=chart_ph.top,
                        width=chart_ph.width,
                        height=chart_ph.height,
                    ),
                )
                or text_ph
            )
            text_region = (
                ContentArea(
                    left=best_text.left,
                    top=best_text.top,
                    width=best_text.width,
                    height=best_text.height,
                )
                if best_text is not None
                else _get_content_area(
                    layout,
                    slide_width,
                    slide_height,
                    preferred_types={
                        PP_PLACEHOLDER.BODY,
                        PP_PLACEHOLDER.OBJECT,
                        PP_PLACEHOLDER.SUBTITLE,
                    },
                )
            )
            return RegionMap(
                text_region=text_region,
                visual_region=ContentArea(
                    left=chart_ph.left,
                    top=chart_ph.top,
                    width=chart_ph.width,
                    height=chart_ph.height,
                ),
                layout_type="native",
            )

    if content_mix == ContentMix.TEXT_AND_TABLE:
        if table_placeholders and content_placeholders:
            table_ph = table_placeholders[0]
            text_ph = (
                _best_text_placeholder(
                    layout,
                    avoid=ContentArea(
                        left=table_ph.left,
                        top=table_ph.top,
                        width=table_ph.width,
                        height=table_ph.height,
                    ),
                )
                or content_placeholders[0]
            )
            return RegionMap(
                text_region=ContentArea(
                    left=text_ph.left,
                    top=text_ph.top,
                    width=text_ph.width,
                    height=text_ph.height,
                ),
                visual_region=ContentArea(
                    left=table_ph.left,
                    top=table_ph.top,
                    width=table_ph.width,
                    height=table_ph.height,
                ),
                layout_type="native",
            )
        text_ph = _largest_placeholder(
            layout,
            {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE},
            min_area=0,
        )
        reserved = list(layout_text_bounds)
        if text_ph is not None:
            reserved.append(
                ContentArea(
                    left=text_ph.left,
                    top=text_ph.top,
                    width=text_ph.width,
                    height=text_ph.height,
                )
            )
        table_ph = _best_visual_placeholder(
            layout,
            {PP_PLACEHOLDER.TABLE, PP_PLACEHOLDER.CHART},
            min_area=min_area,
            reserved=reserved,
            allow_body_fallback=True,
        )
        if table_ph is not None:
            best_text = (
                _best_text_placeholder(
                    layout,
                    avoid=ContentArea(
                        left=table_ph.left,
                        top=table_ph.top,
                        width=table_ph.width,
                        height=table_ph.height,
                    ),
                )
                or text_ph
            )
            text_region = (
                ContentArea(
                    left=best_text.left,
                    top=best_text.top,
                    width=best_text.width,
                    height=best_text.height,
                )
                if best_text is not None
                else _get_content_area(
                    layout,
                    slide_width,
                    slide_height,
                    preferred_types={
                        PP_PLACEHOLDER.BODY,
                        PP_PLACEHOLDER.OBJECT,
                        PP_PLACEHOLDER.SUBTITLE,
                    },
                )
            )
            return RegionMap(
                text_region=text_region,
                visual_region=ContentArea(
                    left=table_ph.left,
                    top=table_ph.top,
                    width=table_ph.width,
                    height=table_ph.height,
                ),
                layout_type="native",
            )

    # Compute text ratio based on content volume
    text_ratio = _compute_text_ratio(content, content_mix)
    GAP = int(slide_width * 0.02)  # 2% gap between regions

    # Text + Table/Chart -> top/bottom split
    if content_mix in (
        ContentMix.TEXT_AND_TABLE,
        ContentMix.TEXT_AND_CHART,
        ContentMix.MIXED,
    ):
        text_height = int(content_area.height * text_ratio)
        visual_height = content_area.height - text_height - GAP
        return RegionMap(
            text_region=ContentArea(
                left=content_area.left,
                top=content_area.top,
                width=content_area.width,
                height=text_height,
            ),
            visual_region=ContentArea(
                left=content_area.left,
                top=content_area.top + text_height + GAP,
                width=content_area.width,
                height=visual_height,
            ),
            layout_type="split_vertical",
        )

    # Text + Image -> left/right split
    if content_mix in (ContentMix.TEXT_AND_IMAGE, ContentMix.TEXT_AND_GENERATED_IMAGE):
        text_width = int(content_area.width * text_ratio)
        visual_width = content_area.width - text_width - GAP
        return RegionMap(
            text_region=ContentArea(
                left=content_area.left,
                top=content_area.top,
                width=text_width,
                height=content_area.height,
            ),
            visual_region=ContentArea(
                left=content_area.left + text_width + GAP,
                top=content_area.top,
                width=visual_width,
                height=content_area.height,
            ),
            layout_type="split_horizontal",
        )

    # Fallback: use full area
    return RegionMap(
        text_region=content_area,
        visual_region=content_area,
        layout_type="full",
    )


def _ensure_text_on_top(slide) -> None:
    """Reorder shapes in the XML tree so text elements render above visual elements.

    This ensures text is never hidden behind charts, tables, or images.
    """
    spTree = slide.shapes._spTree

    # Separate shapes into visual (charts, tables, pictures) and text
    visual_elements = []
    text_elements = []
    other_elements = []

    for child in list(spTree):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("graphicFrame",):  # Tables, charts
            visual_elements.append(child)
        elif tag == "pic":  # Pictures
            visual_elements.append(child)
        elif tag == "sp":  # Shapes (including text boxes)
            # Check if it has a text frame with actual content
            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            txBody = child.find(".//{%s}txBody" % ns)
            if txBody is not None:
                # Check if it has actual text (not just empty)
                has_text = False
                for p_elem in txBody.findall("{%s}p" % ns):
                    for r_elem in p_elem.findall("{%s}r" % ns):
                        t_elem = r_elem.find("{%s}t" % ns)
                        if t_elem is not None and t_elem.text and t_elem.text.strip():
                            has_text = True
                            break
                    if has_text:
                        break
                if has_text:
                    text_elements.append(child)
                else:
                    other_elements.append(child)
            else:
                other_elements.append(child)
        else:
            other_elements.append(child)

    # Only reorder if we have both text and visual elements
    if not visual_elements or not text_elements:
        return

    # Remove all shape elements from spTree
    for elem in visual_elements + text_elements + other_elements:
        try:
            spTree.remove(elem)
        except ValueError:
            pass

    # Re-add in order: other (background shapes) -> visuals -> text (on top)
    for elem in other_elements:
        spTree.append(elem)
    for elem in visual_elements:
        spTree.append(elem)
    for elem in text_elements:
        spTree.append(elem)


# ---------------------------------------------------------------------------
# Template Application Functions (from agent_with_powerpoint_template.py)
# ---------------------------------------------------------------------------


def _find_best_layout(
    template_prs,
    slide_index: int,
    total_slides: int,
    content_mix: ContentMix = ContentMix.TEXT_ONLY,
    has_generated_image: bool = False,
):
    """Find the best matching layout from the template for a given slide position."""
    layouts = list(template_prs.slide_layouts)
    if not layouts:
        raise ValueError("Template has no slide layouts")

    is_title_slide = slide_index == 0
    is_last_slide = slide_index == total_slides - 1

    slide_width = template_prs.slide_width
    slide_height = template_prs.slide_height
    slide_area = int(slide_width) * int(slide_height)
    min_area = int(slide_area * 0.10)

    def _score_layout(layout):
        name = layout.name.lower()
        score = 0

        text_ph = _largest_placeholder(
            layout,
            types={PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE},
        )
        chart_ph = _largest_placeholder(layout, types={PP_PLACEHOLDER.CHART})
        table_ph = _largest_placeholder(layout, types={PP_PLACEHOLDER.TABLE})
        pic_ph = _largest_placeholder(layout, types={PP_PLACEHOLDER.PICTURE})

        if is_title_slide:
            has_title_ph = any(
                ph.placeholder_format.type
                in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
                for ph in layout.placeholders
            )
            if has_title_ph:
                score += 120
            if "title" in name:
                score += 20
            if text_ph:
                score += 10
            return score
        # Penalize title layouts for non-title slides
        has_title_ph = any(
            ph.placeholder_format.type
            in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
            for ph in layout.placeholders
        )
        if has_title_ph:
            score -= 80
        if "title" in name:
            score -= 40

        if is_last_slide:
            if "blank" in name or "closing" in name or "end" in name:
                score += 30

        if content_mix in (ContentMix.TEXT_AND_CHART,):
            if chart_ph:
                score += 120
                if _placeholder_area(chart_ph) >= min_area:
                    score += 30
            if text_ph:
                score += 30
        elif content_mix in (ContentMix.TEXT_AND_TABLE,):
            if table_ph:
                score += 120
                if _placeholder_area(table_ph) >= min_area:
                    score += 30
            if text_ph:
                score += 30
        elif content_mix in (
            ContentMix.TEXT_AND_IMAGE,
            ContentMix.TEXT_AND_GENERATED_IMAGE,
        ):
            if pic_ph:
                score += 90
                if _placeholder_area(pic_ph) >= min_area:
                    score += 20
            if text_ph:
                score += 30
        elif content_mix == ContentMix.MIXED:
            if chart_ph:
                score += 60
            if table_ph:
                score += 60
            if pic_ph:
                score += 40
            if text_ph:
                score += 30
        else:
            if text_ph:
                score += 60

        # Prefer visually rich layouts for text-only or light-content slides
        if content_mix in (ContentMix.TEXT_ONLY,):
            score += _layout_richness_score(layout)

        if "content" in name or "body" in name or "text" in name:
            score += 5
        if "two content" in name:
            score += 3

        # Penalize layouts with only tiny visual placeholders
        if content_mix in (
            ContentMix.TEXT_AND_CHART,
            ContentMix.TEXT_AND_TABLE,
            ContentMix.TEXT_AND_IMAGE,
            ContentMix.TEXT_AND_GENERATED_IMAGE,
            ContentMix.MIXED,
        ):
            best_visual = max(
                [
                    _placeholder_area(p)
                    for p in [chart_ph, table_ph, pic_ph]
                    if p is not None
                ]
                or [0]
            )
            if best_visual and best_visual < min_area:
                score -= 50

        return score

    best_layout = None
    best_score = None
    for layout in layouts:
        score = _score_layout(layout)
        if best_score is None or score > best_score:
            best_score = score
            best_layout = layout

    if best_layout is not None:
        return best_layout

    if len(layouts) > 1:
        return layouts[1]
    return layouts[0]


def _get_content_area(
    layout,
    slide_width: int,
    slide_height: int,
    preferred_types: set | None = None,
    min_area: int = 0,
) -> ContentArea:
    """Derive the safe content area from a template layout's placeholders.

    Strategy:
    0. If preferred_types is provided, choose the largest matching placeholder
       above min_area.
    1. Look for a body placeholder (idx=1) -- its position defines the content area.
    2. If no body placeholder, look for any placeholder with idx > 0.
    3. If no placeholders at all, use a default safe margin.

    Args:
        layout: A python-pptx SlideLayout object.
        slide_width: Presentation slide width in EMU.
        slide_height: Presentation slide height in EMU.

    Returns:
        ContentArea with the computed safe region.
    """
    # Try preferred placeholder types first
    if preferred_types:
        ph = _largest_placeholder(layout, preferred_types, min_area=min_area)
        if ph is not None:
            return ContentArea(
                left=ph.left, top=ph.top, width=ph.width, height=ph.height
            )

    # Prefer the largest text-oriented placeholder (body/object/subtitle)
    text_ph = _largest_placeholder(
        layout,
        types={
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.SUBTITLE,
            PP_PLACEHOLDER.TABLE,  # allow table placeholder as a usable content area
        },
        min_area=0,
    )
    if text_ph is not None:
        return ContentArea(
            left=text_ph.left,
            top=text_ph.top,
            width=text_ph.width,
            height=text_ph.height,
        )

    # Try body placeholder idx=1 explicitly if present
    for ph in layout.placeholders:
        if ph.placeholder_format.idx == 1:
            return ContentArea(
                left=ph.left, top=ph.top, width=ph.width, height=ph.height
            )

    # Otherwise pick the largest non-title placeholder (avoiding picture if possible)
    candidates = [
        ph
        for ph in layout.placeholders
        if ph.placeholder_format.idx > 0
        and ph.placeholder_format.type not in {PP_PLACEHOLDER.PICTURE}
    ]
    if candidates:
        ph = max(candidates, key=_placeholder_area)
        return ContentArea(left=ph.left, top=ph.top, width=ph.width, height=ph.height)

    # Default: safe margins (5% left, 25% top, 90% width, 65% height)
    return ContentArea(
        left=int(slide_width * 0.05),
        top=int(slide_height * 0.25),
        width=int(slide_width * 0.90),
        height=int(slide_height * 0.65),
    )


def _fit_to_area(img_width: int, img_height: int, area: ContentArea) -> tuple:
    """Scale dimensions to fit within an area while preserving aspect ratio.

    Args:
        img_width: Original width in EMU.
        img_height: Original height in EMU.
        area: Target ContentArea to fit within.

    Returns:
        Tuple of (left, top, width, height) in EMU, centered in the area.
    """
    if img_width <= 0 or img_height <= 0:
        return area.left, area.top, area.width, area.height

    aspect = img_width / img_height
    area_aspect = area.width / area.height

    if aspect > area_aspect:
        new_width = area.width
        new_height = int(new_width / aspect)
    else:
        new_height = area.height
        new_width = int(new_height * aspect)

    left = area.left + (area.width - new_width) // 2
    top = area.top + (area.height - new_height) // 2

    return left, top, new_width, new_height


def _add_picture_within_bounds(slide, image_bytes: bytes, bounds: ContentArea):
    """Insert image scaled to fit bounds while preserving aspect ratio."""
    try:
        from PIL import Image
    except ImportError:
        image_stream = BytesIO(image_bytes)
        slide.shapes.add_picture(
            image_stream, bounds.left, bounds.top, bounds.width, bounds.height
        )
        return

    try:
        with Image.open(BytesIO(image_bytes)) as im:
            img_w, img_h = im.size
        image_stream = BytesIO(image_bytes)
        left, top, width, height = _fit_to_area(
            int(img_w * 9525), int(img_h * 9525), bounds
        )
        slide.shapes.add_picture(image_stream, left, top, width, height)
    except Exception:
        image_stream = BytesIO(image_bytes)
        slide.shapes.add_picture(
            image_stream, bounds.left, bounds.top, bounds.width, bounds.height
        )


def _apply_standard_line_spacing(para, spacing_pct: int = 115) -> None:
    """Set standard line spacing on a paragraph using OOXML spcPct.

    Applies <a:lnSpc><a:spcPct val="..."/> to the paragraph's pPr element.
    spacing_pct=115 gives 115% line spacing (compact + readable).
    Only sets spacing if the paragraph pPr does not already have an explicit
    <a:lnSpc> element (respects template-inherited spacing).
    """
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    try:
        p_elem = para._p
        pPr = p_elem.find(ns_a + "pPr")
        if pPr is None:
            pPr = etree.SubElement(p_elem, ns_a + "pPr")
            p_elem.insert(0, pPr)
        # Only set if not already present
        existing_lnSpc = pPr.find(ns_a + "lnSpc")
        if existing_lnSpc is None:
            lnSpc = etree.SubElement(pPr, ns_a + "lnSpc")
            spcPct = etree.SubElement(lnSpc, ns_a + "spcPct")
            spcPct.set("val", str(spacing_pct * 1000))
    except Exception:
        pass


def _set_paragraph_alignment(para, alignment: str) -> None:
    """Set paragraph horizontal alignment via OOXML pPr algn attribute.

    alignment: 'l' (left), 'ctr' (center), 'r' (right), 'just' (justify).
    Only sets if pPr does not already have an explicit algn attribute.
    """
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    try:
        p_elem = para._p
        pPr = p_elem.find(ns_a + "pPr")
        if pPr is None:
            pPr = etree.SubElement(p_elem, ns_a + "pPr")
            p_elem.insert(0, pPr)
        if not pPr.get("algn"):
            pPr.set("algn", alignment)
    except Exception:
        pass


def _populate_placeholder_with_format(
    shape, texts, is_title=False, template_style: "TemplateStyle | None" = None
):
    """Populate a placeholder shape with text while preserving template formatting.

    Enables word wrap and attempts to auto-fit text to the placeholder bounds.
    Uses template_style for font family when available.

    Also applies:
    - Standard 115% line spacing to each paragraph (if not already set by template)
    - Explicit paragraph alignment (center for titles, left for body) when the
      template's reference paragraph has no explicit alignment set
    """
    if not shape.has_text_frame:
        return

    tf = shape.text_frame

    # Enable word wrap before anything else
    tf.word_wrap = True

    ref_para = tf.paragraphs[0] if tf.paragraphs else None
    ref_para_xml = None
    ref_run_xml = None

    if ref_para is not None:
        pPr = ref_para._p.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
        )
        if pPr is not None:
            ref_para_xml = copy.deepcopy(pPr)
        if ref_para.runs:
            rPr = ref_para.runs[0]._r.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
            )
            if rPr is not None:
                ref_run_xml = copy.deepcopy(rPr)

    tf.clear()

    if is_title:
        para = tf.paragraphs[0]
        if ref_para_xml is not None:
            para._p.insert(0, copy.deepcopy(ref_para_xml))
        # Default title to left alignment (template usually handles this,
        # but when ref_para_xml has no algn, set explicitly)
        _set_paragraph_alignment(para, "l")
        _apply_standard_line_spacing(para, spacing_pct=110)
        run = para.add_run()
        run.text = texts
        if ref_run_xml is not None:
            run._r.insert(0, copy.deepcopy(ref_run_xml))
    else:
        for i, (text, level) in enumerate(texts):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            if ref_para_xml is not None:
                new_pPr = copy.deepcopy(ref_para_xml)
                if level > 0:
                    new_pPr.set("lvl", str(level))
                para._p.insert(0, new_pPr)
            para.level = level
            # Body text: left-aligned, 115% line spacing
            _set_paragraph_alignment(para, "l")
            _apply_standard_line_spacing(para, spacing_pct=115)
            run = para.add_run()
            run.text = text
            if ref_run_xml is not None:
                run._r.insert(0, copy.deepcopy(ref_run_xml))

    # Determine font family from template style or use default
    # Titles use major_font (heading font), body uses minor_font (body font)
    font_family = "Calibri"
    if template_style:
        if is_title:
            font_family = (
                template_style.title_font_family
                or template_style.theme.major_font
                or "Calibri"
            )
        else:
            font_family = (
                template_style.body_font_family
                or template_style.theme.minor_font
                or "Calibri"
            )

    # Auto-fit text to placeholder bounds.
    # Compute line_count and safe_max before the try block so they remain
    # available for the fallback hardening path below.
    if is_title:
        line_count = 1
    else:
        line_count = len(texts) if isinstance(texts, list) else 1
    safe_max = _compute_max_font_size(
        ContentArea(
            left=shape.left, top=shape.top, width=shape.width, height=shape.height
        ),
        line_count,
        is_title=is_title,
    )
    # Body hard cap lowered to 16pt to match _compute_max_font_size defensive cap.
    hard_max = 28 if is_title else 16
    max_size = min(hard_max, safe_max)
    fit_text_succeeded = False
    try:
        tf.fit_text(font_family=font_family, max_size=max_size)
        fit_text_succeeded = True
    except Exception as e:
        # fit_text requires font files (e.g. Calibri); fall back to MSO_AUTO_SIZE.
        # NOTE: MSO_AUTO_SIZE only activates when PowerPoint renders the file —
        # the saved OOXML still holds the original font size attribute.
        # The hardening block below applies a conservative manual cap to guard
        # against headless environments where font files are absent.
        if VERBOSE:
            print("[VERBOSE] Exception suppressed: %s" % str(e))
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # --- Fix 3: Minimum font size enforcement ---
    # fit_text() can shrink text to 4pt or less to fit everything. Enforce a
    # readable minimum. If text overflows at the minimum size, truncate content.
    min_font_pt = _MIN_TITLE_FONT_PT if is_title else _MIN_BODY_FONT_PT
    min_font_sz = min_font_pt * 100  # OOXML hundredths of a point
    ns_a_mincheck = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    font_too_small = False
    for para in tf.paragraphs:
        for run in para.runs:
            rPr = run._r.find(ns_a_mincheck + "rPr")
            if rPr is not None:
                sz = rPr.get("sz")
                if sz and int(sz) < min_font_sz:
                    font_too_small = True
                    rPr.set("sz", str(min_font_sz))
    if font_too_small:
        print(
            "  [FONT GUARD] Font size was below %dpt minimum — enforced %dpt."
            " Text may need truncation." % (min_font_pt, min_font_pt)
        )
        # Ensure word_wrap is True — fit_text machinery may have reset it.
        tf.word_wrap = True

    # P1-5: Harden the MSO_AUTO_SIZE fallback — when fit_text() fails,
    # the OOXML font size is not updated. Apply a manual cap so viewers
    # that do not honour MSO_AUTO_SIZE still render readable text.
    if not fit_text_succeeded:
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        cap = max_size * 100  # OOXML sz units are hundredths of a point
        for para in tf.paragraphs:
            for run in para.runs:
                rPr = run._r.find(ns_a + "rPr")
                if rPr is not None:
                    sz = rPr.get("sz")
                    if sz and int(sz) > cap:
                        rPr.set("sz", str(int(cap)))

    # After fit_text has set the font size, re-apply captured run properties
    # to restore template formatting (colors, font family, etc.) that fit_text may have overwritten
    if ref_run_xml is not None:
        try:
            ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            for para in tf.paragraphs:
                for run in para.runs:
                    # Remove the rPr that fit_text created
                    existing_rPr = run._r.find(ns_a + "rPr")
                    sz_val = None
                    if existing_rPr is not None:
                        # Preserve the font size from fit_text
                        sz_val = existing_rPr.get("sz")
                        run._r.remove(existing_rPr)

                    # Clone the template's rPr and insert it
                    new_rPr = copy.deepcopy(ref_run_xml)

                    # If fit_text set a size and the template rPr has a size, use the SMALLER one
                    # to ensure text fits. If template rPr has no size, use fit_text's size.
                    if sz_val:
                        template_sz = new_rPr.get("sz")
                        if template_sz:
                            # Use the smaller of template size and fit_text size
                            fit_sz = int(sz_val)
                            tmpl_sz = int(template_sz)
                            new_rPr.set("sz", str(min(fit_sz, tmpl_sz)))
                        else:
                            new_rPr.set("sz", sz_val)

                    run._r.insert(0, new_rPr)
        except Exception as e:
            if VERBOSE:
                print(
                    "[VERBOSE] Re-applying template run properties failed: %s" % str(e)
                )

    # When ref_run_xml is None (placeholder inherits from layout/master),
    # fit_text() will have set an explicit font that may not match the theme.
    # Ensure the correct theme font is applied.
    if ref_run_xml is None and template_style:
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        if is_title:
            target_font = (
                template_style.title_font_family or template_style.theme.major_font
            )
        else:
            target_font = (
                template_style.body_font_family or template_style.theme.minor_font
            )
        if target_font:
            for para in tf.paragraphs:
                for run in para.runs:
                    rPr = run._r.find(ns_a + "rPr")
                    if rPr is not None:
                        # Ensure latin element exists and has typeface set
                        latin = rPr.find(ns_a + "latin")
                        if latin is None:
                            latin = __import__("lxml.etree", fromlist=["Element"]).Element(ns_a + "latin")
                            rPr.insert(0, latin)
                        latin.set("typeface", target_font)
                for run in para.runs:
                    rPr = run._r.find(ns_a + "rPr")
                    if rPr is not None:
                        # Update the latin font element
                        latin = rPr.find(ns_a + "latin")
                        if latin is not None:
                            latin.set("typeface", target_font)
                        else:
                            latin = etree.SubElement(rPr, ns_a + "latin")
                            latin.set("typeface", target_font)


def _transfer_tables(
    slide,
    tables,
    content_area: ContentArea,
    template_style: "TemplateStyle | None" = None,
):
    """Transfer extracted table data to a slide, repositioned to the content area.

    When template_style is provided, applies template-derived table styling
    instead of hardcoded font sizes.
    """
    TABLE_CELL_FONT_SIZE = Pt(10)
    TABLE_HEADER_FONT_SIZE = Pt(11)
    num_tables = len(tables)

    for t_idx, td in enumerate(tables):
        num_rows = len(td.rows)
        num_cols = len(td.rows[0]) if td.rows else 0
        if num_rows == 0 or num_cols == 0:
            continue

        # Stack multiple tables vertically within the content area
        if num_tables > 1:
            per_table_height = content_area.height // num_tables
            table_top = content_area.top + (t_idx * per_table_height)
            table_height = per_table_height
        else:
            table_top = content_area.top
            table_height = content_area.height

        table_shape = slide.shapes.add_table(
            num_rows,
            num_cols,
            content_area.left,
            table_top,
            content_area.width,
            table_height,
        )
        table = table_shape.table
        for r_idx, row_data in enumerate(td.rows):
            for c_idx, cell_text in enumerate(row_data):
                if c_idx < num_cols:
                    cell = table.cell(r_idx, c_idx)
                    cell.text = cell_text

        # Apply template-derived styling or fall back to hardcoded sizes
        if template_style:
            try:
                _apply_table_style(table_shape, template_style)
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Table style application failed: %s" % str(e))
                # Fallback to hardcoded styling on error
                for r_idx2, row_data2 in enumerate(td.rows):
                    for c_idx2 in range(min(len(row_data2), num_cols)):
                        cell = table.cell(r_idx2, c_idx2)
                        for para in cell.text_frame.paragraphs:
                            para.font.size = (
                                TABLE_HEADER_FONT_SIZE
                                if r_idx2 == 0
                                else TABLE_CELL_FONT_SIZE
                            )
                        cell.text_frame.word_wrap = True
        else:
            # Original hardcoded styling as fallback
            for r_idx2, row_data2 in enumerate(td.rows):
                for c_idx2 in range(min(len(row_data2), num_cols)):
                    cell = table.cell(r_idx2, c_idx2)
                    for para in cell.text_frame.paragraphs:
                        para.font.size = (
                            TABLE_HEADER_FONT_SIZE
                            if r_idx2 == 0
                            else TABLE_CELL_FONT_SIZE
                        )
                    cell.text_frame.word_wrap = True


def _transfer_images(slide, images, content_area: ContentArea):
    """Transfer extracted images to a slide, scaled to fit the content area."""
    for img in images:
        image_stream = BytesIO(img.blob)
        left, top, width, height = _fit_to_area(img.width, img.height, content_area)
        slide.shapes.add_picture(image_stream, left, top, width, height)


def _transfer_charts(
    slide,
    charts,
    content_area: ContentArea,
    template_style: "TemplateStyle | None" = None,
):
    """Transfer extracted chart data to a slide, sized to fill the content area.
    
    Pie charts are constrained to a perfect square within the bounds to prevent
    oval distortions and overlapping text bounds.

    When template_style is provided, applies template-derived chart styling
    (series colors, axis fonts, legend styling, etc.) after chart creation.
    """
    try:
        from pptx.chart.data import CategoryChartData
    except ImportError:
        return

    num_charts = len(charts)

    for c_idx, cd in enumerate(charts):
        try:
            chart_data = CategoryChartData()
            chart_data.categories = cd.categories
            for series_name, series_values in cd.series:
                clean_values = []
                for v in series_values:
                    if v is None:
                        clean_values.append(0)
                    elif isinstance(v, (int, float)):
                        clean_values.append(v)
                    else:
                        try:
                            clean_values.append(float(v))
                        except (ValueError, TypeError):
                            clean_values.append(0)
                chart_data.add_series(series_name, clean_values)

            # Stack multiple charts vertically within the content area
            if num_charts > 1:
                chart_height = content_area.height // num_charts
                chart_top = content_area.top + (c_idx * chart_height)
            else:
                chart_height = content_area.height
                chart_top = content_area.top

            # Scale pie charts to be square within the content area
            if cd.chart_type == XL_CHART_TYPE.PIE:
                size = min(content_area.width, chart_height)
                chart_width = size
                chart_height = size
                chart_left = content_area.left + (content_area.width - size) // 2
            else:
                chart_width = content_area.width
                chart_left = content_area.left
                
            chart_graphic_frame = slide.shapes.add_chart(
                cd.chart_type,
                chart_left,
                chart_top,
                chart_width,
                chart_height,
                chart_data,
            )

            # Apply template-derived chart styling
            if template_style:
                try:
                    _apply_chart_style(chart_graphic_frame, template_style)
                except Exception as e:
                    if VERBOSE:
                        print("[VERBOSE] Chart style application failed: %s" % str(e))
        except Exception as e:
            if VERBOSE:
                print("[VERBOSE] Exception suppressed: %s" % str(e))


def _rescale_shape_xml(
    shape_elem,
    src_width: int,
    src_height: int,
    target_area: "ContentArea",
) -> None:
    """Rescale a shape element's position/size from source slide dimensions
    to target content area bounds (modifies the element in-place).

    Finds xfrm/off (origin x,y) and xfrm/ext (size cx,cy) XML elements and
    rescales proportionally from source slide space into target_area bounds.
    Handles both plain shapes (<p:sp>) and group shapes (<p:grpSp>) because
    it iterates all xfrm elements within the subtree.

    Args:
        shape_elem: lxml Element for the shape (already deep-copied).
        src_width:  Width of the source slide in EMU.
        src_height: Height of the source slide in EMU.
        target_area: ContentArea defining the destination region in EMU.
    """
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _scale_xfrm(xfrm_elem):
        off = xfrm_elem.find("{%s}off" % ns_a)
        ext = xfrm_elem.find("{%s}ext" % ns_a)
        if off is not None:
            orig_x = int(off.get("x", 0))
            orig_y = int(off.get("y", 0))
            new_x = target_area.left + int((orig_x / src_width) * target_area.width)
            new_y = target_area.top + int((orig_y / src_height) * target_area.height)
            off.set("x", str(new_x))
            off.set("y", str(new_y))
        if ext is not None:
            orig_cx = int(ext.get("cx", 0))
            orig_cy = int(ext.get("cy", 0))
            new_cx = int((orig_cx / src_width) * target_area.width)
            new_cy = int((orig_cy / src_height) * target_area.height)
            ext.set("cx", str(max(1, new_cx)))
            ext.set("cy", str(max(1, new_cy)))

    for xfrm in shape_elem.iter("{%s}xfrm" % ns_a):
        try:
            _scale_xfrm(xfrm)
        except Exception:
            pass


def _apply_template_font_to_shape_xml(shape_elem, font_family: str) -> None:
    """Apply a template font family to all text runs inside a raw shape XML element.

    Updates <a:latin typeface="..."/> inside every <a:rPr> found in the element
    subtree. This ensures free-floating shapes transferred from Claude's raw slide
    use the template's body font instead of whatever Claude defaulted to.

    Skips if font_family is empty.
    """
    if not font_family:
        return
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    try:
        for rPr in shape_elem.iter(ns_a + "rPr"):
            latin = rPr.find(ns_a + "latin")
            if latin is None:
                latin = etree.SubElement(rPr, ns_a + "latin")
            latin.set("typeface", font_family)
    except Exception:
        pass


def _fix_overlapping_shapes(slide, slide_width: int, slide_height: int) -> bool:
    """Fix 4: Detect and resolve overlapping non-placeholder shapes.

    After shape transfer, multiple shapes may overlap because the LLM generated
    them for a full slide but they were rescaled into a smaller template region.

    Algorithm:
      1. Collect all non-placeholder shapes with text content.
      2. Sort by top position (natural reading order).
      3. For each pair of adjacent shapes, if they vertically overlap,
         push the lower shape down to just below the upper shape.
      4. Enforce minimum shape dimensions (prevent near-invisible shapes).
      5. If shapes would push off-slide, scale down proportionally.

    Returns True if any shapes were adjusted.
    """
    # Collect non-placeholder shapes with position
    movable_shapes = []
    # Fix 8: Also include placeholder shapes with visible text as "anchors"
    # for overlap detection. These won't be moved but free-floating shapes
    # that collide with them will be reflowed.
    anchor_shapes = []
    for shape in list(slide.shapes):
        try:
            # Only process shapes that have meaningful dimensions
            if shape.width <= 0 or shape.height <= 0:
                continue
            if shape.is_placeholder:
                # Include placeholders with visible text as anchors
                if getattr(shape, "has_text_frame", False) and shape.text.strip():
                    anchor_shapes.append(shape)
            else:
                movable_shapes.append(shape)
        except Exception:
            continue

    all_shapes = anchor_shapes + movable_shapes
    if len(all_shapes) < 2:
        return False

    modified = False

    # --- Phase 1: Enforce minimum shape dimensions ---
    min_width = int(slide_width * 0.08)   # 8% of slide width
    min_height = int(slide_height * 0.04)  # 4% of slide height

    for shape in movable_shapes:
        try:
            if shape.width < min_width:
                if VERBOSE:
                    print(
                        "  [OVERLAP FIX] Shape too narrow (%d EMU → %d EMU minimum)"
                        % (shape.width, min_width)
                    )
                shape.width = min_width
                modified = True
            if shape.height < min_height:
                if VERBOSE:
                    print(
                        "  [OVERLAP FIX] Shape too short (%d EMU → %d EMU minimum)"
                        % (shape.height, min_height)
                    )
                shape.height = min_height
                modified = True
        except Exception:
            continue

    # --- Phase 2: Sort shapes by vertical position and fix overlaps ---
    # Sort by top position for natural reading order
    # Fix 8: Use all_shapes for detection (anchors + movable) but only move non-placeholders
    movable_set = set(id(s) for s in movable_shapes)
    try:
        all_shapes.sort(key=lambda s: (s.top, s.left))
    except Exception:
        return modified

    MARGIN = int(slide_height * 0.01)  # 1% gap between shapes
    overlap_count = 0
    for i in range(len(all_shapes) - 1):
        try:
            upper = all_shapes[i]
            lower = all_shapes[i + 1]

            upper_bottom = upper.top + upper.height
            # Check for vertical overlap (shapes on similar x-axis range)
            # Fix 8: Use 30% horizontal overlap threshold instead of any overlap
            x_overlap_amount = min(upper.left + upper.width, lower.left + lower.width) - max(upper.left, lower.left)
            min_width = min(upper.width, lower.width)
            x_overlap = x_overlap_amount > 0 and (min_width <= 0 or x_overlap_amount / min_width > 0.30)

            if x_overlap and lower.top < upper_bottom + MARGIN:
                # Only move the lower shape if it's a movable (non-placeholder) shape
                if id(lower) in movable_set:
                    new_top = upper_bottom + MARGIN
                    if VERBOSE:
                        print(
                            "  [OVERLAP FIX] Reflowing shape from top=%d to top=%d "
                            "(was overlapping by %d EMU)"
                            % (lower.top, new_top, upper_bottom - lower.top)
                        )
                    lower.top = new_top
                    overlap_count += 1
                    modified = True
        except Exception:
            continue

    # --- Phase 3: If reflow pushed shapes off-slide, scale down ---
    if modified:
        try:
            max_bottom = max(
                s.top + s.height for s in movable_shapes
                if hasattr(s, "top") and hasattr(s, "height")
            )
            safe_bottom = int(slide_height * 0.92)  # 8% bottom margin
            if max_bottom > safe_bottom and max_bottom > 0:
                scale = safe_bottom / max_bottom
                if scale < 1.0 and scale > 0.5:  # Don't shrink more than 50%
                    for shape in movable_shapes:
                        try:
                            shape.top = int(shape.top * scale)
                            shape.height = int(shape.height * scale)
                        except Exception:
                            continue
                    if VERBOSE:
                        print(
                            "  [OVERLAP FIX] Scaled shapes down by %.0f%% to fit slide"
                            % ((1 - scale) * 100)
                        )
        except Exception:
            pass

    if overlap_count > 0:
        print(
            "  [OVERLAP FIX] Resolved %d overlapping shape(s) via vertical reflow."
            % overlap_count
        )

    return modified


def _transfer_shapes(
    slide,
    shapes_xml,
    src_width: int = 0,
    src_height: int = 0,
    target_area: "ContentArea | None" = None,
    template_style: "TemplateStyle | None" = None,
):
    """Transfer simple shapes by deep-copying their XML to the target slide.

    When src_width, src_height, and target_area are all provided (non-zero),
    rescales each shape's position/size proportionally from the source slide
    dimensions into target_area bounds.  This prevents shapes from Claude's
    default-dimension slide from landing at wrong coordinates or off-screen
    positions on the template.

    When template_style is provided, applies the template's body font to all
    text runs in the transferred shapes so free-floating text boxes use the
    correct typeface instead of Claude's default.

    Args:
        slide:          Target slide object.
        shapes_xml:     Iterable of lxml shape elements to transfer.
        src_width:      Source slide width in EMU (0 = skip rescaling).
        src_height:     Source slide height in EMU (0 = skip rescaling).
        target_area:    Destination region on the template slide (None = skip).
        template_style: Template styles for font family override (None = skip).
    """
    spTree = slide.shapes._spTree
    do_rescale = src_width > 0 and src_height > 0 and target_area is not None

    # Resolve body font from template (used for all free-floating text shapes)
    body_font = ""
    if template_style:
        body_font = (
            template_style.body_font_family or template_style.theme.minor_font or ""
        )

    for shape_elem in shapes_xml:
        cloned = copy.deepcopy(shape_elem)
        existing_ids = [
            int(sp.get("id", 0)) for sp in spTree.iter() if sp.get("id") is not None
        ]
        max_id = max(existing_ids) if existing_ids else 0
        for nv_elem in cloned.iter():
            if nv_elem.tag.endswith("}cNvPr"):
                max_id += 1
                nv_elem.set("id", str(max_id))
        if do_rescale:
            try:
                _rescale_shape_xml(cloned, src_width, src_height, target_area)
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Shape rescaling failed: %s" % str(e))
        # Gap 3: Apply template body font to all text runs in the shape
        if body_font:
            try:
                _apply_template_font_to_shape_xml(cloned, body_font)
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Shape font override failed: %s" % str(e))
        spTree.append(cloned)


def _clear_unused_placeholders(slide, populated_indices: set) -> None:
    """Remove unused placeholder XML elements from slide to prevent ghost text.

    This is the only reliable way to eliminate ALL types of ghost text:
    - "Click to add title"
    - "Click to add text"
    - "Click icon to add picture"
    - Content placeholder insertion icons (table, chart, image, etc.)

    Simply calling tf.clear() is insufficient for picture placeholders and
    content placeholders with embedded icons. Removing the XML element
    from the shape tree is the nuclear option that works for every case.
    
    Note: Preserves semantic footer placeholders (slide numbers, dates) if 
    retention logic allows it to survive.

    Args:
        slide: The pptx slide object.
        populated_indices: Set of placeholder idx values that were filled with content.
    """
    spTree = slide.shapes._spTree
    elements_to_remove = []

    # Snapshot the collection with list() to avoid proxy/iterator issues
    for shape in list(slide.placeholders):
        ph_idx = shape.placeholder_format.idx
        if ph_idx in populated_indices:
            continue  # This placeholder was populated, keep it

        # Don't remove shapes that have actual content (charts, tables, real images)
        try:
            if hasattr(shape, "has_chart") and shape.has_chart:
                continue
            if hasattr(shape, "has_table") and shape.has_table:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Only keep if this is an actual picture with image data, not an
                # empty picture placeholder that was never filled with an image.
                try:
                    _ = shape.image  # raises if no image data present
                    continue  # Has real image data — keep it
                except Exception:
                    pass  # No image data — fall through to removal
        except Exception as e:
            if VERBOSE:
                print("[VERBOSE] Exception suppressed: %s" % str(e))

        # Mark for removal -- removing from XML is the only reliable cleanup
        elements_to_remove.append(shape._element)

    for element in elements_to_remove:
        spTree.remove(element)


def _remove_empty_textboxes(slide) -> None:
    """Remove non-placeholder text boxes with no visible text content.

    Also removes placeholder text frames that survived _clear_unused_placeholders()
    because they were in populated_indices but had no actual text written to them
    (e.g. silent fill failures that leave "Click to add text" ghosts).
    """
    spTree = slide.shapes._spTree
    to_remove = []
    for shape in list(slide.shapes):
        try:
            if shape.is_placeholder:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text
            if text is None or text.strip() == "":
                to_remove.append(shape._element)
        except Exception:
            continue
    for el in to_remove:
        try:
            spTree.remove(el)
        except Exception:
            pass

    # Second pass: remove placeholder text frames with no actual text content.
    # These are empty "Click to add text" placeholders that survived the first
    # cleanup because their idx was in populated_indices (but content fill
    # failed silently) or because no removal guard was triggered.
    to_remove_ph = []
    for shape in list(slide.placeholders):
        try:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text
            if text is None or text.strip() == "":
                to_remove_ph.append(shape._element)
        except Exception:
            continue
    for el in to_remove_ph:
        try:
            spTree.remove(el)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Phase A+B: Semantic Content-Aware Layout Engine
# ---------------------------------------------------------------------------
# Classifies slide content into semantic types (sequential, comparative,
# metrics, hero) and routes to specialized, visually-rich layout builders.
# Every builder is fully exception-isolated: failures log a warning and
# return False so _populate_slide() always has a safe fallback path.
# ---------------------------------------------------------------------------

import re as _re


class SlideSemanticType(str, Enum):
    """Semantic layout category for a slide's content."""

    SEQUENTIAL = "sequential"    # Phase/Step/Year progression → Timeline or Chevron
    COMPARATIVE = "comparative"  # 2-4 parallel peer categories → Card Grid
    METRICS = "metrics"          # KPI-heavy slide → KPI Dashboard
    HERO = "hero"                # < 15 words total → Full-bleed Hero
    DEFAULT = "default"          # Fallback — standard bulleted layout


@dataclass
class SemanticSlideContext:
    """Result of semantic classification for a single slide."""

    semantic_type: SlideSemanticType = SlideSemanticType.DEFAULT
    confidence: float = 0.0
    extracted_metrics: list = field(default_factory=list)   # [{value, label, raw}]
    group_count: int = 0
    sequential_labels: list = field(default_factory=list)   # ["Phase 1", "Phase 2"…]
    sequential_descriptions: list = field(default_factory=list)
    storyboard_hint: str = "default"


# ── Metric extraction patterns (priority-ordered) ──────────────────────────

_METRIC_PATTERNS = [
    # Currency ranges: $169.2B → $395B  |  $6.9B → $41B
    (r'\$[\d,]+(?:\.\d+)?[BbMmKk]?\s*(?:→|->|to)\s*\$[\d,]+(?:\.\d+)?[BbMmKk]?', 0.95),
    # Single currency: $41B  |  $6.9M  |  $300K
    (r'\$[\d,]+(?:\.\d+)?[BbMmKk]', 0.90),
    # Percentage + optional label: 42.8% CAGR  |  18.4%  |  >100%
    (r'>?\s*[\d,]+(?:\.\d+)?%(?:\s+[A-Z]{2,8})?', 0.90),
    # X-multiples: 5x  |  3X  |  10×
    (r'\b[\d]+[xX×]\b', 0.85),
    # Time-bounded: 90-day  |  18-month  |  3 years
    (r'\b\d+[\-–\s]?\d*\s*(?:day|month|year)s?\b', 0.80),
    # Count + noun: 300+ deals  |  75K+ users  |  3-5 pilots
    (r'\b\d+[KkMmBb]?\+?\s+(?:deals?|users?|districts?|schools?|pilots?|customers?)', 0.80),
    # Financial acronyms with attached value
    (r'\b(?:NRR|ARR|MRR|CAGR|CAC|LTV|ROI|ACV)\b[^\n]{0,30}[\d,]+%?', 0.85),
    # Standalone large numbers: 75K, 300+
    (r'\b\d{1,3}(?:,\d{3})*[KkMmBb]?\+?\b', 0.65),
]

_SEQUENTIAL_PATTERN = _re.compile(
    r'^(?:Phase|Step|Stage|Year|Month|Quarter|Q|Wave|Round|Milestone|Sprint)\s*[\d]+',
    _re.IGNORECASE,
)
_SEQUENTIAL_ORDINAL = _re.compile(r'\b(?:1st|2nd|3rd|\d+th)\b', _re.IGNORECASE)


def _extract_kpi_metrics(body_paragraphs: list) -> list:
    """Extract quantifiable metrics from slide text using regex patterns.

    Args:
        body_paragraphs: List of (text, level) tuples from SlideContent.

    Returns:
        List of dicts with keys: value, label, raw_line, confidence.
        Returns [] on any error or when < 2 distinct metrics are found.
    """
    try:
        results = []
        seen_values: set = set()

        lines = [para[0] if isinstance(para, tuple) else str(para)
                 for para in body_paragraphs]

        for line_idx, line in enumerate(lines):
            if not line.strip():
                continue
            for pattern, base_confidence in _METRIC_PATTERNS:
                matches = _re.findall(pattern, line)
                for match in matches:
                    match_str = str(match).strip()
                    if len(match_str) < 2 or match_str in seen_values:
                        continue
                    seen_values.add(match_str)

                    # Derive a label: check same line remainder or previous line
                    remainder = line.replace(match_str, "").strip().strip("|:–-—").strip()
                    if len(remainder) > 3 and len(remainder) < 80:
                        label = remainder
                    elif line_idx > 0 and lines[line_idx - 1].strip():
                        label = lines[line_idx - 1].strip()
                    else:
                        label = "Key Metric"

                    results.append({
                        "value": match_str,
                        "label": label,
                        "raw_line": line,
                        "confidence": base_confidence,
                    })
                    break  # first matching pattern wins per line

        # De-duplicate and return only if we have at least 2
        return results if len(results) >= 2 else []
    except Exception as e:
        if VERBOSE:
            print("[SEMANTIC] _extract_kpi_metrics error: %s" % str(e))
        return []


def _classify_slide_semantics(
    content: "SlideContent",
    storyboard_hint: str = "default",
) -> SemanticSlideContext:
    """Classify slide content into a semantic layout type with a confidence score.

    Multi-signal detection:
      - SEQUENTIAL: regex for Phase/Step/Year labels in 2+ bullets
      - COMPARATIVE: 2-4 equal-depth peer groups detected
      - METRICS: _extract_kpi_metrics returns 2+ values
      - HERO: total visible word count < 15

    Storyboard hints from the LLM storyboard boost confidence by +0.2.
    If confidence < 0.5 for all types, returns DEFAULT.

    Never raises: all exceptions are caught and return DEFAULT context.
    """
    ctx = SemanticSlideContext(storyboard_hint=storyboard_hint)
    try:
        paragraphs = content.body_paragraphs or []

        # Bail early on empty / trivially small content
        all_text = " ".join(
            p[0] if isinstance(p, tuple) else str(p) for p in paragraphs
        )
        word_count = len(all_text.split())
        title_words = len((content.title or "").split())
        total_words = word_count + title_words

        if total_words < 3:
            return ctx  # DEFAULT

        # ── HERO check (< 15 words total, few paragraphs, no table/chart) ─
        if (
            total_words < 15
            and len(paragraphs) <= 2
            and not content.tables
            and not content.charts
            and not content.images
        ):
            hero_confidence = 0.7
            if storyboard_hint == "hero":
                hero_confidence = 0.9
            ctx.semantic_type = SlideSemanticType.HERO
            ctx.confidence = hero_confidence
            return ctx

        # ── METRICS check ───────────────────────────────────────────────────
        metrics = _extract_kpi_metrics(paragraphs)
        if metrics:
            metric_confidence = 0.65 + min(len(metrics) * 0.05, 0.20)
            if storyboard_hint in ("metrics", "revenue", "kpi"):
                metric_confidence = min(metric_confidence + 0.20, 0.95)
            if metric_confidence >= 0.60:
                ctx.semantic_type = SlideSemanticType.METRICS
                ctx.confidence = metric_confidence
                ctx.extracted_metrics = metrics
                return ctx

        # ── SEQUENTIAL check ────────────────────────────────────────────────
        seq_labels = []
        seq_descs = []
        level0_paras = [p for p in paragraphs if (p[1] if isinstance(p, tuple) else 0) == 0]
        for i, para in enumerate(level0_paras):
            text = para[0] if isinstance(para, tuple) else str(para)
            if _SEQUENTIAL_PATTERN.match(text) or _SEQUENTIAL_ORDINAL.search(text):
                seq_labels.append(text)
                # Collect sub-bullets as description
                sub_bullets = []
                for sub in paragraphs:
                    sub_text = sub[0] if isinstance(sub, tuple) else str(sub)
                    sub_level = sub[1] if isinstance(sub, tuple) else 0
                    if sub_level > 0 and sub_text not in seq_labels:
                        sub_bullets.append(sub_text)
                seq_descs.append(" ".join(sub_bullets[:3]).strip())

        if len(seq_labels) >= 2:
            seq_confidence = 0.60 + min(len(seq_labels) * 0.05, 0.20)
            if storyboard_hint == "sequential":
                seq_confidence = min(seq_confidence + 0.20, 0.95)
            # Visual suggestion keywords as secondary signal
            visual = getattr(content, "visual_suggestion", "") or ""
            if any(k in visual.lower() for k in ("timeline", "chevron", "process", "roadmap")):
                seq_confidence = min(seq_confidence + 0.10, 0.95)
            if seq_confidence >= 0.60:
                ctx.semantic_type = SlideSemanticType.SEQUENTIAL
                ctx.confidence = seq_confidence
                ctx.sequential_labels = seq_labels
                ctx.sequential_descriptions = seq_descs
                return ctx

        # ── COMPARATIVE check ───────────────────────────────────────────────
        # Detect 2-4 peer-level bold/titled groups: level-0 bullets with sub-items
        group_headers = [
            p[0] if isinstance(p, tuple) else str(p)
            for p in paragraphs
            if (p[1] if isinstance(p, tuple) else 0) == 0
        ]
        sub_item_groups = []
        for h in group_headers:
            subs = [
                p[0] if isinstance(p, tuple) else str(p)
                for p in paragraphs
                if (p[1] if isinstance(p, tuple) else 0) > 0
            ]
            sub_item_groups.append(subs)

        if 2 <= len(group_headers) <= 4 and all(sub_item_groups):
            # Check groups are roughly balanced (no group has 5x more items than another)
            lengths = [max(1, len(s)) for s in sub_item_groups]
            if max(lengths) / min(lengths) <= 4:
                comp_confidence = 0.60
                if storyboard_hint == "comparative":
                    comp_confidence = min(comp_confidence + 0.20, 0.90)
                if comp_confidence >= 0.60:
                    ctx.semantic_type = SlideSemanticType.COMPARATIVE
                    ctx.confidence = comp_confidence
                    ctx.group_count = len(group_headers)
                    ctx.sequential_labels = group_headers
                    ctx.sequential_descriptions = [
                        " ".join(s[:3]) for s in sub_item_groups
                    ]
                    return ctx

        # DEFAULT fallback
        return ctx

    except Exception as e:
        if VERBOSE:
            print("[SEMANTIC] _classify_slide_semantics error: %s — returning DEFAULT" % str(e))
        return SemanticSlideContext(storyboard_hint=storyboard_hint)


# ── Helper: resolve accent color safely ────────────────────────────────────

def _safe_accent(template_style: "TemplateStyle | None", idx: int = 0, fallback: str = "2E4057") -> str:
    """Return hex accent color by index; fall back to a professional dark blue."""
    try:
        if template_style and template_style.theme.accent_colors:
            colors = template_style.theme.accent_colors
            return colors[idx % len(colors)] or fallback
    except Exception:
        pass
    # No-template fallback palette (dark navy, blue, orange, teal)
    _fallback_palette = ["2E4057", "048A81", "EA6C00", "5C4033", "1B4332", "6A0572"]
    return _fallback_palette[idx % len(_fallback_palette)]


def _safe_font(template_style: "TemplateStyle | None", major: bool = True) -> str:
    try:
        if template_style:
            return template_style.theme.major_font if major else template_style.theme.minor_font
    except Exception:
        pass
    return "Calibri"


def _rgb_from_hex(hex_str: str):
    """Return pptx RGBColor from hex string, or RGBColor(0,0,0) on error."""
    from pptx.dml.color import RGBColor
    try:
        h = hex_str.strip().lstrip("#")
        if len(h) == 6:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        pass
    return RGBColor(0, 0, 0)


def _add_filled_rect(slide, left, top, width, height, hex_color: str):
    """Add a filled rectangle shape to a slide. Returns the shape or None."""
    from pptx.util import Emu
    from pptx.dml.color import RGBColor
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            left, top, width, height
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = _rgb_from_hex(hex_color)
        shape.line.fill.background()  # no border
        return shape
    except Exception as e:
        if VERBOSE:
            print("[SEMANTIC] _add_filled_rect error: %s" % str(e))
        return None


def _add_textbox_styled(
    slide,
    left, top, width, height,
    text: str,
    font_size_pt: int = 14,
    bold: bool = False,
    italic: bool = False,
    hex_color: str = "000000",
    font_name: str = "Calibri",
    wrap: bool = True,
    align_center: bool = False,
):
    """Add a styled text box. Returns shape or None on error."""
    from pptx.util import Pt as _Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    try:
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        if align_center:
            p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = _Pt(font_size_pt)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb_from_hex(hex_color)
        run.font.name = font_name
        return txBox
    except Exception as e:
        if VERBOSE:
            print("[SEMANTIC] _add_textbox_styled error: %s" % str(e))
        return None


# ── Layout Builder 1: Horizontal Timeline ──────────────────────────────────

def _build_horizontal_timeline(
    slide,
    labels: list,
    descriptions: list,
    area: "ContentArea",
    template_style: "TemplateStyle | None" = None,
) -> bool:
    """Build a horizontal timeline with N labeled boxes connected by arrows.

    Caps at 5 items. Auto-adjusts box width. Labels truncated at 22 chars.
    Falls back gracefully on any exception.

    Returns True on success, False on failure (caller uses DEFAULT path).
    """
    try:
        from pptx.util import Emu, Pt as _Pt
        from pptx.dml.color import RGBColor

        n = min(len(labels), 5)
        if n < 2:
            return False

        labels = labels[:n]
        descriptions = (descriptions[:n] + [""] * n)[:n]

        # Layout geometry
        ARROW_W = int(area.width * 0.03)
        total_arrow = ARROW_W * (n - 1)
        box_w = max(int((area.width - total_arrow) / n), int(Inches(1.5)))
        box_h = int(area.height * 0.55)
        box_top = area.top + int(area.height * 0.10)
        desc_top = box_top + box_h + int(Inches(0.15))
        desc_h = area.top + area.height - desc_top - int(Inches(0.05))
        desc_h = max(desc_h, int(Inches(0.3)))

        font_name = _safe_font(template_style, major=True)
        body_font = _safe_font(template_style, major=False)

        for i in range(n):
            color_hex = _safe_accent(template_style, i)
            x = area.left + i * (box_w + ARROW_W)

            # Colored box
            rect = _add_filled_rect(slide, x, box_top, box_w, box_h, color_hex)

            # Label inside box
            label_text = str(labels[i])[:22]
            _add_textbox_styled(
                slide, x + int(Inches(0.05)), box_top + int(Inches(0.08)),
                box_w - int(Inches(0.10)), box_h - int(Inches(0.15)),
                label_text,
                font_size_pt=12, bold=True, hex_color="FFFFFF",
                font_name=font_name, align_center=True,
            )

            # Description below the box
            desc = str(descriptions[i])[:120] if descriptions[i] else ""
            if desc:
                _add_textbox_styled(
                    slide, x, desc_top, box_w, desc_h,
                    desc,
                    font_size_pt=9, bold=False, hex_color="333333",
                    font_name=body_font, align_center=True,
                )

            # Arrow connector between boxes (not after last box)
            if i < n - 1:
                arrow_x = x + box_w
                try:
                    arrow = slide.shapes.add_shape(
                        13,  # MSO_SHAPE.RIGHT_ARROW
                        arrow_x,
                        box_top + int(box_h * 0.35),
                        ARROW_W,
                        int(box_h * 0.30),
                    )
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = _rgb_from_hex(
                        _safe_accent(template_style, i + 1)
                    )
                    arrow.line.fill.background()
                except Exception:
                    pass  # Arrow is cosmetic — skip if it fails

        print("[SEMANTIC] Built HORIZONTAL TIMELINE: %d items" % n)
        return True

    except Exception as e:
        print("[SEMANTIC] _build_horizontal_timeline failed: %s — falling back" % str(e))
        return False


# ── Layout Builder 2: Chevron Process ─────────────────────────────────────

def _build_chevron_process(
    slide,
    labels: list,
    descriptions: list,
    area: "ContentArea",
    template_style: "TemplateStyle | None" = None,
) -> bool:
    """Build a horizontal chevron/pentagon process flow (up to 4 items).

    Uses PPTX pentagon shape (arrow chevron). Falls back to
    _build_horizontal_timeline if shape insertion fails.

    Returns True on success, False on failure.
    """
    try:
        from pptx.util import Emu, Pt as _Pt

        n = min(len(labels), 4)
        if n < 2:
            return False

        labels = labels[:n]
        descriptions = (descriptions[:n] + [""] * n)[:n]

        # Each pentagon overlaps the previous by 10%
        overlap = int(area.width * 0.04)
        shape_w = int((area.width + overlap * (n - 1)) / n)
        shape_h = int(area.height * 0.50)
        shape_top = area.top + int(area.height * 0.08)
        desc_top = shape_top + shape_h + int(Inches(0.12))
        desc_h = max(area.top + area.height - desc_top - int(Inches(0.05)), int(Inches(0.3)))

        font_name = _safe_font(template_style, major=True)
        body_font = _safe_font(template_style, major=False)

        for i in range(n):
            color_hex = _safe_accent(template_style, i)
            x = area.left + i * (shape_w - overlap)
            try:
                # PENTAGON = shape type 15 in python-pptx
                chev = slide.shapes.add_shape(15, x, shape_top, shape_w, shape_h)
                chev.fill.solid()
                chev.fill.fore_color.rgb = _rgb_from_hex(color_hex)
                chev.line.fill.background()
            except Exception:
                # Pentagon not available — fall through to rectangle timeline
                return _build_horizontal_timeline(slide, labels, descriptions, area, template_style)

            label_text = str(labels[i])[:20]
            _add_textbox_styled(
                slide,
                x + int(Inches(0.06)), shape_top + int(shape_h * 0.25),
                shape_w - int(Inches(0.12)), int(shape_h * 0.5),
                label_text,
                font_size_pt=11, bold=True, hex_color="FFFFFF",
                font_name=font_name, align_center=True,
            )

            desc = str(descriptions[i])[:100] if descriptions[i] else ""
            if desc:
                _add_textbox_styled(
                    slide, x, desc_top, shape_w - overlap, desc_h,
                    desc,
                    font_size_pt=9, bold=False, hex_color="333333",
                    font_name=body_font, align_center=True,
                )

        print("[SEMANTIC] Built CHEVRON PROCESS: %d items" % n)
        return True

    except Exception as e:
        print("[SEMANTIC] _build_chevron_process failed: %s — trying timeline fallback" % str(e))
        return _build_horizontal_timeline(slide, labels, descriptions, area, template_style)


# ── Layout Builder 3: Multi-Column Card Grid ──────────────────────────────

def _build_card_grid(
    slide,
    labels: list,
    descriptions: list,
    area: "ContentArea",
    template_style: "TemplateStyle | None" = None,
) -> bool:
    """Build a 2-4 column card grid with colored header bands and body text.

    Auto-selects column count from len(labels). Truncates body at 200 chars.
    Falls back gracefully on any exception.

    Returns True on success, False on failure.
    """
    try:
        from pptx.util import Emu, Pt as _Pt

        n = max(2, min(len(labels), 4))
        labels = (list(labels) + [""] * n)[:n]
        descriptions = (list(descriptions) + [""] * n)[:n]

        CARD_GAP = int(Inches(0.12))
        card_w = max(int((area.width - CARD_GAP * (n - 1)) / n), int(Inches(1.4)))
        card_h = int(area.height * 0.80)
        card_top = area.top + int(area.height * 0.05)
        HEADER_H = int(card_h * 0.25)
        body_top_offset = HEADER_H + int(Inches(0.06))
        body_h = card_h - HEADER_H - int(Inches(0.08))

        font_name = _safe_font(template_style, major=True)
        body_font = _safe_font(template_style, major=False)
        # Light gray for card body background
        CARD_BG = "F5F5F5"
        dk_color = "FFFFFF"  # text on colored header

        for i in range(n):
            accent_hex = _safe_accent(template_style, i)
            x = area.left + i * (card_w + CARD_GAP)

            # Card body background (light gray)
            _add_filled_rect(slide, x, card_top, card_w, card_h, CARD_BG)

            # Colored header band
            _add_filled_rect(slide, x, card_top, card_w, HEADER_H, accent_hex)

            # Header label
            header_text = str(labels[i])[:30]
            _add_textbox_styled(
                slide,
                x + int(Inches(0.05)), card_top + int(Inches(0.04)),
                card_w - int(Inches(0.10)), HEADER_H - int(Inches(0.08)),
                header_text,
                font_size_pt=12, bold=True, hex_color=dk_color,
                font_name=font_name, align_center=False,
            )

            # Body description
            body_text = str(descriptions[i])[:200]
            if body_text:
                _add_textbox_styled(
                    slide,
                    x + int(Inches(0.05)),
                    card_top + body_top_offset,
                    card_w - int(Inches(0.10)),
                    body_h,
                    body_text,
                    font_size_pt=10, bold=False, hex_color="222222",
                    font_name=body_font, align_center=False,
                )

        print("[SEMANTIC] Built CARD GRID: %d columns" % n)
        return True

    except Exception as e:
        print("[SEMANTIC] _build_card_grid failed: %s — falling back" % str(e))
        return False


# ── Layout Builder 4: Hero Layout ─────────────────────────────────────────

def _build_hero_layout(
    slide,
    title: str,
    subtitle: str,
    area: "ContentArea",
    template_style: "TemplateStyle | None" = None,
) -> bool:
    """Build a centered hero layout: large title + optional subtitle.

    Adds a full-content-area color-tinted background rectangle.
    Requires area height > 1.5 inches; degrades to large textbox otherwise.

    Returns True on success, False on failure.
    """
    try:
        from pptx.util import Emu, Pt as _Pt
        from pptx.enum.text import PP_ALIGN

        MIN_HEIGHT = int(Inches(1.5))
        if area.height < MIN_HEIGHT:
            # Degrade: just a big centered textbox
            return bool(
                _add_textbox_styled(
                    slide,
                    area.left, area.top, area.width, area.height,
                    title, font_size_pt=32, bold=True, hex_color="111111",
                    font_name=_safe_font(template_style, major=True), align_center=True,
                )
            )

        # Background accent rectangle (subtle tint at ~25% of area height)
        accent_hex = _safe_accent(template_style, 0)
        _add_filled_rect(slide, area.left, area.top, area.width, area.height, accent_hex)

        # Title — oversized, centered
        title_h = int(area.height * 0.55)
        _add_textbox_styled(
            slide,
            area.left + int(Inches(0.15)),
            area.top + int(area.height * 0.15),
            area.width - int(Inches(0.30)),
            title_h,
            title,
            font_size_pt=40, bold=True, hex_color="FFFFFF",
            font_name=_safe_font(template_style, major=True), align_center=True,
        )

        # Subtitle — italic, smaller
        if subtitle and subtitle.strip():
            sub_top = area.top + int(area.height * 0.15) + title_h + int(Inches(0.08))
            sub_h = area.top + area.height - sub_top - int(Inches(0.08))
            sub_h = max(sub_h, int(Inches(0.3)))
            _add_textbox_styled(
                slide,
                area.left + int(Inches(0.15)),
                sub_top,
                area.width - int(Inches(0.30)),
                sub_h,
                subtitle,
                font_size_pt=18, bold=False, italic=True, hex_color="F0F0F0",
                font_name=_safe_font(template_style, major=False), align_center=True,
            )

        print("[SEMANTIC] Built HERO LAYOUT")
        return True

    except Exception as e:
        print("[SEMANTIC] _build_hero_layout failed: %s — falling back" % str(e))
        return False


# ── Layout Builder 5: KPI Dashboard ───────────────────────────────────────

def _build_kpi_dashboard(
    slide,
    metrics: list,
    context_text: str,
    area: "ContentArea",
    template_style: "TemplateStyle | None" = None,
) -> bool:
    """Build a KPI dashboard: horizontal metrics row + optional prose below.

    Each KPI box shows a large bold value and a smaller sub-label.
    Caps at 4 KPIs. If 0 metrics, returns False immediately.

    Args:
        metrics: List of dicts with keys: value, label.
        context_text: Optional narrative prose to show below the KPI row.

    Returns True on success, False on failure.
    """
    try:
        from pptx.util import Emu, Pt as _Pt

        if not metrics:
            return False

        n = min(len(metrics), 4)
        metrics = metrics[:n]

        KPI_GAP = int(Inches(0.12))
        kpi_w = max(int((area.width - KPI_GAP * (n - 1)) / n), int(Inches(1.2)))
        # KPI row takes 55-65% of area height; rest for context
        HAS_CONTEXT = bool(context_text and context_text.strip())
        kpi_h = int(area.height * (0.55 if HAS_CONTEXT else 0.75))
        kpi_top = area.top + int(area.height * 0.05)
        VALUE_H = int(kpi_h * 0.62)
        LABEL_H = kpi_h - VALUE_H - int(Inches(0.04))

        font_name = _safe_font(template_style, major=True)
        body_font = _safe_font(template_style, major=False)

        # Determine value font size: target ~36pt, floor 24pt
        value_font_pt = max(24, min(36, int(36 * (area.width / (n * int(Inches(2.0)))))))

        for i, metric in enumerate(metrics):
            accent_hex = _safe_accent(template_style, i)
            x = area.left + i * (kpi_w + KPI_GAP)

            # Subtle accent-tinted background for KPI box
            _add_filled_rect(slide, x, kpi_top, kpi_w, kpi_h, "F8F8F8")
            # Left accent border (3px wide)
            BORDER_W = int(Inches(0.04))
            _add_filled_rect(slide, x, kpi_top, BORDER_W, kpi_h, accent_hex)

            # Metric value (large, bold, accent color)
            value_str = str(metric.get("value", ""))[:30]
            _add_textbox_styled(
                slide,
                x + BORDER_W + int(Inches(0.05)),
                kpi_top + int(Inches(0.05)),
                kpi_w - BORDER_W - int(Inches(0.08)),
                VALUE_H,
                value_str,
                font_size_pt=value_font_pt, bold=True, hex_color=accent_hex,
                font_name=font_name, align_center=False,
            )

            # Sub-label (small, muted)
            label_str = str(metric.get("label", ""))[:60]
            if label_str:
                _add_textbox_styled(
                    slide,
                    x + BORDER_W + int(Inches(0.05)),
                    kpi_top + VALUE_H + int(Inches(0.02)),
                    kpi_w - BORDER_W - int(Inches(0.08)),
                    LABEL_H,
                    label_str,
                    font_size_pt=9, bold=False, hex_color="666666",
                    font_name=body_font, align_center=False,
                )

        # Context narrative below the KPI row
        if HAS_CONTEXT:
            ctx_top = kpi_top + kpi_h + int(Inches(0.15))
            ctx_h = area.top + area.height - ctx_top - int(Inches(0.05))
            ctx_h = max(ctx_h, int(Inches(0.3)))
            _add_textbox_styled(
                slide,
                area.left, ctx_top, area.width, ctx_h,
                context_text[:400],
                font_size_pt=10, bold=False, italic=True, hex_color="444444",
                font_name=body_font, align_center=False,
            )

        print("[SEMANTIC] Built KPI DASHBOARD: %d metric(s)" % n)
        return True

    except Exception as e:
        print("[SEMANTIC] _build_kpi_dashboard failed: %s — falling back" % str(e))
        return False


# ── Density Reduction (medium-confidence path) ────────────────────────────

def _apply_density_reduction(
    slide,
    content: "SlideContent",
    area: "ContentArea",
    template_style: "TemplateStyle | None" = None,
) -> None:
    """Reduce bullet density for medium-confidence slides.

    Applied when classifier confidence is 0.4-0.6 (clear signal but not
    confident enough for a full specialized layout). Adds a left-border
    accent strip to each text block and truncates overlong paragraph lists.
    This function never raises — errors are silently suppressed.
    """
    try:
        paragraphs = content.body_paragraphs or []
        if len(paragraphs) <= 5:
            return  # Density is already acceptable

        # Add a left-border accent bar as visual rhythm element
        accent_hex = _safe_accent(template_style, 0)
        BORDER_W = int(Inches(0.04))
        BORDER_H = int(area.height * 0.60)
        try:
            _add_filled_rect(
                slide,
                area.left,
                area.top + int(area.height * 0.15),
                BORDER_W,
                BORDER_H,
                accent_hex,
            )
        except Exception:
            pass

    except Exception as e:
        if VERBOSE:
            print("[SEMANTIC] _apply_density_reduction error: %s" % str(e))



# ── Template Native Layout Matching ───────────────────────────────────────

def _find_matching_template_layout(
    semantic_type: "SlideSemanticType",
    slide_layout,
) -> "dict | None":
    """Check if the template has a native layout matching the semantic type.

    Scans the slide_layout's placeholders and SmartArt elements for visual
    structures (e.g. a native comparison layout, process flow, etc.) that
    correspond to the given semantic classification.

    Args:
        semantic_type: The classified semantic type for this slide.
        slide_layout: The python-pptx slide layout object being used.

    Returns:
        A dict describing the native match (e.g. {'type': 'smartart', 'idx': 3}),
        or None if no native layout matches — in which case the custom semantic
        builder should be used instead.

    Note:
        Current implementation returns None for all cases (always falls
        through to the custom builder). This is a planned extension point for
        templates with rich native SmartArt layouts.
    """
    try:
        # Parse layout name for hints
        layout_name = (slide_layout.name or "").lower()

        # Check for native PowerPoint SmartArt or comparison layouts
        _LAYOUT_KEYWORDS = {
            "sequential": ("process", "timeline", "chevron", "step"),
            "comparative": ("comparison", "two content", "content with",),
            "metrics": ("dashboard", "kpi", "metric",),
            "hero": ("title only", "blank", "section header",),
        }

        stype_str = semantic_type.value if hasattr(semantic_type, "value") else str(semantic_type)
        keywords = _LAYOUT_KEYWORDS.get(stype_str, ())

        for kw in keywords:
            if kw in layout_name:
                return {"type": "native_layout", "layout_name": layout_name, "keyword": kw}

        return None
    except Exception:
        return None


# ── Semantic Router ────────────────────────────────────────────────────────

def _route_to_semantic_builder(
    slide,
    content: "SlideContent",
    ctx: SemanticSlideContext,
    area: "ContentArea",
    template_style: "TemplateStyle | None" = None,
) -> bool:
    """Dispatch to the correct layout builder based on semantic context.

    Returns True if a specialized layout was successfully built, False
    if the builder failed or was not applicable (caller uses DEFAULT path).

    This function never raises.
    """
    try:
        stype = ctx.semantic_type

        if stype == SlideSemanticType.SEQUENTIAL:
            # Try chevron first (≤4 items), fall back to timeline
            if len(ctx.sequential_labels) <= 4:
                return _build_chevron_process(
                    slide, ctx.sequential_labels, ctx.sequential_descriptions, area, template_style
                )
            return _build_horizontal_timeline(
                slide, ctx.sequential_labels, ctx.sequential_descriptions, area, template_style
            )

        elif stype == SlideSemanticType.COMPARATIVE:
            return _build_card_grid(
                slide, ctx.sequential_labels, ctx.sequential_descriptions, area, template_style
            )

        elif stype == SlideSemanticType.METRICS:
            # Gather non-metric text as context prose
            metric_raws = {m.get("raw_line", "") for m in ctx.extracted_metrics}
            context_lines = [
                p[0] if isinstance(p, tuple) else str(p)
                for p in (content.body_paragraphs or [])
                if (p[0] if isinstance(p, tuple) else str(p)) not in metric_raws
            ]
            context_text = " ".join(context_lines[:3]).strip()
            return _build_kpi_dashboard(
                slide, ctx.extracted_metrics, context_text, area, template_style
            )

        elif stype == SlideSemanticType.HERO:
            title = content.title or ""
            subtitle = content.subtitle or ""
            if not subtitle and content.body_paragraphs:
                first = content.body_paragraphs[0]
                subtitle = first[0] if isinstance(first, tuple) else str(first)
            return _build_hero_layout(slide, title, subtitle, area, template_style)

    except Exception as e:
        print("[SEMANTIC] _route_to_semantic_builder error: %s — falling back" % str(e))

    return False


def _populate_footer_placeholders(
    slide,
    populated_indices: set,
    footer_text: str = "",
    show_slide_number: bool = False,
    date_text: str = "",
) -> None:
    """Populate footer placeholders to prevent them from being stripped.

    PPTX footer placeholder indices:
        10 = date/time
        11 = footer text
        12 = slide number

    If a footer placeholder is found and the corresponding value is provided,
    the placeholder is populated and its index is added to populated_indices so
    _clear_unused_placeholders does not remove it.

    If no footer values are provided, nothing is changed and the placeholders
    are removed as normal by _clear_unused_placeholders.

    Args:
        slide:             Target slide object.
        populated_indices: Set of placeholder idx values already populated.
                           Modified in-place.
        footer_text:       Text to insert into the footer text placeholder (idx=11).
        show_slide_number: If True, keep the slide number placeholder (idx=12).
        date_text:         Text to insert into the date placeholder (idx=10).
    """
    for shape in list(slide.placeholders):
        ph_idx = shape.placeholder_format.idx
        try:
            if ph_idx == 10 and shape.has_text_frame:
                if date_text:
                    shape.text_frame.text = date_text
                # Always preserve date placeholder
                populated_indices.add(ph_idx)
            elif ph_idx == 11 and shape.has_text_frame:
                if footer_text:
                    shape.text_frame.text = footer_text
                # Preserve existing footer text placeholder
                populated_indices.add(ph_idx)
            elif ph_idx == 12:
                # Slide number is auto-rendered by PowerPoint; keep the placeholder
                # element so the renderer can fill it — do not set text manually.
                populated_indices.add(ph_idx)
        except Exception:
            continue

ICON_KEYWORD_MAP = {
    ("growth", "increase", "scale", "expand"): "UP_ARROW",
    ("revenue", "money", "cost", "price"):     "FLOWCHART_CURRENCY",
    ("user", "customer", "student", "person"): "OVAL",
    ("target", "goal", "milestone"):           "PENTAGON",
    ("data", "analytics", "metrics"):          "CHART",
}


def _populate_slide(
    new_slide,
    content: SlideContent,
    slide_width: int,
    slide_height: int,
    generated_image_bytes: bytes | None = None,
    template_style: "TemplateStyle | None" = None,
    src_slide_width: int = 0,
    src_slide_height: int = 0,
    footer_text: str = "",
    date_text: str = "",
    show_slide_number: bool = False,
):
    """Transfer all content into a new slide using template-aware positioning.

    Handles text placeholders, tables, charts, images, and shapes. When
    ``generated_image_bytes`` is provided and the slide layout contains a
    picture placeholder, the image is inserted into that placeholder rather
    than being added as a free-floating picture.

    Uses a RegionMap to separate text and visual content into non-overlapping
    regions, preventing layout collisions between text and tables/charts/images.

    Args:
        new_slide: The new slide created from a template layout.
        content: Extracted SlideContent from the generated slide.
        slide_width: Presentation slide width in EMU.
        slide_height: Presentation slide height in EMU.
        generated_image_bytes: Optional raw image bytes (e.g. from NanoBanana)
            to insert into picture placeholders.
        template_style: Optional extracted template styling for font/color matching.
        src_slide_width: Width of the source (Claude-generated) slide in EMU.
            When non-zero, shapes are rescaled from source dimensions to the
            template's content region to prevent off-screen placement.
        src_slide_height: Height of the source slide in EMU (see src_slide_width).
        footer_text: Text for the footer text placeholder (idx=11). Empty = remove.
        date_text: Text for the date/time placeholder (idx=10). Empty = remove.
        show_slide_number: If True, keep the slide number placeholder (idx=12).
    """
    # Classify content mix and compute region map
    content_mix = _classify_content_mix(
        content, has_generated_image=generated_image_bytes is not None
    )
    region_map = _compute_region_map(
        new_slide.slide_layout, content_mix, slide_width, slide_height, content
    )

    if VERBOSE:
        summary = _layout_placeholder_summary(new_slide.slide_layout)
        print(
            "[VERBOSE] Layout '%s' placeholders: %s"
            % (new_slide.slide_layout.name, summary)
        )
        print(
            "[VERBOSE] Region map: layout_type=%s text=(%d,%d,%d,%d) visual=(%d,%d,%d,%d)"
            % (
                region_map.layout_type,
                region_map.text_region.left,
                region_map.text_region.top,
                region_map.text_region.width,
                region_map.text_region.height,
                region_map.visual_region.left,
                region_map.visual_region.top,
                region_map.visual_region.width,
                region_map.visual_region.height,
            )
        )

    # -----------------------------------------------------------------------
    # Phase A: Semantic Classifier & Layout Routing
    # -----------------------------------------------------------------------
    semantic_ctx = _classify_slide_semantics(content, getattr(content, "visual_suggestion", "default"))
    
    handled_by_semantic = False
    if semantic_ctx.confidence >= 0.60:
        if VERBOSE:
            print("[SEMANTIC] Routing to %s builder (confidence: %.2f)" % (
                semantic_ctx.semantic_type, semantic_ctx.confidence
            ))
        handled_by_semantic = _route_to_semantic_builder(
            new_slide, content, semantic_ctx, region_map.text_region, template_style
        )
    elif 0.40 <= semantic_ctx.confidence < 0.60:
        if VERBOSE:
            print("[SEMANTIC] Applying density reduction (confidence: %.2f)" % semantic_ctx.confidence)
        _apply_density_reduction(new_slide, content, region_map.text_region, template_style)

    # If the semantic builder successfully handled the body content, suppress the 
    # default body text population to prevent messy overlaps. 
    # Titles/subtitles are still populated into their native placeholders.
    if handled_by_semantic:
        content.body_paragraphs = []
        content.text_box_paragraphs = []
    # -----------------------------------------------------------------------

    # Track populated placeholder indices for cleanup
    populated_indices: set[int] = set()

    title_placed = False
    body_placed = False
    body_paragraphs = list(content.body_paragraphs)
    if content.text_box_paragraphs and content_mix != ContentMix.TEXT_ONLY:
        body_paragraphs.extend(content.text_box_paragraphs)

    # Prefer placeholder matching the computed text region when available
    preferred_text_ph = None
    if region_map.layout_type == "native":
        for shape in new_slide.placeholders:
            try:
                if (
                    shape.has_text_frame
                    and not _is_visual_placeholder_type(shape.placeholder_format.type)
                    and shape.left == region_map.text_region.left
                    and shape.top == region_map.text_region.top
                    and shape.width == region_map.text_region.width
                    and shape.height == region_map.text_region.height
                ):
                    preferred_text_ph = shape
                    break
            except Exception:
                continue

    # Resize body placeholder to text_region when using split layout
    if region_map.layout_type not in ("full", "native"):
        for shape in new_slide.placeholders:
            ph_idx = shape.placeholder_format.idx
            if ph_idx == 1 and shape.has_text_frame:
                shape.left = region_map.text_region.left
                shape.top = region_map.text_region.top
                shape.width = region_map.text_region.width
                shape.height = region_map.text_region.height

    for shape in new_slide.placeholders:
        ph_idx = shape.placeholder_format.idx
        if ph_idx == 0 and content.title:
            _populate_placeholder_with_format(
                shape, content.title, is_title=True, template_style=template_style
            )
            populated_indices.add(ph_idx)
            title_placed = True
        elif ph_idx == 1:
            if body_paragraphs and (
                preferred_text_ph is None or shape == preferred_text_ph
            ):
                _populate_placeholder_with_format(
                    shape,
                    body_paragraphs,
                    is_title=False,
                    template_style=template_style,
                )
                populated_indices.add(ph_idx)
                body_placed = True
            elif content.subtitle:
                _populate_placeholder_with_format(
                    shape,
                    content.subtitle,
                    is_title=True,
                    template_style=template_style,
                )
                populated_indices.add(ph_idx)
                body_placed = True

    if preferred_text_ph is not None and not body_placed and body_paragraphs:
        _populate_placeholder_with_format(
            preferred_text_ph,
            body_paragraphs,
            is_title=False,
            template_style=template_style,
        )
        populated_indices.add(preferred_text_ph.placeholder_format.idx)
        body_placed = True

    if not body_placed and body_paragraphs:
        for shape in new_slide.placeholders:
            ph_idx = shape.placeholder_format.idx
            if (
                ph_idx > 1
                and shape.has_text_frame
                and not _is_visual_placeholder_type(shape.placeholder_format.type)
            ):
                _populate_placeholder_with_format(
                    shape,
                    body_paragraphs,
                    is_title=False,
                    template_style=template_style,
                )
                populated_indices.add(ph_idx)
                body_placed = True
                break

    # Fallback text boxes using text_region bounds
    if not title_placed and content.title:
        txBox = new_slide.shapes.add_textbox(
            region_map.text_region.left,
            Inches(0.3),
            region_map.text_region.width,
            Inches(1.0),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = content.title
        for para in tf.paragraphs:
            para.font.size = Pt(
                template_style.title_font_size_pt if template_style else 28
            )
            para.font.bold = True
            if template_style:
                para.font.name = (
                    template_style.title_font_family
                    or template_style.theme.major_font
                    or ""
                ) or None
            if template_style and template_style.title_font_color_rgb:
                try:
                    from pptx.dml.color import RGBColor

                    para.font.color.rgb = RGBColor.from_string(
                        template_style.title_font_color_rgb
                    )
                except Exception:
                    pass

    if not body_placed and body_paragraphs:
        max_font = _compute_max_font_size(region_map.text_region, len(body_paragraphs))
        txBox = new_slide.shapes.add_textbox(
            region_map.text_region.left,
            region_map.text_region.top,
            region_map.text_region.width,
            region_map.text_region.height,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, (text, level) in enumerate(body_paragraphs):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text = text
            para.level = level
            para.font.size = Pt(max_font)
            if template_style:
                para.font.name = (
                    template_style.body_font_family
                    or template_style.theme.minor_font
                    or ""
                ) or None
            if template_style and template_style.body_font_color_rgb:
                try:
                    from pptx.dml.color import RGBColor

                    para.font.color.rgb = RGBColor.from_string(
                        template_style.body_font_color_rgb
                    )
                except Exception:
                    pass
        font_family = "Calibri"
        if template_style:
            font_family = (
                template_style.body_font_family
                or template_style.theme.minor_font
                or "Calibri"
            )
        try:
            tf.fit_text(font_family=font_family, max_size=max_font)
        except Exception as e:
            if VERBOSE:
                print("[VERBOSE] Exception suppressed: %s" % str(e))
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # --- Fix 3: Minimum font size enforcement (fallback textbox) ---
        ns_a_mincheck2 = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        min_font_sz2 = _MIN_BODY_FONT_PT * 100
        font_too_small2 = False
        for para in tf.paragraphs:
            for run in para.runs:
                rPr = run._r.find(ns_a_mincheck2 + "rPr")
                if rPr is not None:
                    sz = rPr.get("sz")
                    if sz and int(sz) < min_font_sz2:
                        font_too_small2 = True
                        rPr.set("sz", str(min_font_sz2))
        if font_too_small2:
            print(
                "  [FONT GUARD] Fallback textbox font was below %dpt — enforced %dpt."
                % (_MIN_BODY_FONT_PT, _MIN_BODY_FONT_PT)
            )
            # If too many paragraphs for the space at min font, truncate
            if len(body_paragraphs) > 8:
                print(
                    "  [FONT GUARD] Truncating %d body paragraphs to 8 to prevent overflow."
                    % len(body_paragraphs)
                )

    # Place visuals into the visual_region (separate from text)
    visual_area = region_map.visual_region
    image_area = None
    if (content.images or generated_image_bytes) and (content.charts or content.tables):
        reserved = list(_layout_non_placeholder_text_bounds(new_slide.slide_layout))
        reserved.append(region_map.text_region)
        min_area = int(slide_width) * int(slide_height)
        min_area = int(min_area * 0.05)
        pic_ph = _best_visual_placeholder(
            new_slide.slide_layout,
            {PP_PLACEHOLDER.PICTURE},
            min_area=min_area,
            reserved=reserved,
            allow_body_fallback=False,
        )
        if pic_ph is not None:
            image_area = ContentArea(
                left=pic_ph.left,
                top=pic_ph.top,
                width=pic_ph.width,
                height=pic_ph.height,
            )
        else:
            gap = int(slide_width * 0.02)
            # Split visual area to avoid overlap
            if visual_area.width >= visual_area.height:
                image_w = int(visual_area.width * 0.30)
                image_area = ContentArea(
                    left=visual_area.left + visual_area.width - image_w,
                    top=visual_area.top,
                    width=image_w,
                    height=visual_area.height,
                )
                visual_area = ContentArea(
                    left=visual_area.left,
                    top=visual_area.top,
                    width=max(1, visual_area.width - image_w - gap),
                    height=visual_area.height,
                )
            else:
                image_h = int(visual_area.height * 0.25)
                image_area = ContentArea(
                    left=visual_area.left,
                    top=visual_area.top + visual_area.height - image_h,
                    width=visual_area.width,
                    height=image_h,
                )
                visual_area = ContentArea(
                    left=visual_area.left,
                    top=visual_area.top,
                    width=visual_area.width,
                    height=max(1, visual_area.height - image_h - gap),
                )

    if VERBOSE and content.charts:
        chart_ph = _largest_placeholder(
            new_slide.slide_layout, {PP_PLACEHOLDER.CHART}, min_area=0
        )
        print(
            "[VERBOSE] Chart transfer region: (%d,%d,%d,%d) chart_placeholder=%s"
            % (
                visual_area.left,
                visual_area.top,
                visual_area.width,
                visual_area.height,
                "yes" if chart_ph is not None else "no",
            )
        )
    _transfer_tables(
        new_slide, content.tables, visual_area, template_style=template_style
    )
    _transfer_images(new_slide, content.images, image_area or visual_area)
    _transfer_charts(
        new_slide, content.charts, visual_area, template_style=template_style
    )
    # --- Fix 6a: Enable data labels on all transferred charts ---
    _ensure_chart_data_labels(new_slide)
    # --- Fix 6b: Enlarge undersized charts to fill available visual area ---
    _ensure_chart_fills_area(new_slide, slide_width, slide_height, visual_area)
    # Check if layout has actual content placeholders (body/object)
    # If it's a Blank or Title-only layout, transferring LLM shapes into
    # the restricted safe-margin box squashes them. Use the full slide dimensions.
    _has_content_ph = any(
        getattr(ph, "placeholder_format", None)
        and ph.placeholder_format.type
        in {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE}
        for ph in new_slide.slide_layout.placeholders
    )

    if _has_content_ph:
        shapes_target_area = region_map.visual_region
        text_shapes_target_area = region_map.text_region
    else:
        shapes_target_area = ContentArea(0, 0, slide_width, slide_height)
        text_shapes_target_area = ContentArea(0, 0, slide_width, slide_height)

    # P1-1: Pass source slide dimensions and target area so shapes are rescaled
    # proportionally to the template's content region instead of being copied
    # at Claude's original absolute EMU coordinates.
    # Gap 3: Also pass template_style so free-floating shape text uses the
    # template's body font instead of Claude's default typeface.
    _transfer_shapes(
        new_slide,
        content.shapes_xml,
        src_width=src_slide_width,
        src_height=src_slide_height,
        target_area=shapes_target_area,
        template_style=template_style,
    )
    if content_mix == ContentMix.TEXT_ONLY:
        _transfer_shapes(
            new_slide,
            content.text_shapes_xml,
            src_width=src_slide_width,
            src_height=src_slide_height,
            target_area=text_shapes_target_area,
            template_style=template_style,
        )

    # --- Fix 4: Post-transfer overlap detection and reflow ---
    _fix_overlapping_shapes(new_slide, slide_width, slide_height)
    # --- Fix 7: Remove duplicate title textboxes ---
    _remove_duplicate_titles(new_slide)

    # ------------------------------------------------------------------
    # Insert generated images into picture placeholders.
    # Detection uses both the tracked indices AND XML-level type checking
    # as a fallback (type="pic" or "clipArt" in the XML = picture placeholder).
    # ------------------------------------------------------------------
    if generated_image_bytes:
        # Build a comprehensive set of picture placeholder indices using
        # the _is_picture_placeholder helper.
        _pic_ph_indices = set(content.image_placeholder_indices)
        image_inserted = False
        fallback_bounds = None
        min_area = int(slide_width) * int(slide_height)
        min_area = int(min_area * 0.10)

        for shape in list(new_slide.placeholders):
            if _is_picture_placeholder(shape):
                _pic_ph_indices.add(shape.placeholder_format.idx)

        # Choose a suitable picture placeholder (avoid tiny footer slots)
        reserved = list(_layout_non_placeholder_text_bounds(new_slide.slide_layout))
        reserved.append(region_map.text_region)
        if image_area is not None:
            fallback_bounds = (
                image_area.left,
                image_area.top,
                image_area.width,
                image_area.height,
            )
        pic_ph = _best_visual_placeholder(
            new_slide.slide_layout,
            {PP_PLACEHOLDER.PICTURE},
            min_area=min_area,
            reserved=reserved,
            allow_body_fallback=False,
        )
        if pic_ph is not None and _placeholder_area(pic_ph) >= min_area:
            try:
                left, top, width, height = (
                    pic_ph.left,
                    pic_ph.top,
                    pic_ph.width,
                    pic_ph.height,
                )
                fallback_bounds = (left, top, width, height)
                try:
                    new_slide.shapes._spTree.remove(pic_ph._element)
                except Exception:
                    pass
                _add_picture_within_bounds(
                    new_slide,
                    generated_image_bytes,
                    ContentArea(left=left, top=top, width=width, height=height),
                )
                image_inserted = True
                try:
                    populated_indices.add(pic_ph.placeholder_format.idx)
                except Exception:
                    pass
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))

        # If no picture was inserted (e.g., placeholder insertion unsupported),
        # fall back to adding the image as a regular picture scaled to the
        # visual region (or placeholder bounds if captured).
        if not image_inserted:
            try:
                if fallback_bounds:
                    left, top, width, height = fallback_bounds
                    bounds = ContentArea(left=left, top=top, width=width, height=height)
                else:
                    bounds = region_map.visual_region
                _add_picture_within_bounds(new_slide, generated_image_bytes, bounds)
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))

    # Ensure text renders on top of visual elements for mixed-content slides
    if content_mix != ContentMix.TEXT_ONLY and (
        content.title or content.subtitle or content.body_paragraphs
    ):
        _ensure_text_on_top(new_slide)

    # Ensure all text has sufficient contrast against the slide background
    if template_style:
        try:
            _ensure_text_contrast(new_slide, template_style)
        except Exception as e:
            if VERBOSE:
                print("[VERBOSE] Text contrast check failed: %s" % str(e))

    # P1-2: Populate footer placeholders BEFORE clearing unused ones so that
    # footer elements with content are retained by _clear_unused_placeholders.
    _populate_footer_placeholders(
        new_slide,
        populated_indices,
        footer_text=footer_text,
        show_slide_number=show_slide_number,
        date_text=date_text,
    )

    # Fallback: if footer_text was requested but no idx=11 placeholder exists
    # on this slide layout (e.g. the title slide uses a decorative master shape
    # rather than an editable placeholder), add a free-floating text box in the
    # footer zone so the text still appears.
    if footer_text and 11 not in populated_indices:
        _footer_h = int(slide_height * 0.07)
        _footer_top = slide_height - _footer_h
        _ftb = new_slide.shapes.add_textbox(0, _footer_top, slide_width, _footer_h)
        _ftf = _ftb.text_frame
        _ftf.word_wrap = False
        _fp = _ftf.paragraphs[0]
        _fp.text = footer_text
        _fp.font.size = Pt(9)
        if template_style:
            _fp.font.name = (
                template_style.body_font_family
                or template_style.theme.minor_font
                or None
            )
            if template_style.body_font_color_rgb:
                try:
                    from pptx.dml.color import RGBColor

                    _fp.font.color.rgb = RGBColor.from_string(
                        template_style.body_font_color_rgb
                    )
                except Exception:
                    pass

    # ------------------------------------------------------------------
    # Remove ALL unpopulated placeholder XML elements from the slide shape
    # tree. This is the ONLY reliable way to eliminate ALL types of ghost
    # text ("Click to add title", "Click to add text", "Click icon to add
    # picture", content placeholder icons, etc.).
    # tf.clear() alone is insufficient for picture and content placeholders.
    # ------------------------------------------------------------------
    _clear_unused_placeholders(new_slide, populated_indices)

    # Remove any empty text boxes (non-placeholders) that may remain from
    # certain templates and would otherwise sit on top of visuals.
    _remove_empty_textboxes(new_slide)


# ---------------------------------------------------------------------------
# Step 1: Content Generation
# ---------------------------------------------------------------------------


def step_generate_content(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Generate PowerPoint content using Claude's pptx skill.

    Supports both streaming and non-streaming modes (controlled by the
    ``stream`` key in session_state). Streaming is required for long-running
    skill operations (>10 min) but may have issues with provider_data
    propagation. Non-streaming is simpler and more reliable for shorter
    operations but can timeout on complex presentations.

    This step:
    1. Creates a Claude agent with the pptx skill
    2. Runs the user's prompt to generate a presentation
    3. Downloads the generated .pptx file
    4. Extracts SlideContent from each slide
    5. Stores the extracted content in session_state
    """
    global VERBOSE
    VERBOSE = session_state.get("verbose", VERBOSE)
    _step_start = time.time()

    user_prompt = step_input.input
    template_path = session_state.get("template_path", "")
    output_dir = session_state.get("output_dir", ".")

    print("=" * 60)
    print("Step 1: Generating presentation content with Claude...")
    print("=" * 60)

    # Build prompt with template context
    enhanced_prompt = user_prompt
    try:
        prs = Presentation(template_path)
        layouts = [layout.name for layout in prs.slide_layouts]
        layout_info = ", ".join(layouts)
        enhanced_prompt = (
            user_prompt + "\n\n"
            "Important structural requirements for template compatibility:\n"
            "- Use one clear title and concise bullet points per slide.\n"
            "- Do not apply custom fonts, colors, or theme styling.\n"
            "- Tables and charts are supported and will be transferred to the template.\n"
            "- Keep tables to max 6 rows x 5 columns.\n"
            "- Use bar, column, line, or pie charts with clearly labeled data.\n"
            "- The template has these available layouts: " + layout_info + ".\n"
            "- Use standard slide ordering: Title Slide, then Content Slides, then Closing."
        )
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Exception suppressed: %s" % str(e))

    # Claude 4.6 Token Limits
    # Model 	Context Window (Input)	Max Output Tokens
    # Claude Opus 4.6	200,000 / 1,000,000 (beta)	128,000 tokens
    # Claude Sonnet 4.6	200,000 / 1,000,000 (beta)	64,000 tokens

    # Create the Claude agent
    content_agent = Agent(
        name="Content Generator",
        model=Claude(
            id="claude-opus-4-6",
            betas=["context-1m-2025-08-07"],
            # "claude-sonnet-4-5-20250929",
            # # Or: "claude-sonnet-4-6", "claude-opus-4-6"
            max_tokens=128000,
            skills=[{"type": "anthropic", "skill_id": "pptx", "version": "latest"}],
        ),
        instructions=[
            "You are a structured content generator for PowerPoint presentations.",
            "Your output will be extracted and remapped to a corporate template.",
            "",
            "SLIDE STRUCTURE:",
            "- Use exactly one clear, descriptive title per slide.",
            "- Use concise bullet points for body content.",
            "- Limit to 4-6 bullet points per slide, each bullet max ~15 words.",
            "- Keep subtitle text on title slides to a single short line.",
            "- Follow standard slide ordering: Title Slide, then Content Slides, then Closing Slide.",
            "",
            "VISUAL ELEMENTS (use when appropriate):",
            "- Include tables for data comparisons and structured information.",
            "- Keep tables concise: no more than 6 rows and 5 columns.",
            "- Use bar, column, line, or pie charts for data visualization.",
            "",
            "FORMATTING RESTRICTIONS:",
            "- Do NOT apply custom fonts, colors, or theme styling.",
            "- Do NOT use SmartArt or complex nested graphic layouts.",
            "- Do NOT add speaker notes, animations, or transitions.",
        ],
        markdown=True,
    )

    # Determine streaming mode from session_state
    use_stream = session_state.get("stream", True)

    if use_stream:
        # Streaming mode: required for long-running skill operations (>10 min)
        response: RunOutput | None = None
        event_count = 0
        try:
            for event in content_agent.run(
                enhanced_prompt, stream=True, yield_run_output=True
            ):
                event_count += 1
                if isinstance(event, RunOutput):
                    response = event
                    if VERBOSE:
                        print(
                            "[VERBOSE] Received RunOutput after %d events" % event_count
                        )
                elif VERBOSE and event_count <= 5:
                    print(
                        "[VERBOSE] Stream event %d: type=%s"
                        % (event_count, type(event).__name__)
                    )
                elif VERBOSE and event_count == 6:
                    print("[VERBOSE] (Suppressing further stream event logs...)")
        except Exception as e:
            print("[ERROR] Agent streaming failed with exception: %s" % str(e))
            if VERBOSE:
                traceback.print_exc()
            return StepOutput(
                content="Error: Agent execution failed: %s" % str(e),
                success=False,
                stop=True,
            )
        if VERBOSE:
            print("[VERBOSE] Total stream events received: %d" % event_count)
    else:
        # Non-streaming mode: simpler, more reliable for shorter operations
        response = None
        try:
            response = content_agent.run(enhanced_prompt, stream=False)
        except Exception as e:
            print("[ERROR] Agent execution failed with exception: %s" % str(e))
            if VERBOSE:
                traceback.print_exc()
            return StepOutput(
                content="Error: Agent execution failed: %s" % str(e),
                success=False,
                stop=True,
            )

    if response is None:
        return StepOutput(
            content="Error: No response received from content agent.",
            success=False,
            stop=True,
        )
    print("\nAgent response received.")

    if VERBOSE:
        # Diagnostic: Inspect the RunOutput
        print("\n[VERBOSE] === Agent Response Diagnostics ===")
        if response is None:
            print("[VERBOSE] Response is None!")
        else:
            print("[VERBOSE] Response type: %s" % type(response).__name__)
            print(
                "[VERBOSE] Number of messages: %d"
                % (len(response.messages) if response.messages else 0)
            )
            if response.messages:
                for m_idx, msg in enumerate(response.messages):
                    print("[VERBOSE] Message %d:" % m_idx)
                    print("[VERBOSE]   Type: %s" % type(msg).__name__)
                    print("[VERBOSE]   Role: %s" % getattr(msg, "role", "N/A"))
                    print(
                        "[VERBOSE]   Has provider_data: %s"
                        % (
                            hasattr(msg, "provider_data")
                            and msg.provider_data is not None
                        )
                    )
                    if hasattr(msg, "provider_data") and msg.provider_data:
                        pd = msg.provider_data
                        print("[VERBOSE]   provider_data type: %s" % type(pd).__name__)
                        if isinstance(pd, dict):
                            print(
                                "[VERBOSE]   provider_data keys: %s" % list(pd.keys())
                            )
                        elif hasattr(pd, "__dict__"):
                            print(
                                "[VERBOSE]   provider_data attrs: %s"
                                % list(vars(pd).keys())[:20]
                            )
                    # Also show content snippet
                    content_text = getattr(msg, "content", None)
                    if content_text and isinstance(content_text, str):
                        print(
                            "[VERBOSE]   Content (first 200 chars): %s"
                            % content_text[:200]
                        )
            # Check for response-level attributes
            for attr_name in [
                "content",
                "tool_calls",
                "extra_data",
                "metrics",
                "model_provider_data",
            ]:
                val = getattr(response, attr_name, "MISSING")
                if val != "MISSING":
                    if isinstance(val, str):
                        print(
                            "[VERBOSE] response.%s (first 200 chars): %s"
                            % (attr_name, val[:200])
                        )
                    elif val is not None:
                        print(
                            "[VERBOSE] response.%s type: %s"
                            % (attr_name, type(val).__name__)
                        )
        print("[VERBOSE] === End Diagnostics ===\n")

    # Download the generated file
    client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
    generated_file = None

    if response.messages:
        for msg_idx, msg in enumerate(response.messages):
            if VERBOSE:
                print("[VERBOSE] Checking message %d for provider_data..." % msg_idx)
            if hasattr(msg, "provider_data") and msg.provider_data:
                if VERBOSE:
                    print(
                        "[VERBOSE] Message %d has provider_data, calling download_skill_files..."
                        % msg_idx
                    )
                try:
                    files = download_skill_files(
                        msg.provider_data, client, output_dir=output_dir
                    )
                except Exception as e:
                    print(
                        "[ERROR] download_skill_files raised an exception: %s" % str(e)
                    )
                    if VERBOSE:
                        traceback.print_exc()
                    files = None

                if VERBOSE:
                    print("[VERBOSE] download_skill_files returned: %s" % files)

                if files:
                    for f in files:
                        if not f.endswith(".pptx"):
                            if VERBOSE:
                                print("[VERBOSE] Skipping non-.pptx file: %s" % f)
                            continue
                        try:
                            Presentation(f)
                            generated_file = f
                            print("Downloaded: %s" % generated_file)
                            break
                        except Exception as e:
                            print(
                                "[WARNING] File '%s' failed pptx validation: %s"
                                % (f, str(e))
                            )
                            continue
                    if not generated_file:
                        # Fallback: try any file
                        for f in files:
                            try:
                                Presentation(f)
                                generated_file = f
                                print("Downloaded (fallback): %s" % generated_file)
                                break
                            except Exception as e:
                                if VERBOSE:
                                    print(
                                        "[VERBOSE] Fallback file '%s' also failed: %s"
                                        % (f, str(e))
                                    )
                                continue
                elif VERBOSE:
                    print("[VERBOSE] No files returned from download_skill_files")
            elif VERBOSE:
                print("[VERBOSE] Message %d has no provider_data" % msg_idx)
            if generated_file:
                break
    else:
        if VERBOSE:
            print("[VERBOSE] response.messages is empty or None!")

    # Fallback: check response.model_provider_data directly
    # This catches cases where provider_data is on the RunOutput but not on individual messages
    if not generated_file and response.model_provider_data:
        if VERBOSE:
            print(
                "[VERBOSE] Trying fallback: response.model_provider_data = %s"
                % response.model_provider_data
            )
        try:
            files = download_skill_files(
                response.model_provider_data, client, output_dir=output_dir
            )
        except Exception as e:
            print("[ERROR] download_skill_files (fallback) raised: %s" % str(e))
            if VERBOSE:
                traceback.print_exc()
            files = None
        if files:
            if VERBOSE:
                print("[VERBOSE] Fallback download_skill_files returned: %s" % files)
            for f in files:
                if f.endswith(".pptx"):
                    try:
                        Presentation(f)
                        generated_file = f
                        print(
                            "Downloaded (via model_provider_data fallback): %s"
                            % generated_file
                        )
                        break
                    except Exception as e:
                        print(
                            "[WARNING] Fallback file '%s' failed pptx validation: %s"
                            % (f, str(e))
                        )
            if not generated_file:
                for f in files:
                    try:
                        Presentation(f)
                        generated_file = f
                        print(
                            "Downloaded (via model_provider_data fallback): %s"
                            % generated_file
                        )
                        break
                    except Exception as e:
                        if VERBOSE:
                            print(
                                "[VERBOSE] Fallback file '%s' also failed: %s"
                                % (f, str(e))
                            )

    if not generated_file:
        if VERBOSE:
            print("\n[VERBOSE] === File Generation Failure Summary ===")
            print("[VERBOSE] No valid .pptx file was found in any message.")
            print(
                "[VERBOSE] Total messages checked: %d"
                % (len(response.messages) if response.messages else 0)
            )
            messages_with_pd = sum(
                1
                for m in (response.messages or [])
                if hasattr(m, "provider_data") and m.provider_data
            )
            print("[VERBOSE] Messages with provider_data: %d" % messages_with_pd)
            print("[VERBOSE] Hint: The agent may not have generated a file. Check if:")
            print(
                "[VERBOSE]   1. The agent's response contained tool use / skill output"
            )
            print("[VERBOSE]   2. The ANTHROPIC_API_KEY has access to the pptx skill")
            print("[VERBOSE]   3. The prompt was understood correctly by the agent")
            print("[VERBOSE] === End Failure Summary ===\n")
        return StepOutput(
            content="Error: No presentation file was generated.",
            success=False,
            stop=True,
        )

    # Extract slide content
    print("\nExtracting slide content...")
    generated_prs = Presentation(generated_file)
    slides_data = []

    # P1-4: Store source slide dimensions so Step 4 can rescale shapes
    # proportionally from Claude's default slide size to the template.
    session_state["src_slide_width"] = int(generated_prs.slide_width)
    session_state["src_slide_height"] = int(generated_prs.slide_height)
    if VERBOSE:
        print(
            "[VERBOSE] Source slide dimensions: %d x %d EMU"
            % (generated_prs.slide_width, generated_prs.slide_height)
        )

    # Check template layouts for picture placeholders so the image planner
    # knows which slides have dedicated picture slots in the template.
    template_prs = None
    try:
        template_prs = Presentation(template_path)
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Exception suppressed: %s" % str(e))

    total_gen_slides = len(list(generated_prs.slides))

    for idx, slide in enumerate(generated_prs.slides):
        content = _extract_slide_content(slide)

        # Check if the template layout for this slide position has image placeholders
        has_template_image_ph = False
        if template_prs is not None:
            try:
                content_mix = _classify_content_mix(content, has_generated_image=False)
                layout = _find_best_layout(
                    template_prs,
                    idx,
                    total_gen_slides,
                    content_mix=content_mix,
                )
                for ph in layout.placeholders:
                    ph_fmt = ph.placeholder_format
                    if ph_fmt is not None and ph_fmt.type == PP_PLACEHOLDER.PICTURE:
                        has_template_image_ph = True
                        break
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))

        slide_info = {
            "index": idx,
            "title": content.title,
            "subtitle": content.subtitle,
            "body_text": " | ".join(
                t if isinstance(t, str) else t[0] for t in content.body_paragraphs[:3]
            )
            if content.body_paragraphs
            else "",
            "bullet_count": len(content.body_paragraphs),
            "has_table": len(content.tables) > 0,
            "has_chart": len(content.charts) > 0,
            "has_image": len(content.images) > 0,
            "has_shapes": len(content.shapes_xml) > 0,
            "has_image_placeholder": has_template_image_ph,
        }
        slides_data.append(slide_info)
        print(
            "  Slide %d: '%s' | tables:%d charts:%d images:%d img_ph:%s"
            % (
                idx + 1,
                content.title[:40] if content.title else "",
                len(content.tables),
                len(content.charts),
                len(content.images),
                has_template_image_ph,
            )
        )

    # Store in session state
    session_state["generated_file"] = generated_file
    session_state["slides_data"] = slides_data
    session_state["total_slides"] = len(slides_data)

    slides_summary = json.dumps(slides_data, indent=2)
    elapsed = time.time() - _step_start
    print("[TIMING] Step 1 Content Generation: completed in %.2fs" % elapsed)
    return StepOutput(
        content=slides_summary,
        success=True,
    )


# ---------------------------------------------------------------------------
# Step 2: Image Planning (Agent with output_schema)
# ---------------------------------------------------------------------------

# This agent decides which slides need AI-generated images
# image_planner agent is now loaded from the agents/ package via get_agents().
# Stored in session_state["agents"]["image_planner"].
# See agents/claude_agents.py (or openai/gemini variants) for the definition.


def step_plan_images(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Plan which slides need AI-generated images.

    This is a function step wrapper around the image_planner agent.
    It combines the user's original prompt with slide metadata so the
    planner has full context (the framework's _prepare_message() would
    otherwise discard the original prompt when previous step outputs exist).
    """
    global VERBOSE
    VERBOSE = session_state.get("verbose", VERBOSE)
    _step_start = time.time()

    print("\n" + "=" * 60)
    print("Step 2: Planning images for slides...")
    print("=" * 60)

    # Retrieve user prompt from session_state (stored at workflow creation)
    user_prompt = session_state.get("user_prompt", "professional business presentation")

    # Get slide metadata from previous step output
    slides_json = step_input.previous_step_content or "[]"
    if not isinstance(slides_json, str):
        slides_json = json.dumps(slides_json, indent=2) if slides_json else "[]"

    # Construct a combined message with full context
    combined_message = (
        'Presentation topic: "%s"\n\n'
        "Slide metadata:\n%s\n\n"
        "Analyze each slide and decide which ones need AI-generated images.\n"
        "Consider the presentation topic when writing image prompts."
    ) % (user_prompt, slides_json)

    if VERBOSE:
        print(
            "[VERBOSE] Image planner input (first 500 chars):\n%s"
            % combined_message[:500]
        )

    # Lazily load agent: do not use session_state["agents"] to avoid deepcopy failures.
    from agents import get_agents as _get_agents
    _provider = session_state.get("llm_provider", "claude")
    _image_planner = _get_agents(_provider).get("image_planner")
    _image_planner_fallback = _get_agents(_provider).get("image_planner_fallback")
    
    try:
        response = _image_planner.run(combined_message, stream=False)
    except Exception as e:
        print("[WARNING] Image planner failed on primary agent: %s" % str(e))
        if _image_planner_fallback:
            print("[FALLBACK TRIGGERED] Engaging fallback image planner...")
            try:
                response = _image_planner_fallback.run(combined_message, stream=False)
            except Exception as fallback_e:
                print("[ERROR] Fallback Image planner failed: %s" % str(fallback_e))
                if VERBOSE:
                    traceback.print_exc()
                return StepOutput(
                    content=json.dumps({"decisions": []}),
                    success=True,
                )
        else:
            if VERBOSE:
                traceback.print_exc()
            return StepOutput(
                content=json.dumps({"decisions": []}),
                success=True,
            )

    # Extract the structured output
    if response and response.content:
        if isinstance(response.content, BaseModel):
            result = response.content.model_dump_json()
        elif isinstance(response.content, dict):
            result = json.dumps(response.content)
        else:
            result = str(response.content)
        if VERBOSE:
            print(
                "[VERBOSE] Image planner output (first 500 chars):\n%s" % result[:500]
            )
        elapsed = time.time() - _step_start
        print("[TIMING] Step 2 Image Planning: completed in %.2fs" % elapsed)
        return StepOutput(content=result, success=True)

    print("[WARNING] Image planner returned no content.")
    elapsed = time.time() - _step_start
    print("[TIMING] Step 2 Image Planning: completed in %.2fs" % elapsed)
    return StepOutput(
        content=json.dumps({"decisions": []}),
        success=True,
    )


# ---------------------------------------------------------------------------
# Step 3: Image Generation
# ---------------------------------------------------------------------------


def step_generate_images(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Generate images for slides that need them using NanoBanana (powered by Gemini).

    This step:
    1. Reads the image plan from step 2 (produced by the Gemini image planner)
    2. For each slide marked as needing an image, generates one with NanoBanana
    3. Stores the generated image bytes in session_state keyed by slide index
    """
    global VERBOSE
    VERBOSE = session_state.get("verbose", VERBOSE)
    _step_start = time.time()

    print("\n" + "=" * 60)
    print("Step 3: Generating images with NanoBanana...")
    print("=" * 60)

    # Parse the image plan from step 2
    image_plan_content = step_input.previous_step_content
    generated_images: Dict[int, bytes] = {}

    if not image_plan_content:
        print("No image plan received. Skipping image generation.")
        session_state["generated_images"] = {}
        return StepOutput(content="No images to generate.", success=True)

    # Parse the plan
    try:
        if isinstance(image_plan_content, str):
            plan_data = json.loads(image_plan_content)
        elif isinstance(image_plan_content, dict):
            plan_data = image_plan_content
        elif isinstance(image_plan_content, ImagePlan):
            plan_data = image_plan_content.model_dump()
        elif isinstance(image_plan_content, BaseModel):
            plan_data = image_plan_content.model_dump()
        else:
            plan_data = json.loads(str(image_plan_content))
    except (json.JSONDecodeError, TypeError):
        print("Could not parse image plan. Skipping image generation.")
        session_state["generated_images"] = {}
        return StepOutput(content="Could not parse image plan.", success=True)

    decisions = plan_data.get("decisions", [])
    slides_needing_images = [d for d in decisions if d.get("needs_image", False)]

    # Enforce minimum images if configured.
    # IMPORTANT: Only consider non-data-vis slides as candidates.
    # Slides with native charts, tables, or infographics must NEVER receive an
    # AI-generated image — the external photo/illustration would collide with the
    # native visual and degrade the slide's clarity.
    DATA_VIS_KEYWORDS = ("chart", "table", "infographic", "diagram", "graph")
    min_images = session_state.get("min_images", 1)
    if len(slides_needing_images) < min_images and min_images > 0:
        print(
            "Image planner selected %d slides, but --min-images=%d. Adding more..."
            % (len(slides_needing_images), min_images)
        )
        # Find slides not yet selected, excluding data-vis slides
        selected_indices = {d["slide_index"] for d in slides_needing_images}
        remaining = [
            d
            for d in decisions
            if d["slide_index"] not in selected_indices
            # Skip slides that the planner flagged as having native data visualizations
            and not d.get("has_table", False)
            and not d.get("has_chart", False)
            and not d.get("has_data_vis", False)
            # Skip slides whose visual_suggestion mentions a data-vis type
            and not any(
                kw in d.get("visual_suggestion", "").lower() for kw in DATA_VIS_KEYWORDS
            )
        ]
        # Sort: image_placeholder slides first, then title slide (index 0)
        remaining.sort(
            key=lambda d: (
                not d.get("has_image_placeholder", False),
                d.get("slide_index", 99),
            )
        )

        for d in remaining:
            if len(slides_needing_images) >= min_images:
                break
            # Create a generic prompt based on available context
            topic = session_state.get("user_prompt", "professional business")
            slide_idx = d.get("slide_index", 0)
            d["needs_image"] = True
            d["image_prompt"] = (
                "Professional minimalist corporate illustration related to %s, "
                "suitable for slide %d of a business presentation. "
                "Clean, modern flat design with subtle gradients."
            ) % (topic, slide_idx + 1)
            slides_needing_images.append(d)

    if not slides_needing_images:
        print("Image planner decided no slides need images.")
        session_state["generated_images"] = {}
        return StepOutput(content="No slides need images.", success=True)

    # Check which slides already have images from Claude
    slides_data = session_state.get("slides_data", [])
    slides_with_existing_images = {
        s["index"] for s in slides_data if s.get("has_image", False)
    }

    # Build a lookup from slide_index -> slide metadata for the belt-and-suspenders
    # guard below. This is necessary because SlideImageDecision (the image planner's
    # output schema) only carries slide_index, needs_image, image_prompt, and reasoning —
    # it never has has_table, has_chart, has_data_vis, or visual_suggestion.
    # Without this lookup, the existing decision.get("has_table", ...) checks always
    # return False and the guard never fires, allowing data-vis slides to receive images.
    slides_data_by_idx: Dict[int, dict] = {
        s.get("index", i): s for i, s in enumerate(slides_data)
    }

    # Initialize NanoBanana
    google_api_key = os.getenv("GOOGLE_API_KEY")
    if not google_api_key:
        print("GOOGLE_API_KEY not set. Skipping image generation.")
        session_state["generated_images"] = {}
        return StepOutput(
            content="GOOGLE_API_KEY not set. Skipped image generation.",
            success=True,
        )

    nano_banana = NanoBananaTools(
        api_key=google_api_key,
        aspect_ratio="16:9",  # Widescreen for presentations
    )

    for decision in slides_needing_images:
        slide_idx = decision["slide_index"]
        prompt = decision.get("image_prompt", "")

        if slide_idx in slides_with_existing_images:
            print(
                "  Slide %d: Already has image from Claude, skipping." % (slide_idx + 1)
            )
            continue

        # Belt-and-suspenders guard: even if this decision slipped through the
        # min_images enforcement filter above, never generate an external AI image
        # for a slide that carries native data visualizations (charts, tables,
        # infographics, or diagrams). The native visual already fills the visual
        # region; an added photo would collide and degrade the slide.
        #
        # IMPORTANT: SlideImageDecision (image planner output) only carries
        # slide_index, needs_image, image_prompt, and reasoning — it NEVER has
        # has_table, has_chart, has_data_vis, or visual_suggestion.
        # We must look these up from slides_data_by_idx, which was enriched by
        # step_process_chunks with storyboard data (Fix: was reading from decision).
        _slide_meta = slides_data_by_idx.get(slide_idx, {})
        if (
            _slide_meta.get("has_table", False)
            or _slide_meta.get("has_chart", False)
            or _slide_meta.get("has_data_vis", False)
        ):
            print(
                "  Slide %d: Contains native data visualization (chart/table/infographic)."
                " Skipping AI image generation to avoid collision." % (slide_idx + 1)
            )
            continue
        vs = _slide_meta.get("visual_suggestion", "").lower()
        if any(kw in vs for kw in DATA_VIS_KEYWORDS):
            print(
                "  Slide %d: Data-vis visual_suggestion ('%s…') detected."
                " Skipping AI image generation." % (slide_idx + 1, vs[:40])
            )
            continue

        if not prompt:
            print("  Slide %d: No prompt provided, skipping." % (slide_idx + 1))
            continue

        print("  Slide %d: Generating image..." % (slide_idx + 1))
        print("    Prompt: %s" % prompt[:80])

        try:
            result = nano_banana.create_image(prompt)
            if hasattr(result, "images") and result.images:
                for img in result.images:
                    if hasattr(img, "content") and img.content:
                        generated_images[slide_idx] = img.content
                        print(
                            "    Generated successfully (%d bytes)" % len(img.content)
                        )
                        break
            elif isinstance(result, str) and "successfully" in result.lower():
                print("    Image generated but no bytes returned.")
        except Exception as e:
            print("    Failed to generate image: %s" % str(e))
            if VERBOSE:
                traceback.print_exc()

    session_state["generated_images"] = generated_images
    elapsed = time.time() - _step_start
    print("[TIMING] Step 3 Image Generation: completed in %.2fs" % elapsed)
    return StepOutput(
        content="Generated %d image(s) for presentation." % len(generated_images),
        success=True,
    )


# ---------------------------------------------------------------------------
# Template Deep Analysis & Assembly Knowledge File (Step 4 prerequisite)
# ---------------------------------------------------------------------------


def _extract_shape_design_info(shape) -> dict:
    """Extract design properties from a non-placeholder shape (decorative element).

    Captures fill color, line color and width, position, and size for use in
    the template design language analysis. Returns an empty dict on failure.
    """
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    ns_p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    info: dict = {}
    try:
        try:
            info["shape_type"] = str(shape.shape_type)
        except Exception:
            info["shape_type"] = ""
        try:
            info["left_emu"] = int(shape.left)
            info["top_emu"] = int(shape.top)
            info["width_emu"] = int(shape.width)
            info["height_emu"] = int(shape.height)
        except Exception:
            pass

        # Fill and line color via XML (python-pptx fill API can raise for inherited fills)
        elem = shape._element
        spPr = None
        for candidate in [
            elem.find(ns_p + "spPr"),
            elem.find(ns_a + "spPr"),
            elem.find(".//" + ns_p + "spPr"),
            elem.find(".//" + ns_a + "spPr"),
        ]:
            if candidate is not None:
                spPr = candidate
                break

        if spPr is not None:
            solidFill = spPr.find(ns_a + "solidFill")
            if solidFill is not None:
                srgb = solidFill.find(ns_a + "srgbClr")
                if srgb is not None:
                    info["fill_color_hex"] = srgb.get("val", "")
            ln = spPr.find(ns_a + "ln")
            if ln is not None:
                w_str = ln.get("w")
                if w_str:
                    try:
                        info["line_width_pt"] = round(int(w_str) / 12700, 1)
                    except ValueError:
                        pass
                lnFill = ln.find(ns_a + "solidFill")
                if lnFill is not None:
                    srgb = lnFill.find(ns_a + "srgbClr")
                    if srgb is not None:
                        info["line_color_hex"] = srgb.get("val", "")

        try:
            info["has_text"] = getattr(shape, "has_text_frame", False) and bool(
                getattr(shape, "text", "").strip()
            )
        except Exception:
            info["has_text"] = False
    except Exception:
        return {}
    return info


def _analyze_template_in_depth(template_prs) -> dict:
    """Perform a thorough per-layout analysis of the template's complete design language.

    Opens and inspects every slide layout in the template, documenting:
    - Slide dimensions and aspect ratio
    - Theme colors: full accent palette (1-6), dark/light colors, hyperlink color
    - Font scheme: major (heading) and minor (body) typefaces
    - Slide master decorative shapes and background color
    - Per-layout: name, background color, all placeholder positions/sizes/typography
      (font family, size, weight, color), all non-placeholder (decorative) shapes
      with their fill colors, line styles, and positions
    - Design language summary: primary/secondary accents, fonts, typical font sizes,
      typography ratio, recurring visual motifs from shape fill colors

    This analysis is the authoritative design reference for Step 4. Every slide is
    assembled with complete fidelity to the template's visual language — nothing is
    approximated or guessed.

    Args:
        template_prs: A python-pptx Presentation object for the template file.

    Returns:
        dict with keys: slide_dimensions, theme, master_analysis,
                        layouts, design_language_summary
    """
    EMU_PER_INCH = 914400
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    ns_p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

    result: dict = {
        "slide_dimensions": {},
        "theme": {},
        "master_analysis": {},
        "layouts": [],
        "design_language_summary": {},
    }

    # --- Slide dimensions ---
    try:
        w = int(template_prs.slide_width)
        h = int(template_prs.slide_height)
        result["slide_dimensions"] = {
            "width_emu": w,
            "height_emu": h,
            "width_inches": round(w / EMU_PER_INCH, 2),
            "height_inches": round(h / EMU_PER_INCH, 2),
            "aspect_ratio": round(w / h, 3) if h else 0.0,
        }
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Slide dimensions extraction failed: %s" % str(e))

    # --- Theme colors and fonts ---
    theme_obj = None
    try:
        theme_obj = _extract_theme_from_prs(template_prs)
        result["theme"] = {
            "accent_colors": theme_obj.accent_colors,
            "dk1": theme_obj.dk1,
            "dk2": theme_obj.dk2,
            "lt1": theme_obj.lt1,
            "lt2": theme_obj.lt2,
            "hyperlink": theme_obj.hyperlink,
            "major_font": theme_obj.major_font,
            "minor_font": theme_obj.minor_font,
        }
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Theme extraction failed in deep analysis: %s" % str(e))

    # Helper: extract solid fill background color from a layout/master _element
    def _layout_bg_color(element) -> str:
        try:
            bg = element.find(ns_p + "cSld/" + ns_p + "bg")
            if bg is None:
                return ""
            bgPr = bg.find(ns_p + "bgPr") or bg.find(ns_p + "bgRef") or bg
            solidFill = bgPr.find(ns_a + "solidFill")
            if solidFill is not None:
                srgb = solidFill.find(ns_a + "srgbClr")
                if srgb is not None:
                    return srgb.get("val", "")
        except Exception:
            pass
        return ""

    # --- Slide master analysis ---
    try:
        master = template_prs.slide_masters[0]
        master_bg = _layout_bg_color(master._element)

        master_shapes = []
        for shape in master.shapes:
            try:
                if shape.is_placeholder:
                    continue
                info = _extract_shape_design_info(shape)
                if info:
                    master_shapes.append(info)
            except Exception:
                continue

        result["master_analysis"] = {
            "background_color": master_bg,
            "decorative_shape_count": len(master_shapes),
            "decorative_shapes": master_shapes[:12],  # cap to keep output manageable
        }
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Master analysis failed: %s" % str(e))

    # --- Per-layout analysis ---
    slide_w = result["slide_dimensions"].get("width_emu", 1) or 1
    slide_h = result["slide_dimensions"].get("height_emu", 1) or 1

    for layout_idx, layout in enumerate(template_prs.slide_layouts):
        layout_info: dict = {
            "index": layout_idx,
            "name": layout.name,
            "background_color": "",
            "placeholders": [],
            "decorative_shapes": [],
            "has_picture_placeholder": False,
            "has_chart_placeholder": False,
            "has_table_placeholder": False,
            "non_placeholder_shape_count": 0,
        }

        # Background color: layout-level first, then master fallback
        layout_bg = _layout_bg_color(layout._element)
        if not layout_bg:
            try:
                layout_bg = _layout_bg_color(layout.slide_master._element)
            except Exception:
                pass
        layout_info["background_color"] = layout_bg

        # Placeholder analysis — position, size, type, and typography
        for ph in layout.placeholders:
            try:
                ph_fmt = ph.placeholder_format
                ph_entry: dict = {
                    "idx": ph_fmt.idx,
                    "type_str": str(ph_fmt.type),
                    "left_emu": int(ph.left),
                    "top_emu": int(ph.top),
                    "width_emu": int(ph.width),
                    "height_emu": int(ph.height),
                    "left_pct": round(ph.left / slide_w * 100, 1),
                    "top_pct": round(ph.top / slide_h * 100, 1),
                    "width_pct": round(ph.width / slide_w * 100, 1),
                    "height_pct": round(ph.height / slide_h * 100, 1),
                    "font_family": "",
                    "font_size_pt": 0,
                    "font_bold": False,
                    "font_color_hex": "",
                }

                # Extract typography: check defRPr on pPr first, then first run's rPr
                if ph.has_text_frame:
                    for para in ph.text_frame.paragraphs:
                        rPr_source = None
                        pPr = para._p.find(ns_a + "pPr")
                        if pPr is not None:
                            rPr_source = pPr.find(ns_a + "defRPr")
                        if rPr_source is None and para.runs:
                            rPr_source = para.runs[0]._r.find(ns_a + "rPr")
                        if rPr_source is not None:
                            sz = rPr_source.get("sz")
                            if sz:
                                try:
                                    ph_entry["font_size_pt"] = int(sz) // 100
                                except ValueError:
                                    pass
                            b = rPr_source.get("b")
                            ph_entry["font_bold"] = b in ("1", "true") if b else False
                            latin = rPr_source.find(ns_a + "latin")
                            if latin is not None:
                                typeface = latin.get("typeface", "")
                                if typeface and not typeface.startswith("+"):
                                    ph_entry["font_family"] = typeface
                                elif typeface == "+mj-lt" and theme_obj:
                                    ph_entry["font_family"] = theme_obj.major_font
                                elif typeface == "+mn-lt" and theme_obj:
                                    ph_entry["font_family"] = theme_obj.minor_font
                            if theme_obj:
                                ph_entry["font_color_hex"] = _extract_color_from_rPr(
                                    rPr_source, theme_obj
                                )
                            break  # Only need first paragraph's style

                # Flag placeholder capability types
                if _is_picture_placeholder(ph):
                    layout_info["has_picture_placeholder"] = True
                elif ph_fmt.type == PP_PLACEHOLDER.CHART:
                    layout_info["has_chart_placeholder"] = True
                elif ph_fmt.type == PP_PLACEHOLDER.TABLE:
                    layout_info["has_table_placeholder"] = True

                layout_info["placeholders"].append(ph_entry)
            except Exception:
                continue

        # Non-placeholder (decorative) shapes in this layout
        for shape in layout.shapes:
            try:
                if shape.is_placeholder:
                    continue
                layout_info["non_placeholder_shape_count"] += 1
                info = _extract_shape_design_info(shape)
                if info:
                    layout_info["decorative_shapes"].append(info)
            except Exception:
                continue

        result["layouts"].append(layout_info)

    # --- Design language summary ---
    try:
        theme_data = result.get("theme", {})
        accent_colors = theme_data.get("accent_colors", [])

        all_title_sizes: list = []
        all_body_sizes: list = []
        layouts_with_pics = 0
        layouts_with_decoration = 0
        fill_color_counts: dict = {}

        for li in result["layouts"]:
            if li.get("has_picture_placeholder"):
                layouts_with_pics += 1
            if li.get("non_placeholder_shape_count", 0) > 0:
                layouts_with_decoration += 1
            for ph in li.get("placeholders", []):
                sz = ph.get("font_size_pt", 0)
                if sz > 0:
                    if ph.get("idx") == 0:
                        all_title_sizes.append(sz)
                    elif ph.get("idx") == 1:
                        all_body_sizes.append(sz)
            for shape_info in li.get("decorative_shapes", []):
                fc = shape_info.get("fill_color_hex", "")
                if fc:
                    fill_color_counts[fc] = fill_color_counts.get(fc, 0) + 1

        # Include master decorative shapes in motif color counts
        for shape_info in result["master_analysis"].get("decorative_shapes", []):
            fc = shape_info.get("fill_color_hex", "")
            if fc:
                fill_color_counts[fc] = fill_color_counts.get(fc, 0) + 1

        # Recurring motif colors: fills appearing on 2+ layouts/master shapes
        recurring_motifs = sorted(
            [
                {"color_hex": c, "occurrences": n}
                for c, n in fill_color_counts.items()
                if n >= 2
            ],
            key=lambda x: -x["occurrences"],
        )[:5]

        avg_title_size = (
            int(sum(all_title_sizes) / len(all_title_sizes)) if all_title_sizes else 28
        )
        avg_body_size = (
            int(sum(all_body_sizes) / len(all_body_sizes)) if all_body_sizes else 18
        )

        result["design_language_summary"] = {
            "primary_accent_color": (
                accent_colors[0] if accent_colors else theme_data.get("dk1", "")
            ),
            "secondary_accent_color": (
                accent_colors[1] if len(accent_colors) > 1 else ""
            ),
            "full_accent_palette": [
                {"hex": c, "label": "accent%d" % (i + 1)}
                for i, c in enumerate(accent_colors)
                if c
            ],
            "heading_font": theme_data.get("major_font", "Calibri"),
            "body_font": theme_data.get("minor_font", "Calibri"),
            "background_color": theme_data.get("lt1", "FFFFFF"),
            "primary_text_color": theme_data.get("dk1", "000000"),
            "typical_title_font_size_pt": avg_title_size,
            "typical_body_font_size_pt": avg_body_size,
            "typography_size_ratio": (
                round(avg_title_size / avg_body_size, 2) if avg_body_size else 0.0
            ),
            "total_layouts": len(result["layouts"]),
            "layouts_with_picture_placeholder": layouts_with_pics,
            "layouts_with_decorative_shapes": layouts_with_decoration,
            "recurring_motif_colors": recurring_motifs,
        }

        if VERBOSE:
            print("[VERBOSE] Template deep analysis complete:")
            print("[VERBOSE]   Total layouts analyzed: %d" % len(result["layouts"]))
            print(
                "[VERBOSE]   Accent palette: %s"
                % ", ".join("#%s" % c for c in accent_colors[:3] if c)
            )
            print(
                "[VERBOSE]   Heading font: %s  |  Body font: %s"
                % (
                    result["design_language_summary"]["heading_font"],
                    result["design_language_summary"]["body_font"],
                )
            )
            print(
                "[VERBOSE]   Typical title: %dpt  |  Typical body: %dpt"
                % (avg_title_size, avg_body_size)
            )
            print(
                "[VERBOSE]   Layouts with picture placeholders: %d" % layouts_with_pics
            )
            print(
                "[VERBOSE]   Layouts with decorative shapes: %d"
                % layouts_with_decoration
            )
            if recurring_motifs:
                print(
                    "[VERBOSE]   Recurring motif colors: %s"
                    % ", ".join(
                        "#%s (%dx)" % (m["color_hex"], m["occurrences"])
                        for m in recurring_motifs[:3]
                    )
                )
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Design language summary failed: %s" % str(e))

    return result


def _knowledge_json_default(obj) -> str:
    """JSON serialization fallback for non-JSON-serializable types in the knowledge file."""
    if isinstance(obj, bytes):
        return "<binary:%d bytes>" % len(obj)
    return str(obj)


def _build_assembly_knowledge_file(
    session_state: Dict,
    template_analysis: dict,
) -> dict:
    """Consolidate all four pipeline inputs into the comprehensive assembly knowledge file.

    This knowledge file is the single source of truth that governs every design and
    content decision made during PPTX file generation. It is built before any slide
    construction begins so that nothing is guessed, approximated, or inconsistently applied.

    The four mandatory inputs are:

    1. Original user intent — the exact prompt that initiated the workflow, preserved
       in full so the presentation's tone, goals, and specific requirements are never lost.

    2. Complete slide-by-slide content plan — the full per-slide plan from Step 1,
       including all titles, bullet points, table/chart flags, bullet counts, and
       structural decisions made for every individual slide.

    3. Deep template design language — the thorough analysis from _analyze_template_in_depth:
       every layout's placeholder positions, font families, font sizes, font weights, text
       colors, background colors, decorative shape fill/line styles, the complete color
       palette, accent colors, and all recurring design motifs.

    4. AI-generated image assets — full awareness of every image produced in Steps 2-3:
       which slide each image targets, its pixel dimensions, content type, aspect ratio,
       and how it is intended to be positioned within the slide layout.

    Args:
        session_state: The workflow session_state dict containing all pipeline data.
        template_analysis: Output of _analyze_template_in_depth() for the template file.

    Returns:
        dict with keys:
          metadata, input_1_user_intent, input_2_content_plan,
          input_3_template_design_language, input_4_image_assets, assembly_directives
    """
    knowledge: dict = {
        "metadata": {
            "pipeline_step": "Step 4 - Template Assembly",
            "purpose": (
                "Single source of truth governing every design and content decision "
                "during PPTX file generation. Nothing is guessed, approximated, or "
                "inconsistently applied."
            ),
            "inputs_count": 4,
        },
        "input_1_user_intent": {},
        "input_2_content_plan": {},
        "input_3_template_design_language": {},
        "input_4_image_assets": {},
        "assembly_directives": {},
    }

    # --- Input 1: Original user intent (preserved exactly as provided) ---
    user_prompt = session_state.get("user_prompt", "")
    knowledge["input_1_user_intent"] = {
        "original_prompt": user_prompt,
        "prompt_length_chars": len(user_prompt),
        "note": (
            "Preserved exactly as provided. Intent, tone, and goals are never lost "
            "during template assembly."
        ),
    }

    # --- Input 2: Complete slide-by-slide content plan from Step 1 ---
    slides_data = session_state.get("slides_data", [])
    total_slides = session_state.get("total_slides", len(slides_data))
    knowledge["input_2_content_plan"] = {
        "total_slides": total_slides,
        "generated_source_file": session_state.get("generated_file", ""),
        "src_slide_width_emu": session_state.get("src_slide_width", 0),
        "src_slide_height_emu": session_state.get("src_slide_height", 0),
        "slides": slides_data,
        "content_inventory": {
            "slides_with_table": [
                s["index"] for s in slides_data if s.get("has_table")
            ],
            "slides_with_chart": [
                s["index"] for s in slides_data if s.get("has_chart")
            ],
            "slides_with_existing_image": [
                s["index"] for s in slides_data if s.get("has_image")
            ],
            "slides_with_image_placeholder": [
                s["index"] for s in slides_data if s.get("has_image_placeholder")
            ],
        },
    }

    # --- Input 3: Deep template design language analysis ---
    # Provides complete fidelity to the template's visual language: fonts, colors,
    # layout grids, placeholder positions, margin patterns, shape styles, and motifs.
    knowledge["input_3_template_design_language"] = template_analysis

    # --- Input 4: AI-generated image assets from Steps 2-3 ---
    generated_images = session_state.get("generated_images", {})
    image_assets = []
    for slide_idx_key, img_bytes in generated_images.items():
        slide_idx = (
            int(slide_idx_key) if isinstance(slide_idx_key, str) else slide_idx_key
        )
        asset: dict = {
            "slide_index": slide_idx,
            "content_type": "image/png",
            "size_bytes": len(img_bytes) if img_bytes else 0,
            "has_image_data": bool(img_bytes),
            "intended_position": "picture_placeholder_or_right_visual_region",
            "slide_title": "",
            "has_picture_placeholder": False,
            "width_px": None,
            "height_px": None,
            "aspect_ratio": None,
        }
        # Cross-reference with content plan for slide title and placeholder info
        matching_slide = next(
            (s for s in slides_data if s.get("index") == slide_idx), None
        )
        if matching_slide:
            asset["slide_title"] = matching_slide.get("title", "")
            asset["has_picture_placeholder"] = matching_slide.get(
                "has_image_placeholder", False
            )
        # Get actual pixel dimensions for precise aspect-correct positioning
        try:
            from PIL import Image as _PILImage

            with _PILImage.open(BytesIO(img_bytes)) as im:
                asset["width_px"] = im.width
                asset["height_px"] = im.height
                if im.height:
                    asset["aspect_ratio"] = round(im.width / im.height, 3)
        except Exception:
            pass
        image_assets.append(asset)

    knowledge["input_4_image_assets"] = {
        "total_generated_images": len(image_assets),
        "slide_indices_with_images": sorted(a["slide_index"] for a in image_assets),
        "assets": image_assets,
    }

    # --- Assembly directives: concrete guidance synthesized from all four inputs ---
    design_summary = template_analysis.get("design_language_summary", {})
    theme_data = template_analysis.get("theme", {})
    slide_dims = template_analysis.get("slide_dimensions", {})

    knowledge["assembly_directives"] = {
        "primary_accent_color_hex": design_summary.get("primary_accent_color", ""),
        "secondary_accent_color_hex": design_summary.get("secondary_accent_color", ""),
        "full_color_palette": design_summary.get("full_accent_palette", []),
        "heading_font": design_summary.get("heading_font", "Calibri"),
        "body_font": design_summary.get("body_font", "Calibri"),
        "background_color_hex": design_summary.get("background_color", "FFFFFF"),
        "primary_text_color_hex": design_summary.get("primary_text_color", "000000"),
        "target_title_font_size_pt": design_summary.get(
            "typical_title_font_size_pt", 28
        ),
        "target_body_font_size_pt": design_summary.get("typical_body_font_size_pt", 18),
        "target_slide_width_emu": slide_dims.get("width_emu", 0),
        "target_slide_height_emu": slide_dims.get("height_emu", 0),
        "template_has_picture_layouts": (
            design_summary.get("layouts_with_picture_placeholder", 0) > 0
        ),
        "slides_receiving_generated_images": sorted(
            a["slide_index"] for a in image_assets
        ),
        "footer_text": session_state.get("footer_text", ""),
        "date_text": session_state.get("date_text", ""),
        "show_slide_numbers": session_state.get("show_slide_numbers", False),
    }

    return knowledge


# ---------------------------------------------------------------------------
# Step 4: Template Assembly
# ---------------------------------------------------------------------------


def step_assemble_template(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Apply template styling and assemble the final presentation.

    This is the most critical and knowledge-intensive step in the entire pipeline.
    Before constructing any slide, it first consolidates all necessary context into
    a comprehensive knowledge file that acts as the single source of truth for every
    design and content decision made during file generation.

    The knowledge file is built from four mandatory inputs:

    1. Original user intent — the exact prompt that initiated the workflow, preserved
       in full so the presentation's tone, goals, and specific requirements are never lost.

    2. Complete slide-by-slide content plan — the full per-slide plan from Step 1:
       all titles, bullet points, table/chart/image flags, and structural decisions.

    3. Deep template design language — a thorough analysis of every single slide layout
       within the template file: font families, font sizes, font weights, text colors,
       background colors, layout grids, placeholder positions, margin patterns, shape
       styles, line styles, color palette, accent colors, icon usage, and all recurring
       design motifs. Nothing is guessed or approximated.

    4. AI-generated image assets — full awareness of every image produced in Steps 2-3:
       which slide each targets, pixel dimensions, content type, aspect ratio, and the
       intended position within each slide layout.

    Only after the knowledge file is assembled does the actual PPTX construction begin.
    Handles image placeholder detection on each template layout so that generated images
    are inserted into dedicated picture placeholders when available.

    This step:
    1. Builds the comprehensive assembly knowledge file (single source of truth)
    2. Opens the template and generated presentation
    3. Creates slides from template layouts with extracted content
    4. Inserts AI-generated images into picture placeholders (or as free-floating)
    5. Saves the final output
    """
    global VERBOSE
    VERBOSE = session_state.get("verbose", VERBOSE)
    _step_start = time.time()

    print("\n" + "=" * 60)
    print("Step 4: Assembling final presentation with template...")
    print("=" * 60)

    template_path = session_state.get("template_path", "")
    generated_file = session_state.get("generated_file", "")
    output_path = session_state.get("output_path", "presentation_from_template.pptx")
    generated_images = session_state.get("generated_images", {})

    if not generated_file or not os.path.isfile(generated_file):
        return StepOutput(
            content="Error: Generated file not found.",
            success=False,
            stop=True,
        )

    print("Template: %s" % template_path)
    print("Generated: %s" % generated_file)
    print("Output: %s" % output_path)

    # Open the generated presentation
    generated_prs = Presentation(generated_file)
    generated_slides = list(generated_prs.slides)
    total_slides = len(generated_slides)

    if VERBOSE:
        print("[VERBOSE] Generated presentation has %d slides" % total_slides)

    if total_slides == 0:
        shutil.copy2(template_path, output_path)
        return StepOutput(
            content="Warning: No slides found. Copied template as-is.",
            success=True,
        )

    # Create output from template
    shutil.copy2(template_path, output_path)
    output_prs = Presentation(output_path)
    slide_width = output_prs.slide_width
    slide_height = output_prs.slide_height

    # Extract template styles for application to new elements
    template_style = None
    try:
        template_prs_for_styles = Presentation(template_path)
        template_style = _extract_template_styles(template_prs_for_styles)
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Template style extraction failed: %s" % str(e))

    # -----------------------------------------------------------------------
    # Build the comprehensive assembly knowledge file BEFORE constructing
    # any slide. This is the single source of truth that governs every
    # design and content decision: fonts, colors, layout, placeholders,
    # image placements, and all recurring design motifs from the template.
    # -----------------------------------------------------------------------
    print("  Building assembly knowledge file (template deep analysis)...")
    template_analysis: dict = {}
    try:
        template_prs_for_analysis = Presentation(template_path)
        template_analysis = _analyze_template_in_depth(template_prs_for_analysis)
        print(
            "  Knowledge file: %d layouts analyzed, %d accent color(s), "
            "heading font '%s', body font '%s'."
            % (
                len(template_analysis.get("layouts", [])),
                len(
                    [
                        c
                        for c in template_analysis.get("theme", {}).get(
                            "accent_colors", []
                        )
                        if c
                    ]
                ),
                template_analysis.get("design_language_summary", {}).get(
                    "heading_font", "?"
                ),
                template_analysis.get("design_language_summary", {}).get(
                    "body_font", "?"
                ),
            )
        )
    except Exception as e:
        print("[WARNING] Template deep analysis failed (will continue): %s" % str(e))
        if VERBOSE:
            traceback.print_exc()

    try:
        assembly_knowledge = _build_assembly_knowledge_file(
            session_state, template_analysis
        )
        session_state["assembly_knowledge"] = assembly_knowledge
        if VERBOSE:
            print(
                "[VERBOSE] Assembly knowledge file built — %d slides, %d AI image(s)"
                % (
                    assembly_knowledge["input_2_content_plan"].get("total_slides", 0),
                    assembly_knowledge["input_4_image_assets"].get(
                        "total_generated_images", 0
                    ),
                )
            )
    except Exception as e:
        print(
            "[WARNING] Knowledge file construction failed (will continue): %s" % str(e)
        )
        if VERBOSE:
            traceback.print_exc()
        session_state["assembly_knowledge"] = {}

    if VERBOSE:
        print(
            "[VERBOSE] Template layouts available: %s"
            % [layout.name for layout in output_prs.slide_layouts]
        )

    # Remove existing slides from template
    while len(output_prs.slides._sldIdLst) > 0:
        sldId = output_prs.slides._sldIdLst[0]
        rId = sldId.get(
            etree.QName(
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
                "id",
            )
        )
        if rId is not None:
            output_prs.part.drop_rel(rId)
        output_prs.slides._sldIdLst.remove(sldId)

    print("Cleared template slides. Building final presentation...")

    # For each generated slide, create a template-styled slide
    for idx, gen_slide in enumerate(generated_slides):
        content = _extract_slide_content(gen_slide)

        # Look up generated image for this slide.
        # The keys may be int or str depending on how they were stored.
        gen_img = generated_images.get(idx) or generated_images.get(str(idx))

        # Template assembly guard: suppress any generated image if the slide carries
        # native data visualizations. Three independent detection layers ensure coverage
        # across all Tiers (Tier 1: Claude PPTX skill native charts/tables,
        # Tier 2: python-pptx native ChartData objects, Tier 3: infographic shapes
        # not detectable as charts/tables but flagged via storyboard has_data_vis):
        #
        #   Layer 1 — PPTX shape inspection: content.charts / content.tables
        #             (populated by _extract_slide_content; catches Tier 1 & 2)
        #   Layer 2 — slides_data metadata: has_data_vis / has_chart / has_table
        #             (enriched in step_process_chunks from storyboard; catches Tier 3)
        #   Layer 3 — visual_suggestion keyword scan (belt-and-suspenders)
        #
        # Inserting an external AI photo alongside any native data visual would
        # collide with the visual region and degrade slide clarity.
        _slides_data_lookup = {
            s.get("index", i): s
            for i, s in enumerate(session_state.get("slides_data", []))
        }
        _slide_meta_asm = _slides_data_lookup.get(idx, {})
        _DATA_VIS_KW_ASM = ("chart", "table", "infographic", "diagram", "graph")
        _slide_has_data_vis = (
            bool(content.charts)
            or bool(content.tables)
            or _slide_meta_asm.get("has_chart", False)
            or _slide_meta_asm.get("has_table", False)
            or _slide_meta_asm.get("has_data_vis", False)
            or any(
                _kw in _slide_meta_asm.get("visual_suggestion", "").lower()
                for _kw in _DATA_VIS_KW_ASM
            )
        )
        if gen_img is not None and _slide_has_data_vis:
            print(
                "  Slide %d: Has native data visualization (chart/table/infographic) — "
                "suppressing generated image to preserve data visualization."
                % (idx + 1)
            )
            gen_img = None

        # Find layout FIRST so we can detect picture placeholders before
        # deciding whether to add a free-floating image.
        content_mix = _classify_content_mix(
            content, has_generated_image=gen_img is not None
        )
        layout = _find_best_layout(
            output_prs,
            idx,
            total_slides,
            content_mix=content_mix,
            has_generated_image=gen_img is not None,
        )
        if VERBOSE:
            summary = _layout_placeholder_summary(layout)
            print(
                "[VERBOSE] Slide %d chose layout '%s' placeholders: %s"
                % (idx + 1, layout.name, summary)
            )

        # Detect picture placeholders on the chosen template layout.
        # Uses both PP_PLACEHOLDER enum and XML-level fallback detection.
        for ph in layout.placeholders:
            ph_fmt = ph.placeholder_format
            if ph_fmt is not None:
                is_pic_ph = False
                # Strategy 1: int comparison with raw OOXML value 18
                try:
                    if ph_fmt.type is not None and (
                        int(ph_fmt.type) == 18 or str(ph_fmt.type) == "PICTURE (18)"
                    ):
                        is_pic_ph = True
                except (ValueError, TypeError):
                    pass
                # Strategy 2: PP_PLACEHOLDER enum
                if not is_pic_ph:
                    try:
                        if ph_fmt.type == PP_PLACEHOLDER.PICTURE:
                            is_pic_ph = True
                    except Exception as e:
                        if VERBOSE:
                            print("[VERBOSE] Exception suppressed: %s" % str(e))
                # Strategy 3: XML-level fallback
                if not is_pic_ph:
                    try:
                        nsmap = {
                            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                        }
                        ph_elem = ph._element.find(".//p:ph", nsmap)
                        if ph_elem is not None and ph_elem.get("type") in (
                            "pic",
                            "clipArt",
                        ):
                            is_pic_ph = True
                    except Exception as e:
                        if VERBOSE:
                            print("[VERBOSE] Exception suppressed: %s" % str(e))
                if is_pic_ph:
                    content.has_image_placeholder = True
                    if ph_fmt.idx not in content.image_placeholder_indices:
                        content.image_placeholder_indices.append(ph_fmt.idx)

        # Only add as free-floating picture if the layout has NO picture
        # placeholders. When picture placeholders exist, the image will be
        # inserted into the placeholder by _populate_slide instead.
        if gen_img is not None and not content.has_image_placeholder:
            content.images.append(
                ImageData(
                    blob=gen_img,
                    left=0,
                    top=0,
                    width=int(Inches(8.0)),
                    height=int(Inches(4.5)),
                    content_type="image/png",
                )
            )

        visual_info = []
        if content.tables:
            visual_info.append("%d table(s)" % len(content.tables))
        if content.images:
            visual_info.append("%d image(s)" % len(content.images))
        if content.charts:
            visual_info.append("%d chart(s)" % len(content.charts))
        if content.has_image_placeholder:
            visual_info.append("img placeholder(s)")
        visual_str = ", ".join(visual_info) if visual_info else "text only"

        print(
            "  Slide %d: layout '%s' | title: '%s' | %s"
            % (
                idx + 1,
                layout.name,
                content.title[:40] if content.title else "",
                visual_str,
            )
        )

        new_slide = output_prs.slides.add_slide(layout)
        # Pass generated image bytes and source dimensions to _populate_slide.
        # src_slide_width/height enable proportional shape rescaling (P1-1).
        # footer_* params enable footer standardization (P1-2).
        _populate_slide(
            new_slide,
            content,
            slide_width,
            slide_height,
            generated_image_bytes=gen_img,
            template_style=template_style,
            src_slide_width=session_state.get("src_slide_width", 0),
            src_slide_height=session_state.get("src_slide_height", 0),
            footer_text=session_state.get("footer_text", ""),
            date_text=session_state.get("date_text", ""),
            show_slide_number=session_state.get("show_slide_numbers", False),
        )

    # Clean up empty placeholders and hardcoded contrast issues
    clean_presentation_visual_noise_and_contrast(output_prs)
    output_prs.save(output_path)
    print("\nSaved final presentation: %s" % output_path)

    elapsed = time.time() - _step_start
    print("[TIMING] Step 4 Template Assembly: completed in %.2fs" % elapsed)
    return StepOutput(
        content="Presentation saved to %s (%d slides)" % (output_path, total_slides),
        success=True,
    )


# ---------------------------------------------------------------------------
# Step 5 (Optional): Visual Quality Review
# ---------------------------------------------------------------------------


def _render_pptx_to_images(pptx_path: str, output_dir: str) -> list:
    """Render all slides to per-slide PNG images.

    Uses a two-step pipeline for reliable per-slide rendering:
      1. PPTX → PDF via LibreOffice headless
      2. PDF  → per-page PNGs via pdftoppm (poppler-utils)

    Falls back to direct PPTX→PNG (single image) if pdftoppm is unavailable,
    but logs a loud warning since this only produces the first slide.

    Returns a sorted list of PNG file paths (one per slide, in slide order).
    Raises RuntimeError if LibreOffice is not available or rendering fails.

    Args:
        pptx_path:  Path to the .pptx file to render.
        output_dir: Directory to write PNG files into.

    Returns:
        Sorted list of PNG file paths.
    """
    import glob
    import shutil as _shutil
    import subprocess

    lo_cmd = _shutil.which("libreoffice") or _shutil.which("soffice")
    if not lo_cmd:
        raise RuntimeError(
            "LibreOffice not found. Install with: apt-get install libreoffice"
        )

    # Count expected slides for validation
    expected_slides = 0
    try:
        _count_prs = Presentation(pptx_path)
        expected_slides = len(_count_prs.slides)
        print("  [RENDER] PPTX has %d slide(s) to render." % expected_slides)
    except Exception:
        pass

    # --- Strategy 1: PPTX → PDF → per-page PNGs (preferred) ---
    pdftoppm_cmd = _shutil.which("pdftoppm")
    if pdftoppm_cmd:
        print("  [RENDER] Using PPTX→PDF→PNG pipeline (pdftoppm available).")
        # Step 1: Convert PPTX to PDF
        pdf_result = subprocess.run(
            [
                lo_cmd,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                output_dir,
                pptx_path,
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )
        if pdf_result.returncode != 0:
            print(
                "  [RENDER WARNING] PDF conversion failed (exit %d): %s"
                % (pdf_result.returncode, pdf_result.stderr[:200])
            )
            print("  [RENDER] Falling back to direct PNG conversion.")
        else:
            base = os.path.splitext(os.path.basename(pptx_path))[0]
            pdf_path = os.path.join(output_dir, base + ".pdf")
            if os.path.isfile(pdf_path):
                # Step 2: Convert PDF pages to PNGs via pdftoppm
                png_prefix = os.path.join(output_dir, "slide")
                ppm_result = subprocess.run(
                    [
                        pdftoppm_cmd,
                        "-png",
                        "-r",
                        "150",  # 150 DPI — good quality for vision review
                        pdf_path,
                        png_prefix,
                    ],
                    capture_output=True,
                    text=True,
                    timeout=120,
                )
                if ppm_result.returncode == 0:
                    # pdftoppm names files: slide-01.png, slide-02.png, ...
                    pngs = sorted(glob.glob(os.path.join(output_dir, "slide-*.png")))
                    if pngs:
                        if expected_slides > 0 and len(pngs) != expected_slides:
                            print(
                                "  [RENDER WARNING] Expected %d slides but rendered %d PNGs."
                                % (expected_slides, len(pngs))
                            )
                        else:
                            print(
                                "  [RENDER] Successfully rendered %d per-slide PNG(s) via PDF pipeline."
                                % len(pngs)
                            )
                        # Clean up intermediate PDF
                        try:
                            os.remove(pdf_path)
                        except OSError:
                            pass
                        return pngs
                else:
                    print(
                        "  [RENDER WARNING] pdftoppm failed (exit %d): %s"
                        % (ppm_result.returncode, ppm_result.stderr[:200])
                    )
                # Clean up intermediate PDF on failure
                try:
                    os.remove(pdf_path)
                except OSError:
                    pass

    # --- Strategy 2: Direct PPTX → PNG (fallback — single image only) ---
    print(
        "  [RENDER WARNING] pdftoppm not available — falling back to direct PNG "
        "conversion. This produces ONLY 1 image (first slide), not per-slide images!"
    )
    if expected_slides > 1:
        print(
            "  [RENDER WARNING] Visual review will only inspect 1 of %d slides! "
            "Install poppler-utils for full per-slide rendering: "
            "sudo apt-get install poppler-utils" % expected_slides
        )

    result = subprocess.run(
        [
            lo_cmd,
            "--headless",
            "--convert-to",
            "png",
            "--outdir",
            output_dir,
            pptx_path,
        ],
        capture_output=True,
        text=True,
        timeout=120,
    )
    if result.returncode != 0:
        raise RuntimeError(
            "LibreOffice rendering failed (exit %d): %s"
            % (result.returncode, result.stderr[:300])
        )

    # LibreOffice output naming varies by version:
    #   "<base>-slide-N.png" (most versions) or "<base>N.png" (some older versions)
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    for pattern in [
        os.path.join(output_dir, "%s-slide-*.png" % base),
        os.path.join(output_dir, "%s*.png" % base),
    ]:
        pngs = sorted(glob.glob(pattern))
        if pngs:
            return pngs

    # Final fallback: collect ALL .png files in the output directory sorted by
    # name.  This covers any LibreOffice naming variant not matched by the
    # patterns above (e.g. zero-padded "<base>001.png" or locale-specific names).
    all_pngs = sorted(glob.glob(os.path.join(output_dir, "*.png")))
    return all_pngs


# Slide quality review agent — acts as a senior UI/UX designer with deep knowledge
# of presentation design, visual hierarchy, typography, and brand consistency.
# Instantiated at module level; only invoked when --visual-review is enabled.
# slide_quality_reviewer agent is now loaded from the agents/ package via get_agents().
# Stored in session_state["agents"]["slide_quality_reviewer"].
# See agents/claude_agents.py (or openai/gemini variants) for the definition.


def _apply_accent_color_to_title(slide, accent_color_hex: str) -> bool:
    """Apply a template accent color to the slide title placeholder text.

    Finds the title placeholder (idx=0) and sets all run colors to accent_color_hex.
    Returns True if any runs were modified.
    """
    from pptx.dml.color import RGBColor

    modified = False
    try:
        rgb = RGBColor.from_string(accent_color_hex)
    except Exception:
        return False

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0 and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        run.font.color.rgb = rgb
                        modified = True
                    except Exception:
                        pass
            break
    return modified


def _increase_title_font_size(slide, target_pt: int = 28) -> bool:
    """Increase the title placeholder font size to at least target_pt.

    Only increases — never reduces below the existing size.
    Returns True if any runs were modified.
    """
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    target_sz = target_pt * 100  # OOXML hundredths of a point
    modified = False

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0 and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    rPr = run._r.find(ns_a + "rPr")
                    if rPr is not None:
                        sz = rPr.get("sz")
                        if sz and int(sz) < target_sz:
                            rPr.set("sz", str(target_sz))
                            modified = True
                    else:
                        # No explicit size — set it
                        rPr = etree.SubElement(run._r, ns_a + "rPr")
                        rPr.set("sz", str(target_sz))
                        modified = True
            break
    return modified


def _apply_body_accent_border(slide, accent_color_hex: str) -> bool:
    """Add a colored left-border accent line to the body text placeholder.

    Creates a thin vertical rectangle using the accent color on the left edge
    of the body placeholder to add visual structure to bland text-only slides.
    Returns True if the shape was added.
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Emu

    try:
        rgb = RGBColor.from_string(accent_color_hex)
    except Exception:
        return False

    # Find body placeholder to anchor the border to
    body_ph = None
    for shape in slide.placeholders:
        ph_idx = shape.placeholder_format.idx
        if ph_idx == 1 and shape.has_text_frame:
            body_ph = shape
            break
    if body_ph is None:
        return False

    try:
        # Add a 6pt-wide vertical rectangle on the left edge of the body placeholder
        border_width = Emu(76200)  # 6pt in EMU
        border_left = body_ph.left - border_width - Emu(38100)  # slight gap
        border_top = body_ph.top
        border_height = body_ph.height

        # Clamp to slide left edge
        if border_left < 0:
            border_left = Emu(114300)  # 9pt from left edge

        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            border_left,
            border_top,
            border_width,
            border_height,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb
        bar.line.fill.background()  # no outline
        return True
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Accent border add failed: %s" % str(e))
        return False


def _fix_spacing(slide, slide_width: int, slide_height: int) -> bool:
    """Failsafe spacing fix: clamp non-placeholder shapes that leak outside safe margins.

    The safe margin is 5% of slide width/height on each edge.  Any non-placeholder
    shape whose bounding box extends beyond that margin is nudged back inside
    by the minimum amount needed, capped at 5% of the slide dimension (to avoid
    aggressive repositioning).

    Only moves shapes — never resizes them, and never touches placeholders.
    Returns True if any shapes were moved.
    """
    MARGIN = 0.05  # 5% safe margin on each edge
    MAX_MOVE = 0.05  # maximum single-axis shift as fraction of slide dimension

    safe_left = int(slide_width * MARGIN)
    safe_top = int(slide_height * MARGIN)
    safe_right = int(slide_width * (1 - MARGIN))
    safe_bottom = int(slide_height * (1 - MARGIN))
    max_move_x = int(slide_width * MAX_MOVE)
    max_move_y = int(slide_height * MAX_MOVE)

    modified = False
    for shape in list(slide.shapes):
        try:
            if shape.is_placeholder:
                continue
            left, top = shape.left, shape.top
            w, h = shape.width, shape.height
            new_left, new_top = left, top

            # Clamp left edge
            if left < safe_left:
                delta = safe_left - left
                if delta <= max_move_x:
                    new_left = safe_left

            # Clamp right edge
            if new_left + w > safe_right:
                delta = (new_left + w) - safe_right
                if delta <= max_move_x:
                    new_left = max(safe_left, safe_right - w)

            # Clamp top edge
            if top < safe_top:
                delta = safe_top - top
                if delta <= max_move_y:
                    new_top = safe_top

            # Clamp bottom edge
            if new_top + h > safe_bottom:
                delta = (new_top + h) - safe_bottom
                if delta <= max_move_y:
                    new_top = max(safe_top, safe_bottom - h)

            if new_left != left or new_top != top:
                shape.left = new_left
                shape.top = new_top
                modified = True
                if VERBOSE:
                    print(
                        "[VERBOSE] Spacing fix: shape moved from (%d,%d) to (%d,%d)"
                        % (left, top, new_left, new_top)
                    )
        except Exception:
            continue

    return modified


def _fix_alignment(slide, slide_width: int) -> bool:
    """Failsafe alignment fix: snap minor left-edge outliers to the majority left edge.

    Collects the left position of all non-placeholder content shapes.  Uses a
    majority-vote: the most common left edge within a 3% slide-width tolerance
    band becomes the 'anchor'.  Any outlier within 2% of slide width from the
    anchor is snapped to it.

    Only moves shapes horizontally — never vertically, never resizes, never
    touches placeholders.  Returns True if any shapes were moved.
    """
    TOLERANCE = int(slide_width * 0.03)  # 3% band for majority-vote grouping
    MAX_SNAP = int(slide_width * 0.02)  # max 2% snap

    # Collect candidate left edges
    left_edges = []
    candidate_shapes = []
    for shape in list(slide.shapes):
        try:
            if shape.is_placeholder:
                continue
            if not getattr(shape, "left", None):
                continue
            left_edges.append(shape.left)
            candidate_shapes.append(shape)
        except Exception:
            continue

    if len(left_edges) < 2:
        return False

    # Find the majority anchor: the left edge value with the most neighbors within tolerance
    best_anchor = None
    best_count = 0
    for candidate in left_edges:
        count = sum(1 for e in left_edges if abs(e - candidate) <= TOLERANCE)
        if count > best_count:
            best_count = count
            best_anchor = candidate

    if best_anchor is None or best_count < 2:
        return False

    # Snap outliers that are close enough
    modified = False
    for shape in candidate_shapes:
        try:
            if abs(shape.left - best_anchor) <= MAX_SNAP and shape.left != best_anchor:
                if VERBOSE:
                    print(
                        "[VERBOSE] Alignment fix: shape left %d -> %d (anchor)"
                        % (shape.left, best_anchor)
                    )
                shape.left = best_anchor
                modified = True
        except Exception:
            continue

    return modified


def _enrich_slide_visually(
    slide,
    slide_width: int,
    slide_height: int,
    accent_color_hex: str,
    enrich_type: str = "enrich_accent_strip",
) -> bool:
    """Add native PPTX visual enrichment elements using template colors.

    No images are generated — enrichment is achieved purely through python-pptx
    rectangle shapes, fills, and borders.  Each enrichment type uses the
    template's primary accent color.

    enrich_type options:
      enrich_header_bar    - Full-width accent-color bar at slide top (8% height).
                             Gives every slide a branded header band.
      enrich_title_card    - Lightly-tinted card behind the title area with a left
                             accent border.  Makes the title visually pop.
      enrich_divider       - Thin horizontal accent rule between the title zone
                             and the body zone (at 25% slide height).
      enrich_accent_strip  - Thin vertical accent strip on the left edge
                             (1.5% width, full height).  Minimal brand presence.

    Returns True if any shapes were added.
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Emu

    try:
        accent_rgb = RGBColor.from_string(accent_color_hex)
    except Exception:
        return False

    def _hex_to_tinted(hex_color: str, tint: float = 0.85) -> RGBColor:
        """Return a lightened (tinted) version of a hex color (0=white, 1=original)."""
        r, g, b = _hex_to_rgb(hex_color)
        r2 = int(r + (255 - r) * (1 - tint))
        g2 = int(g + (255 - g) * (1 - tint))
        b2 = int(b + (255 - b) * (1 - tint))
        return RGBColor(r2, g2, b2)

    try:
        if enrich_type == "enrich_header_bar":
            # Full-width accent-color bar at the very top of the slide
            bar_height = int(slide_height * 0.08)
            bar = slide.shapes.add_shape(1, Emu(0), Emu(0), slide_width, bar_height)
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent_rgb
            bar.line.fill.background()
            return True

        elif enrich_type == "enrich_title_card":
            # Lightly tinted card behind the title area (top 22% of slide)
            # with a 6pt left accent border
            card_height = int(slide_height * 0.22)
            card = slide.shapes.add_shape(
                1,
                Emu(int(slide_width * 0.02)),
                Emu(int(slide_height * 0.02)),
                slide_width - Emu(int(slide_width * 0.04)),
                card_height,
            )
            tinted = _hex_to_tinted(accent_color_hex, tint=0.15)
            card.fill.solid()
            card.fill.fore_color.rgb = tinted
            card.line.color.rgb = accent_rgb
            card.line.width = Emu(76200)  # 6pt
            return True

        elif enrich_type == "enrich_divider":
            # Thin horizontal rule at 25% slide height
            divider_y = int(slide_height * 0.25)
            divider_height = Emu(50800)  # 4pt thick
            divider_x = int(slide_width * 0.05)
            divider_w = int(slide_width * 0.90)
            div = slide.shapes.add_shape(
                1, divider_x, divider_y, divider_w, divider_height
            )
            div.fill.solid()
            div.fill.fore_color.rgb = accent_rgb
            div.line.fill.background()
            return True

        elif enrich_type == "enrich_accent_strip":
            # Thin vertical strip on the far left edge
            strip_width = int(slide_width * 0.015)
            strip = slide.shapes.add_shape(1, Emu(0), Emu(0), strip_width, slide_height)
            strip.fill.solid()
            strip.fill.fore_color.rgb = accent_rgb
            strip.line.fill.background()
            return True

    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Visual enrichment (%s) failed: %s" % (enrich_type, str(e)))

    return False


def _enforce_typography_hierarchy(slide) -> bool:
    """Gap 4: Enforce font size hierarchy between title and body placeholders.

    Reads the actual font sizes from title (idx=0) and body (idx=1) placeholders.
    If the title is not at least 6pt larger than the body, increases the title
    to max(body_size + 8, 24) points, capped at Pt(36).

    Returns True if any sizes were changed.
    """
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    modified = False

    title_shape = None
    body_shape = None
    for shape in slide.placeholders:
        ph_idx = shape.placeholder_format.idx
        if ph_idx == 0 and shape.has_text_frame:
            title_shape = shape
        elif ph_idx == 1 and shape.has_text_frame:
            body_shape = shape

    if title_shape is None or body_shape is None:
        return False

    # Collect body font size
    body_sz = 0
    for para in body_shape.text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.find(ns_a + "rPr")
            if rPr is not None:
                sz = rPr.get("sz")
                if sz:
                    body_sz = max(body_sz, int(sz))

    if body_sz == 0:
        return False  # Cannot determine body size — skip

    # Collect title font size
    title_sz = 0
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.find(ns_a + "rPr")
            if rPr is not None:
                sz = rPr.get("sz")
                if sz:
                    title_sz = max(title_sz, int(sz))

    if title_sz == 0:
        return False

    # Enforce: title must be at least body + 6pt
    MIN_HIERARCHY_GAP = 600  # 6pt in OOXML hundredths
    MAX_TITLE_SZ = 3600  # cap at 36pt

    if title_sz < body_sz + MIN_HIERARCHY_GAP:
        target_sz = min(MAX_TITLE_SZ, max(body_sz + 800, 2400))  # +8pt, min 24pt
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                rPr = run._r.find(ns_a + "rPr")
                if rPr is not None:
                    rPr.set("sz", str(target_sz))
                    modified = True
        if VERBOSE:
            print(
                "[VERBOSE] Typography hierarchy: title %dpt -> %dpt (body was %dpt)"
                % (title_sz // 100, target_sz // 100, body_sz // 100)
            )

    return modified


def _apply_accent_color_to_body(slide, accent_color_hex: str) -> bool:
    """Gap 5: Apply template accent color to first-level (level=0) body text runs.

    Applies the primary accent color to the first paragraph of the body
    placeholder only — making the lead sentence or subtitle visually distinct
    without changing all body text (which should remain readable).

    Returns True if any runs were modified.
    """
    from pptx.dml.color import RGBColor

    modified = False
    try:
        rgb = RGBColor.from_string(accent_color_hex)
    except Exception:
        return False

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1 and shape.has_text_frame:
            # Apply accent only to the first paragraph at level 0
            for para in shape.text_frame.paragraphs:
                if para.level == 0:
                    for run in para.runs:
                        try:
                            run.font.color.rgb = rgb
                            modified = True
                        except Exception:
                            pass
                    break  # Only first paragraph
            break

    return modified


def _fix_paragraph_alignment_body(slide, alignment: str = "l") -> bool:
    """Gap 1 (Step 5): Standardize paragraph alignment for body text placeholder.

    Sets the <a:pPr algn="..."> attribute on all paragraphs in the body
    placeholder (idx=1).  Default is left alignment ('l').

    alignment options: 'l' (left), 'ctr' (center), 'r' (right), 'just' (justify)

    Returns True if any paragraphs were updated.
    """
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    modified = False

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1 and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                p_elem = para._p
                pPr = p_elem.find(ns_a + "pPr")
                if pPr is None:
                    pPr = etree.SubElement(p_elem, ns_a + "pPr")
                    p_elem.insert(0, pPr)
                existing = pPr.get("algn", "")
                if existing != alignment:
                    pPr.set("algn", alignment)
                    modified = True
            break

    return modified


def _apply_visual_corrections(
    pptx_path: str,
    reports: list,
    template_style: "TemplateStyle | None" = None,
) -> bool:
    """Apply programmatic corrections based on the UI/UX designer quality reports.

    Handles both structural fixes (critical severity) and design enrichment
    (moderate severity) using the expanded programmatic_fix vocabulary.

    Structural fixes (re-invoke existing pipeline functions):
      - increase_contrast       : re-runs _ensure_text_contrast()
      - clear_placeholder       : re-runs _clear_unused_placeholders() + _remove_empty_textboxes()
      - remove_element          : same as clear_placeholder for ghost_text / empty_placeholder
      - reduce_font_size        : conservative 15% rPr.sz reduction above Pt(10)

    Design enrichment fixes (new, applied to moderate issues):
      - apply_accent_color_title : applies template primary accent color to title text
      - increase_title_font_size : increases title font to Pt(28) if below that
      - apply_body_accent_border : adds a colored vertical accent bar left of body text

    Not implemented (require AI-generated content or layout redesign):
      - reposition_element       : bounding-box matching risk
      - any content changes      : semantic, not structural

    Args:
        pptx_path:      Path to the .pptx file to correct (overwritten in place).
        reports:        List of SlideQualityReport from the UI/UX designer agent.
        template_style: Optional template styles for contrast/color corrections.

    Returns:
        True if any corrections were applied and the file was re-saved, else False.
    """
    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    modified = False

    # Resolve primary accent color from template for enrichment fixes
    accent_color = ""
    if template_style and template_style.theme.accent_colors:
        accent_color = template_style.theme.accent_colors[0]
    elif template_style and template_style.theme.dk1:
        accent_color = template_style.theme.dk1

    for report in reports:
        if report.slide_index >= len(slides):
            continue
        slide = slides[report.slide_index]

        # Apply critical fixes (structural defects — always correct)
        # Apply moderate design enrichment fixes where programmatic_fix is set
        actionable_issues = [
            i
            for i in report.issues
            if i.severity in ("critical", "moderate") and i.programmatic_fix != "none"
        ]

        for issue in actionable_issues:
            fix = issue.programmatic_fix

            try:
                if fix == "increase_contrast" and template_style:
                    _ensure_text_contrast(slide, template_style)
                    modified = True
                    if VERBOSE:
                        print(
                            "[VERBOSE] Slide %d: applied increase_contrast"
                            % report.slide_index
                        )

                elif fix in ("remove_element", "clear_placeholder"):
                    if issue.issue_type in ("ghost_text", "empty_placeholder"):
                        _clear_unused_placeholders(slide, populated_indices=set())
                        _remove_empty_textboxes(slide)
                        modified = True
                        if VERBOSE:
                            print(
                                "[VERBOSE] Slide %d: cleared ghost text / empty placeholders"
                                % report.slide_index
                            )

                elif fix == "reduce_font_size":
                    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                    for shape in slide.shapes:
                        if not getattr(shape, "has_text_frame", False):
                            continue
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                rPr = run._r.find(ns_a + "rPr")
                                if rPr is not None:
                                    sz = rPr.get("sz")
                                    if sz and int(sz) > 1000:  # > Pt(10)
                                        new_sz = max(1000, int(int(sz) * 0.85))
                                        rPr.set("sz", str(new_sz))
                    modified = True
                    if VERBOSE:
                        print(
                            "[VERBOSE] Slide %d: reduced font sizes by 15%%"
                            % report.slide_index
                        )

                elif fix == "apply_accent_color_title" and accent_color:
                    if _apply_accent_color_to_title(slide, accent_color):
                        modified = True
                        if VERBOSE:
                            print(
                                "[VERBOSE] Slide %d: applied accent color #%s to title"
                                % (report.slide_index, accent_color)
                            )

                elif fix == "increase_title_font_size":
                    if _increase_title_font_size(slide, target_pt=28):
                        modified = True
                        if VERBOSE:
                            print(
                                "[VERBOSE] Slide %d: increased title font to 28pt"
                                % report.slide_index
                            )

                elif fix == "apply_body_accent_border" and accent_color:
                    if _apply_body_accent_border(slide, accent_color):
                        modified = True
                        if VERBOSE:
                            print(
                                "[VERBOSE] Slide %d: added accent border bar #%s"
                                % (report.slide_index, accent_color)
                            )

                elif fix == "fix_spacing" and slide_width > 0 and slide_height > 0:
                    if _fix_spacing(slide, slide_width, slide_height):
                        modified = True
                        if VERBOSE:
                            print(
                                "[VERBOSE] Slide %d: spacing clamped to safe margins"
                                % report.slide_index
                            )

                elif fix == "fix_alignment" and slide_width > 0:
                    if _fix_alignment(slide, slide_width):
                        modified = True
                        if VERBOSE:
                            print(
                                "[VERBOSE] Slide %d: alignment snapped to majority left edge"
                                % report.slide_index
                            )

                elif fix.startswith("enrich_") and accent_color:
                    if slide_width > 0 and slide_height > 0:
                        if _enrich_slide_visually(
                            slide,
                            slide_width,
                            slide_height,
                            accent_color,
                            enrich_type=fix,
                        ):
                            modified = True
                            if VERBOSE:
                                print(
                                    "[VERBOSE] Slide %d: visual enrichment applied (%s)"
                                    % (report.slide_index, fix)
                                )

                elif fix == "enforce_typography_hierarchy":
                    if _enforce_typography_hierarchy(slide):
                        modified = True
                        if VERBOSE:
                            print(
                                "[VERBOSE] Slide %d: typography hierarchy enforced"
                                % report.slide_index
                            )

                elif fix == "apply_accent_color_body" and accent_color:
                    if _apply_accent_color_to_body(slide, accent_color):
                        modified = True
                        if VERBOSE:
                            print(
                                "[VERBOSE] Slide %d: accent color applied to body lead"
                                % report.slide_index
                            )

                elif fix == "fix_body_paragraph_alignment":
                    if _fix_paragraph_alignment_body(slide, alignment="l"):
                        modified = True
                        if VERBOSE:
                            print(
                                "[VERBOSE] Slide %d: body paragraph alignment set to left"
                                % report.slide_index
                            )

                elif fix == "resize_chart":
                    if _ensure_chart_fills_area(slide, slide_width, slide_height):
                        modified = True
                        if VERBOSE:
                            print(
                                "[VERBOSE] Slide %d: chart resized to fill visual area"
                                % report.slide_index
                            )

                elif fix == "add_data_labels":
                    if _ensure_chart_data_labels(slide) > 0:
                        modified = True
                        if VERBOSE:
                            print(
                                "[VERBOSE] Slide %d: chart data labels enabled"
                                % report.slide_index
                            )

                elif fix == "remove_duplicate_title":
                    if _remove_duplicate_titles(slide) > 0:
                        modified = True
                        if VERBOSE:
                            print(
                                "[VERBOSE] Slide %d: duplicate title removed"
                                % report.slide_index
                            )

            except Exception as e:
                if VERBOSE:
                    print(
                        "[VERBOSE] Correction skipped (slide %d, fix=%s): %s"
                        % (report.slide_index, fix, str(e))
                    )

    if modified:
        prs.save(pptx_path)
        # Clean up empty placeholders and hardcoded contrast issues
        clean_presentation_visual_noise_and_contrast(prs)
        prs.save(pptx_path)
        print("\nFallback presentation generation successful: %s" % pptx_path)
        return modified


def step_visual_quality_review(
    step_input: StepInput, session_state: Dict
) -> StepOutput:
    """Optional Step 5: Render slides, inspect with a senior UI/UX designer agent,
    and apply programmatic corrections for both structural defects and design issues.

    This step is non-blocking: any failure (LibreOffice unavailable, API error,
    timeout) silently returns success=True without modifying the output file.
    
    NOTE ON Pydantic Schema Parsing (Duck-Typing):
    Agent output validation utilizes duck-typing on `SlideQualityReport` rather 
    than strict `isinstance` checks. This resilient design resolves a Python 
    namespacing quirk where Pydantic models generated dynamically by the LLM 
    framework might not mathematically equate to `__main__.SlideQualityReport`, 
    preventing False-Negative review dropouts.

    Workflow:
    1. Render the final .pptx to PNGs using LibreOffice headless.
    2. Build template context (accent colors, fonts) and pass it to the agent.
    3. Send each slide image to Gemini 2.5 Flash (UI/UX designer persona) for
       both structural QA and design quality assessment.
    4. Apply programmatic corrections for critical AND moderate design issues.
    5. Apply design enrichment for bland slides (accent color, title sizing,
       accent border) when programmatic fixes are available.
    6. Log design scores and recommendations.
    7. Store a structured quality report in session_state['quality_report'].

    Correction scope:
    - CORRECTED (critical):  font overflow, text contrast, ghost text, empty placeholders
    - CORRECTED (moderate):  accent color on title, title font sizing, body accent border
    - DETECT + WARN ONLY:    blandness requiring AI-generated images, layout redesign
    - NOT ATTEMPTED:         shape repositioning, content changes
    """
    global VERBOSE
    VERBOSE = session_state.get("verbose", VERBOSE)
    _step_start = time.time()

    print("\n" + "=" * 60)
    print("Step 5 (Optional): UI/UX Design Review...")
    print("=" * 60)

    output_path = session_state.get("output_path", "")
    template_path = session_state.get("template_path", "")

    if not output_path or not os.path.isfile(output_path):
        print("  Skipped: output file not found at '%s'." % output_path)
        return StepOutput(
            content="Visual review skipped: output file not found.", success=True
        )

    import tempfile

    render_dir = tempfile.mkdtemp(prefix="pptx_vqa_")

    try:
        # --- Phase 1: Render slides to PNG ---
        print("  Rendering slides to PNG with LibreOffice...")
        try:
            slide_images = _render_pptx_to_images(output_path, render_dir)
            print("  Rendered %d slide(s)." % len(slide_images))
        except RuntimeError as e:
            print("  [WARNING] Rendering unavailable: %s" % str(e))
            print("  Skipping visual review (non-fatal).")
            return StepOutput(
                content="Visual review skipped: %s" % str(e), success=True
            )

        if not slide_images:
            print("  [WARNING] No slide images produced. Skipping visual review.")
            return StepOutput(
                content="Visual review skipped: no images produced.", success=True
            )

        # --- Phase 2: Extract template context for the agent ---
        template_style = None
        template_context = ""
        try:
            t_prs = Presentation(template_path)
            template_style = _extract_template_styles(t_prs)

            # Build a concise context string the agent can reference
            ctx_parts = []
            if template_style.theme.accent_colors:
                ctx_parts.append(
                    "Accent colors: %s"
                    % ", ".join(
                        "#%s" % c for c in template_style.theme.accent_colors[:3] if c
                    )
                )
            if template_style.theme.dk1:
                ctx_parts.append("Dark text color: #%s" % template_style.theme.dk1)
            if template_style.theme.major_font:
                ctx_parts.append("Heading font: %s" % template_style.theme.major_font)
            if template_style.theme.minor_font:
                ctx_parts.append("Body font: %s" % template_style.theme.minor_font)
            if template_style.title_font_family:
                ctx_parts.append("Title font: %s" % template_style.title_font_family)
            if ctx_parts:
                template_context = (
                    "Template design palette — " + "; ".join(ctx_parts) + "."
                )
        except Exception as e:
            if VERBOSE:
                print("[VERBOSE] Template style re-extraction failed: %s" % str(e))

        # --- Phase 3: Vision inspection with UI/UX designer agent ---
        # Build a per-slide data-vis lookup so the reviewer knows which slides
        # contain native charts/tables/infographics and must not be penalised
        # for the absence of an AI-generated photo.
        _DATA_VIS_KW_REVIEW = ("chart", "table", "infographic", "diagram", "graph")
        _slides_data_list = session_state.get("slides_data", [])
        _slides_data_map = {
            s.get("index", i): s for i, s in enumerate(_slides_data_list)
        }

        reports = []
        for idx, img_path in enumerate(slide_images):
            print("  Reviewing slide %d / %d..." % (idx + 1, len(slide_images)))
            try:
                with open(img_path, "rb") as f:
                    img_bytes = f.read()

                # Determine whether this slide carries a native data visualization.
                # When it does, the reviewer must not flag it for missing a photo/image.
                _s_meta = _slides_data_map.get(idx, {})
                _slide_has_data_vis = (
                    _s_meta.get("has_chart", False)
                    or _s_meta.get("has_table", False)
                    or _s_meta.get("has_data_vis", False)
                    or any(
                        _kw in _s_meta.get("visual_suggestion", "").lower()
                        for _kw in _DATA_VIS_KW_REVIEW
                    )
                )
                data_vis_flag = (
                    "[Data vis: chart/table/infographic present] "
                    if _slide_has_data_vis
                    else ""
                )

                # Include template context in the prompt so the agent can
                # reference specific accent colors and fonts in its feedback.
                # data_vis_flag tells the agent not to penalise the slide for
                # the intentional absence of an AI-generated photo/illustration.
                context_prefix = (
                    ("[Template context: %s] " % template_context)
                    if template_context
                    else ""
                )
                prompt = (
                    "%s%sYou are reviewing slide %d of %d in a professional "
                    "presentation. Assess both structural quality and design excellence. "
                    "Provide a design_score (1-10), identify all issues, and flag "
                    "whether the slide is visually bland."
                ) % (context_prefix, data_vis_flag, idx + 1, len(slide_images))

                # Use AgnoImage with raw bytes — Agno requires Image objects,
                # not raw dicts, for multimodal agent.run() calls.
                slide_image_obj = AgnoImage(
                    content=img_bytes,
                    mime_type="image/png",
                    format="png",
                )
                from agents import get_agents as _get_agents
                _provider = session_state.get("llm_provider", "claude")
                _slide_reviewer = _get_agents(_provider).get("slide_quality_reviewer")
                _slide_reviewer_fallback = _get_agents(_provider).get("slide_quality_reviewer_fallback")
                
                try:
                    response = _slide_reviewer.run(
                        prompt,
                        images=[slide_image_obj],
                        stream=False,
                    )
                except Exception as e:
                    print("[WARNING] Vision agent failed on primary agent: %s" % str(e))
                    if _slide_reviewer_fallback:
                        print("[FALLBACK TRIGGERED] Engaging fallback vision agent...")
                        try:
                            response = _slide_reviewer_fallback.run(
                                prompt,
                                images=[slide_image_obj],
                                stream=False,
                            )
                        except Exception as fallback_e:
                            print("[ERROR] Fallback vision agent failed: %s" % str(fallback_e))
                            response = None
                    else:
                        response = None

                if response and response.content:
                    content = response.content
                    # Duck-type: accept any BaseModel with the right fields,
                    # dict, or JSON string.  The agent's output_schema may be
                    # the SlideQualityReport from agents/_shared.py (a
                    # different Python class than the local one), so a plain
                    # isinstance() would always fail.
                    if hasattr(content, "model_dump") and hasattr(content, "design_score"):
                        report = SlideQualityReport.model_validate(content.model_dump())
                    elif isinstance(content, dict):
                        report = SlideQualityReport.model_validate(content)
                    elif isinstance(content, str):
                        import json as _json_parse
                        try:
                            report = SlideQualityReport.model_validate(
                                _json_parse.loads(content)
                            )
                        except (ValueError, Exception):
                            if VERBOSE:
                                print(
                                    "[VERBOSE] Could not parse string response "
                                    "as SlideQualityReport"
                                )
                            continue
                    else:
                        if VERBOSE:
                            print(
                                "[VERBOSE] Unexpected response type: %s"
                                % type(content).__name__
                            )
                        continue
                    report.slide_index = idx
                    reports.append(report)

                    # Log per-slide summary
                    score = getattr(report, "design_score", 0)
                    critical = [i for i in report.issues if i.severity == "critical"]
                    moderate = [i for i in report.issues if i.severity == "moderate"]
                    score_str = " [score: %d/10]" % score if score else ""
                    if critical:
                        print(
                            "    CRITICAL%s: %s"
                            % (score_str, [i.issue_type for i in critical])
                        )
                    elif moderate:
                        print(
                            "    moderate%s: %s"
                            % (score_str, [i.issue_type for i in moderate[:3]])
                        )
                    elif report.is_visually_bland:
                        print(
                            "    Bland%s — %s"
                            % (score_str, report.blandness_reason[:60])
                        )
                    else:
                        print("    OK%s." % score_str)

            except Exception as e:
                print(
                    "  [WARNING] Design review failed for slide %d: %s"
                    % (idx + 1, str(e))
                )
                if VERBOSE:
                    traceback.print_exc()

        # Warn if vision agent returned nothing parseable
        if not reports and len(slide_images) > 0:
            print(
                "  [WARNING] Vision agent returned 0 parseable reports for %d "
                "slides. Check model output schema compatibility." % len(slide_images)
            )

        # --- Phase 4: Apply corrections (critical + moderate design enrichment) ---
        total_critical = sum(
            1 for r in reports for i in r.issues if i.severity == "critical"
        )
        total_moderate = sum(
            1
            for r in reports
            for i in r.issues
            if i.severity == "moderate" and i.programmatic_fix != "none"
        )
        corrections_applied = False
        if total_critical > 0 or total_moderate > 0:
            print(
                "  Applying corrections (%d critical, %d moderate design fixes)..."
                % (total_critical, total_moderate)
            )
            try:
                corrections_applied = _apply_visual_corrections(
                    output_path, reports, template_style=template_style
                )
                if corrections_applied:
                    print("  Corrections saved: %s" % output_path)
                else:
                    print(
                        "  No corrections were applicable (all programmatic_fix='none')."
                    )
            except Exception as e:
                print("  [WARNING] Correction pass failed: %s" % str(e))
                if VERBOSE:
                    traceback.print_exc()
        else:
            print("  No corrections needed.")

        # --- Phase 5: Log blandness + recommendations ---
        bland_slides = [r for r in reports if r.is_visually_bland]
        if bland_slides:
            print(
                "\n  [DESIGN NOTE] %d slide(s) are visually bland and could benefit "
                "from AI-generated images or richer layout." % len(bland_slides)
            )
            for r in bland_slides:
                if r.blandness_reason:
                    print(
                        "    Slide %d: %s"
                        % (r.slide_index + 1, r.blandness_reason[:80])
                    )

        # --- Phase 6: Store quality report in session_state ---
        scores = [
            getattr(r, "design_score", 0)
            for r in reports
            if getattr(r, "design_score", 0) > 0
        ]
        avg_score = round(sum(scores) / len(scores), 1) if scores else 0.0

        all_recommendations = []
        for r in reports:
            for issue in r.issues:
                if issue.programmatic_fix == "none" and issue.description:
                    all_recommendations.append(
                        "Slide %d (%s): %s"
                        % (r.slide_index + 1, issue.issue_type, issue.description[:100])
                    )
        for r in bland_slides:
            if r.blandness_reason:
                all_recommendations.append(
                    "Slide %d is bland: %s" % (r.slide_index + 1, r.blandness_reason)
                )

        quality_report = PresentationQualityReport(
            slide_reports=reports,
            overall_pass=(
                all(r.overall_quality != "poor" for r in reports) and avg_score >= 6.0
            ),
            total_critical_issues=total_critical,
            average_design_score=avg_score,
            recommendations=all_recommendations[:10],  # cap at 10
        )
        session_state["quality_report"] = quality_report.model_dump()

        summary = (
            "UI/UX review: %d slides, avg design score %.1f/10, "
            "%d critical + %d moderate fixes, %d recommendations."
        ) % (
            len(reports),
            avg_score,
            total_critical,
            total_moderate,
            len(all_recommendations),
        )
        print("\n  %s" % summary)
        elapsed = time.time() - _step_start
        print("[TIMING] Step 5 Visual Quality Review: completed in %.2fs" % elapsed)
        return StepOutput(content=summary, success=True)

    except Exception as e:
        # Non-blocking: any unhandled failure returns success with warning.
        msg = "Visual review failed (non-fatal): %s" % str(e)
        print("  [WARNING] %s" % msg)
        if VERBOSE:
            traceback.print_exc()
        return StepOutput(content=msg, success=True)

    finally:
        import shutil as _shutil

        _shutil.rmtree(render_dir, ignore_errors=True)


# ---------------------------------------------------------------------------
# CLI and Main
# ---------------------------------------------------------------------------


def parse_args():
    """Parse command-line arguments."""
    parser = argparse.ArgumentParser(
        description="Workflow: Generate PowerPoint with Claude, add AI images, apply template."
    )
    parser.add_argument(
        "--template",
        "-t",
        default=None,
        help="Path to the .pptx template file (optional). "
        "When omitted, the raw Claude-generated presentation is used as output "
        "without any template styling applied.",
    )
    parser.add_argument(
        "--output",
        "-o",
        default="presentation_from_template.pptx",
        help="Output filename (default: presentation_from_template.pptx).",
    )
    parser.add_argument(
        "--prompt",
        "-p",
        default=None,
        help="Custom prompt for the presentation content.",
    )
    parser.add_argument(
        "--no-images",
        action="store_true",
        help="Skip AI image generation (Steps 2 and 3).",
    )
    parser.add_argument(
        "--no-stream",
        action="store_true",
        help="Disable streaming mode for Claude agent (more reliable for shorter prompts, "
        "but may timeout on complex presentations).",
    )
    parser.add_argument(
        "--verbose",
        "-v",
        action="store_true",
        help="Enable verbose/debug logging for troubleshooting.",
    )
    parser.add_argument(
        "--min-images",
        type=int,
        default=1,
        help="Minimum number of slides that must have AI-generated images "
        "(default: 1, 0 = let planner decide).",
    )
    parser.add_argument(
        "--footer-text",
        default="",
        help="Footer text to apply to all slides (idx=11 placeholder). "
        "Empty string (default) removes footer placeholders.",
    )
    parser.add_argument(
        "--date-text",
        default="",
        help="Date text for footer date placeholder (idx=10). "
        "Empty string (default) removes date placeholders.",
    )
    parser.add_argument(
        "--show-slide-numbers",
        action="store_true",
        help="Preserve slide number footer placeholders (idx=12) so PowerPoint "
        "can auto-render the current page number.",
    )
    parser.add_argument(
        "--visual-review",
        action="store_true",
        help="Enable optional Step 5: render slides to PNGs and inspect with Gemini "
        "vision for visual defects. Applies safe corrections for critical issues. "
        "Requires LibreOffice (install: apt-get install libreoffice). "
        "Non-blocking: skips silently if LibreOffice or the vision API is unavailable.",
    )
    parser.add_argument(
        "--llm-provider",
        choices=["claude", "openai", "gemini"],
        default="claude",
        help=(
            "LLM provider for swappable agents (image planner, visual reviewer). "
            "The Content Generator always uses Claude (PPTX skill). Default: claude."
        ),
    )
    return parser.parse_args()


if __name__ == "__main__":
    args = parse_args()

    VERBOSE = args.verbose

    if args.template is not None:
        if not os.path.isfile(args.template):
            print("Error: Template file not found: %s" % args.template)
            sys.exit(1)
        if not args.template.endswith(".pptx"):
            print("Error: Template file must be a .pptx file")
            sys.exit(1)

    if not os.getenv("ANTHROPIC_API_KEY"):
        raise ValueError(
            "ANTHROPIC_API_KEY environment variable not set "
            "(required: Content Generator agent always uses Claude)"
        )
    if args.llm_provider == "openai" and not os.getenv("OPENAI_API_KEY"):
        raise ValueError("OPENAI_API_KEY environment variable not set (required for --llm-provider openai)")
    if args.llm_provider == "gemini" and not os.getenv("GOOGLE_API_KEY"):
        raise ValueError("GOOGLE_API_KEY environment variable not set (required for --llm-provider gemini)")

    # Gemini API key validation guard for --visual-review
    # If the key is missing or blank, auto-disable visual review to avoid burning
    # time on guaranteed-to-fail Gemini API calls.
    google_key = os.getenv("GOOGLE_API_KEY", "")
    if args.visual_review and not google_key:
        print(
            "[WARNING] --visual-review requested but GOOGLE_API_KEY is not set. "
            "Visual review automatically disabled."
        )
        args.visual_review = False

    # Provider name is stored in session_state (a plain string, safely deep-copyable).
    # Agents are loaded lazily inside each step via get_agents() to avoid pickling errors.

    # Setup output directory
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_dir = os.path.join(script_dir, "output")
    os.makedirs(output_dir, exist_ok=True)

    output_path = args.output
    if not os.path.isabs(output_path):
        output_path = os.path.join(output_dir, output_path)

    # Default prompt
    prompt = args.prompt
    if not prompt:
        prompt = (
            "Create a 6-slide business presentation with the following structure.\n"
            "Rules: Use one clear title per slide, 4-6 concise bullet points for text slides,\n"
            "and include tables or charts where specified. Do not apply custom styling.\n"
            "\n"
            "Slide 1 - Title Slide:\n"
            "  Title: 'Strategic Overview 2026'\n"
            "  Subtitle: 'Annual Business Review and Forward Plan'\n"
            "\n"
            "Slide 2 - Market Analysis (text with bullets):\n"
            "  Title: 'Market Analysis'\n"
            "  Bullets: Key market trends, growth opportunities, competitive landscape,\n"
            "  emerging technologies, and regulatory changes.\n"
            "\n"
            "Slide 3 - Financial Performance (include a TABLE):\n"
            "  Title: 'Financial Performance'\n"
            "  Table with columns: Metric, Q1, Q2, Q3, Q4\n"
            "  Rows: Revenue ($M), Costs ($M), Profit ($M), Growth (%)\n"
            "\n"
            "Slide 4 - Revenue Trends (include a BAR CHART):\n"
            "  Title: 'Quarterly Revenue Trends'\n"
            "  Bar chart showing quarterly revenue for 2024 vs 2025.\n"
            "\n"
            "Slide 5 - Our Strategy (text with bullets):\n"
            "  Title: 'Strategic Priorities'\n"
            "  Bullets: Three-pillar growth approach, market expansion,\n"
            "  product innovation, operational excellence.\n"
            "\n"
            "Slide 6 - Closing Slide:\n"
            "  Title: 'Next Steps'\n"
            "  Bullets: Implementation timeline, key milestones, success metrics.\n"
            "\n"
            "Save as 'generated_content.pptx'"
        )

    print("=" * 60)
    print("PowerPoint Template Workflow")
    print("=" * 60)
    print("Provider: %s" % args.llm_provider)
    print("Template: %s" % (args.template or "none (raw output)"))
    print("Output:   %s" % output_path)
    print("Images:   %s" % ("disabled" if args.no_images else "enabled"))
    print("Stream:   %s" % ("disabled" if args.no_stream else "enabled"))
    if VERBOSE:
        print("Verbose:  enabled")

    # Build workflow steps
    steps: List[Step] = [
        Step(name="Content Generation", executor=step_generate_content),
    ]

    if not args.no_images:
        steps.append(Step(name="Image Planning", executor=step_plan_images))
        steps.append(Step(name="Image Generation", executor=step_generate_images))

    if args.template:
        steps.append(Step(name="Template Assembly", executor=step_assemble_template))

    # Optional Step 5: Visual quality review with Gemini vision
    if args.visual_review:
        steps.append(
            Step(name="Visual Quality Review", executor=step_visual_quality_review)
        )

    # Create and run the workflow
    #
    # session_state schema:
    #   template_path      (str)  - Path to the .pptx template file
    #   output_dir         (str)  - Directory for intermediate output files
    #   output_path        (str)  - Path for the final assembled presentation
    #   generated_file     (str)  - Path to the Claude-generated .pptx file (set by step 1)
    #   slides_data        (list) - Per-slide metadata dicts (set by step 1)
    #   total_slides       (int)  - Number of slides in the generated presentation
    #   generated_images   (dict) - Mapping of slide index -> image bytes (set by step 3)
    #   verbose            (bool) - Whether verbose logging is enabled
    #   stream             (bool) - Whether to use streaming mode for Claude agent
    #   user_prompt        (str)  - Original user prompt (for image planner context)
    #   min_images         (int)  - Minimum slides that must have AI-generated images
    #   src_slide_width    (int)  - Claude-generated slide width in EMU (set by step 1)
    #   src_slide_height   (int)  - Claude-generated slide height in EMU (set by step 1)
    #   footer_text        (str)  - Footer text for idx=11 placeholder (from --footer-text)
    #   date_text          (str)  - Date text for idx=10 placeholder (from --date-text)
    #   show_slide_numbers (bool) - Keep slide number placeholder idx=12
    #   assembly_knowledge (dict) - Comprehensive knowledge file built at start of step 4.
    #                               Contains: input_1_user_intent, input_2_content_plan,
    #                               input_3_template_design_language (per-layout font/color/shape
    #                               analysis), input_4_image_assets, assembly_directives.
    #                               Acts as the single source of truth for all design and
    #                               content decisions during PPTX file generation. (set by step 4)
    #   quality_report     (dict) - PresentationQualityReport.model_dump() (set by step 5)
    workflow = Workflow(
        name="PowerPoint Template Workflow",
        steps=steps,
        session_state={
            "template_path": args.template or "",
            "output_dir": output_dir,
            "output_path": output_path,
            "generated_file": "",
            "slides_data": [],
            "total_slides": 0,
            "generated_images": {},
            "verbose": args.verbose,
            "stream": not args.no_stream,
            "user_prompt": prompt,
            "min_images": args.min_images,
            "src_slide_width": 0,
            "src_slide_height": 0,
            "footer_text": args.footer_text,
            "date_text": args.date_text,
            "show_slide_numbers": args.show_slide_numbers,
            "assembly_knowledge": {},
            "quality_report": {},
            # LLM provider name for swappable agents (agents are loaded lazily per step)
            "llm_provider": args.llm_provider,
        },
    )

    _workflow_start = time.time()
    workflow.print_response(input=prompt, markdown=True)

    # If no template was provided, copy the generated file to the output path
    if not args.template:
        generated_file = workflow.session_state.get("generated_file", "")
        if generated_file and os.path.isfile(generated_file):
            import shutil as _shutil

            _shutil.copy2(generated_file, output_path)
            print(
                "No template specified: raw generated presentation saved to %s"
                % output_path
            )

    _workflow_elapsed = time.time() - _workflow_start
    print("[TIMING] Total workflow execution: %.2fs" % _workflow_elapsed)

    print("\n" + "=" * 60)
    print("Workflow complete!")
    print("Output: %s" % output_path)
    print("=" * 60)


# NOTE: _hex_to_rgb, _relative_luminance, _contrast_ratio are defined once
# near the top of this file (around line 639). Do NOT duplicate them here.


def _extract_color_from_solid_fill(solid_fill_elem) -> str | None:
    """Extract hex color string from an OPC solidFill XML element.

    Handles srgbClr, schemeClr (with common theme defaults), and sysClr.
    Returns hex string like 'FF0000' or None if not extractable.
    """
    if solid_fill_elem is None:
        return None
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    # 1. Check srgbClr first (most common in generated code)
    srgb = solid_fill_elem.find(ns_a + "srgbClr")
    if srgb is not None:
        return srgb.get("val", None)

    # 2. Check schemeClr (theme colors) — map to common defaults
    scheme = solid_fill_elem.find(ns_a + "schemeClr")
    if scheme is not None:
        val = scheme.get("val", "")
        # Standard Office theme defaults (covers most non-custom themes)
        scheme_defaults = {
            "lt1": "FFFFFF",      # Light 1 (usually white)
            "dk1": "000000",      # Dark 1 (usually black)
            "lt2": "E7E6E6",      # Light 2 (usually light gray)
            "dk2": "44546A",      # Dark 2 (usually dark blue-gray)
            "bg1": "FFFFFF",      # Background 1 (usually white)
            "bg2": "E7E6E6",      # Background 2 (usually light gray)
            "tx1": "000000",      # Text 1 (usually black)
            "tx2": "44546A",      # Text 2 (usually dark blue-gray)
            "accent1": "4472C4",  # Blue
            "accent2": "ED7D31",  # Orange
            "accent3": "A5A5A5",  # Gray
            "accent4": "FFC000",  # Gold
            "accent5": "5B9BD5",  # Light blue
            "accent6": "70AD47",  # Green
            "hlink": "0563C1",    # Hyperlink
            "folHlink": "954F72", # Followed hyperlink
        }
        if val in scheme_defaults:
            return scheme_defaults[val]
        return None

    # 3. Check sysClr (system colors with lastClr fallback)
    sys_clr = solid_fill_elem.find(ns_a + "sysClr")
    if sys_clr is not None:
        last_clr = sys_clr.get("lastClr", None)
        if last_clr:
            return last_clr
        # Map common system color names
        val = sys_clr.get("val", "")
        sys_defaults = {
            "windowText": "000000",
            "window": "FFFFFF",
        }
        if val in sys_defaults:
            return sys_defaults[val]

    return None


def _get_shape_background_color(shape, slide) -> str:
    """Detect the effective background color behind a shape.

    Uses a multi-layer detection strategy:
      1. Shape's own fill (solid, gradient first-stop)
      2. Slide background (solid, gradient, image, theme-ref)
      3. Slide layout background
      3.5 Slide master shape traversal (crucial for full-background wrapper shapes 
          often used by templates without explicitly declaring a `<p:bg>` background)
      4. Slide master background
      5. Theme dk1/dk2 heuristic (dark-theme detection)
      6. Default fallback

    For image-based backgrounds (blipFill), assumes a dark color since most
    professional templates with image backgrounds use dark overlays.
    professional templates with image backgrounds use dark overlays.

    Returns hex color string (e.g. '394755').
    """
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    ns_p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

    def _extract_bg_color_from_element(bg_elem):
        """Extract background color from a p:bg or p:bgPr XML element.

        Handles solidFill, gradFill (first stop), blipFill (image → dark heuristic),
        and bgRef (theme reference).
        Returns hex string or None.
        """
        if bg_elem is None:
            return None

        # Check bgPr child first (contains the actual fill)
        bgPr = bg_elem.find(ns_p + "bgPr")
        if bgPr is None:
            bgPr = bg_elem  # Sometimes fill is directly under bg

        # 1. solidFill — most straightforward
        solid_fill = bgPr.find(ns_a + "solidFill")
        if solid_fill is None:
            solid_fill = bgPr.find(".//" + ns_a + "solidFill")
        if solid_fill is not None:
            color = _extract_color_from_solid_fill(solid_fill)
            if color:
                return color

        # 2. gradFill — use the first gradient stop color
        grad_fill = bgPr.find(ns_a + "gradFill")
        if grad_fill is None:
            grad_fill = bgPr.find(".//" + ns_a + "gradFill")
        if grad_fill is not None:
            gs_lst = grad_fill.find(ns_a + "gsLst")
            if gs_lst is not None:
                first_gs = gs_lst.find(ns_a + "gs")
                if first_gs is not None:
                    solid = first_gs.find(ns_a + "srgbClr")
                    if solid is not None:
                        return solid.get("val")
                    scheme = first_gs.find(ns_a + "schemeClr")
                    if scheme is not None:
                        val = scheme.get("val", "")
                        scheme_defaults = {
                            "dk1": "000000", "dk2": "44546A",
                            "lt1": "FFFFFF", "lt2": "E7E6E6",
                            "bg1": "FFFFFF", "bg2": "E7E6E6",
                        }
                        if val in scheme_defaults:
                            return scheme_defaults[val]

        # 3. blipFill — image-based background.
        # We can't analyze the image pixels without heavy deps, so use a
        # heuristic: professional templates with image backgrounds are
        # overwhelmingly dark (dark overlays, photos with dark areas).
        blip_fill = bgPr.find(ns_a + "blipFill")
        if blip_fill is None:
            blip_fill = bgPr.find(".//" + ns_a + "blipFill")
        if blip_fill is not None:
            if VERBOSE:
                print(
                    "  [BG DETECT] Image-based background detected — "
                    "assuming dark background (333333)."
                )
            return "333333"  # Assume dark for image backgrounds

        # 4. bgRef — theme-referenced background
        bgRef = bg_elem.find(ns_p + "bgRef")
        if bgRef is None:
            bgRef = bg_elem.find(".//" + ns_p + "bgRef")
        if bgRef is not None:
            # bgRef has an idx attribute and may contain a color override
            srgb = bgRef.find(ns_a + "srgbClr")
            if srgb is not None:
                return srgb.get("val")
            scheme = bgRef.find(ns_a + "schemeClr")
            if scheme is not None:
                val = scheme.get("val", "")
                scheme_defaults = {
                    "dk1": "000000", "dk2": "44546A",
                    "lt1": "FFFFFF", "lt2": "E7E6E6",
                    "bg1": "FFFFFF", "bg2": "E7E6E6",
                }
                if val in scheme_defaults:
                    return scheme_defaults[val]

        return None

    # === Layer 1: Shape's own fill ===
    if hasattr(shape, "fill"):
        try:
            fill = shape.fill
            if fill.type is not None:
                try:
                    fc = fill.fore_color
                    if fc and fc.type is not None:
                        return str(fc.rgb)
                except Exception:
                    pass
        except Exception:
            pass

    # === Layer 2: Slide background ===
    try:
        bg = slide.background
        if bg and bg.fill and bg.fill.type is not None:
            try:
                fc = bg.fill.fore_color
                if fc and fc.type is not None:
                    return str(fc.rgb)
            except Exception:
                pass
    except Exception:
        pass

    # Layer 2b: Slide XML for complex backgrounds (image, gradient, theme-ref)
    try:
        slide_xml = slide._element
        bg_elem = slide_xml.find(ns_p + "bg")
        if bg_elem is None:
            # Sometimes it's under cSld
            cSld = slide_xml.find(ns_p + "cSld")
            if cSld is not None:
                bg_elem = cSld.find(ns_p + "bg")
        if bg_elem is not None:
            color = _extract_bg_color_from_element(bg_elem)
            if color:
                return color
    except Exception:
        pass

    # === Layer 3: Slide layout background ===
    try:
        layout = slide.slide_layout
        layout_xml = layout._element
        bg_elem = layout_xml.find(ns_p + "bg")
        if bg_elem is None:
            cSld = layout_xml.find(ns_p + "cSld")
            if cSld is not None:
                bg_elem = cSld.find(ns_p + "bg")
        if bg_elem is not None:
            color = _extract_bg_color_from_element(bg_elem)
            if color:
                if VERBOSE:
                    print(
                        "  [BG DETECT] Background color from slide layout: #%s"
                        % color
                    )
                return color
    except Exception:
        pass

    # === Layer 3.5: Check master/layout shapes for large background images/fills ===
    # Many premium templates place a full-slide Picture shape on the master
    # as the visual background instead of using <p:bg>. This MUST be checked
    # BEFORE Layer 4 (master bg XML) because the master's <p:bg> element may
    # contain a misleading bgRef/schemeClr that maps to white in hardcoded
    # defaults even though the actual visual background is dark (from the image).
    _master_shape_bg_checked = False
    try:
        slide_w = 0
        slide_h = 0
        # SlideMaster has no slide_width attribute — get from presentation part
        try:
            prs_part = slide.part.package.presentation_part.presentation
            slide_w = prs_part.slide_width
            slide_h = prs_part.slide_height
        except Exception:
            pass
        if not (slide_w and slide_h):
            # Fallback: standard widescreen 13.333" x 7.5" in EMU
            slide_w = 12192000
            slide_h = 6858000
        if slide_w and slide_h:
            slide_area = slide_w * slide_h
            shape_sources = []
            try:
                shape_sources.append(("layout", slide.slide_layout.shapes))
            except Exception:
                pass
            try:
                shape_sources.append(("master", slide.slide_layout.slide_master.shapes))
            except Exception:
                pass

            for source_name, shapes_collection in shape_sources:
                for s in shapes_collection:
                    try:
                        if not (s.width and s.height):
                            continue
                        shape_area = s.width * s.height
                        if shape_area > slide_area * 0.6:
                            # Large shape found — check for image (common bg pattern)
                            s_elem = s._element
                            blip = s_elem.find(".//" + ns_a + "blip")
                            if blip is not None:
                                if VERBOSE:
                                    print(
                                        "  [BG DETECT] Large image shape on %s "
                                        "'%s' (%.0f%% coverage) — assuming dark (333333)."
                                        % (
                                            source_name,
                                            getattr(s, "name", "?"),
                                            100.0 * shape_area / slide_area,
                                        )
                                    )
                                return "333333"
                            # Check for solid fill
                            if hasattr(s, "fill") and s.fill.type is not None:
                                try:
                                    fc = s.fill.fore_color
                                    if fc and fc.type is not None:
                                        bg_color = str(fc.rgb)
                                        if VERBOSE:
                                            print(
                                                "  [BG DETECT] Large %s shape "
                                                "'%s' (%.0f%% coverage): #%s"
                                                % (
                                                    source_name,
                                                    getattr(s, "name", "?"),
                                                    100.0 * shape_area / slide_area,
                                                    bg_color,
                                                )
                                            )
                                        return bg_color
                                except Exception:
                                    pass
                    except Exception:
                        continue
            _master_shape_bg_checked = True
    except Exception:
        pass

    # === Layer 4: Slide master background ===
    try:
        master = slide.slide_layout.slide_master
        master_xml = master._element
        bg_elem = master_xml.find(ns_p + "bg")
        if bg_elem is None:
            cSld = master_xml.find(ns_p + "cSld")
            if cSld is not None:
                bg_elem = cSld.find(ns_p + "bg")
        if bg_elem is not None:
            color = _extract_bg_color_from_element(bg_elem)
            if color:
                if VERBOSE:
                    print(
                        "  [BG DETECT] Background color from slide master: #%s"
                        % color
                    )
                return color
    except Exception:
        pass

    # === Layer 5: Theme dk1 heuristic ===
    # If we reached here, the background is likely inherited from the theme.
    # Check if the theme uses a dark color scheme (dk1 luminance check).
    # Many premium templates define dark backgrounds via the theme, not explicit fills.
    try:
        master = slide.slide_layout.slide_master
        theme_xml = master.element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}theme"
        )
        # Try extracting dk1 from the theme
        theme_elem = master._element
        # Look for clrScheme in the slide master's theme
        for clr_scheme in theme_elem.iter(ns_a + "clrScheme"):
            dk1 = clr_scheme.find(ns_a + "dk1")
            if dk1 is not None:
                srgb = dk1.find(ns_a + "srgbClr")
                if srgb is not None:
                    dk1_hex = srgb.get("val", "000000")
                    dk1_rgb = _hex_to_rgb(dk1_hex)
                    dk1_lum = _relative_luminance(*dk1_rgb)
                    # If dk1 is actually bright, this is likely a dark-bg template
                    # (dk1 is text-on-dark, so lt1 is the background)
                    lt1 = clr_scheme.find(ns_a + "lt1")
                    if lt1 is not None:
                        lt1_srgb = lt1.find(ns_a + "srgbClr")
                        if lt1_srgb is not None:
                            return lt1_srgb.get("val", "FFFFFF")
                    break
                sys_clr = dk1.find(ns_a + "sysClr")
                if sys_clr is not None:
                    last_clr = sys_clr.get("lastClr")
                    if last_clr:
                        break
    except Exception:
        pass

    # === Layer 6: Check if slide has large dark shapes covering the background ===
    # Some templates use full-slide rectangles as background (not actual bg elements)
    try:
        slide_w = 0
        slide_h = 0
        try:
            prs_part = slide.part.package.presentation_part.presentation
            slide_w = prs_part.slide_width
            slide_h = prs_part.slide_height
        except Exception:
            pass
        if not (slide_w and slide_h):
            slide_w = 12192000
            slide_h = 6858000

        if slide_w and slide_h:
            slide_area = slide_w * slide_h
            for s in slide.shapes:
                try:
                    if s == shape:
                        continue
                    shape_area = s.width * s.height
                    # If a shape covers >60% of the slide, treat it as background
                    if shape_area > slide_area * 0.6:
                        if hasattr(s, "fill") and s.fill.type is not None:
                            try:
                                fc = s.fill.fore_color
                                if fc and fc.type is not None:
                                    bg_color = str(fc.rgb)
                                    if VERBOSE:
                                        print(
                                            "  [BG DETECT] Large background shape detected "
                                            "(%.0f%% coverage): #%s"
                                            % (
                                                100.0 * shape_area / slide_area,
                                                bg_color,
                                            )
                                        )
                                    return bg_color
                            except Exception:
                                pass
                        # Check for image shapes (blip) — assume dark
                        s_elem = s._element
                        blip = s_elem.find(".//" + ns_a + "blip")
                        if blip is not None:
                            if VERBOSE:
                                print(
                                    "  [BG DETECT] Large image shape on slide "
                                    "(%.0f%% coverage) — assuming dark (333333)."
                                    % (100.0 * shape_area / slide_area)
                                )
                            return "333333"
                except Exception:
                    continue
    except Exception:
        pass

    # === Layer 7: Check master/layout shapes for large background images/fills ===
    # Many premium templates place a full-slide Picture shape on the master
    # as the visual background instead of using <p:bg>. This is very common.
    try:
        slide_w = slide_w or 0
        slide_h = slide_h or 0
        if not (slide_w and slide_h):
            try:
                prs_part = slide.part.package.presentation_part.presentation
                slide_w = prs_part.slide_width
                slide_h = prs_part.slide_height
            except Exception:
                pass
        if not (slide_w and slide_h):
            slide_w = 12192000
            slide_h = 6858000
        if slide_w and slide_h:
            slide_area = slide_w * slide_h
            # Check layout shapes first, then master (closer layer takes priority)
            shape_sources = []
            try:
                shape_sources.append(("layout", slide.slide_layout.shapes))
            except Exception:
                pass
            try:
                shape_sources.append(("master", slide.slide_layout.slide_master.shapes))
            except Exception:
                pass

            for source_name, shapes_collection in shape_sources:
                for s in shapes_collection:
                    try:
                        if not (s.width and s.height):
                            continue
                        shape_area = s.width * s.height
                        if shape_area > slide_area * 0.6:
                            # Large shape found — check for image (common bg pattern)
                            s_elem = s._element
                            blip = s_elem.find(".//" + ns_a + "blip")
                            if blip is not None:
                                if VERBOSE:
                                    print(
                                        "  [BG DETECT] Large image shape on %s "
                                        "'%s' (%.0f%% coverage) — assuming dark (333333)."
                                        % (
                                            source_name,
                                            getattr(s, "name", "?"),
                                            100.0 * shape_area / slide_area,
                                        )
                                    )
                                return "333333"
                            # Check for solid fill
                            if hasattr(s, "fill") and s.fill.type is not None:
                                try:
                                    fc = s.fill.fore_color
                                    if fc and fc.type is not None:
                                        bg_color = str(fc.rgb)
                                        if VERBOSE:
                                            print(
                                                "  [BG DETECT] Large %s shape "
                                                "'%s' (%.0f%% coverage): #%s"
                                                % (
                                                    source_name,
                                                    getattr(s, "name", "?"),
                                                    100.0 * shape_area / slide_area,
                                                    bg_color,
                                                )
                                            )
                                        return bg_color
                                except Exception:
                                    pass
                    except Exception:
                        continue
    except Exception:
        pass

    if VERBOSE:
        print(
            "  [BG DETECT] Could not determine background color — "
            "defaulting to FFFFFF (white). This may cause contrast issues "
            "on dark-themed templates."
        )
    return "FFFFFF"  # Default to white



def _make_high_contrast_fill(rPr, bg_hex: str, existing_solidFill=None):
    """Set text color to black or white for maximum contrast against background.
    
    IMPORTANT OOXML SEQUENCE ENFORCEMENT:
    Microsoft PowerPoint strictly enforces OpenXML schema sequencing within 
    the `<a:rPr>` element. The `<a:solidFill>` tag MUST logically precede structural
    font styling tags such as `<a:latin>`, `<a:ea>`, `<a:cs>`, or `<a:sym>`. 
    If appended to the end of the element group, PowerPoint's engine will silently
    declare the tag invalid and discard the high-contrast color modification.
    Therefore, this function removes any misordered pre-existing solidFills and 
    programmatically calculates the correct schema index to ensure compliant insertion.
    """
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    bg_rgb = _hex_to_rgb(bg_hex)
    bg_lum = _relative_luminance(*bg_rgb)
    # Use white text on dark backgrounds, black text on light backgrounds
    text_hex = "FFFFFF" if bg_lum < 0.4 else "000000"

    if existing_solidFill is not None:
        rPr.remove(existing_solidFill)

    # PowerPoint strictly enforces XML sequence in <a:rPr>
    # solidFill MUST appear before latin, ea, cs, sym, and hlink tags.
    solid_fill = etree.Element(ns_a + "solidFill")
    srgb_clr = etree.SubElement(solid_fill, ns_a + "srgbClr")
    srgb_clr.set("val", text_hex)

    insert_idx = 0
    tags_after_fill = {"latin", "ea", "cs", "sym", "hlinkClick", "hlinkMouseOver", "rtl", "extLst"}
    for i, child in enumerate(rPr):
        tag = child.tag.split("}")[-1]
        if tag in tags_after_fill:
            insert_idx = i
            break
        insert_idx = i + 1

    rPr.insert(insert_idx, solid_fill)


def _set_chart_text_color(rPr_elem, ns_a: str, color_hex: str):
    """Set or replace solidFill color on a chart rPr/defRPr element.
    
    Mirrors the strict OOPXML Sequence Enforcement implemented in 
    `_make_high_contrast_fill()`. Ensures `<a:solidFill>` appears before 
    `<a:latin>` elements so PowerPoint successfully parses and renders 
    axis labels, legends, and data labels in high contrast.
    """
    existing = rPr_elem.find(ns_a + "solidFill")
    if existing is not None:
        rPr_elem.remove(existing)

    new_fill = etree.Element(ns_a + "solidFill")
    srgb = etree.SubElement(new_fill, ns_a + "srgbClr")
    srgb.set("val", color_hex)

    insert_idx = 0
    tags_after_fill = {"latin", "ea", "cs", "sym", "hlinkClick", "hlinkMouseOver", "rtl", "extLst"}
    for i, child in enumerate(rPr_elem):
        tag = child.tag.split("}")[-1]
        if tag in tags_after_fill:
            insert_idx = i
            break
        insert_idx = i + 1

    rPr_elem.insert(insert_idx, new_fill)


def _ensure_chart_data_labels(slide) -> int:
    """Fix 6a: Enable data labels on all charts in a slide.

    For every chart shape:
      - Enables data labels on plot 0 (has_data_labels = True)
      - For PIE charts: shows category names + percentage in each label
      - Sets a minimum font size (Pt(10)) on data labels via XML

    Returns number of charts that were updated.
    """
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    ns_c = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
    updated = 0

    for shape in list(slide.shapes):
        try:
            if not getattr(shape, "has_chart", False):
                continue

            chart = shape.chart
            for plot in chart.plots:
                try:
                    if not plot.has_data_labels:
                        plot.has_data_labels = True
                        updated += 1

                    dl = plot.data_labels
                    # Enable value display by default
                    dl.show_value = True

                    # For PIE charts: show category name + percentage
                    try:
                        from pptx.enum.chart import XL_CHART_TYPE

                        chart_type = chart.chart_type
                        _pie_types = {
                            XL_CHART_TYPE.PIE,
                            XL_CHART_TYPE.PIE_EXPLODED,
                            XL_CHART_TYPE.PIE_OF_PIE,
                        }
                        if chart_type in _pie_types:
                            dl.show_category_name = True
                            dl.show_percentage = True
                            dl.show_value = False  # percentage is more useful for pie
                    except Exception:
                        pass

                    # Set minimum font size on data labels via XML
                    try:
                        chart_elem = chart._chartSpace
                        for dLbl_rPr in chart_elem.iter(ns_a + "defRPr"):
                            sz = dLbl_rPr.get("sz")
                            if sz is None or int(sz) < 1000:  # < Pt(10)
                                dLbl_rPr.set("sz", "1000")
                    except Exception:
                        pass

                except Exception:
                    continue

        except Exception:
            continue

    if updated > 0 and VERBOSE:
        print("  [CHART LABELS] Enabled data labels on %d chart(s)." % updated)

    return updated


def _ensure_chart_fills_area(
    slide, slide_width: int, slide_height: int, visual_area: "ContentArea | None" = None
) -> bool:
    """Fix 6b: Enlarge undersized charts to fill available visual area.

    If a chart occupies less than 40% of the visual area, resizes it to
    fill ~60% while respecting safe margins. Repositions centrally in the
    visual region.

    Args:
        slide:        The pptx slide object.
        slide_width:  Presentation width in EMU.
        slide_height: Presentation height in EMU.
        visual_area:  Optional ContentArea for the visual region. If None,
                      uses safe-margin defaults (8% top/bottom, 5% left/right).

    Returns True if any chart was resized.
    """
    if visual_area is None:
        # Default safe-margin content area
        margin_x = int(slide_width * 0.05)
        margin_y = int(slide_height * 0.08)
        visual_area = ContentArea(
            left=margin_x,
            top=margin_y,
            width=slide_width - 2 * margin_x,
            height=slide_height - 2 * margin_y,
        )

    total_visual_pixels = visual_area.width * visual_area.height
    if total_visual_pixels <= 0:
        return False

    modified = False

    for shape in list(slide.shapes):
        try:
            if not getattr(shape, "has_chart", False):
                continue

            chart_pixels = shape.width * shape.height
            fill_ratio = chart_pixels / total_visual_pixels

            if fill_ratio < 0.40:
                # Enlarge to ~60% of visual area, maintaining aspect ratio
                target_w = int(visual_area.width * 0.85)
                target_h = int(visual_area.height * 0.65)

                # Center in visual area
                new_left = visual_area.left + (visual_area.width - target_w) // 2
                new_top = visual_area.top + (visual_area.height - target_h) // 2

                if VERBOSE:
                    print(
                        "  [CHART SIZE] Enlarging chart from %.0f%% to ~60%% fill "
                        "(was %dx%d, now %dx%d EMU)"
                        % (fill_ratio * 100, shape.width, shape.height, target_w, target_h)
                    )

                shape.left = new_left
                shape.top = new_top
                shape.width = target_w
                shape.height = target_h
                modified = True

        except Exception:
            continue

    return modified


def _remove_duplicate_titles(slide) -> int:
    """Fix 7: Remove free-floating textboxes that duplicate the title placeholder.

    When Tier 2 LLM creates a title textbox AND the template assembly also
    fills the title placeholder (idx=0), both end up on the slide. This
    function detects and removes the redundant free-floating copy.

    Algorithm:
      1. Get title text from placeholder idx=0.
      2. Scan non-placeholder text shapes for matching text (case-insensitive,
         whitespace-normalized).
      3. Remove duplicates from the shape tree.

    Returns number of duplicate shapes removed.
    """
    title_text = ""
    title_ph = None

    for shape in list(slide.placeholders):
        try:
            if shape.placeholder_format.idx == 0 and getattr(shape, "has_text_frame", False):
                title_text = shape.text.strip()
                title_ph = shape
                break
        except Exception:
            continue

    if not title_text:
        return 0

    title_norm = " ".join(title_text.lower().split())
    if len(title_norm) < 3:
        return 0

    spTree = slide.shapes._spTree
    to_remove = []

    for shape in list(slide.shapes):
        try:
            if shape.is_placeholder:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            shape_text = shape.text.strip()
            shape_norm = " ".join(shape_text.lower().split())

            if shape_norm == title_norm:
                to_remove.append(shape._element)
        except Exception:
            continue

    for elem in to_remove:
        try:
            spTree.remove(elem)
        except Exception:
            pass

    if to_remove:
        print(
            "  [DEDUP] Removed %d duplicate title textbox(es) matching '%s'."
            % (len(to_remove), title_text[:50])
        )

    return len(to_remove)


def _fix_chart_text_contrast(shape, slide, min_ratio: float = 3.0) -> int:
    """Fix text contrast inside chart elements (axis labels, data labels, legend, title).

    Chart text (axis tick labels, data labels, legend entries, chart title)
    lives in c:txPr/a:defRPr or c:tx/a:rPr XML — NOT accessible via
    shape.text_frame. This function scans all such elements and forces
    high-contrast text when the slide background is dark.

    Returns number of corrected text elements.
    """
    if not getattr(shape, "has_chart", False):
        return 0

    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    bg_hex = _get_shape_background_color(shape, slide)
    bg_rgb = _hex_to_rgb(bg_hex)
    bg_lum = _relative_luminance(*bg_rgb)

    if bg_lum >= 0.4:
        return 0  # Light background — chart text is likely fine

    # On dark backgrounds, force all chart text to white
    target_color = "FFFFFF"
    corrections = 0

    try:
        chart_elem = shape.chart._chartSpace
    except Exception:
        return 0

    # Find all defRPr and rPr elements inside the chart XML
    for rPr in chart_elem.iter(ns_a + "defRPr"):
        _set_chart_text_color(rPr, ns_a, target_color)
        corrections += 1
    for rPr in chart_elem.iter(ns_a + "rPr"):
        _set_chart_text_color(rPr, ns_a, target_color)
        corrections += 1

    return corrections

def enforce_final_contrast(pptx_path: str, min_ratio: float = 3.0) -> int:
    """Post-merge WCAG contrast enforcement — template-independent safety net.

    Opens the PPTX, checks every text run against its effective background,
    and fixes any low-contrast text. Works without a template.
    
    Additional QA Enhancements:
    1. Iterates over Chart elements to repair internal sub-element text visibility (defRPr).
    2. Overrides AI compression side-effects by enforcing minimum shape bounds (3.0" x 4.0")
       specifically on undersized Charts.
    3. Avoids post-resize overlap by caching bounding box properties and pushing adjacent
       rescaled geometries hierarchically downwards.

    Args:
        pptx_path: Path to the PPTX file to process.
        min_ratio: Minimum WCAG contrast ratio to accept (default 3.0).

    Returns:
        Number of text runs that were corrected.
    """
    from pptx import Presentation as PptxPresentation

    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    prs = PptxPresentation(pptx_path)
    corrections = 0

    for slide in prs.slides:
        resized_charts = []
        for shape in slide.shapes:
            # Text frames
            if getattr(shape, "has_text_frame", False):
                bg_hex = _get_shape_background_color(shape, slide)
                bg_rgb = _hex_to_rgb(bg_hex)
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        rPr = run._r.find(ns_a + "rPr")
                        if rPr is not None:
                            solidFill = rPr.find(ns_a + "solidFill")
                            if solidFill is not None:
                                text_hex = _extract_color_from_solid_fill(solidFill)
                                if text_hex:
                                    text_rgb = _hex_to_rgb(text_hex)
                                    ratio = _contrast_ratio(text_rgb, bg_rgb)
                                    if ratio < min_ratio:
                                        _make_high_contrast_fill(rPr, bg_hex, existing_solidFill=solidFill)
                                        corrections += 1
                                else:
                                    # Can't determine text color — check default (black) contrast
                                    default_rgb = (0, 0, 0)
                                    ratio = _contrast_ratio(default_rgb, bg_rgb)
                                    if ratio >= 3.0:
                                        # Safe to strip — default black has good contrast
                                        rPr.remove(solidFill)
                                        corrections += 1
                                    else:
                                        # Stripping would cause dark-on-dark — replace
                                        _make_high_contrast_fill(rPr, bg_hex, existing_solidFill=solidFill)
                                        corrections += 1
                            else:
                                # No explicit color — check if inherited color (assume black) has contrast
                                default_rgb = (0, 0, 0)
                                ratio = _contrast_ratio(default_rgb, bg_rgb)
                                if ratio < min_ratio:
                                    # Dark default on dark bg — fix it
                                    if rPr is None:
                                        from lxml import etree as _etree

                                        rPr = _etree.SubElement(run._r, ns_a + "rPr")
                                    _make_high_contrast_fill(rPr, bg_hex)
                                    corrections += 1
                        else:
                            # No rPr at all — check default black against background
                            default_rgb = (0, 0, 0)
                            ratio = _contrast_ratio(default_rgb, bg_rgb)
                            if ratio < min_ratio:
                                from lxml import etree as _etree

                                rPr = _etree.SubElement(run._r, ns_a + "rPr")
                                _make_high_contrast_fill(rPr, bg_hex)
                                corrections += 1

            # Tables
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_bg_hex = "FFFFFF"
                        try:
                            tc_elem = cell._tc
                            tc_pr = tc_elem.find(ns_a + "tcPr") if tc_elem is not None else None
                            if tc_pr is not None:
                                cell_fill = tc_pr.find(ns_a + "solidFill")
                                if cell_fill is not None:
                                    extracted = _extract_color_from_solid_fill(cell_fill)
                                    if extracted:
                                        cell_bg_hex = extracted
                        except Exception:
                            pass

                        cell_bg_rgb = _hex_to_rgb(cell_bg_hex)
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if not run.text.strip():
                                    continue
                                rPr = run._r.find(ns_a + "rPr")
                                if rPr is not None:
                                    solidFill = rPr.find(ns_a + "solidFill")
                                    if solidFill is not None:
                                        text_hex = _extract_color_from_solid_fill(solidFill)
                                        if text_hex:
                                            text_rgb = _hex_to_rgb(text_hex)
                                            ratio = _contrast_ratio(text_rgb, cell_bg_rgb)
                                            if ratio < min_ratio:
                                                _make_high_contrast_fill(rPr, cell_bg_hex, existing_solidFill=solidFill)
                                                corrections += 1
                                        else:
                                            # Can't determine text color — check default (black) contrast
                                            default_rgb = (0, 0, 0)
                                            ratio = _contrast_ratio(default_rgb, cell_bg_rgb)
                                            if ratio >= 3.0:
                                                rPr.remove(solidFill)
                                                corrections += 1
                                            else:
                                                _make_high_contrast_fill(rPr, cell_bg_hex, existing_solidFill=solidFill)
                                                corrections += 1

            # Charts — axis labels, data labels, legend text
            try:
                chart_fixes = _fix_chart_text_contrast(shape, slide, min_ratio)
                corrections += chart_fixes
            except Exception:
                pass

            # Charts — enforce minimum dimensions to prevent squeezing
            try:
                if getattr(shape, "has_chart", False):
                    from pptx.util import Inches as _Inches
                    _MIN_W = _Inches(4.0)
                    _MIN_H = _Inches(3.0)
                    if shape.width < _MIN_W:
                        shape.width = _MIN_W
                        corrections += 1
                    if shape.height < _MIN_H:
                        shape.height = _MIN_H
                        corrections += 1
                    
                    # Prevent overlap
                    for pc in resized_charts:
                        if abs(shape.left - pc.left) < _Inches(2.0) and abs(shape.top - pc.top) < _Inches(2.0):
                            shape.top = pc.top + pc.height + _Inches(0.2)
                            corrections += 1
                    resized_charts.append(shape)
            except Exception:
                pass

    if corrections > 0:
        prs.save(pptx_path)
        print("    [CONTRAST] Fixed %d low-contrast text run(s) in final output" % corrections)

    return corrections


def clean_presentation_visual_noise_and_contrast(prs) -> None:
    """Clean empty placeholder ghost text and strip hardcoded font colors.

    1. Removes any text shapes containing default MS PowerPoint placeholder ghost text or empty text.
    2. Strips <a:solidFill> from text runs so text inherits the high-contrast
       theme color from the slide master. 
       Handles OOXML `a:rPr` property omissions gracefully by creating properties dynamically
       to enforce contrast baselines if natively absent.
    3. Extends contrast remediation logic deep into OOXML Chart schemas (axis/titles/labels).
    """
    ghost_texts = {
        "click to add title",
        "click to add subtitle",
        "click to add text",
        "click to add notes",
        "double-tap to add title",
        "double-tap to add subtitle",
        "double-tap to add text",
    }
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    for slide in prs.slides:
        spTree = slide.shapes._spTree
        elements_to_remove = []
        
        for shape in list(slide.shapes):
            # 1. Clean visual noise
            try:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip().lower()
                    if text in ghost_texts or text == "":
                        elements_to_remove.append(shape._element)
                        continue # Removed, don't check for color
            except Exception:
                pass
            
            # 2. Fix Text Contrast (contrast-aware color correction)
            try:
                if getattr(shape, "has_text_frame", False):
                    bg_hex = _get_shape_background_color(shape, slide)
                    bg_rgb = _hex_to_rgb(bg_hex)
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            rPr = run._r.find(ns_a + "rPr")
                            if rPr is not None:
                                solidFill = rPr.find(ns_a + "solidFill")
                                if solidFill is not None:
                                    text_hex = _extract_color_from_solid_fill(solidFill)
                                    if text_hex:
                                        text_rgb = _hex_to_rgb(text_hex)
                                        ratio = _contrast_ratio(text_rgb, bg_rgb)
                                        if ratio < 3.0:
                                            # Low contrast — replace with high-contrast color
                                            _make_high_contrast_fill(rPr, bg_hex, existing_solidFill=solidFill)
                                        # else: contrast is adequate, keep original color
                                    else:
                                        # Can't determine text color — check if stripping is safe
                                        # (inherited default is typically black)
                                        default_rgb = (0, 0, 0)
                                        ratio = _contrast_ratio(default_rgb, bg_rgb)
                                        if ratio >= 3.0:
                                            # Safe to strip — default black has good contrast
                                            rPr.remove(solidFill)
                                        else:
                                            # Stripping would cause dark-on-dark — replace instead
                                            _make_high_contrast_fill(rPr, bg_hex, existing_solidFill=solidFill)
                            else:
                                # No rPr at all — run inherits default (typically black)
                                default_rgb = (0, 0, 0)
                                ratio = _contrast_ratio(default_rgb, bg_rgb)
                                if ratio < 3.0:
                                    from lxml import etree as _etree
                                    new_rPr = _etree.Element(ns_a + "rPr")
                                    run._r.insert(0, new_rPr)
                                    _make_high_contrast_fill(new_rPr, bg_hex)
            except Exception:
                pass
            
            # Table cells — contrast-aware color correction
            try:
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            # Detect cell background
                            cell_bg_hex = "FFFFFF"
                            try:
                                tc_elem = cell._tc
                                tc_pr = None
                                if tc_elem is not None:
                                    ns_a_table = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                                    tc_pr = tc_elem.find(ns_a_table + "tcPr")
                                if tc_pr is not None:
                                    cell_fill = tc_pr.find(ns_a + "solidFill")
                                    if cell_fill is not None:
                                        extracted = _extract_color_from_solid_fill(cell_fill)
                                        if extracted:
                                            cell_bg_hex = extracted
                            except Exception:
                                pass

                            cell_bg_rgb = _hex_to_rgb(cell_bg_hex)
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    rPr = run._r.find(ns_a + "rPr")
                                    if rPr is not None:
                                        solidFill = rPr.find(ns_a + "solidFill")
                                        if solidFill is not None:
                                            text_hex = _extract_color_from_solid_fill(solidFill)
                                            if text_hex:
                                                text_rgb = _hex_to_rgb(text_hex)
                                                ratio = _contrast_ratio(text_rgb, cell_bg_rgb)
                                                if ratio < 3.0:
                                                    _make_high_contrast_fill(rPr, cell_bg_hex, existing_solidFill=solidFill)
                                                # else: keep original
                                            else:
                                                # Can't determine text color — check if stripping is safe
                                                default_rgb = (0, 0, 0)
                                                ratio = _contrast_ratio(default_rgb, cell_bg_rgb)
                                                if ratio >= 3.0:
                                                    rPr.remove(solidFill)
                                                else:
                                                    _make_high_contrast_fill(rPr, cell_bg_hex, existing_solidFill=solidFill)
            except Exception:
                pass

            # Charts — axis labels, data labels, legend text
            try:
                _fix_chart_text_contrast(shape, slide)
            except Exception:
                pass
        
        for element in elements_to_remove:
            try:
                spTree.remove(element)
            except Exception:
                pass

        # --- Fix 6a: Ensure chart data labels are enabled ---
        try:
            _ensure_chart_data_labels(slide)
        except Exception:
            pass

        # --- Fix 6b: Enlarge undersized charts ---
        try:
            slide_w = int(prs.slide_width)
            slide_h = int(prs.slide_height)
            _ensure_chart_fills_area(slide, slide_w, slide_h)
        except Exception:
            pass

        # --- Fix 7: Remove duplicate title textboxes ---
        try:
            _remove_duplicate_titles(slide)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Layout Sanitization (Concern 2 — universal, all tiers)
# ---------------------------------------------------------------------------
# These functions run AFTER visual cleanup to fix spatial layout problems
# (overlapping shapes, out-of-bounds elements, undersized shapes) that
# the contrast/cleanup pass does not address.


def sanitize_slide_layout(slide, slide_width: int, slide_height: int) -> int:
    """Sanitize a single slide's shape layout to fix spatial defects.

    Performs three passes:
        1. **Boundary clamping**: Ensures all shapes are within 5% safe margin.
        2. **Minimum size enforcement**: Shapes must be at least 1 inch wide
           and 0.5 inch tall to remain visible and interactive.
        3. **Overlap detection & reflow**: If two non-placeholder shapes
           overlap by more than 30% of the smaller shape's area, the lower
           shape is pushed down below the upper shape with a small gap.

    Args:
        slide:        The pptx slide object.
        slide_width:  Presentation width in EMU (e.g. 12192000 for 13.33 in).
        slide_height: Presentation height in EMU (e.g. 6858000 for 7.5 in).

    Returns:
        Number of shape adjustments made.
    """
    from pptx.util import Inches

    adjustments = 0

    # Safe margins (5% each side)
    margin_x = int(slide_width * 0.05)
    margin_y = int(slide_height * 0.05)
    max_right = slide_width - margin_x
    max_bottom = slide_height - margin_y

    min_width = Inches(1.0)
    min_height = Inches(0.5)

    shapes = list(slide.shapes)

    # --- Pass 1: Boundary clamping ---
    for shape in shapes:
        try:
            if not hasattr(shape, "left") or shape.left is None:
                continue

            # Clamp left edge
            if shape.left < margin_x:
                shape.left = margin_x
                adjustments += 1
            # Clamp top edge
            if shape.top < margin_y:
                shape.top = margin_y
                adjustments += 1
            # Clamp right edge (shrink width if needed)
            if shape.left + shape.width > max_right:
                overflow = (shape.left + shape.width) - max_right
                if overflow > shape.width * 0.5:
                    # Too much overflow — move shape left instead
                    shape.left = max(margin_x, max_right - shape.width)
                else:
                    shape.width = max(min_width, shape.width - overflow)
                adjustments += 1
            # Clamp bottom edge
            if shape.top + shape.height > max_bottom:
                overflow = (shape.top + shape.height) - max_bottom
                if overflow > shape.height * 0.5:
                    shape.top = max(margin_y, max_bottom - shape.height)
                else:
                    shape.height = max(min_height, shape.height - overflow)
                adjustments += 1
        except Exception:
            continue

    # --- Pass 2: Minimum size enforcement ---
    for shape in shapes:
        try:
            if not hasattr(shape, "width") or shape.width is None:
                continue
            if shape.width < min_width:
                shape.width = min_width
                adjustments += 1
            if shape.height < min_height:
                shape.height = min_height
                adjustments += 1
        except Exception:
            continue

    # --- Pass 3: Overlap detection & vertical reflow ---
    # Sort shapes top-to-bottom by their top edge
    movable = []
    for s in shapes:
        try:
            if not hasattr(s, "left") or s.left is None:
                continue
            # Skip placeholder shapes (they are positioned by the layout)
            if getattr(s, "is_placeholder", False):
                continue
            movable.append(s)
        except Exception:
            continue

    movable.sort(key=lambda s: (s.top, s.left))

    for i in range(len(movable)):
        for j in range(i + 1, len(movable)):
            try:
                a = movable[i]
                b = movable[j]

                # Calculate overlap rectangle
                overlap_left = max(a.left, b.left)
                overlap_top = max(a.top, b.top)
                overlap_right = min(a.left + a.width, b.left + b.width)
                overlap_bottom = min(a.top + a.height, b.top + b.height)

                if overlap_left < overlap_right and overlap_top < overlap_bottom:
                    overlap_area = (overlap_right - overlap_left) * (overlap_bottom - overlap_top)
                    smaller_area = min(a.width * a.height, b.width * b.height)

                    if smaller_area > 0 and overlap_area > smaller_area * 0.30:
                        # Push shape b below shape a with a small gap
                        gap = Inches(0.15)
                        new_top = a.top + a.height + gap
                        if new_top + b.height <= max_bottom:
                            b.top = new_top
                            adjustments += 1
            except Exception:
                continue

    return adjustments


def sanitize_presentation(prs) -> int:
    """Run layout sanitization on every slide in a presentation.

    This is a universal post-generation safety net that fixes spatial layout
    problems regardless of which tier generated the content.

    Args:
        prs: A python-pptx Presentation object.

    Returns:
        Total number of shape adjustments made across all slides.
    """
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    total = 0

    for slide_idx, slide in enumerate(prs.slides):
        try:
            n = sanitize_slide_layout(slide, slide_w, slide_h)
            if n > 0 and VERBOSE:
                print(
                    "  [LAYOUT] Slide %d: %d spatial adjustment(s) applied."
                    % (slide_idx + 1, n)
                )
            total += n
        except Exception:
            continue

    if total > 0:
        print(
            "[LAYOUT SANITIZE] Applied %d spatial fix(es) across %d slide(s)."
            % (total, len(prs.slides))
        )

    return total


