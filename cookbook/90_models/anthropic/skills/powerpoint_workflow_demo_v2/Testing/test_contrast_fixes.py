#!/usr/bin/env python3
"""Test script to verify visual quality fixes against the existing PPTX output.

This script:
1. Copies the existing output PPTX to a new file
2. Runs enforce_final_contrast() (which now includes chart text contrast and chart sizing)
3. Runs clean_presentation_visual_noise_and_contrast()  
4. Reports what was fixed
5. Saves the result for manual inspection

Usage:
    python test_contrast_fixes.py
"""

import os
import shutil
import sys

# Setup path
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, SCRIPT_DIR)

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Emu

# Import the fixed functions
from powerpoint_template_workflow import (
    enforce_final_contrast,
    clean_presentation_visual_noise_and_contrast,
    _get_shape_background_color,
    _hex_to_rgb,
    _relative_luminance,
    _contrast_ratio,
    _extract_color_from_solid_fill,
    _fix_chart_text_contrast,
)


def diagnose_pptx(pptx_path: str) -> dict:
    """Analyze a PPTX file and report all contrast issues found."""
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    prs = PptxPresentation(pptx_path)
    
    report = {
        "slides": [],
        "total_low_contrast": 0,
        "total_no_rpr": 0,
        "total_chart_text": 0,
        "total_chart_small": 0,
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_report = {
            "index": slide_idx + 1,
            "shapes": [],
            "issues": [],
        }
        
        for shape in slide.shapes:
            shape_info = {
                "name": getattr(shape, "name", "unknown"),
                "type": type(shape).__name__,
                "has_chart": getattr(shape, "has_chart", False),
                "has_table": getattr(shape, "has_table", False),
                "has_text_frame": getattr(shape, "has_text_frame", False),
                "width_inches": shape.width / Inches(1) if shape.width else 0,
                "height_inches": shape.height / Inches(1) if shape.height else 0,
            }
            
            bg_hex = _get_shape_background_color(shape, slide)
            bg_rgb = _hex_to_rgb(bg_hex)
            bg_lum = _relative_luminance(*bg_rgb)
            shape_info["background"] = f"#{bg_hex} (lum={bg_lum:.2f})"
            
            # Check text frame contrast
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        rPr = run._r.find(ns_a + "rPr")
                        text_hex = None
                        color_source = "none"
                        
                        if rPr is not None:
                            solidFill = rPr.find(ns_a + "solidFill")
                            if solidFill is not None:
                                text_hex = _extract_color_from_solid_fill(solidFill)
                                color_source = "explicit"
                            else:
                                color_source = "rPr-no-fill"
                        else:
                            color_source = "no-rPr"
                            report["total_no_rpr"] += 1
                        
                        if text_hex:
                            text_rgb = _hex_to_rgb(text_hex)
                            ratio = _contrast_ratio(text_rgb, bg_rgb)
                            if ratio < 3.0:
                                issue = (
                                    f"  LOW CONTRAST: text=#{text_hex} bg=#{bg_hex} "
                                    f"ratio={ratio:.1f} src={color_source} "
                                    f"text='{run.text[:40]}'"
                                )
                                slide_report["issues"].append(issue)
                                report["total_low_contrast"] += 1
                        elif color_source in ("no-rPr", "rPr-no-fill"):
                            # Default black on dark bg?
                            ratio = _contrast_ratio((0, 0, 0), bg_rgb)
                            if ratio < 3.0:
                                issue = (
                                    f"  INHERITED BLACK on dark bg: bg=#{bg_hex} "
                                    f"ratio={ratio:.1f} src={color_source} "
                                    f"text='{run.text[:40]}'"
                                )
                                slide_report["issues"].append(issue)
                                report["total_low_contrast"] += 1
            
            # Check chart text
            if getattr(shape, "has_chart", False):
                try:
                    chart_elem = shape.chart._chartSpace
                    chart_text_count = 0
                    for rPr in chart_elem.iter(ns_a + "defRPr"):
                        chart_text_count += 1
                    for rPr in chart_elem.iter(ns_a + "rPr"):
                        chart_text_count += 1
                    shape_info["chart_text_elements"] = chart_text_count
                    if bg_lum < 0.4:
                        report["total_chart_text"] += chart_text_count
                        slide_report["issues"].append(
                            f"  CHART TEXT on dark bg: {chart_text_count} text elements need contrast fix"
                        )
                except Exception:
                    pass
                
                # Check chart size
                w_inches = shape.width / Inches(1) if shape.width else 0
                h_inches = shape.height / Inches(1) if shape.height else 0
                if w_inches < 4.0 or h_inches < 3.0:
                    report["total_chart_small"] += 1
                    slide_report["issues"].append(
                        f"  SQUEEZED CHART: {w_inches:.1f}x{h_inches:.1f} inches "
                        f"(min 4.0x3.0)"
                    )
            
            slide_report["shapes"].append(shape_info)
        
        report["slides"].append(slide_report)
    
    return report


def print_report(report: dict, label: str):
    """Pretty-print a diagnostic report."""
    print(f"\n{'='*60}")
    print(f"  {label}")
    print(f"{'='*60}")
    print(f"  Low contrast text runs: {report['total_low_contrast']}")
    print(f"  Runs with no rPr:       {report['total_no_rpr']}")
    print(f"  Chart text on dark bg:  {report['total_chart_text']}")
    print(f"  Squeezed charts:        {report['total_chart_small']}")
    
    for slide in report["slides"]:
        if slide["issues"]:
            print(f"\n  Slide {slide['index']}:")
            for issue in slide["issues"]:
                print(f"    {issue}")
    
    if not any(s["issues"] for s in report["slides"]):
        print("\n  ✅ No issues found!")
    print()


def main():
    # Source file
    source_pptx = os.path.join(
        SCRIPT_DIR,
        "output_chunked/chunked_workflow_work/session_e3fa4b1d_20260305_051723/"
        "software_deck_template.pptx"
    )
    
    if not os.path.exists(source_pptx):
        print(f"ERROR: Source PPTX not found: {source_pptx}")
        sys.exit(1)
    
    # Copy to test file
    output_pptx = os.path.join(SCRIPT_DIR, "software_deck_template_FIXED.pptx")
    shutil.copy2(source_pptx, output_pptx)
    print(f"Copied source to: {output_pptx}")
    
    # Diagnose BEFORE fixes
    print("\n\n" + "#" * 60)
    print("  STEP 1: Diagnosing ORIGINAL file (before fixes)")
    print("#" * 60)
    before_report = diagnose_pptx(output_pptx)
    print_report(before_report, "BEFORE FIXES")
    
    # Apply fix: enforce_final_contrast (includes chart text + chart sizing)
    print("\n" + "#" * 60)
    print("  STEP 2: Running enforce_final_contrast()")
    print("#" * 60)
    corrections = enforce_final_contrast(output_pptx, min_ratio=3.0)
    print(f"  enforce_final_contrast returned: {corrections} corrections")
    
    # Apply fix: clean_presentation_visual_noise_and_contrast
    print("\n" + "#" * 60)
    print("  STEP 3: Running clean_presentation_visual_noise_and_contrast()")
    print("#" * 60)
    prs = PptxPresentation(output_pptx)
    clean_presentation_visual_noise_and_contrast(prs)
    prs.save(output_pptx)
    print("  clean_presentation_visual_noise_and_contrast completed.")
    
    # Diagnose AFTER fixes
    print("\n\n" + "#" * 60)
    print("  STEP 4: Diagnosing FIXED file (after fixes)")
    print("#" * 60)
    after_report = diagnose_pptx(output_pptx)
    print_report(after_report, "AFTER FIXES")
    
    # Summary
    fixed_contrast = before_report["total_low_contrast"] - after_report["total_low_contrast"]
    fixed_chart = before_report["total_chart_text"] - after_report["total_chart_text"]
    fixed_sizing = before_report["total_chart_small"] - after_report["total_chart_small"]
    
    print("=" * 60)
    print("  SUMMARY")
    print("=" * 60)
    print(f"  Contrast issues fixed:  {fixed_contrast}")
    print(f"  Chart text fixed:       {fixed_chart}")
    print(f"  Charts resized:         {fixed_sizing}")
    print(f"\n  Output file: {output_pptx}")
    print("  Please open this file to verify the visual quality.")
    

if __name__ == "__main__":
    main()
