"""
Offline unit tests for Brand/Style parsing logic.

Tests the BrandStyleIntent model, extract_style_from_template(),
_format_brand_context_for_prompt(), and _build_brand_override_log().

No API calls are made — all tests are fully offline.

This test file re-imports only the brand-specific code to avoid requiring
the full agno/anthropic dependency chain.

Usage:
    python test_brand_style_parsing.py
"""

import os
import sys
import tempfile

from pptx import Presentation
from pptx.util import Pt
from pydantic import BaseModel, Field
from typing import List

# --- Copy of BrandStyleIntent model (to avoid importing the full module chain) ---

class BrandStyleIntent(BaseModel):
    """Parsed branding/styling intent extracted from a user query or template file."""
    has_branding: bool = Field(False)
    brand_name: str = Field("")
    style_keywords: List[str] = Field(default_factory=list)
    color_palette: List[str] = Field(default_factory=list)
    tone_override: str = Field("")
    typography_hints: List[str] = Field(default_factory=list)
    content_query: str = Field("")
    source: str = Field("query")
    source_detail: str = Field("")


# --- Copy of helper functions (to avoid importing the full module chain) ---

VERBOSE = False  # noqa: F841


def _format_brand_context_for_prompt(brand_intent):
    if not brand_intent or not brand_intent.has_branding:
        return ""
    sections = ["## Brand/Style Guidance\n"]
    if brand_intent.brand_name:
        sections.append("**Brand:** %s" % brand_intent.brand_name)
    if brand_intent.style_keywords:
        sections.append("**Style:** %s" % ", ".join(brand_intent.style_keywords[:5]))
    if brand_intent.color_palette:
        sections.append("**Color Palette:** %s" % ", ".join(brand_intent.color_palette[:6]))
    if brand_intent.tone_override:
        sections.append("**Tone:** %s" % brand_intent.tone_override)
    if brand_intent.typography_hints:
        sections.append("**Typography:** %s" % ", ".join(brand_intent.typography_hints[:3]))
    sections.append(
        "\nUse these brand guidelines to inform visual direction, tone, terminology, "
        "and content framing throughout the presentation. Reflect the brand's identity "
        "in slide language, suggested color references, and overall aesthetic.\n"
    )
    return "\n".join(sections)


def _build_brand_override_log(query_intent, template_intent):
    template_name = template_intent.source_detail or "provided template"
    lines = [
        "[BRAND OVERRIDE] User specified '%s branding' in query, but a template file "
        "was provided (%s)." % (query_intent.brand_name, template_name),
        "[BRAND OVERRIDE] Styling will be derived from the template file. "
        "Query-level branding intent has been disregarded.",
        "[BRAND OVERRIDE] Reason: Explicit template file takes precedence over "
        "natural language branding directives per workflow specification.",
    ]
    if template_intent.color_palette:
        lines.append(
            "[BRAND OVERRIDE] Template colors: %s" % ", ".join(template_intent.color_palette[:6])
        )
    if template_intent.typography_hints:
        lines.append(
            "[BRAND OVERRIDE] Template fonts: %s" % ", ".join(template_intent.typography_hints)
        )
    return "\n".join(lines)


def extract_style_from_template(template_path):
    intent = BrandStyleIntent(
        source="template",
        source_detail=os.path.basename(template_path),
        has_branding=True,
    )
    try:
        prs = Presentation(template_path)
        colors = []
        try:
            theme_part = prs.slide_masters[0].part.slide_master.element
            ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            for clr_scheme in theme_part.iter(ns_a + "clrScheme"):
                for child in clr_scheme:
                    tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                    if tag in ("dk1", "dk2", "lt1", "lt2", "accent1", "accent2",
                               "accent3", "accent4", "accent5", "accent6"):
                        for color_el in child:
                            val = color_el.get("val", "")
                            last_clr = color_el.get("lastClr", "")
                            hex_val = val if len(val) == 6 else last_clr
                            if hex_val and len(hex_val) == 6:
                                colors.append("#%s" % hex_val.upper())
        except Exception:
            pass
        if colors:
            intent.color_palette = list(dict.fromkeys(colors))[:8]

        fonts = []
        try:
            theme_part = prs.slide_masters[0].part.slide_master.element
            ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            for font_scheme in theme_part.iter(ns_a + "fontScheme"):
                for font_tag in ("majorFont", "minorFont"):
                    font_el = font_scheme.find(ns_a + font_tag)
                    if font_el is not None:
                        latin = font_el.find(ns_a + "latin")
                        if latin is not None:
                            typeface = latin.get("typeface", "")
                            if typeface and typeface not in fonts:
                                fonts.append(typeface)
        except Exception:
            pass
        if fonts:
            intent.typography_hints = fonts[:3]

        try:
            if prs.slides:
                first_slide = prs.slides[0]
                for shape in first_slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if 0 < len(text) < 60 and len(text.split()) <= 4:
                            lower = text.lower()
                            skip = {
                                "click to add title", "click to add subtitle",
                                "click to add text", "title", "subtitle",
                            }
                            if lower not in skip:
                                intent.brand_name = text
                                break
        except Exception:
            pass

    except Exception as e:
        print("[WARNING] Template style extraction failed: %s" % str(e))
        return BrandStyleIntent(source="template", source_detail=os.path.basename(template_path))

    return intent


# === TESTS ===


def test_brand_style_intent_defaults():
    """BrandStyleIntent should have sensible defaults when no branding is present."""
    intent = BrandStyleIntent()
    assert intent.has_branding is False, "Default has_branding should be False"
    assert intent.brand_name == "", "Default brand_name should be empty"
    assert intent.style_keywords == [], "Default style_keywords should be empty list"
    assert intent.color_palette == [], "Default color_palette should be empty list"
    assert intent.tone_override == "", "Default tone_override should be empty"
    assert intent.typography_hints == [], "Default typography_hints should be empty list"
    assert intent.content_query == "", "Default content_query should be empty"
    assert intent.source == "query", "Default source should be 'query'"
    assert intent.source_detail == "", "Default source_detail should be empty"
    print("  PASS: test_brand_style_intent_defaults")


def test_brand_style_intent_with_values():
    """BrandStyleIntent should accept and store all fields correctly."""
    intent = BrandStyleIntent(
        has_branding=True,
        brand_name="Nike",
        style_keywords=["bold", "sporty", "dynamic"],
        color_palette=["#FF6600", "#000000", "#FFFFFF"],
        tone_override="empowering",
        typography_hints=["Futura", "Helvetica Neue"],
        content_query="Create a presentation about AI trends in healthcare",
        source="query",
        source_detail="user query",
    )
    assert intent.has_branding is True
    assert intent.brand_name == "Nike"
    assert len(intent.style_keywords) == 3
    assert "#FF6600" in intent.color_palette
    assert intent.tone_override == "empowering"
    assert "Futura" in intent.typography_hints
    assert "AI trends" in intent.content_query
    print("  PASS: test_brand_style_intent_with_values")


def test_brand_style_intent_json_roundtrip():
    """BrandStyleIntent should survive JSON serialization/deserialization."""
    intent = BrandStyleIntent(
        has_branding=True,
        brand_name="Apple",
        style_keywords=["minimalist", "elegant"],
        color_palette=["#333333", "#F5F5F7"],
        tone_override="innovative",
        source="query",
    )
    json_str = intent.model_dump_json()
    restored = BrandStyleIntent.model_validate_json(json_str)
    assert restored.brand_name == "Apple"
    assert restored.color_palette == ["#333333", "#F5F5F7"]
    assert restored.tone_override == "innovative"
    print("  PASS: test_brand_style_intent_json_roundtrip")


def test_format_brand_context_empty():
    """_format_brand_context_for_prompt should return empty for no-brand intent."""
    intent = BrandStyleIntent()
    result = _format_brand_context_for_prompt(intent)
    assert result == "", "Should return empty string when has_branding is False"

    result_none = _format_brand_context_for_prompt(None)
    assert result_none == "", "Should return empty string for None input"
    print("  PASS: test_format_brand_context_empty")


def test_format_brand_context_with_brand():
    """_format_brand_context_for_prompt should produce a markdown section."""
    intent = BrandStyleIntent(
        has_branding=True,
        brand_name="Tesla",
        style_keywords=["futuristic", "clean", "premium"],
        color_palette=["#CC0000", "#000000"],
        tone_override="visionary",
        typography_hints=["Gotham"],
    )
    result = _format_brand_context_for_prompt(intent)
    assert "## Brand/Style Guidance" in result, "Should have header"
    assert "Tesla" in result, "Should mention brand name"
    assert "#CC0000" in result, "Should include hex color"
    assert "visionary" in result, "Should include tone"
    assert "Gotham" in result, "Should include typography"
    assert "futuristic" in result, "Should include style keywords"
    print("  PASS: test_format_brand_context_with_brand")


def test_build_brand_override_log():
    """_build_brand_override_log should produce descriptive log output."""
    query_intent = BrandStyleIntent(
        has_branding=True,
        brand_name="Nike",
        source="query",
    )
    template_intent = BrandStyleIntent(
        has_branding=True,
        brand_name="Acme Corp",
        source="template",
        source_detail="corporate_template.pptx",
        color_palette=["#003366", "#FF9900"],
        typography_hints=["Arial", "Calibri"],
    )
    log = _build_brand_override_log(query_intent, template_intent)

    assert "[BRAND OVERRIDE]" in log, "Should contain override prefix"
    assert "Nike branding" in log, "Should mention query brand"
    assert "corporate_template.pptx" in log, "Should mention template name"
    assert "disregarded" in log.lower(), "Should state query branding was disregarded"
    assert "#003366" in log, "Should include template colors"
    assert "Arial" in log, "Should include template fonts"
    print("  PASS: test_build_brand_override_log")


def test_extract_style_from_template_basic():
    """extract_style_from_template should read a minimal template without crashing."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name

    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = "Acme Corp"
        prs.save(tmp_path)

        intent = extract_style_from_template(tmp_path)
        assert intent.source == "template", "Source should be 'template'"
        assert intent.has_branding is True, "Should have branding from template"
        assert intent.source_detail == os.path.basename(tmp_path), \
            "source_detail should be the filename"
        # Should extract theme fonts and colors from the default template
        print("  PASS: test_extract_style_from_template_basic")
    finally:
        os.unlink(tmp_path)


def test_extract_style_from_template_company_name():
    """extract_style_from_template should detect company name from title slide."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name

    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = "TechCo"
        prs.save(tmp_path)

        intent = extract_style_from_template(tmp_path)
        assert intent.brand_name == "TechCo", \
            "Should detect 'TechCo' as company name, got: '%s'" % intent.brand_name
        print("  PASS: test_extract_style_from_template_company_name")
    finally:
        os.unlink(tmp_path)


def test_extract_style_from_template_nonexistent():
    """extract_style_from_template should handle missing file gracefully."""
    intent = extract_style_from_template("/tmp/nonexistent_template_abc123.pptx")
    assert intent.source == "template", "Source should still be 'template'"
    print("  PASS: test_extract_style_from_template_nonexistent")


def test_template_override_flow():
    """Full override scenario: query has brand + template present -> template wins."""
    # Simulate query-level brand intent
    query_intent = BrandStyleIntent(
        has_branding=True,
        brand_name="Nike",
        style_keywords=["sporty", "bold"],
        color_palette=["#FF6600"],
        tone_override="empowering",
        source="query",
    )

    # Create a template and extract its style
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name

    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = "Globex Corp"
        prs.save(tmp_path)

        template_intent = extract_style_from_template(tmp_path)

        # Template should override query
        assert query_intent.has_branding is True
        assert template_intent.source == "template"
        assert template_intent.brand_name == "Globex Corp"

        # Build override log
        log = _build_brand_override_log(query_intent, template_intent)
        assert "Nike branding" in log
        assert os.path.basename(tmp_path) in log

        # The effective intent should be template_intent, not query_intent
        effective = template_intent
        assert effective.source == "template"
        assert effective.brand_name == "Globex Corp"
        assert "Nike" not in effective.brand_name

        print("  PASS: test_template_override_flow")
    finally:
        os.unlink(tmp_path)


def run_all_tests():
    """Run all brand/style parsing tests."""
    print("=" * 60)
    print("Running Brand/Style Parsing Tests")
    print("=" * 60)

    tests = [
        test_brand_style_intent_defaults,
        test_brand_style_intent_with_values,
        test_brand_style_intent_json_roundtrip,
        test_format_brand_context_empty,
        test_format_brand_context_with_brand,
        test_build_brand_override_log,
        test_extract_style_from_template_basic,
        test_extract_style_from_template_company_name,
        test_extract_style_from_template_nonexistent,
        test_template_override_flow,
    ]

    passed = 0
    failed = 0
    for test_fn in tests:
        try:
            test_fn()
            passed += 1
        except Exception as e:
            print("  FAIL: %s — %s" % (test_fn.__name__, e))
            import traceback
            traceback.print_exc()
            failed += 1

    print("\n" + "=" * 60)
    print("Results: %d passed, %d failed" % (passed, failed))
    print("=" * 60)

    if failed > 0:
        sys.exit(1)
    else:
        print("ALL BRAND/STYLE PARSING TESTS PASSED")


if __name__ == "__main__":
    run_all_tests()
