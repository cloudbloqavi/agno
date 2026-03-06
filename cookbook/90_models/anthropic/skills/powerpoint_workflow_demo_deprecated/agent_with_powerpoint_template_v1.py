"""
Agno Agent with PowerPoint Skills + Custom Template.

This cookbook demonstrates how to use Claude's pptx skill to create PowerPoint
presentations and then apply a custom .pptx template's styling (theme, layouts,
fonts, colors) to the generated content.

Prerequisites:
- uv pip install agno anthropic python-pptx
- export ANTHROPIC_API_KEY="your_api_key_here"
- A .pptx template file

Usage:
    python agent_with_powerpoint_template.py --template my_template.pptx
    python agent_with_powerpoint_template.py -t my_template.pptx -o output.pptx
    python agent_with_powerpoint_template.py -t my_template.pptx -p "Create a 5-slide presentation about AI"
"""

import argparse
import copy
import os
import shutil
import sys
from dataclasses import dataclass, field
from io import BytesIO

from agno.agent import Agent
from agno.models.anthropic import Claude
from anthropic import Anthropic
from file_download_helper import download_skill_files
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Data Classes for Extracted Content
# ---------------------------------------------------------------------------


@dataclass
class TableData:
    """Extracted table data with position."""

    rows: list  # list of list of str (2D cell text grid)
    left: int  # EMU position
    top: int
    width: int
    height: int


@dataclass
class ImageData:
    """Extracted image data with position."""

    blob: bytes
    left: int
    top: int
    width: int
    height: int
    content_type: str = "image/png"


@dataclass
class ChartExtract:
    """Extracted chart data with position."""

    chart_type: int  # XL_CHART_TYPE enum value
    categories: list  # list of str
    series: list  # list of (name, values) tuples
    left: int
    top: int
    width: int
    height: int


@dataclass
class SlideContent:
    """All extracted content from a single slide."""

    title: str = ""
    subtitle: str = ""
    body_paragraphs: list = field(default_factory=list)  # list of (text, level) tuples
    tables: list = field(default_factory=list)  # list of TableData
    images: list = field(default_factory=list)  # list of ImageData
    charts: list = field(default_factory=list)  # list of ChartExtract
    shapes_xml: list = field(default_factory=list)  # list of lxml Element (deep copies)


@dataclass
class ContentArea:
    """Defines the safe content region on a template slide (all values in EMU)."""

    left: int
    top: int
    width: int
    height: int


# ---------------------------------------------------------------------------
# Template Application Logic
# ---------------------------------------------------------------------------


def _find_best_layout(template_prs, slide_index: int, total_slides: int):
    """Find the best matching layout from the template for a given slide position.

    Uses a heuristic based on layout name keywords and slide position:
    - First slide maps to title layouts
    - Last slide maps to Ending title/note or blank layouts
    - Middle slides map to content layouts

    Args:
        template_prs: The template Presentation object.
        slide_index: Zero-based index of the current slide.
        total_slides: Total number of slides in the generated presentation.

    Returns:
        The best matching slide layout from the template.
    """
    layouts = list(template_prs.slide_layouts)
    if not layouts:
        raise ValueError("Template has no slide layouts")

    layout_names = [(i, layout.name.lower()) for i, layout in enumerate(layouts)]

    # Determine what kind of slide we need
    is_title_slide = slide_index == 0
    is_last_slide = slide_index == total_slides - 1

    if is_title_slide:
        # Look for title slide layout
        for i, name in layout_names:
            if "title slide" in name or (
                "title" in name and "content" not in name and "only" not in name
            ):
                return layouts[i]
        # Fallback: any layout with "title" in the name
        for i, name in layout_names:
            if "title" in name:
                return layouts[i]
        return layouts[0]

    if is_last_slide:
        # Look for a closing/thank you layout, or fall back to content
        for i, name in layout_names:
            if "blank" in name or "closing" in name or "end" in name:
                return layouts[i]

    # Content slides: look for content or body layouts
    for i, name in layout_names:
        if "content" in name or "body" in name or "text" in name:
            return layouts[i]

    # Fallback: look for "two content" or similar
    for i, name in layout_names:
        if "object" in name or "list" in name:
            return layouts[i]

    # Last resort: use the second layout if available (typically a content layout),
    # otherwise use the first
    if len(layouts) > 1:
        return layouts[1]
    return layouts[0]


def _get_content_area(layout, slide_width: int, slide_height: int) -> ContentArea:
    """Derive the safe content area from a template layout's placeholders.

    Strategy:
    1. Look for a body placeholder (idx=1) — its position defines the content area.
    2. If no body placeholder, look for any placeholder with idx > 0.
    3. If no placeholders at all, use a default safe margin.

    Args:
        layout: A python-pptx SlideLayout object.
        slide_width: Presentation slide width in EMU.
        slide_height: Presentation slide height in EMU.

    Returns:
        ContentArea with the computed safe region.
    """
    # Try body placeholder first (idx=1)
    for ph in layout.placeholders:
        if ph.placeholder_format.idx == 1:
            return ContentArea(
                left=ph.left, top=ph.top, width=ph.width, height=ph.height
            )

    # Try any non-title placeholder
    for ph in layout.placeholders:
        if ph.placeholder_format.idx > 0:
            return ContentArea(
                left=ph.left, top=ph.top, width=ph.width, height=ph.height
            )

    # Default: safe margins (5% left, 25% top, 90% width, 65% height)
    return ContentArea(
        left=int(slide_width * 0.05),
        top=int(slide_height * 0.25),
        width=int(slide_width * 0.90),
        height=int(slide_height * 0.65),
    )


def _fit_to_area(img_width: int, img_height: int, area: ContentArea) -> tuple:
    """Scale dimensions to fit within an area while preserving aspect ratio.

    Args:
        img_width: Original width in EMU.
        img_height: Original height in EMU.
        area: Target ContentArea to fit within.

    Returns:
        Tuple of (left, top, width, height) in EMU, centered in the area.
    """
    if img_width <= 0 or img_height <= 0:
        return area.left, area.top, area.width, area.height

    aspect = img_width / img_height
    area_aspect = area.width / area.height

    if aspect > area_aspect:
        new_width = area.width
        new_height = int(new_width / aspect)
    else:
        new_height = area.height
        new_width = int(new_height * aspect)

    left = area.left + (area.width - new_width) // 2
    top = area.top + (area.height - new_height) // 2

    return left, top, new_width, new_height


def _extract_slide_content(slide) -> SlideContent:
    """Extract all content from a slide including text, tables, images, charts, and shapes.

    Args:
        slide: A python-pptx slide object.

    Returns:
        SlideContent with all extracted elements.
    """
    content = SlideContent()

    for shape in slide.shapes:
        # --- Tables ---
        if shape.has_table:
            table = shape.table
            rows_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                rows_data.append(row_data)
            content.tables.append(
                TableData(
                    rows=rows_data,
                    left=shape.left,
                    top=shape.top,
                    width=shape.width,
                    height=shape.height,
                )
            )
            continue

        # --- Charts ---
        if shape.has_chart:
            try:
                chart = shape.chart
                chart_type_val = chart.chart_type
                categories = []
                series_data = []

                # Extract chart data
                plot = chart.plots[0] if chart.plots else None
                if plot:
                    # Get categories
                    if plot.categories:
                        categories = list(plot.categories)

                    # Get series
                    for series in plot.series:
                        name = series.name if hasattr(series, "name") else ""
                        values = (
                            list(series.values) if hasattr(series, "values") else []
                        )
                        series_data.append((name or "", values))

                if categories or series_data:
                    content.charts.append(
                        ChartExtract(
                            chart_type=chart_type_val,
                            categories=categories,
                            series=series_data,
                            left=shape.left,
                            top=shape.top,
                            width=shape.width,
                            height=shape.height,
                        )
                    )
            except Exception:
                # Skip charts we can't parse
                pass
            continue

        # --- Images/Pictures ---
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                content_type = shape.image.content_type
                content.images.append(
                    ImageData(
                        blob=blob,
                        left=shape.left,
                        top=shape.top,
                        width=shape.width,
                        height=shape.height,
                        content_type=content_type,
                    )
                )
            except Exception:
                pass
            continue

        # --- Group Shapes ---
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Deep copy the group XML element
            group_xml = copy.deepcopy(shape._element)
            content.shapes_xml.append(group_xml)
            # Also extract images from within the group
            try:
                for grp_shape in shape.shapes:
                    if grp_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            blob = grp_shape.image.blob
                            ct = grp_shape.image.content_type
                            content.images.append(
                                ImageData(
                                    blob=blob,
                                    left=grp_shape.left,
                                    top=grp_shape.top,
                                    width=grp_shape.width,
                                    height=grp_shape.height,
                                    content_type=ct,
                                )
                            )
                        except Exception:
                            pass
            except Exception:
                pass
            continue

        # --- Text shapes (existing logic) ---
        if shape.has_text_frame:
            text_frame = shape.text_frame
            if shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx
                if ph_idx == 0:
                    content.title = text_frame.text.strip()
                elif ph_idx == 1:
                    paragraphs = text_frame.paragraphs
                    if len(paragraphs) == 1 and paragraphs[0].level == 0:
                        content.subtitle = paragraphs[0].text.strip()
                    else:
                        for para in paragraphs:
                            text = para.text.strip()
                            if text:
                                content.body_paragraphs.append((text, para.level))
                else:
                    # Other placeholder - treat as body
                    for para in text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            content.body_paragraphs.append((text, para.level))
            else:
                # Non-placeholder simple shapes - deep copy XML for transfer
                shape_xml = copy.deepcopy(shape._element)
                content.shapes_xml.append(shape_xml)
        else:
            # Non-text shapes (connectors, freeforms, etc.) - deep copy XML
            if not shape.is_placeholder:
                shape_xml = copy.deepcopy(shape._element)
                content.shapes_xml.append(shape_xml)

    return content


def _populate_placeholder_with_format(shape, texts, is_title=False):
    """Populate a placeholder shape with text while preserving template formatting.

    Captures the template paragraph's formatting before modifying anything,
    uses clear() on the text frame, then recreates paragraphs with cloned
    formatting from the original reference paragraph/run. Enables word wrap
    and attempts to auto-fit text to the placeholder bounds.

    Args:
        shape: A placeholder shape from the template.
        texts: For titles, a single string. For body, a list of (text, level) tuples.
        is_title: If True, texts is a string; if False, texts is a list of (text, level).
    """
    if not shape.has_text_frame:
        return

    tf = shape.text_frame
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    # Enable word wrap before anything else
    tf.word_wrap = True

    # Capture reference formatting from the first paragraph and its first run
    ref_para = tf.paragraphs[0] if tf.paragraphs else None
    ref_para_xml = None
    ref_run_xml = None

    if ref_para is not None:
        # Save paragraph properties XML
        pPr = ref_para._p.find(ns + "pPr")
        if pPr is not None:
            ref_para_xml = copy.deepcopy(pPr)

        # Save run properties XML from first run
        if ref_para.runs:
            rPr = ref_para.runs[0]._r.find(ns + "rPr")
            if rPr is not None:
                ref_run_xml = copy.deepcopy(rPr)

    # Clear the text frame
    tf.clear()

    if is_title:
        # Single text for title
        para = tf.paragraphs[0]
        # Restore paragraph formatting
        if ref_para_xml is not None:
            para._p.insert(0, copy.deepcopy(ref_para_xml))
        run = para.add_run()
        run.text = texts
        if ref_run_xml is not None:
            run._r.insert(0, copy.deepcopy(ref_run_xml))
    else:
        # Multiple paragraphs for body content
        for i, (text, level) in enumerate(texts):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()

            # Restore paragraph formatting from reference
            if ref_para_xml is not None:
                new_pPr = copy.deepcopy(ref_para_xml)
                # Update indent level if needed
                if level > 0:
                    new_pPr.set("lvl", str(level))
                para._p.insert(0, new_pPr)

            para.level = level

            run = para.add_run()
            run.text = text
            if ref_run_xml is not None:
                run._r.insert(0, copy.deepcopy(ref_run_xml))

    # Auto-fit text to placeholder bounds
    try:
        max_size = 28 if is_title else 18
        tf.fit_text(font_family="Calibri", max_size=max_size)
    except Exception:
        # fit_text requires font metrics; fall back to MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _transfer_tables(slide, tables, content_area: ContentArea):
    """Transfer extracted table data to a slide, repositioned to the content area.

    Args:
        slide: Target slide object.
        tables: List of TableData objects.
        content_area: Safe content region from the template layout.
    """
    TABLE_CELL_FONT_SIZE = Pt(10)
    TABLE_HEADER_FONT_SIZE = Pt(11)
    num_tables = len(tables)

    for t_idx, td in enumerate(tables):
        num_rows = len(td.rows)
        num_cols = len(td.rows[0]) if td.rows else 0
        if num_rows == 0 or num_cols == 0:
            continue

        # Stack multiple tables vertically within the content area
        if num_tables > 1:
            per_table_height = content_area.height // num_tables
            table_top = content_area.top + (t_idx * per_table_height)
            table_height = per_table_height
        else:
            table_top = content_area.top
            table_height = content_area.height

        table_shape = slide.shapes.add_table(
            num_rows,
            num_cols,
            content_area.left,
            table_top,
            content_area.width,
            table_height,
        )
        table = table_shape.table

        for r_idx, row_data in enumerate(td.rows):
            for c_idx, cell_text in enumerate(row_data):
                if c_idx < num_cols:
                    cell = table.cell(r_idx, c_idx)
                    cell.text = cell_text
                    # Control font size to prevent overflow
                    for para in cell.text_frame.paragraphs:
                        para.font.size = (
                            TABLE_HEADER_FONT_SIZE
                            if r_idx == 0
                            else TABLE_CELL_FONT_SIZE
                        )
                    cell.text_frame.word_wrap = True


def _transfer_images(slide, images, content_area: ContentArea):
    """Transfer extracted images to a slide, scaled to fit the content area.

    Args:
        slide: Target slide object.
        images: List of ImageData objects.
        content_area: Safe content region from the template layout.
    """
    for img in images:
        image_stream = BytesIO(img.blob)
        left, top, width, height = _fit_to_area(img.width, img.height, content_area)
        slide.shapes.add_picture(image_stream, left, top, width, height)


def _transfer_charts(slide, charts, content_area: ContentArea):
    """Transfer extracted chart data to a slide, sized to fill the content area.

    Args:
        slide: Target slide object.
        charts: List of ChartExtract objects.
        content_area: Safe content region from the template layout.
    """
    try:
        from pptx.chart.data import CategoryChartData
    except ImportError:
        return

    num_charts = len(charts)

    for c_idx, cd in enumerate(charts):
        try:
            chart_data = CategoryChartData()
            chart_data.categories = cd.categories

            for series_name, series_values in cd.series:
                # Ensure values are numeric
                clean_values = []
                for v in series_values:
                    if v is None:
                        clean_values.append(0)
                    elif isinstance(v, (int, float)):
                        clean_values.append(v)
                    else:
                        try:
                            clean_values.append(float(v))
                        except (ValueError, TypeError):
                            clean_values.append(0)
                chart_data.add_series(series_name, clean_values)

            # Stack multiple charts vertically within the content area
            if num_charts > 1:
                chart_height = content_area.height // num_charts
                chart_top = content_area.top + (c_idx * chart_height)
            else:
                chart_height = content_area.height
                chart_top = content_area.top

            slide.shapes.add_chart(
                cd.chart_type,
                content_area.left,
                chart_top,
                content_area.width,
                chart_height,
                chart_data,
            )
        except Exception:
            # Skip charts that can't be recreated
            pass


def _transfer_shapes(slide, shapes_xml):
    """Transfer simple shapes by deep-copying their XML to the target slide."""
    spTree = slide.shapes._spTree

    for shape_elem in shapes_xml:
        cloned = copy.deepcopy(shape_elem)
        # Re-assign shape ID to avoid collisions
        # Find the max existing shape ID and increment
        existing_ids = [
            int(sp.get("id", 0)) for sp in spTree.iter() if sp.get("id") is not None
        ]
        max_id = max(existing_ids) if existing_ids else 0

        # Update the shape's nvSpPr/nvGrpSpPr id
        for nv_elem in cloned.iter():
            if nv_elem.tag.endswith("}cNvPr"):
                max_id += 1
                nv_elem.set("id", str(max_id))

        spTree.append(cloned)


def _populate_slide(
    new_slide, content: SlideContent, slide_width: int, slide_height: int
):
    """Transfer all content into a new slide using template-aware positioning.

    Args:
        new_slide: The target slide (created from a template layout).
        content: SlideContent object with all extracted elements.
        slide_width: Presentation slide width in EMU.
        slide_height: Presentation slide height in EMU.
    """
    # Compute content area from the slide's layout
    content_area = _get_content_area(new_slide.slide_layout, slide_width, slide_height)

    # --- Text ---
    title_placed = False
    body_placed = False

    for shape in new_slide.placeholders:
        ph_idx = shape.placeholder_format.idx

        if ph_idx == 0 and content.title:
            _populate_placeholder_with_format(shape, content.title, is_title=True)
            title_placed = True

        elif ph_idx == 1:
            if content.body_paragraphs:
                _populate_placeholder_with_format(
                    shape, content.body_paragraphs, is_title=False
                )
                body_placed = True
            elif content.subtitle:
                _populate_placeholder_with_format(
                    shape, content.subtitle, is_title=True
                )
                body_placed = True

    if not body_placed and content.body_paragraphs:
        for shape in new_slide.placeholders:
            ph_idx = shape.placeholder_format.idx
            if ph_idx > 1 and shape.has_text_frame:
                _populate_placeholder_with_format(
                    shape, content.body_paragraphs, is_title=False
                )
                body_placed = True
                break

    # Fallback text boxes using content area bounds
    if not title_placed and content.title:
        txBox = new_slide.shapes.add_textbox(
            content_area.left,
            Inches(0.3),
            content_area.width,
            Inches(1.0),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = content.title
        for para in tf.paragraphs:
            para.font.size = Pt(28)
            para.font.bold = True

    if not body_placed and content.body_paragraphs:
        txBox = new_slide.shapes.add_textbox(
            content_area.left,
            content_area.top,
            content_area.width,
            content_area.height,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, (text, level) in enumerate(content.body_paragraphs):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text = text
            para.level = level
            para.font.size = Pt(18)
        try:
            tf.fit_text(font_family="Calibri", max_size=18)
        except Exception:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # --- Visual elements ---
    _transfer_tables(new_slide, content.tables, content_area)
    _transfer_images(new_slide, content.images, content_area)
    _transfer_charts(new_slide, content.charts, content_area)
    _transfer_shapes(new_slide, content.shapes_xml)


def apply_template(template_path: str, generated_path: str, output_path: str) -> str:
    """Apply a template's styling to a generated presentation's content.

    This function creates a new presentation based on the template file
    (preserving its theme, master slides, layouts, fonts, and colors),
    then transfers the text content from the generated presentation into
    slides created from the template's layouts.

    Args:
        template_path: Path to the user's .pptx template file.
        generated_path: Path to the AI-generated .pptx file.
        output_path: Path where the final output .pptx will be saved.

    Returns:
        The output file path.
    """
    print("Opening template: %s" % template_path)
    print("Opening generated presentation: %s" % generated_path)

    generated_prs = Presentation(generated_path)
    generated_slides = list(generated_prs.slides)
    total_slides = len(generated_slides)

    if total_slides == 0:
        print("Warning: Generated presentation has no slides. Copying template as-is.")
        shutil.copy2(template_path, output_path)
        return output_path

    print("Found %d slides in generated presentation" % total_slides)

    # Create a new presentation from the template (preserves theme/masters)
    # We copy the template file first, then open and modify it
    shutil.copy2(template_path, output_path)
    output_prs = Presentation(output_path)
    slide_width = output_prs.slide_width
    slide_height = output_prs.slide_height

    # Remove all existing slides from the template copy
    while len(output_prs.slides._sldIdLst) > 0:
        sldId = output_prs.slides._sldIdLst[0]
        rId = sldId.get(
            etree.QName(
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
                "id",
            )
        )
        if rId is not None:
            output_prs.part.drop_rel(rId)
        output_prs.slides._sldIdLst.remove(sldId)

    print("Cleared template slides. Transferring content...")

    # For each generated slide, find matching layout and transfer content
    for idx, gen_slide in enumerate(generated_slides):
        content = _extract_slide_content(gen_slide)
        layout = _find_best_layout(output_prs, idx, total_slides)

        visual_info = []
        if content.tables:
            visual_info.append("%d table(s)" % len(content.tables))
        if content.images:
            visual_info.append("%d image(s)" % len(content.images))
        if content.charts:
            visual_info.append("%d chart(s)" % len(content.charts))
        if content.shapes_xml:
            visual_info.append("%d shape(s)" % len(content.shapes_xml))
        visual_str = ", ".join(visual_info) if visual_info else "text only"

        print(
            "  Slide %d: layout '%s' | title: '%s' | %s"
            % (
                idx + 1,
                layout.name,
                content.title[:50] if content.title else "",
                visual_str,
            )
        )

        new_slide = output_prs.slides.add_slide(layout)
        _populate_slide(new_slide, content, slide_width, slide_height)

    output_prs.save(output_path)
    print("Saved output presentation: %s" % output_path)
    return output_path


def _build_prompt_with_template_context(user_prompt: str, template_path: str) -> str:
    """Enhance the prompt with template layout information.

    Inspects the template file to discover available slide layouts and appends
    structural guidance to the user prompt so the agent generates content that
    maps cleanly to the template's placeholders.

    Args:
        user_prompt: The original user-provided (or default) prompt.
        template_path: Path to the .pptx template file.

    Returns:
        The enhanced prompt string with template context appended.
    """
    try:
        prs = Presentation(template_path)
        layouts = [layout.name for layout in prs.slide_layouts]
        layout_info = ", ".join(layouts)
    except Exception:
        # If the template cannot be read, return the prompt unchanged
        return user_prompt

    enhanced = (
        user_prompt + "\n\n"
        "Important structural requirements for template compatibility:\n"
        "- Use one clear title and concise bullet points per slide.\n"
        "- Do not apply custom fonts, colors, or theme styling.\n"
        "- Tables and charts are supported and will be transferred to the template.\n"
        "- Keep tables to max 6 rows x 5 columns.\n"
        "- Use bar, column, line, or pie charts with clearly labeled data.\n"
        "- The template has these available layouts: " + layout_info + ".\n"
        "- Use standard slide ordering: Title Slide, then Content Slides, then Closing."
    )
    return enhanced


# ---------------------------------------------------------------------------
# Create Agent
# ---------------------------------------------------------------------------

# Create a simple agent with PowerPoint skills
powerpoint_agent = Agent(
    name="PowerPoint Creator",
    model=Claude(
        id="claude-sonnet-4-5-20250929",
        skills=[
            {"type": "anthropic", "skill_id": "pptx", "version": "latest"}
        ],  # Enable PowerPoint presentation skill
    ),
    instructions=[
        "You are a structured content generator for PowerPoint presentations.",
        "Your output will be extracted and remapped to a corporate template.",
        "Follow these rules strictly:",
        "",
        "SLIDE STRUCTURE:",
        "- Use exactly one clear, descriptive title per slide (this maps to the title placeholder).",
        "- Use concise bullet points for body content (this maps to the body placeholder).",
        "- Limit to 4-6 bullet points per slide, each bullet max ~15 words.",
        "- Keep subtitle text on title slides to a single short line.",
        "- Follow standard slide ordering: Title Slide, then Content Slides, then Closing Slide.",
        "- Each slide should be self-contained; do not reference other slides visually.",
        "",
        "VISUAL ELEMENTS (use when appropriate):",
        "- Include tables for data comparisons and structured information.",
        "- Keep tables concise: no more than 6 rows and 5 columns.",
        "- Use bar, column, line, or pie charts for data visualization when numbers are involved.",
        "- Keep charts simple with clearly labeled data and categories.",
        "- Images may be included when explicitly requested by the user.",
        "",
        "FORMATTING RESTRICTIONS:",
        "- Do NOT apply custom fonts, colors, or theme styling (a template will be applied).",
        "- Do NOT use SmartArt or complex nested graphic layouts.",
        "- Do NOT add speaker notes, animations, or transitions.",
        "- Keep visual elements positioned in the center/lower portion of slides.",
    ],
    markdown=True,
)


# ---------------------------------------------------------------------------
# CLI and Main
# ---------------------------------------------------------------------------


def parse_args():
    """Parse command-line arguments."""
    parser = argparse.ArgumentParser(
        description="Create a PowerPoint presentation using Claude and apply a custom template."
    )
    parser.add_argument(
        "--template",
        "-t",
        required=True,
        help="Path to the user's .pptx template file.",
    )
    parser.add_argument(
        "--output",
        "-o",
        default="presentation_from_template.pptx",
        help="Output filename (default: presentation_from_template.pptx).",
    )
    parser.add_argument(
        "--prompt",
        "-p",
        default=None,
        help="Custom prompt for the presentation content.",
    )
    return parser.parse_args()


if __name__ == "__main__":
    args = parse_args()

    # Set up output directory relative to script location
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_dir = os.path.join(script_dir, "output")
    os.makedirs(output_dir, exist_ok=True)

    # Validate template file exists
    if not os.path.isfile(args.template):
        print("Error: Template file not found: %s" % args.template)
        sys.exit(1)

    if not args.template.endswith(".pptx"):
        print("Error: Template file must be a .pptx file")
        sys.exit(1)

    # Check for API key
    if not os.getenv("ANTHROPIC_API_KEY"):
        raise ValueError("ANTHROPIC_API_KEY environment variable not set")

    print("=" * 60)
    print("Agno Agent with PowerPoint Skills + Custom Template")
    print("=" * 60)
    print("Template: %s" % args.template)
    print("Output:   %s" % args.output)

    # Use custom prompt or a sensible default
    prompt = args.prompt
    if not prompt:
        prompt = (
            "Create a 6-slide business presentation with the following structure.\n"
            "Rules: Use one clear title per slide, 4-6 concise bullet points for text slides,\n"
            "and include tables or charts where specified. Do not apply custom styling.\n"
            "\n"
            "Slide 1 - Title Slide:\n"
            "  Title: 'Strategic Overview 2026'\n"
            "  Subtitle: 'Annual Business Review and Forward Plan'\n"
            "\n"
            "Slide 2 - Market Analysis (text with bullets):\n"
            "  Title: 'Market Analysis'\n"
            "  Bullets: Key market trends, growth opportunities, competitive landscape,\n"
            "  emerging technologies, and regulatory changes.\n"
            "\n"
            "Slide 3 - Financial Performance (include a TABLE):\n"
            "  Title: 'Financial Performance'\n"
            "  Table with columns: Metric, Q1, Q2, Q3, Q4\n"
            "  Rows: Revenue ($M), Costs ($M), Profit ($M), Growth (%)\n"
            "  Use realistic sample data.\n"
            "\n"
            "Slide 4 - Revenue Trends (include a BAR CHART):\n"
            "  Title: 'Quarterly Revenue Trends'\n"
            "  Bar chart showing quarterly revenue for 2024 vs 2025.\n"
            "  Categories: Q1, Q2, Q3, Q4. Two series: '2024' and '2025'.\n"
            "  Use realistic sample data.\n"
            "\n"
            "Slide 5 - Our Strategy (text with bullets):\n"
            "  Title: 'Strategic Priorities'\n"
            "  Bullets: Three-pillar growth approach, market expansion,\n"
            "  product innovation, operational excellence, talent development.\n"
            "\n"
            "Slide 6 - Closing Slide:\n"
            "  Title: 'Next Steps'\n"
            "  Bullets: Implementation timeline, key milestones,\n"
            "  resource allocation, success metrics.\n"
            "\n"
            "Save as 'generated_content.pptx'"
        )

    # Enhance the prompt with template layout context
    prompt = _build_prompt_with_template_context(prompt, args.template)

    # -----------------------------------------------------------------
    # Phase 1: Content Generation
    # -----------------------------------------------------------------
    print("\n" + "-" * 60)
    print("Phase 1: Generating presentation content with Claude...")
    print("-" * 60 + "\n")

    response = powerpoint_agent.run(prompt)
    print(response.content)

    # Download the generated file to a temp location
    print("\n" + "-" * 60)
    print("Downloading generated file...")
    print("-" * 60)

    client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
    generated_file = None

    if response.messages:
        for msg in response.messages:
            if hasattr(msg, "provider_data") and msg.provider_data:
                files = download_skill_files(
                    msg.provider_data,
                    client,
                    output_dir=output_dir,
                )
                if files:
                    # Find the first valid .pptx file
                    for f in files:
                        if not f.endswith(".pptx"):
                            continue
                        try:
                            # Validate it's a real pptx
                            Presentation(f)
                            generated_file = f
                            print("Using generated presentation: %s" % generated_file)
                            break
                        except Exception:
                            print("Skipping invalid file: %s" % f)
                            continue

                    if not generated_file:
                        # Try all files regardless of extension
                        for f in files:
                            try:
                                Presentation(f)
                                generated_file = f
                                print(
                                    "Using generated presentation: %s" % generated_file
                                )
                                break
                            except Exception:
                                continue

                if generated_file:
                    break

    if not generated_file:
        print("Error: No presentation file was generated by the agent.")
        sys.exit(1)

    # -----------------------------------------------------------------
    # Phase 2: Template Application
    # -----------------------------------------------------------------
    print("\n" + "-" * 60)
    print("Phase 2: Applying template styling to generated content...")
    print("-" * 60 + "\n")

    output_path = args.output
    if not os.path.isabs(output_path):
        output_path = os.path.join(output_dir, output_path)

    output_file = apply_template(args.template, generated_file, output_path)

    print("\n" + "=" * 60)
    print("Done! Output saved to: %s" % output_file)
    print("=" * 60)
