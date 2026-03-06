import sys
import os

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def clean_presentation_visual_noise_and_contrast(prs) -> None:
    """Clean empty placeholder ghost text and strip hardcoded font colors."""
    ghost_texts = {
        "click to add title",
        "click to add subtitle",
        "click to add text",
        "click to add notes",
        "double-tap to add title",
        "double-tap to add subtitle",
        "double-tap to add text",
    }
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    for slide in prs.slides:
        spTree = slide.shapes._spTree
        elements_to_remove = []
        
        for shape in list(slide.shapes):
            # 1. Clean visual noise
            try:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip().lower()
                    if text in ghost_texts or text == "":
                        elements_to_remove.append(shape._element)
                        continue # Removed
            except Exception:
                pass
            
            # 2. Fix Text Contrast (strip hardcoded colors)
            try:
                if getattr(shape, "has_text_frame", False):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            rPr = run._r.find(ns_a + "rPr")
                            if rPr is not None:
                                solidFill = rPr.find(ns_a + "solidFill")
                                if solidFill is not None:
                                    rPr.remove(solidFill)
            except Exception:
                pass
            
            # Fix contrast in tables
            try:
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    rPr = run._r.find(ns_a + "rPr")
                                    if rPr is not None:
                                        solidFill = rPr.find(ns_a + "solidFill")
                                        if solidFill is not None:
                                            rPr.remove(solidFill)
            except Exception:
                pass
        
        for element in elements_to_remove:
            try:
                spTree.remove(element)
            except Exception:
                pass

def run_test():
    # Create test presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    
    # Slide 1: Test Visual Noise Cleanup
    slide1 = prs.slides.add_slide(blank_slide_layout)
    txBox1 = slide1.shapes.add_textbox(0, 0, 100, 100)
    tf1 = txBox1.text_frame
    tf1.text = "Click to add title"
    
    txBox2 = slide1.shapes.add_textbox(100, 100, 100, 100)
    tf2 = txBox2.text_frame
    tf2.text = "Valid text"

    # Slide 2: Test Contrast Cleanup (SolidFill stripping)
    slide2 = prs.slides.add_slide(blank_slide_layout)
    txBox3 = slide2.shapes.add_textbox(0, 0, 100, 100)
    tf3 = txBox3.text_frame
    
    p = tf3.add_paragraph()
    run = p.add_run()
    run.text = "This had bad contrast"
    # Set explicit font color
    run.font.color.rgb = RGBColor(0, 0, 0)
    
    # Run our function
    clean_presentation_visual_noise_and_contrast(prs)
    
    # Assertions for Slide 1
    # text frame 1 should be gone, text frame 2 should remain
    texts = []
    for shape in slide1.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    
    assert len(texts) == 1, f"Expected 1 text box, found {len(texts)}"
    assert texts[0] == "Valid text", f"Expected 'Valid text', got {texts[0]}"
    
    # Assertions for Slide 2
    # Ensure solidFill was removed
    shape3 = slide2.shapes[0]
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for para in shape3.text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.find(ns_a + "rPr")
            if rPr is not None:
                assert rPr.find(ns_a + "solidFill") is None, "solidFill was not removed!"
    
    print("ALL TESTS PASSED: Empty placeholders removed and solidFill stripped.")
    prs.save("test_output.pptx")

if __name__ == "__main__":
    run_test()
