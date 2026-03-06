"""
Agno Workflow: PowerPoint Template Generation Pipeline.

A sequential workflow that generates presentations using Claude's pptx skill,
intelligently adds AI-generated images via NanoBanana (powered by Gemini), and
applies a custom .pptx template for professional styling.

Prerequisites:
- uv pip install agno anthropic python-pptx google-genai pillow
- export ANTHROPIC_API_KEY="your_api_key_here"
- export GOOGLE_API_KEY="your_google_api_key_here"
- A .pptx template file

Usage:
    # Basic usage with a template:
    .venvs/demo/bin/python cookbook/90_models/anthropic/skills/powerpoint_template_workflow.py \\
        --template my_template.pptx

    # Full options: custom prompt, output path, verbose logging:
    .venvs/demo/bin/python cookbook/90_models/anthropic/skills/powerpoint_template_workflow.py \\
        -t my_template.pptx -o report.pptx -p "Create a 5-slide AI trends presentation" -v

    # Disable streaming (more reliable for shorter prompts, may timeout on complex ones):
    .venvs/demo/bin/python cookbook/90_models/anthropic/skills/powerpoint_template_workflow.py \\
        -t my_template.pptx --no-stream

    # Skip AI image generation entirely:
    .venvs/demo/bin/python cookbook/90_models/anthropic/skills/powerpoint_template_workflow.py \\
        -t my_template.pptx --no-images

CLI Flags:
    --template, -t   Path to the .pptx template file (required).
    --output, -o     Output filename (default: presentation_from_template.pptx).
    --prompt, -p     Custom prompt for the presentation content.
    --no-images      Skip AI image generation (Steps 2 and 3).
    --no-stream      Disable streaming mode for Claude agent (more reliable for
                     shorter prompts, but may timeout on complex presentations).
    --verbose, -v    Enable verbose/debug logging for troubleshooting.
"""

import argparse
import copy
import json
import os
import shutil
import sys
import traceback
from dataclasses import dataclass, field
from enum import Enum
from io import BytesIO
from typing import Dict, List

from agno.agent import Agent
from agno.models.anthropic import Claude
from agno.models.google import Gemini
from agno.run.agent import RunOutput
from agno.tools.nano_banana import NanoBananaTools
from agno.workflow.step import Step
from agno.workflow.types import StepInput, StepOutput
from agno.workflow.workflow import Workflow
from anthropic import Anthropic
from file_download_helper import download_skill_files
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field

# Module-level verbose flag (set from CLI args)
VERBOSE = False


# ---------------------------------------------------------------------------
# Pydantic Models for Structured Data Flow
# ---------------------------------------------------------------------------


class SlideImageDecision(BaseModel):
    """Decision about whether a slide needs an AI-generated image."""

    slide_index: int = Field(description="Zero-based index of the slide")
    needs_image: bool = Field(
        description="Whether this slide would benefit from an image"
    )
    image_prompt: str = Field(
        default="",
        description="If needs_image is True, a detailed prompt for generating the image. "
        "Should describe a professional, clean illustration suitable for a business presentation. "
        "If needs_image is False, leave empty.",
    )
    reasoning: str = Field(
        default="",
        description="Brief explanation of why the slide does or does not need an image",
    )


class ImagePlan(BaseModel):
    """Plan for which slides need AI-generated images."""

    decisions: List[SlideImageDecision] = Field(
        description="List of image decisions, one per slide"
    )


# ---------------------------------------------------------------------------
# Dataclasses for Slide Content (reused from agent_with_powerpoint_template)
# ---------------------------------------------------------------------------


@dataclass
class TableData:
    """Extracted table data with position."""

    rows: list
    left: int
    top: int
    width: int
    height: int


@dataclass
class ImageData:
    """Extracted image data with position."""

    blob: bytes
    left: int
    top: int
    width: int
    height: int
    content_type: str = "image/png"


@dataclass
class ChartExtract:
    """Extracted chart data with position."""

    chart_type: int
    categories: list
    series: list
    left: int
    top: int
    width: int
    height: int


@dataclass
class SlideContent:
    """All extracted content from a single slide."""

    title: str = ""
    subtitle: str = ""
    body_paragraphs: list = field(default_factory=list)
    tables: list = field(default_factory=list)
    images: list = field(default_factory=list)
    charts: list = field(default_factory=list)
    shapes_xml: list = field(default_factory=list)
    text_shapes_xml: list = field(default_factory=list)
    text_box_paragraphs: list = field(default_factory=list)
    # Track picture placeholders detected in the slide layout
    has_image_placeholder: bool = False
    image_placeholder_indices: list = field(default_factory=list)
    # Tracks whether the slide contains any non-placeholder text boxes
    has_text_box: bool = False


@dataclass
class ContentArea:
    """Defines the safe content region on a template slide (all values in EMU)."""

    left: int
    top: int
    width: int
    height: int


class ContentMix(Enum):
    """Classification of what element types a slide contains."""

    TEXT_ONLY = "text_only"
    TEXT_AND_TABLE = "text_and_table"
    TEXT_AND_CHART = "text_and_chart"
    TEXT_AND_IMAGE = "text_and_image"
    TEXT_AND_GENERATED_IMAGE = "text_and_generated_image"
    VISUAL_ONLY = "visual_only"
    MIXED = "mixed"  # Multiple visual types


@dataclass
class RegionMap:
    """Defines separate regions for text and visual elements on a slide."""

    text_region: ContentArea
    visual_region: ContentArea
    layout_type: str = "full"  # "full", "split_horizontal", "split_vertical", "native"


@dataclass
class TemplateTheme:
    """Extracted theme colors and font scheme from a template."""

    accent_colors: list = field(default_factory=list)  # Up to 6 accent hex colors
    dk1: str = ""  # Dark 1 color (hex)
    dk2: str = ""  # Dark 2 color (hex)
    lt1: str = ""  # Light 1 color (hex)
    lt2: str = ""  # Light 2 color (hex)
    hyperlink: str = ""
    followed_hyperlink: str = ""
    major_font: str = "Calibri"  # Heading font
    minor_font: str = "Calibri"  # Body font


@dataclass
class TemplateTableStyle:
    """Styling extracted from a reference template table."""

    header_fill_rgb: str = ""  # hex like "4472C4"
    header_font_color_rgb: str = ""
    header_font_size_pt: int = 11
    header_font_bold: bool = True
    header_font_family: str = ""
    cell_font_size_pt: int = 10
    cell_font_color_rgb: str = ""
    cell_font_family: str = ""
    band_row_color_rgb: str = ""  # Alternating row fill
    border_color_rgb: str = ""
    border_width_emu: int = 12700  # 1pt default
    cell_fill_rgb: str = ""  # Default cell fill
    # Raw XML for reliable deep-copy transfer
    table_properties_xml: object = None  # lxml Element or None
    table_style_id: str = ""


@dataclass
class TemplateChartStyle:
    """Styling extracted from a reference template chart."""

    chart_type: int = -1
    series_fill_colors: list = field(default_factory=list)  # list of hex RGB
    series_line_colors: list = field(default_factory=list)
    axis_font_size_pt: int = 10
    axis_font_color_rgb: str = ""
    axis_font_family: str = ""
    axis_line_color_rgb: str = ""
    legend_font_size_pt: int = 10
    legend_font_color_rgb: str = ""
    legend_font_family: str = ""
    data_label_font_size_pt: int = 9
    data_label_font_color_rgb: str = ""
    plot_area_fill_rgb: str = ""
    chart_bg_fill_rgb: str = ""
    # Store the raw chartStyle XML element for full fidelity transfer
    chart_style_xml: object = None


@dataclass
class TemplateStyle:
    """Complete style information extracted from a template presentation."""

    theme: TemplateTheme = field(default_factory=TemplateTheme)
    table_styles: list = field(default_factory=list)
    chart_styles: list = field(default_factory=list)
    body_font_size_pt: int = 18
    body_font_color_rgb: str = ""
    body_font_family: str = ""
    title_font_size_pt: int = 28
    title_font_color_rgb: str = ""
    title_font_family: str = ""


# ---------------------------------------------------------------------------
# Template Style Extraction Functions
# ---------------------------------------------------------------------------


def _extract_theme_from_prs(prs) -> TemplateTheme:
    """Extract theme colors and font scheme from a presentation."""
    theme = TemplateTheme()
    try:
        # Access the theme through the slide master's relationships
        slide_master = prs.slide_masters[0]
        theme_part = None
        for rel in slide_master.part.rels.values():
            if "theme" in rel.reltype:
                theme_part = rel.target_part
                break
        if theme_part is None:
            return theme

        theme_xml = theme_part.element
        ns = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        }

        # Extract color scheme
        clrScheme = theme_xml.find(".//a:clrScheme", ns)
        if clrScheme is not None:
            accent_colors = []
            for child in clrScheme:
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                # Get the color value from srgbClr or sysClr
                color_val = ""
                srgb = child.find("a:srgbClr", ns)
                if srgb is not None:
                    color_val = srgb.get("val", "")
                else:
                    sys_clr = child.find("a:sysClr", ns)
                    if sys_clr is not None:
                        color_val = sys_clr.get("lastClr", "") or sys_clr.get("val", "")

                if tag.startswith("accent"):
                    accent_colors.append(color_val)
                elif tag == "dk1":
                    theme.dk1 = color_val
                elif tag == "dk2":
                    theme.dk2 = color_val
                elif tag == "lt1":
                    theme.lt1 = color_val
                elif tag == "lt2":
                    theme.lt2 = color_val
                elif tag == "hlink":
                    theme.hyperlink = color_val
                elif tag == "folHlink":
                    theme.followed_hyperlink = color_val
            theme.accent_colors = accent_colors

        # Extract font scheme
        fontScheme = theme_xml.find(".//a:fontScheme", ns)
        if fontScheme is not None:
            major_font_elem = fontScheme.find(".//a:majorFont/a:latin", ns)
            if major_font_elem is not None:
                theme.major_font = major_font_elem.get("typeface", "Calibri")
            minor_font_elem = fontScheme.find(".//a:minorFont/a:latin", ns)
            if minor_font_elem is not None:
                theme.minor_font = minor_font_elem.get("typeface", "Calibri")
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Theme extraction error: %s" % str(e))
    return theme


def _resolve_scheme_color(scheme_val: str, theme: TemplateTheme) -> str:
    """Resolve a scheme color name to an RGB hex value using the theme."""
    scheme_map = {
        "dk1": theme.dk1,
        "dk2": theme.dk2,
        "lt1": theme.lt1,
        "lt2": theme.lt2,
        "tx1": theme.dk1,  # tx1 typically maps to dk1
        "tx2": theme.dk2,  # tx2 typically maps to dk2
        "bg1": theme.lt1,  # bg1 typically maps to lt1
        "bg2": theme.lt2,  # bg2 typically maps to lt2
        "accent1": theme.accent_colors[0] if len(theme.accent_colors) > 0 else "",
        "accent2": theme.accent_colors[1] if len(theme.accent_colors) > 1 else "",
        "accent3": theme.accent_colors[2] if len(theme.accent_colors) > 2 else "",
        "accent4": theme.accent_colors[3] if len(theme.accent_colors) > 3 else "",
        "accent5": theme.accent_colors[4] if len(theme.accent_colors) > 4 else "",
        "accent6": theme.accent_colors[5] if len(theme.accent_colors) > 5 else "",
        "hlink": theme.hyperlink,
        "folHlink": theme.followed_hyperlink,
    }
    return scheme_map.get(scheme_val, "")


def _extract_color_from_rPr(rPr_elem, theme: TemplateTheme) -> str:
    """Extract RGB hex color from an rPr or defRPr element, resolving scheme colors."""
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    solidFill = rPr_elem.find(ns_a + "solidFill")
    if solidFill is not None:
        srgb = solidFill.find(ns_a + "srgbClr")
        if srgb is not None:
            return srgb.get("val", "")
        schemeClr = solidFill.find(ns_a + "schemeClr")
        if schemeClr is not None:
            scheme_val = schemeClr.get("val", "")
            return _resolve_scheme_color(scheme_val, theme)
    return ""


def _hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color string to (R, G, B) tuple (0-255)."""
    hex_color = hex_color.strip().lstrip("#")
    if len(hex_color) != 6:
        return (0, 0, 0)
    try:
        return (
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )
    except ValueError:
        return (0, 0, 0)


def _relative_luminance(r: int, g: int, b: int) -> float:
    """Calculate relative luminance per WCAG 2.0 formula.

    Returns a value between 0.0 (black) and 1.0 (white).
    """

    def _linearize(c):
        c_srgb = c / 255.0
        if c_srgb <= 0.03928:
            return c_srgb / 12.92
        return ((c_srgb + 0.055) / 1.055) ** 2.4

    return 0.2126 * _linearize(r) + 0.7152 * _linearize(g) + 0.0722 * _linearize(b)


def _contrast_ratio(color1_hex: str, color2_hex: str) -> float:
    """Calculate WCAG contrast ratio between two hex colors.

    Returns a ratio >= 1.0. WCAG AA requires >= 4.5 for normal text.
    """
    r1, g1, b1 = _hex_to_rgb(color1_hex)
    r2, g2, b2 = _hex_to_rgb(color2_hex)
    l1 = _relative_luminance(r1, g1, b1)
    l2 = _relative_luminance(r2, g2, b2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _get_slide_background_color(
    slide, template_style: "TemplateStyle | None" = None
) -> str:
    """Determine the effective background color of a slide.

    Checks (in priority order):
    1. Slide's own background fill (solid fill)
    2. Slide layout's background fill
    3. Slide master's background fill
    4. Template theme lt1 (light 1, typically white)
    5. Default white "FFFFFF"

    Returns hex RGB string (e.g., "FFFFFF").
    """
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    ns_p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

    # Helper to extract color from a background element
    def _extract_bg_color(element):
        if element is None:
            return ""
        # Look for <p:bg><p:bgPr><a:solidFill>
        bgPr = element.find(ns_p + "bgPr")
        if bgPr is None:
            bgPr = element.find(ns_p + "bgRef")
        if bgPr is None:
            bgPr = element

        solidFill = bgPr.find(ns_a + "solidFill")
        if solidFill is not None:
            srgb = solidFill.find(ns_a + "srgbClr")
            if srgb is not None:
                return srgb.get("val", "")
            schemeClr = solidFill.find(ns_a + "schemeClr")
            if schemeClr is not None and template_style:
                return _resolve_scheme_color(
                    schemeClr.get("val", ""), template_style.theme
                )
        return ""

    try:
        # 1. Check slide background
        bg = slide._element.find(ns_p + "cSld/" + ns_p + "bg")
        color = _extract_bg_color(bg)
        if color:
            return color
    except Exception:
        pass

    try:
        # 2. Check slide layout background
        bg = slide.slide_layout._element.find(ns_p + "cSld/" + ns_p + "bg")
        color = _extract_bg_color(bg)
        if color:
            return color
    except Exception:
        pass

    try:
        # 3. Check slide master background
        bg = slide.slide_layout.slide_master._element.find(ns_p + "cSld/" + ns_p + "bg")
        color = _extract_bg_color(bg)
        if color:
            return color
    except Exception:
        pass

    # 4. Use theme lt1 (usually white or near-white)
    if template_style and template_style.theme.lt1:
        return template_style.theme.lt1

    # 5. Default white
    return "FFFFFF"


def _ensure_text_contrast(slide, template_style: "TemplateStyle | None" = None) -> None:
    """Check all text elements on a slide and fix poor contrast.

    For each text run, checks if the text color has sufficient contrast
    against the slide's background. If contrast ratio < 3.0 (below WCAG AA),
    swaps the text to a high-contrast alternative.

    Uses theme dk1 (dark color) for text on light backgrounds, and
    theme lt1 (light color) for text on dark backgrounds.
    """
    if template_style is None:
        return

    bg_color = _get_slide_background_color(slide, template_style)
    if not bg_color:
        return

    bg_r, bg_g, bg_b = _hex_to_rgb(bg_color)
    bg_luminance = _relative_luminance(bg_r, bg_g, bg_b)

    # Determine the contrasting color to use
    # If background is light (luminance > 0.5), use dark text
    # If background is dark (luminance <= 0.5), use light text
    theme = template_style.theme
    if bg_luminance > 0.5:
        # Light background: use dark text
        contrast_color = theme.dk1 or "000000"
    else:
        # Dark background: use light text
        contrast_color = theme.lt1 or "FFFFFF"

    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    CONTRAST_THRESHOLD = 3.0  # Minimum acceptable contrast ratio

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                # Get the current text color
                text_color = ""
                rPr = run._r.find(ns_a + "rPr")
                if rPr is not None:
                    solidFill = rPr.find(ns_a + "solidFill")
                    if solidFill is not None:
                        srgb = solidFill.find(ns_a + "srgbClr")
                        if srgb is not None:
                            text_color = srgb.get("val", "")
                        elif template_style:
                            schemeClr = solidFill.find(ns_a + "schemeClr")
                            if schemeClr is not None:
                                text_color = _resolve_scheme_color(
                                    schemeClr.get("val", ""), theme
                                )

                if not text_color:
                    continue  # No explicit color, skip (inherits from theme which should be correct)

                # Check contrast
                ratio = _contrast_ratio(text_color, bg_color)

                if ratio < CONTRAST_THRESHOLD:
                    if VERBOSE:
                        print(
                            "[VERBOSE] Low contrast detected: text=%s bg=%s ratio=%.1f, fixing to %s"
                            % (text_color, bg_color, ratio, contrast_color)
                        )

                    # Replace the text color with the contrasting color
                    if rPr is not None:
                        solidFill = rPr.find(ns_a + "solidFill")
                        if solidFill is not None:
                            rPr.remove(solidFill)
                        # Create new solidFill with contrast color
                        new_fill = etree.SubElement(rPr, ns_a + "solidFill")
                        srgb = etree.SubElement(new_fill, ns_a + "srgbClr")
                        srgb.set("val", contrast_color)


def _extract_table_styles_from_prs(prs) -> list:
    """Scan template slides for tables and extract their styling."""
    styles = []
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            style = TemplateTableStyle()

            # Extract tblPr XML (contains tblStyle reference and banding info)
            tbl_elem = table._tbl
            tblPr = tbl_elem.find("a:tblPr", ns)
            if tblPr is not None:
                style.table_properties_xml = copy.deepcopy(tblPr)
                style.table_style_id = tblPr.get("tblStyle", "") or ""

            # Extract header row styling (first row)
            if len(table.rows) > 0:
                first_row = table.rows[0]
                for cell in first_row.cells:
                    tf = cell.text_frame
                    if tf and tf.paragraphs:
                        for para in tf.paragraphs:
                            if para.runs:
                                run = para.runs[0]
                                rPr = run._r.find("a:rPr", ns)
                                if rPr is not None:
                                    # Font size
                                    sz = rPr.get("sz")
                                    if sz:
                                        style.header_font_size_pt = int(sz) // 100
                                    # Bold
                                    b = rPr.get("b")
                                    if b:
                                        style.header_font_bold = b in ("1", "true")
                                    # Font color
                                    solidFill = rPr.find("a:solidFill", ns)
                                    if solidFill is not None:
                                        srgb = solidFill.find("a:srgbClr", ns)
                                        if srgb is not None:
                                            style.header_font_color_rgb = srgb.get(
                                                "val", ""
                                            )
                                    # Font family
                                    latin = rPr.find("a:latin", ns)
                                    if latin is not None:
                                        style.header_font_family = latin.get(
                                            "typeface", ""
                                        )

                    # Cell fill
                    tc_elem = cell._tc
                    tcPr = tc_elem.find("a:tcPr", ns)
                    if tcPr is not None:
                        solidFill = tcPr.find("a:solidFill", ns)
                        if solidFill is not None:
                            srgb = solidFill.find("a:srgbClr", ns)
                            if srgb is not None:
                                style.header_fill_rgb = srgb.get("val", "")
                    break  # Only need first cell of header

            # Extract body cell styling (second row if exists)
            if len(table.rows) > 1:
                second_row = table.rows[1]
                for cell in second_row.cells:
                    tf = cell.text_frame
                    if tf and tf.paragraphs:
                        for para in tf.paragraphs:
                            if para.runs:
                                run = para.runs[0]
                                rPr = run._r.find("a:rPr", ns)
                                if rPr is not None:
                                    sz = rPr.get("sz")
                                    if sz:
                                        style.cell_font_size_pt = int(sz) // 100
                                    solidFill = rPr.find("a:solidFill", ns)
                                    if solidFill is not None:
                                        srgb = solidFill.find("a:srgbClr", ns)
                                        if srgb is not None:
                                            style.cell_font_color_rgb = srgb.get(
                                                "val", ""
                                            )
                                    latin = rPr.find("a:latin", ns)
                                    if latin is not None:
                                        style.cell_font_family = latin.get(
                                            "typeface", ""
                                        )

                    # Cell fill
                    tc_elem = cell._tc
                    tcPr = tc_elem.find("a:tcPr", ns)
                    if tcPr is not None:
                        solidFill = tcPr.find("a:solidFill", ns)
                        if solidFill is not None:
                            srgb = solidFill.find("a:srgbClr", ns)
                            if srgb is not None:
                                style.cell_fill_rgb = srgb.get("val", "")
                    break

            # Extract border styling from first cell
            if len(table.rows) > 0:
                for cell in table.rows[0].cells:
                    tc_elem = cell._tc
                    tcPr = tc_elem.find("a:tcPr", ns)
                    if tcPr is not None:
                        # Look for border elements (lnL, lnR, lnT, lnB)
                        for border_tag in ["a:lnB", "a:lnT", "a:lnL", "a:lnR"]:
                            ln = tcPr.find(border_tag, ns)
                            if ln is not None:
                                w = ln.get("w")
                                if w:
                                    style.border_width_emu = int(w)
                                solidFill = ln.find("a:solidFill", ns)
                                if solidFill is not None:
                                    srgb = solidFill.find("a:srgbClr", ns)
                                    if srgb is not None:
                                        style.border_color_rgb = srgb.get("val", "")
                                break  # One border is enough
                    break

            styles.append(style)
    return styles


def _extract_chart_styles_from_prs(prs) -> list:
    """Scan template slides for charts and extract their styling."""
    styles = []
    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_chart:
                continue
            chart = shape.chart
            style = TemplateChartStyle()

            try:
                style.chart_type = chart.chart_type
            except Exception:
                pass

            # Extract series colors from the chart XML
            chart_elem = chart._chartSpace

            # Find series elements and extract their fill colors
            series_elems = chart_elem.findall(".//c:ser", ns)
            for ser in series_elems:
                # Check for solidFill in the series
                spPr = ser.find("c:spPr", ns)
                if spPr is not None:
                    solidFill = spPr.find("a:solidFill", ns)
                    if solidFill is not None:
                        srgb = solidFill.find("a:srgbClr", ns)
                        if srgb is not None:
                            style.series_fill_colors.append(srgb.get("val", ""))
                        else:
                            schemeClr = solidFill.find("a:schemeClr", ns)
                            if schemeClr is not None:
                                style.series_fill_colors.append(
                                    "scheme:" + schemeClr.get("val", "")
                                )
                    # Line color
                    ln = spPr.find("a:ln", ns)
                    if ln is not None:
                        lnFill = ln.find("a:solidFill", ns)
                        if lnFill is not None:
                            srgb = lnFill.find("a:srgbClr", ns)
                            if srgb is not None:
                                style.series_line_colors.append(srgb.get("val", ""))

            # Extract axis styling
            for axis_tag in ["c:catAx", "c:valAx", "c:dateAx"]:
                ax = chart_elem.find(".//" + axis_tag, ns)
                if ax is not None:
                    # Axis text properties
                    txPr = ax.find("c:txPr", ns)
                    if txPr is not None:
                        defRPr = txPr.find(".//a:defRPr", ns)
                        if defRPr is not None:
                            sz = defRPr.get("sz")
                            if sz:
                                style.axis_font_size_pt = int(sz) // 100
                            solidFill = defRPr.find("a:solidFill", ns)
                            if solidFill is not None:
                                srgb = solidFill.find("a:srgbClr", ns)
                                if srgb is not None:
                                    style.axis_font_color_rgb = srgb.get("val", "")
                            latin = defRPr.find("a:latin", ns)
                            if latin is not None:
                                style.axis_font_family = latin.get("typeface", "")
                    # Axis line
                    spPr = ax.find("c:spPr", ns)
                    if spPr is not None:
                        ln = spPr.find("a:ln", ns)
                        if ln is not None:
                            sf = ln.find("a:solidFill", ns)
                            if sf is not None:
                                srgb = sf.find("a:srgbClr", ns)
                                if srgb is not None:
                                    style.axis_line_color_rgb = srgb.get("val", "")
                    break  # Use first axis found

            # Extract legend styling
            legend = chart_elem.find(".//c:legend", ns)
            if legend is not None:
                txPr = legend.find("c:txPr", ns)
                if txPr is not None:
                    defRPr = txPr.find(".//a:defRPr", ns)
                    if defRPr is not None:
                        sz = defRPr.get("sz")
                        if sz:
                            style.legend_font_size_pt = int(sz) // 100
                        solidFill = defRPr.find("a:solidFill", ns)
                        if solidFill is not None:
                            srgb = solidFill.find("a:srgbClr", ns)
                            if srgb is not None:
                                style.legend_font_color_rgb = srgb.get("val", "")
                        latin = defRPr.find("a:latin", ns)
                        if latin is not None:
                            style.legend_font_family = latin.get("typeface", "")

            # Extract data label styling
            dLbls = chart_elem.find(".//c:dLbls", ns)
            if dLbls is not None:
                txPr = dLbls.find("c:txPr", ns)
                if txPr is not None:
                    defRPr = txPr.find(".//a:defRPr", ns)
                    if defRPr is not None:
                        sz = defRPr.get("sz")
                        if sz:
                            style.data_label_font_size_pt = int(sz) // 100
                        solidFill = defRPr.find("a:solidFill", ns)
                        if solidFill is not None:
                            srgb = solidFill.find("a:srgbClr", ns)
                            if srgb is not None:
                                style.data_label_font_color_rgb = srgb.get("val", "")

            # Plot area fill
            plotArea = chart_elem.find(".//c:plotArea", ns)
            if plotArea is not None:
                spPr = plotArea.find("c:spPr", ns)
                if spPr is not None:
                    solidFill = spPr.find("a:solidFill", ns)
                    if solidFill is not None:
                        srgb = solidFill.find("a:srgbClr", ns)
                        if srgb is not None:
                            style.plot_area_fill_rgb = srgb.get("val", "")

            styles.append(style)
    return styles


def _extract_template_styles(template_prs) -> TemplateStyle:
    """Extract all styling information from a template presentation."""
    ts = TemplateStyle()
    try:
        ts.theme = _extract_theme_from_prs(template_prs)
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Theme extraction failed: %s" % str(e))
    try:
        ts.table_styles = _extract_table_styles_from_prs(template_prs)
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Table style extraction failed: %s" % str(e))
    try:
        ts.chart_styles = _extract_chart_styles_from_prs(template_prs)
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Chart style extraction failed: %s" % str(e))

    # Extract body/title font sizes from template placeholders
    try:
        for slide in template_prs.slides:
            for shape in slide.placeholders:
                ph_idx = shape.placeholder_format.idx
                if ph_idx == 0 and shape.has_text_frame:  # Title
                    for para in shape.text_frame.paragraphs:
                        if para.runs:
                            rPr = para.runs[0]._r.find(
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
                            )
                            if rPr is not None:
                                sz = rPr.get("sz")
                                if sz and not ts.title_font_size_pt:
                                    ts.title_font_size_pt = int(sz) // 100
                                color_rgb = _extract_color_from_rPr(rPr, ts.theme)
                                if color_rgb and not ts.title_font_color_rgb:
                                    ts.title_font_color_rgb = color_rgb
                                latin = rPr.find(
                                    "{http://schemas.openxmlformats.org/drawingml/2006/main}latin"
                                )
                                font_name = ""
                                if latin is not None:
                                    typeface = latin.get("typeface", "")
                                    if typeface and not typeface.startswith("+"):
                                        font_name = typeface
                                    elif typeface == "+mj-lt":
                                        font_name = ts.theme.major_font
                                    elif typeface == "+mn-lt":
                                        font_name = ts.theme.minor_font
                                if font_name and not ts.title_font_family:
                                    ts.title_font_family = font_name
                elif ph_idx == 1 and shape.has_text_frame:  # Body
                    for para in shape.text_frame.paragraphs:
                        if para.runs:
                            rPr = para.runs[0]._r.find(
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
                            )
                            if rPr is not None:
                                sz = rPr.get("sz")
                                if sz and not ts.body_font_size_pt:
                                    ts.body_font_size_pt = int(sz) // 100
                                color_rgb = _extract_color_from_rPr(rPr, ts.theme)
                                if color_rgb and not ts.body_font_color_rgb:
                                    ts.body_font_color_rgb = color_rgb
                                latin = rPr.find(
                                    "{http://schemas.openxmlformats.org/drawingml/2006/main}latin"
                                )
                                font_name = ""
                                if latin is not None:
                                    typeface = latin.get("typeface", "")
                                    if typeface and not typeface.startswith("+"):
                                        font_name = typeface
                                    elif typeface == "+mj-lt":
                                        font_name = ts.theme.major_font
                                    elif typeface == "+mn-lt":
                                        font_name = ts.theme.minor_font
                                if font_name and not ts.body_font_family:
                                    ts.body_font_family = font_name
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Placeholder style extraction failed: %s" % str(e))

    # If slide-level extraction didn't find explicit colors, scan slide layouts
    # Template placeholders often inherit their styling from the layout/master level
    if not ts.title_font_color_rgb or not ts.body_font_color_rgb:
        try:
            ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            for layout in template_prs.slide_layouts:
                for ph in layout.placeholders:
                    ph_idx = ph.placeholder_format.idx
                    if ph_idx in (0, 1) and ph.has_text_frame:
                        for para in ph.text_frame.paragraphs:
                            # Check default run properties on the paragraph level
                            pPr = para._p.find(ns_a + "pPr")
                            if pPr is not None:
                                defRPr = pPr.find(ns_a + "defRPr")
                                if defRPr is not None:
                                    color_rgb = _extract_color_from_rPr(
                                        defRPr, ts.theme
                                    )
                                    sz = defRPr.get("sz")
                                    if (
                                        ph_idx == 0
                                        and not ts.title_font_color_rgb
                                        and color_rgb
                                    ):
                                        ts.title_font_color_rgb = color_rgb
                                    if ph_idx == 0 and not ts.title_font_size_pt and sz:
                                        ts.title_font_size_pt = int(sz) // 100
                                    if (
                                        ph_idx == 1
                                        and not ts.body_font_color_rgb
                                        and color_rgb
                                    ):
                                        ts.body_font_color_rgb = color_rgb
                                    if ph_idx == 1 and not ts.body_font_size_pt and sz:
                                        ts.body_font_size_pt = int(sz) // 100
                                    # Also extract font family from defRPr
                                    latin = defRPr.find(ns_a + "latin")
                                    font_name = ""
                                    if latin is not None:
                                        typeface = latin.get("typeface", "")
                                        if typeface and not typeface.startswith("+"):
                                            font_name = typeface
                                        elif typeface == "+mj-lt":
                                            font_name = ts.theme.major_font
                                        elif typeface == "+mn-lt":
                                            font_name = ts.theme.minor_font
                                    if (
                                        ph_idx == 0
                                        and not ts.title_font_family
                                        and font_name
                                    ):
                                        ts.title_font_family = font_name
                                    if (
                                        ph_idx == 1
                                        and not ts.body_font_family
                                        and font_name
                                    ):
                                        ts.body_font_family = font_name
                            # Also check actual runs
                            if para.runs:
                                rPr = para.runs[0]._r.find(ns_a + "rPr")
                                if rPr is not None:
                                    color_rgb = _extract_color_from_rPr(rPr, ts.theme)
                                    sz = rPr.get("sz")
                                    if (
                                        ph_idx == 0
                                        and not ts.title_font_color_rgb
                                        and color_rgb
                                    ):
                                        ts.title_font_color_rgb = color_rgb
                                    if ph_idx == 0 and not ts.title_font_size_pt and sz:
                                        ts.title_font_size_pt = int(sz) // 100
                                    if (
                                        ph_idx == 1
                                        and not ts.body_font_color_rgb
                                        and color_rgb
                                    ):
                                        ts.body_font_color_rgb = color_rgb
                                    if ph_idx == 1 and not ts.body_font_size_pt and sz:
                                        ts.body_font_size_pt = int(sz) // 100
                                    # Also extract font family from run rPr
                                    latin = rPr.find(ns_a + "latin")
                                    font_name = ""
                                    if latin is not None:
                                        typeface = latin.get("typeface", "")
                                        if typeface and not typeface.startswith("+"):
                                            font_name = typeface
                                        elif typeface == "+mj-lt":
                                            font_name = ts.theme.major_font
                                        elif typeface == "+mn-lt":
                                            font_name = ts.theme.minor_font
                                    if (
                                        ph_idx == 0
                                        and not ts.title_font_family
                                        and font_name
                                    ):
                                        ts.title_font_family = font_name
                                    if (
                                        ph_idx == 1
                                        and not ts.body_font_family
                                        and font_name
                                    ):
                                        ts.body_font_family = font_name
        except Exception as e:
            if VERBOSE:
                print(
                    "[VERBOSE] Layout placeholder style extraction failed: %s" % str(e)
                )

    if VERBOSE:
        print("[VERBOSE] Extracted template styles:")
        print("[VERBOSE]   Theme accent colors: %s" % ts.theme.accent_colors)
        print(
            "[VERBOSE]   Theme fonts: major=%s minor=%s"
            % (ts.theme.major_font, ts.theme.minor_font)
        )
        print("[VERBOSE]   Reference tables found: %d" % len(ts.table_styles))
        print("[VERBOSE]   Reference charts found: %d" % len(ts.chart_styles))
        print("[VERBOSE]   Title font family: %s" % ts.title_font_family)
        print("[VERBOSE]   Body font family: %s" % ts.body_font_family)

    return ts


# ---------------------------------------------------------------------------
# Template Style Application Functions
# ---------------------------------------------------------------------------


def _apply_table_style(table_shape, template_style: TemplateStyle):
    """Apply template-derived styling to a newly created table."""
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    table = table_shape.table
    tbl = table._tbl
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

    # Get the best matching table style from the template
    ref_style = template_style.table_styles[0] if template_style.table_styles else None
    theme = template_style.theme

    # Apply tblPr from reference if available (includes tblStyle ID + banding flags)
    if ref_style and ref_style.table_properties_xml is not None:
        existing_tblPr = tbl.find(qn("a:tblPr"))
        if existing_tblPr is not None:
            tbl.remove(existing_tblPr)
        new_tblPr = copy.deepcopy(ref_style.table_properties_xml)
        tbl.insert(0, new_tblPr)

    # Determine font settings (reference table > theme > defaults)
    header_font_family = (
        ref_style.header_font_family
        if ref_style and ref_style.header_font_family
        else theme.minor_font or "Calibri"
    )
    cell_font_family = (
        ref_style.cell_font_family
        if ref_style and ref_style.cell_font_family
        else theme.minor_font or "Calibri"
    )
    header_font_size = Pt(
        ref_style.header_font_size_pt
        if ref_style and ref_style.header_font_size_pt
        else 11
    )
    cell_font_size = Pt(
        ref_style.cell_font_size_pt if ref_style and ref_style.cell_font_size_pt else 10
    )
    header_bold = ref_style.header_font_bold if ref_style else True

    # Color resolution helper
    def _resolve_color(rgb_hex, fallback=""):
        if rgb_hex:
            try:
                return RGBColor.from_string(rgb_hex)
            except Exception:
                pass
        if fallback:
            try:
                return RGBColor.from_string(fallback)
            except Exception:
                pass
        return None

    # Apply styling to each cell
    for r_idx, row in enumerate(table.rows):
        is_header = r_idx == 0
        for cell in row.cells:
            # Font formatting
            for para in cell.text_frame.paragraphs:
                para.font.size = header_font_size if is_header else cell_font_size
                para.font.bold = header_bold if is_header else False
                para.font.name = header_font_family if is_header else cell_font_family

                # Font color
                if is_header and ref_style and ref_style.header_font_color_rgb:
                    color = _resolve_color(ref_style.header_font_color_rgb)
                    if color:
                        para.font.color.rgb = color
                elif not is_header and ref_style and ref_style.cell_font_color_rgb:
                    color = _resolve_color(ref_style.cell_font_color_rgb)
                    if color:
                        para.font.color.rgb = color
                elif theme.dk1:
                    color = _resolve_color(theme.dk1)
                    if color:
                        para.font.color.rgb = color

            cell.text_frame.word_wrap = True

            # Cell fill (via XML for reliability)
            tc_elem = cell._tc
            tcPr = tc_elem.find(qn("a:tcPr"))
            if tcPr is None:
                tcPr = etree.SubElement(tc_elem, qn("a:tcPr"))

            if is_header and ref_style and ref_style.header_fill_rgb:
                # Remove existing fill
                for old_fill in tcPr.findall(qn("a:solidFill")):
                    tcPr.remove(old_fill)
                fill_elem = etree.SubElement(tcPr, qn("a:solidFill"))
                srgb = etree.SubElement(fill_elem, qn("a:srgbClr"))
                srgb.set("val", ref_style.header_fill_rgb)
            elif is_header and theme.accent_colors:
                # Use first accent color for header
                for old_fill in tcPr.findall(qn("a:solidFill")):
                    tcPr.remove(old_fill)
                fill_elem = etree.SubElement(tcPr, qn("a:solidFill"))
                srgb = etree.SubElement(fill_elem, qn("a:srgbClr"))
                srgb.set("val", theme.accent_colors[0])

            if not is_header and ref_style and ref_style.cell_fill_rgb:
                for old_fill in tcPr.findall(qn("a:solidFill")):
                    tcPr.remove(old_fill)
                fill_elem = etree.SubElement(tcPr, qn("a:solidFill"))
                srgb = etree.SubElement(fill_elem, qn("a:srgbClr"))
                srgb.set("val", ref_style.cell_fill_rgb)

            # Borders
            if ref_style and ref_style.border_color_rgb:
                border_color = ref_style.border_color_rgb
                border_w = str(ref_style.border_width_emu)
                for border_name in ["lnL", "lnR", "lnT", "lnB"]:
                    ln = tcPr.find(qn("a:" + border_name))
                    if ln is None:
                        ln = etree.SubElement(tcPr, qn("a:" + border_name))
                    ln.set("w", border_w)
                    for old_fill in ln.findall(qn("a:solidFill")):
                        ln.remove(old_fill)
                    bf = etree.SubElement(ln, qn("a:solidFill"))
                    srgb = etree.SubElement(bf, qn("a:srgbClr"))
                    srgb.set("val", border_color)


def _apply_chart_style(chart_shape, template_style: TemplateStyle):
    """Apply template-derived styling to a newly created chart."""
    from pptx.oxml.ns import qn

    chart = chart_shape.chart
    theme = template_style.theme

    # Find the best matching chart style (prefer same chart type)
    ref_style = None
    for cs in template_style.chart_styles:
        try:
            if cs.chart_type == chart.chart_type:
                ref_style = cs
                break
        except Exception:
            pass
    if ref_style is None and template_style.chart_styles:
        ref_style = template_style.chart_styles[0]

    chart_space = chart._chartSpace

    # Apply series colors
    colors_to_use = []
    if ref_style and ref_style.series_fill_colors:
        colors_to_use = ref_style.series_fill_colors
    elif theme.accent_colors:
        colors_to_use = theme.accent_colors

    if colors_to_use:
        series_elems = chart_space.findall(".//" + qn("c:ser"))
        for i, ser in enumerate(series_elems):
            color = colors_to_use[i % len(colors_to_use)]
            if color.startswith("scheme:"):
                continue  # Skip scheme colors for now

            spPr = ser.find(qn("c:spPr"))
            if spPr is None:
                spPr = etree.SubElement(ser, qn("c:spPr"))

            # Set fill color
            for old in spPr.findall(qn("a:solidFill")):
                spPr.remove(old)
            fill = etree.SubElement(spPr, qn("a:solidFill"))
            srgb = etree.SubElement(fill, qn("a:srgbClr"))
            srgb.set("val", color)

    # Apply axis styling
    font_family = (
        ref_style.axis_font_family
        if ref_style and ref_style.axis_font_family
        else theme.minor_font or "Calibri"
    )
    font_size = str(
        (
            ref_style.axis_font_size_pt
            if ref_style and ref_style.axis_font_size_pt
            else 10
        )
        * 100
    )
    font_color = ref_style.axis_font_color_rgb if ref_style else (theme.dk1 or "000000")

    for axis_tag in [qn("c:catAx"), qn("c:valAx"), qn("c:dateAx")]:
        for ax in chart_space.findall(".//" + axis_tag):
            # Create or update txPr (text properties)
            txPr = ax.find(qn("c:txPr"))
            if txPr is None:
                txPr = etree.SubElement(ax, qn("c:txPr"))
            # Clear existing
            for child in list(txPr):
                txPr.remove(child)
            # Build text properties
            bodyPr = etree.SubElement(txPr, qn("a:bodyPr"))
            lstStyle = etree.SubElement(txPr, qn("a:lstStyle"))
            p = etree.SubElement(txPr, qn("a:p"))
            pPr = etree.SubElement(p, qn("a:pPr"))
            defRPr = etree.SubElement(pPr, qn("a:defRPr"))
            defRPr.set("sz", font_size)
            if font_color:
                sf = etree.SubElement(defRPr, qn("a:solidFill"))
                srgb = etree.SubElement(sf, qn("a:srgbClr"))
                srgb.set("val", font_color)
            latin = etree.SubElement(defRPr, qn("a:latin"))
            latin.set("typeface", font_family)
            endParaRPr = etree.SubElement(p, qn("a:endParaRPr"))
            endParaRPr.set("lang", "en-US")

            # Axis line color
            if ref_style and ref_style.axis_line_color_rgb:
                spPr = ax.find(qn("c:spPr"))
                if spPr is None:
                    spPr = etree.SubElement(ax, qn("c:spPr"))
                ln = spPr.find(qn("a:ln"))
                if ln is None:
                    ln = etree.SubElement(spPr, qn("a:ln"))
                for old in ln.findall(qn("a:solidFill")):
                    ln.remove(old)
                sf = etree.SubElement(ln, qn("a:solidFill"))
                srgb = etree.SubElement(sf, qn("a:srgbClr"))
                srgb.set("val", ref_style.axis_line_color_rgb)

    # Apply legend styling
    legend = chart_space.find(".//" + qn("c:legend"))
    if legend is not None:
        leg_font_family = (
            ref_style.legend_font_family
            if ref_style and ref_style.legend_font_family
            else font_family
        )
        leg_font_size = str(
            (
                ref_style.legend_font_size_pt
                if ref_style and ref_style.legend_font_size_pt
                else 10
            )
            * 100
        )
        leg_font_color = ref_style.legend_font_color_rgb if ref_style else font_color

        txPr = legend.find(qn("c:txPr"))
        if txPr is None:
            txPr = etree.SubElement(legend, qn("c:txPr"))
        for child in list(txPr):
            txPr.remove(child)
        bodyPr = etree.SubElement(txPr, qn("a:bodyPr"))
        lstStyle = etree.SubElement(txPr, qn("a:lstStyle"))
        p = etree.SubElement(txPr, qn("a:p"))
        pPr = etree.SubElement(p, qn("a:pPr"))
        defRPr = etree.SubElement(pPr, qn("a:defRPr"))
        defRPr.set("sz", leg_font_size)
        if leg_font_color:
            sf = etree.SubElement(defRPr, qn("a:solidFill"))
            srgb = etree.SubElement(sf, qn("a:srgbClr"))
            srgb.set("val", leg_font_color)
        latin = etree.SubElement(defRPr, qn("a:latin"))
        latin.set("typeface", leg_font_family)
        endParaRPr = etree.SubElement(p, qn("a:endParaRPr"))
        endParaRPr.set("lang", "en-US")

    # Plot area fill
    if ref_style and ref_style.plot_area_fill_rgb:
        plotArea = chart_space.find(".//" + qn("c:plotArea"))
        if plotArea is not None:
            spPr = plotArea.find(qn("c:spPr"))
            if spPr is None:
                spPr = etree.SubElement(plotArea, qn("c:spPr"))
            for old in spPr.findall(qn("a:solidFill")):
                spPr.remove(old)
            fill = etree.SubElement(spPr, qn("a:solidFill"))
            srgb = etree.SubElement(fill, qn("a:srgbClr"))
            srgb.set("val", ref_style.plot_area_fill_rgb)


# ---------------------------------------------------------------------------
# Content Extraction Functions (from agent_with_powerpoint_template.py)
# ---------------------------------------------------------------------------


def _extract_slide_content(slide) -> SlideContent:
    """Extract all content from a slide including text, tables, images, charts, and shapes.

    Detects picture placeholders using multiple strategies for robustness:
      1. int(type) == 18  (the raw OOXML value for PICTURE placeholders)
      2. str(type) contains 'PICTURE (18)'
      3. PP_PLACEHOLDER.PICTURE enum comparison
      4. XML-level fallback: <p:ph type="pic"/>
    """
    content = SlideContent()

    # Detect picture placeholders in the slide layout
    for shape in slide.placeholders:
        ph_fmt = shape.placeholder_format
        if ph_fmt is not None:
            ph_type_val = ph_fmt.type
            # Strategy 1 & 2: enum / int / string comparison
            is_picture_ph = False
            try:
                if ph_type_val is not None and (
                    int(ph_type_val) == 18 or str(ph_type_val) == "PICTURE (18)"
                ):
                    is_picture_ph = True
            except (ValueError, TypeError):
                pass
            # Also try the enum directly
            if not is_picture_ph:
                try:
                    if ph_type_val == PP_PLACEHOLDER.PICTURE:
                        is_picture_ph = True
                except Exception as e:
                    if VERBOSE:
                        print("[VERBOSE] Exception suppressed: %s" % str(e))
            # Strategy 3: XML-level fallback
            if not is_picture_ph:
                nsmap = {
                    "p": "http://schemas.openxmlformats.org/presentationml/2006/main"
                }
                ph_elem = shape._element.find(".//p:ph", nsmap)
                if ph_elem is not None and ph_elem.get("type") == "pic":
                    is_picture_ph = True
            if is_picture_ph:
                content.has_image_placeholder = True
                content.image_placeholder_indices.append(ph_fmt.idx)

    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                rows_data.append(row_data)
            content.tables.append(
                TableData(
                    rows=rows_data,
                    left=shape.left,
                    top=shape.top,
                    width=shape.width,
                    height=shape.height,
                )
            )
            continue

        if shape.has_chart:
            try:
                chart = shape.chart
                chart_type_val = chart.chart_type
                categories = []
                series_data = []
                plot = chart.plots[0] if chart.plots else None
                if plot:
                    if plot.categories:
                        categories = list(plot.categories)
                    for series in plot.series:
                        name = series.name if hasattr(series, "name") else ""
                        values = (
                            list(series.values) if hasattr(series, "values") else []
                        )
                        series_data.append((name or "", values))
                if categories or series_data:
                    content.charts.append(
                        ChartExtract(
                            chart_type=chart_type_val,
                            categories=categories,
                            series=series_data,
                            left=shape.left,
                            top=shape.top,
                            width=shape.width,
                            height=shape.height,
                        )
                    )
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                ct = shape.image.content_type
                content.images.append(
                    ImageData(
                        blob=blob,
                        left=shape.left,
                        top=shape.top,
                        width=shape.width,
                        height=shape.height,
                        content_type=ct,
                    )
                )
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_xml = copy.deepcopy(shape._element)
            content.shapes_xml.append(group_xml)
            try:
                for grp_shape in shape.shapes:
                    if grp_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            blob = grp_shape.image.blob
                            ct = grp_shape.image.content_type
                            content.images.append(
                                ImageData(
                                    blob=blob,
                                    left=grp_shape.left,
                                    top=grp_shape.top,
                                    width=grp_shape.width,
                                    height=grp_shape.height,
                                    content_type=ct,
                                )
                            )
                        except Exception as e:
                            if VERBOSE:
                                print("[VERBOSE] Exception suppressed: %s" % str(e))
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))
            continue

        if shape.has_text_frame:
            text_frame = shape.text_frame
            if shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx
                if ph_idx == 0:
                    content.title = text_frame.text.strip()
                elif ph_idx == 1:
                    paragraphs = text_frame.paragraphs
                    if len(paragraphs) == 1 and paragraphs[0].level == 0:
                        content.subtitle = paragraphs[0].text.strip()
                    else:
                        for para in paragraphs:
                            text = para.text.strip()
                            if text:
                                content.body_paragraphs.append((text, para.level))
                else:
                    for para in text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            content.body_paragraphs.append((text, para.level))
            else:
                # Preserve non-placeholder text boxes as shapes to keep styling.
                text_val = text_frame.text.strip()
                if text_val:
                    content.has_text_box = True
                    for para in text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            content.text_box_paragraphs.append((t, para.level))
                shape_xml = copy.deepcopy(shape._element)
                content.text_shapes_xml.append(shape_xml)
        elif not shape.is_placeholder:
            shape_xml = copy.deepcopy(shape._element)
            content.shapes_xml.append(shape_xml)

    return content


# ---------------------------------------------------------------------------
# Layout Collision Resolution Helpers
# ---------------------------------------------------------------------------


def _is_picture_placeholder(shape) -> bool:
    """Check if a shape is a picture placeholder using multiple strategies."""
    ph_fmt = shape.placeholder_format
    if ph_fmt is None:
        return False
    ph_type_val = ph_fmt.type
    # Strategy 1: int/string comparison
    try:
        if ph_type_val is not None and (
            int(ph_type_val) == 18 or str(ph_type_val) == "PICTURE (18)"
        ):
            return True
    except (ValueError, TypeError):
        pass
    # Strategy 2: PP_PLACEHOLDER enum
    try:
        if ph_type_val == PP_PLACEHOLDER.PICTURE:
            return True
    except Exception:
        pass
    # Strategy 3: XML-level fallback
    try:
        nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        ph_elem = shape._element.find(".//p:ph", nsmap)
        if ph_elem is not None and ph_elem.get("type") in ("pic", "clipArt"):
            return True
    except Exception:
        pass
    return False


def _placeholder_area(ph) -> int:
    """Return placeholder area in EMU^2."""
    try:
        return int(ph.width) * int(ph.height)
    except Exception:
        return 0


def _is_text_placeholder_type(ph_type) -> bool:
    """Return True if placeholder type is intended for text content."""
    return ph_type in (
        PP_PLACEHOLDER.BODY,
        PP_PLACEHOLDER.OBJECT,
        PP_PLACEHOLDER.SUBTITLE,
    )


def _is_visual_placeholder_type(ph_type) -> bool:
    """Return True if placeholder type is intended for visual content."""
    return ph_type in (
        PP_PLACEHOLDER.PICTURE,
        PP_PLACEHOLDER.CHART,
        PP_PLACEHOLDER.TABLE,
    )


def _largest_placeholder(layout, types=None, min_area: int = 0):
    """Return the largest placeholder matching types and min_area, or None."""
    best = None
    best_area = 0
    for ph in layout.placeholders:
        try:
            ph_type = ph.placeholder_format.type
        except Exception:
            continue
        if types is not None and ph_type not in types:
            continue
        area = _placeholder_area(ph)
        if area < min_area:
            continue
        if area > best_area:
            best = ph
            best_area = area
    return best


def _layout_placeholder_summary(layout) -> dict:
    """Return counts of placeholder types for logging."""
    summary = {
        "title": 0,
        "subtitle": 0,
        "body": 0,
        "object": 0,
        "picture": 0,
        "chart": 0,
        "table": 0,
        "other": 0,
    }
    for ph in layout.placeholders:
        try:
            ph_type = ph.placeholder_format.type
        except Exception:
            summary["other"] += 1
            continue
        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            summary["title"] += 1
        elif ph_type == PP_PLACEHOLDER.SUBTITLE:
            summary["subtitle"] += 1
        elif ph_type == PP_PLACEHOLDER.BODY:
            summary["body"] += 1
        elif ph_type == PP_PLACEHOLDER.OBJECT:
            summary["object"] += 1
        elif ph_type == PP_PLACEHOLDER.PICTURE:
            summary["picture"] += 1
        elif ph_type == PP_PLACEHOLDER.CHART:
            summary["chart"] += 1
        elif ph_type == PP_PLACEHOLDER.TABLE:
            summary["table"] += 1
        else:
            summary["other"] += 1
    return summary


def _layout_richness_score(layout) -> int:
    """Heuristic for visually rich layouts (cards, shapes, picture slots)."""
    score = 0
    # Prefer layouts with multiple body placeholders (card-style layouts)
    body_count = 0
    for ph in layout.placeholders:
        try:
            if ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
                body_count += 1
        except Exception:
            continue
    if body_count >= 2:
        score += 10
    if body_count >= 3:
        score += 10

    # Prefer layouts with at least one picture placeholder (visual interest)
    has_pic = False
    for ph in layout.placeholders:
        if _is_picture_placeholder(ph):
            has_pic = True
            break
    if has_pic:
        score += 15

    # Prefer layouts with non-placeholder shapes (cards, colored blocks)
    non_placeholder_shapes = 0
    for shape in layout.shapes:
        try:
            if shape.is_placeholder:
                continue
            # Ignore lines
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                continue
            non_placeholder_shapes += 1
        except Exception:
            continue
    if non_placeholder_shapes >= 2:
        score += 15
    elif non_placeholder_shapes == 1:
        score += 5

    return score


def _rect_overlap(a: ContentArea, b: ContentArea) -> int:
    """Return overlap area between two rectangles (or 0)."""
    x1 = max(a.left, b.left)
    y1 = max(a.top, b.top)
    x2 = min(a.left + a.width, b.left + b.width)
    y2 = min(a.top + a.height, b.top + b.height)
    if x2 <= x1 or y2 <= y1:
        return 0
    return (x2 - x1) * (y2 - y1)


def _layout_non_placeholder_text_bounds(layout) -> list[ContentArea]:
    """Collect bounds of non-placeholder text boxes defined in the layout."""
    bounds = []
    for shape in layout.shapes:
        try:
            if shape.is_placeholder:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = getattr(shape, "text", "")
            if text and text.strip():
                bounds.append(
                    ContentArea(
                        left=shape.left,
                        top=shape.top,
                        width=shape.width,
                        height=shape.height,
                    )
                )
        except Exception:
            continue
    return bounds


def _best_visual_placeholder(
    layout,
    desired_types: set,
    min_area: int,
    reserved: list[ContentArea],
    allow_body_fallback: bool = True,
):
    """Pick a visual placeholder that minimizes overlap with reserved text areas."""
    candidates = []
    for ph in layout.placeholders:
        try:
            ph_type = ph.placeholder_format.type
        except Exception:
            continue
        if ph_type in desired_types:
            candidates.append(ph)
    if allow_body_fallback and not candidates:
        for ph in layout.placeholders:
            try:
                ph_type = ph.placeholder_format.type
            except Exception:
                continue
            if ph_type in (
                PP_PLACEHOLDER.BODY,
                PP_PLACEHOLDER.OBJECT,
                PP_PLACEHOLDER.SUBTITLE,
            ):
                candidates.append(ph)

    best = None
    best_score = None
    for ph in candidates:
        area = _placeholder_area(ph)
        if area < min_area and desired_types != {PP_PLACEHOLDER.PICTURE}:
            # allow small picture placeholders but penalize small chart/table
            pass
        ph_bounds = ContentArea(
            left=ph.left, top=ph.top, width=ph.width, height=ph.height
        )
        overlap = 0
        for r in reserved:
            overlap += _rect_overlap(ph_bounds, r)
        # Prefer zero overlap, then minimal overlap, then largest area
        score = (overlap == 0, -overlap, area)
        if best_score is None or score > best_score:
            best_score = score
            best = ph
    return best


def _best_text_placeholder(layout, avoid: ContentArea | None = None):
    """Pick a text placeholder with minimal overlap to avoid region."""
    candidates = []
    for ph in layout.placeholders:
        try:
            ph_type = ph.placeholder_format.type
        except Exception:
            continue
        if ph_type in (
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.SUBTITLE,
        ):
            candidates.append(ph)
    if not candidates:
        return None
    best = None
    best_score = None
    for ph in candidates:
        area = _placeholder_area(ph)
        overlap = 0
        if avoid is not None:
            ph_bounds = ContentArea(
                left=ph.left, top=ph.top, width=ph.width, height=ph.height
            )
            overlap = _rect_overlap(ph_bounds, avoid)
        score = (overlap == 0, -overlap, area)
        if best_score is None or score > best_score:
            best_score = score
            best = ph
    return best


def _classify_content_mix(
    content: SlideContent, has_generated_image: bool = False
) -> ContentMix:
    """Classify the mix of element types in a slide."""
    has_text = bool(
        content.title
        or content.body_paragraphs
        or content.subtitle
        or content.has_text_box
    )
    has_table = len(content.tables) > 0
    has_chart = len(content.charts) > 0
    has_image = len(content.images) > 0

    visual_count = sum([has_table, has_chart, has_image, has_generated_image])

    if visual_count == 0:
        return ContentMix.TEXT_ONLY
    # If visuals exist but no text, still treat generated/user images as paired with text
    # so we reserve space for titles/subtitles in the template.
    if not has_text and has_generated_image:
        return ContentMix.TEXT_AND_GENERATED_IMAGE
    if not has_text and has_image:
        return ContentMix.TEXT_AND_IMAGE
    if not has_text and visual_count > 0:
        return ContentMix.VISUAL_ONLY
    if visual_count > 1:
        return ContentMix.MIXED
    if has_table:
        return ContentMix.TEXT_AND_TABLE
    if has_chart:
        return ContentMix.TEXT_AND_CHART
    if has_generated_image:
        return ContentMix.TEXT_AND_GENERATED_IMAGE
    if has_image:
        return ContentMix.TEXT_AND_IMAGE
    return ContentMix.TEXT_ONLY


def _compute_text_ratio(content: SlideContent, content_mix: ContentMix) -> float:
    """Compute how much of the content area should be allocated to text.

    Returns a ratio (0.0-1.0) representing the text portion.
    """
    num_paragraphs = len(content.body_paragraphs)

    if content_mix in (
        ContentMix.TEXT_AND_TABLE,
        ContentMix.TEXT_AND_CHART,
        ContentMix.MIXED,
    ):
        # Top/bottom split: give text 25-45% based on paragraph count
        if num_paragraphs <= 2:
            return 0.25
        elif num_paragraphs <= 4:
            return 0.35
        else:
            return 0.45

    if content_mix in (ContentMix.TEXT_AND_IMAGE, ContentMix.TEXT_AND_GENERATED_IMAGE):
        # Left/right split: give text 40-55% based on paragraph count
        if num_paragraphs <= 2:
            return 0.40
        elif num_paragraphs <= 4:
            return 0.45
        else:
            return 0.55

    return 0.50  # Default even split


def _compute_max_font_size(
    region: ContentArea, num_paragraphs: int, is_title: bool = False
) -> int:
    """Compute the maximum font size that fits text in the given region.

    Returns font size in points, clamped between 10-28pt for titles and 10-18pt for body.
    """
    if num_paragraphs == 0:
        return 28 if is_title else 18

    # Rough estimate: each point of font size ~ 1 * EMU_PER_PT in EMU height per line
    # Plus ~50% for line spacing
    EMU_PER_PT = 12700  # 1 point = 12700 EMU
    LINE_SPACING_FACTOR = 1.5

    available_height = region.height
    lines_needed = max(num_paragraphs, 1)

    # font_size_pt <= available_height / (lines_needed * EMU_PER_PT * LINE_SPACING_FACTOR)
    max_size_pt = int(
        available_height / (lines_needed * EMU_PER_PT * LINE_SPACING_FACTOR)
    )

    if is_title:
        return max(10, min(28, max_size_pt))
    else:
        return max(10, min(18, max_size_pt))


def _compute_region_map(
    layout,
    content_mix: ContentMix,
    slide_width: int,
    slide_height: int,
    content: SlideContent,
) -> RegionMap:
    """Compute separate text and visual regions based on content mix and layout."""
    slide_area = int(slide_width) * int(slide_height)
    min_area = int(slide_area * 0.10)
    content_area = _get_content_area(layout, slide_width, slide_height)
    layout_text_bounds = _layout_non_placeholder_text_bounds(layout)

    # For text-only or visual-only slides, use the full content area
    if content_mix in (ContentMix.TEXT_ONLY, ContentMix.VISUAL_ONLY):
        return RegionMap(
            text_region=content_area,
            visual_region=content_area,
            layout_type="full",
        )

    # Check if the layout natively has separate content placeholders
    # (e.g., "Two Content", "Content + Picture")
    content_placeholders = []
    picture_placeholders = []
    chart_placeholders = []
    table_placeholders = []
    for ph in layout.placeholders:
        if _is_picture_placeholder(ph):
            picture_placeholders.append(ph)
        elif ph.placeholder_format.type == PP_PLACEHOLDER.CHART:
            chart_placeholders.append(ph)
        elif ph.placeholder_format.type == PP_PLACEHOLDER.TABLE:
            table_placeholders.append(ph)
        elif ph.placeholder_format.idx > 0 and ph.has_text_frame:
            content_placeholders.append(ph)

    if content_mix in (ContentMix.TEXT_AND_IMAGE, ContentMix.TEXT_AND_GENERATED_IMAGE):
        # If layout has a picture placeholder + content placeholder, use native regions
        if picture_placeholders and content_placeholders:
            text_ph = content_placeholders[0]
            pic_ph = picture_placeholders[0]
            return RegionMap(
                text_region=ContentArea(
                    left=text_ph.left,
                    top=text_ph.top,
                    width=text_ph.width,
                    height=text_ph.height,
                ),
                visual_region=ContentArea(
                    left=pic_ph.left,
                    top=pic_ph.top,
                    width=pic_ph.width,
                    height=pic_ph.height,
                ),
                layout_type="native",
            )
        text_ph = _largest_placeholder(
            layout,
            {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE},
            min_area=0,
        )
        reserved = list(layout_text_bounds)
        if text_ph is not None:
            reserved.append(
                ContentArea(
                    left=text_ph.left,
                    top=text_ph.top,
                    width=text_ph.width,
                    height=text_ph.height,
                )
            )
        pic_ph = _best_visual_placeholder(
            layout,
            {PP_PLACEHOLDER.PICTURE},
            min_area=min_area,
            reserved=reserved,
            allow_body_fallback=True,
        )
        if pic_ph is not None:
            text_region = (
                ContentArea(
                    left=text_ph.left,
                    top=text_ph.top,
                    width=text_ph.width,
                    height=text_ph.height,
                )
                if text_ph is not None
                else _get_content_area(
                    layout,
                    slide_width,
                    slide_height,
                    preferred_types={
                        PP_PLACEHOLDER.BODY,
                        PP_PLACEHOLDER.OBJECT,
                        PP_PLACEHOLDER.SUBTITLE,
                    },
                )
            )
            return RegionMap(
                text_region=text_region,
                visual_region=ContentArea(
                    left=pic_ph.left,
                    top=pic_ph.top,
                    width=pic_ph.width,
                    height=pic_ph.height,
                ),
                layout_type="native",
            )

    # If content is chart/table/image and layout has a matching placeholder,
    # use that placeholder for the visual region and derive text elsewhere.
    if content_mix == ContentMix.TEXT_AND_CHART:
        if chart_placeholders and content_placeholders:
            chart_ph = chart_placeholders[0]
            text_ph = (
                _best_text_placeholder(
                    layout,
                    avoid=ContentArea(
                        left=chart_ph.left,
                        top=chart_ph.top,
                        width=chart_ph.width,
                        height=chart_ph.height,
                    ),
                )
                or content_placeholders[0]
            )
            return RegionMap(
                text_region=ContentArea(
                    left=text_ph.left,
                    top=text_ph.top,
                    width=text_ph.width,
                    height=text_ph.height,
                ),
                visual_region=ContentArea(
                    left=chart_ph.left,
                    top=chart_ph.top,
                    width=chart_ph.width,
                    height=chart_ph.height,
                ),
                layout_type="native",
            )
        text_ph = _largest_placeholder(
            layout,
            {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE},
            min_area=0,
        )
        reserved = list(layout_text_bounds)
        if text_ph is not None:
            reserved.append(
                ContentArea(
                    left=text_ph.left,
                    top=text_ph.top,
                    width=text_ph.width,
                    height=text_ph.height,
                )
            )
        chart_ph = _best_visual_placeholder(
            layout,
            {PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.TABLE},
            min_area=min_area,
            reserved=reserved,
            allow_body_fallback=True,
        )
        if chart_ph is not None:
            best_text = (
                _best_text_placeholder(
                    layout,
                    avoid=ContentArea(
                        left=chart_ph.left,
                        top=chart_ph.top,
                        width=chart_ph.width,
                        height=chart_ph.height,
                    ),
                )
                or text_ph
            )
            text_region = (
                ContentArea(
                    left=best_text.left,
                    top=best_text.top,
                    width=best_text.width,
                    height=best_text.height,
                )
                if best_text is not None
                else _get_content_area(
                    layout,
                    slide_width,
                    slide_height,
                    preferred_types={
                        PP_PLACEHOLDER.BODY,
                        PP_PLACEHOLDER.OBJECT,
                        PP_PLACEHOLDER.SUBTITLE,
                    },
                )
            )
            return RegionMap(
                text_region=text_region,
                visual_region=ContentArea(
                    left=chart_ph.left,
                    top=chart_ph.top,
                    width=chart_ph.width,
                    height=chart_ph.height,
                ),
                layout_type="native",
            )

    if content_mix == ContentMix.TEXT_AND_TABLE:
        if table_placeholders and content_placeholders:
            table_ph = table_placeholders[0]
            text_ph = (
                _best_text_placeholder(
                    layout,
                    avoid=ContentArea(
                        left=table_ph.left,
                        top=table_ph.top,
                        width=table_ph.width,
                        height=table_ph.height,
                    ),
                )
                or content_placeholders[0]
            )
            return RegionMap(
                text_region=ContentArea(
                    left=text_ph.left,
                    top=text_ph.top,
                    width=text_ph.width,
                    height=text_ph.height,
                ),
                visual_region=ContentArea(
                    left=table_ph.left,
                    top=table_ph.top,
                    width=table_ph.width,
                    height=table_ph.height,
                ),
                layout_type="native",
            )
        text_ph = _largest_placeholder(
            layout,
            {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE},
            min_area=0,
        )
        reserved = list(layout_text_bounds)
        if text_ph is not None:
            reserved.append(
                ContentArea(
                    left=text_ph.left,
                    top=text_ph.top,
                    width=text_ph.width,
                    height=text_ph.height,
                )
            )
        table_ph = _best_visual_placeholder(
            layout,
            {PP_PLACEHOLDER.TABLE, PP_PLACEHOLDER.CHART},
            min_area=min_area,
            reserved=reserved,
            allow_body_fallback=True,
        )
        if table_ph is not None:
            best_text = (
                _best_text_placeholder(
                    layout,
                    avoid=ContentArea(
                        left=table_ph.left,
                        top=table_ph.top,
                        width=table_ph.width,
                        height=table_ph.height,
                    ),
                )
                or text_ph
            )
            text_region = (
                ContentArea(
                    left=best_text.left,
                    top=best_text.top,
                    width=best_text.width,
                    height=best_text.height,
                )
                if best_text is not None
                else _get_content_area(
                    layout,
                    slide_width,
                    slide_height,
                    preferred_types={
                        PP_PLACEHOLDER.BODY,
                        PP_PLACEHOLDER.OBJECT,
                        PP_PLACEHOLDER.SUBTITLE,
                    },
                )
            )
            return RegionMap(
                text_region=text_region,
                visual_region=ContentArea(
                    left=table_ph.left,
                    top=table_ph.top,
                    width=table_ph.width,
                    height=table_ph.height,
                ),
                layout_type="native",
            )

    # Compute text ratio based on content volume
    text_ratio = _compute_text_ratio(content, content_mix)
    GAP = int(slide_width * 0.02)  # 2% gap between regions

    # Text + Table/Chart -> top/bottom split
    if content_mix in (
        ContentMix.TEXT_AND_TABLE,
        ContentMix.TEXT_AND_CHART,
        ContentMix.MIXED,
    ):
        text_height = int(content_area.height * text_ratio)
        visual_height = content_area.height - text_height - GAP
        return RegionMap(
            text_region=ContentArea(
                left=content_area.left,
                top=content_area.top,
                width=content_area.width,
                height=text_height,
            ),
            visual_region=ContentArea(
                left=content_area.left,
                top=content_area.top + text_height + GAP,
                width=content_area.width,
                height=visual_height,
            ),
            layout_type="split_vertical",
        )

    # Text + Image -> left/right split
    if content_mix in (ContentMix.TEXT_AND_IMAGE, ContentMix.TEXT_AND_GENERATED_IMAGE):
        text_width = int(content_area.width * text_ratio)
        visual_width = content_area.width - text_width - GAP
        return RegionMap(
            text_region=ContentArea(
                left=content_area.left,
                top=content_area.top,
                width=text_width,
                height=content_area.height,
            ),
            visual_region=ContentArea(
                left=content_area.left + text_width + GAP,
                top=content_area.top,
                width=visual_width,
                height=content_area.height,
            ),
            layout_type="split_horizontal",
        )

    # Fallback: use full area
    return RegionMap(
        text_region=content_area,
        visual_region=content_area,
        layout_type="full",
    )


def _ensure_text_on_top(slide) -> None:
    """Reorder shapes in the XML tree so text elements render above visual elements.

    This ensures text is never hidden behind charts, tables, or images.
    """
    spTree = slide.shapes._spTree

    # Separate shapes into visual (charts, tables, pictures) and text
    visual_elements = []
    text_elements = []
    other_elements = []

    for child in list(spTree):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("graphicFrame",):  # Tables, charts
            visual_elements.append(child)
        elif tag == "pic":  # Pictures
            visual_elements.append(child)
        elif tag == "sp":  # Shapes (including text boxes)
            # Check if it has a text frame with actual content
            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            txBody = child.find(".//{%s}txBody" % ns)
            if txBody is not None:
                # Check if it has actual text (not just empty)
                has_text = False
                for p_elem in txBody.findall("{%s}p" % ns):
                    for r_elem in p_elem.findall("{%s}r" % ns):
                        t_elem = r_elem.find("{%s}t" % ns)
                        if t_elem is not None and t_elem.text and t_elem.text.strip():
                            has_text = True
                            break
                    if has_text:
                        break
                if has_text:
                    text_elements.append(child)
                else:
                    other_elements.append(child)
            else:
                other_elements.append(child)
        else:
            other_elements.append(child)

    # Only reorder if we have both text and visual elements
    if not visual_elements or not text_elements:
        return

    # Remove all shape elements from spTree
    for elem in visual_elements + text_elements + other_elements:
        try:
            spTree.remove(elem)
        except ValueError:
            pass

    # Re-add in order: other (background shapes) -> visuals -> text (on top)
    for elem in other_elements:
        spTree.append(elem)
    for elem in visual_elements:
        spTree.append(elem)
    for elem in text_elements:
        spTree.append(elem)


# ---------------------------------------------------------------------------
# Template Application Functions (from agent_with_powerpoint_template.py)
# ---------------------------------------------------------------------------


def _find_best_layout(
    template_prs,
    slide_index: int,
    total_slides: int,
    content_mix: ContentMix = ContentMix.TEXT_ONLY,
    has_generated_image: bool = False,
):
    """Find the best matching layout from the template for a given slide position."""
    layouts = list(template_prs.slide_layouts)
    if not layouts:
        raise ValueError("Template has no slide layouts")

    layout_names = [(i, layout.name.lower()) for i, layout in enumerate(layouts)]
    is_title_slide = slide_index == 0
    is_last_slide = slide_index == total_slides - 1

    slide_width = template_prs.slide_width
    slide_height = template_prs.slide_height
    slide_area = int(slide_width) * int(slide_height)
    min_area = int(slide_area * 0.10)

    def _score_layout(layout):
        name = layout.name.lower()
        score = 0

        text_ph = _largest_placeholder(
            layout,
            types={PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE},
        )
        chart_ph = _largest_placeholder(layout, types={PP_PLACEHOLDER.CHART})
        table_ph = _largest_placeholder(layout, types={PP_PLACEHOLDER.TABLE})
        pic_ph = _largest_placeholder(layout, types={PP_PLACEHOLDER.PICTURE})

        if is_title_slide:
            has_title_ph = any(
                ph.placeholder_format.type
                in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
                for ph in layout.placeholders
            )
            if has_title_ph:
                score += 120
            if "title" in name:
                score += 20
            if text_ph:
                score += 10
            return score
        # Penalize title layouts for non-title slides
        has_title_ph = any(
            ph.placeholder_format.type
            in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
            for ph in layout.placeholders
        )
        if has_title_ph:
            score -= 80
        if "title" in name:
            score -= 40

        if is_last_slide:
            if "blank" in name or "closing" in name or "end" in name:
                score += 30

        if content_mix in (ContentMix.TEXT_AND_CHART,):
            if chart_ph:
                score += 120
                if _placeholder_area(chart_ph) >= min_area:
                    score += 30
            if text_ph:
                score += 30
        elif content_mix in (ContentMix.TEXT_AND_TABLE,):
            if table_ph:
                score += 120
                if _placeholder_area(table_ph) >= min_area:
                    score += 30
            if text_ph:
                score += 30
        elif content_mix in (
            ContentMix.TEXT_AND_IMAGE,
            ContentMix.TEXT_AND_GENERATED_IMAGE,
        ):
            if pic_ph:
                score += 90
                if _placeholder_area(pic_ph) >= min_area:
                    score += 20
            if text_ph:
                score += 30
        elif content_mix == ContentMix.MIXED:
            if chart_ph:
                score += 60
            if table_ph:
                score += 60
            if pic_ph:
                score += 40
            if text_ph:
                score += 30
        else:
            if text_ph:
                score += 60

        # Prefer visually rich layouts for text-only or light-content slides
        if content_mix in (ContentMix.TEXT_ONLY,):
            score += _layout_richness_score(layout)

        if "content" in name or "body" in name or "text" in name:
            score += 5
        if "two content" in name:
            score += 3

        # Penalize layouts with only tiny visual placeholders
        if content_mix in (
            ContentMix.TEXT_AND_CHART,
            ContentMix.TEXT_AND_TABLE,
            ContentMix.TEXT_AND_IMAGE,
            ContentMix.TEXT_AND_GENERATED_IMAGE,
            ContentMix.MIXED,
        ):
            best_visual = max(
                [
                    _placeholder_area(p)
                    for p in [chart_ph, table_ph, pic_ph]
                    if p is not None
                ]
                or [0]
            )
            if best_visual and best_visual < min_area:
                score -= 50

        return score

    best_layout = None
    best_score = None
    for layout in layouts:
        score = _score_layout(layout)
        if best_score is None or score > best_score:
            best_score = score
            best_layout = layout

    if best_layout is not None:
        return best_layout

    if len(layouts) > 1:
        return layouts[1]
    return layouts[0]


def _get_content_area(
    layout,
    slide_width: int,
    slide_height: int,
    preferred_types: set | None = None,
    min_area: int = 0,
) -> ContentArea:
    """Derive the safe content area from a template layout's placeholders.

    Strategy:
    0. If preferred_types is provided, choose the largest matching placeholder
       above min_area.
    1. Look for a body placeholder (idx=1) -- its position defines the content area.
    2. If no body placeholder, look for any placeholder with idx > 0.
    3. If no placeholders at all, use a default safe margin.

    Args:
        layout: A python-pptx SlideLayout object.
        slide_width: Presentation slide width in EMU.
        slide_height: Presentation slide height in EMU.

    Returns:
        ContentArea with the computed safe region.
    """
    # Try preferred placeholder types first
    if preferred_types:
        ph = _largest_placeholder(layout, preferred_types, min_area=min_area)
        if ph is not None:
            return ContentArea(
                left=ph.left, top=ph.top, width=ph.width, height=ph.height
            )

    # Prefer the largest text-oriented placeholder (body/object/subtitle)
    text_ph = _largest_placeholder(
        layout,
        types={
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.SUBTITLE,
            PP_PLACEHOLDER.TABLE,  # allow table placeholder as a usable content area
        },
        min_area=0,
    )
    if text_ph is not None:
        return ContentArea(
            left=text_ph.left,
            top=text_ph.top,
            width=text_ph.width,
            height=text_ph.height,
        )

    # Try body placeholder idx=1 explicitly if present
    for ph in layout.placeholders:
        if ph.placeholder_format.idx == 1:
            return ContentArea(
                left=ph.left, top=ph.top, width=ph.width, height=ph.height
            )

    # Otherwise pick the largest non-title placeholder (avoiding picture if possible)
    candidates = [
        ph
        for ph in layout.placeholders
        if ph.placeholder_format.idx > 0
        and ph.placeholder_format.type not in {PP_PLACEHOLDER.PICTURE}
    ]
    if candidates:
        ph = max(candidates, key=_placeholder_area)
        return ContentArea(left=ph.left, top=ph.top, width=ph.width, height=ph.height)

    # Default: safe margins (5% left, 25% top, 90% width, 65% height)
    return ContentArea(
        left=int(slide_width * 0.05),
        top=int(slide_height * 0.25),
        width=int(slide_width * 0.90),
        height=int(slide_height * 0.65),
    )


def _fit_to_area(img_width: int, img_height: int, area: ContentArea) -> tuple:
    """Scale dimensions to fit within an area while preserving aspect ratio.

    Args:
        img_width: Original width in EMU.
        img_height: Original height in EMU.
        area: Target ContentArea to fit within.

    Returns:
        Tuple of (left, top, width, height) in EMU, centered in the area.
    """
    if img_width <= 0 or img_height <= 0:
        return area.left, area.top, area.width, area.height

    aspect = img_width / img_height
    area_aspect = area.width / area.height

    if aspect > area_aspect:
        new_width = area.width
        new_height = int(new_width / aspect)
    else:
        new_height = area.height
        new_width = int(new_height * aspect)

    left = area.left + (area.width - new_width) // 2
    top = area.top + (area.height - new_height) // 2

    return left, top, new_width, new_height


def _add_picture_within_bounds(slide, image_bytes: bytes, bounds: ContentArea):
    """Insert image scaled to fit bounds while preserving aspect ratio."""
    try:
        from PIL import Image
    except ImportError:
        image_stream = BytesIO(image_bytes)
        slide.shapes.add_picture(
            image_stream, bounds.left, bounds.top, bounds.width, bounds.height
        )
        return

    try:
        with Image.open(BytesIO(image_bytes)) as im:
            img_w, img_h = im.size
        image_stream = BytesIO(image_bytes)
        left, top, width, height = _fit_to_area(
            int(img_w * 9525), int(img_h * 9525), bounds
        )
        slide.shapes.add_picture(image_stream, left, top, width, height)
    except Exception:
        image_stream = BytesIO(image_bytes)
        slide.shapes.add_picture(
            image_stream, bounds.left, bounds.top, bounds.width, bounds.height
        )


def _populate_placeholder_with_format(
    shape, texts, is_title=False, template_style: "TemplateStyle | None" = None
):
    """Populate a placeholder shape with text while preserving template formatting.

    Enables word wrap and attempts to auto-fit text to the placeholder bounds.
    Uses template_style for font family when available.
    """
    if not shape.has_text_frame:
        return

    tf = shape.text_frame

    # Enable word wrap before anything else
    tf.word_wrap = True

    ref_para = tf.paragraphs[0] if tf.paragraphs else None
    ref_para_xml = None
    ref_run_xml = None

    if ref_para is not None:
        pPr = ref_para._p.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
        )
        if pPr is not None:
            ref_para_xml = copy.deepcopy(pPr)
        if ref_para.runs:
            rPr = ref_para.runs[0]._r.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
            )
            if rPr is not None:
                ref_run_xml = copy.deepcopy(rPr)

    tf.clear()

    if is_title:
        para = tf.paragraphs[0]
        if ref_para_xml is not None:
            para._p.insert(0, copy.deepcopy(ref_para_xml))
        run = para.add_run()
        run.text = texts
        if ref_run_xml is not None:
            run._r.insert(0, copy.deepcopy(ref_run_xml))
    else:
        for i, (text, level) in enumerate(texts):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            if ref_para_xml is not None:
                new_pPr = copy.deepcopy(ref_para_xml)
                if level > 0:
                    new_pPr.set("lvl", str(level))
                para._p.insert(0, new_pPr)
            para.level = level
            run = para.add_run()
            run.text = text
            if ref_run_xml is not None:
                run._r.insert(0, copy.deepcopy(ref_run_xml))

    # Determine font family from template style or use default
    # Titles use major_font (heading font), body uses minor_font (body font)
    font_family = "Calibri"
    if template_style:
        if is_title:
            font_family = (
                template_style.title_font_family
                or template_style.theme.major_font
                or "Calibri"
            )
        else:
            font_family = (
                template_style.body_font_family
                or template_style.theme.minor_font
                or "Calibri"
            )

    # Auto-fit text to placeholder bounds
    try:
        # Compute a safe max size based on placeholder height and line count
        if is_title:
            line_count = 1
        else:
            line_count = len(texts) if isinstance(texts, list) else 1
        safe_max = _compute_max_font_size(
            ContentArea(
                left=shape.left, top=shape.top, width=shape.width, height=shape.height
            ),
            line_count,
            is_title=is_title,
        )
        hard_max = 28 if is_title else 18
        max_size = min(hard_max, safe_max)
        tf.fit_text(font_family=font_family, max_size=max_size)
    except Exception as e:
        # fit_text requires font metrics; fall back to MSO_AUTO_SIZE
        if VERBOSE:
            print("[VERBOSE] Exception suppressed: %s" % str(e))
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # After fit_text has set the font size, re-apply captured run properties
    # to restore template formatting (colors, font family, etc.) that fit_text may have overwritten
    if ref_run_xml is not None:
        try:
            ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            for para in tf.paragraphs:
                for run in para.runs:
                    # Remove the rPr that fit_text created
                    existing_rPr = run._r.find(ns_a + "rPr")
                    sz_val = None
                    if existing_rPr is not None:
                        # Preserve the font size from fit_text
                        sz_val = existing_rPr.get("sz")
                        run._r.remove(existing_rPr)

                    # Clone the template's rPr and insert it
                    new_rPr = copy.deepcopy(ref_run_xml)

                    # If fit_text set a size and the template rPr has a size, use the SMALLER one
                    # to ensure text fits. If template rPr has no size, use fit_text's size.
                    if sz_val:
                        template_sz = new_rPr.get("sz")
                        if template_sz:
                            # Use the smaller of template size and fit_text size
                            fit_sz = int(sz_val)
                            tmpl_sz = int(template_sz)
                            new_rPr.set("sz", str(min(fit_sz, tmpl_sz)))
                        else:
                            new_rPr.set("sz", sz_val)

                    run._r.insert(0, new_rPr)
        except Exception as e:
            if VERBOSE:
                print(
                    "[VERBOSE] Re-applying template run properties failed: %s" % str(e)
                )

    # When ref_run_xml is None (placeholder inherits from layout/master),
    # fit_text() will have set an explicit font that may not match the theme.
    # Ensure the correct theme font is applied.
    if ref_run_xml is None and template_style:
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        if is_title:
            target_font = (
                template_style.title_font_family or template_style.theme.major_font
            )
        else:
            target_font = (
                template_style.body_font_family or template_style.theme.minor_font
            )
        if target_font:
            for para in tf.paragraphs:
                for run in para.runs:
                    rPr = run._r.find(ns_a + "rPr")
                    if rPr is not None:
                        # Update the latin font element
                        latin = rPr.find(ns_a + "latin")
                        if latin is not None:
                            latin.set("typeface", target_font)
                        else:
                            latin = etree.SubElement(rPr, ns_a + "latin")
                            latin.set("typeface", target_font)


def _transfer_tables(
    slide,
    tables,
    content_area: ContentArea,
    template_style: "TemplateStyle | None" = None,
):
    """Transfer extracted table data to a slide, repositioned to the content area.

    When template_style is provided, applies template-derived table styling
    instead of hardcoded font sizes.
    """
    TABLE_CELL_FONT_SIZE = Pt(10)
    TABLE_HEADER_FONT_SIZE = Pt(11)
    num_tables = len(tables)

    for t_idx, td in enumerate(tables):
        num_rows = len(td.rows)
        num_cols = len(td.rows[0]) if td.rows else 0
        if num_rows == 0 or num_cols == 0:
            continue

        # Stack multiple tables vertically within the content area
        if num_tables > 1:
            per_table_height = content_area.height // num_tables
            table_top = content_area.top + (t_idx * per_table_height)
            table_height = per_table_height
        else:
            table_top = content_area.top
            table_height = content_area.height

        table_shape = slide.shapes.add_table(
            num_rows,
            num_cols,
            content_area.left,
            table_top,
            content_area.width,
            table_height,
        )
        table = table_shape.table
        for r_idx, row_data in enumerate(td.rows):
            for c_idx, cell_text in enumerate(row_data):
                if c_idx < num_cols:
                    cell = table.cell(r_idx, c_idx)
                    cell.text = cell_text

        # Apply template-derived styling or fall back to hardcoded sizes
        if template_style:
            try:
                _apply_table_style(table_shape, template_style)
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Table style application failed: %s" % str(e))
                # Fallback to hardcoded styling on error
                for r_idx2, row_data2 in enumerate(td.rows):
                    for c_idx2 in range(min(len(row_data2), num_cols)):
                        cell = table.cell(r_idx2, c_idx2)
                        for para in cell.text_frame.paragraphs:
                            para.font.size = (
                                TABLE_HEADER_FONT_SIZE
                                if r_idx2 == 0
                                else TABLE_CELL_FONT_SIZE
                            )
                        cell.text_frame.word_wrap = True
        else:
            # Original hardcoded styling as fallback
            for r_idx2, row_data2 in enumerate(td.rows):
                for c_idx2 in range(min(len(row_data2), num_cols)):
                    cell = table.cell(r_idx2, c_idx2)
                    for para in cell.text_frame.paragraphs:
                        para.font.size = (
                            TABLE_HEADER_FONT_SIZE
                            if r_idx2 == 0
                            else TABLE_CELL_FONT_SIZE
                        )
                    cell.text_frame.word_wrap = True


def _transfer_images(slide, images, content_area: ContentArea):
    """Transfer extracted images to a slide, scaled to fit the content area."""
    for img in images:
        image_stream = BytesIO(img.blob)
        left, top, width, height = _fit_to_area(img.width, img.height, content_area)
        slide.shapes.add_picture(image_stream, left, top, width, height)


def _transfer_charts(
    slide,
    charts,
    content_area: ContentArea,
    template_style: "TemplateStyle | None" = None,
):
    """Transfer extracted chart data to a slide, sized to fill the content area.

    When template_style is provided, applies template-derived chart styling
    (series colors, axis fonts, legend styling, etc.) after chart creation.
    """
    try:
        from pptx.chart.data import CategoryChartData
    except ImportError:
        return

    num_charts = len(charts)

    for c_idx, cd in enumerate(charts):
        try:
            chart_data = CategoryChartData()
            chart_data.categories = cd.categories
            for series_name, series_values in cd.series:
                clean_values = []
                for v in series_values:
                    if v is None:
                        clean_values.append(0)
                    elif isinstance(v, (int, float)):
                        clean_values.append(v)
                    else:
                        try:
                            clean_values.append(float(v))
                        except (ValueError, TypeError):
                            clean_values.append(0)
                chart_data.add_series(series_name, clean_values)

            # Stack multiple charts vertically within the content area
            if num_charts > 1:
                chart_height = content_area.height // num_charts
                chart_top = content_area.top + (c_idx * chart_height)
            else:
                chart_height = content_area.height
                chart_top = content_area.top

            chart_graphic_frame = slide.shapes.add_chart(
                cd.chart_type,
                content_area.left,
                chart_top,
                content_area.width,
                chart_height,
                chart_data,
            )

            # Apply template-derived chart styling
            if template_style:
                try:
                    _apply_chart_style(chart_graphic_frame, template_style)
                except Exception as e:
                    if VERBOSE:
                        print("[VERBOSE] Chart style application failed: %s" % str(e))
        except Exception as e:
            if VERBOSE:
                print("[VERBOSE] Exception suppressed: %s" % str(e))


def _transfer_shapes(slide, shapes_xml):
    """Transfer simple shapes by deep-copying their XML to the target slide."""
    spTree = slide.shapes._spTree
    for shape_elem in shapes_xml:
        cloned = copy.deepcopy(shape_elem)
        existing_ids = [
            int(sp.get("id", 0)) for sp in spTree.iter() if sp.get("id") is not None
        ]
        max_id = max(existing_ids) if existing_ids else 0
        for nv_elem in cloned.iter():
            if nv_elem.tag.endswith("}cNvPr"):
                max_id += 1
                nv_elem.set("id", str(max_id))
        spTree.append(cloned)


def _clear_unused_placeholders(slide, populated_indices: set) -> None:
    """Remove unused placeholder XML elements from slide to prevent ghost text.

    This is the only reliable way to eliminate ALL types of ghost text:
    - "Click to add title"
    - "Click to add text"
    - "Click icon to add picture"
    - Content placeholder insertion icons (table, chart, image, etc.)

    Simply calling tf.clear() is insufficient for picture placeholders and
    content placeholders with embedded icons.  Removing the XML element
    from the shape tree is the nuclear option that works for every case.

    Args:
        slide: The pptx slide object.
        populated_indices: Set of placeholder idx values that were filled with content.
    """
    spTree = slide.shapes._spTree
    elements_to_remove = []

    # Snapshot the collection with list() to avoid proxy/iterator issues
    for shape in list(slide.placeholders):
        ph_idx = shape.placeholder_format.idx
        if ph_idx in populated_indices:
            continue  # This placeholder was populated, keep it

        # Don't remove shapes that have actual content (charts, tables, real images)
        try:
            if hasattr(shape, "has_chart") and shape.has_chart:
                continue
            if hasattr(shape, "has_table") and shape.has_table:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception as e:
            if VERBOSE:
                print("[VERBOSE] Exception suppressed: %s" % str(e))

        # Mark for removal -- removing from XML is the only reliable cleanup
        elements_to_remove.append(shape._element)

    for element in elements_to_remove:
        spTree.remove(element)


def _remove_empty_textboxes(slide) -> None:
    """Remove non-placeholder text boxes with no visible text content."""
    spTree = slide.shapes._spTree
    to_remove = []
    for shape in list(slide.shapes):
        try:
            if shape.is_placeholder:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text
            if text is None or text.strip() == "":
                to_remove.append(shape._element)
        except Exception:
            continue
    for el in to_remove:
        try:
            spTree.remove(el)
        except Exception:
            pass


def _populate_slide(
    new_slide,
    content: SlideContent,
    slide_width: int,
    slide_height: int,
    generated_image_bytes: bytes | None = None,
    template_style: "TemplateStyle | None" = None,
):
    """Transfer all content into a new slide using template-aware positioning.

    Handles text placeholders, tables, charts, images, and shapes. When
    ``generated_image_bytes`` is provided and the slide layout contains a
    picture placeholder, the image is inserted into that placeholder rather
    than being added as a free-floating picture.

    Uses a RegionMap to separate text and visual content into non-overlapping
    regions, preventing layout collisions between text and tables/charts/images.

    Args:
        new_slide: The new slide created from a template layout.
        content: Extracted SlideContent from the generated slide.
        slide_width: Presentation slide width in EMU.
        slide_height: Presentation slide height in EMU.
        generated_image_bytes: Optional raw image bytes (e.g. from NanoBanana)
            to insert into picture placeholders.
        template_style: Optional extracted template styling for font/color matching.
    """
    # Classify content mix and compute region map
    content_mix = _classify_content_mix(
        content, has_generated_image=generated_image_bytes is not None
    )
    region_map = _compute_region_map(
        new_slide.slide_layout, content_mix, slide_width, slide_height, content
    )

    if VERBOSE:
        summary = _layout_placeholder_summary(new_slide.slide_layout)
        print(
            "[VERBOSE] Layout '%s' placeholders: %s"
            % (new_slide.slide_layout.name, summary)
        )
        print(
            "[VERBOSE] Region map: layout_type=%s text=(%d,%d,%d,%d) visual=(%d,%d,%d,%d)"
            % (
                region_map.layout_type,
                region_map.text_region.left,
                region_map.text_region.top,
                region_map.text_region.width,
                region_map.text_region.height,
                region_map.visual_region.left,
                region_map.visual_region.top,
                region_map.visual_region.width,
                region_map.visual_region.height,
            )
        )

    # Track populated placeholder indices for cleanup
    populated_indices: set[int] = set()

    title_placed = False
    body_placed = False
    body_paragraphs = list(content.body_paragraphs)
    if content.text_box_paragraphs and content_mix != ContentMix.TEXT_ONLY:
        body_paragraphs.extend(content.text_box_paragraphs)

    # Prefer placeholder matching the computed text region when available
    preferred_text_ph = None
    if region_map.layout_type == "native":
        for shape in new_slide.placeholders:
            try:
                if (
                    shape.has_text_frame
                    and not _is_visual_placeholder_type(shape.placeholder_format.type)
                    and shape.left == region_map.text_region.left
                    and shape.top == region_map.text_region.top
                    and shape.width == region_map.text_region.width
                    and shape.height == region_map.text_region.height
                ):
                    preferred_text_ph = shape
                    break
            except Exception:
                continue

    # Resize body placeholder to text_region when using split layout
    if region_map.layout_type not in ("full", "native"):
        for shape in new_slide.placeholders:
            ph_idx = shape.placeholder_format.idx
            if ph_idx == 1 and shape.has_text_frame:
                shape.left = region_map.text_region.left
                shape.top = region_map.text_region.top
                shape.width = region_map.text_region.width
                shape.height = region_map.text_region.height

    for shape in new_slide.placeholders:
        ph_idx = shape.placeholder_format.idx
        if ph_idx == 0 and content.title:
            _populate_placeholder_with_format(
                shape, content.title, is_title=True, template_style=template_style
            )
            populated_indices.add(ph_idx)
            title_placed = True
        elif ph_idx == 1:
            if body_paragraphs and (
                preferred_text_ph is None or shape == preferred_text_ph
            ):
                _populate_placeholder_with_format(
                    shape,
                    body_paragraphs,
                    is_title=False,
                    template_style=template_style,
                )
                populated_indices.add(ph_idx)
                body_placed = True
            elif content.subtitle:
                _populate_placeholder_with_format(
                    shape,
                    content.subtitle,
                    is_title=True,
                    template_style=template_style,
                )
                populated_indices.add(ph_idx)
                body_placed = True

    if preferred_text_ph is not None and not body_placed and body_paragraphs:
        _populate_placeholder_with_format(
            preferred_text_ph,
            body_paragraphs,
            is_title=False,
            template_style=template_style,
        )
        populated_indices.add(preferred_text_ph.placeholder_format.idx)
        body_placed = True

    if not body_placed and body_paragraphs:
        for shape in new_slide.placeholders:
            ph_idx = shape.placeholder_format.idx
            if (
                ph_idx > 1
                and shape.has_text_frame
                and not _is_visual_placeholder_type(shape.placeholder_format.type)
            ):
                _populate_placeholder_with_format(
                    shape,
                    body_paragraphs,
                    is_title=False,
                    template_style=template_style,
                )
                populated_indices.add(ph_idx)
                body_placed = True
                break

    # Fallback text boxes using text_region bounds
    if not title_placed and content.title:
        txBox = new_slide.shapes.add_textbox(
            region_map.text_region.left,
            Inches(0.3),
            region_map.text_region.width,
            Inches(1.0),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = content.title
        for para in tf.paragraphs:
            para.font.size = Pt(
                template_style.title_font_size_pt if template_style else 28
            )
            para.font.bold = True
            if template_style:
                para.font.name = (
                    template_style.title_font_family
                    or template_style.theme.major_font
                    or ""
                ) or None
            if template_style and template_style.title_font_color_rgb:
                try:
                    from pptx.dml.color import RGBColor

                    para.font.color.rgb = RGBColor.from_string(
                        template_style.title_font_color_rgb
                    )
                except Exception:
                    pass

    if not body_placed and body_paragraphs:
        max_font = _compute_max_font_size(region_map.text_region, len(body_paragraphs))
        txBox = new_slide.shapes.add_textbox(
            region_map.text_region.left,
            region_map.text_region.top,
            region_map.text_region.width,
            region_map.text_region.height,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, (text, level) in enumerate(body_paragraphs):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text = text
            para.level = level
            para.font.size = Pt(max_font)
            if template_style:
                para.font.name = (
                    template_style.body_font_family
                    or template_style.theme.minor_font
                    or ""
                ) or None
            if template_style and template_style.body_font_color_rgb:
                try:
                    from pptx.dml.color import RGBColor

                    para.font.color.rgb = RGBColor.from_string(
                        template_style.body_font_color_rgb
                    )
                except Exception:
                    pass
        font_family = "Calibri"
        if template_style:
            font_family = (
                template_style.body_font_family
                or template_style.theme.minor_font
                or "Calibri"
            )
        try:
            tf.fit_text(font_family=font_family, max_size=max_font)
        except Exception as e:
            if VERBOSE:
                print("[VERBOSE] Exception suppressed: %s" % str(e))
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # Place visuals into the visual_region (separate from text)
    visual_area = region_map.visual_region
    image_area = None
    if (content.images or generated_image_bytes) and (content.charts or content.tables):
        reserved = list(_layout_non_placeholder_text_bounds(new_slide.slide_layout))
        reserved.append(region_map.text_region)
        min_area = int(slide_width) * int(slide_height)
        min_area = int(min_area * 0.05)
        pic_ph = _best_visual_placeholder(
            new_slide.slide_layout,
            {PP_PLACEHOLDER.PICTURE},
            min_area=min_area,
            reserved=reserved,
            allow_body_fallback=False,
        )
        if pic_ph is not None:
            image_area = ContentArea(
                left=pic_ph.left,
                top=pic_ph.top,
                width=pic_ph.width,
                height=pic_ph.height,
            )
        else:
            gap = int(slide_width * 0.02)
            # Split visual area to avoid overlap
            if visual_area.width >= visual_area.height:
                image_w = int(visual_area.width * 0.30)
                image_area = ContentArea(
                    left=visual_area.left + visual_area.width - image_w,
                    top=visual_area.top,
                    width=image_w,
                    height=visual_area.height,
                )
                visual_area = ContentArea(
                    left=visual_area.left,
                    top=visual_area.top,
                    width=max(1, visual_area.width - image_w - gap),
                    height=visual_area.height,
                )
            else:
                image_h = int(visual_area.height * 0.25)
                image_area = ContentArea(
                    left=visual_area.left,
                    top=visual_area.top + visual_area.height - image_h,
                    width=visual_area.width,
                    height=image_h,
                )
                visual_area = ContentArea(
                    left=visual_area.left,
                    top=visual_area.top,
                    width=visual_area.width,
                    height=max(1, visual_area.height - image_h - gap),
                )

    if VERBOSE and content.charts:
        chart_ph = _largest_placeholder(
            new_slide.slide_layout, {PP_PLACEHOLDER.CHART}, min_area=0
        )
        print(
            "[VERBOSE] Chart transfer region: (%d,%d,%d,%d) chart_placeholder=%s"
            % (
                visual_area.left,
                visual_area.top,
                visual_area.width,
                visual_area.height,
                "yes" if chart_ph is not None else "no",
            )
        )
    _transfer_tables(
        new_slide, content.tables, visual_area, template_style=template_style
    )
    _transfer_images(new_slide, content.images, image_area or visual_area)
    _transfer_charts(
        new_slide, content.charts, visual_area, template_style=template_style
    )
    _transfer_shapes(new_slide, content.shapes_xml)
    if content_mix == ContentMix.TEXT_ONLY:
        _transfer_shapes(new_slide, content.text_shapes_xml)

    # ------------------------------------------------------------------
    # Insert generated images into picture placeholders.
    # Detection uses both the tracked indices AND XML-level type checking
    # as a fallback (type="pic" or "clipArt" in the XML = picture placeholder).
    # ------------------------------------------------------------------
    if generated_image_bytes:
        # Build a comprehensive set of picture placeholder indices using
        # the _is_picture_placeholder helper.
        _pic_ph_indices = set(content.image_placeholder_indices)
        image_inserted = False
        fallback_bounds = None
        min_area = int(slide_width) * int(slide_height)
        min_area = int(min_area * 0.10)

        for shape in list(new_slide.placeholders):
            if _is_picture_placeholder(shape):
                _pic_ph_indices.add(shape.placeholder_format.idx)

        # Choose a suitable picture placeholder (avoid tiny footer slots)
        reserved = list(_layout_non_placeholder_text_bounds(new_slide.slide_layout))
        reserved.append(region_map.text_region)
        if image_area is not None:
            fallback_bounds = (
                image_area.left,
                image_area.top,
                image_area.width,
                image_area.height,
            )
        pic_ph = _best_visual_placeholder(
            new_slide.slide_layout,
            {PP_PLACEHOLDER.PICTURE},
            min_area=min_area,
            reserved=reserved,
            allow_body_fallback=False,
        )
        if pic_ph is not None and _placeholder_area(pic_ph) >= min_area:
            try:
                left, top, width, height = (
                    pic_ph.left,
                    pic_ph.top,
                    pic_ph.width,
                    pic_ph.height,
                )
                fallback_bounds = (left, top, width, height)
                try:
                    new_slide.shapes._spTree.remove(pic_ph._element)
                except Exception:
                    pass
                _add_picture_within_bounds(
                    new_slide,
                    generated_image_bytes,
                    ContentArea(left=left, top=top, width=width, height=height),
                )
                image_inserted = True
                try:
                    populated_indices.add(pic_ph.placeholder_format.idx)
                except Exception:
                    pass
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))

        # If no picture was inserted (e.g., placeholder insertion unsupported),
        # fall back to adding the image as a regular picture scaled to the
        # visual region (or placeholder bounds if captured).
        if not image_inserted:
            try:
                if fallback_bounds:
                    left, top, width, height = fallback_bounds
                    bounds = ContentArea(left=left, top=top, width=width, height=height)
                else:
                    bounds = region_map.visual_region
                _add_picture_within_bounds(new_slide, generated_image_bytes, bounds)
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))

    # Ensure text renders on top of visual elements for mixed-content slides
    if content_mix != ContentMix.TEXT_ONLY and (
        content.title or content.subtitle or content.body_paragraphs
    ):
        _ensure_text_on_top(new_slide)

    # Ensure all text has sufficient contrast against the slide background
    if template_style:
        try:
            _ensure_text_contrast(new_slide, template_style)
        except Exception as e:
            if VERBOSE:
                print("[VERBOSE] Text contrast check failed: %s" % str(e))

    # ------------------------------------------------------------------
    # Remove ALL unpopulated placeholder XML elements from the slide shape
    # tree. This is the ONLY reliable way to eliminate ALL types of ghost
    # text ("Click to add title", "Click to add text", "Click icon to add
    # picture", content placeholder icons, etc.).
    # tf.clear() alone is insufficient for picture and content placeholders.
    # ------------------------------------------------------------------
    _clear_unused_placeholders(new_slide, populated_indices)

    # Remove any empty text boxes (non-placeholders) that may remain from
    # certain templates and would otherwise sit on top of visuals.
    _remove_empty_textboxes(new_slide)


# ---------------------------------------------------------------------------
# Step 1: Content Generation
# ---------------------------------------------------------------------------


def step_generate_content(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Generate PowerPoint content using Claude's pptx skill.

    Supports both streaming and non-streaming modes (controlled by the
    ``stream`` key in session_state). Streaming is required for long-running
    skill operations (>10 min) but may have issues with provider_data
    propagation. Non-streaming is simpler and more reliable for shorter
    operations but can timeout on complex presentations.

    This step:
    1. Creates a Claude agent with the pptx skill
    2. Runs the user's prompt to generate a presentation
    3. Downloads the generated .pptx file
    4. Extracts SlideContent from each slide
    5. Stores the extracted content in session_state
    """
    global VERBOSE
    VERBOSE = session_state.get("verbose", VERBOSE)

    user_prompt = step_input.input
    template_path = session_state.get("template_path", "")
    output_dir = session_state.get("output_dir", ".")

    print("=" * 60)
    print("Step 1: Generating presentation content with Claude...")
    print("=" * 60)

    # Build prompt with template context
    enhanced_prompt = user_prompt
    try:
        prs = Presentation(template_path)
        layouts = [layout.name for layout in prs.slide_layouts]
        layout_info = ", ".join(layouts)
        enhanced_prompt = (
            user_prompt + "\n\n"
            "Important structural requirements for template compatibility:\n"
            "- Use one clear title and concise bullet points per slide.\n"
            "- Do not apply custom fonts, colors, or theme styling.\n"
            "- Tables and charts are supported and will be transferred to the template.\n"
            "- Keep tables to max 6 rows x 5 columns.\n"
            "- Use bar, column, line, or pie charts with clearly labeled data.\n"
            "- The template has these available layouts: " + layout_info + ".\n"
            "- Use standard slide ordering: Title Slide, then Content Slides, then Closing."
        )
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Exception suppressed: %s" % str(e))

    # Create the Claude agent
    content_agent = Agent(
        name="Content Generator",
        model=Claude(
            id="claude-sonnet-4-5-20250929",  # Or: "claude-sonnet-4-6", "claude-opus-4-6"
            skills=[{"type": "anthropic", "skill_id": "pptx", "version": "latest"}],
        ),
        instructions=[
            "You are a structured content generator for PowerPoint presentations.",
            "Your output will be extracted and remapped to a corporate template.",
            "",
            "SLIDE STRUCTURE:",
            "- Use exactly one clear, descriptive title per slide.",
            "- Use concise bullet points for body content.",
            "- Limit to 4-6 bullet points per slide, each bullet max ~15 words.",
            "- Keep subtitle text on title slides to a single short line.",
            "- Follow standard slide ordering: Title Slide, then Content Slides, then Closing Slide.",
            "",
            "VISUAL ELEMENTS (use when appropriate):",
            "- Include tables for data comparisons and structured information.",
            "- Keep tables concise: no more than 6 rows and 5 columns.",
            "- Use bar, column, line, or pie charts for data visualization.",
            "",
            "FORMATTING RESTRICTIONS:",
            "- Do NOT apply custom fonts, colors, or theme styling.",
            "- Do NOT use SmartArt or complex nested graphic layouts.",
            "- Do NOT add speaker notes, animations, or transitions.",
        ],
        markdown=True,
    )

    # Determine streaming mode from session_state
    use_stream = session_state.get("stream", True)

    if use_stream:
        # Streaming mode: required for long-running skill operations (>10 min)
        response: RunOutput | None = None
        event_count = 0
        try:
            for event in content_agent.run(
                enhanced_prompt, stream=True, yield_run_output=True
            ):
                event_count += 1
                if isinstance(event, RunOutput):
                    response = event
                    if VERBOSE:
                        print(
                            "[VERBOSE] Received RunOutput after %d events" % event_count
                        )
                elif VERBOSE and event_count <= 5:
                    print(
                        "[VERBOSE] Stream event %d: type=%s"
                        % (event_count, type(event).__name__)
                    )
                elif VERBOSE and event_count == 6:
                    print("[VERBOSE] (Suppressing further stream event logs...)")
        except Exception as e:
            print("[ERROR] Agent streaming failed with exception: %s" % str(e))
            if VERBOSE:
                traceback.print_exc()
            return StepOutput(
                content="Error: Agent execution failed: %s" % str(e),
                success=False,
                stop=True,
            )
        if VERBOSE:
            print("[VERBOSE] Total stream events received: %d" % event_count)
    else:
        # Non-streaming mode: simpler, more reliable for shorter operations
        response = None
        try:
            response = content_agent.run(enhanced_prompt, stream=False)
        except Exception as e:
            print("[ERROR] Agent execution failed with exception: %s" % str(e))
            if VERBOSE:
                traceback.print_exc()
            return StepOutput(
                content="Error: Agent execution failed: %s" % str(e),
                success=False,
                stop=True,
            )

    if response is None:
        return StepOutput(
            content="Error: No response received from content agent.",
            success=False,
            stop=True,
        )
    print("\nAgent response received.")

    if VERBOSE:
        # Diagnostic: Inspect the RunOutput
        print("\n[VERBOSE] === Agent Response Diagnostics ===")
        if response is None:
            print("[VERBOSE] Response is None!")
        else:
            print("[VERBOSE] Response type: %s" % type(response).__name__)
            print(
                "[VERBOSE] Number of messages: %d"
                % (len(response.messages) if response.messages else 0)
            )
            if response.messages:
                for m_idx, msg in enumerate(response.messages):
                    print("[VERBOSE] Message %d:" % m_idx)
                    print("[VERBOSE]   Type: %s" % type(msg).__name__)
                    print("[VERBOSE]   Role: %s" % getattr(msg, "role", "N/A"))
                    print(
                        "[VERBOSE]   Has provider_data: %s"
                        % (
                            hasattr(msg, "provider_data")
                            and msg.provider_data is not None
                        )
                    )
                    if hasattr(msg, "provider_data") and msg.provider_data:
                        pd = msg.provider_data
                        print("[VERBOSE]   provider_data type: %s" % type(pd).__name__)
                        if isinstance(pd, dict):
                            print(
                                "[VERBOSE]   provider_data keys: %s" % list(pd.keys())
                            )
                        elif hasattr(pd, "__dict__"):
                            print(
                                "[VERBOSE]   provider_data attrs: %s"
                                % list(vars(pd).keys())[:20]
                            )
                    # Also show content snippet
                    content_text = getattr(msg, "content", None)
                    if content_text and isinstance(content_text, str):
                        print(
                            "[VERBOSE]   Content (first 200 chars): %s"
                            % content_text[:200]
                        )
            # Check for response-level attributes
            for attr_name in [
                "content",
                "tool_calls",
                "extra_data",
                "metrics",
                "model_provider_data",
            ]:
                val = getattr(response, attr_name, "MISSING")
                if val != "MISSING":
                    if isinstance(val, str):
                        print(
                            "[VERBOSE] response.%s (first 200 chars): %s"
                            % (attr_name, val[:200])
                        )
                    elif val is not None:
                        print(
                            "[VERBOSE] response.%s type: %s"
                            % (attr_name, type(val).__name__)
                        )
        print("[VERBOSE] === End Diagnostics ===\n")

    # Download the generated file
    client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
    generated_file = None

    if response.messages:
        for msg_idx, msg in enumerate(response.messages):
            if VERBOSE:
                print("[VERBOSE] Checking message %d for provider_data..." % msg_idx)
            if hasattr(msg, "provider_data") and msg.provider_data:
                if VERBOSE:
                    print(
                        "[VERBOSE] Message %d has provider_data, calling download_skill_files..."
                        % msg_idx
                    )
                try:
                    files = download_skill_files(
                        msg.provider_data, client, output_dir=output_dir
                    )
                except Exception as e:
                    print(
                        "[ERROR] download_skill_files raised an exception: %s" % str(e)
                    )
                    if VERBOSE:
                        traceback.print_exc()
                    files = None

                if VERBOSE:
                    print("[VERBOSE] download_skill_files returned: %s" % files)

                if files:
                    for f in files:
                        if not f.endswith(".pptx"):
                            if VERBOSE:
                                print("[VERBOSE] Skipping non-.pptx file: %s" % f)
                            continue
                        try:
                            Presentation(f)
                            generated_file = f
                            print("Downloaded: %s" % generated_file)
                            break
                        except Exception as e:
                            print(
                                "[WARNING] File '%s' failed pptx validation: %s"
                                % (f, str(e))
                            )
                            continue
                    if not generated_file:
                        # Fallback: try any file
                        for f in files:
                            try:
                                Presentation(f)
                                generated_file = f
                                print("Downloaded (fallback): %s" % generated_file)
                                break
                            except Exception as e:
                                if VERBOSE:
                                    print(
                                        "[VERBOSE] Fallback file '%s' also failed: %s"
                                        % (f, str(e))
                                    )
                                continue
                elif VERBOSE:
                    print("[VERBOSE] No files returned from download_skill_files")
            elif VERBOSE:
                print("[VERBOSE] Message %d has no provider_data" % msg_idx)
            if generated_file:
                break
    else:
        if VERBOSE:
            print("[VERBOSE] response.messages is empty or None!")

    # Fallback: check response.model_provider_data directly
    # This catches cases where provider_data is on the RunOutput but not on individual messages
    if not generated_file and response.model_provider_data:
        if VERBOSE:
            print(
                "[VERBOSE] Trying fallback: response.model_provider_data = %s"
                % response.model_provider_data
            )
        try:
            files = download_skill_files(
                response.model_provider_data, client, output_dir=output_dir
            )
        except Exception as e:
            print("[ERROR] download_skill_files (fallback) raised: %s" % str(e))
            if VERBOSE:
                traceback.print_exc()
            files = None
        if files:
            if VERBOSE:
                print("[VERBOSE] Fallback download_skill_files returned: %s" % files)
            for f in files:
                if f.endswith(".pptx"):
                    try:
                        Presentation(f)
                        generated_file = f
                        print(
                            "Downloaded (via model_provider_data fallback): %s"
                            % generated_file
                        )
                        break
                    except Exception as e:
                        print(
                            "[WARNING] Fallback file '%s' failed pptx validation: %s"
                            % (f, str(e))
                        )
            if not generated_file:
                for f in files:
                    try:
                        Presentation(f)
                        generated_file = f
                        print(
                            "Downloaded (via model_provider_data fallback): %s"
                            % generated_file
                        )
                        break
                    except Exception as e:
                        if VERBOSE:
                            print(
                                "[VERBOSE] Fallback file '%s' also failed: %s"
                                % (f, str(e))
                            )

    if not generated_file:
        if VERBOSE:
            print("\n[VERBOSE] === File Generation Failure Summary ===")
            print("[VERBOSE] No valid .pptx file was found in any message.")
            print(
                "[VERBOSE] Total messages checked: %d"
                % (len(response.messages) if response.messages else 0)
            )
            messages_with_pd = sum(
                1
                for m in (response.messages or [])
                if hasattr(m, "provider_data") and m.provider_data
            )
            print("[VERBOSE] Messages with provider_data: %d" % messages_with_pd)
            print("[VERBOSE] Hint: The agent may not have generated a file. Check if:")
            print(
                "[VERBOSE]   1. The agent's response contained tool use / skill output"
            )
            print("[VERBOSE]   2. The ANTHROPIC_API_KEY has access to the pptx skill")
            print("[VERBOSE]   3. The prompt was understood correctly by the agent")
            print("[VERBOSE] === End Failure Summary ===\n")
        return StepOutput(
            content="Error: No presentation file was generated.",
            success=False,
            stop=True,
        )

    # Extract slide content
    print("\nExtracting slide content...")
    generated_prs = Presentation(generated_file)
    slides_data = []

    # Check template layouts for picture placeholders so the image planner
    # knows which slides have dedicated picture slots in the template.
    template_prs = None
    try:
        template_prs = Presentation(template_path)
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Exception suppressed: %s" % str(e))

    total_gen_slides = len(list(generated_prs.slides))

    for idx, slide in enumerate(generated_prs.slides):
        content = _extract_slide_content(slide)

        # Check if the template layout for this slide position has image placeholders
        has_template_image_ph = False
        if template_prs is not None:
            try:
                content_mix = _classify_content_mix(content, has_generated_image=False)
                layout = _find_best_layout(
                    template_prs,
                    idx,
                    total_gen_slides,
                    content_mix=content_mix,
                )
                for ph in layout.placeholders:
                    ph_fmt = ph.placeholder_format
                    if ph_fmt is not None and ph_fmt.type == PP_PLACEHOLDER.PICTURE:
                        has_template_image_ph = True
                        break
            except Exception as e:
                if VERBOSE:
                    print("[VERBOSE] Exception suppressed: %s" % str(e))

        slide_info = {
            "index": idx,
            "title": content.title,
            "subtitle": content.subtitle,
            "body_text": " | ".join(
                t if isinstance(t, str) else t[0] for t in content.body_paragraphs[:3]
            )
            if content.body_paragraphs
            else "",
            "bullet_count": len(content.body_paragraphs),
            "has_table": len(content.tables) > 0,
            "has_chart": len(content.charts) > 0,
            "has_image": len(content.images) > 0,
            "has_shapes": len(content.shapes_xml) > 0,
            "has_image_placeholder": has_template_image_ph,
        }
        slides_data.append(slide_info)
        print(
            "  Slide %d: '%s' | tables:%d charts:%d images:%d img_ph:%s"
            % (
                idx + 1,
                content.title[:40] if content.title else "",
                len(content.tables),
                len(content.charts),
                len(content.images),
                has_template_image_ph,
            )
        )

    # Store in session state
    session_state["generated_file"] = generated_file
    session_state["slides_data"] = slides_data
    session_state["total_slides"] = len(slides_data)

    slides_summary = json.dumps(slides_data, indent=2)
    return StepOutput(
        content=slides_summary,
        success=True,
    )


# ---------------------------------------------------------------------------
# Step 2: Image Planning (Agent with output_schema)
# ---------------------------------------------------------------------------

# This agent decides which slides need AI-generated images
image_planner = Agent(
    name="Image Planner",
    model=Gemini(id="gemini-3-flash-preview"),
    instructions=[
        "You are an image planning specialist for PowerPoint presentations.",
        "You will receive slide metadata AND the user's presentation topic/request.",
        "",
        "RULES (follow strictly):",
        "- If has_image_placeholder is true: ALWAYS set needs_image to true.",
        "- Title slides (index 0): ALWAYS generate an image — first impressions matter.",
        "- Slides with existing images: NEVER generate (already have visuals).",
        "- Data slides (has_table or has_chart is true): Usually skip unless topic is very visual.",
        "- All other content slides: Default to YES. Visuals enhance every presentation.",
        "- If the user explicitly requested 'visuals', 'images', or 'with pictures',",
        "  generate images for at LEAST half of eligible slides.",
        "",
        "IMPORTANT: When in doubt, generate an image. It is better to have too many",
        "images than too few. Empty picture placeholders look unprofessional.",
        "",
        "When writing image prompts:",
        "- Use the PRESENTATION TOPIC to create relevant imagery.",
        "- Describe professional, clean, modern illustrations.",
        "- Use abstract or metaphorical imagery, not literal text depictions.",
        "- Specify style: 'minimalist corporate illustration', 'flat design', etc.",
        "- Keep prompts under 100 words.",
        "- Make images suitable for a professional business presentation.",
    ],
    output_schema=ImagePlan,
    markdown=False,
)


def step_plan_images(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Plan which slides need AI-generated images.

    This is a function step wrapper around the image_planner agent.
    It combines the user's original prompt with slide metadata so the
    planner has full context (the framework's _prepare_message() would
    otherwise discard the original prompt when previous step outputs exist).
    """
    global VERBOSE
    VERBOSE = session_state.get("verbose", VERBOSE)

    print("\n" + "=" * 60)
    print("Step 2: Planning images for slides...")
    print("=" * 60)

    # Retrieve user prompt from session_state (stored at workflow creation)
    user_prompt = session_state.get("user_prompt", "professional business presentation")

    # Get slide metadata from previous step output
    slides_json = step_input.previous_step_content or "[]"
    if not isinstance(slides_json, str):
        slides_json = json.dumps(slides_json, indent=2) if slides_json else "[]"

    # Construct a combined message with full context
    combined_message = (
        'Presentation topic: "%s"\n\n'
        "Slide metadata:\n%s\n\n"
        "Analyze each slide and decide which ones need AI-generated images.\n"
        "Consider the presentation topic when writing image prompts."
    ) % (user_prompt, slides_json)

    if VERBOSE:
        print(
            "[VERBOSE] Image planner input (first 500 chars):\n%s"
            % combined_message[:500]
        )

    # Run the image_planner agent directly with the combined message
    try:
        response = image_planner.run(combined_message, stream=False)
    except Exception as e:
        print("[ERROR] Image planner failed: %s" % str(e))
        if VERBOSE:
            traceback.print_exc()
        return StepOutput(
            content=json.dumps({"decisions": []}),
            success=True,
        )

    # Extract the structured output
    if response and response.content:
        if isinstance(response.content, BaseModel):
            result = response.content.model_dump_json()
        elif isinstance(response.content, dict):
            result = json.dumps(response.content)
        else:
            result = str(response.content)
        if VERBOSE:
            print(
                "[VERBOSE] Image planner output (first 500 chars):\n%s" % result[:500]
            )
        return StepOutput(content=result, success=True)

    print("[WARNING] Image planner returned no content.")
    return StepOutput(
        content=json.dumps({"decisions": []}),
        success=True,
    )


# ---------------------------------------------------------------------------
# Step 3: Image Generation
# ---------------------------------------------------------------------------


def step_generate_images(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Generate images for slides that need them using NanoBanana (powered by Gemini).

    This step:
    1. Reads the image plan from step 2 (produced by the Gemini image planner)
    2. For each slide marked as needing an image, generates one with NanoBanana
    3. Stores the generated image bytes in session_state keyed by slide index
    """
    global VERBOSE
    VERBOSE = session_state.get("verbose", VERBOSE)

    print("\n" + "=" * 60)
    print("Step 3: Generating images with NanoBanana...")
    print("=" * 60)

    # Parse the image plan from step 2
    image_plan_content = step_input.previous_step_content
    generated_images: Dict[int, bytes] = {}

    if not image_plan_content:
        print("No image plan received. Skipping image generation.")
        session_state["generated_images"] = {}
        return StepOutput(content="No images to generate.", success=True)

    # Parse the plan
    try:
        if isinstance(image_plan_content, str):
            plan_data = json.loads(image_plan_content)
        elif isinstance(image_plan_content, dict):
            plan_data = image_plan_content
        elif isinstance(image_plan_content, ImagePlan):
            plan_data = image_plan_content.model_dump()
        elif isinstance(image_plan_content, BaseModel):
            plan_data = image_plan_content.model_dump()
        else:
            plan_data = json.loads(str(image_plan_content))
    except (json.JSONDecodeError, TypeError):
        print("Could not parse image plan. Skipping image generation.")
        session_state["generated_images"] = {}
        return StepOutput(content="Could not parse image plan.", success=True)

    decisions = plan_data.get("decisions", [])
    slides_needing_images = [d for d in decisions if d.get("needs_image", False)]

    # Enforce minimum images if configured
    min_images = session_state.get("min_images", 1)
    if len(slides_needing_images) < min_images and min_images > 0:
        print(
            "Image planner selected %d slides, but --min-images=%d. Adding more..."
            % (len(slides_needing_images), min_images)
        )
        # Find slides not yet selected, preferring those with image placeholders
        selected_indices = {d["slide_index"] for d in slides_needing_images}
        remaining = [d for d in decisions if d["slide_index"] not in selected_indices]
        # Sort: image_placeholder slides first, then title slide (index 0)
        remaining.sort(
            key=lambda d: (
                not d.get("has_image_placeholder", False),
                d.get("slide_index", 99),
            )
        )

        for d in remaining:
            if len(slides_needing_images) >= min_images:
                break
            # Create a generic prompt based on available context
            topic = session_state.get("user_prompt", "professional business")
            slide_idx = d.get("slide_index", 0)
            d["needs_image"] = True
            d["image_prompt"] = (
                "Professional minimalist corporate illustration related to %s, "
                "suitable for slide %d of a business presentation. "
                "Clean, modern flat design with subtle gradients."
            ) % (topic, slide_idx + 1)
            slides_needing_images.append(d)

    if not slides_needing_images:
        print("Image planner decided no slides need images.")
        session_state["generated_images"] = {}
        return StepOutput(content="No slides need images.", success=True)

    # Check which slides already have images from Claude
    slides_data = session_state.get("slides_data", [])
    slides_with_existing_images = {
        s["index"] for s in slides_data if s.get("has_image", False)
    }

    # Initialize NanoBanana
    google_api_key = os.getenv("GOOGLE_API_KEY")
    if not google_api_key:
        print("GOOGLE_API_KEY not set. Skipping image generation.")
        session_state["generated_images"] = {}
        return StepOutput(
            content="GOOGLE_API_KEY not set. Skipped image generation.",
            success=True,
        )

    nano_banana = NanoBananaTools(
        api_key=google_api_key,
        aspect_ratio="16:9",  # Widescreen for presentations
    )

    for decision in slides_needing_images:
        slide_idx = decision["slide_index"]
        prompt = decision.get("image_prompt", "")

        if slide_idx in slides_with_existing_images:
            print(
                "  Slide %d: Already has image from Claude, skipping." % (slide_idx + 1)
            )
            continue

        if not prompt:
            print("  Slide %d: No prompt provided, skipping." % (slide_idx + 1))
            continue

        print("  Slide %d: Generating image..." % (slide_idx + 1))
        print("    Prompt: %s" % prompt[:80])

        try:
            result = nano_banana.create_image(prompt)
            if hasattr(result, "images") and result.images:
                for img in result.images:
                    if hasattr(img, "content") and img.content:
                        generated_images[slide_idx] = img.content
                        print(
                            "    Generated successfully (%d bytes)" % len(img.content)
                        )
                        break
            elif isinstance(result, str) and "successfully" in result.lower():
                print("    Image generated but no bytes returned.")
        except Exception as e:
            print("    Failed to generate image: %s" % str(e))
            if VERBOSE:
                traceback.print_exc()

    session_state["generated_images"] = generated_images
    return StepOutput(
        content="Generated %d image(s) for presentation." % len(generated_images),
        success=True,
    )


# ---------------------------------------------------------------------------
# Step 4: Template Assembly
# ---------------------------------------------------------------------------


def step_assemble_template(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Apply template styling and assemble the final presentation.

    Handles image placeholder detection on each template layout so that
    generated images are inserted into dedicated picture placeholders when
    available, rather than being added as free-floating pictures.

    This step:
    1. Opens the template and generated presentation
    2. Creates slides from template layouts with extracted content
    3. Inserts AI-generated images into picture placeholders (or as free-floating)
    4. Saves the final output
    """
    global VERBOSE
    VERBOSE = session_state.get("verbose", VERBOSE)

    print("\n" + "=" * 60)
    print("Step 4: Assembling final presentation with template...")
    print("=" * 60)

    template_path = session_state.get("template_path", "")
    generated_file = session_state.get("generated_file", "")
    output_path = session_state.get("output_path", "presentation_from_template.pptx")
    generated_images = session_state.get("generated_images", {})

    if not generated_file or not os.path.isfile(generated_file):
        return StepOutput(
            content="Error: Generated file not found.",
            success=False,
            stop=True,
        )

    print("Template: %s" % template_path)
    print("Generated: %s" % generated_file)
    print("Output: %s" % output_path)

    # Open the generated presentation
    generated_prs = Presentation(generated_file)
    generated_slides = list(generated_prs.slides)
    total_slides = len(generated_slides)

    if VERBOSE:
        print("[VERBOSE] Generated presentation has %d slides" % total_slides)

    if total_slides == 0:
        shutil.copy2(template_path, output_path)
        return StepOutput(
            content="Warning: No slides found. Copied template as-is.",
            success=True,
        )

    # Create output from template
    shutil.copy2(template_path, output_path)
    output_prs = Presentation(output_path)
    slide_width = output_prs.slide_width
    slide_height = output_prs.slide_height

    # Extract template styles for application to new elements
    template_style = None
    try:
        template_prs_for_styles = Presentation(template_path)
        template_style = _extract_template_styles(template_prs_for_styles)
    except Exception as e:
        if VERBOSE:
            print("[VERBOSE] Template style extraction failed: %s" % str(e))

    if VERBOSE:
        print(
            "[VERBOSE] Template layouts available: %s"
            % [layout.name for layout in output_prs.slide_layouts]
        )

    # Remove existing slides from template
    while len(output_prs.slides._sldIdLst) > 0:
        sldId = output_prs.slides._sldIdLst[0]
        rId = sldId.get(
            etree.QName(
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
                "id",
            )
        )
        if rId is not None:
            output_prs.part.drop_rel(rId)
        output_prs.slides._sldIdLst.remove(sldId)

    print("Cleared template slides. Building final presentation...")

    # For each generated slide, create a template-styled slide
    for idx, gen_slide in enumerate(generated_slides):
        content = _extract_slide_content(gen_slide)

        # Look up generated image for this slide.
        # The keys may be int or str depending on how they were stored.
        gen_img = generated_images.get(idx) or generated_images.get(str(idx))

        # Find layout FIRST so we can detect picture placeholders before
        # deciding whether to add a free-floating image.
        content_mix = _classify_content_mix(
            content, has_generated_image=gen_img is not None
        )
        layout = _find_best_layout(
            output_prs,
            idx,
            total_slides,
            content_mix=content_mix,
            has_generated_image=gen_img is not None,
        )
        if VERBOSE:
            summary = _layout_placeholder_summary(layout)
            print(
                "[VERBOSE] Slide %d chose layout '%s' placeholders: %s"
                % (idx + 1, layout.name, summary)
            )

        # Detect picture placeholders on the chosen template layout.
        # Uses both PP_PLACEHOLDER enum and XML-level fallback detection.
        for ph in layout.placeholders:
            ph_fmt = ph.placeholder_format
            if ph_fmt is not None:
                is_pic_ph = False
                # Strategy 1: int comparison with raw OOXML value 18
                try:
                    if ph_fmt.type is not None and (
                        int(ph_fmt.type) == 18 or str(ph_fmt.type) == "PICTURE (18)"
                    ):
                        is_pic_ph = True
                except (ValueError, TypeError):
                    pass
                # Strategy 2: PP_PLACEHOLDER enum
                if not is_pic_ph:
                    try:
                        if ph_fmt.type == PP_PLACEHOLDER.PICTURE:
                            is_pic_ph = True
                    except Exception as e:
                        if VERBOSE:
                            print("[VERBOSE] Exception suppressed: %s" % str(e))
                # Strategy 3: XML-level fallback
                if not is_pic_ph:
                    try:
                        nsmap = {
                            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                        }
                        ph_elem = ph._element.find(".//p:ph", nsmap)
                        if ph_elem is not None and ph_elem.get("type") in (
                            "pic",
                            "clipArt",
                        ):
                            is_pic_ph = True
                    except Exception as e:
                        if VERBOSE:
                            print("[VERBOSE] Exception suppressed: %s" % str(e))
                if is_pic_ph:
                    content.has_image_placeholder = True
                    if ph_fmt.idx not in content.image_placeholder_indices:
                        content.image_placeholder_indices.append(ph_fmt.idx)

        # Only add as free-floating picture if the layout has NO picture
        # placeholders. When picture placeholders exist, the image will be
        # inserted into the placeholder by _populate_slide instead.
        if gen_img is not None and not content.has_image_placeholder:
            content.images.append(
                ImageData(
                    blob=gen_img,
                    left=0,
                    top=0,
                    width=int(Inches(8.0)),
                    height=int(Inches(4.5)),
                    content_type="image/png",
                )
            )

        visual_info = []
        if content.tables:
            visual_info.append("%d table(s)" % len(content.tables))
        if content.images:
            visual_info.append("%d image(s)" % len(content.images))
        if content.charts:
            visual_info.append("%d chart(s)" % len(content.charts))
        if content.has_image_placeholder:
            visual_info.append("img placeholder(s)")
        visual_str = ", ".join(visual_info) if visual_info else "text only"

        print(
            "  Slide %d: layout '%s' | title: '%s' | %s"
            % (
                idx + 1,
                layout.name,
                content.title[:40] if content.title else "",
                visual_str,
            )
        )

        new_slide = output_prs.slides.add_slide(layout)
        # Pass generated image bytes to _populate_slide for insertion
        # into picture placeholders
        _populate_slide(
            new_slide,
            content,
            slide_width,
            slide_height,
            generated_image_bytes=gen_img,
            template_style=template_style,
        )

    output_prs.save(output_path)
    print("\nSaved final presentation: %s" % output_path)

    return StepOutput(
        content="Presentation saved to %s (%d slides)" % (output_path, total_slides),
        success=True,
    )


# ---------------------------------------------------------------------------
# CLI and Main
# ---------------------------------------------------------------------------


def parse_args():
    """Parse command-line arguments."""
    parser = argparse.ArgumentParser(
        description="Workflow: Generate PowerPoint with Claude, add AI images, apply template."
    )
    parser.add_argument(
        "--template",
        "-t",
        required=True,
        help="Path to the .pptx template file.",
    )
    parser.add_argument(
        "--output",
        "-o",
        default="presentation_from_template.pptx",
        help="Output filename (default: presentation_from_template.pptx).",
    )
    parser.add_argument(
        "--prompt",
        "-p",
        default=None,
        help="Custom prompt for the presentation content.",
    )
    parser.add_argument(
        "--no-images",
        action="store_true",
        help="Skip AI image generation (Steps 2 and 3).",
    )
    parser.add_argument(
        "--no-stream",
        action="store_true",
        help="Disable streaming mode for Claude agent (more reliable for shorter prompts, "
        "but may timeout on complex presentations).",
    )
    parser.add_argument(
        "--verbose",
        "-v",
        action="store_true",
        help="Enable verbose/debug logging for troubleshooting.",
    )
    parser.add_argument(
        "--min-images",
        type=int,
        default=1,
        help="Minimum number of slides that must have AI-generated images "
        "(default: 1, 0 = let planner decide).",
    )
    return parser.parse_args()


if __name__ == "__main__":
    args = parse_args()

    VERBOSE = args.verbose

    if not os.path.isfile(args.template):
        print("Error: Template file not found: %s" % args.template)
        sys.exit(1)

    if not args.template.endswith(".pptx"):
        print("Error: Template file must be a .pptx file")
        sys.exit(1)

    if not os.getenv("ANTHROPIC_API_KEY"):
        raise ValueError("ANTHROPIC_API_KEY environment variable not set")

    # Setup output directory
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_dir = os.path.join(script_dir, "output")
    os.makedirs(output_dir, exist_ok=True)

    output_path = args.output
    if not os.path.isabs(output_path):
        output_path = os.path.join(output_dir, output_path)

    # Default prompt
    prompt = args.prompt
    if not prompt:
        prompt = (
            "Create a 6-slide business presentation with the following structure.\n"
            "Rules: Use one clear title per slide, 4-6 concise bullet points for text slides,\n"
            "and include tables or charts where specified. Do not apply custom styling.\n"
            "\n"
            "Slide 1 - Title Slide:\n"
            "  Title: 'Strategic Overview 2026'\n"
            "  Subtitle: 'Annual Business Review and Forward Plan'\n"
            "\n"
            "Slide 2 - Market Analysis (text with bullets):\n"
            "  Title: 'Market Analysis'\n"
            "  Bullets: Key market trends, growth opportunities, competitive landscape,\n"
            "  emerging technologies, and regulatory changes.\n"
            "\n"
            "Slide 3 - Financial Performance (include a TABLE):\n"
            "  Title: 'Financial Performance'\n"
            "  Table with columns: Metric, Q1, Q2, Q3, Q4\n"
            "  Rows: Revenue ($M), Costs ($M), Profit ($M), Growth (%)\n"
            "\n"
            "Slide 4 - Revenue Trends (include a BAR CHART):\n"
            "  Title: 'Quarterly Revenue Trends'\n"
            "  Bar chart showing quarterly revenue for 2024 vs 2025.\n"
            "\n"
            "Slide 5 - Our Strategy (text with bullets):\n"
            "  Title: 'Strategic Priorities'\n"
            "  Bullets: Three-pillar growth approach, market expansion,\n"
            "  product innovation, operational excellence.\n"
            "\n"
            "Slide 6 - Closing Slide:\n"
            "  Title: 'Next Steps'\n"
            "  Bullets: Implementation timeline, key milestones, success metrics.\n"
            "\n"
            "Save as 'generated_content.pptx'"
        )

    print("=" * 60)
    print("PowerPoint Template Workflow")
    print("=" * 60)
    print("Template: %s" % args.template)
    print("Output:   %s" % output_path)
    print("Images:   %s" % ("disabled" if args.no_images else "enabled"))
    print("Stream:   %s" % ("disabled" if args.no_stream else "enabled"))
    if VERBOSE:
        print("Verbose:  enabled")

    # Build workflow steps
    steps: List[Step] = [
        Step(name="Content Generation", executor=step_generate_content),
    ]

    if not args.no_images:
        steps.append(Step(name="Image Planning", executor=step_plan_images))
        steps.append(Step(name="Image Generation", executor=step_generate_images))

    steps.append(Step(name="Template Assembly", executor=step_assemble_template))

    # Create and run the workflow
    #
    # session_state schema:
    #   template_path      (str)  - Path to the .pptx template file
    #   output_dir         (str)  - Directory for intermediate output files
    #   output_path        (str)  - Path for the final assembled presentation
    #   generated_file     (str)  - Path to the Claude-generated .pptx file (set by step 1)
    #   slides_data        (list) - Per-slide metadata dicts (set by step 1)
    #   total_slides       (int)  - Number of slides in the generated presentation
    #   generated_images   (dict) - Mapping of slide index -> image bytes (set by step 3)
    #   verbose            (bool) - Whether verbose logging is enabled
    #   stream             (bool) - Whether to use streaming mode for Claude agent
    #   user_prompt        (str)  - Original user prompt (for image planner context)
    #   min_images         (int)  - Minimum slides that must have AI-generated images
    workflow = Workflow(
        name="PowerPoint Template Workflow",
        steps=steps,
        session_state={
            "template_path": args.template,
            "output_dir": output_dir,
            "output_path": output_path,
            "generated_file": "",
            "slides_data": [],
            "total_slides": 0,
            "generated_images": {},
            "verbose": args.verbose,
            "stream": not args.no_stream,
            "user_prompt": prompt,
            "min_images": args.min_images,
        },
    )

    workflow.print_response(input=prompt, markdown=True)

    print("\n" + "=" * 60)
    print("Workflow complete!")
    print("Output: %s" % output_path)
    print("=" * 60)
