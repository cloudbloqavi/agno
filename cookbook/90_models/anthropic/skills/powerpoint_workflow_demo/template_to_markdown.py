"""
Template to Markdown Utility.

Extracts comprehensive styling from a PPTX template and returns it as a
markdown-formatted string suitable for injection into Claude's prompt.

Usage as a module:
    from template_to_markdown import extract_template_to_markdown
    markdown = extract_template_to_markdown("my_template.pptx")

Usage as a CLI tool:
    python template_to_markdown.py my_template.pptx
    python template_to_markdown.py my_template.pptx --output style_guide.md

Prerequisites:
    pip install python-pptx lxml
"""

import argparse
from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

# Semantic role names for color scheme slots (in order they appear in XML)
COLOR_ROLE_NAMES = [
    "Dark 1 (Text)",
    "Dark 2",
    "Light 1 (Background)",
    "Light 2",
    "Accent 1 (Primary)",
    "Accent 2 (Secondary)",
    "Accent 3",
    "Accent 4",
    "Accent 5",
    "Accent 6",
    "Hyperlink",
    "Followed Hyperlink",
]

# Color role tags in the order they appear in <a:clrScheme>
COLOR_ROLE_TAGS = [
    "dk1",
    "dk2",
    "lt1",
    "lt2",
    "accent1",
    "accent2",
    "accent3",
    "accent4",
    "accent5",
    "accent6",
    "hlink",
    "folHlink",
]

# Usage hints per role
COLOR_USAGE_HINTS: dict[str, str] = {
    "Dark 1 (Text)": "Main body text",
    "Dark 2": "Subtitles, secondary text",
    "Light 1 (Background)": "Slide backgrounds",
    "Light 2": "Subtle backgrounds, dividers",
    "Accent 1 (Primary)": "Title backgrounds, key shapes, first data series",
    "Accent 2 (Secondary)": "Secondary highlights, second data series",
    "Accent 3": "Third data series, supporting elements",
    "Accent 4": "Fourth data series",
    "Accent 5": "Fifth data series",
    "Accent 6": "Sixth data series",
    "Hyperlink": "Hyperlinks",
    "Followed Hyperlink": "Visited hyperlinks",
}

# Layout name to purpose hint mapping
LAYOUT_HINTS: dict[str, str] = {
    "Title Slide": "Opening slide with title and subtitle",
    "Title and Content": "Content with title header",
    "Section Header": "Section divider slides",
    "Two Content": "Side-by-side content comparison",
    "Comparison": "Two-column comparison with labels",
    "Title Only": "Title with free-form content area",
    "Blank": "Fully custom slide with no placeholders",
    "Content with Caption": "Content area with caption text",
    "Picture with Caption": "Image placeholder with caption",
    "Title and Vertical Text": "Vertical text layout",
    "Vertical Title and Text": "Vertical title with text",
}

# ---------------------------------------------------------------------------
# Helper: get theme XML root from a presentation
# ---------------------------------------------------------------------------


def _get_theme_element(prs: Presentation) -> Optional[etree._Element]:
    """Return the root lxml element of the first slide master's theme."""
    try:
        slide_master = prs.slide_masters[0]
        for rel in slide_master.part.rels.values():
            if "theme" in rel.reltype:
                theme_part = rel.target_part
                try:
                    return theme_part.element
                except AttributeError:
                    return etree.fromstring(theme_part.blob)
    except Exception:
        pass
    return None


def _hex_from_clr_element(clr_el: etree._Element) -> str:
    """Extract a hex color string from a color scheme child element."""
    srgb = clr_el.find("a:srgbClr", NSMAP)
    if srgb is not None:
        return (srgb.get("val") or "").upper()
    sys_clr = clr_el.find("a:sysClr", NSMAP)
    if sys_clr is not None:
        val = sys_clr.get("lastClr") or sys_clr.get("val") or ""
        return val.upper()
    return ""


# ---------------------------------------------------------------------------
# Extraction functions
# ---------------------------------------------------------------------------


def extract_slide_dimensions(prs: Presentation) -> dict[str, float]:
    """Return slide width/height in EMU and inches.

    Returns:
        Dict with keys: width_emu, height_emu, width_in, height_in.
    """
    width_emu = prs.slide_width
    height_emu = prs.slide_height
    return {
        "width_emu": int(width_emu),
        "height_emu": int(height_emu),
        "width_in": round(Emu(width_emu).inches, 2),
        "height_in": round(Emu(height_emu).inches, 2),
    }


def extract_theme_colors(prs: Presentation) -> dict[str, str]:
    """Extract theme color hex values keyed by semantic role name.

    Returns:
        Dict mapping role name (e.g. "Accent 1 (Primary)") to hex string (e.g. "4472C4").
    """
    result: dict[str, str] = {}
    theme_el = _get_theme_element(prs)
    if theme_el is None:
        return result

    clr_scheme = theme_el.find(".//a:clrScheme", NSMAP)
    if clr_scheme is None:
        return result

    # Build a tag -> hex lookup from the scheme children
    tag_to_hex: dict[str, str] = {}
    for child in clr_scheme:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        tag_to_hex[tag] = _hex_from_clr_element(child)

    for tag, role in zip(COLOR_ROLE_TAGS, COLOR_ROLE_NAMES):
        hex_val = tag_to_hex.get(tag, "")
        if hex_val:
            result[role] = hex_val

    return result


def extract_font_scheme(prs: Presentation) -> dict[str, str]:
    """Extract major (heading) and minor (body) font names from theme.

    Returns:
        Dict with keys: "major" (heading font) and "minor" (body font).
    """
    fonts = {"major": "Calibri Light", "minor": "Calibri"}
    theme_el = _get_theme_element(prs)
    if theme_el is None:
        return fonts

    font_scheme = theme_el.find(".//a:fontScheme", NSMAP)
    if font_scheme is None:
        return fonts

    major_el = font_scheme.find(".//a:majorFont/a:latin", NSMAP)
    if major_el is not None:
        typeface = major_el.get("typeface", "")
        if typeface and typeface not in ("+mj-lt", "+mn-lt"):
            fonts["major"] = typeface

    minor_el = font_scheme.find(".//a:minorFont/a:latin", NSMAP)
    if minor_el is not None:
        typeface = minor_el.get("typeface", "")
        if typeface and typeface not in ("+mj-lt", "+mn-lt"):
            fonts["minor"] = typeface

    return fonts


def extract_placeholder_fonts(layout) -> dict[str, dict]:
    """Extract font info from placeholders in a slide layout.

    Args:
        layout: A pptx SlideLayout object.

    Returns:
        Dict keyed by placeholder type string ("title", "body", "subtitle", etc.).
        Each value is a dict: {family, size_pt, bold, color_hex, align}.
    """
    result: dict[str, dict] = {}

    for ph in layout.placeholders:
        ph_type = ph.placeholder_format.type
        # Map numeric type to label
        type_label = _ph_type_label(ph_type, ph.placeholder_format.idx)

        if not ph.has_text_frame:
            continue

        font_info: dict = {
            "family": None,
            "size_pt": None,
            "bold": None,
            "color_hex": None,
            "align": None,
        }

        tf = ph.text_frame
        for para in tf.paragraphs:
            if font_info["align"] is None and para.alignment is not None:
                font_info["align"] = str(para.alignment)
            for run in para.runs:
                f = run.font
                if font_info["family"] is None and f.name:
                    font_info["family"] = f.name
                if font_info["size_pt"] is None and f.size:
                    font_info["size_pt"] = round(f.size.pt)
                if font_info["bold"] is None and f.bold is not None:
                    font_info["bold"] = f.bold
                if font_info["color_hex"] is None:
                    try:
                        rgb = f.color.rgb
                        font_info["color_hex"] = str(rgb)
                    except Exception:
                        pass

        # Also check paragraph-level default run properties via XML
        if font_info["size_pt"] is None or font_info["family"] is None:
            for txBody in ph._element.findall(".//a:txBody", NSMAP):
                for lstStyle in txBody.findall(".//a:lstStyle", NSMAP):
                    defPPr = lstStyle.find(".//a:defPPr", NSMAP)
                    if defPPr is not None:
                        defRPr = defPPr.find("a:defRPr", NSMAP)
                        if defRPr is not None:
                            if font_info["size_pt"] is None:
                                sz = defRPr.get("sz")
                                if sz:
                                    font_info["size_pt"] = int(sz) // 100
                            if font_info["family"] is None:
                                latin = defRPr.find("a:latin", NSMAP)
                                if latin is not None:
                                    font_info["family"] = latin.get("typeface")

        result[type_label] = font_info

    return result


def extract_layout_info(prs: Presentation) -> list[dict]:
    """Extract all layout names, placeholder info, and background colors.

    Returns:
        List of dicts, one per layout, with keys:
        index, name, hint, bg_color, placeholders (from extract_placeholder_fonts).
    """
    layouts = []
    for idx, layout in enumerate(prs.slide_layouts):
        name = layout.name or f"Layout {idx}"
        hint = LAYOUT_HINTS.get(name, "Custom layout")

        # Try to extract background color
        bg_color = _extract_background_color(layout)

        ph_fonts = extract_placeholder_fonts(layout)

        layouts.append(
            {
                "index": idx,
                "name": name,
                "hint": hint,
                "bg_color": bg_color,
                "placeholders": ph_fonts,
            }
        )

    return layouts


# ---------------------------------------------------------------------------
# Private helpers
# ---------------------------------------------------------------------------


def _ph_type_label(ph_type, idx: int) -> str:
    """Convert a placeholder type enum to a human-readable string."""
    type_map = {
        1: "title",
        2: "body",
        3: "center_title",
        4: "subtitle",
        5: "body",
        6: "title",
        7: "date",
        8: "footer",
        9: "slide_number",
        10: "header",
        11: "media",
        12: "picture",
        13: "graphic_frame",
        14: "table",
        15: "chart",
    }
    try:
        type_int = int(ph_type)
    except (TypeError, ValueError):
        type_int = -1

    label = type_map.get(type_int, f"placeholder_{idx}")
    # Disambiguate multiple body placeholders
    if label == "body" and idx > 1:
        label = f"body_{idx}"
    return label


def _extract_background_color(layout) -> str:
    """Attempt to extract a solid background fill color from a slide layout."""
    try:
        bg = layout.background
        fill = bg.fill
        if fill.type is not None:
            fg = fill.fore_color
            return str(fg.rgb)
    except Exception:
        pass
    # Try XML path
    try:
        bg_el = layout._element.find(".//p:bg", NSMAP)
        if bg_el is not None:
            srgb = bg_el.find(".//a:srgbClr", NSMAP)
            if srgb is not None:
                return (srgb.get("val") or "").upper()
    except Exception:
        pass
    return ""


# ---------------------------------------------------------------------------
# Markdown section builders
# ---------------------------------------------------------------------------


def _build_dimensions_section(dims: dict[str, float]) -> str:
    width_in = dims["width_in"]
    height_in = dims["height_in"]
    width_emu = dims["width_emu"]
    height_emu = dims["height_emu"]

    # Compute aspect ratio
    from math import gcd

    w = int(round(width_in * 100))
    h = int(round(height_in * 100))
    d = gcd(w, h)
    ar_w, ar_h = w // d, h // d

    return (
        "## Slide Dimensions\n"
        f"- Width: {width_in:.2f} inches ({width_emu} EMU)\n"
        f"- Height: {height_in:.2f} inches ({height_emu} EMU)\n"
        f"- Aspect ratio: {ar_w}:{ar_h}\n"
    )


def _build_colors_section(theme: dict[str, str]) -> str:
    if not theme:
        return "## Color Palette\n\n_No theme colors found._\n"

    lines = [
        "## Color Palette",
        "",
        "| Role | Hex | Usage |",
        "|------|-----|-------|",
    ]
    for role, hex_val in theme.items():
        usage = COLOR_USAGE_HINTS.get(role, "")
        lines.append(f"| {role} | #{hex_val} | {usage} |")
    lines.append("")
    return "\n".join(lines)


def _build_fonts_section(fonts: dict[str, str]) -> str:
    return (
        "## Font Scheme\n"
        f"- **Heading font**: {fonts['major']} (use for all slide titles and headers)\n"
        f"- **Body font**: {fonts['minor']} (use for all body text, bullets, captions)\n"
    )


def _build_typography_section(layouts: list[dict], fonts: dict[str, str]) -> str:
    """Build a typography scale table from common placeholder font data."""
    # Collect representative entries from the first few layouts
    seen: dict[str, dict] = {}
    for layout in layouts:
        for ph_type, info in layout["placeholders"].items():
            if ph_type not in seen:
                seen[ph_type] = {**info, "layout_name": layout["name"]}

    if not seen:
        return (
            "## Typography Scale\n\n"
            f"- **Headings**: {fonts['major']}\n"
            f"- **Body**: {fonts['minor']}\n"
        )

    lines = [
        "## Typography Scale",
        "",
        "| Element | Font Family | Size | Bold | Color |",
        "|---------|-------------|------|------|-------|",
    ]
    for ph_type, info in seen.items():
        family = info.get("family") or fonts.get(
            "major" if "title" in ph_type else "minor", ""
        )
        size = f"{info['size_pt']}pt" if info.get("size_pt") else "—"
        bold = "Yes" if info.get("bold") else "No"
        color = f"#{info['color_hex']}" if info.get("color_hex") else "—"
        lines.append(
            f"| {ph_type.replace('_', ' ').title()} | {family} | {size} | {bold} | {color} |"
        )

    lines.append("")
    return "\n".join(lines)


def _build_layouts_overview_section(layouts: list[dict]) -> str:
    lines = [
        "## Available Slide Layouts",
        "",
        "| Index | Name | Best For |",
        "|-------|------|----------|",
    ]
    for layout in layouts:
        lines.append(f"| {layout['index']} | {layout['name']} | {layout['hint']} |")
    lines.append("")
    return "\n".join(lines)


def _build_layout_details_section(layouts: list[dict], fonts: dict[str, str]) -> str:
    lines = ["## Layout Details", ""]
    for layout in layouts:
        lines.append(f"### Layout {layout['index']}: {layout['name']}")
        if layout["bg_color"]:
            lines.append(f"- Background: #{layout['bg_color']}")
        else:
            lines.append("- Background: (inherits from slide master)")

        phs = layout["placeholders"]
        if phs:
            lines.append("- Placeholders:")
            for ph_type, info in phs.items():
                family = info.get("family") or fonts.get(
                    "major" if "title" in ph_type else "minor", ""
                )
                parts = [ph_type.replace("_", " ").title()]
                if family:
                    parts.append(f"{family}")
                if info.get("size_pt"):
                    parts.append(f"{info['size_pt']}pt")
                if info.get("bold"):
                    parts.append("Bold")
                if info.get("color_hex"):
                    parts.append(f"#{info['color_hex']}")
                lines.append(f"  - {', '.join(parts)}")
        else:
            lines.append("- Placeholders: (none detected)")
        lines.append("")

    return "\n".join(lines)


def build_pptxgenjs_directives(
    theme: dict[str, str],
    fonts: dict[str, str],
    dims: dict[str, float],
    layouts: list[dict],
) -> str:
    """Build a clear PptxGenJS-specific instruction block for Claude.

    Tells Claude exactly how to apply the template's styling in PptxGenJS code.

    Args:
        theme: Dict of role->hex from extract_theme_colors.
        fonts: Dict with "major" and "minor" keys from extract_font_scheme.
        dims: Dict from extract_slide_dimensions.
        layouts: List of layout dicts from extract_layout_info.

    Returns:
        Markdown string with numbered PptxGenJS directives.
    """
    major_font = fonts.get("major", "Calibri Light")
    minor_font = fonts.get("minor", "Calibri")
    primary_hex = theme.get("Accent 1 (Primary)", "4472C4")
    title_color = theme.get("Dark 2", theme.get("Dark 1 (Text)", "000000"))
    body_color = theme.get("Dark 1 (Text)", "000000")
    bg_color = theme.get("Light 1 (Background)", "FFFFFF")
    width_emu = dims["width_emu"]
    height_emu = dims["height_emu"]
    width_in = dims["width_in"]
    height_in = dims["height_in"]

    lines = [
        "## Design Guidelines for PptxGenJS",
        "",
        "When generating the PPTX using PptxGenJS, apply these exact specifications:",
        "",
        f'1. Use font family `"{major_font}"` for all slide titles '
        f'(`slide.addText(title, {{ fontFace: "{major_font}" }})`)',
        f'2. Use font family `"{minor_font}"` for all body text and bullet points',
        f"3. Apply color `#{primary_hex}` as the primary accent for shapes and highlights",
        f'4. Set slide dimensions: width {width_emu} EMU ({width_in}"), '
        f'height {height_emu} EMU ({height_in}")',
        f"5. For title text color use `#{title_color}`; for body text use `#{body_color}`",
        f"6. Use `#{bg_color}` as the default slide background",
        "7. Match chart series colors to the accent palette in order:",
    ]

    accent_order = [
        "Accent 1 (Primary)",
        "Accent 2 (Secondary)",
        "Accent 3",
        "Accent 4",
        "Accent 5",
        "Accent 6",
    ]
    accent_list = [f"`#{theme[k]}`" for k in accent_order if k in theme]
    if accent_list:
        lines.append(f"   {', '.join(accent_list)}")

    lines.append(
        "8. For layout selection, use the layout index that matches the slide purpose "
        "(see Available Slide Layouts table above)"
    )
    lines.append("")
    lines.append("### PptxGenJS Snippet Reference")
    lines.append("")
    lines.append("```javascript")
    lines.append("// Title slide")
    lines.append("const slide = prs.addSlide();")
    lines.append(f'slide.addText("Presentation Title", {{')
    lines.append(f"    x: 0.5, y: 2.0, w: 9.0, h: 1.5,")
    lines.append(f'    fontFace: "{major_font}",')
    lines.append(f"    fontSize: 40,")
    lines.append(f'    color: "{title_color}",')
    lines.append(f'    align: "center"')
    lines.append(f"}});")
    lines.append("")
    lines.append("// Body text")
    lines.append(f'slide.addText("Bullet point", {{')
    lines.append(f"    x: 0.5, y: 1.5, w: 9.0, h: 4.0,")
    lines.append(f'    fontFace: "{minor_font}",')
    lines.append(f"    fontSize: 18,")
    lines.append(f'    color: "{body_color}",')
    lines.append(f"    bullet: true")
    lines.append(f"}});")
    lines.append("")
    lines.append("// Accent shape")
    lines.append(f"slide.addShape(prs.ShapeType.rect, {{")
    lines.append(f"    x: 0, y: 0, w: 10, h: 0.1,")
    lines.append(f'    fill: {{ color: "{primary_hex}" }}')
    lines.append(f"}});")
    lines.append("```")
    lines.append("")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Main public function
# ---------------------------------------------------------------------------


def extract_template_to_markdown(template_path: "str | Path") -> str:
    """Parse a PPTX template and return a comprehensive markdown styling guide.

    The returned string can be injected directly into Claude's prompt to give
    it precise instructions for matching the template's visual design when
    generating presentations via PptxGenJS.

    Args:
        template_path: Path to the .pptx template file.

    Returns:
        Markdown-formatted string with all styling information.
    """
    prs = Presentation(str(template_path))

    theme = extract_theme_colors(prs)
    fonts = extract_font_scheme(prs)
    dims = extract_slide_dimensions(prs)
    layouts = extract_layout_info(prs)

    sections = [
        "# Template Styling Guide",
        "",
        _build_dimensions_section(dims),
        _build_colors_section(theme),
        _build_fonts_section(fonts),
        _build_typography_section(layouts, fonts),
        _build_layouts_overview_section(layouts),
        _build_layout_details_section(layouts, fonts),
        build_pptxgenjs_directives(theme, fonts, dims, layouts),
    ]

    return "\n".join(sections)


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Extract PPTX template styling as a markdown guide"
    )
    parser.add_argument("template", help="Path to the .pptx template file")
    parser.add_argument(
        "--output",
        "-o",
        help="Save markdown to this file instead of printing to stdout",
    )
    args = parser.parse_args()

    markdown = extract_template_to_markdown(args.template)

    if args.output:
        Path(args.output).write_text(markdown, encoding="utf-8")
        print(f"Styling guide saved to {args.output}")
    else:
        print(markdown)
