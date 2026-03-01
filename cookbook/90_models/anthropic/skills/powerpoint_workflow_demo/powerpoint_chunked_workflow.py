"""
Agno Workflow: Chunked PowerPoint Generation Pipeline.

A chunked workflow that overcomes Claude API limitations for large presentations
by splitting generation into manageable chunks, then merging the results.

Problem solved: Single Claude API calls fail for 10+ slide presentations;
               Claude PPTX skill is also prone to throttling and timeouts.
Solution: Generate slides in configurable chunks (default: 3 slides per call),
          then merge all chunks into one final presentation.
          A 3-tier fallback ensures production reliability when the primary
          Claude PPTX skill is unavailable or too slow.

Architecture:
  This file is a thin orchestration layer built on top of powerpoint_template_workflow.py.
  It imports all helpers, agents, Pydantic models, and step functions from that file via
  a wildcard import, then adds the chunked orchestration logic on top.

  powerpoint_template_workflow.py  — Core pipeline: content gen, images, template assembly,
                                    visual review, all helper functions (~6500 lines)
  powerpoint_chunked_workflow.py   — Chunked orchestration layer (~2600 lines, this file)

Chunk generation uses a 3-tier fallback hierarchy per chunk:
  Tier 1  Claude PPTX Skill    - Primary; native charts, tables, rich visuals
  Tier 2  LLM Code Generation  - Fallback; LLM writes + executes python-pptx
                                  code (native charts only); 80-92% quality parity
  Tier 3  python-pptx Direct   - Last resort; text-only slides; 100% reliable

Relationship between the two files:
  - powerpoint_template_workflow.py is self-contained and can run standalone (single Claude
    API call, suitable for short presentations of up to ~7 slides).
  - powerpoint_chunked_workflow.py wraps the same template/image/review logic via wildcard
    import so that large presentations (8-15+ slides) are split into chunks and merged.
  - Do NOT modify powerpoint_template_workflow.py to add chunking logic; keep them separate.

Workflow steps:
  Step 1  Optimize & Plan    - LLM analyzes prompt, decides slide count, creates storyboard
  Step 2  Generate Chunks    - Call Claude pptx skill (Tier 1) for each chunk;
                               auto-escalates to Tier 2 (LLM code gen) on timeout/
                               failure, then Tier 3 (text-only) if Tier 2 fails.
  Step 3  Process Chunks     - Apply template + image pipeline per chunk (if template provided)
  Step 4  Visual Review      - Optional per-chunk visual QA (if --visual-review + template)
  Step 5  Merge Chunks       - Merge all processed chunks into the final PPTX

Prerequisites:
- uv pip install agno anthropic python-pptx google-genai pillow
- export ANTHROPIC_API_KEY="your_api_key_here"
- export GOOGLE_API_KEY="your_google_api_key_here" (for image generation)
- A .pptx template file (optional)

Usage:
    # Basic usage (auto-decide slide count, 3 slides per chunk):
    .venvs/demo/bin/python powerpoint_chunked_workflow.py \\
        -p "Create a presentation about AI in healthcare"

    # With template, 4 slides per chunk:
    .venvs/demo/bin/python powerpoint_chunked_workflow.py \\
        -t my_template.pptx --chunk-size 4

    # Large presentation with visual review (5 passes max):
    .venvs/demo/bin/python powerpoint_chunked_workflow.py \\
        -t my_template.pptx -p "12-slide enterprise AI strategy deck" \\
        --chunk-size 3 --visual-review --visual-passes 5

    # Quick generation without images or template:
    .venvs/demo/bin/python powerpoint_chunked_workflow.py \\
        -p "Startup pitch deck for SaaS product" --no-images

CLI Flags:
    --template, -t       Path to .pptx template (optional). Without it, skips
                         template assembly and visual review; just merges raw chunks.
    --output, -o         Output filename (default: presentation_chunked.pptx).
    --prompt, -p         User prompt describing the presentation.
    --no-images          Skip AI image generation.
    --no-stream          Disable streaming mode for Claude agent.
    --min-images         Minimum slides that must have images (default: 1).
    --visual-review      Enable visual QA with Gemini vision per chunk.
    --footer-text        Footer text for all slides.
    --date-text          Date text for footer date placeholder.
    --show-slide-numbers Preserve slide number placeholder on all slides.
    --verbose, -v        Enable verbose/debug logging.
    --chunk-size         Number of slides per Claude API chunk call (default: 3).
    --max-retries        Max retries per chunk on failure (default: 2).
    NOTE: When all retries fail or a timeout (300s) occurs, the system
          automatically switches to Tier 2 (LLM code gen) fallback,
          then Tier 3 (text-only) if Tier 2 also fails.
    --visual-passes      Maximum visual inspection passes per chunk (default: 3).
    --start-tier         Starting tier for chunk generation (default: 1):
                         1 = Claude PPTX skill (best quality, native charts/tables)
                         2 = LLM code generation (80-92% quality, faster, python-pptx native charts)
                         3 = Text-only (structural only, instant, no API calls)
                         Fallback chain continues from selected tier (e.g., Tier 2 → Tier 3).

Logging conventions:
    Always printed:
        [STEP_NAME] Starting / result messages
        [TIMING] step_XXX completed in X.Xs
        [ERROR] ...
        [WARNING] ...
        [VISUAL REVIEW MISSING FIX] ...  (always, per spec)
    Verbose-only (requires --verbose / -v):
        [VERBOSE] detailed debug information
"""

import argparse
import concurrent.futures
import copy
import json
import os
import shutil
import sys
import time
import traceback
import zipfile
from pathlib import Path
from typing import Dict, List, Optional

from agno.run.agent import RunOutput

# === WILDCARD IMPORT: Reuse all helpers, agents, models, and step functions ===
# This gives us access to all ~6500 lines of helper logic without duplication.
# Specifically imports: SlideImageDecision, ImagePlan, ShapeIssue, SlideQualityReport,
# PresentationQualityReport, all dataclasses, image_planner, slide_quality_reviewer,
# step_plan_images, step_generate_images, step_assemble_template, step_visual_quality_review,
# and all _helper functions.
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from agno.agent import Agent
from agno.models.anthropic import Claude
from agno.tools.python import PythonTools
from agno.workflow.step import Step
from agno.workflow.types import StepInput, StepOutput
from agno.workflow.workflow import Workflow
from anthropic import Anthropic
from file_download_helper import download_skill_files
from powerpoint_template_workflow import *  # noqa: F401, F403, E402
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field

# === NEW PYDANTIC MODELS FOR CHUNKED WORKFLOW ===


class SlideStoryboard(BaseModel):
    """Storyboard entry for a single slide in the presentation."""

    slide_number: int = Field(..., description="1-based slide number")
    slide_title: str = Field(..., description="Title for this slide")
    slide_type: str = Field(..., description="Type: title/agenda/content/data/closing")
    key_points: List[str] = Field(..., description="3-5 bullet points for this slide")
    visual_suggestion: str = Field(..., description="Visual element recommendation")
    transition_note: str = Field(..., description="How this slide connects to next")


class StoryboardPlan(BaseModel):
    """Complete storyboard plan for the presentation, produced by the query optimizer."""

    total_slides: int = Field(..., description="Total number of slides")
    presentation_title: str = Field(..., description="Main presentation title")
    search_topic: str = Field(
        "",
        description=(
            "Primary web research topic phrase extracted from user request; used to guide searches"
        ),
    )
    target_audience: str = Field(..., description="Target audience")
    tone: str = Field(
        ...,
        description="Presentation tone, e.g. 'professional', 'inspiring', or 'technical'.",
    )
    brand_voice: str = Field(
        ...,
        description="Brand voice style, e.g. 'authoritative', 'conversational', or 'data-driven'.",
    )
    global_context: str = Field(
        ...,
        description="2-3 sentence shared context covering company, product, and central theme.",
    )
    slides: List[SlideStoryboard] = Field(
        ..., description="Per-slide storyboard entries"
    )


# === MODULE-LEVEL AGENTS ===
# Do NOT create agents in loops — define them here at module level.

# output_schema is intentionally omitted: claude-opus-4-6 does not support structured
# outputs, which causes Agno to make an internal non-streaming extraction call that the
# context-1m beta rejects ("Streaming is required for operations that may take longer
# than 10 minutes").  The storyboard JSON is instead requested via prompt instructions
# and parsed manually below.
query_optimizer = Agent(
    name="Presentation Strategist",
    model=Claude(
        id="claude-opus-4-6", betas=["context-1m-2025-08-07"], max_tokens=128000
    ),
    description=(
        "You are a presentation strategist who first searches the web for current, "
        "relevant facts and data about the topic, then creates an optimized presentation "
        "plan with a per-slide storyboard grounded in that research."
    ),
    tools=[
        {
            "type": "web_search_20250305",
            "name": "web_search",
            "max_uses": 5,
        }
    ],
    markdown=False,
)


# === HELPER: STORYBOARD MARKDOWN FORMATTING ===


def _format_slide_markdown(slide: SlideStoryboard) -> str:
    """Format a SlideStoryboard as a markdown string for the pptx agent.

    Excludes transition_note (planning meta-info, not useful for content generation)
    to keep context size lean. Includes type, key points, and visual suggestion only.

    Args:
        slide: SlideStoryboard instance to format.

    Returns:
        Markdown string with slide number, title, type, key points, and visual suggestion.
    """
    points = "\n".join("- %s" % p for p in slide.key_points)
    return ("# Slide %d: %s\n\n**Type:** %s\n\n## Content\n%s\n\n**Visual:** %s\n") % (
        slide.slide_number,
        slide.slide_title,
        slide.slide_type,
        points,
        slide.visual_suggestion,
    )


def _format_global_context_markdown(plan: StoryboardPlan) -> str:
    """Format the global context as a markdown string for the pptx agent.

    Kept concise: title, audience, tone, brand voice, and the 2-3 sentence global context.
    This file is included in every chunk prompt so brevity matters for context size.

    Args:
        plan: StoryboardPlan instance containing global presentation metadata.

    Returns:
        Markdown string with presentation title, audience, tone, brand voice, and context.
    """
    return (
        "# Presentation: %s\n\n"
        "Audience: %s | Tone: %s | Brand Voice: %s\n\n"
        "## Context\n%s\n"
    ) % (
        plan.presentation_title,
        plan.target_audience,
        plan.tone,
        plan.brand_voice,
        plan.global_context,
    )


# === HELPER: SAVE PROMPT TO FILE ===


def _save_prompt_to_file(
    prompt: str, step_name: str, output_dir: str, extra: str = ""
) -> str:
    """Save a prompt string to a timestamped .txt file inside output_dir.

    Files are written to output_dir directly (which is output_chunked/chunked_workflow_work/).
    Filenames follow the pattern: prompt_<step_name>[_<extra>]_<timestamp_ms>.txt

    Args:
        prompt: The full prompt text to save.
        step_name: Short identifier for the workflow step (e.g. 'optimize_and_plan', 'chunk').
        output_dir: Directory in which to write the file.
        extra: Optional extra qualifier appended between step_name and timestamp.

    Returns:
        Absolute path of the saved file, or empty string if saving failed.
    """
    os.makedirs(output_dir, exist_ok=True)
    timestamp_ms = int(time.time() * 1000)
    if extra:
        filename = "prompt_%s_%s_%d.txt" % (step_name, extra, timestamp_ms)
    else:
        filename = "prompt_%s_%d.txt" % (step_name, timestamp_ms)
    filepath = os.path.join(output_dir, filename)
    try:
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(prompt)
    except Exception as e:
        print("[WARNING] Failed to save prompt file %s: %s" % (filepath, e))
        return ""
    return filepath


# === HELPER: EXTRACT SLIDES DATA FROM A CHUNK PPTX ===


def _extract_chunk_slides_data(chunk_file: str) -> List[dict]:
    """Extract basic slide metadata from a PPTX chunk file.

    Returns a list of dicts compatible with the session_state['slides_data'] format
    used by step_plan_images and step_generate_images.
    """
    slides_data = []
    try:
        prs = Presentation(chunk_file)
        for idx, slide in enumerate(prs.slides):
            slide_info: dict = {
                "index": idx,
                "title": "",
                "body": "",
                "has_table": False,
                "has_chart": False,
                "has_image": False,
            }
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for ph_attr in ["placeholder_format"]:
                        ph = getattr(shape, ph_attr, None)
                        if ph and hasattr(ph, "idx") and ph.idx == 0:
                            slide_info["title"] = shape.text_frame.text.strip()
                if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                    slide_info["has_table"] = True
                if shape.shape_type == 3:  # CHART
                    slide_info["has_chart"] = True
            slides_data.append(slide_info)
    except Exception as e:
        print("[WARNING] Could not extract slides data from %s: %s" % (chunk_file, e))
    return slides_data


# === WORKFLOW STEP 1: OPTIMIZE AND PLAN ===


def step_optimize_and_plan(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 1: Enhance the user prompt, decide slide count, and generate a per-slide storyboard.

    Uses the query_optimizer agent (Claude Opus) to produce a StoryboardPlan with:
    - Optimal slide count (respects user-specified count, otherwise picks 8-15)
    - Global context applicable to all slides
    - Per-slide storyboard with title, type, key points, visual suggestions
    - Presentation tone and brand voice

    Saves storyboard to individual markdown files in {output_dir}/storyboard/.

    Args:
        step_input: Workflow step input (not used directly; context comes from session_state).
        session_state: Shared workflow state containing user_prompt, output_dir, chunk_size,
                       and max_retries.

    Returns:
        StepOutput with success=True and a summary string when a valid storyboard is produced,
        or success=False with an error message if the optimizer fails or returns invalid JSON.
    """
    step_start = time.time()

    user_prompt = session_state.get("user_prompt", "")
    output_dir = session_state.get("output_dir", ".")
    chunk_size = session_state.get("chunk_size", 3)
    max_retries = session_state.get("max_retries", 2)

    print("=" * 60)
    print("Step 1: Optimizing query and generating storyboard...")
    print("=" * 60)
    print("User prompt: %s" % user_prompt[:200])

    storyboard_dir = os.path.join(output_dir, "storyboard")
    os.makedirs(storyboard_dir, exist_ok=True)

    optimizer_prompt = (
        "Analyze the following user request for a PowerPoint presentation and create an optimized storyboard.\n\n"
        "USER REQUEST:\n%s\n\n"
        "STEP 1 — RESEARCH FIRST:\n"
        "STEP 1A — DEFINE SEARCH TOPIC (MANDATORY, DO THIS FIRST):\n"
        "Extract and define ONE clear Search Topic from the user request before calling web_search.\n"
        "Search Topic should be a short phrase (6-14 words) that captures the exact core focus.\n"
        "Example format (for your internal planning only):\n"
        "Search Topic: 'AI adoption trends in US healthcare providers (2023-2025)'\n"
        "Then generate 4 focused search queries from that Search Topic:\n"
        "- Query 1: market size or growth\n"
        "- Query 2: adoption or usage rates\n"
        "- Query 3: key challenges or risks\n"
        "- Query 4: case studies or real-world examples\n"
        "Add geography, industry, and year constraints when relevant.\n\n"
        "STEP 1B — RUN WEB SEARCH USING THAT SEARCH TOPIC:\n"
        "Before planning slides, use web_search with those queries to find 2-4 relevant facts, statistics, "
        "or examples for the Search Topic. Prioritize recent, credible sources and specific numbers "
        "(e.g., market size, CAGR, adoption rates, trend changes). Use these findings to ground the storyboard "
        "in real, specific information. For each fact you use, internally track source name + publication year "
        "and prefer the newest credible source when sources conflict.\n"
        "If web_search returns weak or conflicting data, do NOT invent facts. Use conservative language "
        "and reduce numeric specificity rather than fabricating values.\n\n"
        "STEP 2 — BUILD THE STORYBOARD:\n"
        "1. If the user specifies a slide count (e.g. '12 slides', '10-slide deck'), honor it exactly.\n"
        "2. If not specified, decide the optimal count: typically 8-15 slides for professional decks.\n"
        "   - Simple topics: 8-10 slides. Complex/technical topics: 12-15 slides.\n"
        "   - Do NOT add unnecessary slides; quality over quantity.\n"
        "3. Define a clear tone and brand voice appropriate to the topic.\n"
        "4. Write global_context as 2-3 focused sentences covering: the core topic, target audience, "
        "and central theme or key message. Include a specific fact or statistic from your research.\n"
        "5. For each slide, provide:\n"
        "   - A concise, descriptive title (5-8 words max)\n"
        "   - slide_type: one of title, agenda, content, data, closing\n"
        "   - 3-4 key_points: each a single sentence (10-20 words), specific and actionable. "
        "     Include real data or examples where relevant. Avoid vague filler bullets.\n"
        "   - visual_suggestion: one concrete line — specify chart type + data, image concept, or 'none'. "
        "     Example: 'bar chart: AI adoption rate by industry 2023' not just 'chart'.\n"
        "   - transition_note: one brief sentence connecting this slide to the next.\n"
        "6. Ensure continuity: the storyboard should feel like a coherent narrative arc.\n"
        "7. Use professional language. Do not add emojis or overly casual language.\n\n"
        "STEP 3 — OUTPUT FORMAT:\n"
        "Respond with ONLY a valid JSON object matching this exact schema (no markdown fences, "
        "no extra commentary before or after the JSON).\n"
        "The JSON must be syntactically valid and parseable with strict JSON parsers.\n"
        "Ensure slides length == total_slides and slide_number values are contiguous from 1..total_slides:\n"
        "{\n"
        '  "total_slides": <integer>,\n'
        '  "presentation_title": "<string>",\n'
        '  "search_topic": "<string>",\n'
        '  "target_audience": "<string>",\n'
        '  "tone": "<string>",\n'
        '  "brand_voice": "<string>",\n'
        '  "global_context": "<string>",\n'
        '  "slides": [\n'
        "    {\n"
        '      "slide_number": <integer>,\n'
        '      "slide_title": "<string>",\n'
        '      "slide_type": "<title|agenda|content|data|closing>",\n'
        '      "key_points": ["<string>", ...],\n'
        '      "visual_suggestion": "<string>",\n'
        '      "transition_note": "<string>"\n'
        "    },\n"
        "    ...\n"
        "  ]\n"
        "}\n"
    ) % user_prompt

    prompt_file = _save_prompt_to_file(
        optimizer_prompt, "optimize_and_plan", output_dir
    )
    if prompt_file:
        print("[PROMPT] Optimizer prompt saved to: %s" % prompt_file)

    try:
        response = None
        for event in query_optimizer.run(
            optimizer_prompt, stream=True, yield_run_output=True
        ):
            if isinstance(event, RunOutput):
                response = event
    except Exception as e:
        print("[ERROR] Query optimizer failed: %s" % str(e))
        traceback.print_exc()
        return StepOutput(
            content="Query optimization failed: %s" % str(e), success=False
        )

    # Parse the StoryboardPlan from response.
    # Without output_schema the agent returns plain text; extract JSON from it.
    plan: Optional[StoryboardPlan] = None
    if response and response.content:
        content = response.content
        if isinstance(content, StoryboardPlan):
            plan = content
        elif isinstance(content, dict):
            try:
                plan = StoryboardPlan(**content)
            except Exception as e:
                print("[ERROR] Failed to parse StoryboardPlan from dict: %s" % e)
        elif isinstance(content, str):
            # Strip markdown code fences if present (```json ... ``` or ``` ... ```)
            import re as _re

            json_text = content.strip()
            fence_match = _re.search(r"```(?:json)?\s*([\s\S]+?)```", json_text)
            if fence_match:
                json_text = fence_match.group(1).strip()
            # Locate the outermost JSON object in the text
            obj_match = _re.search(r"\{[\s\S]+\}", json_text)
            if obj_match:
                json_text = obj_match.group(0)
            try:
                plan = StoryboardPlan.model_validate_json(json_text)
            except Exception as e:
                print("[ERROR] Failed to parse StoryboardPlan from JSON string: %s" % e)
                if VERBOSE:  # noqa: F405
                    print(
                        "[VERBOSE] Raw optimizer response (first 2000 chars):\n%s"
                        % content[:2000]
                    )

    if not plan:
        print("[ERROR] No valid storyboard plan produced.")
        return StepOutput(content="No storyboard plan produced.", success=False)

    print(
        "Storyboard plan: '%s' (%d slides, tone: %s)"
        % (plan.presentation_title, plan.total_slides, plan.tone)
    )

    if VERBOSE:  # noqa: F405
        print("[VERBOSE] Full storyboard JSON:\n%s" % plan.model_dump_json(indent=2))

    # Save global context markdown
    global_context_path = os.path.join(storyboard_dir, "global_context.md")
    with open(global_context_path, "w", encoding="utf-8") as f:
        f.write(_format_global_context_markdown(plan))
    print("Saved global context: %s" % global_context_path)

    # Save per-slide storyboard markdown files
    for slide in plan.slides:
        slide_path = os.path.join(storyboard_dir, "slide_%03d.md" % slide.slide_number)
        slide_md = _format_slide_markdown(slide)
        if VERBOSE:  # noqa: F405
            print("[VERBOSE] Slide %d storyboard:\n%s" % (slide.slide_number, slide_md))
        with open(slide_path, "w", encoding="utf-8") as f:
            f.write(slide_md)
    print("Saved %d slide storyboard files to: %s" % (len(plan.slides), storyboard_dir))

    # Store in session_state
    session_state["storyboard"] = plan
    session_state["total_slides"] = plan.total_slides
    session_state["storyboard_dir"] = storyboard_dir
    session_state["chunk_size"] = chunk_size
    session_state["max_retries"] = max_retries

    step_elapsed = time.time() - step_start
    print("[TIMING] step_optimize_and_plan completed in %.1fs" % step_elapsed)

    summary = (
        "Storyboard created: '%s' | %d slides | tone: %s | brand voice: %s | chunk size: %d | Duration: %.1fs"
    ) % (
        plan.presentation_title,
        plan.total_slides,
        plan.tone,
        plan.brand_voice,
        chunk_size,
        step_elapsed,
    )
    return StepOutput(content=summary, success=True)


CHUNK_TIMEOUT_SECONDS = 300  # 5 minutes


def _run_chunk_agent(chunk_agent, chunk_prompt):
    """Run the chunk agent with streaming and collect the final RunOutput.

    Intended to be submitted to a ThreadPoolExecutor so the caller can enforce a
    wall-clock timeout via Future.result(timeout=...).

    Args:
        chunk_agent: Configured Agent instance to run.
        chunk_prompt: Full prompt string to send to the agent.

    Returns:
        Tuple of (RunOutput or None, int) where the second element is the total
        number of streaming events received. RunOutput is None if no RunOutput
        event was emitted during the stream.
    """
    response = None
    event_count = 0
    for event in chunk_agent.run(chunk_prompt, stream=True, yield_run_output=True):
        event_count += 1
        if isinstance(event, RunOutput):
            response = event
    return response, event_count


# === HELPER: GENERATE A SINGLE CHUNK VIA CLAUDE PPTX SKILL ===


def generate_chunk_pptx(
    chunk_slides: List[SlideStoryboard],
    session_state: Dict,
    chunk_idx: int,
) -> Optional[str]:
    """Call the Claude pptx skill for a chunk of slides with retry logic (Tier 1).

    Creates a fresh agent per call (not reused across chunks) and applies
    exponential backoff on retries.

    Applies a 300-second (CHUNK_TIMEOUT_SECONDS) wall-clock timeout per attempt
    via ThreadPoolExecutor. On timeout, activates the session-level fallback flag
    and returns None immediately (no further retries for this chunk).

    Args:
        chunk_slides: List of SlideStoryboard objects for this chunk.
        session_state: Shared workflow session state.
        chunk_idx: 0-based chunk index (used for file naming and logging).

    Returns:
        Path to the generated chunk PPTX file, or None if all attempts failed.
        When None is returned, session_state["use_fallback_generator"] is set to
        True, causing all subsequent chunks to bypass Tier 1 and use Tier 2/3.
    """
    storyboard: StoryboardPlan = session_state["storyboard"]
    storyboard_dir = session_state["storyboard_dir"]
    output_dir = session_state["output_dir"]
    max_retries = session_state.get("max_retries", 2)

    # Load global context
    global_context_path = os.path.join(storyboard_dir, "global_context.md")
    global_context = ""
    if os.path.exists(global_context_path):
        with open(global_context_path, encoding="utf-8") as f:
            global_context = f.read()

    # Load per-slide markdown for this chunk
    slide_details = []
    for s in chunk_slides:
        md_path = os.path.join(storyboard_dir, "slide_%03d.md" % s.slide_number)
        if os.path.exists(md_path):
            with open(md_path, encoding="utf-8") as f:
                slide_details.append(f.read())
        else:
            # Fallback: format inline
            slide_details.append(_format_slide_markdown(s))

    first_slide = chunk_slides[0].slide_number
    last_slide = chunk_slides[-1].slide_number

    chunk_prompt = (
        "## Global Presentation Context\n"
        "%s\n\n"
        "## Task: Generate slides %d through %d of %d\n\n"
        "You are generating a CHUNK of a larger presentation. "
        "This chunk contains %d slides.\n"
        "Maintain the presentation's tone (%s) and brand voice (%s).\n"
        'These are slides %d-%d of the full %d-slide deck titled "%s".\n\n'
        "## Per-Slide Content for This Chunk:\n\n"
        "%s\n\n"
        "Please generate EXACTLY %d slides for this chunk with the content described above.\n"
        "Do not add extra slides. Do not include slide numbers outside the range %d-%d.\n"
        "Use clean formatting without custom fonts or colors. "
        "Include native data-vis (tables/charts/infographics/diagrams) only where explicitly suggested.\n"
        "For any chart: use native PPTX chart objects only (bar, column, line, or pie) — "
        "do NOT use matplotlib or embed chart images.\n"
        "For any table: use native PPTX table objects only — do NOT embed a table as an image "
        "or use matplotlib/PIL to render one.\n"
        "For infographics or diagrams: use native PowerPoint shapes (rectangles, arrows, text boxes) "
        "or a native table to approximate the visual — do NOT insert images for infographics or diagrams.\n"
        "If a requested visual cannot be represented exactly, preserve the slide structure and add a concise "
        "native textbox note; do NOT fail the chunk and do NOT use image-based substitutes.\n"
        "Save the output as 'chunk_%03d.pptx'."
    ) % (
        global_context,
        first_slide,
        last_slide,
        storyboard.total_slides,
        len(chunk_slides),
        storyboard.tone,
        storyboard.brand_voice,
        first_slide,
        last_slide,
        storyboard.total_slides,
        storyboard.presentation_title,
        "\n\n---\n\n".join(slide_details),
        len(chunk_slides),
        first_slide,
        last_slide,
        chunk_idx,
    )

    prompt_file = _save_prompt_to_file(
        chunk_prompt, "chunk", output_dir, "chunk_%03d" % chunk_idx
    )
    if prompt_file:
        print("[PROMPT] Chunk %d prompt saved to: %s" % (chunk_idx, prompt_file))

    chunk_output_path = os.path.join(output_dir, "chunk_%03d.pptx" % chunk_idx)

    # Create a fresh agent per chunk call — do NOT reuse across calls.
    # betas + max_tokens=128000: the PPTX skill generates full multi-slide decks whose
    # output can be large; 128k tokens and the context-1m beta are both safe here because
    # this agent is always invoked with stream=True (see _run_chunk_agent).
    chunk_agent = Agent(
        name="Chunk Generator %d" % chunk_idx,
        model=Claude(
            id="claude-opus-4-6",
            betas=["context-1m-2025-08-07"],
            skills=[{"type": "anthropic", "skill_id": "pptx", "version": "latest"}],
            max_tokens=128000,
        ),
        instructions=[
            "You are a structured content generator for PowerPoint presentations.",
            "Generate EXACTLY the number of slides specified in the task.",
            "Use one clear title per slide with concise bullet points.",
            "Do NOT apply custom fonts, colors, or theme styling.",
            "Do NOT add animations, transitions, or speaker notes.",
            "Keep tables to max 6 rows x 5 columns.",
            "For charts: use only native PPTX chart objects (bar, column, line, or pie). Do NOT use matplotlib, PIL, or any image-based approach for charts.",
            "For any chart mentioned in a visual_suggestion: produce a native Office chart with synthesized data — never embed a chart as an image.",
            "For tables: use ONLY native PPTX table objects — never embed a table as an image or use matplotlib/PIL to render one.",
            "For infographics or diagrams: use native PowerPoint shapes (rectangles, arrows, text boxes) or a native table to approximate the visual — never insert an image for an infographic or diagram.",
            "When a visual request is ambiguous or over-specified, prefer deterministic native data-vis output and keep slide completeness over visual perfection.",
        ],
        markdown=True,
    )

    client = Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

    for attempt in range(max_retries + 1):
        attempt_start = time.time()

        if attempt > 0:
            delay_ms = int(1000 * (2 ** (attempt - 1)))  # exponential: 1000ms, 2000ms
            print(
                "[CHUNK %d] Retry %d/%d after %dms delay..."
                % (chunk_idx, attempt, max_retries, delay_ms)
            )
            time.sleep(delay_ms / 1000.0)

        print(
            "[CHUNK %d] API call attempt %d/%d (slides %d-%d)..."
            % (chunk_idx, attempt + 1, max_retries + 1, first_slide, last_slide)
        )

        try:
            response = None
            event_count = 0
            timed_out = False
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
                future = executor.submit(_run_chunk_agent, chunk_agent, chunk_prompt)
                try:
                    response, event_count = future.result(timeout=CHUNK_TIMEOUT_SECONDS)
                except concurrent.futures.TimeoutError:
                    print(
                        "[CHUNK %d] Attempt %d/%d timed out after %ds. Activating fallback generator."
                        % (
                            chunk_idx,
                            attempt + 1,
                            max_retries + 1,
                            CHUNK_TIMEOUT_SECONDS,
                        )
                    )
                    session_state["use_fallback_generator"] = True
                    timed_out = True
                except Exception as e:
                    raise

            if timed_out:
                return None

            if response is None:
                print(
                    "[CHUNK %d] No RunOutput received after %d events."
                    % (chunk_idx, event_count)
                )
                attempt_elapsed = time.time() - attempt_start
                print(
                    "[TIMING] Chunk %d attempt %d/%d: %.1fs (no output)"
                    % (chunk_idx, attempt + 1, max_retries + 1, attempt_elapsed)
                )
                continue

            if VERBOSE:  # noqa: F405
                msg_count = len(response.messages) if response.messages else 0
                print(
                    "[VERBOSE] Chunk %d attempt %d: received %d events, %d messages"
                    % (chunk_idx, attempt + 1, event_count, msg_count)
                )
                if response.messages:
                    for m_idx, msg in enumerate(response.messages):
                        print(
                            "[VERBOSE] Chunk %d message %d: type=%s role=%s has_provider_data=%s"
                            % (
                                chunk_idx,
                                m_idx,
                                type(msg).__name__,
                                getattr(msg, "role", "N/A"),
                                bool(getattr(msg, "provider_data", None)),
                            )
                        )

        except Exception as e:
            print(
                "[CHUNK %d] Attempt %d/%d failed with error: %s"
                % (chunk_idx, attempt + 1, max_retries + 1, e)
            )
            attempt_elapsed = time.time() - attempt_start
            print(
                "[TIMING] Chunk %d attempt %d/%d: %.1fs (error)"
                % (chunk_idx, attempt + 1, max_retries + 1, attempt_elapsed)
            )
            if attempt == max_retries:
                print(
                    "[CHUNK %d] All %d attempts failed. Skipping chunk."
                    % (chunk_idx, max_retries + 1)
                )
            continue

        # Try to download the generated file from message provider_data
        generated_file = None

        if response.messages:
            for msg in response.messages:
                if hasattr(msg, "provider_data") and msg.provider_data:
                    if VERBOSE:  # noqa: F405
                        print(
                            "[VERBOSE] Chunk %d: attempting file download from message provider_data..."
                            % chunk_idx
                        )
                    try:
                        files = download_skill_files(
                            msg.provider_data, client, output_dir=output_dir
                        )
                    except Exception as e:
                        print(
                            "[CHUNK %d] download_skill_files (message) failed: %s"
                            % (chunk_idx, e)
                        )
                        files = []

                    if VERBOSE:  # noqa: F405
                        print(
                            "[VERBOSE] Chunk %d: download returned files: %s"
                            % (chunk_idx, files)
                        )

                    if files:
                        for f in files:
                            if not f.endswith(".pptx"):
                                continue
                            try:
                                Presentation(f)
                                generated_file = f
                                break
                            except Exception:
                                continue
                    if generated_file:
                        break

        # Fallback: try response.model_provider_data
        if (
            not generated_file
            and hasattr(response, "model_provider_data")
            and response.model_provider_data
        ):
            if VERBOSE:  # noqa: F405
                print(
                    "[VERBOSE] Chunk %d: trying fallback model_provider_data download..."
                    % chunk_idx
                )
            try:
                files = download_skill_files(
                    response.model_provider_data, client, output_dir=output_dir
                )
                if VERBOSE:  # noqa: F405
                    print(
                        "[VERBOSE] Chunk %d: fallback download returned files: %s"
                        % (chunk_idx, files)
                    )
                for f in files:
                    if not f.endswith(".pptx"):
                        continue
                    try:
                        Presentation(f)
                        generated_file = f
                        break
                    except Exception:
                        continue
            except Exception as e:
                print(
                    "[CHUNK %d] download_skill_files (fallback) failed: %s"
                    % (chunk_idx, e)
                )

        attempt_elapsed = time.time() - attempt_start

        if generated_file and os.path.exists(generated_file):
            # Normalize to standard chunk name
            if generated_file != chunk_output_path:
                shutil.copy2(generated_file, chunk_output_path)
                generated_file = chunk_output_path
            print(
                "[TIMING] Chunk %d attempt %d/%d: %.1fs (success)"
                % (chunk_idx, attempt + 1, max_retries + 1, attempt_elapsed)
            )
            print("[CHUNK %d] Successfully generated: %s" % (chunk_idx, generated_file))
            return generated_file
        else:
            print(
                "[TIMING] Chunk %d attempt %d/%d: %.1fs (no file returned)"
                % (chunk_idx, attempt + 1, max_retries + 1, attempt_elapsed)
            )
            print(
                "[CHUNK %d] Attempt %d/%d produced no file."
                % (chunk_idx, attempt + 1, max_retries + 1)
            )

    print(
        "[CHUNK %d] All %d attempts failed. Skipping chunk."
        % (chunk_idx, max_retries + 1)
    )
    session_state["use_fallback_generator"] = True
    return None


# Slide type -> python-pptx layout index mapping for the fallback generator.
# Index 0: Title Slide (large title + subtitle, used for opening/closing)
# Index 1: Title and Content (standard bullet slide, the most common layout)
# Index 2: Section Header (bold title only, ideal for agenda/divider slides)
# Index 3: Two Content (two side-by-side content areas, good for data/comparison)
# Fallback for unknown types: index 1 (Title and Content)
FALLBACK_SLIDE_LAYOUT_MAP = {
    "title": 0,
    "agenda": 2,
    "content": 1,
    "data": 3,
    "closing": 0,
}


def _detect_chart_type_from_suggestion(visual: str) -> Optional[object]:
    """Return an XL_CHART_TYPE enum value when the visual_suggestion describes a chart.

    Scans the visual_suggestion string for chart-type keywords.
    Returns None if no recognizable chart keyword is found (e.g. an image,
    icon, or diagram description), in which case the caller should fall back
    to a textbox annotation.

    Keyword mapping:
      pie                        -> XL_CHART_TYPE.PIE
      line                       -> XL_CHART_TYPE.LINE
      bar (horizontal context)   -> XL_CHART_TYPE.BAR_CLUSTERED
      column / chart (generic)   -> XL_CHART_TYPE.COLUMN_CLUSTERED
    """
    from pptx.enum.chart import XL_CHART_TYPE

    v = visual.lower()
    if "pie" in v:
        return XL_CHART_TYPE.PIE
    if "line" in v:
        return XL_CHART_TYPE.LINE
    if "bar" in v:
        return XL_CHART_TYPE.BAR_CLUSTERED
    if "column" in v or "chart" in v:
        return XL_CHART_TYPE.COLUMN_CLUSTERED
    return None


def generate_chunk_pptx_fallback(
    chunk_slides: List[SlideStoryboard],
    session_state: Dict,
    chunk_idx: int,
) -> Optional[str]:
    """Tier 3 (last-resort) chunk generator using python-pptx directly.

    No Claude API call — generates slides programmatically from SlideStoryboard
    data with zero network I/O. Always produces a valid .pptx or returns None
    only on an extreme exception (e.g., disk full).

    This is the last tier in the 3-tier fallback hierarchy:
      Tier 1: Claude PPTX skill (generate_chunk_pptx)
      Tier 2: LLM code generation (generate_chunk_pptx_v2)
      Tier 3: This function — python-pptx direct, <100ms

    Output slides contain title text + bullet points. When a slide's
    visual_suggestion contains a chart-type keyword (bar, column, line, pie, chart),
    a native python-pptx chart is inserted using CategoryChartData + add_chart()
    (no matplotlib or image embedding). Non-chart visuals receive a small textbox
    annotation. The output is a structurally valid .pptx compatible with
    step_process_chunks() template assembly and _merge_pptx_zip_level() merging.

    Slide layout mapping:
    - slide_type == "title"   -> layout index 0 (TITLE slide)
    - all others              -> layout index 1 (TITLE_AND_CONTENT)

    Args:
        chunk_slides: List of SlideStoryboard objects for this chunk.
        session_state: Shared workflow session state.
        chunk_idx: 0-based chunk index (used for file naming and logging).

    Returns:
        Path to the generated chunk PPTX file, or None if generation failed.
    """
    output_dir = session_state.get("output_dir", ".")
    output_path = os.path.join(output_dir, "chunk_%03d.pptx" % chunk_idx)

    try:
        prs = Presentation()

        for slide in chunk_slides:
            layout_idx = FALLBACK_SLIDE_LAYOUT_MAP.get(slide.slide_type, 1)
            # Guard: some presentations may have fewer layouts than expected.
            # Fall back to index 1 (Title and Content) if the chosen index is out of range.
            if layout_idx >= len(prs.slide_layouts):
                layout_idx = min(1, len(prs.slide_layouts) - 1)
            slide_layout = prs.slide_layouts[layout_idx]

            pptx_slide = prs.slides.add_slide(slide_layout)

            # Set title
            if pptx_slide.shapes.title:
                pptx_slide.shapes.title.text = slide.slide_title

            if slide.slide_type == "title":
                # Set subtitle placeholder
                subtitle_text = (
                    slide.key_points[0] if slide.key_points else slide.slide_title
                )
                if subtitle_text:
                    for ph in pptx_slide.placeholders:
                        if ph.placeholder_format.idx == 1:
                            ph.text = subtitle_text
                            break
            else:
                # Set body content with key points
                body_ph = None
                for ph in pptx_slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        body_ph = ph
                        break
                if body_ph:
                    tf = body_ph.text_frame
                    tf.word_wrap = True
                    tf.clear()
                    for i, point in enumerate(slide.key_points):
                        if i == 0:
                            tf.paragraphs[0].text = point
                        else:
                            p = tf.add_paragraph()
                            p.text = point
                        tf.paragraphs[i].level = 0

            # Add visual element if applicable — prefer native chart over textbox label
            visual = slide.visual_suggestion
            if visual and visual.lower() != "none":
                chart_xl_type = _detect_chart_type_from_suggestion(visual)
                if chart_xl_type is not None:
                    # Insert a native python-pptx chart with synthesized sample data.
                    # CategoryChartData / add_chart() produce editable Office chart objects;
                    # no matplotlib or image embedding is used.
                    try:
                        from pptx.chart.data import CategoryChartData

                        chart_data = CategoryChartData()
                        chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
                        chart_data.add_series("Series 1", (25, 40, 35, 55))
                        pptx_slide.shapes.add_chart(
                            chart_xl_type,
                            Inches(1),
                            Inches(2.5),
                            Inches(8),
                            Inches(3.5),
                            chart_data,
                        )
                    except Exception as chart_err:
                        print(
                            "[CHUNK FALLBACK] Native chart insertion failed: %s; "
                            "falling back to textbox label." % chart_err
                        )
                        txBox = pptx_slide.shapes.add_textbox(
                            Inches(7), Inches(5.5), Inches(2.5), Inches(0.5)
                        )
                        tf = txBox.text_frame
                        tf.text = "[Chart: %s]" % visual[:60]
                        if tf.paragraphs and tf.paragraphs[0].runs:
                            run = tf.paragraphs[0].runs[0]
                        elif tf.paragraphs:
                            run = tf.paragraphs[0].add_run()
                        else:
                            run = None
                        if run is not None:
                            run.font.size = Pt(8)
                            run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
                else:
                    v_lower = visual.lower()
                    if "table" in v_lower:
                        # Insert a native 2x2 placeholder table using python-pptx add_table().
                        # No matplotlib or image embedding — all native OOXML.
                        try:
                            table_shape = pptx_slide.shapes.add_table(
                                2, 2, Inches(1), Inches(2.5), Inches(8), Inches(3.0)
                            )
                            tbl = table_shape.table
                            tbl.cell(0, 0).text = "Item"
                            tbl.cell(0, 1).text = "Value"
                            tbl.cell(1, 0).text = "—"
                            tbl.cell(1, 1).text = "—"
                        except Exception as tbl_err:
                            print(
                                "[CHUNK FALLBACK] Native table insertion failed: %s; "
                                "falling back to textbox label." % tbl_err
                            )
                            txBox = pptx_slide.shapes.add_textbox(
                                Inches(7), Inches(5.5), Inches(2.5), Inches(0.5)
                            )
                            tf = txBox.text_frame
                            tf.text = "[Table: %s]" % visual[:60]
                            if tf.paragraphs and tf.paragraphs[0].runs:
                                run = tf.paragraphs[0].runs[0]
                            elif tf.paragraphs:
                                run = tf.paragraphs[0].add_run()
                            else:
                                run = None
                            if run is not None:
                                run.font.size = Pt(8)
                                run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
                    elif "infographic" in v_lower or "diagram" in v_lower:
                        # Insert labeled rectangle shapes to approximate an infographic/diagram.
                        # Uses native python-pptx shapes — no images inserted.
                        try:
                            from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

                            labels = ["Step 1", "Step 2", "Step 3"]
                            box_w = Inches(2.2)
                            box_h = Inches(1.0)
                            top = Inches(2.8)
                            for li, label in enumerate(labels):
                                left = Inches(0.8 + li * 2.8)
                                shape = pptx_slide.shapes.add_shape(
                                    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                    left,
                                    top,
                                    box_w,
                                    box_h,
                                )
                                shape.text = label
                                shape.text_frame.paragraphs[0].font.size = Pt(11)
                        except Exception as inf_err:
                            print(
                                "[CHUNK FALLBACK] Infographic shape insertion failed: %s; "
                                "falling back to textbox label." % inf_err
                            )
                            txBox = pptx_slide.shapes.add_textbox(
                                Inches(7), Inches(5.5), Inches(2.5), Inches(0.5)
                            )
                            tf = txBox.text_frame
                            tf.text = "[Visual: %s]" % visual[:60]
                            if tf.paragraphs and tf.paragraphs[0].runs:
                                run = tf.paragraphs[0].runs[0]
                            elif tf.paragraphs:
                                run = tf.paragraphs[0].add_run()
                            else:
                                run = None
                            if run is not None:
                                run.font.size = Pt(8)
                                run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
                    else:
                        # Non-chart, non-table, non-diagram visual: add a label annotation
                        txBox = pptx_slide.shapes.add_textbox(
                            Inches(7), Inches(5.5), Inches(2.5), Inches(0.5)
                        )
                        tf = txBox.text_frame
                        tf.text = "[Visual: %s]" % visual[:60]
                        if tf.paragraphs and tf.paragraphs[0].runs:
                            run = tf.paragraphs[0].runs[0]
                        elif tf.paragraphs:
                            run = tf.paragraphs[0].add_run()
                        else:
                            run = None
                        if run is not None:
                            run.font.size = Pt(8)
                            run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

        prs.save(output_path)
        print(
            "[CHUNK %d FALLBACK] Generated %d slides via python-pptx fallback"
            % (chunk_idx, len(chunk_slides))
        )
        return output_path

    except Exception as e:
        print(
            "[CHUNK %d FALLBACK] Failed to generate fallback PPTX: %s" % (chunk_idx, e)
        )
        return None


# Instructions for the Tier 2 fallback code-generation agent.
# Kept as a module-level constant to avoid string duplication.
PPTX_CODE_GEN_INSTRUCTIONS = [
    "You are a Python code generator that creates PowerPoint presentations using python-pptx.",
    "When given slide specifications, write a COMPLETE, SELF-CONTAINED Python script that generates a .pptx file.",
    "The script must import all required libraries at the top.",
    "ALLOWED imports only: pptx, pptx.util, pptx.chart.data, pptx.enum.chart, pptx.dml.color, pptx.enum.text, io, os, os.path, collections, math.",
    "FORBIDDEN imports: matplotlib, matplotlib.pyplot, subprocess, socket, requests, urllib, httpx, shutil, glob, sys, importlib, __import__.",
    "For each slide, create one slide in the presentation using prs.slides.add_slide(prs.slide_layouts[N]).",
    "Slide layout indices: 0=Title Slide, 1=Title and Content, 2=Section Header, 3=Two Content.",
    "For CHARTS: ALWAYS use python-pptx native CategoryChartData or ChartData (creates editable Office charts). NEVER use matplotlib, PIL, or any image-based approach for charts.",
    "For python-pptx ChartData bar/column charts: from pptx.chart.data import ChartData; from pptx.enum.chart import XL_CHART_TYPE.",
    "ChartData example: chart_data = ChartData(); chart_data.categories = ['A','B','C']; chart_data.add_series('Series1', (10, 20, 30)); slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(1), Inches(1.5), Inches(8), Inches(4.5), chart_data).",
    "For line charts use XL_CHART_TYPE.LINE, for pie charts use XL_CHART_TYPE.PIE — always via ChartData, never via image embedding.",
    "For TABLES: ALWAYS use slide.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(4.5)). Fill cells via table.cell(row, col).text = 'value'. NEVER embed a table as an image or use matplotlib/PIL to render one.",
    "For INFOGRAPHICS or DIAGRAMS: use native python-pptx shapes (add_shape with MSO_AUTO_SHAPE_TYPE.RECTANGLE, add_textbox, add_connector) or a native table to approximate the layout. NEVER insert an image for an infographic or diagram.",
    "Treat charts/tables/infographics/diagrams as native data-vis. Preserve data-vis intent even when exact styling cannot be replicated.",
    "Synthesize plausible, specific data values from the visual_suggestion and key_points descriptions. Do NOT use generic placeholder data.",
    "CHART AXIS SAFETY: When styling chart axes, always wrap axis.format.line.fill.background(), axis.format.line.color.rgb, and axis.major_gridlines.format.line.color.rgb calls in try/except blocks to handle cases where the underlying XML element does not yet exist. Example: try: chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0x80,0x80,0x80)\nexcept Exception: pass",
    "CHART AXIS SAFETY: Similarly wrap axis.tick_labels.font.color.rgb, chart.legend.font.color.rgb, and series.data_labels.font.color.rgb in try/except blocks — these can raise 'NoneType object has no attribute attrib' if the XML element is absent.",
    "Save the final presentation using prs.save('EXACT_OUTPUT_PATH') where EXACT_OUTPUT_PATH is the path given in the task.",
    "Execute the script using the save_to_file_and_run tool immediately after writing it.",
    "If the script has an error, fix it and re-run. Maximum 2 fix attempts.",
    "If a visual cannot be implemented exactly, keep the slide and add a concise native textbox note. Do not skip slides.",
    "Do not add speaker notes, animations, or transitions.",
    "Do not print to stdout or write any files other than the final prs.save() call.",
]

# Tier 2 fallback agent: generates python-pptx code with native charts and executes it.
# Created at module level — NOT inside any function or loop (per project rules).
# Uses Claude without the PPTX skill, with PythonTools for code execution.
# PythonTools base_dir is set to current directory; the generated script saves
# to the absolute output path passed in the prompt.
# NOTE: betas=["context-1m-2025-08-07"] is intentionally omitted because this agent
# is invoked with stream=False (see generate_chunk_pptx_v2), and the context-1m beta
# mandates streaming for every call — using it with stream=False raises:
# "Streaming is required for operations that may take longer than 10 minutes."
fallback_code_agent = Agent(
    name="PPTX Code Generator",
    model=Claude(id="claude-opus-4-6", max_tokens=128000),
    instructions=PPTX_CODE_GEN_INSTRUCTIONS,
    tools=[
        PythonTools(
            base_dir=Path("."),
        )
    ],
    markdown=False,
)


def generate_chunk_pptx_v2(
    chunk_slides: List[SlideStoryboard],
    session_state: Dict,
    chunk_idx: int,
) -> Optional[str]:
    """Tier 2 fallback chunk generator using LLM code generation + PythonTools execution.

    Prompts the fallback_code_agent (Claude Opus without PPTX skill, equipped with
    PythonTools) to write and execute a python-pptx script that creates the chunk
    slides with native charts (CategoryChartData/ChartData), tables, and rich visual content.
    matplotlib is FORBIDDEN — all charts must use python-pptx native chart objects.

    Quality level: 80-92% parity with Tier 1 (Claude PPTX skill). Charts and tables
    are generated via python-pptx native objects rather than the native PPTX skill,
    so visual fidelity may differ in edge cases.

    This is Tier 2 in the three-tier fallback hierarchy:
      Tier 1: Claude PPTX skill (generate_chunk_pptx)            — 100% quality
      Tier 2: LLM code generation (this function)                 — 80-92% quality
      Tier 3: Text-only python-pptx (generate_chunk_pptx_fallback) — structural only

    No retries are attempted within this function — callers escalate to Tier 3 on
    failure. Use generate_chunk_pptx() for retry logic (Tier 1 only).

    Args:
        chunk_slides: List of SlideStoryboard objects for this chunk.
        session_state: Shared workflow session state.
        chunk_idx: 0-based chunk index (used for file naming and logging).

    Returns:
        Path to the generated chunk PPTX file, or None if generation failed.
    """
    storyboard: StoryboardPlan = session_state.get("storyboard")
    output_dir = session_state.get("output_dir", ".")
    chunk_output_path = os.path.join(output_dir, "chunk_%03d.pptx" % chunk_idx)

    first_slide = chunk_slides[0].slide_number
    last_slide = chunk_slides[-1].slide_number

    print(
        "[CHUNK %d TIER2] Starting LLM code generation fallback (slides %d-%d)..."
        % (chunk_idx, first_slide, last_slide)
    )

    # Build the code generation prompt with full slide specifications
    slide_specs = []
    for slide in chunk_slides:
        spec = (
            "Slide %d (type=%s): title='%s'\n  Key points: %s\n  Visual suggestion: %s"
        ) % (
            slide.slide_number,
            slide.slide_type,
            slide.slide_title,
            "; ".join(slide.key_points),
            slide.visual_suggestion,
        )
        slide_specs.append(spec)

    global_ctx = ""
    if storyboard:
        global_ctx = (
            "Presentation: '%s' | Audience: %s | Tone: %s | Brand voice: %s\n"
            "Context: %s"
        ) % (
            storyboard.presentation_title,
            storyboard.target_audience,
            storyboard.tone,
            storyboard.brand_voice,
            storyboard.global_context,
        )

    code_gen_prompt = (
        "Generate a complete Python script using python-pptx to create "
        "a PowerPoint file at this EXACT path: %s\n\n"
        "GLOBAL CONTEXT:\n%s\n\n"
        "SLIDES TO GENERATE (%d slides):\n%s\n\n"
        "REQUIREMENTS:\n"
        "- Create exactly %d slides in the exact order listed above.\n"
        "- For 'title' type slides: use prs.slide_layouts[0] (Title Slide layout).\n"
        "- For 'agenda'/'section' type slides: use prs.slide_layouts[2] (Section Header).\n"
        "- For 'content' type slides: use prs.slide_layouts[1] (Title and Content).\n"
        "- For 'data' type slides: use prs.slide_layouts[3] if available, else [1].\n"
        "- For 'closing' type slides: use prs.slide_layouts[0].\n"
        "- For any slide with a chart visual_suggestion: generate a REAL chart using "
        "  python-pptx native ChartData ONLY (e.g. CategoryChartData + slide.shapes.add_chart()).\n"
        "  Do NOT use matplotlib, PIL, or any image-based approach for charts.\n"
        "- Synthesize specific, plausible data values matching the visual_suggestion topic.\n"
        "- For any slide with a table visual_suggestion: generate a REAL native table using slide.shapes.add_table(). NEVER use matplotlib or embed a table as an image.\n"
        "- For any slide with an infographic or diagram visual_suggestion: use native python-pptx shapes (add_shape with MSO_AUTO_SHAPE_TYPE.RECTANGLE, add_textbox) or a native table to approximate it. NEVER insert an image for an infographic or diagram.\n"
        "- If exact visual styling is not feasible, preserve content and structure with a concise native textbox note. Never skip a requested slide.\n"
        "- Save the file to: %s\n"
        "- Then immediately execute the script using save_to_file_and_run.\n"
        "- Return only tool execution needed to write and run the script; avoid extra narrative output.\n"
    ) % (
        chunk_output_path,
        global_ctx,
        len(chunk_slides),
        "\n\n".join(slide_specs),
        len(chunk_slides),
        chunk_output_path,
    )

    if VERBOSE:  # noqa: F405
        print(
            "[VERBOSE] Chunk %d Tier 2 code-gen prompt length: %d chars"
            % (chunk_idx, len(code_gen_prompt))
        )

    # stream=True is required: max_tokens=128000 causes the Anthropic SDK to enforce
    # streaming for all calls that may take longer than 10 minutes, even without betas.
    # Tier 2 only cares whether the file was created on disk, not about the response
    # content, so we iterate through the stream and discard events.
    t2_start = time.time()
    try:
        for _ in fallback_code_agent.run(code_gen_prompt, stream=True):
            pass
        t2_elapsed = time.time() - t2_start
        print(
            "[TIMING] Chunk %d Tier 2 code generation: %.1fs" % (chunk_idx, t2_elapsed)
        )
    except Exception as e:
        t2_elapsed = time.time() - t2_start
        print(
            "[CHUNK %d TIER2] Code generation agent failed after %.1fs: %s"
            % (chunk_idx, t2_elapsed, str(e))
        )
        print("[CHUNK %d TIER2] Falling back to Tier 3 (text-only)." % chunk_idx)
        return None

    # Verify the file was actually written by the executed code
    if os.path.exists(chunk_output_path):
        try:
            Presentation(chunk_output_path)
            print(
                "[CHUNK %d TIER2] Successfully generated via LLM code execution: %s"
                % (chunk_idx, chunk_output_path)
            )
            return chunk_output_path
        except Exception as e:
            print(
                "[CHUNK %d TIER2] Generated file is invalid PPTX: %s — falling back to Tier 3."
                % (chunk_idx, str(e))
            )
            return None
    else:
        print(
            "[CHUNK %d TIER2] No file produced at %s — falling back to Tier 3."
            % (chunk_idx, chunk_output_path)
        )
        return None


# === WORKFLOW STEP 2: GENERATE CHUNKS ===


def step_generate_chunks(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 2: Orchestrate chunked PPTX generation across all slide groups.

    Splits the full storyboard into chunks of {chunk_size} slides, dispatches
    generation for each chunk using the 3-tier fallback hierarchy, and stores results.

    Each chunk is generated using a 3-tier fallback:
      - Tier 1 (Claude PPTX skill): attempted unless session flag is active
      - Tier 2 (LLM code generation): used when Tier 1 fails or is bypassed
      - Tier 3 (text-only python-pptx): used when Tier 2 also fails
    The session-level flag "use_fallback_generator" persists across chunks —
    once set, Tier 1 is skipped for all remaining chunks in the run.

    A 1-second inter-chunk delay is applied between calls to avoid rate limits.

    Args:
        step_input: Workflow step input (not used directly).
        session_state: Shared workflow state; must contain storyboard, chunk_size, output_dir,
                       max_retries, start_tier, and use_fallback_generator.

    Returns:
        StepOutput with success=True when at least one chunk was generated successfully,
        or success=False if no storyboard was found or all chunks failed.
    """
    step_start = time.time()

    storyboard: Optional[StoryboardPlan] = session_state.get("storyboard")
    if not storyboard:
        print("[ERROR] No storyboard found in session_state.")
        return StepOutput(content="No storyboard found.", success=False)

    chunk_size = session_state.get("chunk_size", 3)
    slides = storyboard.slides

    # Build chunk list
    chunks = [slides[i : i + chunk_size] for i in range(0, len(slides), chunk_size)]

    print("\n" + "=" * 60)
    print("Step 2: Generating presentation chunks...")
    print("=" * 60)
    print(
        "Total slides: %d | Chunk size: %d | Number of chunks: %d"
        % (len(slides), chunk_size, len(chunks))
    )

    if VERBOSE:  # noqa: F405
        for ci, chunk in enumerate(chunks):
            slide_nums = [s.slide_number for s in chunk]
            print("[VERBOSE] Chunk %d: slides %s" % (ci, slide_nums))

    chunk_files: List[Optional[str]] = []
    successful = 0
    total_chunks = len(chunks)
    start_tier = session_state.get("start_tier", 1)

    for chunk_idx, chunk_slides in enumerate(chunks):
        chunk_start = time.time()
        print(
            "[GENERATE] Chunk %d/%d: slides %d-%d"
            % (
                chunk_idx + 1,
                total_chunks,
                chunk_slides[0].slide_number,
                chunk_slides[-1].slide_number,
            )
        )

        # Determine effective starting tier for this chunk
        # If use_fallback_generator flag is set (Tier 1 failed in a previous chunk),
        # effective tier is max(start_tier, 2) to skip Tier 1 permanently.
        effective_tier = (
            max(start_tier, 2)
            if session_state.get("use_fallback_generator")
            else start_tier
        )

        chunk_file = None

        if effective_tier == 3:
            # Start with Tier 3: text-only python-pptx (no fallback needed)
            print(
                "[GENERATE] Chunk %d/%d: Starting at Tier 3 (text-only, instant)."
                % (chunk_idx + 1, total_chunks)
            )
            chunk_file = generate_chunk_pptx_fallback(
                chunk_slides, session_state, chunk_idx
            )

        elif effective_tier == 2:
            # Start with Tier 2: LLM code generation, fallback to Tier 3
            print(
                "[GENERATE] Chunk %d/%d: Starting at Tier 2 (LLM code generation)."
                % (chunk_idx + 1, total_chunks)
            )
            chunk_file = generate_chunk_pptx_v2(chunk_slides, session_state, chunk_idx)
            if chunk_file is None:
                print(
                    "[GENERATE] Chunk %d/%d: Tier 2 failed. Falling back to Tier 3 (text-only)."
                    % (chunk_idx + 1, total_chunks)
                )
                chunk_file = generate_chunk_pptx_fallback(
                    chunk_slides, session_state, chunk_idx
                )

        else:  # effective_tier == 1
            # Start with Tier 1: Claude PPTX skill, fallback to Tier 2 → Tier 3
            print(
                "[GENERATE] Chunk %d/%d: Starting at Tier 1 (Claude PPTX skill)."
                % (chunk_idx + 1, total_chunks)
            )
            chunk_file = generate_chunk_pptx(chunk_slides, session_state, chunk_idx)

            if chunk_file is None and session_state.get("use_fallback_generator"):
                # Tier 1 failed for this chunk (and activated the session flag).
                # Try Tier 2 immediately for this chunk.
                print(
                    "[GENERATE] Chunk %d/%d: Tier 1 failed. Attempting Tier 2 (LLM code generation)..."
                    % (chunk_idx + 1, total_chunks)
                )
                chunk_file = generate_chunk_pptx_v2(
                    chunk_slides, session_state, chunk_idx
                )

                if chunk_file is None:
                    # Tier 2 also failed — fall through to Tier 3.
                    print(
                        "[GENERATE] Chunk %d/%d: Tier 2 also failed. "
                        "Using Tier 3 (text-only fallback)."
                        % (chunk_idx + 1, total_chunks)
                    )
                    chunk_file = generate_chunk_pptx_fallback(
                        chunk_slides, session_state, chunk_idx
                    )

        chunk_elapsed = time.time() - chunk_start
        if chunk_file:
            print(
                "[TIMING] Chunk %d/%d done in %.1fs -> %s"
                % (chunk_idx + 1, total_chunks, chunk_elapsed, chunk_file)
            )
        else:
            print(
                "[TIMING] Chunk %d/%d FAILED after %.1fs (skipping)"
                % (chunk_idx + 1, total_chunks, chunk_elapsed)
            )
            print(
                "[WARNING] Chunk %d failed (slides %d-%d). Continuing..."
                % (
                    chunk_idx,
                    chunk_slides[0].slide_number,
                    chunk_slides[-1].slide_number,
                )
            )

        chunk_files.append(chunk_file)

        if chunk_file:
            successful += 1

        # Inter-chunk delay (not for last chunk) to avoid rate limits
        if chunk_idx < total_chunks - 1:
            inter_delay = 1.0
            print("[GENERATE] Waiting %.1fs before next chunk..." % inter_delay)
            time.sleep(inter_delay)

    session_state["chunk_files"] = chunk_files
    session_state["chunk_slide_groups"] = chunks

    step_elapsed = time.time() - step_start
    failed = total_chunks - successful
    print(
        "\n[TIMING] step_generate_chunks completed in %.1fs (%d chunks: %d succeeded, %d failed)"
        % (step_elapsed, total_chunks, successful, failed)
    )

    summary = "%d of %d chunks generated successfully. Duration: %.1fs" % (
        successful,
        total_chunks,
        step_elapsed,
    )
    return StepOutput(content=summary, success=successful > 0)


# === WORKFLOW STEP 3: PROCESS CHUNKS (TEMPLATE + IMAGES) ===


def step_process_chunks(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 3: Apply template assembly and image pipeline to each chunk.

    For each successfully generated chunk, runs:
    1. Image planning (which slides need AI-generated images)
    2. Image generation (NanoBanana)
    3. Template assembly (if --template is provided)

    Each chunk is processed with a temporary session_state copy that adapts
    the existing step functions to work on individual chunk files.

    Args:
        step_input: Workflow step input (not used directly).
        session_state: Shared workflow state; must contain chunk_files, chunk_slide_groups,
                       template_path, output_dir, and no_images.

    Returns:
        StepOutput with success=True always (failures are logged but do not abort the step).
        Content string reports how many chunks were processed.
    """
    step_start = time.time()

    print("\n" + "=" * 60)
    print("Step 3: Processing chunks (images + template assembly)...")
    print("=" * 60)

    chunk_files: List[Optional[str]] = session_state.get("chunk_files", [])
    chunk_slide_groups: List[List[SlideStoryboard]] = session_state.get(
        "chunk_slide_groups", []
    )
    template_path = session_state.get("template_path", "")
    output_dir = session_state.get("output_dir", ".")
    no_images = session_state.get("no_images", False)

    processed_chunks: Dict[int, Optional[str]] = {}
    total_process_chunks = len(chunk_files)

    for chunk_idx, chunk_file in enumerate(chunk_files):
        chunk_proc_start = time.time()

        if chunk_file is None:
            print(
                "[PROCESS] Chunk %d (%d/%d): skipped (no file)."
                % (chunk_idx, chunk_idx + 1, total_process_chunks)
            )
            processed_chunks[chunk_idx] = None
            continue

        print(
            "\n[PROCESS] Chunk %d (%d/%d): processing %s"
            % (chunk_idx, chunk_idx + 1, total_process_chunks, chunk_file)
        )

        # Determine which slides are in this chunk
        chunk_slides = (
            chunk_slide_groups[chunk_idx] if chunk_idx < len(chunk_slide_groups) else []
        )
        slides_data = _extract_chunk_slides_data(chunk_file)

        # Enrich slides_data entries with storyboard visual_suggestion and has_data_vis.
        # _extract_chunk_slides_data only sees raw PPTX shapes (has_chart, has_table).
        # For Tier 2/3 chunks the storyboard visual_suggestion is the authoritative
        # signal for "this slide will carry an infographic/diagram/chart/table and must
        # NOT receive an external AI-generated image."
        # Without this enrichment the image_planner receives no keyword signal for
        # data-vis slides, so it may incorrectly plan images for them.
        _DATA_VIS_KW = ("chart", "table", "infographic", "diagram", "graph")
        for _sd_i, _sd in enumerate(slides_data):
            if _sd_i < len(chunk_slides):
                _vs = chunk_slides[_sd_i].visual_suggestion or ""
                _sd["visual_suggestion"] = _vs
                # has_data_vis: true when the storyboard declares a data visual
                # OR the actual PPTX shape inspection found a chart or table.
                _sd["has_data_vis"] = (
                    _sd.get("has_chart", False)
                    or _sd.get("has_table", False)
                    or any(_kw in _vs.lower() for _kw in _DATA_VIS_KW)
                )

        total_chunk_slides = len(slides_data)

        assembled_path = os.path.join(
            output_dir, "chunk_%03d_assembled.pptx" % chunk_idx
        )

        # Build a temporary session_state for the existing step functions
        chunk_session = dict(session_state)
        chunk_session["generated_file"] = chunk_file
        chunk_session["total_slides"] = total_chunk_slides
        chunk_session["slides_data"] = slides_data
        chunk_session["output_path"] = assembled_path
        chunk_session["generated_images"] = {}
        # Use chunk-specific output subdirectory to avoid collisions
        chunk_out_subdir = os.path.join(output_dir, "chunk_%03d_work" % chunk_idx)
        os.makedirs(chunk_out_subdir, exist_ok=True)
        chunk_session["output_dir"] = chunk_out_subdir

        if VERBOSE:  # noqa: F405
            print(
                "[VERBOSE] Chunk %d session state keys: %s"
                % (chunk_idx, sorted(chunk_session.keys()))
            )

        # Adjust src_slide dimensions if not set
        if not chunk_session.get("src_slide_width"):
            try:
                prs = Presentation(chunk_file)
                chunk_session["src_slide_width"] = prs.slide_width
                chunk_session["src_slide_height"] = prs.slide_height
            except Exception:
                pass

        current_file = chunk_file

        # --- Image planning ---
        if not no_images:
            print("[PROCESS] Chunk %d: running image planning..." % chunk_idx)
            try:
                # Build slides JSON for image planner
                slides_json = json.dumps(slides_data, indent=2)
                user_prompt = session_state.get(
                    "user_prompt", "professional presentation"
                )
                combined_message = (
                    'Presentation topic: "%s"\n\nSlide metadata:\n%s\n\n'
                    "Analyze each slide and decide which ones need AI-generated images.\n"
                    "Treat chart/table/infographic/diagram slides as native data-vis and do NOT request image generation for those slides.\n"
                    "Consider the presentation topic when writing image prompts."
                ) % (user_prompt, slides_json)

                img_plan_response = image_planner.run(combined_message, stream=False)  # noqa: F405

                if img_plan_response and img_plan_response.content:
                    content = img_plan_response.content

                    if VERBOSE:  # noqa: F405
                        if isinstance(content, BaseModel):
                            print(
                                "[VERBOSE] Chunk %d image plan:\n%s"
                                % (chunk_idx, content.model_dump_json(indent=2))
                            )
                        else:
                            print(
                                "[VERBOSE] Chunk %d image plan content: %s"
                                % (chunk_idx, str(content)[:500])
                            )

                    if isinstance(content, BaseModel):
                        plan_json = content.model_dump_json()
                    elif isinstance(content, dict):
                        plan_json = json.dumps(content)
                    else:
                        plan_json = str(content)

                    # Create a mock StepInput for step_generate_images
                    mock_input = StepInput(
                        input=user_prompt,
                        previous_step_content=plan_json,
                    )
                    step_generate_images(mock_input, chunk_session)  # noqa: F405
                    print(
                        "[PROCESS] Chunk %d: images generated. Count: %d"
                        % (chunk_idx, len(chunk_session.get("generated_images", {})))
                    )
                else:
                    print(
                        "[PROCESS] Chunk %d: image planner returned no plan."
                        % chunk_idx
                    )

            except Exception as e:
                print("[PROCESS] Chunk %d: image pipeline failed: %s" % (chunk_idx, e))
                if session_state.get("verbose"):
                    traceback.print_exc()

        # --- Template assembly ---
        if template_path and os.path.isfile(template_path):
            print("[PROCESS] Chunk %d: running template assembly..." % chunk_idx)
            try:
                # Propagate generated images back to chunk_session
                mock_assemble_input = StepInput(
                    input=session_state.get("user_prompt", ""),
                    previous_step_content=json.dumps(slides_data),
                )
                step_assemble_template(mock_assemble_input, chunk_session)  # noqa: F405

                assembled_output = chunk_session.get("output_path", assembled_path)
                if assembled_output and os.path.isfile(assembled_output):
                    current_file = assembled_output
                    print(
                        "[PROCESS] Chunk %d: assembled -> %s"
                        % (chunk_idx, current_file)
                    )
                else:
                    print(
                        "[PROCESS] Chunk %d: template assembly produced no file; "
                        "keeping raw chunk." % chunk_idx
                    )
            except Exception as e:
                print(
                    "[PROCESS] Chunk %d: template assembly failed: %s" % (chunk_idx, e)
                )
                if session_state.get("verbose"):
                    traceback.print_exc()
        else:
            # No template: just copy raw chunk to assembled path name for consistency
            shutil.copy2(chunk_file, assembled_path)
            current_file = assembled_path
            print(
                "[PROCESS] Chunk %d: no template; raw chunk copied to %s"
                % (chunk_idx, assembled_path)
            )

        processed_chunks[chunk_idx] = current_file

        chunk_proc_elapsed = time.time() - chunk_proc_start
        print(
            "[TIMING] Chunk %d processing done in %.1fs"
            % (chunk_idx, chunk_proc_elapsed)
        )
        print("[PROCESS] Chunk %d: result -> %s" % (chunk_idx, current_file))

    session_state["processed_chunks"] = processed_chunks
    successful = sum(1 for v in processed_chunks.values() if v is not None)

    step_elapsed = time.time() - step_start
    print(
        "\n[TIMING] step_process_chunks completed in %.1fs (%d chunks processed)"
        % (step_elapsed, successful)
    )

    return StepOutput(
        content="%d of %d chunks processed. Duration: %.1fs"
        % (successful, len(chunk_files), step_elapsed),
        success=True,
    )


# === WORKFLOW STEP 4 (OPTIONAL): VISUAL REVIEW PER CHUNK ===


def step_visual_review_chunks(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 4 (Optional): Run visual inspection on each chunk's assembled PPTX.

    For each chunk:
    1. Render slides to PNG via LibreOffice.
    2. Call slide_quality_reviewer for each slide image.
    3. Apply programmatic corrections if needed.
    4. Repeat up to max_passes passes until no further changes are needed.

    This step is non-blocking: any failure silently returns success=True.
    If a programmatic fix is missing in Python, logs it to console regardless
    of --verbose setting.

    Args:
        step_input: Workflow step input (not used directly).
        session_state: Shared workflow state; must contain processed_chunks, output_dir,
                       template_path, and visual_passes.

    Returns:
        StepOutput with success=True always; errors per chunk are logged and do not abort.
        Content string reports how many chunks were reviewed.
    """
    step_start = time.time()

    print("\n" + "=" * 60)
    print("Step 4 (Optional): Visual review per chunk...")
    print("=" * 60)

    processed_chunks: Dict[int, Optional[str]] = session_state.get(
        "processed_chunks", {}
    )
    output_dir = session_state.get("output_dir", ".")
    template_path = session_state.get("template_path", "")
    max_passes = session_state.get("visual_passes", 3)
    reviewed_chunks: Dict[int, Optional[str]] = {}

    for chunk_idx, assembled_path in sorted(processed_chunks.items()):
        chunk_review_start = time.time()

        if assembled_path is None or not os.path.isfile(assembled_path):
            print(
                "[VISUAL] Chunk %d: skipped (file not found: %s)."
                % (chunk_idx, assembled_path)
            )
            reviewed_chunks[chunk_idx] = None
            continue

        print(
            "\n[VISUAL REVIEW] Chunk %d: starting review of %s"
            % (chunk_idx, assembled_path)
        )

        # Build a per-chunk session_state for the visual review step
        chunk_session = dict(session_state)
        chunk_session["output_path"] = assembled_path
        chunk_session["template_path"] = template_path

        current_path = assembled_path

        for pass_num in range(max_passes):
            pass_start = time.time()
            print(
                "[VISUAL REVIEW] Chunk %d: pass %d/%d starting..."
                % (chunk_idx, pass_num + 1, max_passes)
            )
            chunk_session["output_path"] = current_path

            try:
                mock_input = StepInput(
                    input=session_state.get("user_prompt", ""),
                    previous_step_content="",
                )
                result = step_visual_quality_review(mock_input, chunk_session)  # noqa: F405

                # Check if any actionable issues exist (mirrors _apply_visual_corrections logic).
                # PresentationQualityReport has 'total_critical_issues', NOT 'total_corrections_applied'.
                # We check slide_reports for critical/moderate issues with a real programmatic_fix.
                quality_report = chunk_session.get("quality_report", {})
                slide_reports_data = quality_report.get("slide_reports", [])

                if VERBOSE:  # noqa: F405
                    for r in slide_reports_data:
                        issues = r.get("issues", [])
                        print(
                            "[VERBOSE] Chunk %d pass %d slide %s: %d issues"
                            % (
                                chunk_idx,
                                pass_num + 1,
                                r.get("slide_index", "?"),
                                len(issues),
                            )
                        )
                        for issue in issues:
                            print(
                                "[VERBOSE]   severity=%s fix=%s desc=%s"
                                % (
                                    issue.get("severity", "?"),
                                    issue.get("programmatic_fix", "?"),
                                    str(issue.get("description", ""))[:80],
                                )
                            )

                changes_applied = any(
                    any(
                        i.get("severity") in ("critical", "moderate")
                        and i.get("programmatic_fix") != "none"
                        for i in r.get("issues", [])
                    )
                    for r in slide_reports_data
                )

                pass_elapsed = time.time() - pass_start
                print(
                    "[TIMING] Chunk %d pass %d: %.1fs"
                    % (chunk_idx, pass_num + 1, pass_elapsed)
                )

                if not changes_applied:
                    print(
                        "[VISUAL REVIEW] Chunk %d: pass %d/%d — no changes needed. Done."
                        % (chunk_idx, pass_num + 1, max_passes)
                    )
                    break
                else:
                    print(
                        "[VISUAL REVIEW] Chunk %d: pass %d/%d — corrections applied. Re-checking..."
                        % (chunk_idx, pass_num + 1, max_passes)
                    )

            except Exception as e:
                pass_elapsed = time.time() - pass_start
                print(
                    "[TIMING] Chunk %d pass %d: %.1fs (error)"
                    % (chunk_idx, pass_num + 1, pass_elapsed)
                )
                print(
                    "[VISUAL] Chunk %d, pass %d: review failed: %s"
                    % (chunk_idx, pass_num + 1, e)
                )
                # Log missing programmatic fix to console regardless of verbose mode (per spec).
                # Any exception here means the visual review or correction logic is broken/missing.
                print(
                    "[VISUAL REVIEW MISSING FIX] Chunk %d, pass %d: exception during "
                    "visual correction step: %s" % (chunk_idx, pass_num + 1, str(e))
                )
                print(
                    "[SUGGESTION] Review step_visual_quality_review() and "
                    "_apply_visual_corrections() for the issue type that raised this error. "
                    "Add handling logic if a programmatic_fix type is missing."
                )
                break

        reviewed_chunks[chunk_idx] = chunk_session.get("output_path", current_path)

        chunk_review_elapsed = time.time() - chunk_review_start
        print(
            "[TIMING] Chunk %d total review: %.1fs" % (chunk_idx, chunk_review_elapsed)
        )
        print(
            "[VISUAL REVIEW] Chunk %d: reviewed -> %s"
            % (chunk_idx, reviewed_chunks[chunk_idx])
        )

    session_state["reviewed_chunks"] = reviewed_chunks
    reviewed_count = sum(1 for v in reviewed_chunks.values() if v is not None)

    step_elapsed = time.time() - step_start
    print(
        "\n[TIMING] step_visual_review_chunks completed in %.1fs (%d chunks reviewed)"
        % (step_elapsed, reviewed_count)
    )

    return StepOutput(
        content="%d of %d chunks visually reviewed. Duration: %.1fs"
        % (reviewed_count, len(processed_chunks), step_elapsed),
        success=True,
    )


# === HELPER: MERGE MULTIPLE PPTX FILES ===


def _merge_pptx_zip_level(pptx_paths: List[str], output_path: str) -> bool:
    """Merge multiple PPTX files by manipulating the ZIP structure directly.

    This is the most reliable approach — avoids OPC context issues that cause
    PowerPoint to report "found a problem with content" on the merged file.
    Binary parts (images, charts, workbooks) are copied at the raw bytes level,
    so there is no risk of the OPC package context dropping or corrupting data.

    Args:
        pptx_paths: List of valid, existing PPTX file paths to merge in order.
        output_path: Destination path for the merged presentation.

    Returns:
        True if merge succeeded, False otherwise.
    """
    import posixpath
    import re

    from lxml import etree

    valid_paths = [p for p in pptx_paths if p and os.path.exists(p)]
    if not valid_paths:
        print("[MERGE] No valid PPTX files to merge")
        return False

    if len(valid_paths) == 1:
        shutil.copy(valid_paths[0], output_path)
        print("[MERGE] Single file, copied directly: %s" % output_path)
        return True

    # Use first file as base — copy it to output
    shutil.copy(valid_paths[0], output_path)

    # Read the base presentation XML
    with zipfile.ZipFile(output_path, "r") as base_zip:
        base_prs_xml = base_zip.read("ppt/presentation.xml")
        base_prs_rels_xml = base_zip.read("ppt/_rels/presentation.xml.rels")
        base_content_types_xml = base_zip.read("[Content_Types].xml")

    base_prs_tree = etree.fromstring(base_prs_xml)
    base_prs_rels_tree = etree.fromstring(base_prs_rels_xml)
    base_ct_tree = etree.fromstring(base_content_types_xml)

    # Namespaces
    NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    NS_RELS = "http://schemas.openxmlformats.org/package/2006/relationships"
    NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"

    def _get_slide_numbers(prs_rels_tree):
        """Extract slide XML numbers from ``presentation.xml.rels``.

        Args:
            prs_rels_tree: Parsed lxml tree for ``ppt/_rels/presentation.xml.rels``.

        Returns:
            List of integer slide numbers extracted from Target attributes of slide
            relationships in prs_rels_tree. Non-slide relationships are ignored.
        """
        slide_nums = []
        for rel in prs_rels_tree.findall("{%s}Relationship" % NS_RELS):
            target = rel.get("Target", "")
            m = re.match(r"slides/slide(\d+)\.xml", target)
            if m:
                slide_nums.append(int(m.group(1)))
        return slide_nums

    existing_slide_nums = _get_slide_numbers(base_prs_rels_tree)
    next_slide_num = max(existing_slide_nums) + 1 if existing_slide_nums else 1

    # Track next rel ID for presentation.xml.rels
    existing_rel_ids = [
        int(rel.get("Id", "rId0").replace("rId", ""))
        for rel in base_prs_rels_tree.findall("{%s}Relationship" % NS_RELS)
        if rel.get("Id", "").startswith("rId")
    ]
    next_rel_id = max(existing_rel_ids) + 1 if existing_rel_ids else 100

    # Open output as writable archive
    with zipfile.ZipFile(
        output_path, "a", compression=zipfile.ZIP_DEFLATED, allowZip64=True
    ) as out_zip:
        for src_path in valid_paths[1:]:
            with zipfile.ZipFile(src_path, "r") as src_zip:
                src_names = set(src_zip.namelist())

                # Read source presentation rels to find slides
                src_prs_rels_xml = src_zip.read("ppt/_rels/presentation.xml.rels")
                src_prs_rels = etree.fromstring(src_prs_rels_xml)
                src_slide_nums = sorted(_get_slide_numbers(src_prs_rels))

                for src_slide_num in src_slide_nums:
                    new_slide_num = next_slide_num
                    next_slide_num += 1

                    old_slide_name = "ppt/slides/slide%d.xml" % src_slide_num
                    new_slide_name = "ppt/slides/slide%d.xml" % new_slide_num
                    old_slide_rels_name = (
                        "ppt/slides/_rels/slide%d.xml.rels" % src_slide_num
                    )
                    new_slide_rels_name = (
                        "ppt/slides/_rels/slide%d.xml.rels" % new_slide_num
                    )

                    if old_slide_name not in src_names:
                        continue

                    # Copy slide XML
                    slide_xml_bytes = src_zip.read(old_slide_name)
                    out_zip.writestr(new_slide_name, slide_xml_bytes)

                    # Copy slide rels and rewrite media/chart refs with unique names
                    if old_slide_rels_name in src_names:
                        slide_rels_xml = src_zip.read(old_slide_rels_name)
                        slide_rels_tree = etree.fromstring(slide_rels_xml)

                        for rel in slide_rels_tree.findall(
                            "{%s}Relationship" % NS_RELS
                        ):
                            rel_type = rel.get("Type", "")
                            target = rel.get("Target", "")

                            # Skip slide layout — keep reference as-is
                            if "slideLayout" in rel_type:
                                continue

                            if not target.startswith(".."):
                                continue  # absolute or external refs

                            # Resolve the actual part path in the source zip
                            # e.g. "../media/image1.png" -> "ppt/media/image1.png"
                            # NOTE: lstrip("../") is character-based and strips individual
                            # '.' and '/' chars, giving wrong results for nested paths like
                            # "../charts/chart1.xml" (produces "ppt/slides/charts/chart1.xml"
                            # instead of "ppt/charts/chart1.xml"). Use posixpath.normpath
                            # to handle the ".." parent-dir segment correctly.
                            actual_old = posixpath.normpath("ppt/slides/" + target)

                            if actual_old not in src_names:
                                continue

                            # Generate unique name for target archive
                            basename = os.path.basename(actual_old)
                            stem, ext = os.path.splitext(basename)
                            new_part_name = actual_old
                            counter = 1
                            all_names = set(out_zip.namelist())
                            while new_part_name in all_names:
                                new_part_name = (
                                    os.path.dirname(actual_old)
                                    + "/"
                                    + stem
                                    + "_s%d_%d" % (new_slide_num, counter)
                                    + ext
                                )
                                counter += 1

                            # Copy the part
                            part_bytes = src_zip.read(actual_old)
                            out_zip.writestr(new_part_name, part_bytes)

                            # If chart, also copy its rels and embedded workbook
                            if "chart" in actual_old:
                                chart_basename = os.path.basename(actual_old)
                                chart_rels_old = (
                                    os.path.dirname(actual_old)
                                    + "/_rels/"
                                    + chart_basename
                                    + ".rels"
                                )
                                if chart_rels_old in src_names:
                                    cr_bytes = src_zip.read(chart_rels_old)
                                    cr_new = (
                                        os.path.dirname(new_part_name)
                                        + "/_rels/"
                                        + os.path.basename(new_part_name)
                                        + ".rels"
                                    )
                                    # Copy chart's embedded xlsx workbook and
                                    # build an updated chart rels XML that points
                                    # to the newly-renamed workbook file.
                                    # NOTE: wb_old uses posixpath.normpath to
                                    # correctly resolve "../embeddings/file.xlsx"
                                    # relative to ppt/charts/ (not ppt/slides/).
                                    # lstrip("../") would incorrectly strip leading
                                    # '.' and '/' characters individually.
                                    cr_tree = etree.fromstring(cr_bytes)
                                    for cr_rel in cr_tree.findall(
                                        "{%s}Relationship" % NS_RELS
                                    ):
                                        cr_target = cr_rel.get("Target", "")
                                        if cr_target.startswith(".."):
                                            wb_old = posixpath.normpath(
                                                os.path.dirname(actual_old)
                                                + "/"
                                                + cr_target
                                            )
                                            if wb_old in src_names:
                                                wb_bytes = src_zip.read(wb_old)
                                                wb_stem, wb_ext = os.path.splitext(
                                                    os.path.basename(wb_old)
                                                )
                                                wb_new = wb_old.replace(
                                                    os.path.basename(wb_old),
                                                    wb_stem
                                                    + "_s%d_wb" % new_slide_num
                                                    + wb_ext,
                                                )
                                                if wb_new not in set(
                                                    out_zip.namelist()
                                                ):
                                                    out_zip.writestr(wb_new, wb_bytes)
                                                # Update the chart rels entry to point
                                                # to the renamed workbook.
                                                # The new target must be relative to the
                                                # chart's directory (ppt/charts/), so we
                                                # build a "../<subdir>/<name>" path.
                                                new_wb_rel_target = "../" + "/".join(
                                                    wb_new.split("/")[1:]
                                                )
                                                cr_rel.set("Target", new_wb_rel_target)
                                    # Serialise the (possibly updated) rels tree.
                                    updated_cr_bytes = etree.tostring(
                                        cr_tree,
                                        xml_declaration=True,
                                        encoding="UTF-8",
                                        standalone=True,
                                    )
                                    out_zip.writestr(cr_new, updated_cr_bytes)

                            # Update the slide relationship target to point to the new
                            # unique part name.
                            # NOTE: split("/")[2:] was wrong — it drops the subdirectory,
                            # e.g. "ppt/charts/chart1_s5_1.xml".split("/")[2:] yields
                            # ["chart1_s5_1.xml"] (missing "charts"), giving the incorrect
                            # relative target "../chart1_s5_1.xml".
                            # split("/")[1:] yields ["charts", "chart1_s5_1.xml"],
                            # producing the correct "../charts/chart1_s5_1.xml".
                            new_rel_target = "../" + "/".join(
                                new_part_name.split("/")[1:]
                            )
                            rel.set("Target", new_rel_target)

                        updated_rels_bytes = etree.tostring(
                            slide_rels_tree,
                            xml_declaration=True,
                            encoding="UTF-8",
                            standalone=True,
                        )
                        out_zip.writestr(new_slide_rels_name, updated_rels_bytes)

                    # Register slide in presentation.xml.rels
                    new_rel_id = "rId%d" % next_rel_id
                    next_rel_id += 1
                    new_prs_rel = etree.SubElement(
                        base_prs_rels_tree, "{%s}Relationship" % NS_RELS
                    )
                    new_prs_rel.set("Id", new_rel_id)
                    new_prs_rel.set(
                        "Type",
                        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
                    )
                    new_prs_rel.set("Target", "slides/slide%d.xml" % new_slide_num)

                    # Register slide in presentation.xml sldIdLst
                    sld_id_lst = base_prs_tree.find(".//{%s}sldIdLst" % NS_P)
                    if sld_id_lst is None:
                        sld_id_lst = etree.SubElement(
                            base_prs_tree, "{%s}sldIdLst" % NS_P
                        )

                    existing_ids = [
                        int(el.get("id", 256))
                        for el in sld_id_lst.findall("{%s}sldId" % NS_P)
                    ]
                    new_id = max(existing_ids) + 1 if existing_ids else 256
                    sld_id_el = etree.SubElement(sld_id_lst, "{%s}sldId" % NS_P)
                    sld_id_el.set("id", str(new_id))
                    sld_id_el.set("{%s}id" % NS_R, new_rel_id)

                    # Add content type entry for new slide
                    existing_ct_parts = {
                        el.get("PartName", "")
                        for el in base_ct_tree.findall("{%s}Override" % NS_CT)
                    }
                    new_part_uri = "/ppt/slides/slide%d.xml" % new_slide_num
                    if new_part_uri not in existing_ct_parts:
                        ct_el = etree.SubElement(base_ct_tree, "{%s}Override" % NS_CT)
                        ct_el.set("PartName", new_part_uri)
                        ct_el.set(
                            "ContentType",
                            "application/vnd.openxmlformats-officedocument.presentationml.slide+xml",
                        )

    # Write updated presentation.xml, rels, and content types back.
    # Python's zipfile does not support overwriting entries, so we copy to a temp file.
    import tempfile

    tmp_path = output_path + ".tmp"
    with zipfile.ZipFile(output_path, "r") as old_zip:
        with zipfile.ZipFile(
            tmp_path, "w", compression=zipfile.ZIP_DEFLATED, allowZip64=True
        ) as new_zip:
            for item in old_zip.namelist():
                if item == "ppt/presentation.xml":
                    new_zip.writestr(
                        item,
                        etree.tostring(
                            base_prs_tree,
                            xml_declaration=True,
                            encoding="UTF-8",
                            standalone=True,
                        ),
                    )
                elif item == "ppt/_rels/presentation.xml.rels":
                    new_zip.writestr(
                        item,
                        etree.tostring(
                            base_prs_rels_tree,
                            xml_declaration=True,
                            encoding="UTF-8",
                            standalone=True,
                        ),
                    )
                elif item == "[Content_Types].xml":
                    new_zip.writestr(
                        item,
                        etree.tostring(
                            base_ct_tree,
                            xml_declaration=True,
                            encoding="UTF-8",
                            standalone=True,
                        ),
                    )
                else:
                    new_zip.writestr(item, old_zip.read(item))

    os.replace(tmp_path, output_path)
    print("[MERGE] Saved merged presentation: %s" % output_path)
    return True


def _try_auto_repair_with_libreoffice(pptx_path: str) -> bool:
    """Attempt to auto-repair a PPTX by converting it through LibreOffice.

    Only runs if LibreOffice is available. Non-destructive on failure —
    the original file is left unchanged if the conversion fails.

    Args:
        pptx_path: Path to the PPTX file to repair in-place.

    Returns:
        True if LibreOffice repair succeeded, False otherwise.
    """
    import glob
    import subprocess
    import tempfile

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return False

    tmp_dir = tempfile.mkdtemp(prefix="pptx_repair_")
    try:
        result = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "pptx",
                "--outdir",
                tmp_dir,
                pptx_path,
            ],
            capture_output=True,
            timeout=180,
        )
        if result.returncode == 0:
            converted = glob.glob(os.path.join(tmp_dir, "*.pptx"))
            if converted:
                shutil.copy(converted[0], pptx_path)
                print("[MERGE] Auto-repair via LibreOffice succeeded: %s" % pptx_path)
                return True
    except Exception as e:
        print("[MERGE] Auto-repair via LibreOffice failed: %s" % e)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)
    return False


def merge_pptx_files(pptx_paths: List[str], output_path: str) -> bool:
    """Merge multiple PPTX files into a single presentation using ZIP-level manipulation.

    Uses _merge_pptx_zip_level() which copies all binary parts (images, charts,
    workbooks) at the raw bytes level, avoiding OPC package context issues that
    cause PowerPoint to report "found a problem with content" on the merged file.

    Args:
        pptx_paths: List of PPTX file paths to merge in order.
        output_path: Destination path for the merged presentation.

    Returns:
        True if merge succeeded, False otherwise.
    """
    merge_start = time.time()
    valid_paths = [p for p in pptx_paths if p and os.path.exists(p)]
    print("[MERGE] Merging %d PPTX files into %s" % (len(valid_paths), output_path))
    if VERBOSE:  # noqa: F405
        for i, p in enumerate(valid_paths):
            print("[VERBOSE][MERGE] Source %d: %s" % (i, p))
    result = _merge_pptx_zip_level(valid_paths, output_path)
    merge_elapsed = time.time() - merge_start
    print("[TIMING] merge_pptx_files completed in %.1fs" % merge_elapsed)
    return result


# === WORKFLOW STEP 5 (FINAL): MERGE ALL CHUNKS ===


def step_merge_chunks(step_input: StepInput, session_state: Dict) -> StepOutput:
    """Step 5 (Final): Merge all processed/reviewed chunk PPTX files into the final output.

    Source selection priority (explicit, robust):
    1. Template + visual review + reviewed_chunks present -> use reviewed_chunks
    2. Template + processed_chunks present             -> use processed_chunks
    3. No template (raw mode)                          -> use raw chunk_files

    Chunks are merged in order (by chunk_idx).
    """
    step_start = time.time()

    print("\n" + "=" * 60)
    print("Step 5 (Final): Merging chunks into final presentation...")
    print("=" * 60)

    output_path = session_state.get("output_path", "presentation_chunked.pptx")
    has_template = bool(session_state.get("template_path"))
    visual_review = session_state.get("visual_review", False)
    chunk_files: List[Optional[str]] = session_state.get("chunk_files", [])
    processed_chunks: Dict[int, Optional[str]] = session_state.get(
        "processed_chunks", {}
    )
    reviewed_chunks: Dict[int, Optional[str]] = session_state.get("reviewed_chunks", {})

    # Determine which chunk paths to use (priority: reviewed > processed > raw)
    if has_template and visual_review and reviewed_chunks:
        source_label = "reviewed (template + visual review)"
        ordered_paths = [reviewed_chunks.get(i) for i in sorted(reviewed_chunks.keys())]
    elif has_template and processed_chunks:
        source_label = "processed (template-assembled)"
        ordered_paths = [
            processed_chunks.get(i) for i in sorted(processed_chunks.keys())
        ]
    else:
        # No template path: use raw chunk files directly
        source_label = "raw (no template)"
        ordered_paths = [f for f in chunk_files if f is not None]

    if not ordered_paths:
        print("[MERGE] No chunk files found to merge")
        return StepOutput(
            content="No files to merge",
            success=False,
        )

    print(
        "Merging from: %s (%d total, %d valid)"
        % (
            source_label,
            len(ordered_paths),
            sum(1 for p in ordered_paths if p and os.path.exists(p)),
        )
    )

    if VERBOSE:  # noqa: F405
        print("[VERBOSE] Ordered chunk files for merge:")
        for i, p in enumerate(ordered_paths):
            print("[VERBOSE]   %d. %s" % (i, p))

    success = merge_pptx_files(
        [p for p in ordered_paths if p],
        output_path,
    )

    # Attempt optional auto-repair (only if LibreOffice is available)
    if success:
        _try_auto_repair_with_libreoffice(output_path)

    step_elapsed = time.time() - step_start
    final_file = os.path.basename(output_path)
    print(
        "[TIMING] step_merge_chunks completed in %.1fs (final: %s)"
        % (step_elapsed, final_file)
    )

    if success:
        summary = "Merged %d chunks (%s) -> %s. Duration: %.1fs" % (
            len([p for p in ordered_paths if p]),
            source_label,
            output_path,
            step_elapsed,
        )
        print("[MERGE] %s" % summary)
        return StepOutput(
            content=summary,
            success=True,
        )
    else:
        return StepOutput(
            content="Merge failed. No output file produced.",
            success=False,
        )


# === WORKFLOW BUILDER ===


def build_chunked_workflow(session_state: Dict) -> Workflow:
    """Build the chunked PPTX workflow with the appropriate set of steps.

    Steps included:
    - Step 1: Optimize & Plan   (always)
    - Step 2: Generate Chunks   (always)
    - Step 3: Process Chunks    (only when template_path is set)
    - Step 4: Visual Review     (only when template_path AND visual_review are both set)
    - Step 5: Merge Chunks      (always)

    No-template pipeline: Step 1 -> Step 2 -> Step 5
    Template pipeline:    Step 1 -> Step 2 -> Step 3 [-> Step 4] -> Step 5

    Args:
        session_state: Shared workflow state; must contain template_path and visual_review
                       to determine which optional steps to include.

    Returns:
        Configured Workflow instance with the appropriate step sequence.
    """
    has_template = bool(session_state.get("template_path"))
    do_visual_review = has_template and bool(session_state.get("visual_review"))

    steps = [
        Step(name="Optimize and Plan", executor=step_optimize_and_plan),
        Step(name="Generate Chunks", executor=step_generate_chunks),
    ]

    # Template assembly + image pipeline only runs when a template is provided
    if has_template:
        steps.append(Step(name="Process Chunks", executor=step_process_chunks))

    # Visual review requires both --visual-review flag AND a template
    if do_visual_review:
        steps.append(
            Step(name="Visual Review Chunks", executor=step_visual_review_chunks)
        )

    steps.append(Step(name="Merge Chunks", executor=step_merge_chunks))

    return Workflow(
        name="Chunked PPTX Workflow",
        steps=steps,
        session_state=session_state,
    )


# === MAIN ENTRY POINT ===


def main() -> None:
    """Parse CLI flags, build session state, and execute the chunked workflow.

    Validates required environment variables and the optional template path, creates
    output working directories under output_chunked/, runs the workflow end-to-end,
    and writes the final PPTX to the path specified by --output. Step-level progress
    and timing diagnostics are printed to stdout.
    """
    parser = argparse.ArgumentParser(
        description="Chunked PPTX generation workflow — overcomes Claude API limits for large presentations."
    )

    # Existing args (compatible with powerpoint_template_workflow.py)
    parser.add_argument(
        "--template",
        "-t",
        default=None,
        help="Path to .pptx template file (optional). Without it, skips template assembly.",
    )
    parser.add_argument(
        "--output",
        "-o",
        default="presentation_chunked.pptx",
        help="Output filename (default: presentation_chunked.pptx).",
    )
    parser.add_argument(
        "--prompt",
        "-p",
        default=None,
        help="User prompt describing the presentation topic.",
    )
    parser.add_argument(
        "--no-images",
        action="store_true",
        help="Skip AI image generation.",
    )
    parser.add_argument(
        "--no-stream",
        action="store_true",
        help="Disable streaming mode for Claude agent.",
    )
    parser.add_argument(
        "--min-images",
        type=int,
        default=1,
        help="Minimum slides that must have AI-generated images (default: 1).",
    )
    parser.add_argument(
        "--visual-review",
        action="store_true",
        help="Enable visual QA with Gemini vision per chunk (requires LibreOffice + template).",
    )
    parser.add_argument(
        "--footer-text",
        default="",
        help="Footer text for all slides (idx=11 placeholder).",
    )
    parser.add_argument(
        "--date-text",
        default="",
        help="Date text for footer date placeholder (idx=10).",
    )
    parser.add_argument(
        "--show-slide-numbers",
        action="store_true",
        help="Preserve slide number placeholder (idx=12) on all slides.",
    )
    parser.add_argument(
        "--verbose",
        "-v",
        action="store_true",
        help="Enable verbose/debug logging.",
    )

    # New args for chunked workflow
    parser.add_argument(
        "--chunk-size",
        type=int,
        default=3,
        help="Number of slides per Claude API chunk call (default: 3).",
    )
    parser.add_argument(
        "--max-retries",
        type=int,
        default=2,
        help="Max retries per chunk on failure (default: 2).",
    )
    parser.add_argument(
        "--visual-passes",
        type=int,
        default=3,
        help="Maximum visual inspection passes per chunk (default: 3).",
    )
    parser.add_argument(
        "--start-tier",
        type=int,
        choices=[1, 2, 3],
        default=1,
        help=(
            "Starting tier for chunk generation (default: 1). "
            "1=Claude PPTX skill (best quality), "
            "2=LLM code generation (80-92%% quality, faster, python-pptx native charts), "
            "3=text-only (structural, instant). "
            "Fallback continues from selected tier."
        ),
    )

    args = parser.parse_args()

    # Update module-level VERBOSE (imported from powerpoint_template_workflow via *)
    global VERBOSE  # noqa: F405
    VERBOSE = args.verbose  # noqa: F405

    # Validate API key
    if not os.getenv("ANTHROPIC_API_KEY"):
        print("Error: ANTHROPIC_API_KEY environment variable not set.")
        sys.exit(1)

    # Validate template if provided
    if args.template is not None:
        if not os.path.isfile(args.template):
            print("Error: Template file not found: %s" % args.template)
            sys.exit(1)
        if not args.template.endswith(".pptx"):
            print("Error: Template file must be a .pptx file.")
            sys.exit(1)

    # Warn when visual flags are passed without a template (they will be ignored)
    if not args.template and args.visual_review:
        print("[WARNING] --visual-review is ignored when --template is not provided")
    if not args.template and args.visual_passes != 3:
        print("[WARNING] --visual-passes is ignored when --template is not provided")

    # Effective values: visual review and passes are forced off when no template
    effective_visual_review = bool(args.visual_review) and bool(args.template)
    effective_visual_passes = args.visual_passes if args.template else 0

    # Setup output directory
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_base = os.path.join(script_dir, "output_chunked")
    os.makedirs(output_base, exist_ok=True)

    # Resolve output path
    output_path = args.output
    if not os.path.isabs(output_path):
        output_path = os.path.join(output_base, output_path)

    # Chunked workflow uses its own working directory
    output_dir = os.path.join(output_base, "chunked_workflow_work")
    os.makedirs(output_dir, exist_ok=True)

    default_prompt = "Create a professional business presentation about AI transformation in enterprise companies"

    session_state = {
        # Core paths
        "template_path": args.template or "",
        "output_path": output_path,
        "output_dir": output_dir,
        # User inputs
        "user_prompt": args.prompt or default_prompt,
        "verbose": args.verbose,
        "stream": not args.no_stream,
        "no_images": args.no_images,
        "min_images": args.min_images,
        # visual_review and visual_passes are forced False/0 when no template
        "visual_review": effective_visual_review,
        "footer_text": args.footer_text,
        "date_text": args.date_text,
        "show_slide_numbers": args.show_slide_numbers,
        # Chunked workflow settings
        "chunk_size": args.chunk_size,
        "max_retries": args.max_retries,
        "visual_passes": effective_visual_passes,
        "start_tier": args.start_tier,
        # Fields populated by steps
        "storyboard": None,
        "storyboard_dir": None,
        "total_slides": 0,
        "chunk_files": [],
        "chunk_slide_groups": [],
        "processed_chunks": {},
        "reviewed_chunks": {},
        "use_fallback_generator": False,
        # Fields used by existing step helpers
        "generated_file": "",
        "slides_data": [],
        "generated_images": {},
        "src_slide_width": 0,
        "src_slide_height": 0,
        "assembly_knowledge": {},
        "quality_report": {},
    }

    workflow = build_chunked_workflow(session_state)

    print("=" * 60)
    print("Chunked PPTX Workflow")
    print("=" * 60)
    print("Prompt:     %s" % (args.prompt or default_prompt)[:80])
    print("Output:     %s" % output_path)
    if args.template:
        print("Mode:       template-assisted generation")
        print("Template:   %s" % args.template)
        if effective_visual_review:
            print("Visual review: enabled (%d passes max)" % args.visual_passes)
        else:
            print("Visual review: disabled")
    else:
        print("Mode:       raw generation (no template)")
        print("Visual review: skipped (no template)")
    print("Chunk size: %d slides per API call" % args.chunk_size)
    print("Max retries per chunk: %d" % args.max_retries)
    print(
        "Start tier: %d (%s)"
        % (
            args.start_tier,
            {1: "Claude PPTX skill", 2: "LLM code generation", 3: "text-only"}[
                args.start_tier
            ],
        )
    )
    print("Images:     %s" % ("disabled" if args.no_images else "enabled"))
    if args.verbose:
        print("Verbose:    enabled")

    start_time = time.time()

    workflow.run()

    elapsed = time.time() - start_time
    print("\n" + "=" * 60)
    print("[TIMING] Total workflow: %.1fs" % elapsed)
    print("Output: %s" % output_path)
    print("=" * 60)


if __name__ == "__main__":
    main()
